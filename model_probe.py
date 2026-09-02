"""What the model was asked and what it answered: printed, and kept.

Every run leaves a folder holding `transcript.txt` and the slide images the
model was shown. No .pptx is ever written.

A separate file on purpose. The tool's own passes are written so that a
designer never sees a prompt: qc.llm takes the question, the schema and the
images, and hands back a dict that code re-verifies against real geometry
before it becomes a record. That is right for the product and useless when the
question is "why did it say THAT" - the prompt, the payload and the raw answer
are exactly what the passes hide.

So this script runs the real passes, through the real qc.llm, with a logging
wrapper installed over `ask_json`, and prints every call end to end:

    label -> system prompt -> user prompt -> response schema -> raw JSON answer

Two probes, in the order the tool does them.

PASS 1, LAYOUT DETECTION. What layout is this slide, and which of the master's
layouts does it belong on. Note that the shipped tool does NOT ask a model
this any more (qc.layoutpick, 31/08/2026): it ranks the master's layouts by
arithmetic and a designer picks. The probe asks BOTH - it prints the
deterministic ranking, then puts the same slide and the same layout menu to the
model, then says whether they agree. That comparison is the point: it is how
you find out whether the retired judgment was worth retiring, on a real deck,
without putting a model back in the rebuild path.

PASS 2, APPLY MASTER AND AUDIT DESIGN. Applying the master is pure code and no
model is asked (qc.applymaster via qc.prep) - the plans, the per-slide rules
and the migration are printed so the model's later answers can be read against
what it was actually looking at. Then the rebuilt deck is audited three ways:
qc.design.scan (measured, no model), qc.copilot (vision: what a designer would
adjust) and qc.components (vision: what the things on the slide ARE and which
line they belong on). Both vision passes print their prompts through the
wrapper and their verified records afterwards, so you can see the whole chain -
what the model named, and what survived code's re-measurement.

NO .pptx IS PRODUCED. The rebuilt deck stays in memory and dies with the
process. What is written is the transcript and the slide images: every run
makes a folder holding `transcript.txt` and one PNG per slide the model was
shown, because the answer and the picture are one piece of evidence. Reading
"the model missed the misaligned label" without the image it was reading is
how a rendering fault gets debugged as a prompt fault.

CONFIDENTIALITY. Those PNGs are pictures of the client's slides, and the
transcript quotes the shape inventory. Under out/ by default, which this repo
gitignores along with everything else run-derived. Do not move a run folder
somewhere that gets committed or synced.

Usage
-----
    .venv\\Scripts\\python model_probe.py --master MASTER.pptx --deck MESSY.pptx

    --slides 0,4,7      which slides pass 1 asks about (default: all of them)
    --max-slides 3      cap on slides each vision pass reviews
                        (default 0, meaning no cap)
    --pass layout|design|all
    --profile NAME      a saved profile, for the palette and the frame
    --out DIR           where the run folder goes (default: out/probe/<stamp>)
    --brief             shorten the prompts and schemas to 800 characters
                        (the answers are never shortened)

Every call costs a vision request against a dense image, and the default now
reviews the WHOLE deck: a probe that silently looked at three slides of forty
answered a question nobody asked. Narrow it with --slides and --max-slides
when a run is only about one slide.
"""

import argparse
import importlib
import io
import json
import sys
import threading
import time
from pathlib import Path

# --- output ---------------------------------------------------------------
#
# Layout names and slide text can be Arabic, and the Windows console default
# codepage cannot encode it. A probe that dies on a print in the middle of a
# transcript is worse than one that prints a replacement character.
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:                                        # pragma: no cover
    pass

WIDTH = 78
_PRINT_LOCK = threading.Lock()
_LOG_FILE = None
BRIEF = False

# The run folder, and what went into it. `_IMAGES` is keyed by the sha1 of the
# PNG bytes because that is the only handle the wrapper has: qc.copilot and
# qc.components choose their own slides internally and hand qc.llm bytes, so by
# the time an image reaches the log there is nothing on it saying which slide it
# came from. Hashing what render() wrote puts the slide number back.
RUN_DIR: Path | None = None
_IMAGES: dict = {}
_SAVED: list = []


def out(text: str = "") -> None:
    """One print, under a lock, tee'd to the run's transcript.txt. The lock
    matters because the vision passes go through a thread pool: without it two
    calls' prompts interleave line by line and neither the terminal nor the
    file is readable."""
    with _PRINT_LOCK:
        print(text, flush=True)
        if _LOG_FILE is not None:
            _LOG_FILE.write(text + "\n")
            _LOG_FILE.flush()


def _digest(png: bytes) -> str:
    import hashlib

    return hashlib.sha1(png).hexdigest()


def save_image(png: bytes, name: str) -> str:
    """Write one PNG into the run folder and return its filename.

    Written once per distinct image however many passes ask about it: a slide
    shown to the copilot and again to the component review is the same picture,
    and two copies of it on disk under two call numbers would invite reading a
    difference into them that is not there. The transcript names the file at
    every call instead.
    """
    if RUN_DIR is None:                                  # pragma: no cover
        return ""
    known = _IMAGES.get(_digest(png))
    if known:
        return known
    path = RUN_DIR / name
    path.write_bytes(png)
    _IMAGES[_digest(png)] = name
    _SAVED.append(name)
    return name


def rule(char: str = "-") -> str:
    return char * WIDTH


def banner(title: str, char: str = "=") -> None:
    out("")
    out(rule(char))
    out(f" {title}")
    out(rule(char))


def block(label: str, body: str, truncate: bool = True) -> None:
    """A labelled slab of text, indented so it cannot be mistaken for the
    probe's own commentary.

    --brief shortens the prompts and the schema, never the answer: the answer
    is the thing being read, and a transcript that cut its own subject off to
    save room would be worth less than no transcript.
    """
    out("")
    out(f"  [{label}]")
    text = body if body is not None else ""
    if truncate and BRIEF and len(text) > 800:
        text = text[:800] + f"\n... [{len(text) - 800} more characters]"
    for line in text.splitlines() or [""]:
        out(f"    {line}")


def kv(key: str, value) -> None:
    out(f"  {key:<32}{value}")


# --- the wrapper ----------------------------------------------------------

_STATE = {"n": 0, "label": "unlabelled", "calls": []}


def label(text: str) -> None:
    """Name the pass that is about to ask something. Set by each section
    rather than inferred from a stack walk: the caller knows why it is asking
    and a frame does not."""
    _STATE["label"] = text


def _logged(real):
    def ask_json(*, system: str, prompt: str, schema: dict,
                 images: list | None = None, max_tokens: int = 8192):
        _STATE["n"] += 1
        n = _STATE["n"]
        pngs = images or []
        from qc.config import LLM_MODEL

        # Named, not just counted. The image is half of what was asked, and a
        # transcript that says "images=1" cannot be read six hours later
        # against the picture that produced the answer.
        shown = [_IMAGES.get(_digest(png))
                 or save_image(png, f"call-{n:02d}_img{k}.png")
                 for k, png in enumerate(pngs, 1)]

        out("")
        out(rule("~"))
        out(f" MODEL CALL #{n}  -  {_STATE['label']}")
        out(f" model={LLM_MODEL}  images={len(pngs)}"
            f"  bytes={sum(len(p) for p in pngs):,}"
            f"  max_output_tokens={max_tokens}")
        if shown:
            out(f" the model saw: {', '.join(shown)}")
        out(rule("~"))
        block("SYSTEM PROMPT", system)
        block("USER PROMPT", prompt)
        block("RESPONSE SCHEMA", json.dumps(schema, indent=2, sort_keys=True))

        started = time.perf_counter()
        try:
            answer = real(system=system, prompt=prompt, schema=schema,
                          images=images, max_tokens=max_tokens)
        except Exception as exc:
            took = time.perf_counter() - started
            block("NO ANSWER", f"{type(exc).__name__}: {exc}",
                  truncate=False)
            out(f"    ({took:.1f}s)")
            _STATE["calls"].append({"n": n, "label": _STATE["label"],
                                    "ok": False, "seconds": took})
            raise
        took = time.perf_counter() - started
        block("RAW ANSWER (parsed JSON)",
              json.dumps(answer, indent=2, ensure_ascii=False),
              truncate=False)
        out(f"    ({took:.1f}s)")
        _STATE["calls"].append({"n": n, "label": _STATE["label"],
                                "ok": True, "seconds": took})
        return answer

    return ask_json


def install() -> None:
    """Wrap `ask_json` everywhere it is bound, and force the calls serial.

    Two bindings, not one. qc.llm owns the function, and the passes that ask
    per slide did `from .llm import ask_json` at import time - so patching only
    qc.llm leaves qc.copilot and qc.components holding the original and their
    calls print nothing at all.

    Concurrency goes to 1 for the duration. qc.llm.ask_in_parallel asks four
    slides at once, which is right in the product and wrong here: four
    overlapping transcripts cannot be read, and the wrapper's lock would only
    keep the blocks intact, not in order. The name is patched on qc.llm rather
    than qc.config because that is where ask_in_parallel reads it.
    """
    llm = importlib.import_module("qc.llm")
    wrapped = _logged(llm.ask_json)
    llm.ask_json = wrapped
    llm.LLM_CONCURRENCY = 1
    for name in ("qc.copilot", "qc.components", "qc.layoutsuggest"):
        mod = importlib.import_module(name)
        if getattr(mod, "ask_json", None) is not None:
            mod.ask_json = wrapped


# --- pass 1: what layout is this slide ------------------------------------
#
# The vocabularies below are CLOSED, and the layout name is an enum built from
# the master's own layouts. That is the house rule for every judgment pass in
# this package (qc.llm's docstring): the model names things by ids it was
# handed, never in prose of its own, so an answer either matches something real
# or is rejected by the schema before it reaches Python.

NO_FIT = "__none_of_these__"

ARCHETYPES = (
    "title", "section divider", "agenda or contents", "one column content",
    "two column comparison", "three or more column grid", "quote",
    "full bleed image", "image with caption", "chart", "table", "timeline",
    "process flow", "team or bio grid", "closing", "other",
)

LAYOUT_SYSTEM = """You are a senior presentation designer at Prezlab. You are
handed one slide from a client deck and the list of layouts that exist in the
brand master it is about to be rebuilt on. Say what the slide IS and which
layout it belongs on.

You see the rendered slide image, an inventory of its shapes (id, kind,
position and size as fractions of the slide, whether it holds text), and the
master's layout menu with what each layout can hold.

Answer four things.

1. reads: what is on this slide, in one sentence, as a designer would describe
   its structure to another designer. Not its subject matter, its shape.
2. archetype: the one entry from the closed list that fits it best.
3. layout: the EXACT name of the layout from the menu it should be rebuilt on,
   or the literal string "{no_fit}" when this master has nothing that can hold
   this slide without content being orphaned or crushed.
4. why: the evidence, in one or two sentences. Name the columns, the blocks and
   the title you are counting.

Also give runner_up (the next best layout name from the menu, or "{no_fit}")
and confidence.

Rules that matter more than a confident answer:
- The menu is the whole world. A layout name not on it is a wrong answer.
- Count what has to be PLACED, not what is decorative. A background panel and
  a page number are not content blocks.
- Columns first. A two column comparison put on a one column layout is the
  defect that sends a slide back to the designer; a block count off by one is a
  text box somebody nudges.
- "{no_fit}" is a real answer and the most useful one when it is true: it says
  the master is missing a layout, which is a fact about the master rather than
  a fault on the slide. Do not reach for the nearest thing to avoid saying it.
- Write in clear US English without em dashes."""


def layout_schema(names: list) -> dict:
    """The answer's shape, with the layout names as an enum.

    Built per master rather than declared once, because the closed vocabulary
    IS this master's layout list. `additionalProperties` is declared and
    stripped by qc.llm before it reaches the endpoint (Gemini rejects the key);
    it stays here for the same reason it stays in the shipped schemas - it is
    the probe's statement that an invented key is not an answer.
    """
    return {
        "type": "object", "additionalProperties": False,
        "required": ["reads", "archetype", "layout", "runner_up",
                     "confidence", "why"],
        "properties": {
            "reads": {"type": "string"},
            "archetype": {"type": "string", "enum": list(ARCHETYPES)},
            "layout": {"type": "string", "enum": list(names) + [NO_FIT]},
            "runner_up": {"type": "string", "enum": list(names) + [NO_FIT]},
            "confidence": {"type": "string",
                           "enum": ["high", "medium", "low"]},
            "why": {"type": "string"},
        },
    }


def layout_menu(layouts: list, slide_w: int) -> list:
    """The master's layouts as the model sees them: name, archetype token, and
    what the layout can hold.

    Read with the tool's OWN readers (qc.layoutgap.layout_signature and
    qc.layoutpick._offers_sentence) rather than a second description invented
    here. Like for like is the only way the comparison at the end of the
    section means anything: the model and the ranking are looking at the same
    reading of the same layout.
    """
    from qc.layoutgap import layout_signature
    from qc.layoutpick import _offers_sentence

    menu = []
    for entry in layouts:
        name = entry.get("name")
        if not name:
            continue
        lay = layout_signature(entry, slide_w)
        menu.append({"name": name, "archetype": entry.get("type") or "",
                     "holds": _offers_sentence(lay)})
    return menu


def probe_layouts(deck_bytes: bytes, master_bytes: bytes,
                  wanted: list) -> None:
    """Pass 1, per slide: the deterministic ranking, the model's answer, and
    whether they agree."""
    from pptx import Presentation

    from qc.applymaster import plan_assignments
    from qc.copilot import inventory
    from qc.layoutgap import describe
    from qc.layoutpick import rank
    from qc.llm import LLMUnavailable, ask_json
    from qc.prep import read_master

    banner("PASS 1 - LAYOUT DETECTION")

    _, _, layouts, _space, master_size = read_master(master_bytes)
    deck = Presentation(io.BytesIO(deck_bytes))
    slide_w = int(deck.slide_width or 0) or 1
    slide_h = int(deck.slide_height or 0) or 1
    slides = list(deck.slides)
    plans = {p.slide_index: p for p in plan_assignments(deck, layouts)}
    menu = layout_menu(layouts, slide_w)

    kv("master layouts", f"{len(menu)}")
    kv("deck slides", f"{len(slides)}")
    kv("slide size (EMU)", f"{slide_w} x {slide_h}  "
                           f"(master states {master_size[0]} x {master_size[1]})")
    block("THE MASTER'S LAYOUT MENU (as handed to the model)",
          "\n".join(f"{m['name']}   [{m['archetype'] or 'no archetype'}]  "
                    f"-  {m['holds']}" for m in menu))

    # Rendered once, for every slide the probe will ask about. The vision call
    # is the expensive half; the render is the slow half, and doing them in one
    # batch keeps PowerPoint being started once rather than per slide.
    pngs = render(deck_bytes, wanted, "upload")

    agreements = []
    for idx in wanted:
        if idx >= len(slides):
            out(f"\n  slide {idx + 1}: not in this deck, skipped")
            continue
        slide = slides[idx]
        plan = plans.get(idx)
        banner(f"SLIDE {idx + 1} of {len(slides)}", "-")

        candidates, sig = rank(
            slide, layouts, slide_w, slide_h,
            source_type=getattr(plan, "source_type", None),
            source_name=getattr(plan, "source_layout", "") or "")

        kv("[code] slide holds", describe(sig))
        kv("[code] signature", json.dumps(sig, sort_keys=True))
        if plan is not None:
            kv("[code] the file's own pick",
               f"{plan.target_layout!r} by {plan.match_rule}"
               + (f"  ({plan.note})" if plan.note else ""))
        out("  [code] layoutpick ranking (what the designer would see):")
        for i, cand in enumerate(candidates[:5], 1):
            out(f"      {i}. {cand.name!r}  score={cand.score:.2f}  "
                f"fits={'yes' if cand.fits else 'no'}")
            out(f"         holds {cand.offers}"
                + (f"; {cand.why}" if cand.why else ""))
        code_pick = candidates[0].name if candidates else None

        png = pngs.get(idx)
        if png is None:
            out("\n  [model] not asked: this slide did not render, and the "
                "question is about a picture.")
            continue

        inv = inventory(slide, slide_w, slide_h)
        label(f"pass 1 layout detection, slide {idx + 1}")
        prompt = (
            "The master's layout menu:\n"
            + json.dumps(menu, indent=2, ensure_ascii=False)
            + "\n\nThis slide's shapes:\n"
            + json.dumps(inv, sort_keys=True)
            + "\n\nHow this file's own code reads the slide, for reference "
              "only - disagree with it if the picture says otherwise: "
            + describe(sig) + ".")
        try:
            answer = ask_json(system=LAYOUT_SYSTEM.format(no_fit=NO_FIT),
                              prompt=prompt,
                              schema=layout_schema([m["name"] for m in menu]),
                              images=[png])
        except LLMUnavailable as exc:
            out(f"\n  [model] could not be asked: {exc}")
            continue

        model_pick = answer.get("layout")
        out("")
        kv("[model] reads it as", answer.get("reads"))
        kv("[model] archetype", answer.get("archetype"))
        kv("[model] layout", f"{model_pick!r} "
                             f"(confidence {answer.get('confidence')})")
        kv("[model] runner up", repr(answer.get("runner_up")))
        kv("[model] why", answer.get("why"))

        verdict = ("AGREE" if model_pick == code_pick
                   else "DISAGREE" if model_pick != NO_FIT
                   else "MODEL SAYS THE MASTER HAS NO HOME FOR THIS SLIDE")
        out("")
        kv("[compare]", f"model={model_pick!r}  ranking={code_pick!r}  "
                        f"file={getattr(plan, 'target_layout', None)!r}"
                        f"  ->  {verdict}")
        agreements.append((idx, verdict))

    if agreements:
        agreed = sum(1 for _, v in agreements if v == "AGREE")
        out("")
        kv("pass 1 summary", f"{agreed} of {len(agreements)} slides: the model "
                             f"picked what the ranking put first")


# --- pass 2: apply the master, then audit the design ----------------------


def render(deck_bytes: bytes, indices: list, tag: str) -> dict:
    """{slide_index: png}, every one of them also written into the run folder.

    Empty, with a printed reason, when no renderer answers - the vision passes
    are then skipped rather than fed nothing.

    `tag` names the deck the slides came from and goes in the filename, because
    pass 1 renders THE UPLOAD and pass 2 renders THE REBUILD. Same slide
    numbers, two different pictures, and telling them apart is most of what a
    designer wants from these files.
    """
    from qc.render import export_decks_png

    if not indices:
        return {}
    try:
        images = export_decks_png({tag: deck_bytes}, list(indices), width=1100)
    except Exception as exc:
        out(f"\n  !! slides did not render ({type(exc).__name__}: {exc}). "
            f"Every pass that needs a picture is skipped.")
        return {}
    pngs = {int(k.split(":", 1)[1]): v for k, v in images.items()}
    for idx in sorted(pngs):
        save_image(pngs[idx], f"{tag}_slide-{idx + 1:02d}.png")
    return pngs


def apply_master(deck_bytes: bytes, master_bytes: bytes, filename: str):
    """Plan and rebuild, printing every decision. Returns (deck_bytes, prep) -
    the ORIGINAL bytes when the rebuild could not run, so the audit below has
    something to read either way.

    No model is asked in this section and that is worth seeing in a transcript
    about model behaviour: the layout each slide lands on, the frame written
    into it and the content migration are all arithmetic (qc.prep.build), so
    the same two files rebuild the same way every time.
    """
    from qc.prep import PrepError, build, plan
    from qc.unify import com_available

    banner("PASS 2a - APPLY THE MASTER (pure code, no model asked)")
    kv("PowerPoint COM here", "yes" if com_available() else "NO")

    prepared = plan(deck_bytes, filename, master_bytes)
    kv("slides planned", prepared.slides)
    kv("needing a decision", prepared.undecided)
    out("")
    out("  per-slide plan (slide: source layout -> target, how it was chosen):")
    for p in prepared.plans:
        out(f"      {p.slide_index + 1:>3}: {p.source_layout!r}"
            f" [{p.source_type or '-'}]  ->  {p.target_layout!r}"
            f"  by {p.match_rule}"
            + (f"  ({p.note})" if p.note else ""))

    if prepared.choices:
        out("")
        out("  the slides qc.layoutpick would put in front of a designer:")
        for choice in prepared.choices:
            out(f"      slide {choice.slide_index + 1}: {choice.reason}")
            out(f"         it holds {choice.wants}; suggested "
                f"{choice.suggested!r}")
        out("")
        out("  (nothing is picked here - the probe lets every suggestion "
            "stand, which is what run() does with an untouched plan)")

    try:
        prep = build(prepared, master_bytes)
    except PrepError as exc:
        out("")
        out(f"  !! the master could not be applied: {exc}")
        out("  The design audit below therefore reads THE UPLOAD, not a "
            "rebuild. Every finding it reports may be one the master would "
            "have reset.")
        return deck_bytes, None

    out("")
    kv("slides rebuilt", f"{prep.applied} of {prep.slides}")
    kv("slides that failed", len(prep.errors) or "none")
    for idx, err in (prep.errors or {}).items():
        out(f"      slide {idx + 1}: {err}")
    kv("slide masters in output", prep.masters)
    if prep.stragglers:
        kv("still on the old design", prep.stragglers)
    for note in prep.space_notes or []:
        out(f"      frame: {note}")
    if prep.changes:
        out("")
        out(f"  content migration ({len(prep.changes)} changes):")
        for change in prep.changes[:40]:
            out(f"      slide {getattr(change, 'slide_index', 0) + 1}: "
                f"{getattr(change, 'action', '')} - "
                f"{getattr(change, 'detail', '')}")
        if len(prep.changes) > 40:
            out(f"      ... {len(prep.changes) - 40} more")
    if prep.coverage is not None:
        from qc.layoutgap import headline

        out("")
        block("LAYOUT COVERAGE OF THE MASTER", headline(prep.coverage))
    return prep.deck or deck_bytes, prep


def audit_design(deck_bytes: bytes, palette_cfg: dict, space,
                 max_slides: int, tag: str = "rebuilt") -> None:
    """The three audits of the rebuilt deck: measured, then the two vision
    passes, each printed with its prompts and with what survived code's
    re-measurement.

    `tag` is what the saved slide images are named after, and it is "upload"
    rather than "rebuilt" when the master could not be applied. A folder of
    PNGs labelled as a rebuild that is actually the untouched upload is worse
    than no PNGs.
    """
    from qc.design import scan, summary

    banner("PASS 2b - DESIGN AUDIT, MEASURED (no model asked)")
    findings = scan(deck_bytes, palette_cfg)
    kv("findings", json.dumps(summary(findings), sort_keys=True))
    for f in findings:
        out("")
        out(f"      [{f.severity}] {f.kind}: {f.headline}")
        out(f"         {f.detail}")
        out(f"         slides {[i + 1 for i in f.slides]}, "
            f"{f.places} place(s), id {f.finding_id}")
        for opt in f.options or []:
            out(f"         option {opt.remedy_id}: {opt.label}"
                f"  [{opt.op or 'leave it'}]  {opt.note}")

    # The two vision passes want a manifest to de-duplicate against. A real run
    # hands them the audit's records so a vision finding does not restate a
    # measured one; here the deck has not been through qc.engine, so they are
    # told there is nothing yet and every judgment they make is printed.
    from pptx import Presentation

    n = len(Presentation(io.BytesIO(deck_bytes)).slides)
    manifest = {"records": [], "slides": n}
    wanted = vision_slides(deck_bytes, max_slides)
    # The passes read their own budget off a module global, and 0 there means
    # "review nothing" rather than "no cap" - so an uncapped run is handed the
    # length of the list it is being given, not the flag.
    budget = max(max_slides or len(wanted), 1)
    kv("slides in the deck", n)
    kv("slides worth a call",
       f"{[i + 1 for i in wanted]} ("
       + (f"the first {max_slides}" if max_slides else "every slide")
       + " with three or more shapes)")
    thumbs = render(deck_bytes, wanted, tag) if wanted else {}
    if not wanted:
        out("  No slide in this deck has three shapes on it, so neither "
            "vision pass has anything to ask about.")

    banner("PASS 2c - DESIGN COPILOT (vision: what a designer would adjust)")
    if not thumbs:
        out("  skipped: no rendered slides.")
    else:
        import qc.copilot as copilot

        copilot.MAX_SLIDES = budget
        label("pass 2c design copilot")
        records, reviewed = copilot.run_copilot(deck_bytes, thumbs, manifest)
        out("")
        kv("slides reviewed", reviewed)
        kv("records after verification", len(records))
        print_records(records)

    banner("PASS 2d - COMPONENT REVIEW (vision: what the things ARE)")
    if not thumbs:
        out("  skipped: no rendered slides.")
    else:
        import qc.components as components

        components.MAX_SLIDES = budget
        kv("frame handed to the model",
           "the master's presentation space" if space is not None
           else "none stated, so component-to-component lines only")
        label("pass 2d component review")
        records, reviewed = components.run_components(deck_bytes, thumbs,
                                                      manifest, space)
        out("")
        kv("slides reviewed", reviewed)
        kv("records after verification", len(records))
        print_records(records)


def vision_slides(deck_bytes: bytes, cap: int) -> list:
    """The slides a vision pass would actually ask about. A `cap` of 0 (the
    default) means every one of them.

    Not "the first `cap` slides", which is what this did on its first run and
    why a budget of one slide bought zero calls: both passes skip a slide with
    fewer than three shapes on it (a cover, a divider), so rendering blind
    spends the whole budget on the slides neither of them will ask about. The
    threshold is read from the passes rather than restated - three is their
    number, and a probe that guessed its own would report a pass that never
    ran as a pass that found nothing.
    """
    from pptx import Presentation

    from qc.copilot import inventory

    prs = Presentation(io.BytesIO(deck_bytes))
    slide_w, slide_h = prs.slide_width, prs.slide_height
    picked = []
    for idx, slide in enumerate(prs.slides):
        if len(inventory(slide, slide_w, slide_h)) >= 3:
            picked.append(idx)
        if cap and len(picked) >= cap:
            break
    return picked


def print_records(records: list) -> None:
    """What the model's judgment turned into once geometry had its say.

    The gap between an answer and these rows is the whole reason the contract
    is worth reading in a transcript: an observation that does not check out
    against the file is dropped silently in the product, and here you can see
    exactly which ones went (qc.copilot.synthesize, qc.components.synthesize).
    """
    if not records:
        out("      nothing survived verification.")
        return
    for r in records:
        out("")
        out(f"      slide {r['slide_index'] + 1} shape {r['shape_id']}: "
            f"{r['issue_type']}  [{r['severity']}/{r['confidence']}"
            f"/{r.get('source')}]")
        out(f"         {r['message']}")
        if r.get("property"):
            out(f"         {r['property']}: {r.get('old_value')} -> "
                f"{r.get('new_value')}")


def frame_of(deck_bytes: bytes, profile_obj):
    """(left, top, right, bottom) EMU, or None. The profile's stated rectangle
    first and the deck's own master second, which is the order the audit uses
    (qc.modules.margin_alignment._space_box) - two readings of one frame is two
    numbers to disagree about."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(deck_bytes))
    if profile_obj is not None:
        try:
            from qc.modules.margin_alignment import _space_box

            return _space_box(profile_obj, prs)
        except Exception:
            pass
    try:
        from qc.stylespec import dominant_master, read_presentation_space

        space = read_presentation_space(prs, dominant_master(prs))
    except Exception:
        return None
    if not space or space.get("problem") or not space.get("box_emu"):
        return None
    return tuple(space["box_emu"])


# --- main -----------------------------------------------------------------


def parse_slides(text: str, total: int) -> list:
    """The slides pass 1 asks about: every slide in the deck unless --slides
    names some.

    The default used to be the first three, which reads as a broken flag
    rather than as a budget - a forty-slide deck reported on three of them
    and nothing on the command line said three. `total` is the deck's own
    slide count, so an empty --slides means the whole deck at whatever length
    it happens to be."""
    if not text or not text.strip():
        return list(range(total))
    out_ = []
    for bit in text.replace(" ", "").split(","):
        if not bit:
            continue
        if "-" in bit:
            lo, _, hi = bit.partition("-")
            out_.extend(range(int(lo), int(hi) + 1))
        else:
            out_.append(int(bit))
    return sorted(set(out_))


def main(argv=None) -> int:
    global _LOG_FILE, BRIEF, RUN_DIR

    ap = argparse.ArgumentParser(
        description="Print the prompts and the raw answers of every model "
                    "call this tool makes. Writes no .pptx.")
    ap.add_argument("--master", required=True, help="the brand master .pptx")
    ap.add_argument("--deck", required=True, help="the messy client .pptx")
    ap.add_argument("--slides", default="",
                    help="slides pass 1 asks about, 0-based: 0,4,7 or 0-2 "
                         "(default: every slide in the deck)")
    ap.add_argument("--max-slides", type=int, default=0,
                    help="cap on slides each vision pass reviews "
                         "(default 0, meaning no cap)")
    ap.add_argument("--pass", dest="which", default="all",
                    choices=("layout", "design", "all"))
    ap.add_argument("--profile", default="",
                    help="a saved profile name, for the palette and the frame")
    ap.add_argument("--out", default="",
                    help="where the run folder goes (default: "
                         "out/probe/<timestamp>)")
    ap.add_argument("--brief", action="store_true",
                    help="shorten the prompts and schemas to 800 characters; "
                         "the answers are never shortened")
    args = ap.parse_args(argv)

    BRIEF = args.brief

    master_path, deck_path = Path(args.master), Path(args.deck)
    for path in (master_path, deck_path):
        if not path.exists():
            print(f"no such file: {path}")
            return 2

    # A folder per run, named for when it ran, ALWAYS - there is no flag to
    # turn it off. A probe whose output has to be re-run to be re-read is a
    # probe that gets run twice, and the second run is against a model that
    # has moved on. Under out/, which this repo gitignores: these files are
    # pictures of a client's slides.
    RUN_DIR = Path(args.out) if args.out else (
        Path("out") / "probe" / time.strftime("%Y-%m-%d_%H%M%S"))
    RUN_DIR.mkdir(parents=True, exist_ok=True)
    _LOG_FILE = open(RUN_DIR / "transcript.txt", "w", encoding="utf-8")
    master_bytes = master_path.read_bytes()
    deck_bytes = deck_path.read_bytes()

    from qc.config import (AI_ENABLED, LLM_CONCURRENCY, LLM_MODEL, RENDERER)
    from qc.llm import api_configured, configuration_note

    banner("MODEL PROBE", "=")
    kv("master", f"{master_path.name}  ({len(master_bytes):,} bytes)")
    kv("deck", f"{deck_path.name}  ({len(deck_bytes):,} bytes)")
    kv("model", LLM_MODEL)
    kv("key present", "yes" if api_configured() else "NO")
    kv("settings note", configuration_note() or "none")
    kv("QC_AI", "on" if AI_ENABLED else "off (the web routes would refuse; "
                                        "this probe asks anyway)")
    kv("renderer", RENDERER)
    kv("concurrency", f"{LLM_CONCURRENCY} in the product, forced to 1 here so "
                      f"the transcript stays in order")
    kv("writing to", f"{RUN_DIR}  (transcript.txt plus one PNG per slide "
                     f"the model is shown)")
    if not api_configured():
        out("")
        out("  Set GEMINI_API_KEY in the .env at the project root. Nothing "
            "below can be asked without it.")
        return 1

    install()
    profile_obj = None
    palette_cfg = {}
    if args.profile:
        from qc.profile import Profile

        try:
            profile_obj = Profile.load(args.profile)
            palette_cfg = profile_obj.module_config("color_palette")
            kv("profile", args.profile)
        except Exception as exc:
            out(f"  profile {args.profile!r} did not load "
                f"({type(exc).__name__}: {exc}); carrying on without one")

    started = time.perf_counter()
    if args.which in ("layout", "all"):
        from pptx import Presentation

        n_slides = len(Presentation(io.BytesIO(deck_bytes)).slides)
        probe_layouts(deck_bytes, master_bytes,
                      parse_slides(args.slides, n_slides))

    if args.which in ("design", "all"):
        rebuilt, prep = apply_master(deck_bytes, master_bytes, deck_path.name)
        audit_design(rebuilt, palette_cfg, frame_of(rebuilt, profile_obj),
                     args.max_slides,
                     tag="rebuilt" if prep is not None else "upload")

    banner("EVERY CALL THIS RUN MADE", "=")
    for call in _STATE["calls"]:
        out(f"  #{call['n']:<3} {'ok  ' if call['ok'] else 'FAIL'}  "
            f"{call['seconds']:>6.1f}s  {call['label']}")
    ok = sum(1 for c in _STATE["calls"] if c["ok"])
    out("")
    kv("calls", f"{len(_STATE['calls'])} ({ok} answered, "
                f"{len(_STATE['calls']) - ok} failed)")
    kv("model time", f"{sum(c['seconds'] for c in _STATE['calls']):.1f}s")
    kv("wall clock", f"{time.perf_counter() - started:.1f}s")
    out("")
    kv("run folder", RUN_DIR)
    kv("transcript", "transcript.txt")
    kv("slide images", f"{len(_SAVED)}"
                       + (f": {', '.join(_SAVED)}" if _SAVED else
                          " (nothing rendered, so the model saw no pictures)"))
    kv("decks written", "none - no .pptx is produced by this script")
    if _LOG_FILE is not None:
        _LOG_FILE.close()
    return 0


if __name__ == "__main__":
    sys.exit(main())
