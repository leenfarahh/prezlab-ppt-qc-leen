"""Annotated deck export: writes findings and comments into speaker notes.

build_annotated() takes the original deck bytes plus the audit manifest and
inline comments, and returns a new deck whose flagged slides carry a QC
block in their speaker notes. The block is delimited by MARKER so a second
annotation pass replaces it instead of stacking duplicates.
"""

import io
from collections import defaultdict

from pptx import Presentation

MARKER = "--- Prezlab QC ---"
_SEVERITY_ORDER = {"error": 0, "warning": 1, "info": 2}
_MESSAGE_LIMIT = 200


def _qc_block(records: list[dict], slide_comments: list[dict]) -> str:
    lines = [MARKER]
    for rec in sorted(records,
                      key=lambda r: _SEVERITY_ORDER.get(r.get("severity"), 3)):
        message = (rec.get("message") or "")[:_MESSAGE_LIMIT]
        lines.append(f"[{rec['severity'].upper()}] {rec['issue_type']}: {message}")
    if slide_comments:
        lines.append("Comments:")
        for c in slide_comments:
            lines.append(f"- {c['author']}: {c['text']}")
    return "\n".join(lines)


def build_annotated(deck_bytes: bytes, manifest: dict, comments: list[dict]) -> bytes:
    prs = Presentation(io.BytesIO(deck_bytes))

    findings: dict[int, list[dict]] = defaultdict(list)
    for rec in manifest.get("records", []):
        # Preflight records are engine bookkeeping, not designer-facing issues.
        if rec.get("module") == "preflight":
            continue
        findings[rec["slide_index"]].append(rec)

    slide_comments: dict[int, list[dict]] = defaultdict(list)
    for c in comments:
        slide_comments[c["slide_index"]].append(c)

    for s_idx, slide in enumerate(prs.slides):
        recs = findings.get(s_idx, [])
        notes = slide_comments.get(s_idx, [])
        if not recs and not notes:
            # Do not touch notes_slide here: accessing it would create an
            # empty notes part on slides we have nothing to say about.
            continue

        tf = slide.notes_slide.notes_text_frame
        # Everything from the marker onward is ours from a previous run;
        # keep only the user's pre-marker text so re-annotation is idempotent.
        base = tf.text.split(MARKER, 1)[0].rstrip()
        block = _qc_block(recs, notes)
        tf.text = f"{base}\n{block}" if base else block

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
