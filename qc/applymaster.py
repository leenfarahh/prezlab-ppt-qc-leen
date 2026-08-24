"""Apply a submitted master's layouts to every slide of a content deck.

    copy slide -> apply the master's layout to the COPY -> delete the original

one slide at a time, in place, never batching. That ordering is deliberate and
was specified by the design lead: duplicating the whole deck first and then
restyling the copies leaves the deck at 2n slides mid-run, so any failure
halfway through strands a half-branded double-length file that a designer then
has to untangle by hand. Doing it per slide keeps the deck at n slides at every
moment, and a failure leaves the remaining slides simply untouched.

Verified against desktop PowerPoint (19/08/2026):

    designs = pres.Designs.Load(master_path)   adds the foreign master WITHOUT
                                               restyling existing slides
    dup = pres.Slides(i).Duplicate()(1)        the copy lands at index i+1
    dup.CustomLayout = target                  PowerPoint's own placeholder
                                               matching runs here
    pres.Slides(i).Delete()                    the copy shifts back into slot i

After each iteration slot i holds the finished slide and slot i+1 holds the
next untouched original, so a plain 1..n walk is correct and the slide count
never changes.

Why PowerPoint and not package surgery: assigning CustomLayout runs
PowerPoint's placeholder-matching engine, which is what actually moves a
slide's content into the new layout's placeholders. Rewriting the OOXML
directly would repoint the layout relationship without moving any content,
which looks right in the XML and wrong on screen. That means this is a
Windows-plus-PowerPoint operation; there is no cloud path for it yet, and the
caller gets a clear refusal rather than a silent no-op.
"""

import io
import os
import tempfile
from dataclasses import dataclass, field

from pptx import Presentation

from .unify import (com_available, com_failure_advice, force_quit,
                    sweep_automation)

# Archetypes that can stand in for an unmatched slide, best first. A content
# layout is the safest home for a slide whose own archetype has no counterpart:
# it has both a title and a body placeholder, so the least content is orphaned.
_FALLBACK_TYPES = ("obj", "titleOnly", "blank")


@dataclass
class SlidePlan:
    slide_index: int          # zero-based
    source_layout: str
    source_type: str | None
    target_layout: str | None
    match_rule: str           # how the target was chosen, for the report
    note: str = ""


@dataclass
class ApplyResult:
    deck: bytes | None
    plans: list[SlidePlan]
    errors: dict[int, str] = field(default_factory=dict)
    fatal: str | None = None
    # Slides still on a design other than the one applied, and how many slide
    # masters the output therefore carries.
    #
    # This is the consequence nobody could see: a slide that cannot be rebuilt
    # keeps the deck's ORIGINAL design alive, so the output has two masters, and
    # PowerPoint's master view lists the original first. A designer opens it,
    # sees a master without the new guides or the presentation-space rectangle,
    # and reads that as "the master was not copied" - when it was, onto the
    # other one (design lead, 21/08/2026).
    stragglers: list = field(default_factory=list)
    masters: int = 1

    @property
    def applied(self) -> int:
        """Slides actually rebuilt. Zero on a fatal, where the plans exist but
        nothing ever ran: counting planned-but-unrun slides as applied would
        report a successful format for a deck that was never touched."""
        if self.fatal or self.deck is None:
            return 0
        return sum(1 for p in self.plans
                   if p.target_layout and p.slide_index not in self.errors)


# PowerPoint rejects calls while it is busy (RPC_E_CALL_REJECTED / SERVERCALL_
# RETRYLATER), which is transient and not a real failure. Observed in a test
# run on 19/08/2026: a call failed once and the identical run passed straight
# after. Left unhandled, a designer silently loses a slide to a race.
_TRANSIENT_HRESULTS = (-2147418111, -2147417846)  # call rejected, server busy
_COM_ATTEMPTS = 4
_COM_BACKOFF_SEC = 0.4


def _is_transient(exc) -> bool:
    args = getattr(exc, "args", ())
    return bool(args) and args[0] in _TRANSIENT_HRESULTS


def _com_retry(fn, attempts: int = _COM_ATTEMPTS):
    """Run a COM call, retrying only the transient busy/rejected errors. Any
    other failure raises immediately: retrying a real error just delays it."""
    import time

    last = None
    for attempt in range(attempts):
        try:
            return fn()
        except Exception as exc:
            if not _is_transient(exc):
                raise
            last = exc
            time.sleep(_COM_BACKOFF_SEC * (attempt + 1))
    raise last


def _norm(name: str) -> str:
    return " ".join((name or "").split()).casefold()


def plan_assignments(deck_prs, target_layouts: list[dict]) -> list[SlidePlan]:
    """Choose a target layout for every slide, and record WHY.

    Matching runs name first, then archetype, because a designer who named a
    layout "Section Header" in both files meant those to correspond, and that
    intent outranks any structural guess. Archetype is the second rule because
    the OOXML type token (secHead, twoObj, picTx) is the master's own statement
    of what a layout is for. Only then does it fall back, and a fallback is
    always labelled as one so nobody reads a guess as a match."""
    by_name = {_norm(l["name"]): l for l in target_layouts}
    by_type: dict[str, list[dict]] = {}
    for lay in target_layouts:
        if lay.get("type"):
            by_type.setdefault(lay["type"], []).append(lay)

    fallback = None
    for want in _FALLBACK_TYPES:
        if by_type.get(want):
            fallback = by_type[want][0]
            break
    if fallback is None and target_layouts:
        fallback = target_layouts[0]

    plans = []
    for idx, slide in enumerate(deck_prs.slides):
        try:
            layout = slide.slide_layout
            src_name, src_type = layout.name, layout._element.get("type")
        except Exception:
            src_name, src_type = "(unreadable)", None

        hit = by_name.get(_norm(src_name))
        if hit is not None:
            rule, note = "name", ""
        elif src_type and by_type.get(src_type):
            candidates = by_type[src_type]
            if len(candidates) == 1:
                hit, rule, note = candidates[0], "archetype", ""
            else:
                # Several layouts share the archetype; prefer the one whose
                # placeholder count is closest, so a two-content slide does not
                # land on a layout with one box.
                want = len(layout.placeholders)
                hit = min(candidates,
                          key=lambda c: abs(len(c["placeholders"]) - want))
                rule, note = "archetype", f"{len(candidates)} layouts share this type"
        else:
            hit, rule = fallback, "fallback"
            note = (f"no layout named '{src_name}'"
                    + (f" and none of type '{src_type}'" if src_type else "")
                    + "; content may be orphaned, check this slide")

        plans.append(SlidePlan(
            slide_index=idx, source_layout=src_name, source_type=src_type,
            target_layout=hit["name"] if hit else None,
            match_rule=rule if hit else "none",
            note=note if hit else "the master defines no usable layout"))
    return plans


def apply_master(deck_bytes: bytes, master_bytes: bytes,
                 plans: list[SlidePlan]) -> ApplyResult:
    """Run the per-slide copy/apply/delete loop over desktop PowerPoint."""
    if not com_available():
        return ApplyResult(
            deck=None, plans=plans,
            fatal="Applying a master needs desktop PowerPoint on this machine. "
                  "Run it on the Windows box; there is no cloud path for this "
                  "step yet.")

    import pythoncom
    import win32com.client

    from .render import _RENDER_LOCK  # PowerPoint is single-instance

    errors: dict[int, str] = {}
    with _RENDER_LOCK:
        pythoncom.CoInitialize()
        paths = []
        app = pres = None
        try:
            for blob, suffix in ((deck_bytes, "_deck.pptx"),
                                 (master_bytes, "_master.pptx")):
                fd, p = tempfile.mkstemp(suffix=suffix)
                os.close(fd)
                with open(p, "wb") as f:
                    f.write(blob)
                paths.append(p)
            deck_path, master_path = paths
            out_fd, out_path = tempfile.mkstemp(suffix="_branded.pptx")
            os.close(out_fd)
            os.remove(out_path)  # SaveAs will not overwrite silently
            paths.append(out_path)

            # Leftovers from an interrupted run go first: DispatchEx attaches to
            # a running /AUTOMATION instance rather than starting a clean one,
            # and the snapshot below cannot see one that predates this run
            # (qc.unify.sweep_automation).
            # What it hands back is what survived, which is exactly the
            # snapshot the teardown needs (qc.unify.force_quit).
            started = sweep_automation()
            try:
                app = _com_retry(
                    lambda: win32com.client.DispatchEx("PowerPoint.Application"))
                app.DisplayAlerts = 1
            except Exception as exc:
                return ApplyResult(None, plans,
                                   fatal=com_failure_advice(exc))

            pres = app.Presentations.Open(deck_path, 0, 0, 0)
            try:
                design = pres.Designs.Load(master_path)
            except Exception as exc:
                return ApplyResult(None, plans,
                                   fatal=f"Could not load the master's design: {exc}")

            master = design.SlideMaster
            layouts = {master.CustomLayouts(i).Name: master.CustomLayouts(i)
                       for i in range(1, master.CustomLayouts.Count + 1)}

            for plan in plans:
                i = plan.slide_index + 1  # COM is 1-based
                if i > pres.Slides.Count:
                    errors[plan.slide_index] = "slide index out of range"
                    continue
                target = layouts.get(plan.target_layout)
                if target is None:
                    errors[plan.slide_index] = (
                        f"master has no layout named '{plan.target_layout}'")
                    continue
                try:
                    # copy -> apply to the copy -> delete the original, so the
                    # deck is never longer than it started for more than one
                    # statement at a time
                    copy = _com_retry(lambda: pres.Slides(i).Duplicate()(1))
                    _com_retry(lambda: setattr(copy, "CustomLayout", target))
                    _com_retry(lambda: pres.Slides(i).Delete())
                except Exception as exc:
                    errors[plan.slide_index] = f"apply failed: {exc}"
                    # Leave the deck consistent: if the duplicate survived the
                    # failure, drop it rather than shipping a doubled slide.
                    try:
                        if pres.Slides.Count > len(plans):
                            pres.Slides(i + 1).Delete()
                    except Exception:
                        pass

            # Which slides did NOT end up on the applied design. Read before
            # the cleanup below, since deleting a design shifts the indexes.
            stragglers = []
            try:
                applied_index = design.Index
                stragglers = [i - 1 for i in range(1, pres.Slides.Count + 1)
                              if pres.Slides(i).Design.Index != applied_index]
            except Exception:
                stragglers = []

            # designs left carrying no slides are dead weight in the file
            try:
                used = {pres.Slides(i).Design.Index
                        for i in range(1, pres.Slides.Count + 1)}
                for i in range(pres.Designs.Count, 0, -1):
                    if i not in used:
                        try:
                            pres.Designs(i).Delete()
                        except Exception:
                            pass
            except Exception:
                pass

            try:
                n_masters = int(pres.Designs.Count)
            except Exception:
                n_masters = 1

            pres.SaveAs(out_path)
            pres.Close()
            pres = None
            with open(out_path, "rb") as f:
                out = f.read()
            return ApplyResult(deck=out, plans=plans, errors=errors,
                               stragglers=stragglers, masters=n_masters)
        finally:
            try:
                if pres is not None:
                    pres.Close()
            except Exception:
                pass
            # Quit, and make sure: an instance wedged mid-run survives Quit with
            # no window, and every later automation attempt on the host fails
            # against it (qc.unify.force_quit).
            if app is not None:
                force_quit(app, started)
            pythoncom.CoUninitialize()
            for p in paths:
                try:
                    os.remove(p)
                except OSError:
                    pass


def apply_master_to_deck(deck_bytes: bytes, master_bytes: bytes,
                         target_layouts: list[dict]) -> ApplyResult:
    """Plan, then run. target_layouts is the Style Spec's `layouts` list."""
    plans = plan_assignments(Presentation(io.BytesIO(deck_bytes)), target_layouts)
    return apply_master(deck_bytes, master_bytes, plans)
