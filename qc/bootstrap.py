"""Bootstrap a formatting profile FROM a reference deck.

    python -m qc.bootstrap "deck.pptx" --id client_x --name "Client X (from deck)"

Reads the deck's own dominant conventions (fonts per role, complex-script
typefaces, palette, layouts, canvas-proportional safe zones) and writes a
profile JSON that treats those conventions as the rules. The output is a
STARTING POINT for the design lead to review, not a final brand definition;
provenance is stamped into the profile name.

Bootstrap choices, deliberately conservative:
- A latin/cs family is "allowed" when it covers at least 5% of runs in its
  role (minimum 2 runs), so one-off pasted fonts still get flagged.
- Title/subtitle size targets come from the modal placeholder size; the body
  role gets NO size target, because v1 maps all non-placeholder text to
  "body" and real decks legitimately vary body text sizes (a single target
  would flood the report with noise).
- Palette = solid colors (shape fills + text colors, theme refs resolved)
  seen at least 3 times, capped at 12.
- Layout allowlist = layouts used by at least 2 slides.

Two blocks describe what the deck's DESIGN SYSTEM declares, as opposed to
what its slides happen to do. Both are read by qc/stylespec.py, which owns
the design surface; this module just carries them into the profile:
- config.theme    the theme part's own truth: role-mapped scheme colors, the
                  master's colour map, and the major/minor font per script.
                  Read from theme1.xml, never inferred from usage, so it
                  survives a deck whose slides all override the theme.
- config.brand    the recurring brand mark (logo) and where it sits.

Both are absent from the hand-written seeded profiles, which have no source
deck to read a theme from; every consumer must treat them as optional.

Note the division of labour: this module answers "what conventions does this
DECK follow?" by surveying slides. qc/stylespec.py answers "what does this
MASTER define?" by reading the master, layouts, and theme, and never looks at
slide content. A submitted master with no slides is Stage 1's normal case and
this module's degenerate one.
"""

import argparse
import json
from collections import Counter
from datetime import date
from pathlib import Path

from pptx import Presentation

from spike.arabic import contains_arabic, cs_typeface
from spike.color_resolver import resolve_solid_fill
from spike.ns import find
from spike.resolver import resolve_run
from .profile import PROFILES_DIR
from .stylespec import (MIN_LOGO_SLIDE_STAMPS, dominant_master, extract_brand,
                        extract_theme, find_brand_marks)
from .util import font_role, iter_shapes_deep

# Re-exported so callers and tests keep one import site for these; the
# implementations live in stylespec.py, which owns the design surface.
__all__ = ["build_profile", "learn_margins", "survey", "dominant_master",
           "extract_theme", "extract_brand", "find_brand_marks"]

MIN_FAMILY_SHARE = 0.05
MIN_COLOR_COUNT = 3
MAX_PALETTE = 12
MIN_MARGIN_SAMPLES = 20

# Alignment tolerances for bootstrapped profiles sit at a PERCEPTUAL floor
# (~0.8mm), not machine precision: calibrating against a designer-finished
# ground-truth deck (19/07/2026) showed that professionally formatted decks
# carry many sub-millimeter offsets no eye can see, and flagging them is
# pure noise. Deviations >= ~1mm still escalate to error in the alignment
# module, so the visible problems keep their teeth.
VISUAL_EDGE_TOLERANCE_EMU = 28575     # 0.79mm
VISUAL_SPACING_TOLERANCE_EMU = 28575
INTENT_WINDOW_EMU = 137160            # 0.15in, mirrors the module default


def _allowed(counter: Counter) -> list[str]:
    total = sum(counter.values())
    if not total:
        return []
    return [fam for fam, n in counter.most_common()
            if n >= max(2, total * MIN_FAMILY_SHARE)]


def survey(prs) -> dict:
    latin = {r: Counter() for r in ("title", "subtitle", "body", "caption")}
    cs = {r: Counter() for r in ("title", "subtitle", "body", "caption")}
    sizes = {r: Counter() for r in ("title", "subtitle", "body", "caption")}
    colors = Counter()
    layouts = Counter()

    for slide in prs.slides:
        layouts[slide.slide_layout.name] += 1
        master = slide.slide_layout.slide_master
        for shape, _path in iter_shapes_deep(slide.shapes):
            spPr = find(shape._element, "p:spPr")
            if spPr is not None:
                rgb = resolve_solid_fill(spPr, master)
                if rgb:
                    colors["%02X%02X%02X" % rgb] += 1
            if not getattr(shape, "has_text_frame", False):
                continue
            role = font_role(shape)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text or not run.text.strip():
                        continue
                    eff = resolve_run(run, para, shape, slide, prs)
                    rPr = find(run._r, "a:rPr")
                    if rPr is not None:
                        rgb = resolve_solid_fill(rPr, master)
                        if rgb:
                            colors["%02X%02X%02X" % rgb] += 1
                    if contains_arabic(run.text):
                        c = cs_typeface(run)
                        if c:
                            cs[role][c] += 1
                    else:
                        latin[role][eff.family.value] += 1
                    if eff.size_pt.value:
                        sizes[role][round(eff.size_pt.value * 2) / 2] += 1
    return {"latin": latin, "cs": cs, "sizes": sizes,
            "colors": colors, "layouts": layouts}


MARGIN_BUFFER_EMU = 28575  # shapes AT the learned margin must not flag


def learn_margins(prs) -> dict | None:
    """The margins the deck's TEXT actually respects: the 5th percentile of
    text-shape edge distances per side, minus a small buffer so shapes at
    the observed margin never flag. None when the deck is too sparse to
    trust (falls back to canvas ratios).

    Ground truth (19/07/2026, designer-formatted deck): the content grid is
    defined by text blocks (~0.48in margin there); decorative graphics roam
    past it by design, so they neither teach the margin here nor answer to
    it in the safe-zone check."""
    sw, sh = prs.slide_width, prs.slide_height
    lefts, tops, rights, bottoms = [], [], [], []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False) \
                    or not shape.text_frame.text.strip():
                continue
            l, t = shape.left, shape.top
            w, h = shape.width, shape.height
            if None in (l, t, w, h) or getattr(shape, "rotation", 0):
                continue
            if w * h >= 0.7 * sw * sh or w >= 0.95 * sw or h >= 0.95 * sh:
                continue  # full-bleed content crosses margins by design
            lefts.append(l)
            tops.append(t)
            rights.append(sw - (l + w))
            bottoms.append(sh - (t + h))
    if len(lefts) < MIN_MARGIN_SAMPLES:
        return None

    def pct5(vals):
        vals = sorted(vals)
        return vals[max(0, len(vals) * 5 // 100)]

    # clamp to [0, 15% of the dimension]: percentile handles outliers, the
    # buffer keeps at-the-margin shapes quiet, the cap keeps sparse decks
    # from producing absurd margins
    return {
        "left": max(0, min(pct5(lefts) - MARGIN_BUFFER_EMU, int(0.15 * sw))),
        "right": max(0, min(pct5(rights) - MARGIN_BUFFER_EMU, int(0.15 * sw))),
        "top": max(0, min(pct5(tops) - MARGIN_BUFFER_EMU, int(0.15 * sh))),
        "bottom": max(0, min(pct5(bottoms) - MARGIN_BUFFER_EMU, int(0.15 * sh))),
    }


def build_profile(prs, profile_id: str, name: str) -> dict:
    s = survey(prs)
    w, h = prs.slide_width, prs.slide_height

    roles = {}
    for role in ("title", "subtitle", "body", "caption"):
        allowed_latin = _allowed(s["latin"][role]) or _allowed(s["latin"]["body"])
        allowed_cs = _allowed(s["cs"][role]) or _allowed(s["cs"]["body"])
        entry = {"latin": allowed_latin, "complex_script": allowed_cs,
                 "allowed_weights": ["regular", "bold"]}
        # size targets only where role inference is reliable (placeholders)
        if role in ("title", "subtitle") and s["sizes"][role]:
            entry["size_pt"] = s["sizes"][role].most_common(1)[0][0]
        roles[role] = entry

    theme = extract_theme(prs)
    # An observed color that IS a theme slot gets labelled with its role, so a
    # later theme change carries the palette with it instead of stranding
    # literal RGB. Exact hex matches only: perceptual matching is the color
    # module's job at audit time, under its own tolerances.
    slot_by_hex = {}
    for slot, hexval in (theme.get("colors") or {}).items():
        slot_by_hex.setdefault(hexval, slot)

    named = [{"name": f"color_{i + 1:02d}", "hex": hexval,
              "theme_ref": slot_by_hex.get(hexval),
              "allowed_tints": [], "allowed_shades": []}
             for i, (hexval, n) in enumerate(s["colors"].most_common(MAX_PALETTE))
             if n >= MIN_COLOR_COUNT]

    allowlist = [name_ for name_, n in s["layouts"].items() if n >= 2]

    return {
        "id": profile_id,
        "name": name,
        "client_scope": None, "project_scope": None,
        "is_default": False, "version": 1,
        "owner": f"bootstrapped {date.today().strftime('%d/%m/%Y')}, review with design lead",
        "config": {
            # What the design system declares, as opposed to what slides do.
            "theme": theme,
            "brand": extract_brand(prs),
            "master_slide": {
                "enforce_existing_only": True, "pinned_layout_id": None,
                "layout_allowlist": allowlist,
                "geometry_tolerance_emu": 9525,
            },
            "font": {"roles": roles, "theme_font_refs_allowed": True,
                     "size_tolerance_pt": 0.5},
            "color_palette": {
                "theme_color_slots": ["dk1", "lt1", "dk2", "lt2", "accent1", "accent2",
                                      "accent3", "accent4", "accent5", "accent6"],
                "named_colors": named,
                "on_palette_mode": "by_name",
                "match_tolerance_deltaE": 2.0,
                "auto_replace_max_deltaE": 5.0,
                "ambiguity_band_deltaE": 10.0,
            },
            "geometry": {
                # margins learned from the deck's own content edges; canvas
                # ratios only when the deck is too sparse to trust
                "safe_zone_margins_emu": learn_margins(prs) or {
                    "left": round(w * 0.0375), "right": round(w * 0.0375),
                    "top": round(h * 0.040), "bottom": round(h * 0.053),
                },
                "grid": {"columns": 12, "gutter_emu": 0, "enabled": False},
                "alignment": {"edge_tolerance_emu": VISUAL_EDGE_TOLERANCE_EMU,
                              "center_tolerance_emu": VISUAL_EDGE_TOLERANCE_EMU,
                              "spacing_tolerance_emu": VISUAL_SPACING_TOLERANCE_EMU,
                              "intent_window_emu": INTENT_WINDOW_EMU},
            },
            "shape_size": {"size_tolerance_emu": 9525, "min_cohort_size": 3,
                           "preserve_picture_aspect": True,
                           "dominant_size_strategy": "median"},
            "header_footer": {"template": {"footer_text": None, "slide_number": False,
                                           "date": {"enabled": False, "format": "DD/MM/YYYY"},
                                           "position_emu": None, "font_role": "caption"}},
            # the deck's own label-case convention, when it clearly has one
            "typography": {"label_case": _learned_label_case(prs)},
        },
    }


def _learned_label_case(prs):
    from qc.modules.typography import learn_label_case

    return learn_label_case(prs)


def main(argv=None):
    ap = argparse.ArgumentParser(description="Generate a profile from a reference deck.")
    ap.add_argument("deck")
    ap.add_argument("--id", required=True, help="Profile id (filename in qc/profiles/)")
    ap.add_argument("--name", default=None)
    args = ap.parse_args(argv)

    prs = Presentation(args.deck)
    name = args.name or f"{args.id} (bootstrapped from {Path(args.deck).name})"
    profile = build_profile(prs, args.id, name)
    out = PROFILES_DIR / f"{args.id}.json"
    out.write_text(json.dumps(profile, indent=2, ensure_ascii=False), encoding="utf-8")

    cfg = profile["config"]
    print(f"Profile written: {out}")
    for role in ("title", "subtitle", "body", "caption"):
        r = cfg["font"]["roles"][role]
        size = r.get("size_pt", "no target")
        print(f"  {role:9} latin={r['latin']} cs={r['complex_script']} size={size}")
    print(f"  palette: {[c['hex'] for c in cfg['color_palette']['named_colors']]}")
    print(f"  layout allowlist: {cfg['master_slide']['layout_allowlist']}")
    print(f"  safe zones (EMU): {cfg['geometry']['safe_zone_margins_emu']}")

    theme = cfg.get("theme") or {}
    if theme.get("colors"):
        accents = {k: v for k, v in theme["colors"].items() if k.startswith("accent")}
        print(f"  theme accents: {accents}")
    tf = (theme.get("fonts") or {})
    if tf:
        print(f"  theme fonts: major={tf['major']['latin']}/"
              f"{tf['major']['complex_script']} "
              f"minor={tf['minor']['latin']}/{tf['minor']['complex_script']}")
    logo = (cfg.get("brand") or {}).get("logo")
    if logo:
        where = (f"master" if logo["scope"] == "master"
                 else f"{len(logo['layouts'])} layout(s)" if logo["layouts"]
                 else f"{logo['slide_count']} slides (not on the master)")
        drift = ", position varies" if logo["position_varies"] else ""
        print(f"  logo: {logo['image_sha1'][:12]} on {where}{drift}")
    else:
        print("  logo: none detected (see find_brand_marks docstring for scope)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
