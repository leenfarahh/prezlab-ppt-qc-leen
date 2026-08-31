"""One place to ask about the deck, and to ask for it to be done.

The tool has five doors - read a master, check coverage, apply a master, audit,
design QC - and a designer holding a client deck does not arrive knowing which
one answers their question. "Is this master missing anything?" is the coverage
report. "Why is this navy different from that navy?" is the palette. "What is
going to break if I download this now?" is three of them at once. Learning the
map is not the job.

So this answers instead. It reads what the passes already produced and says what
it found, in prose, with links to the exact card where the answer gets acted on.

AND IT ACTS (design lead, 27/08/2026). Pointing a designer at a card was the
right answer while the alternative was a chat that quietly rewrote a client
deck, and it is the wrong one now that there is a way to have both: "fix the
fonts on slide 7" comes back as a PLAN - what would happen, to which records,
in a sentence built from what actually resolved - and one press performs it.
The vocabulary and the resolving are qc.actions; the performing is the route,
through the very same function the button on the page calls.

THE CONFIRMATION GATE DOES NOT MOVE. Nothing this module returns has changed
anything. A plan is a proposal with a handle, it expires with the job, and the
designer presses the button. What got faster is finding the button, which was
the actual complaint; what did not get faster is changing a client's deck on
the strength of an ambiguous sentence.

THREE RULES MAKE THE ANSWERS TRUSTWORTHY, and they are the same three the rest of
the tool runs on:

FACTS FIRST. The model is handed a fact sheet assembled here, deterministically,
out of what the passes recorded - the coverage report, the palette roll-up, the
findings, the plans, the proposals waiting to be ticked. It is not handed the
deck and asked to look. Nothing it says can be newer or different from what the
pages show, because it is reading the same records the pages render.

A CLOSED SET OF PLACES TO SEND SOMEBODY, AND A CLOSED SET OF THINGS TO DO. The
model picks a link KIND and a slide number from the facts; this file builds the
URL. It cannot produce a link to a page that does not exist, to another job, or
to anything that POSTs. An action is the same shape: a name from a fixed list
and ids drawn from this job's own records, resolved by code before a designer is
shown a word of it, so a request naming a finding that is not open is refused
with what IS open rather than performed against something else.

AND EVERY CLAIM IS CHECKED BEFORE IT IS SHOWN. An answer naming a colour the
deck does not use, or a slide past the end of it, is a fabrication however
plausible it reads, and it is discarded rather than shown with a hedge
(qc.layoutpick drops a layout name the master does not have for the same
reason). What comes back then is a plain "that could not be answered from what
the passes recorded", which is a true statement and a useful one.

WHAT IS SENT: what the pages already display - headlines, counts, slide numbers,
hex codes, layout names, shape labels. Never the deck's body copy. A designer
asking what slide 4 says gets told that is not something this can see, which is
the honest answer and the same rule qc.assist follows.
"""

import json
import re

from .actions import ACTION_SCHEMA, Refused
from .actions import plan as plan_action
from .actions import vocabulary
from .llm import ask_json

# Enough of each list for a real answer without spending the context window on
# the tail of a 2000-record audit. Every cap that bites is stated in the facts
# themselves, so the model can say "and more beyond these" rather than implying
# the list is complete.
MAX_DESIGN_FINDINGS = 40
MAX_AUDIT_KINDS = 30
MAX_COLOURS = 24
MAX_PROPOSALS = 30
MAX_GAPS = 12
MAX_QUESTION_CHARS = 500

# Where a designer can be sent, per job kind. The model names a kind; this file
# builds the URL, so a link is always to a page that exists and always a GET.
_LINKS = {
    "audit": {
        "slide_card": ("/design/{job}?n={slide0}",
                       "the design card for that slide"),
        "deck_cards": ("/design/{job}?view=deck",
                       "the deck-wide decisions"),
        "audit_report": ("/audit/{job}",
                         "the full audit report"),
        "checklist": ("/checklist/{job}",
                      "the colour and type checklist for this deck"),
    },
    "format": {
        "review_deck": ("/format/{job}/review?view=deck",
                        "the before/after review, slide by slide"),
        "review_master": ("/format/{job}/review?view=master",
                          "the before/after review of the layouts"),
        "checklist": ("/checklist/{job}",
                      "the colour and type checklist for this deck"),
    },
    # A prep job IS both, so it can be sent to either set. Spelled out rather
    # than unioned at import: a link kind a page does not have is a dead button,
    # and the two sets above are free to diverge.
    "prep": {
        "slide_card": ("/design/{job}?n={slide0}",
                       "the design card for that slide"),
        "deck_cards": ("/design/{job}?view=deck",
                       "the deck-wide decisions"),
        "audit_report": ("/audit/{job}",
                         "the full audit report"),
        "review_deck": ("/format/{job}/review?view=deck",
                        "the before/after review, slide by slide"),
        "review_master": ("/format/{job}/review?view=master",
                          "the before/after review of the layouts"),
        "checklist": ("/checklist/{job}",
                      "the colour and type checklist for this deck"),
    },
}

ANSWER_SCHEMA = {
    "type": "object", "additionalProperties": False,
    "required": ["answer", "links", "colours", "slides", "answerable"],
    "properties": {
        "answer": {"type": "string"},
        # What the designer asked to have DONE, when they asked for anything.
        # Absent is the normal case: most questions are questions. The names
        # and the ids come from the vocabulary in the fact sheet, and qc.actions
        # resolves both against the job before a plan reaches the page.
        "action": ACTION_SCHEMA,
        # The kind, and the slide it applies to when it needs one. Never a URL.
        "links": {
            "type": "array",
            "items": {
                "type": "object", "additionalProperties": False,
                "required": ["kind", "label"],
                "properties": {
                    "kind": {"type": "string"},
                    "label": {"type": "string"},
                    "slide": {"type": "integer"},
                },
            },
        },
        # Every hex the answer names, and every slide number, listed separately
        # so they can be checked against the facts without parsing the prose.
        "colours": {"type": "array", "items": {"type": "string"}},
        "slides": {"type": "array", "items": {"type": "integer"}},
        "answerable": {"type": "boolean"},
    },
}

_SYSTEM = """You are the assistant inside Prezlab's PowerPoint QC tool, talking
to the designer who is using it.

You are given a FACT SHEET about one deck, assembled from what the tool's passes
actually recorded, and a question. Answer from the fact sheet and from nothing
else. You cannot see the deck, and the fact sheet does not contain the words on
the slides.

Set `answerable` to false when the fact sheet does not contain the answer, and
say so in one sentence naming what would answer it (which pass, or which page).
A designer told "the palette report does not cover that" can go and get it; a
designer given a confident guess cannot tell it from a fact.

WHEN THE DESIGNER ASKS FOR SOMETHING TO BE DONE, set `action`. Pick a `name`
from the `actions` list in the fact sheet, and fill in ONLY ids that appear in
the vocabulary beside it: issue types from `issue_types_you_may_name`, a
`finding` and `remedy` from `open_design_findings`, ids from
`decisions_already_applied` or `pending_removals`. Slide numbers are as a
designer says them, counting from 1. Never invent an id, and never name a shape,
a colour value or a coordinate: you are choosing from a menu, not writing an
instruction.

Choose `decide` for the broad asks - "sort this out", "do what you can", "clean
up slide 7" - and `fix_findings` when they named what to fix. Use `take_remedy`
only when they named the finding or chose between its options, because a design
finding has more than one right answer and the choice is theirs.

Leave `action` out for a question. "Which colours were typed in by hand?" is a
question. Asking for something you were not asked for is worse than asking for
nothing: a designer reads the plan and presses the button.

Setting `action` does NOT change anything by itself. The tool resolves it
against the real records and shows the designer exactly what would happen; they
press the button. So say what you are proposing in the answer, in the same
sentence you would use to a colleague, and never write as though it is already
done.

`links` are the other half. Pick a `kind` from the list the fact sheet gives you and,
where the kind needs one, the `slide` number it applies to. Do not write URLs.
At most three links, and only ones a designer would actually use next.

List every hex code your answer names in `colours`, and every slide number in
`slides`. They are checked against the fact sheet, and an answer naming a colour
this deck does not use is thrown away rather than shown.

Write like a senior designer talking to a peer: two or three sentences, US
English, active voice, no em dashes, no filler openers. Lead with the answer.
Numbers only where the fact sheet has them."""


# --- the fact sheet -------------------------------------------------------


def _coverage_block(cov) -> dict:
    """A Coverage as facts. One shape whichever pass produced it, so an answer
    about layouts reads the same on the design page and on the format result."""
    return {
        "placed": cov.matched, "no_layout": cov.unplaced,
        "by_rule": cov.by_rule,
        "looked_at": cov.reviewed, "not_looked_at": cov.not_reviewed,
        "slides_were_looked_at": cov.review_ran,
        "unused_layouts": cov.unused_layouts[:MAX_GAPS],
        "gaps": [{"wants": g.label, "slides": [i + 1 for i in g.slides][:12],
                  "how_many": g.places, "closest_layout": g.closest,
                  "why_it_does_not_fit": g.closest_note,
                  "checked_and_refused": g.refused}
                 for g in (cov.gaps or [])[:MAX_GAPS]],
    }


def _master_layouts(job) -> list | None:
    """The layouts of whatever master this job was judged against, or None.

    Two ways in, because the audit takes two: a master uploaded for one run
    leaves its Style Spec on the job, and a saved profile keeps the master file
    itself. Either gives the same list qc.applymaster plans from."""
    spec = job.get("master_spec") or {}
    layouts = spec.get("layouts")
    if layouts:
        return layouts
    profile = job.get("profile")
    if not profile:
        return None
    try:
        from .stylespec import dominant_master, extract_layouts
        from .templates import load_master

        blob = load_master(profile)
        if not blob:
            return None
        import io

        from pptx import Presentation

        master = dominant_master(Presentation(io.BytesIO(blob)))
        if master is None:
            return None
        return extract_layouts(master, embed_assets=False)
    except Exception:
        return None


def _coverage_facts(job) -> dict | None:
    """Whether this master can build this deck, computed if nobody has yet.

    A format job already has it. An audit job does not, and the question is the
    same one - so it is planned here rather than answered with "not available"
    about something the files state. Deterministic: the planner and the layout
    list, no render and no model, which is why it is cheap enough to do on the
    way to answering a question.
    """
    cov = job.get("coverage")
    if cov is not None:
        return _coverage_block(cov)
    deck = job.get("deck")
    layouts = _master_layouts(job)
    if not deck or not layouts:
        return None
    try:
        import io

        from pptx import Presentation

        from .applymaster import plan_assignments
        from .layoutgap import report

        prs = Presentation(io.BytesIO(deck))
        plans = plan_assignments(prs, layouts)
        return _coverage_block(report(prs, layouts, plans))
    except Exception:
        return None


def _audit_facts(job) -> dict:
    """What the audit and the design pass found, as the pages show it.

    Records are rolled up by kind rather than listed: a deck with 900 font
    records has one answer to "what is wrong with the type", and it is not 900
    rows. Design findings ARE listed, because each one is a decision a designer
    has to make and they are few."""
    manifest = job.get("manifest") or {}
    records = [r for r in (manifest.get("records") or [])
               if r.get("module") != "preflight"]

    kinds: dict[tuple, dict] = {}
    for rec in records:
        key = (rec.get("module"), rec.get("issue_type"), rec.get("severity"))
        entry = kinds.setdefault(key, {
            "module": key[0], "issue": key[1], "severity": key[2],
            "count": 0, "slides": set()})
        entry["count"] += 1
        entry["slides"].add(rec.get("slide_index"))

    ranked = sorted(kinds.values(), key=lambda e: -e["count"])
    audit = [{"module": e["module"], "issue": e["issue"],
              "severity": e["severity"], "count": e["count"],
              "slides": sorted(s + 1 for s in e["slides"] if s is not None)[:12]}
             for e in ranked[:MAX_AUDIT_KINDS]]

    answered = {a.finding_id for a in (job.get("design_applied") or [])}
    findings = [f for f in (job.get("design") or [])
                if f.finding_id not in answered]
    design = [{"id": f.finding_id, "kind": f.kind, "severity": f.severity,
               "headline": f.headline,
               "slides": [i + 1 for i in (f.slides or [])][:12],
               "places": f.evidence.get("places") if f.evidence else None,
               "options": [o.label for o in (f.options or [])]}
              for f in findings[:MAX_DESIGN_FINDINGS]]

    return {
        "kind": "audit",
        "slides": manifest.get("slides"),
        "profile": job.get("profile"),
        "audit_findings_by_kind": audit,
        "audit_kinds_omitted": max(0, len(ranked) - MAX_AUDIT_KINDS),
        "design_decisions": design,
        "design_decisions_omitted": max(0, len(findings) - MAX_DESIGN_FINDINGS),
        "design_decisions_answered": len(answered),
    }


def _format_facts(job) -> dict:
    """What the format pass did, what it could not place, and what is waiting
    for a tick."""
    plans = job.get("plans") or []
    by_rule: dict[str, int] = {}
    for plan in plans:
        by_rule[plan.match_rule] = by_rule.get(plan.match_rule, 0) + 1

    proposals = [c for c in (job.get("changes") or [])
                 if getattr(c, "remove_op", None)]
    done = set(job.get("removed") or [])
    waiting = [{"slide": c.slide_index + 1, "what": c.action,
                "detail": c.detail}
               for c in proposals if c.remove_id not in done][:MAX_PROPOSALS]

    facts = {
        "kind": "format",
        "slides": len(plans),
        "profile": job.get("profile"),
        "layouts_by_rule": by_rule,
        "slides_that_failed": sorted(i + 1 for i in (job.get("errors") or {})),
        "masters_in_output": job.get("masters"),
        "removals_waiting_for_a_tick": waiting,
        "removals_waiting_omitted": max(0, len(proposals) - len(waiting)
                                        - len(done)),
        "removals_already_performed": len(done),
    }

    return facts


def _palette_facts(job) -> dict | None:
    """The deck's own colours, and which of them were typed in by hand.

    Read rather than remembered, because it is the one thing a designer asks
    about that no pass records on its own (qc.extract.palette_inventory) - but
    read ONCE PER DECK, not once per question. The extraction is a full parse
    plus a colour and font resolution of every run on every slide, and a
    designer asks three or four things in a row about a document that is not
    changing between them. It is cached on the job under the same key the
    checklist page uses, and dropped by the same _invalidate_renders that drops
    the thumbnails when a fix lands.
    """
    deck_bytes = job.get("deck")
    if not deck_bytes:
        return None
    try:
        from .extract import extract_deck, palette_inventory

        extracted = job.get("extracted")
        if extracted is None:
            extracted = job["extracted"] = extract_deck(deck_bytes)
        pal = palette_inventory(extracted)
    except Exception:
        return None
    return {
        "theme_slots": pal.get("theme_slots") or {},
        "distinct_colours": pal.get("distinct_count"),
        "hand_typed_colours": pal.get("explicit_count"),
        "colours": [{"hex": c["hex"], "uses": c["uses"],
                     "roles": sorted(c["roles"]),
                     "written_as": sorted(c["written"]),
                     "slides": [i + 1 for i in c["slides"]][:8]}
                    for c in (pal.get("colours") or [])[:MAX_COLOURS]],
        "colours_omitted": max(0, (pal.get("distinct_count") or 0)
                               - MAX_COLOURS),
    }


def facts(job, kind: str) -> dict:
    """Everything known about this deck, small enough to send and honest about
    what it leaves out."""
    if kind == "prep":
        # A prep job ran both passes on one deck, so it gets both sheets rather
        # than a choice between them. The audit half goes on top because the
        # format half's keys are about the rebuild and the audit's are about
        # the file as it now stands, and a collision would silently answer the
        # second question with the first one's number.
        out = _format_facts(job)
        out.update(_audit_facts(job))
        out["kind"] = "prep"
        out["slides"] = (out.get("slides")
                         or len(job.get("plans") or []) or None)
    else:
        out = _audit_facts(job) if kind == "audit" else _format_facts(job)
    out["deck_name"] = job.get("filename")
    # Asked of both kinds, because a designer's question does not know which
    # page they are on. On a format job it is already computed; on an audit job
    # it is planned here, which is what makes this one box rather than two.
    coverage = _coverage_facts(job)
    if coverage is not None:
        out["coverage"] = coverage
    else:
        out["coverage"] = ("not available: this run was not judged against a "
                           "master file, so there is nothing to compare the "
                           "deck's layouts with")
    palette = _palette_facts(job)
    if palette is not None:
        out["palette"] = palette
    else:
        out["palette"] = ("not available: the deck's bytes are no longer held "
                          "in memory, so its colours cannot be read")
    out["link_kinds"] = {name: why
                         for name, (_url, why) in _LINKS.get(kind, {}).items()}
    # What may be ASKED FOR, in the ids this job actually has. Assembled per
    # job for the reason the rest of the sheet is: a model shown a finding that
    # is not open will ask for it, and the designer gets a refusal instead of
    # the thing they wanted (qc.actions.vocabulary).
    try:
        out["you_may_ask_for"] = vocabulary(job)
    except Exception:
        out["you_may_ask_for"] = ("not available: this job's records could not "
                                  "be read, so nothing can be asked for")
    out["what_you_cannot_see"] = (
        "the words on the slides, the pictures, and anything no pass recorded")
    return out


# --- checking the answer --------------------------------------------------


_HEX = re.compile(r"#?\b([0-9A-Fa-f]{6})\b")
_SLIDE = re.compile(r"\bslide\s+(\d+)\b", re.I)


def _known_colours(sheet: dict) -> set[str]:
    palette = sheet.get("palette")
    if not isinstance(palette, dict):
        return set()
    known = {str(v).upper() for v in (palette.get("theme_slots") or {}).values()}
    known |= {str(c.get("hex") or "").upper()
              for c in (palette.get("colours") or [])}
    return {h for h in known if h}


def unverified(answer: dict, sheet: dict) -> list[str]:
    """Claims in this answer that the fact sheet does not support.

    Checked from the prose AND from the model's own lists, because either one
    alone can be gamed by accident: a hex in the sentence but not in `colours`
    is still a claim a designer will read as a fact.

    Colours are checked only when the palette was readable at all - with the
    deck's bytes gone there is nothing to check against, and inventing a
    complaint is as bad as missing one.
    """
    problems = []
    total = sheet.get("slides") or 0

    numbers = {int(n) for n in _SLIDE.findall(answer.get("answer") or "")}
    numbers |= {int(n) for n in (answer.get("slides") or [])
                if isinstance(n, int)}
    for n in sorted(numbers):
        if total and not (1 <= n <= total):
            problems.append(f"slide {n}, in a deck of {total}")

    known = _known_colours(sheet)
    if known:
        named = {h.upper() for h in _HEX.findall(answer.get("answer") or "")}
        named |= {str(h).lstrip("#").upper()
                  for h in (answer.get("colours") or [])}
        for h in sorted(named - known):
            problems.append(f"#{h}, which this deck does not use")
    return problems


def _links(answer: dict, kind: str, job_id: str, slides: int) -> list[dict]:
    """Validated links, as {label, href}. Anything the closed set does not
    cover is dropped rather than guessed at."""
    allowed = _LINKS.get(kind, {})
    out, seen = [], set()
    for item in (answer.get("links") or [])[:6]:
        template = allowed.get(str(item.get("kind") or ""))
        if template is None:
            continue
        url, _why = template
        slide = item.get("slide")
        if "{slide0}" in url:
            if not isinstance(slide, int) or not (1 <= slide <= (slides or 0)):
                continue
        href = url.format(job=job_id, slide0=(slide or 1) - 1)
        if href in seen:
            continue
        seen.add(href)
        label = " ".join(str(item.get("label") or "").split())[:60]
        out.append({"label": label or "Open", "href": href})
        if len(out) == 3:
            break
    return out


# --- the one call ---------------------------------------------------------


def _plan(answer: dict, job) -> tuple[object | None, str]:
    """(Plan, refusal) for whatever the answer asked to have done.

    A refusal is returned as PROSE, not as an error, and it replaces nothing:
    the answer still shows, with "that particular thing could not be done, and
    here is what was actually there" under it. The two most common reasons a
    plan does not resolve - the finding was already answered, the slide has
    nothing fixable on it - are both facts a designer wants to be told.
    """
    request = answer.get("action")
    if not isinstance(request, dict) or not request.get("name"):
        return None, ""
    try:
        return plan_action(job, request), ""
    except Refused as refusal:
        return None, str(refusal)
    except Exception as exc:
        return None, (f"That could not be turned into something to do "
                      f"({type(exc).__name__}), so nothing was proposed and "
                      f"nothing was changed.")


def ask(job, kind: str, job_id: str, question: str) -> dict:
    """Answer one question about one deck, and propose the action it asked for.

    Returns {answer, links, answerable, checked, plan, refusal}. `plan` is a
    qc.actions.Plan the caller holds and performs on confirmation, or None;
    NOTHING has been changed by the time this returns.

    `checked` is False when the answer named something the facts do not support,
    and the answer is replaced rather than shown with a warning: a claim a
    designer cannot tell from a fact is worse than no claim. A discarded answer
    takes its plan with it - an answer that invented a colour has no business
    proposing a change on the strength of the same reasoning.
    """
    question = " ".join((question or "").split())[:MAX_QUESTION_CHARS]
    if not question:
        return {"answer": "Ask me something about this deck, or ask for it to "
                          "be done.", "links": [], "answerable": False,
                "checked": True, "plan": None, "refusal": ""}

    sheet = facts(job, kind)
    answer = ask_json(
        system=_SYSTEM,
        prompt=("Fact sheet:\n" + json.dumps(sheet, sort_keys=True, default=str)
                + "\n\nQuestion:\n" + question),
        schema=ANSWER_SCHEMA,
        max_tokens=2048,
    )

    problems = unverified(answer, sheet)
    if problems:
        return {
            "answer": ("That answer referred to " + ", and ".join(problems)
                       + ", so it was discarded rather than shown. Nothing was "
                       "changed. Ask again, or open the report and read the "
                       "rows directly."),
            "links": [], "answerable": False, "checked": False,
            "plan": None, "refusal": ""}

    plan, refusal = _plan(answer, job)
    text = " ".join(str(answer.get("answer") or "").split())
    return {
        "answer": text or "There is nothing recorded about that.",
        "links": _links(answer, kind, job_id, sheet.get("slides") or 0),
        "answerable": bool(answer.get("answerable")),
        "checked": True,
        "plan": plan,
        "refusal": refusal,
    }
