"""Component orchestration: the model decides WHAT the things on a slide are and
WHICH LINE they belong on; code measures, computes targets and applies them.

This exists because two decisions in the geometry pipeline are not geometry
questions, and every alignment bug this tool has shipped traces back to one of
them being answered by arithmetic.

WHAT IS ONE THING. A card, its icon and its label are one entity: move the card
without them and the slide breaks. Today that is guessed from overlap and
adjacency (qc.util.rides_with, qc.fixer._carried_contents), and the guess fails
in both directions - a corner rule welded to a photo gets left behind, while a
column of stacked blocks decides it is carrying its own neighbours and each one
moves twice. Adjacency is not composition, and no threshold makes it one.

WHICH ONE IS WRONG. Given four shapes that nearly share an edge, the audit takes
the median and calls the minority misaligned (margin_alignment._edge_misaligned;
qc.copilot does the same with its own median). On a slide where one element sits
on the master's stated line and three drifted off it, the majority wins and the
tool proposes pulling the one correct element off the line to join them (design
lead, 24/08/2026). A median is a vote about where things happen to be. What is
needed is a REFERENCE, and choosing one is a judgment about the design.

So the split is the same one qc.copilot established, one level up:

  The model (vision) - names the components, says which of them share a line, and
                     says what that line IS (the master's frame, or one named
                     component). Chooses from a closed vocabulary; may only
                     reference shape ids from the inventory it was handed.
  Code             - verifies every component against real geometry, resolves
                     the anchor to an EMU coordinate, measures who is off it,
                     computes each target, and emits ordinary FindingRecords so
                     they flow through the pipeline unchanged: tickable, never
                     pre-selected, collision-guarded, before/after rendered.
  The designer     - approves, as with every other fix.

The model never emits a coordinate. Ask a language model for EMU and you get
plausible EMU; ask it which of two lines is the intended one and you get the
answer geometry cannot compute.

Confidentiality: this sends SLIDE IMAGES to whichever model qc.llm is pointed
at, like the design copilot. Use it only on decks approved for cloud processing.
"""

import io
import json

from pptx import Presentation

from .design import _dimensions
from .llm import ask_in_parallel, ask_json
from .records import make_record

MAX_SLIDES = 20
MAX_COMPONENTS = 24
MAX_ALIGNMENTS = 6

# The perceptual floor and ceiling, shared with qc.copilot so two passes cannot
# disagree about what "off the line" means. Below TOL a designer cannot see it;
# above WINDOW the gap is a layout decision rather than a drift.
TOL_EMU = 28575           # 0.03in
WINDOW_EMU = 457200       # 0.5in

AXES = ("top", "left", "right")

LAYOUT_SCHEMA = {
    "type": "object", "additionalProperties": False,
    "required": ["components", "alignments"],
    "properties": {
        "components": {
            "type": "array",
            "items": {
                "type": "object", "additionalProperties": False,
                "required": ["name", "shape_ids"],
                "properties": {
                    "name": {"type": "string"},
                    "shape_ids": {"type": "array",
                                  "items": {"type": "string"}},
                    # Whether this component should be a REAL group in the
                    # file, not merely treated as one for alignment.
                    "group": {"type": "boolean"},
                },
            },
        },
        "alignments": {
            "type": "array",
            "items": {
                "type": "object", "additionalProperties": False,
                "required": ["axis", "components", "anchor", "rationale"],
                "properties": {
                    "axis": {"enum": list(AXES)},
                    "components": {"type": "array",
                                   "items": {"type": "string"}},
                    "anchor": {"type": "string"},
                    "rationale": {"type": "string"},
                },
            },
        },
    },
}

_SYSTEM = """You are a senior presentation designer at Prezlab looking at one
slide. You have the rendered image and an inventory of its shapes (id, kind,
position and size as fractions of the slide, whether it holds text).

Answer two questions and nothing else.

1. COMPONENTS. Group the shape ids into the things a designer would select and
drag as a unit: a card with its icon and label, a photo with its caption, a
number chip welded to the box it labels, a heading with its rule. Every shape
belongs to exactly one component; a shape that composes with nothing is a
component of one. Name each component in two or three plain words. At most
{max_components}.

2. ALIGNMENTS. Say which components were meant to share an edge, on which axis
(top, left or right), and what the reference is - the `anchor`. The anchor is
either the exact name of one of your components, or the literal string "frame"
when the intended line is the presentation space edge described below. At most
{max_alignments}, and only where a designer would agree the elements were meant
to line up: a row of cards sharing a top, a column of blocks sharing a left
edge, a panel meant to start where the frame starts.

Rules that matter more than coverage:
- Do NOT say which components are misaligned or by how much. Name the intended
  line and the code will measure it. You are choosing the reference, not the
  offender.
- Prefer "frame" as the anchor whenever the master states an edge the elements
  were plainly meant to sit on. A stated line outranks whatever the majority of
  elements happen to be doing.
- Adjacency is not composition. Two cards side by side are two components. A
  label sitting ON a card is part of it.
- A deliberate indent, a stagger, or an off-grid composition is not a
  misalignment. If nothing was meant to share an edge, return an empty
  alignments list.
- Reference shapes ONLY by ids from the inventory. Write rationales in clear US
  English without em dashes.

Set `group` on a component when its members should be a REAL group in the file,
so the next person to open the deck drags one object instead of three. Say yes
when the members only make sense together: a card with its icon and its label, a
chart with its own key, a badge with its number. Say no when they are separate
things that merely align - two cards in a row, a title above a body, anything a
designer would want to move on its own. A group nobody wanted is worse than no
group, because it puts every element inside it one click further away."""


def inventory(slide, slide_w: int, slide_h: int, prs=None) -> list[dict]:
    """Shape inventory the model reasons over: ids, kind, normalized geometry.

    No text content - the image already shows the words, and the question here
    is about boxes. Rotated shapes are listed (they are part of a component and
    travel with it) but flagged, because their stored box is not their rendered
    one and code refuses to measure them.

    One exception to "the image already shows the words", because sometimes it
    does not. White text on a white ground is not in the picture at all, so the
    box holding it reads as empty and gets left out of the component it belongs
    to - and asking a model to read what was never rendered is not a fixable
    prompt problem. Those shapes carry `invisible`, and the prompt says what it
    means (qc.design.invisible_text). `prs` is what resolving the run colours
    takes; without it the flag is absent rather than wrong.
    """
    from .design import invisible_text

    hidden = {}
    if prs is not None:
        try:
            hidden = invisible_text(slide, prs)
        except Exception:
            # A slide whose colours will not resolve is not a slide with no
            # hidden text; it is one this cannot answer for. The inventory still
            # goes, without the flag.
            #
            # The import is OUTSIDE this guard on purpose. It used to be inside,
            # and when a rewrite elsewhere deleted invisible_text the ImportError
            # was caught here and every slide reported no hidden text - a feature
            # silently gone, with green tests everywhere except its own.
            hidden = {}

    out = []
    for shape in slide.shapes:
        left, top, width, height = _dimensions(shape)
        if None in (left, top, width, height):
            continue
        entry = {
            "id": str(shape.shape_id),
            "kind": str(shape.shape_type).split(" ")[0].lower(),
            "x": round(left / slide_w, 4), "y": round(top / slide_h, 4),
            "w": round(width / slide_w, 4), "h": round(height / slide_h, 4),
            "text": bool(getattr(shape, "has_text_frame", False)
                         and shape.text_frame.text.strip()),
            "rotated": bool(getattr(shape, "rotation", 0)),
        }
        if str(shape.shape_id) in hidden:
            entry["invisible"] = True
        out.append(entry)
    return out


def _frame_note(space, slide_w: int, slide_h: int) -> str:
    """The presentation space, in the same normalized units as the inventory,
    so "frame" is a line the model can actually see rather than a word."""
    if space is None:
        return ("The master states no presentation space, so \"frame\" is not "
                "an available anchor on this slide.")
    left, top, right, bottom = space
    return ("The master states a presentation space - the rectangle the "
            "designer drew for content. In the inventory's units its edges "
            f"are left x={left / slide_w:.4f}, top y={top / slide_h:.4f}, "
            f"right x={right / slide_w:.4f}, bottom y={bottom / slide_h:.4f}. "
            "Use \"frame\" as the anchor when elements were meant to start on "
            "one of those edges.")


def _hidden_note(inv: list[dict]) -> str:
    """Tell the model about the words it cannot see.

    A flag in the payload that nothing explains is noise. This is the one place
    where the picture and the file disagree, and the file is right: white text
    on a white panel is really there, and a component built without it tears the
    label off the card it belongs to."""
    hidden = [entry["id"] for entry in inv if entry.get("invisible")]
    if not hidden:
        return ""
    return ("\n\nThe image does not show everything. These shapes hold text "
            "whose colour is within 1.5:1 of the colour behind it, so it is not "
            "in the picture even though it is in the file: "
            + ", ".join(hidden)
            + ". Treat them as text blocks that belong to whatever they sit on, "
              "not as empty boxes.")


def _ask_vision(png: bytes, inv: list[dict], frame_note: str) -> dict:
    """The one call this module makes. Which model answers is qc.llm's business
    and no part of this file's - the question and the schema are."""
    parsed = ask_json(
        system=_SYSTEM.format(max_components=MAX_COMPONENTS,
                             max_alignments=MAX_ALIGNMENTS),
        prompt=frame_note + "\n\nShape inventory:\n"
        + json.dumps(inv, sort_keys=True) + _hidden_note(inv),
        schema=LAYOUT_SCHEMA,
        images=[png],
    )
    return {"components": (parsed.get("components") or [])[:MAX_COMPONENTS],
            "alignments": (parsed.get("alignments") or [])[:MAX_ALIGNMENTS]}


# --- the precision gate ---------------------------------------------------


_EDGE = {
    "top": (lambda box: box[1], "spPr.xfrm.off.y"),
    "left": (lambda box: box[0], "spPr.xfrm.off.x"),
    "right": (lambda box: box[2], "spPr.xfrm.off.x"),
}


def _boxes(slide) -> dict:
    """Top-level shapes only, and that is the question this module asks: which
    LOOSE shapes form one card, and should they become a group. A shape already
    inside a group is excluded from a grouping proposal by _grouped_ids for the
    same reason, so descending into groups here would build components out of
    members that can never be offered."""
    out = {}
    for shape in slide.shapes:
        # qc.design's reader, so a box measured here and a box measured by the
        # audit are the same number arrived at the same way (and an inheriting
        # placeholder does not pay python-pptx's xpath scan four times).
        left, top, width, height = _dimensions(shape)
        if None in (left, top, width, height):
            continue
        out[str(shape.shape_id)] = {
            "shape": shape, "box": (left, top, left + width, top + height),
            "rot": bool(getattr(shape, "rotation", 0)),
        }
    return out


def _component_boxes(components: list[dict], by_id: dict) -> dict:
    """{name: {"ids": [...], "box": bounding box, "measurable": bool}}.

    A component's edge is its BOUNDING BOX's edge, which is what a designer
    sees and drags. Members that cannot be measured (rotation) still belong to
    it and still travel with it, but a component containing one is not held to
    a line - its rendered extent is not the box on file."""
    out = {}
    for comp in components:
        name = str(comp.get("name") or "").strip()
        ids = [str(i) for i in (comp.get("shape_ids") or [])]
        members = [by_id[i] for i in ids if i in by_id]
        if not name or not members or name in out:
            continue
        boxes = [m["box"] for m in members]
        out[name] = {
            "ids": [i for i in ids if i in by_id],
            "box": (min(b[0] for b in boxes), min(b[1] for b in boxes),
                    max(b[2] for b in boxes), max(b[3] for b in boxes)),
            "measurable": not any(m["rot"] for m in members),
        }
    return out


def _grouped_ids(slide) -> set:
    """Ids of shapes already inside a group. Regrouping a group is a different
    decision and not this one."""
    from .util import iter_shapes_deep

    return {str(shape.shape_id)
            for shape, path in iter_shapes_deep(slide.shapes) if path}


def _grouping_records(slide, s_idx: int, layout: dict, comps: dict,
                      seen: set) -> list[dict]:
    """One record per component the model says should be a real group.

    The model already answers "what is one thing" - that is what this pass
    exists for - and that answer was only ever used to move things TOGETHER. It
    was never written into the file, so the next person to open the deck drags
    the card and leaves the icon behind (design lead, 26/08/2026).

    Verified before a group is offered: two members at least, every id real,
    none rotated, and nothing already inside a group. A group takes the bounding
    box of its members, and a rotated member's stored box is not its rendered
    one, so grouping around one would move what it contains.
    """
    already = _grouped_ids(slide)
    out = []
    for spec in layout.get("components") or []:
        if not spec.get("group"):
            continue
        name = str(spec.get("name") or "").strip()
        comp = comps.get(name)
        if comp is None or len(comp["ids"]) < 2 or not comp["measurable"]:
            continue
        if any(i in already for i in comp["ids"]):
            continue
        ids = ",".join(comp["ids"])
        rec = make_record(
            slide_index=s_idx, shape_id=comp["ids"][0], shape_path=None,
            module="margin_alignment", issue_type="margin_alignment.should_be_grouped",
            severity="info", confidence="medium", action="flagged",
            source="vision",
            locator=f"group:{ids}", property="spTree.grpSp",
            old_value=f"{len(comp['ids'])} separate shapes", new_value=ids,
            profile_rule_id=None,
            message=(f"{name!r} is one object made of {len(comp['ids'])} "
                     f"shapes. Grouping them means the next person to open this "
                     f"deck drags the whole thing rather than leaving part of "
                     f"it behind. Component review."),
        )
        key = (rec.issue_type, rec.shape_id, rec.locator)
        if key in seen:
            continue
        seen.add(key)
        out.append(rec.to_dict())
    return out


def synthesize(slide, s_idx: int, layout: dict, space, existing: list[dict],
               slide_w: int, slide_h: int) -> list[dict]:
    """Verify what the model said against the real geometry and emit records.

    Everything that does not check out is dropped silently: an unknown shape id,
    a component naming nothing, an anchor that is neither "frame" nor a
    component, a claimed line the geometry says is already held, a gap too big
    to be drift. The pass is allowed to find nothing - that is a normal answer
    for a well-composed slide, and inventing a finding to look useful is the
    one thing it must not do.
    """
    by_id = _boxes(slide)
    comps = _component_boxes(layout.get("components") or [], by_id)
    seen = {(r["issue_type"], str(r["shape_id"]), r.get("locator"))
            for r in existing if r["slide_index"] == s_idx}
    out: list[dict] = []
    out.extend(_grouping_records(slide, s_idx, layout, comps, seen))

    for spec in layout.get("alignments") or []:
        axis = spec.get("axis")
        if axis not in _EDGE:
            continue
        edge_of, prop = _EDGE[axis]
        names = [str(n) for n in (spec.get("components") or [])]
        members = [comps[n] for n in names if n in comps]
        if len(members) < 2:
            continue

        anchor_name = str(spec.get("anchor") or "").strip()
        if anchor_name == "frame":
            if space is None:
                continue          # no frame stated: the anchor does not exist
            anchor = edge_of(space)
            anchor_ids: set = set()
            anchor_is_frame = True
            why = "the presentation space the master states"
        elif anchor_name in comps:
            ref = comps[anchor_name]
            if not ref["measurable"]:
                continue
            anchor = edge_of(ref["box"])
            anchor_ids = set(ref["ids"])
            anchor_is_frame = False
            why = f"{anchor_name!r}, which sits on the intended line"
        else:
            continue              # an anchor naming nothing is not an anchor

        rationale = str(spec.get("rationale") or "").strip()[:160]
        for name, comp in ((n, c) for n, c in zip(names, members)
                           if set(c["ids"]) - anchor_ids):
            if not comp["measurable"]:
                continue
            off = edge_of(comp["box"]) - anchor
            if axis == "right":
                off = -off
            # WHICH SIDE COUNTS DEPENDS ON WHAT THE ANCHOR IS, and conflating
            # the two cases cost this pass half its findings.
            #
            # Against the FRAME the test is one-sided on purpose: a component
            # sticking OUT past the presentation space is a margin breach and
            # belongs to that rule, so only inboard drift is this one's.
            #
            # Against another COMPONENT there is no outboard. A label 9mm to
            # the LEFT of the block it heads is exactly as misaligned as one
            # 9mm to the right, and no margin rule owns it because the shape is
            # still well inside the frame. The signed test silently dropped
            # every one of those: the model named the intended line correctly
            # and the measurement threw the answer away (31/08/2026, on a
            # two-column comparison slide whose left heading sat 9mm out).
            gap = off if anchor_is_frame else abs(off)
            # A gap past WINDOW is a composition, not a mistake.
            if not (TOL_EMU < gap <= WINDOW_EMU):
                continue
            ids = ",".join(comp["ids"])
            rec = make_record(
                slide_index=s_idx, shape_id=comp["ids"][0], shape_path=None,
                module="margin_alignment",
                issue_type="margin_alignment.component_edge_misaligned",
                severity="warning", confidence="medium", action="flagged",
                source="vision",
                locator=f"comp:{axis}:{ids}",
                property=prop,
                old_value=edge_of(comp["box"]), new_value=int(anchor),
                profile_rule_id="geometry.alignment.edge_tolerance_emu",
                message=(f"{name!r} sits {gap / 36000:.1f}mm off the {axis} "
                         f"edge of {why}; the fix moves its "
                         f"{len(comp['ids'])} element(s) together. "
                         f"Component review: {rationale}"),
            )
            key = (rec.issue_type, rec.shape_id, rec.locator)
            if key in seen:
                continue
            seen.add(key)
            out.append(rec.to_dict())
    return out


def run_components(deck_bytes: bytes, thumbs: dict[int, bytes],
                   manifest: dict, space=None) -> tuple[list[dict], int]:
    """Review up to MAX_SLIDES rendered slides; returns (new_records,
    slides_reviewed).

    A failure on one slide is skipped rather than raised: this is an advisory
    layer over a pipeline that works without it, and one bad API call must not
    cost a designer the whole run."""
    prs = Presentation(io.BytesIO(deck_bytes))
    slide_w, slide_h = prs.slide_width, prs.slide_height
    frame_note = _frame_note(space, slide_w, slide_h)
    existing = manifest.get("records") or []

    # WHICH SLIDES ARE WORTH A CALL IS DECIDED FIRST, and the budget is spent on
    # the calls rather than on the successes. The loop here used to stop once
    # MAX_SLIDES had ANSWERED, which reads as generous and is not: under an
    # outage nothing answers, so a 200-slide deck made 200 failing calls to
    # review nothing. Capping the candidates caps the cost at MAX_SLIDES however
    # the afternoon goes.
    candidates = []
    for s_idx, slide in enumerate(prs.slides):
        if len(candidates) >= MAX_SLIDES:
            break
        png = thumbs.get(s_idx)
        if png is None:
            continue
        inv = inventory(slide, slide_w, slide_h, prs)
        if len(inv) < 3:
            continue
        candidates.append((s_idx, slide, png, inv))

    answers = ask_in_parallel(
        candidates, lambda c: _ask_vision(c[2], c[3], frame_note))

    new_records: list[dict] = []
    reviewed = 0
    for (s_idx, slide, _png, _inv), layout in zip(candidates, answers):
        if isinstance(layout, Exception):
            # A slide the model could not answer for is SKIPPED, never counted
            # as reviewed and never recorded as clean: the two mean opposite
            # things to a designer reading "0 findings". LLMUnavailable arrives
            # here too - one bad call must not cost the whole run.
            continue
        reviewed += 1
        # Synthesized in slide order, on the main thread. Verification reads
        # geometry off the Presentation, which is not thread-safe to share, and
        # the records must come out in a fixed order anyway.
        new_records.extend(synthesize(slide, s_idx, layout, space,
                                      existing + new_records,
                                      slide_w, slide_h))
    return new_records, reviewed
