"""Design copilot: a model looks at rendered slides and proposes layout
actions a designer would make; code verifies and executes them.

Strict division of labor (the lesson of the ground-truth calibration):
- The model (vision) supplies JUDGMENT: which shapes should be distributed,
  aligned, or size-matched. It chooses from a CLOSED action vocabulary and
  may only reference shape ids from the inventory we hand it.
- Code supplies PRECISION: every observation is re-verified against the
  actual geometry and materialized as an ordinary FindingRecord with a
  computed target, so it flows through the existing pipeline - tickable
  (never pre-selected), collision-guarded, before/after rendered.
- The designer supplies APPROVAL, as with every other fix.

WHICH model is qc.llm's business, not this file's (30/08/2026). This pass named
its own vendor and its own model id until then, which made it the one judgment
pass in the package without a timeout, without retries and without a truncation
guard - see _ask_vision.

Confidentiality: unlike the metadata-only assistant, this sends SLIDE
IMAGES to the configured provider. The UI says so explicitly; use it only on
decks approved for cloud processing.
"""

import io
import json

from pptx import Presentation

from .llm import ask_in_parallel, ask_json
from .records import make_record

MAX_SLIDES = 20
MAX_OBS_PER_SLIDE = 3
TOL_EMU = 28575           # perceptual floor, same as calibrated profiles

# How far a shape may be off a shared edge and still be SNAPPED to it in one
# press. Past this the finding is still reported; it is just not offered as a
# computed move, because a move this big is more likely to be the model having
# grouped the wrong shapes than a designer having missed by that much.
#
# This replaced a 0.15in ceiling that DROPPED anything further out (31/08/2026).
# 0.15in is a perceptual threshold - it answers "were these meant to line up?",
# which is the right question in qc.modules.margin_alignment where nothing else
# supplies intent. Here the model has ALREADY supplied the intent, so reusing it
# inverted the pass: the further off a shape was, the more certain the silence,
# and the defects a designer sees from across the room were exactly the ones
# thrown away. A label sitting 0.35in left of the block it heads was reported by
# the model on every run and discarded by this line every time.
SNAP_MAX_EMU = 914400     # 1in: a designer's miss, not a mis-grouping

ACTIONS = ("distribute_row", "distribute_col", "align_left", "align_top",
           "match_widths", "match_heights")

OBSERVATIONS_SCHEMA = {
    "type": "object", "additionalProperties": False,
    "required": ["observations"],
    "properties": {
        "observations": {
            "type": "array",
            "items": {
                "type": "object", "additionalProperties": False,
                "required": ["action", "shape_ids", "rationale"],
                "properties": {
                    "action": {"enum": list(ACTIONS)},
                    "shape_ids": {"type": "array",
                                  "items": {"type": "string"}},
                    "rationale": {"type": "string"},
                },
            },
        },
    },
}

_SYSTEM = """You are a senior presentation designer at Prezlab reviewing a
slide for layout consistency. You see the rendered slide image and an
inventory of its shapes (id, kind, position and size as fractions of the
slide, whether it holds text).

Propose at most {max_obs} layout actions a designer would make, chosen ONLY
from: distribute_row / distribute_col (spread a line of shapes to equal
gaps), align_left / align_top (line up edges that should match),
match_widths / match_heights (make sibling cards the same size). Reference
shapes ONLY by ids from the inventory.

distribute_* and match_* need at least three shapes. align_left and
align_top take TWO or more, because the commonest real misalignment is a
pair: a label sitting over the block it heads, a caption under its image, a
column heading over its column. If one column's label sits square with its
body and another's does not, that is exactly the case to report.

Judge what you SEE. Do not skip something because the gap looks large; a
badly placed element is still a misalignment, and how far to move it is not
your problem. Equally, one item in an otherwise even row that sits lower, or
smaller, than its neighbours is worth reporting even when the rest are
perfect.

Rules: propose an action only when the improvement would be clearly
visible and safe; analogous elements (a row of cards, a set of columns)
are good targets, decorative compositions are not. If the slide already
looks professionally composed, return an empty list. Write rationales in
clear US English without em dashes."""


def inventory(slide, slide_w: int, slide_h: int) -> list[dict]:
    """Shape inventory the model reasons over: ids, kind, normalized geometry.
    No text content - the image already shows it; the inventory stays
    minimal."""
    out = []
    for shape in slide.shapes:
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
        if None in (l, t, w, h):
            continue
        has_text = bool(getattr(shape, "has_text_frame", False)
                        and shape.text_frame.text.strip())
        out.append({
            "id": str(shape.shape_id),
            "kind": str(shape.shape_type).split(" ")[0].lower(),
            "x": round(l / slide_w, 3), "y": round(t / slide_h, 3),
            "w": round(w / slide_w, 3), "h": round(h / slide_h, 3),
            "text": has_text,
        })
    return out


def _ask_vision(png: bytes, inv: list[dict]) -> list[dict]:
    """The one call this module makes. Which model answers is qc.llm's business
    and no part of this file's - the question and the schema are.

    THIS USED TO BUILD ITS OWN MODEL CLIENT (until 30/08/2026), which is the
    exact drift qc.llm exists to prevent and which its docstring already claimed
    was not happening. Going direct meant this pass alone had: no timeout, so a
    half-open connection pinned a worker until the process died; no retry, so a
    429 on slide 7 of 20 read as "slide 7 is fine"; no MAX_TOKENS guard, so a
    truncated reply came out of the route as a 500 rather than a skipped slide;
    and a second place where a model id lived.
    """
    parsed = ask_json(
        system=_SYSTEM.format(max_obs=MAX_OBS_PER_SLIDE),
        prompt="Shape inventory:\n" + json.dumps(inv, sort_keys=True),
        schema=OBSERVATIONS_SCHEMA,
        images=[png],
    )
    return (parsed.get("observations") or [])[:MAX_OBS_PER_SLIDE]


def _shared_edge(shapes, edge) -> int:
    """The edge the set is meant to share, in EMU.

    Three or more: the MEDIAN, because a majority already sitting on a line is
    the strongest statement of where the line is.

    Exactly two: the edge of the LARGER shape. A pair carries no majority, and
    taking the median of two would pick whichever value happens to sort second
    - so half the time a label would drag the block it heads sideways instead
    of moving to it. Size is the tie-break that matches how the pair reads: a
    caption belongs to its image, an eyebrow to its heading, a label to its
    column. The big one is the spine and the small one is the satellite.
    """
    if len(shapes) > 2:
        vals = sorted(edge(s) for s in shapes)
        return vals[len(vals) // 2]
    spine = max(shapes, key=lambda s: s.width * s.height)
    return edge(spine)


def _claim(issue: str, shape_id, prop, locator) -> tuple:
    """What makes two records the same finding, for the de-duplication against
    what the audit already recorded.

    An EDGE record is claimed by its shape and its AXIS. "Shape 8's top edge is
    off the line it shares" is one finding however many passes notice it, and
    the axis is in `property` (both the measured module and this one write
    spPr.xfrm.off.y for a top). The locator cannot be the key any more because
    it now carries the alignment CLUSTER rather than the record's identity: a
    measured record states no cluster, so keying on it would let the same shape
    be reported twice on the same edge and offer the designer two cards that
    move it to the same place.

    Everything else keeps the locator, which is where a distribution states
    which row it is about.
    """
    if issue == "margin_alignment.edge_misaligned":
        return (issue, str(shape_id), prop)
    return (issue, str(shape_id), locator)


def synthesize(slide, s_idx: int, observations: list[dict],
               existing: list[dict]) -> list[dict]:
    """Code is the precision gate: re-verify each observation against the
    real geometry and emit ordinary FindingRecords with computed targets.
    Anything that does not check out is dropped silently."""
    by_id = {str(sh.shape_id): sh for sh in slide.shapes}
    seen = {_claim(r["issue_type"], r["shape_id"], r.get("property"),
                   r.get("locator"))
            for r in existing if r["slide_index"] == s_idx}
    out: list[dict] = []

    def emit(**kw):
        # Defaults, not fixed values: a caller that states its own confidence
        # (the alignment branch drops to "low" past the snap rail) has to win.
        kw = {"severity": "warning", "confidence": "medium", **kw}
        rec = make_record(slide_index=s_idx, action="flagged",
                          source="vision", **kw)
        key = _claim(rec.issue_type, rec.shape_id, rec.property, rec.locator)
        if key not in seen:
            seen.add(key)
            out.append(rec.to_dict())

    for obs in observations:
        shapes = [by_id.get(str(i)) for i in obs.get("shape_ids", [])]
        shapes = [s for s in shapes if s is not None
                  and None not in (s.left, s.top, s.width, s.height)
                  and not getattr(s, "rotation", 0)]
        action = obs.get("action")
        # Two shapes is a real alignment case and the commonest one there is: a
        # label over the block it heads, a caption under its image. Distributing
        # or size-matching a PAIR is meaningless - there is no rhythm in two
        # gaps and no sibling set in two shapes - so those keep the floor of
        # three. Requiring three everywhere is why the comparison slide went
        # unreported: the defect was a pair and a pair could not be expressed.
        if len(shapes) < (2 if action in ("align_left", "align_top") else 3):
            continue
        note = f"Design copilot: {obs.get('rationale', '').strip()[:160]}"

        if action in ("distribute_row", "distribute_col"):
            row = action == "distribute_row"
            run = (lambda s: s.left) if row else (lambda s: s.top)
            size = (lambda s: s.width) if row else (lambda s: s.height)
            shapes = sorted(shapes, key=run)
            gaps = [run(shapes[i + 1]) - (run(shapes[i]) + size(shapes[i]))
                    for i in range(len(shapes) - 1)]
            if any(g <= 0 for g in gaps):
                continue  # overlapping: not a distribution case
            if max(gaps) - min(gaps) <= TOL_EMU:
                continue  # already even to the eye
            ids = ",".join(str(s.shape_id) for s in shapes)
            emit(shape_id=shapes[0].shape_id, module="margin_alignment",
                 issue_type="margin_alignment.uneven_spacing",
                 locator=f"dist-{'row' if row else 'col'}:{ids}",
                 property="spPr.xfrm.off",
                 old_value=", ".join(str(g) for g in gaps),
                 new_value=sum(gaps) // len(gaps),
                 profile_rule_id="geometry.alignment.spacing_tolerance_emu",
                 message=f"{note} Distribute evenly; first and last stay put.")

        elif action in ("align_left", "align_top"):
            left = action == "align_left"
            edge = (lambda s: s.left) if left else (lambda s: s.top)
            target = _shared_edge(shapes, edge)
            # THE PEER SET TRAVELS WITH THE RECORD, and this is the whole
            # difference between a fix that works and one that makes the slide
            # worse. qc.fixer infers what must travel with a moving shape from
            # overlap and adjacency, and for a vertical move that means "what
            # sits beside it in the same row" - which, on a row of cards being
            # aligned to each other, is every other card. So the stray was
            # seated on the line and its neighbours, already on that line, were
            # dragged the same distance off it (reproduced 01/09/2026 on a
            # ten-circle grid: one circle came onto the line and two left it).
            #
            # Geometry cannot tell a satellite from a peer here. The model
            # already did: these ids are the set it said shares a line. Naming
            # them in the locator is how that answer reaches the fix, exactly
            # as qc.components names a component's members in "comp:".
            peers = ",".join(str(s.shape_id)
                             for s in sorted(shapes, key=lambda s: s.shape_id))
            cluster = f"align-{'x' if left else 'y'}:{peers}"
            for s in shapes:
                off = abs(edge(s) - target)
                if off <= TOL_EMU:
                    continue          # already on the line, to the eye
                # Past the snap rail the finding is REPORTED, never dropped.
                # Low confidence keeps it out of the one-click set
                # (qc.fixer.is_fixable), so a designer sees it and decides,
                # which is the right answer for a move this big.
                snappable = off <= SNAP_MAX_EMU
                emit(shape_id=s.shape_id, module="margin_alignment",
                     issue_type="margin_alignment.edge_misaligned",
                     confidence="medium" if snappable else "low",
                     property="spPr.xfrm.off.x" if left
                     else "spPr.xfrm.off.y",
                     locator=cluster,
                     old_value=edge(s),
                     new_value=int(target) if snappable else None,
                     profile_rule_id="geometry.alignment.edge_tolerance_emu",
                     message=(f"{note} Snap to the shared "
                              f"{'left' if left else 'top'} edge."
                              if snappable else
                              f"{note} It sits {off // 36000}mm off the "
                              f"{'left' if left else 'top'} edge the others "
                              f"share, which is too far to move for you: check "
                              f"whether it belongs on that line at all."))

        elif action in ("match_widths", "match_heights"):
            widths = action == "match_widths"
            dim = (lambda s: s.width) if widths else (lambda s: s.height)
            vals = sorted(dim(s) for s in shapes)
            median = vals[len(vals) // 2]
            for s in shapes:
                if abs(dim(s) - median) <= TOL_EMU:
                    continue
                target_w = int(median) if widths else s.width
                target_h = s.height if widths else int(median)
                emit(shape_id=s.shape_id, module="shape_size",
                     issue_type="shape_size.size_mismatch",
                     property="spPr.xfrm.ext",
                     old_value=f"{s.width}x{s.height}",
                     new_value=f"{target_w}x{target_h}",
                     profile_rule_id="shape_size.size_tolerance_emu",
                     message=f"{note} Match sibling "
                             f"{'widths' if widths else 'heights'}.")
    return out


def run_copilot(deck_bytes: bytes, thumbs: dict[int, bytes],
                manifest: dict) -> tuple[list[dict], int]:
    """Review up to MAX_SLIDES rendered slides; returns (new_records,
    slides_reviewed). API failures on individual slides are skipped so one
    bad call never sinks the run."""
    prs = Presentation(io.BytesIO(deck_bytes))
    existing = manifest.get("records") or []

    # The candidates are chosen first and the budget is spent on CALLS, not on
    # successes. The loop here used to stop once MAX_SLIDES had answered, which
    # meant that under an outage - when nothing answers - a 200-slide deck made
    # 200 failing calls to review nothing.
    candidates = []
    for s_idx, slide in enumerate(prs.slides):
        if len(candidates) >= MAX_SLIDES:
            break
        png = thumbs.get(s_idx)
        if png is None:
            continue
        inv = inventory(slide, prs.slide_width, prs.slide_height)
        if len(inv) < 3:
            continue
        candidates.append((s_idx, slide, png, inv))

    # Independent questions, asked together (qc.llm.ask_in_parallel).
    answers = ask_in_parallel(candidates, lambda c: _ask_vision(c[2], c[3]))

    new_records: list[dict] = []
    reviewed = 0
    for (s_idx, slide, _png, _inv), observations in zip(candidates, answers):
        if isinstance(observations, Exception):
            continue
        reviewed += 1
        # Verification runs in slide order on the main thread: it reads geometry
        # off the Presentation, which is not safe to share between threads.
        new_records.extend(
            synthesize(slide, s_idx, observations, existing + new_records))
    return new_records, reviewed
