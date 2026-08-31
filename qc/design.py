"""Design QC: the checks that only become answerable once a master is on.

The formatting pass rebuilds a deck on the client's master. That is the point at
which a new class of defect appears, and none of it existed in the file the
designer submitted:

    the master's palette meets the deck's own colours, and a navy that was
    consistent within the old design is now three near-identical navies;

    the master's backgrounds go under text that was written for a different
    ground, so a line that read fine on white is now grey on grey (the content
    migration deliberately refuses to touch colour - qc.migrate._resolve_collisions
    - because choosing one is a design judgment, and this is where that judgment
    gets made);

    and the master's own furniture lands on top of content the deck already
    had, or beside it, or outside the frame the master states.

Three properties separate this from the audit modules in qc/modules, and each is
the reason this is its own pass rather than another module there:

A FINDING IS A DECISION, NOT AN OCCURRENCE. Three spellings of the same navy on
forty shapes is ONE question - "which navy is the navy?" - and a designer answers
it once. FindingRecord is per shape by construction (Appendix A.2), so forty rows
is what it would produce, and a page that asks the same question forty times gets
the same answer forty times or, more likely, none.

EVERY FINDING CARRIES ITS OWN WAYS OUT. Low-contrast text can be fixed by
recolouring the text or by recolouring the ground, and those are different
decisions with different consequences - one changes a word, the other changes a
panel. The tool is not entitled to pick. It states the options, with what each
costs, and the designer picks (qc.remedy applies the pick).

AND NOTHING HERE IS EVER DONE UNASKED. Every finding's option list ends in
"leave it", which is a real answer and is recorded as one. A check that fires and
offers no way to decline it is a check that gets ignored.

Detection and proposal only: applying a choice, and taking it back, is qc.remedy.
"""

import hashlib
import io
import json
import math
from contextlib import contextmanager
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from spike.color_resolver import (ciede2000, clr_map, color_scheme,
                                  resolve_color_element)
from spike.ns import find
from spike.resolver import layout_ph_index

EMU_IN = 914400

# --- calibration ----------------------------------------------------------
#
# Every number here is a threshold a designer can be shown, because every one of
# them decides whether something appears on their screen or not.

# Two literal colours this close are the same colour spelled twice. CIEDE2000,
# so it is a perceptual distance and not a per-channel one: #1F3864 and #203965
# are 8 RGB units apart and indistinguishable, while #1F3864 and #1F3864 with a
# lumMod are far apart in RGB and obviously different on a slide.
SAME_INTENT_DELTA_E = 3.0
# Neutrals get a much tighter one, because greys are a CONTINUUM and a designer
# steps along it on purpose: #F2F2F2 is 2.6 deltaE from white and is a light
# panel, not a mis-typed white, while #FEFEFE is 0.2 away and is nothing but a
# mis-typed white. One threshold for both called the panel a typo.
NEUTRAL_SAME_INTENT_DELTA_E = 1.0
# Past this a colour is not a variant of anything in the palette; it is a colour
# the palette does not contain.
OFF_PALETTE_DELTA_E = 10.0

# Below this chroma a colour is a NEUTRAL - a grey, a near-black, an off-white -
# and it is not reported as off-palette however far it sits from the brand
# colours. A palette states the brand; the greys between them are the designer's
# to choose, and every deck has a dozen. Measured on the first smoke run: a
# panel at #888888 and a label at #999999 were both reported as "not a palette
# colour, nearest accent5 20 deltaE away", which is true, useless, and exactly
# the kind of row that trains a designer to stop reading the page.
#
# Two near-identical greys are still reported: that is an inconsistency, not a
# palette question, and it is the case this whole check exists for.
NEUTRAL_CHROMA = 6.0

# WCAG 2.1 AAA (design lead, 26/08/2026). AA - 4.5:1 body, 3:1 display - is
# written for a screen a reader controls: their room, their brightness, their
# distance. This work gets projected in a lit room onto a wall, at the back of
# which somebody is reading 11pt legal type, and AA passes plenty of pairings
# that fail there. AAA is the published next tier, so it is the bar to move to
# rather than a number invented here that nobody could look up.
#
# The large-text allowance stays, because it is part of the standard rather
# than generosity: display type at 18pt (or 14pt bold) clears at 4.5:1 where
# body copy needs 7:1. Without it every deck's cover headline is a finding, and
# a check that fires on every cover is a check that gets switched off.
#
# This does fire on brand pairings a client signed off under AA. That is the
# point of raising it, and every finding still ends in "leave it as it is",
# recorded as a decision.
CONTRAST_BODY = 7.0
CONTRAST_LARGE = 4.5
LARGE_PT = 18.0
LARGE_BOLD_PT = 14.0

# Where text stops being hard to read and starts being unreadable. Held at 3:1
# and NOT tied to CONTRAST_LARGE, which it used to be: with the bar at AA the
# two numbers happened to be equal, and raising the bar would have promoted
# every AA-passing warning to an error by accident. Severity is a claim about
# legibility, not about which standard is being applied.
UNREADABLE_RATIO = 3.0

# A fill this transparent is a scrim, not a ground: the colour under it shows
# through and no single background colour exists. Contrast is not judged there
# rather than judged wrongly - a 40% wash over a photograph is the standard way
# a designer FIXES contrast, and flagging it would be exactly backwards.
OPAQUE_ALPHA = 0.90

# How much of a text box another shape has to cover before it counts as that
# text's background (looking down) or as hiding it (looking up).
GROUND_COVER = 0.60
HIDES_COVER = 0.50
# Two text boxes sharing this much of the smaller one are printing on top of each
# other. Same figure as qc.migrate.MIN_TEXT_OVERLAP_SHARE, and deliberately: the
# migration reports these and declines to move them, and a designer who reads
# "text overlaps text" there and finds nothing here would be right to distrust
# both pages.
TEXT_OVERLAP_SHARE = 0.30

# A shape covering this much of the canvas is the ground the slide sits on. It
# overlaps everything by design and is never reported for it.
BACKDROP_SHARE = 0.35

# Position bin for grouping the same piece of furniture across slides: 0.1in.
# Fine enough that two different corner marks stay separate, coarse enough that
# the same badge nudged by a rounding error on slide 14 does not become its own
# finding.
POS_BIN = EMU_IN // 10

# A shape whose box is at least this far outside the stated frame is outside it.
# The slack is the same 2mm qc.util allows for bleed: designers park things ON a
# line, and a marker drawn by hand is never exactly on it either.
FRAME_SLACK = 72000

# Overlap is O(n^2), so a limit has to exist - but 120 was set before the pass
# was profiled and it was cutting the answer off, not the cost. Measured on a
# dense slide, 24/08/2026: at 200 shapes a cap of 120 reported 54 overlaps and a
# cap of 250 reported 204, so three quarters of them were being hidden to save
# 0.3s. Past ~250 the curve turns (300 shapes costs 4.5s), which is where the
# limit belongs. When it does bite it now SAYS SO on the slide: a silent cut
# reads as "nothing else here", which on the densest slides in a deck - the ones
# most likely to have overlaps - is the worst place to be quietly wrong.
MAX_PAIRWISE = 250


# --- the finding ----------------------------------------------------------


@dataclass
class Remedy:
    """One way out, and what taking it costs.

    `note` is not decoration. "Recolour the text" and "recolour the panel" are
    both correct and they are not interchangeable, and the only person who can
    say which is right for this slide is looking at the slide.
    """
    remedy_id: str
    label: str
    note: str
    op: str = ""          # what qc.remedy performs; "" for the leave-it option
    params: dict = field(default_factory=dict)


@dataclass
class DesignFinding:
    finding_id: str
    kind: str             # palette | contrast | overlap | frame
    headline: str
    detail: str
    severity: str         # error | warning | info
    slides: list[int]     # every slide it appears on, zero-based
    options: list[Remedy]
    evidence: dict = field(default_factory=dict)

    @property
    def places(self) -> int:
        return int(self.evidence.get("places") or len(self.slides))


def _plural(n: int, one: str, many: str | None = None) -> str:
    """"1 surface" / "3 surfaces". Every count on this page is read by a person
    deciding whether to click something, and "1 surface(s)" is the sound of a
    tool talking to itself."""
    return f"{n} {one if n == 1 else (many or one + 's')}"


def _finding_id(kind: str, key) -> str:
    """A finding's identity, derived from WHAT IT IS rather than from where it
    fell in a list.

    This is what lets the page re-run detection after a fix and still know which
    row the designer already answered: a sequence number would renumber every
    row the moment one finding disappeared, and the "applied" mark would slide
    onto its neighbour.
    """
    blob = json.dumps([kind, key], sort_keys=True, default=str)
    return hashlib.sha1(blob.encode("utf-8")).hexdigest()[:10]


# --- colour ---------------------------------------------------------------


def parse_hex(text) -> tuple[int, int, int] | None:
    h = str(text or "").lstrip("#").strip()
    if len(h) != 6:
        return None
    try:
        return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))
    except ValueError:
        return None


def hex_of(rgb) -> str:
    return "{:02X}{:02X}{:02X}".format(*rgb)


def luminance(rgb) -> float:
    """WCAG relative luminance."""
    def channel(c):
        c /= 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

    r, g, b = (channel(c) for c in rgb)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def contrast_ratio(a, b) -> float:
    """WCAG contrast ratio, 1.0 (identical) to 21.0 (black on white)."""
    la, lb = luminance(a), luminance(b)
    hi, lo = max(la, lb), min(la, lb)
    return (hi + 0.05) / (lo + 0.05)


def chroma(rgb) -> float:
    """How much colour a colour has, in Lab. Zero for every grey from black to
    white, and it is the only honest way to ask: RGB saturation calls #808080 and
    #FFFFFF equally unsaturated while treating a dark navy as vivid."""
    from spike.color_resolver import _srgb_to_lab

    _L, a, b = _srgb_to_lab(tuple(rgb))
    return (a * a + b * b) ** 0.5


def hue_distance(a, b) -> float:
    """Distance between two colours IGNORING how light each one is.

    The measure a recolour suggestion has to be ranked by, and the reason is the
    shape of the problem: fixing contrast means changing lightness on purpose, so
    ranking candidates by full perceptual distance ranks them by the very thing
    being changed. Asked to darken a mid grey it proposed a mid purple, which is
    genuinely the nearest colour overall and obviously the wrong answer. Keeping
    the hue and moving the lightness is what a designer would do.
    """
    from spike.color_resolver import _srgb_to_lab

    _l1, a1, b1 = _srgb_to_lab(tuple(a))
    _l2, a2, b2 = _srgb_to_lab(tuple(b))
    return ((a1 - a2) ** 2 + (b1 - b2) ** 2) ** 0.5


def _same_intent_tolerance(a, b) -> float:
    """How close two colours have to be before they count as one colour spelled
    twice. Tighter for neutrals - see NEUTRAL_SAME_INTENT_DELTA_E."""
    if chroma(a) < NEUTRAL_CHROMA and chroma(b) < NEUTRAL_CHROMA:
        return NEUTRAL_SAME_INTENT_DELTA_E
    return SAME_INTENT_DELTA_E


def _alpha_of(color_el) -> float:
    el = find(color_el, "a:alpha")
    if el is None or not el.get("val"):
        return 1.0
    try:
        return int(el.get("val")) / 100_000.0
    except ValueError:
        return 1.0


def _color_child(parent):
    """The colour element under a solidFill-shaped parent, or None."""
    for tag in ("a:srgbClr", "a:schemeClr", "a:sysClr", "a:prstClr",
                "a:scrgbClr", "a:hslClr"):
        el = find(parent, tag)
        if el is not None:
            return el
    return None


def _solid(parent, master):
    """(rgb, color_element) for a solid fill under `parent`, else (None, None).

    None also comes back for a solid fill that is not opaque: see OPAQUE_ALPHA.
    """
    solid = find(parent, "a:solidFill")
    if solid is None:
        return None, None
    el = _color_child(solid)
    if el is None or _alpha_of(el) < OPAQUE_ALPHA:
        return None, None
    return resolve_color_element(el, master), el


# --- fills, and what is behind them --------------------------------------


def _theme_bg_is_solid(master, idx: int) -> bool:
    """Whether the theme's background fill style at this p:bgRef index is a flat
    colour. Anything else (a gradient, most often) has no single value to judge
    contrast against."""
    from spike.color_resolver import _theme_element

    lst = find(_theme_element(master), ".//a:bgFillStyleLst")
    if lst is None:
        return False
    # p:bgRef idx is 1-based into bgFillStyleLst, offset by 1000.
    entries = list(lst)
    n = idx - 1000 if idx >= 1000 else idx
    if not (1 <= n <= len(entries)):
        return False
    return entries[n - 1].tag.endswith("}solidFill")


def _bg_of(container, master):
    """(rgb, why) for one container's own p:bg. rgb None means "this container
    states no usable background", and `why` says which kind of nothing."""
    from pptx.oxml.ns import qn

    element = getattr(container, "_element", None)
    if element is None:
        element = getattr(container, "element", None)
    cSld = element.find(qn("p:cSld")) if element is not None else None
    bg = find(cSld, "p:bg")
    if bg is None:
        return None, "none stated"

    ref = find(bg, "p:bgRef")
    if ref is not None:
        try:
            idx = int(ref.get("idx") or 0)
        except ValueError:
            idx = 0
        el = _color_child(ref)
        if el is None or not _theme_bg_is_solid(master, idx):
            return None, "a theme background that is not a flat colour"
        return resolve_color_element(el, master), "the theme background"

    pr = find(bg, "p:bgPr")
    if pr is None:
        return None, "none stated"
    if find(pr, "a:noFill") is not None:
        return None, "none stated"
    rgb, _el = _solid(pr, master)
    if rgb is not None:
        return rgb, "a flat background"
    return None, "a picture or gradient background"


def slide_ground(slide, master):
    """(rgb, where) for what a slide's own background is, resolved down the
    cascade the way PowerPoint paints it: the slide's own p:bg beats the
    layout's, which beats the master's, and a deck that states none anywhere
    lands on the theme's bg1 slot.

    rgb None means no flat colour exists at any level - a photographic or
    gradient ground - and contrast is not judged over it rather than judged
    against a guess.
    """
    for container, where in ((slide, "the slide"),
                             (slide.slide_layout, "the layout"),
                             (master, "the master")):
        rgb, why = _bg_of(container, master)
        if rgb is not None:
            return rgb, f"{where}'s background"
        if why != "none stated":
            return None, f"{where} carries {why}"
    slot = clr_map(master).get("bg1", "lt1")
    rgb = color_scheme(master).get(slot)
    if rgb is not None:
        return rgb, f"the theme's {slot}"
    return (255, 255, 255), "no stated background (assumed white)"


def _style_fill(shape, master):
    """The fill an autoshape gets from its p:style fillRef rather than from its
    own spPr. Every rectangle drawn in PowerPoint and left on its theme style is
    filled this way and has nothing in spPr at all, so a check that reads only
    spPr sees a transparent shape where the slide shows a solid accent panel."""
    style = find(shape._element, "p:style")
    ref = find(style, "a:fillRef")
    if ref is None:
        return None
    try:
        if int(ref.get("idx") or 0) == 0:
            return None  # idx 0 is explicitly "no fill"
    except ValueError:
        return None
    el = _color_child(ref)
    if el is None or _alpha_of(el) < OPAQUE_ALPHA:
        return None
    return resolve_color_element(el, master)


def _inherited_ph_fill(shape, slide, master):
    """A placeholder's fill from the layout and then the master placeholder it
    inherits from."""
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        idx = shape.placeholder_format.idx
    except Exception:
        return None
    for container in (slide.slide_layout, master):
        for ph in container.placeholders:
            try:
                if ph.placeholder_format.idx != idx:
                    continue
            except Exception:
                continue
            rgb, _el = _solid(find(ph._element, "p:spPr"), master)
            if rgb is not None:
                return rgb
    return None


def shape_fill(shape, slide, master):
    """(rgb, kind) for the colour a shape is actually painted.

    kind is "solid" for a flat colour this tool could repaint, "opaque" for a
    picture or pattern that hides what is under it but has no single value, and
    "clear" for a shape that lets the ground through.
    """
    spPr = find(shape._element, "p:spPr")
    if spPr is not None:
        if find(spPr, "a:noFill") is not None:
            return None, "clear"
        rgb, _el = _solid(spPr, master)
        if rgb is not None:
            return rgb, "solid"
        for tag, kind in (("a:blipFill", "opaque"), ("a:gradFill", "opaque"),
                          ("a:pattFill", "opaque")):
            if find(spPr, tag) is not None:
                return None, kind
        if find(spPr, "a:solidFill") is not None:
            return None, "clear"  # a solid fill, but see-through: not a ground
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return None, "opaque"
    rgb = _style_fill(shape, master)
    if rgb is not None:
        return rgb, "solid"
    rgb = _inherited_ph_fill(shape, slide, master)
    if rgb is not None:
        return rgb, "solid"
    return None, "clear"


# --- geometry -------------------------------------------------------------


@dataclass
class Placed:
    """One shape, in SLIDE coordinates, at its place in the paint order."""
    z: int                 # document order == paint order; higher is on top
    top: int               # index of its top-level ancestor
    shape: object
    box: tuple | None      # (l, t, r, b) in slide EMU
    grouped: bool = False


def _dimensions(shape):
    """(left, top, width, height), resolving placeholder inheritance ourselves.

    THE HOTTEST PATH IN THE PACKAGE. Reading .left on a placeholder that states
    no position of its own sends python-pptx up the inheritance chain:
    _effective_value -> _inherited_value -> _base_placeholder ->
    layout.placeholders.get(idx) - and that get() is a LINEAR SCAN which runs an
    lxml xpath ("./*[1]/p:nvPr/p:ph") on every candidate it passes. Four reads
    per box, every box, every pass. Measured on fixtures/large_200.pptx
    (30/08/2026): 96,760 xpath calls, 7.7s of a 19s design scan.

    Same answer python-pptx gives - a directly-applied value wins, otherwise the
    layout placeholder's - reached through the index spike.resolver already
    keeps for the font cascade. ONE index, shared, because this and the font
    cascade are asking the identical question ("which layout placeholder does
    this shape inherit from") and two maps would be two things to keep in step.
    A shape that states its own geometry never touches it at all, which is the
    common case; only an inheriting placeholder looks anything up.

    THE FIRST READ MUST BE OFF THE ELEMENT, not off the shape. `shape.left` is
    the very property that walks the inheritance chain, so asking it first and
    checking for None afterwards would pay the whole cost before deciding not
    to. `_element.x` is what BaseShape.left returns for an ordinary shape and is
    the directly-applied value for a placeholder: no inheritance, no xpath.
    """
    el = shape._element
    try:
        l, t, w, h = el.x, el.y, el.cx, el.cy
    except AttributeError:
        # A connector or graphic frame that does not carry the usual xfrm.
        return shape.left, shape.top, shape.width, shape.height
    if not (l is None or t is None or w is None or h is None):
        return l, t, w, h
    if not getattr(shape, "is_placeholder", False):
        return l, t, w, h
    try:
        base = layout_ph_index(shape.part.slide_layout).get(el.ph_idx)
    except (AttributeError, ValueError):
        return l, t, w, h
    if base is None:
        return l, t, w, h
    # getattr on the LAYOUT placeholder, so layout-to-master inheritance is
    # still python-pptx's to resolve. Only the slide-to-layout scan is ours.
    return (l if l is not None else base.left,
            t if t is not None else base.top,
            w if w is not None else base.width,
            h if h is not None else base.height)


def _raw_box(shape):
    l, t, w, h = _dimensions(shape)
    if None in (l, t, w, h) or w <= 0 or h <= 0:
        return None
    return (l, t, l + w, t + h)


def _group_transform(group):
    """(dx, dy, sx, sy) mapping this group's CHILD coordinates into its parent's.

    A shape inside a group carries its own a:off, and that offset is in the
    group's child coordinate space, not the slide's. Comparing it with a
    top-level shape's box - which is what an overlap check does - compares two
    different coordinate systems and produces overlaps that are not on the
    slide and misses ones that are. The group states both spaces (a:chOff /
    a:chExt against a:off / a:ext) and the mapping between them is this.
    """
    from pptx.oxml.ns import qn

    xfrm = find(group._element.find(qn("p:grpSpPr")), "a:xfrm")
    if xfrm is None:
        return (0, 0, 1.0, 1.0)
    off, ext = find(xfrm, "a:off"), find(xfrm, "a:ext")
    ch_off, ch_ext = find(xfrm, "a:chOff"), find(xfrm, "a:chExt")
    if None in (off, ext, ch_off, ch_ext):
        return (0, 0, 1.0, 1.0)

    def num(el, *names):
        try:
            return tuple(int(el.get(n)) for n in names)
        except (TypeError, ValueError):
            return None

    o, e = num(off, "x", "y"), num(ext, "cx", "cy")
    co, ce = num(ch_off, "x", "y"), num(ch_ext, "cx", "cy")
    if None in (o, e, co, ce) or ce[0] <= 0 or ce[1] <= 0:
        return (0, 0, 1.0, 1.0)
    sx, sy = e[0] / ce[0], e[1] / ce[1]
    return (o[0] - co[0] * sx, o[1] - co[1] * sy, sx, sy)


def _compose(outer, inner):
    """inner applied first, then outer."""
    odx, ody, osx, osy = outer
    idx, idy, isx, isy = inner
    return (odx + idx * osx, ody + idy * osy, osx * isx, osy * isy)


def _mapped(box, xf):
    if box is None:
        return None
    dx, dy, sx, sy = xf
    l, t, r, b = box
    return (int(l * sx + dx), int(t * sy + dy),
            int(r * sx + dx), int(b * sy + dy))


def placed_shapes(slide) -> list[Placed]:
    """Every shape on a slide, in paint order, with its box in slide
    coordinates. Group children are included and transformed; the group itself
    is listed too, because a group is what a designer selects and moves.

    Memoized only inside _placed_cache() - see the note on _PLACED_MEMO for why
    a cache of boxes has to be scoped rather than ambient. The identity check on
    read is the usual one: id() is reused once an object is collected."""
    if _PLACED_MEMO is not None:
        hit = _PLACED_MEMO.get(id(slide))
        if hit is not None and hit[0] is slide:
            return hit[1]
    out: list[Placed] = []

    def walk(shapes, xf, top, grouped):
        for shape in shapes:
            out.append(Placed(len(out), top, shape, _mapped(_raw_box(shape), xf),
                              grouped))
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes, _compose(xf, _group_transform(shape)),
                     top, True)

    for i, shape in enumerate(slide.shapes):
        walk([shape], (0, 0, 1.0, 1.0), i, False)
    if _PLACED_MEMO is not None:
        _PLACED_MEMO[id(slide)] = (slide, out)
    return out


# Set to a dict for the duration of ONE scan, and None the rest of the time.
#
# scan() runs five checks and four of them walk every slide, so the same
# traversal happens four times over a deck nothing is mutating. Caching it is
# worth about a third of the run - but a cache of BOXES is not the same animal
# as _MARKER_MEMO's cache of identities: a box goes stale the moment anything
# moves a shape, and qc.remedy and qc.fixer both move shapes.
#
# So it is not ambient. It exists only inside _placed_cache(), which scan holds
# open and closes in a finally, and every other caller of placed_shapes gets the
# uncached function. That is not a loss: qc.extract and qc.layoutgap ask once per
# slide and would never have hit the cache anyway. It also means nothing pins a
# Presentation in memory after the scan that opened it has returned.
_PLACED_MEMO: dict | None = None


@contextmanager
def _placed_cache():
    """Memoize placed_shapes per slide for the body of one scan."""
    global _PLACED_MEMO
    outer = _PLACED_MEMO           # nested scans keep the outer one honest
    _PLACED_MEMO = {}
    try:
        yield
    finally:
        _PLACED_MEMO = outer


def _cover(inner, outer) -> float:
    """How much of `inner` sits inside `outer`, 0.0 to 1.0."""
    if inner is None or outer is None:
        return 0.0
    x0, y0 = max(inner[0], outer[0]), max(inner[1], outer[1])
    x1, y1 = min(inner[2], outer[2]), min(inner[3], outer[3])
    if x1 <= x0 or y1 <= y0:
        return 0.0
    area = (inner[2] - inner[0]) * (inner[3] - inner[1])
    return ((x1 - x0) * (y1 - y0)) / area if area else 0.0


def _is_backdrop(box, slide_w, slide_h) -> bool:
    if box is None:
        return False
    return ((box[2] - box[0]) * (box[3] - box[1])
            >= BACKDROP_SHARE * slide_w * slide_h)


def _text_of(shape) -> str:
    try:
        return shape.text_frame.text.strip() if shape.has_text_frame else ""
    except Exception:
        return ""


def _label(shape, limit: int = 28) -> str:
    text = _text_of(shape)
    if text:
        flat = " ".join(text.split())
        return flat[:limit] + ("…" if len(flat) > limit else "")
    return shape.name or "a shape"


# Is this shape one of the master's frame markers? A pure read of its name and
# alt text - and the single most expensive question the pass asks, because the
# pairwise loops ask it about the same shape over and over: 68,484 calls for
# 1,300 shapes, 2.6s of an 8.1s scan (measured 24/08/2026 on a 26-slide deck).
# space_marker() reaches it through a descendant XPath for p:cNvPr, so every
# one of those calls walks a subtree.
#
# Memoized on the element rather than on the shape, because python-pptx builds a
# NEW wrapper object each time .shapes is iterated - caching on the wrapper
# would never hit. The stored element is compared by identity on read: ids are
# recycled after a garbage collection, and a stale id answering for a different
# shape would be a wrong answer rather than a slow one.
_MARKER_MEMO: dict[int, tuple] = {}


def _is_frame_marker(shape) -> bool:
    from .stylespec import space_marker

    element = shape._element
    hit = _MARKER_MEMO.get(id(element))
    if hit is not None and hit[0] is element:
        return hit[1]
    try:
        value = space_marker(shape) is not None
    except Exception:
        value = False
    _MARKER_MEMO[id(element)] = (element, value)
    return value


# A text box is almost never full of text, and the difference matters here in a
# way it does not to a rule check. A body placeholder is sized for the longest
# slide in the deck and then carries two lines; the empty four inches underneath
# are box, not words, and a caption placed there overlaps NOTHING a reader can
# see. On the spike's own "clean" fixture - a deck built to have no findings -
# that exact shape pair was reported as text printing over text, because the
# placeholder's BOX covers 37% of the caption while its text stops three inches
# above it.
#
# So the words get their own rectangle, estimated: how many lines the text wraps
# to at its resolved size, from whichever edge the frame anchors to. It is an
# estimate and it is stated as one - PowerPoint's own line breaking depends on
# the font's metrics, which are not in the file - but "the top 1.2in of a 4.95in
# box" is a far better answer than "all of it", and it errs toward the box: any
# paragraph whose size will not resolve keeps the full height.
# Average advance width as a share of the point size, and the number the whole
# estimate turns on: it decides how many characters fit on a line, and one
# invented line adds a WHOLE line height to the answer.
#
# 0.5 was a guess. Calibrated 24/08/2026 by asking the fonts: 5,040 cases of
# (string x face x size x box width) over the title and body strings from real
# decks, in Arial, Calibri, Georgia, Trebuchet, Segoe UI, Verdana and Tahoma,
# regular and bold, comparing the line count the FONT gives against the line
# count this model gives.
#
#     em/char   invented wrap   missed wrap   exact
#      0.44          0.5%          11.9%      87.7%
#      0.46          1.1%           9.2%      89.7%
#      0.48          1.6%           7.8%      90.5%
#      0.50 (old)    4.0%           5.3%      90.8%
#      0.52          5.8%           4.3%      89.9%
#
# 0.46 because the two errors do not cost the same. An INVENTED wrap adds a
# whole line height and reports a title that fits as one that spills, which is
# the finding a designer stops trusting the page over (design lead, 24/08/2026:
# "but nothing is spilling out?"). A MISSED wrap loses one finding on a box the
# designer is looking at anyway, on a card that says out loud it is an estimate.
# Moving 0.50 -> 0.46 cuts invented wraps by nearly four times and costs about a
# point of overall accuracy. Same direction the text-extent estimate below
# already errs in.
CHAR_WIDTH_EM = 0.46
# Single spacing as PowerPoint renders it, used only when the paragraph does not
# state its own (see _line_spacing_pt).
LINE_SPACING = 1.2
FRAME_INSET_EMU = 91440  # the default 0.05in left+right inset, doubled


def _line_spacing_pt(para, size_pt: float) -> float:
    """How tall ONE line of this paragraph is, in points.

    Reads the paragraph's own a:lnSpc, which the estimate used to ignore
    entirely - it assumed 1.2 for everything. Display titles are routinely set
    at 85-90%, and a title measured at 1.2 when it renders at 0.9 is estimated a
    third taller than it is, which is enough on its own to report a title that
    fits as one that spills (design lead, 24/08/2026).

    Two forms in the schema and both appear in real decks: a:spcPct val is
    thousandths of a percent (90000 = 90%), a:spcPts val is hundredths of a
    point (2400 = an exact 24pt line).

    Only the paragraph's OWN setting is read, not one inherited from the layout
    or master list styles. That leaves the inherited case on the 1.2 default,
    which over-estimates exactly as before - the safe direction here, since this
    number only ever grows the estimate, and a missed overflow beats an invented
    one. Worth extending through spike.resolver's cascade when there is evidence
    a client master sets it there.
    """
    pPr = para._p.find(qn("a:pPr"))
    lnSpc = pPr.find(qn("a:lnSpc")) if pPr is not None else None
    if lnSpc is not None:
        pct = lnSpc.find(qn("a:spcPct"))
        if pct is not None and pct.get("val"):
            try:
                return size_pt * (int(pct.get("val")) / 100000.0)
            except ValueError:
                pass
        pts = lnSpc.find(qn("a:spcPts"))
        if pts is not None and pts.get("val"):
            try:
                return int(pts.get("val")) / 100.0
            except ValueError:
                pass
    return size_pt * LINE_SPACING


def natural_text_height(shape, box, slide, prs):
    """How tall this shape's text WANTS to be, in EMU, ignoring the box.

    None when it cannot be estimated, which is a real answer and is treated as
    one everywhere: a paragraph whose size will not resolve through the cascade
    gets no guess made about it.

    Unclamped on purpose. _text_extent needs it clipped to the box (words cannot
    be drawn outside the box they are in), and the fit check needs the raw number
    (the difference between the two IS the overflow).
    """
    if box is None or not getattr(shape, "has_text_frame", False):
        return None
    frame = shape.text_frame
    paragraphs = [p for p in frame.paragraphs if p.text.strip()]
    if not paragraphs:
        return None

    from spike.resolver import resolve_run

    width = max(1, (box[2] - box[0]) - FRAME_INSET_EMU)
    wraps = getattr(frame, "word_wrap", None) is not False
    total_pt = 0.0
    for para in paragraphs:
        size = 0.0
        for run in para.runs:
            if not run.text.strip():
                continue
            try:
                size = max(size, float(resolve_run(run, para, shape, slide,
                                                   prs).size_pt.value))
            except Exception:
                return None
        if not size:
            return None
        # A box with wrap="none" holds one line per paragraph however long the
        # text is, and that line runs PAST the box sideways. The height estimate
        # is right; the width is understated, because guessing how far the
        # overflow reaches is a worse guess than the box itself. Understating
        # means a missed overlap rather than an invented one, which is the
        # direction to be wrong in on a page that offers to move things.
        # Float, not int(). Truncating the capacity added a second downward bias
        # on top of the character-width one: a line with room for 17.9
        # characters became room for 17, so a 19-character title was declared to
        # wrap and picked up a whole extra line height. The capacity is an
        # estimate either way, and rounding it down is not caution - it is a
        # thumb on the scale in the direction of inventing lines.
        per_line = max(1.0, width / (size * CHAR_WIDTH_EM * 12700))
        chars = len(para.text.strip())
        lines = 1 if not wraps else max(1, math.ceil(chars / per_line))
        total_pt += lines * _line_spacing_pt(para, size)
    return int(total_pt * 12700)


def _text_extent(shape, box, slide, prs):
    """The part of a text box its words actually occupy, or `box` when that
    cannot be estimated."""
    from pptx.enum.text import MSO_ANCHOR

    natural = natural_text_height(shape, box, slide, prs)
    if natural is None or box is None:
        return box
    frame = shape.text_frame
    height = min(box[3] - box[1], natural)
    if height <= 0:
        return box
    try:
        anchor = frame.vertical_anchor
    except Exception:
        anchor = None
    if anchor == MSO_ANCHOR.BOTTOM:
        return (box[0], box[3] - height, box[2], box[3])
    if anchor == MSO_ANCHOR.MIDDLE:
        mid = (box[1] + box[3]) // 2
        return (box[0], mid - height // 2, box[2], mid - height // 2 + height)
    return (box[0], box[1], box[2], box[1] + height)


# --- the ground under a piece of text ------------------------------------


def _ground_under(item: Placed, stack: list[Placed], slide, master, bg,
                  own_box=None):
    """(rgb, where) for what shows through behind one text shape.

    Read the way the slide is painted: the shape's own fill first, then
    whatever is drawn underneath it, then the slide's background. The first
    thing that actually stops the light is the answer.

    rgb None means the ground cannot be reduced to one colour - a photograph, a
    gradient, a translucent wash - and no contrast finding is made. That is a
    deliberate refusal, not a gap: the ratio over a photograph varies across the
    photograph, and one number for it would be fiction.
    """
    own, kind = shape_fill(item.shape, slide, master)
    if own is not None:
        return own, "its own fill"
    if kind == "opaque":
        return None, "its own picture or gradient fill"

    mine = own_box or item.box
    for other in sorted((p for p in stack if p.z < item.z),
                        key=lambda p: -p.z):
        if other.box is None or other.shape is item.shape:
            continue
        if _is_frame_marker(other.shape):
            continue
        if _cover(mine, other.box) < GROUND_COVER:
            continue
        rgb, other_kind = shape_fill(other.shape, slide, master)
        if rgb is not None:
            return rgb, f"{_label(other.shape)!r} behind it"
        if other_kind == "opaque":
            return None, f"{_label(other.shape)!r}, a picture or gradient"
    return bg


# --- text colour ---------------------------------------------------------


def _run_color(run, para, shape, slide, prs, master):
    """(rgb, source) for the colour a run is actually drawn in, resolved through
    the same cascade the font resolver walks (spike.resolver.rpr_layers), then
    through the theme's clrMap and any tint/shade on the reference.

    python-pptx answers None for every inherited colour, which is most of them:
    a master that states its body colour once is the normal case, and a check
    reading only the run would find no colour on a deck that plainly has one.
    """
    from spike.resolver import rpr_layers

    for source, rpr in rpr_layers(run, para, shape, slide, prs):
        if rpr is None:
            continue
        solid = find(rpr, "a:solidFill")
        if solid is None:
            continue
        el = _color_child(solid)
        if el is None:
            continue
        rgb = resolve_color_element(el, master)
        if rgb is not None:
            return rgb, source
    slot = clr_map(master).get("tx1", "dk1")
    rgb = color_scheme(master).get(slot)
    if rgb is not None:
        return rgb, f"theme {slot}"
    return (0, 0, 0), "assumed black"


def _run_highlight(run, master):
    """A run's own highlight, which IS its background when it has one."""
    rpr = find(run._r, "a:rPr")
    hl = find(rpr, "a:highlight")
    if hl is None:
        return None
    el = _color_child(hl)
    return resolve_color_element(el, master) if el is not None else None


def _threshold(size_pt: float, bold: bool) -> float:
    if size_pt >= LARGE_PT or (bold and size_pt >= LARGE_BOLD_PT):
        return CONTRAST_LARGE
    return CONTRAST_BODY


# --- palette -------------------------------------------------------------


def palette_of(profile_cfg: dict | None, master) -> dict[str, tuple]:
    """{label: rgb} for everything a colour on a slide could be trying to be:
    the profile's named colours first, then the master theme's own slots.

    The theme is included because a deck formatted onto a master is meant to use
    that master's colours whether or not anybody wrote them into a profile - and
    an ephemeral profile read straight off a master file names none of them.
    """
    out: dict[str, tuple] = {}
    for entry in (profile_cfg or {}).get("named_colors", []) or []:
        rgb = parse_hex(entry.get("hex"))
        if rgb:
            out[entry.get("name") or hex_of(rgb)] = rgb
    for slot, rgb in (color_scheme(master) or {}).items():
        out.setdefault(f"theme {slot}", rgb)
    return out


@dataclass
class ColorUse:
    rgb: tuple
    surface: str          # fill | line | text
    slide_index: int
    shape_id: str
    locator: str | None   # "p0/r1" for a run, None for a shape surface
    label: str


def _collect_colors(prs) -> list[ColorUse]:
    """Every LITERAL colour on the slides, with where it is.

    Literal only, and that is the whole discipline of this check. A surface
    painted with a theme reference is on-palette BY CONSTRUCTION - it moves when
    the theme moves, which is what a theme is for - and rewriting one to a hex
    would sever exactly that. The colours worth arguing about are the ones typed
    in by hand.
    """
    uses: list[ColorUse] = []
    for s_idx, slide in enumerate(prs.slides):
        master = slide.slide_layout.slide_master
        for item in placed_shapes(slide):
            shape = item.shape
            if _is_frame_marker(shape):
                continue
            spPr = find(shape._element, "p:spPr")
            for parent, surface in ((spPr, "fill"),
                                    (find(spPr, "a:ln"), "line")):
                solid = find(parent, "a:solidFill")
                el = find(solid, "a:srgbClr") if solid is not None else None
                if el is None:
                    continue
                rgb = resolve_color_element(el, master)
                if rgb:
                    uses.append(ColorUse(rgb, surface, s_idx,
                                         str(shape.shape_id), None,
                                         _label(shape)))
            if not getattr(shape, "has_text_frame", False):
                continue
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                for r_idx, run in enumerate(para.runs):
                    if not run.text.strip():
                        continue
                    el = find(find(run._r, "a:rPr"), "a:solidFill")
                    el = find(el, "a:srgbClr") if el is not None else None
                    if el is None:
                        continue
                    rgb = resolve_color_element(el, master)
                    if rgb:
                        uses.append(ColorUse(rgb, "text", s_idx,
                                             str(shape.shape_id),
                                             f"p{p_idx}/r{r_idx}",
                                             _label(shape)))
    return uses


def _targets(uses) -> list[dict]:
    return [{"slide_index": u.slide_index, "shape_id": u.shape_id,
             "surface": u.surface, "locator": u.locator} for u in uses]


def _theme_option(anchor_name: str, uses) -> list[Remedy]:
    """The better fix, offered when it exists: point the surface at the theme
    slot rather than at the slot's current hex.

    Only when the palette entry that matched IS a theme slot (palette_of labels
    those "theme accent1"). Writing the hex fixes today's colour; writing the
    reference means the surface follows the theme the next time the client
    rebrands, which is the difference between a deck that survives a palette
    change and one that has to be hunted through again.
    """
    if not anchor_name.startswith("theme "):
        return []
    slot = anchor_name.split(" ", 1)[1]
    return [Remedy(
        "theme", f"Point them at the theme's {slot} instead",
        f"Better than a hex where it applies: the surface then moves with the "
        f"theme, so a rebrand reaches it without anyone looking for it. Same "
        f"colour on screen today.",
        op="set_theme_color",
        params={"slot": slot, "targets": _targets(uses)})]


def _palette_findings(prs, palette) -> list[DesignFinding]:
    uses = _collect_colors(prs)
    if not uses:
        return []

    by_hex: dict[str, list[ColorUse]] = {}
    for u in uses:
        by_hex.setdefault(hex_of(u.rgb), []).append(u)
    # Most-used first: on a tie of perceptual distance the spelling the deck
    # actually prefers is the one to keep, and ordering here is what makes that
    # deterministic rather than dictionary order.
    ranked = sorted(by_hex.items(), key=lambda kv: (-len(kv[1]), kv[0]))

    findings = []
    for h, group in ranked:
        rgb = group[0].rgb
        near = sorted(((ciede2000(rgb, p), name, p)
                       for name, p in (palette or {}).items()),
                      key=lambda t: t[0])
        best = near[0] if near else None

        # This colour IS a palette colour. Nothing below may fire for it, and
        # the early exit is load-bearing rather than an optimisation: without
        # it, Brand Navy found its own near-identical twin further down and the
        # page offered to rewrite the palette value INTO the typo. Applying both
        # that and the variant fix swapped the two shapes and left the deck with
        # exactly the two spellings it started with (first apply/undo run).
        if best and hex_of(best[2]) == h:
            continue

        # Case 1: a hair off a palette colour. This is the "three navies"
        # finding and it is the common one: someone typed the hex from a PDF.
        if best and best[0] <= _same_intent_tolerance(rgb, best[2]):
            name, anchor = best[1], best[2]
            findings.append(DesignFinding(
                finding_id=_finding_id("palette", ["variant", h]),
                kind="palette", severity="warning",
                headline=f"#{h} is {name} spelled differently",
                detail=(f"#{h} is used on {_plural(len(group), 'surface')} and is "
                        f"visually the same color as {name} (#{hex_of(anchor)}), "
                        f"{best[0]:.1f} deltaE away - a difference no screen "
                        f"shows. Two spellings of one color is what makes a "
                        f"deck impossible to recolor later."),
                slides=sorted({u.slide_index for u in group}),
                evidence={"hex": h, "anchor": hex_of(anchor),
                          "anchor_name": name, "delta_e": round(best[0], 2),
                          "places": len(group),
                          "surfaces": sorted({u.surface for u in group})},
                options=[
                    Remedy("snap", f"Use {name} (#{hex_of(anchor)}) everywhere",
                           f"Rewrites {_plural(len(group), 'surface')} to the palette "
                           f"value. Nothing moves and nothing else changes "
                           f"color; on screen this is invisible, which is the "
                           f"point.",
                           op="set_color",
                           params={"hex": hex_of(anchor),
                                   "targets": _targets(group)}),
                ] + _theme_option(name, group) + [
                    Remedy("leave", "Leave #%s as it is" % h,
                           "Recorded as a decision. The deck keeps both "
                           "spellings and a later palette change has to be "
                           "made twice."),
                ]))
            continue

        # Case 2: near ANOTHER literal spelling but not near the palette. Two
        # hand-typed colours agreeing with each other and with nothing else.
        #
        # The comparison is strict and ordered - (uses, hex) has to be GREATER,
        # not greater-or-equal - so that of a pair used equally often, exactly
        # one of the two raises the finding. Greater-or-equal reported the pair
        # twice, once from each side, with each row proposing to unify onto the
        # other; a designer who accepted both got the two colours swapped.
        twin = next((other_h for other_h, other in ranked
                     if other_h != h
                     and (len(other), other_h) > (len(group), h)
                     and ciede2000(rgb, other[0].rgb)
                     <= _same_intent_tolerance(rgb, other[0].rgb)),
                    None)
        if twin:
            twin_uses = by_hex[twin]
            findings.append(DesignFinding(
                finding_id=_finding_id("palette", ["twin", h, twin]),
                kind="palette", severity="warning",
                headline=f"#{h} and #{twin} are the same color twice",
                detail=(f"#{h} ({_plural(len(group), 'surface')}) and #{twin} "
                        f"({_plural(len(twin_uses), 'surface')}) are perceptually "
                        f"identical and neither is in the palette. One of them "
                        f"is a typo or a leftover from another deck."),
                slides=sorted({u.slide_index for u in group}),
                evidence={"hex": h, "anchor": twin, "places": len(group),
                          "surfaces": sorted({u.surface for u in group})},
                options=[
                    Remedy("snap", f"Use #{twin} everywhere",
                           f"The spelling this deck uses more often "
                           f"({len(twin_uses)} against {len(group)}). Rewrites "
                           f"{_plural(len(group), 'surface')}.",
                           op="set_color",
                           params={"hex": twin, "targets": _targets(group)}),
                ] + ([Remedy("palette",
                             f"Use {best[1]} (#{hex_of(best[2])}) instead",
                             f"Brings both onto the palette rather than onto "
                             f"each other, at a visible shift of "
                             f"{best[0]:.1f} deltaE.",
                             op="set_color",
                             params={"hex": hex_of(best[2]),
                                     "targets": _targets(group + twin_uses)})]
                     if best and best[0] <= OFF_PALETTE_DELTA_E else [])
                + [Remedy("leave", "Leave both", "Recorded as a decision.")]))
            continue

        # Case 3: nowhere near anything the palette states. Neutrals excepted -
        # see NEUTRAL_CHROMA.
        if best and best[0] > OFF_PALETTE_DELTA_E \
                and chroma(rgb) >= NEUTRAL_CHROMA:
            findings.append(DesignFinding(
                finding_id=_finding_id("palette", ["off", h]),
                kind="palette", severity="warning",
                headline=f"#{h} is not a palette color",
                detail=(f"#{h} is used on {_plural(len(group), 'surface')}. The nearest "
                        f"thing the palette states is {best[1]} "
                        f"(#{hex_of(best[2])}), {best[0]:.1f} deltaE away - far "
                        f"enough to be a different color, not a variant of it. "
                        f"It may be deliberate; the palette does not say so."),
                slides=sorted({u.slide_index for u in group}),
                evidence={"hex": h, "anchor": hex_of(best[2]),
                          "anchor_name": best[1],
                          "delta_e": round(best[0], 2), "places": len(group),
                          "surfaces": sorted({u.surface for u in group})},
                options=[
                    Remedy("snap",
                           f"Replace with {best[1]} (#{hex_of(best[2])})",
                           f"A visible change: {best[0]:.1f} deltaE. Check it "
                           f"on the slide before you download.",
                           op="set_color",
                           params={"hex": hex_of(best[2]),
                                   "targets": _targets(group)}),
                    Remedy("leave", "Keep it off-palette",
                           "Recorded as a decision, so the next audit of this "
                           "deck is not a rerun of this conversation."),
                ]))
    return findings


# --- contrast ------------------------------------------------------------


def _contrast_findings(prs, palette) -> list[DesignFinding]:
    findings = []
    for s_idx, slide in enumerate(prs.slides):
        master = slide.slide_layout.slide_master
        stack = placed_shapes(slide)
        bg = slide_ground(slide, master)
        for item in stack:
            shape = item.shape
            if not getattr(shape, "has_text_frame", False):
                continue
            if not _text_of(shape) or _is_frame_marker(shape):
                continue
            words = _text_extent(shape, item.box, slide, prs)
            ground, where = _ground_under(item, stack, slide, master, bg, words)
            if ground is None:
                continue

            worst = None
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                for r_idx, run in enumerate(para.runs):
                    if not run.text.strip():
                        continue
                    behind = _run_highlight(run, master) or ground
                    rgb, _src = _run_color(run, para, shape, slide, prs, master)
                    try:
                        from spike.resolver import resolve_run
                        font = resolve_run(run, para, shape, slide, prs)
                        size = float(font.size_pt.value)
                        bold = bool(font.bold.value)
                    except Exception:
                        size, bold = 18.0, False
                    need = _threshold(size, bold)
                    ratio = contrast_ratio(rgb, behind)
                    if ratio >= need:
                        continue
                    if worst is None or ratio < worst["ratio"]:
                        worst = {"ratio": ratio, "need": need, "rgb": rgb,
                                 "ground": behind, "size": size, "bold": bold,
                                 "locator": f"p{p_idx}/r{r_idx}"}
            if worst is None:
                continue

            findings.append(_contrast_finding(
                s_idx, shape, worst, where, palette,
                _ground_owner(item, stack, slide, master, words)))
    return findings


def _ground_owner(item, stack, slide, master, own_box=None):
    """The shape whose fill IS this text's background, when one is - the thing a
    "recolour the panel instead" option would have to repaint. None when the
    ground is the slide's own background, which no per-shape fix can change."""
    own, _kind = shape_fill(item.shape, slide, master)
    if own is not None:
        return item.shape
    mine = own_box or item.box
    for other in sorted((p for p in stack if p.z < item.z), key=lambda p: -p.z):
        if other.box is None or _is_frame_marker(other.shape):
            continue
        if _cover(mine, other.box) < GROUND_COVER:
            continue
        rgb, kind = shape_fill(other.shape, slide, master)
        if rgb is not None:
            return other.shape
        if kind == "opaque":
            return None
    return None


def _passing(against, palette, need, prefer=None) -> list[tuple]:
    """(label, rgb, ratio) for palette colours that clear the bar against
    `against`, SAME HUE AS THE COLOUR BEING REPLACED first.

    Two wrong orderings came before this one and both are worth keeping written
    down, because each looked obviously right.

    Highest contrast first: asked to fix grey-on-grey it answered black, every
    time, on every finding. Maximum contrast is a correct answer to the WCAG
    question and the wrong answer to the design one.

    Nearest overall (CIEDE2000) next: asked to darken a mid grey it answered a
    mid purple, which really is the nearest palette colour to #BBBBBB and is
    absurd. Fixing contrast means moving lightness deliberately, so ranking by a
    distance that includes lightness ranks by the thing being changed.

    So: keep the hue, move the lightness (hue_distance), and among candidates of
    the same hue prefer the one the PROFILE names over a theme slot - a profile
    that states an Ink is stating the brand's ink, and the theme's dk1 is only
    ever the fallback. palette_of puts named colours first, so position in the
    palette is that preference.
    """
    rank = {name: i for i, name in enumerate(palette or {})}
    out = []
    for name, rgb in (palette or {}).items():
        ratio = contrast_ratio(rgb, against)
        if ratio >= need:
            drift = hue_distance(rgb, tuple(prefer)) if prefer else 0.0
            out.append((name, rgb, ratio, drift))
    out.sort(key=lambda t: (round(t[3], 1), rank.get(t[0], 999), -t[2]))
    return [(name, rgb, ratio) for name, rgb, ratio, _d in out]


def _best_available(against, palette, prefer=None) -> list[tuple]:
    """(label, rgb, ratio) for palette colours ranked by contrast against
    `against`, best first, whether or not they clear any bar.

    The counterpart to _passing, for the case _passing returns nothing. Ranked
    by ratio here rather than by hue, because this list is only ever consulted
    when nothing clears the bar and the question has narrowed from "which colour
    belongs here" to "how close can this get".
    """
    rank = {name: i for i, name in enumerate(palette or {})}
    out = [(name, rgb, contrast_ratio(rgb, against))
           for name, rgb in (palette or {}).items()]
    # Best contrast first; on a tie the palette's own order, which puts a
    # profile's named colours ahead of theme slots (_passing's preference, and
    # for the same reason).
    out.sort(key=lambda t: (-t[2], rank.get(t[0], 999)))
    return out


def _text_target(s_idx, shape, locator, rgb):
    return {"hex": hex_of(rgb),
            "targets": [{"slide_index": s_idx, "shape_id": str(shape.shape_id),
                         "surface": "text", "locator": locator}]}


def _contrast_finding(s_idx, shape, worst, where, palette, owner):
    """The card for one unreadable piece of text: recolour the words, repaint
    what they sit on, or decide to live with it.

    BOTH KINDS OF DECISION ARE ALWAYS OFFERED when they would change anything,
    and that took a second pass to get right. Every option used to be gated on
    clearing the bar, which was invisible at AA and wrong at AAA (design lead,
    26/08/2026): against a mid grey nothing in any palette reaches 7:1 and black
    itself only reaches about 5.3, so a grey-on-grey card came back offering to
    paint the panel black and nothing else. "Recolour the text" and "repaint the
    panel" change different things on the slide and only the designer can say
    which is right here, so a card that drops one of them has made the decision
    for them.

    A short option is still an option, as long as it says so: 1.2:1 becoming
    5.3:1 is most of the way, and a designer who can read "still under the 7:1
    the standard asks for" can weigh that against repainting a panel. What is
    never offered is an option that changes nothing.
    """
    text_hex, ground_hex = hex_of(worst["rgb"]), hex_of(worst["ground"])
    ratio, need = worst["ratio"], worst["need"]
    size_note = (f"{worst['size']:.0f}pt{' bold' if worst['bold'] else ''}")
    # Two lists, joined at the end: everything that CLEARS the bar, then
    # everything that only gets closer. Within each, cheapest first. A short
    # option ahead of a clearing one would leave the finding standing, and
    # auto_choice takes the first option.
    clears: list = []
    closer: list = []
    # Enough of a gain to be worth a click and worth a re-render.
    floor = ratio + 0.1

    # --- recolour the text ------------------------------------------------
    #
    # First because it changes the smallest thing on the slide that can fix the
    # problem.
    ink = _passing(worst["ground"], palette, need, prefer=worst["rgb"])
    ink_clears = bool(ink)
    short = ""
    if not ink:
        # Nothing clears the bar on this ground. Rank by contrast instead of by
        # hue, because the question has narrowed from "which colour belongs
        # here" to "how close can this get".
        ink = [(name, rgb, got)
               for name, rgb, got in _best_available(worst["ground"], palette)
               if got > floor]
        short = (f" Still under the {need:.1f}:1 the standard asks for: nothing "
                 f"in the palette clears it on this ground, which is a fact "
                 f"about the ground rather than about the text.")
    if ink:
        name, rgb, got = ink[0]
        (clears if ink_clears else closer).append(Remedy(
            "ink", f"Set the text to {name} (#{hex_of(rgb)})",
            f"Takes the ratio from {ratio:.1f}:1 to {got:.1f}:1. Recolors only "
            f"the runs that fail, so any deliberate accent color in this box "
            f"survives.{short}",
            op="set_color",
            params=_text_target(s_idx, shape, worst["locator"], rgb)))

    mono = max(((255, 255, 255), (0, 0, 0)),
               key=lambda c: contrast_ratio(c, worst["ground"]))
    mono_ratio = contrast_ratio(mono, worst["ground"])
    if mono_ratio > floor and (not ink or mono != ink[0][1]):
        mono_clears = mono_ratio >= need
        (clears if mono_clears else closer).append(Remedy(
            "mono", f"Set the text to #{hex_of(mono)}",
            f"{mono_ratio:.1f}:1, "
            + ("the safest available reading" if mono_clears
               else "the most any colour gets on this ground, and still under "
                    f"the {need:.1f}:1 bar")
            + ". Not a palette color"
            + (", and the palette has nothing here that passes."
               if not _passing(worst["ground"], palette, need) else "."),
            op="set_color",
            params=_text_target(s_idx, shape, worst["locator"], mono)))

    # --- repaint the ground ----------------------------------------------
    #
    # Only when a shape owns it: repainting a slide background to fix one label
    # is not a fix, it is a new deck.
    if owner is not None:
        grounds = _passing(worst["rgb"], palette, need, prefer=worst["ground"])
        ground_clears = bool(grounds)
        ground_short = ""
        if not grounds:
            grounds = [(name, rgb, got)
                       for name, rgb, got in _best_available(worst["rgb"],
                                                             palette)
                       if got > floor]
            ground_short = (f" Still under {need:.1f}:1; nothing in the palette "
                            f"clears it behind this text.")
        if grounds:
            name, rgb, got = grounds[0]
            (clears if ground_clears else closer).append(Remedy(
                "ground", f"Set {_label(owner)!r} to {name} (#{hex_of(rgb)})",
                f"Keeps the text color and repaints what it sits on: "
                f"{got:.1f}:1. Everything else in that shape changes ground "
                f"too, so check what else is on it.{ground_short}",
                op="set_color",
                params={"hex": hex_of(rgb),
                        "targets": [{"slide_index": s_idx,
                                     "shape_id": str(owner.shape_id),
                                     "surface": "fill", "locator": None}]}))

    options = clears + closer + [
        Remedy("leave", "Leave it as it is",
               "Recorded as a decision. Brand colors sometimes fail WCAG on "
               "purpose and that is the client's call, not this tool's.")]

    return DesignFinding(
        finding_id=_finding_id("contrast", [s_idx, str(shape.shape_id),
                                            text_hex, ground_hex]),
        kind="contrast",
        severity="error" if ratio < UNREADABLE_RATIO else "warning",
        headline=f"{_label(shape)!r} reads at {ratio:.1f}:1",
        detail=(f"#{text_hex} on #{ground_hex} ({where}) is {ratio:.1f}:1, "
                f"under the {need:.1f}:1 that {size_note} text needs. "
                + ("Below 3:1 text stops being readable rather than merely "
                   "being hard to read." if ratio < UNREADABLE_RATIO else
                   "Legible on a good screen in a dark room; not on a "
                   "projector in a lit one.")),
        slides=[s_idx],
        evidence={"ratio": round(ratio, 2), "need": need, "text": text_hex,
                  "ground": ground_hex, "pair": f"{text_hex}>{ground_hex}",
                  "size_pt": worst["size"], "places": 1, "where": where},
        options=options)


# --- what a render will not show ------------------------------------------
#
# Text this close in colour to the thing it sits on does not appear in a render
# AT ALL. Every vision pass in this tool is shown pictures, so a box of white
# text on a white ground is an EMPTY BOX to it, and every judgment about that
# slide is then made without the words: a layout question decides the slide has
# one block where it has two, a component question leaves a card's label out of
# the card, and both answers look reasonable.
#
# A model cannot be asked to read what is not in the image. It can be handed the
# fact, which is what qc.components.inventory does with this.
#
# 1.5:1 rather than the contrast check's bar, because this is a different
# question. The contrast check asks whether a reader can read it; this asks
# whether it is in the picture at all.
INVISIBLE_RATIO = 1.5


def invisible_text(slide, prs, master=None) -> dict:
    """{shape_id: ratio} for every text shape a render will not show.

    Resolved through exactly the same cascade the contrast check uses - the
    run's effective colour, the ground actually under it, highlights included -
    so the two passes cannot disagree about what colour a word is. The worst run
    in a shape decides it: a heading in white over a white panel is invisible
    even if the caption under it is not.
    """
    if master is None:
        master = slide.slide_layout.slide_master
    stack = placed_shapes(slide)
    bg = slide_ground(slide, master)
    out = {}
    for item in stack:
        shape = item.shape
        if not getattr(shape, "has_text_frame", False):
            continue
        if not _text_of(shape) or _is_frame_marker(shape):
            continue
        words = _text_extent(shape, item.box, slide, prs)
        ground, _where = _ground_under(item, stack, slide, master, bg, words)
        if ground is None:
            # No single ground colour - a photograph or a gradient. Whether the
            # words show against it is not a question this can answer, and
            # guessing "invisible" would tell the model a box is empty when it
            # may be perfectly legible.
            continue
        worst = None
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                behind = _run_highlight(run, master) or ground
                rgb, _src = _run_color(run, para, shape, slide, prs, master)
                ratio = contrast_ratio(rgb, behind)
                if worst is None or ratio < worst:
                    worst = ratio
        if worst is not None and worst < INVISIBLE_RATIO:
            out[str(shape.shape_id)] = round(worst, 2)
    return out


# --- overlap -------------------------------------------------------------


def _clearance(box, other, slide_w, slide_h):
    """(dx, dy) the smallest moves that would take `box` clear of `other`, each
    None when that direction runs off the canvas."""
    down = other[3] - box[1] + POS_BIN
    right = other[2] - box[0] + POS_BIN
    up = other[1] - box[3] - POS_BIN
    left = other[0] - box[2] - POS_BIN
    dy = None
    if box[3] + down <= slide_h:
        dy = down
    elif box[1] + up >= 0:
        dy = up
    dx = None
    if box[2] + right <= slide_w:
        dx = right
    elif box[0] + left >= 0:
        dx = left
    return dx, dy


def _move_options(s_idx, shape, box, other_box, slide_w, slide_h) -> list[Remedy]:
    dx, dy = _clearance(box, other_box, slide_w, slide_h)
    out = []
    if dy is not None:
        out.append(Remedy(
            "move_y", f"Move {_label(shape)!r} {abs(dy) / EMU_IN:.2f}in "
                      f"{'down' if dy > 0 else 'up'}",
            "The shortest move that clears it vertically. Only this shape "
            "moves; anything composed with it stays where it is, so check the "
            "result.",
            op="offset", params={"slide_index": s_idx,
                                 "shape_id": str(shape.shape_id),
                                 "dx": 0, "dy": int(dy)}))
    if dx is not None:
        out.append(Remedy(
            "move_x", f"Move {_label(shape)!r} {abs(dx) / EMU_IN:.2f}in "
                      f"{'right' if dx > 0 else 'left'}",
            "The shortest move that clears it horizontally.",
            op="offset", params={"slide_index": s_idx,
                                 "shape_id": str(shape.shape_id),
                                 "dx": int(dx), "dy": 0}))
    return out


def _pairwise_capped(s_idx, dropped: int, compared: int) -> DesignFinding:
    """The slide is too dense to compare every pair, and says so.

    Not a defect and not fixable - it is the pass telling a designer how far it
    got. A slide of 300 shapes silently checked to 120 of them looks exactly
    like a slide with nothing wrong, and the slides that trip this are the
    crowded ones where overlaps actually live."""
    return DesignFinding(
        finding_id=_finding_id("error", ["capped", s_idx, dropped]),
        kind="error", severity="info",
        headline=f"This slide is too crowded to check every pair",
        detail=(f"{compared + dropped} elements is past the {MAX_PAIRWISE} this "
                f"check compares pair by pair, so {dropped} of them were left "
                f"out and the overlap list for this slide is incomplete. The "
                f"ones carrying text were kept, because a pair with no text on "
                f"either side could not be reported anyway. Everything else on "
                f"this page - colour, contrast, fit, the frame - looked at the "
                f"whole slide."),
        slides=[s_idx], options=[], evidence={"places": dropped})


def _overlap_findings(prs) -> list[DesignFinding]:
    findings = []
    slide_w, slide_h = prs.slide_width, prs.slide_height
    for s_idx, slide in enumerate(prs.slides):
        master = slide.slide_layout.slide_master
        stack = [p for p in placed_shapes(slide)
                 if p.box is not None and not _is_frame_marker(p.shape)]
        # A group and its own children overlap by construction, and so do the
        # children with each other: that is what a group IS. Only shapes with
        # different top-level ancestors can be said to collide.
        stack = [p for p in stack
                 if not (p.shape.shape_type == MSO_SHAPE_TYPE.GROUP)]
        if len(stack) > MAX_PAIRWISE:
            # WHICH shapes get dropped matters, and taking the first N in
            # z-order was an arbitrary answer. This check only ever reports a
            # pair where at least one side carries text, so text-bearing shapes
            # are kept first: every pair dropped that way is a pair that could
            # not have produced a finding anyway.
            dropped = len(stack) - MAX_PAIRWISE
            stack = sorted(stack, key=lambda p: not _text_of(p.shape))
            stack = stack[:MAX_PAIRWISE]
            findings.append(_pairwise_capped(s_idx, dropped, len(stack)))
        extent = {p.z: (_text_extent(p.shape, p.box, slide, prs)
                        if _text_of(p.shape) else p.box) for p in stack}

        for i, a in enumerate(stack):
            for b in stack[i + 1:]:
                if a.top == b.top:
                    continue
                lower, upper = (a, b) if a.z < b.z else (b, a)
                low_text, up_text = _text_of(lower.shape), _text_of(upper.shape)
                if not low_text and not up_text:
                    continue  # two graphics overlapping is composition

                # Measured on where the WORDS are, not on the boxes that hold
                # them: see _text_extent.
                low_box = extent.get(lower.z, lower.box)
                up_box = extent.get(upper.z, upper.box)

                if low_text and up_text:
                    share = max(_cover(low_box, up_box),
                                _cover(up_box, low_box))
                    if share < TEXT_OVERLAP_SHARE:
                        continue
                    findings.append(_text_on_text(
                        s_idx, lower, upper, share, slide_w, slide_h))
                    continue

                if not low_text:
                    continue  # text drawn over a graphic is the normal case
                if _is_backdrop(upper.box, slide_w, slide_h):
                    continue  # a full-canvas wash over everything is a device
                _rgb, kind = shape_fill(upper.shape, slide, master)
                if kind not in ("solid", "opaque"):
                    continue
                if _cover(low_box, upper.box) < HIDES_COVER:
                    continue
                findings.append(_hidden_text(
                    s_idx, lower, upper, slide_w, slide_h))
    return findings


def _text_on_text(s_idx, lower, upper, share, slide_w, slide_h):
    # The shape drawn LAST is the one on top and the one to move: it is the one
    # a reader sees, and moving what is underneath it just swaps which line is
    # unreadable.
    return DesignFinding(
        finding_id=_finding_id("overlap", ["text", s_idx,
                                           str(lower.shape.shape_id),
                                           str(upper.shape.shape_id)]),
        kind="overlap", severity="warning",
        headline=f"{_label(upper.shape)!r} prints over "
                 f"{_label(lower.shape)!r}",
        detail=(f"The two boxes share {share * 100:.0f}% of the smaller one. "
                f"Which of them should give way is a layout decision, so "
                f"nothing was moved when the master was applied."),
        slides=[s_idx],
        evidence={"share": round(share, 3), "places": 1},
        options=_move_options(s_idx, upper.shape, upper.box, lower.box,
                              slide_w, slide_h)
        + [Remedy("leave", "Leave the arrangement alone",
                  "Recorded as a decision. Overlapping type is sometimes the "
                  "composition.")])


def _hidden_text(s_idx, lower, upper, slide_w, slide_h):
    cover = _cover(lower.box, upper.box)
    return DesignFinding(
        finding_id=_finding_id("overlap", ["hidden", s_idx,
                                           str(lower.shape.shape_id),
                                           str(upper.shape.shape_id)]),
        kind="overlap", severity="error",
        headline=f"{_label(lower.shape)!r} is hidden behind "
                 f"{_label(upper.shape)!r}",
        detail=(f"{_label(upper.shape)!r} is drawn after the text and covers "
                f"{cover * 100:.0f}% of it with a solid fill, so the words are "
                f"on the slide and cannot be read. Text that is present but "
                f"invisible survives every proofread."),
        slides=[s_idx],
        evidence={"cover": round(cover, 3), "places": 1},
        options=[
            Remedy("behind", f"Send {_label(upper.shape)!r} behind the text",
                   "Changes the drawing order only: nothing moves, nothing "
                   "resizes. Usually the right answer when the cover is a "
                   "panel the text belongs on.",
                   op="zorder", params={"slide_index": s_idx,
                                        "shape_id": str(upper.shape.shape_id),
                                        "below": str(lower.shape.shape_id)}),
            # The same repair from the other end, and not the same fix: sending
            # the cover back puts it behind EVERYTHING between them, while
            # bringing the text forward puts it in front of everything. On a
            # slide of stacked panels those are visibly different results, and
            # which one is wanted depends on what else is in the pile.
            Remedy("front", f"Bring {_label(lower.shape)!r} to the front",
                   "Also order-only. Choose this when the cover belongs where "
                   "it is in the stack and it is the text that was left "
                   "underneath.",
                   op="front", params={"slide_index": s_idx,
                                       "shape_id": str(lower.shape.shape_id)}),
        ] + _move_options(s_idx, lower.shape, lower.box, upper.box,
                          slide_w, slide_h)
        + [Remedy("leave", "Leave it covered", "Recorded as a decision.")])


# --- fit: text against the box, and the box against its card -------------
#
# The check the client's own slides asked for (design lead, 23/08/2026, three
# real slides). Every defect on them is a FIT defect and not one the existing
# passes look for:
#
#   a card's bullets running past the bottom of the card, the last line clipped
#   mid-sentence, because the text was written for a taller box;
#
#   a heading printing across the card's own edge into the column beside it,
#   because the box is wider than the card that holds it;
#
#   and a name sitting half outside the panel it belongs to.
#
# None of it is an overlap in the sense the overlap check means (nothing is
# hidden), none of it is an alignment error in the sense qc.modules.
# margin_alignment means (the CARDS line up fine), and all of it is the first
# thing a reader sees. It is about containment: does the text fit the box, and
# does the box fit the thing it sits in.

# Overflow has to clear both bars before it is worth a designer's attention: a
# tenth of the box AND a tenth of an inch. The estimate is not precise enough to
# report two millimetres, and a slide full of two-millimetre findings is a slide
# nobody reads.
OVERFLOW_SHARE = 0.10
OVERFLOW_MIN_EMU = 73152          # 0.08in
# How much of a text box has to sit inside a filled shape before that shape is
# the CARD the text belongs to rather than something it happens to cross.
CONTAINED_SHARE = 0.55
# And how far outside its card the text has to reach to be worth reporting. A
# designer parks type on a card's edge deliberately; a tenth of an inch past it
# is a decision, not an accident.
ESCAPE_MIN_EMU = 91440            # 0.10in
# A shape at least this much of the slide in BOTH directions is the ground the
# slide sits on, not a card sitting on it. See _container_of.
GROUND_SPAN = 0.80

# The smallest shrink worth offering. Below this the option is a fix that cannot
# fix: qc.remedy floors type at 8pt (unreadable is not a repair for overflowing),
# so a box needing 50% would come back at the floor and still overflow, having
# wrecked the type scale on the way. An overflow this bad is a copy-length
# problem, and the honest options are autofit, a bigger box, or shorter text.
MIN_SHRINK = 0.70

# PowerPoint's own defaults for a text box's internal padding. A box stating
# more than this was padded by hand, and that padding is room the text is not
# getting. Read from python-pptx, which resolves the default when a box states
# none: 0.1in left and right, 0.05in top and bottom.
DEFAULT_INSETS = {"left": 91440, "right": 91440, "top": 45720, "bottom": 45720}


def _stated_insets(shape) -> dict | None:
    try:
        frame = shape.text_frame
        return {"left": int(frame.margin_left), "right": int(frame.margin_right),
                "top": int(frame.margin_top),
                "bottom": int(frame.margin_bottom)}
    except Exception:
        return None


def _inset_relief(shape, box, slide, prs):
    """What resetting a hand-padded box's internal margins would give the text.

    (extra_width, extra_height, natural_after) in EMU, or None when the box's
    padding is already the default.

    Padding is the one dimension of a fit problem that costs nothing: it is
    invisible on the slide, so returning it changes neither the type size nor
    the shape the designer positioned. That is why it is offered before either.

    The overflow estimate does not count vertical padding at all - it compares
    the text's natural height against the whole box - so the height returned
    here is room the finding never charged for, and the remedy says so rather
    than quietly claiming a fix the number cannot show.
    """
    insets = _stated_insets(shape)
    if insets is None or box is None:
        return None
    extra_w = max(0, (insets["left"] + insets["right"])
                  - (DEFAULT_INSETS["left"] + DEFAULT_INSETS["right"]))
    extra_h = max(0, (insets["top"] + insets["bottom"])
                  - (DEFAULT_INSETS["top"] + DEFAULT_INSETS["bottom"]))
    if not extra_w and not extra_h:
        return None
    after = natural_text_height(shape, (box[0], box[1], box[2] + extra_w,
                                       box[3]), slide, prs)
    return extra_w, extra_h, after


def _widen_to_fit(shape, box, slide, prs, slide_w):
    """(dw, room) - the smallest widening that makes the text fit, and how much
    room there is beside the box. dw is None when nothing within that room fits
    it.

    Searched rather than solved. Line count is a step function of width: a box
    gets no shorter until it gains enough width to pull one more word up, so the
    honest answer is the first step that clears the box, found by walking the
    room in tenths. Solving it as if height were continuous in width would name
    a width that does not actually fit.
    """
    height = box[3] - box[1]
    room = max(0, int(slide_w) - box[2])
    if room <= 0:
        return None, 0
    step = max(EMU_IN // 10, room // 10)
    grown = step
    while grown <= room:
        after = natural_text_height(shape, (box[0], box[1], box[2] + grown,
                                           box[3]), slide, prs)
        if after is not None and after <= height:
            return int(grown), room
        grown += step
    return None, room




def _shrinks_to_fit(shape) -> bool:
    """Whether PowerPoint is already handling the overflow itself.

    a:normAutofit means "shrink text on overflow", and a box carrying it cannot
    overflow: PowerPoint scales the type down at render time. Reporting one would
    be reporting a problem the file already solves - and offering to shrink the
    type would be the second thing doing the same job."""
    try:
        body = find(shape.text_frame._txBody, "a:bodyPr")
    except Exception:
        return False
    return find(body, "a:normAutofit") is not None


def _container_of(item, stack, slide, master, slide_w, slide_h):
    """The filled shape a text box sits IN, or None.

    The same notion the contrast check calls the ground, asked with a different
    threshold and a different purpose: there, "what colour is behind these
    words"; here, "what are these words supposed to fit inside".

    A near-full-canvas ground is not a card. Everything on a slide with a bleed
    image sits on it, so counting it as a container makes every shape near an
    edge a finding: a white heading deliberately bleeding 0.2in off the left of
    a full-bleed navy panel was reported as running outside it (test, first fit
    run). The test is BOTH dimensions, not area - a full-height sidebar or a
    half-slide card is a real container that text can really escape, and text
    crossing the visible edge of one is exactly what this check is for.
    """
    for other in sorted((p for p in stack if p.z < item.z), key=lambda p: -p.z):
        if other.box is None or _is_frame_marker(other.shape):
            continue
        if getattr(other.shape, "has_text_frame", False) and _text_of(other.shape):
            continue  # a text box is not a card, whatever it is behind
        if (other.box[2] - other.box[0] >= GROUND_SPAN * slide_w
                and other.box[3] - other.box[1] >= GROUND_SPAN * slide_h):
            continue  # the ground the slide sits on, not a card on it
        if _cover(item.box, other.box) < CONTAINED_SHARE:
            continue
        _rgb, kind = shape_fill(other.shape, slide, master)
        if kind in ("solid", "opaque"):
            return other.shape, other.box
    return None, None


def _fit_findings(prs) -> list[DesignFinding]:
    findings = []
    slide_w, slide_h = prs.slide_width, prs.slide_height
    for s_idx, slide in enumerate(prs.slides):
        master = slide.slide_layout.slide_master
        stack = [p for p in placed_shapes(slide)
                 if p.box is not None and not _is_frame_marker(p.shape)
                 and p.shape.shape_type != MSO_SHAPE_TYPE.GROUP]
        for item in stack:
            shape, box = item.shape, item.box
            if not getattr(shape, "has_text_frame", False) or not _text_of(shape):
                continue
            natural = natural_text_height(shape, box, slide, prs)
            height = box[3] - box[1]

            # 1. The words do not fit the box they are in.
            if natural is not None and not _shrinks_to_fit(shape):
                over = natural - height
                if over > max(OVERFLOW_MIN_EMU, OVERFLOW_SHARE * height):
                    findings.append(_overflow_finding(
                        s_idx, shape, box, natural, over, slide, prs,
                        slide_w, slide_h))

            # 2. The box does not fit the card it sits in.
            card, card_box = _container_of(item, stack, slide, master,
                                           slide_w, slide_h)
            if card is None:
                continue
            words = _text_extent(shape, box, slide, prs) or box
            escape = {
                "left": card_box[0] - box[0],
                "right": box[2] - card_box[2],
                "top": card_box[1] - words[1],
                "bottom": words[3] - card_box[3],
            }
            worst = max(escape.items(), key=lambda kv: kv[1])
            if worst[1] > ESCAPE_MIN_EMU:
                findings.append(_escape_finding(
                    s_idx, shape, box, card, card_box, escape, worst,
                    slide_w, slide_h))
    return findings


def _overflow_finding(s_idx, shape, box, natural, over, slide, prs,
                      slide_w, slide_h):
    """Every way out of a box that will not hold its text, cheapest first.

    The order is the answer to "what does this cost the design", and it was
    wrong until 26/08/2026: autofit came first, so the tool's own
    recommendation - and the "let the tool decide" path, which takes the first
    option - was to make the type smaller. Shrinking type is the most expensive
    fix on the list. It breaks the deck's type scale, it is the one change a
    reader notices, and on a slide with an inch of white space beside the box it
    is unnecessary (design lead, 26/08/2026).

    So, cheapest first:

        1. the box's own internal padding, when it was set by hand. Invisible
           on the slide, so returning it costs nothing at all.
        2. more height, when there is room below. Keeps the measure - the line
           length the designer chose - and keeps the type.
        3. more width, when there is room beside and it actually fits. Keeps
           the type, but changes the measure and can break a column.
        4. autofit, which keeps working when the copy is edited later.
        5. an explicit shrink, when even autofit would be too much.

    Each is offered only when it would work, so the first option on the card is
    always the cheapest one that fixes THIS box.
    """
    height = box[3] - box[1]
    scale = height / natural
    room = slide_h - box[3]
    options = []

    relief = _inset_relief(shape, box, slide, prs)
    if relief is not None:
        extra_w, extra_h, after = relief
        if after is not None and after <= height + extra_h:
            options.append(Remedy(
                "insets",
                "Reset the box's internal margins to the default",
                f"This box pads its text by "
                f"{(extra_w + DEFAULT_INSETS['left'] + DEFAULT_INSETS['right']) / EMU_IN:.2f}in "
                f"across and "
                f"{(extra_h + DEFAULT_INSETS['top'] + DEFAULT_INSETS['bottom']) / EMU_IN:.2f}in "
                f"down, where PowerPoint's default is 0.20in and 0.10in. "
                f"Resetting returns {extra_w / EMU_IN:.2f}in of width and "
                f"{extra_h / EMU_IN:.2f}in of height to the words. Nothing "
                f"moves, nothing changes size on screen, and the type scale is "
                f"untouched, which is why it is first. The overflow figure "
                f"above does not count padding, so this fix is worth more than "
                f"that number suggests.",
                op="set_insets",
                params={"slide_index": s_idx,
                        "shape_id": str(shape.shape_id),
                        **DEFAULT_INSETS}))

    if room > over:
        options.append(Remedy(
            "grow", f"Make the box {over / EMU_IN:.2f}in taller",
            f"Keeps the type size and the line length, and takes the space "
            f"below, of which there is {room / EMU_IN:.2f}in before the slide "
            f"edge. Check what is under it first: growing a box does not move "
            f"its neighbours.",
            op="resize", params={"slide_index": s_idx,
                                 "shape_id": str(shape.shape_id),
                                 "dh": int(over)}))

    dw, side_room = _widen_to_fit(shape, box, slide, prs, slide_w)
    if dw:
        options.append(Remedy(
            "widen", f"Make the box {dw / EMU_IN:.2f}in wider",
            f"Keeps the type size and pulls the copy up into fewer lines: "
            f"measured, not assumed, so this width is one the text really fits "
            f"in. There is {side_room / EMU_IN:.2f}in to the slide edge. It "
            f"changes the line length, so on a column grid check the column "
            f"beside it before taking this one.",
            op="resize", params={"slide_index": s_idx,
                                 "shape_id": str(shape.shape_id),
                                 "dw": int(dw), "anchor": "left"}))

    options.append(Remedy(
        "autofit", "Let PowerPoint shrink the text to fit",
        "Writes the shrink-on-overflow setting, so the type scales down in the "
        "box and keeps scaling if the copy is edited later. It makes the type "
        "smaller, which is why it comes after the fixes that do not.",
        op="autofit", params={"slide_index": s_idx,
                              "shape_id": str(shape.shape_id)}))
    if scale >= MIN_SHRINK:
        options.append(Remedy(
            "shrink", f"Shrink the type to {scale * 100:.0f}% of its size",
            f"Sets each run's size explicitly. More predictable than autofit "
            f"and it survives a copy-paste into another deck, but it breaks the "
            f"deck's type scale on this one shape.",
            op="scale_text", params={"slide_index": s_idx,
                                     "shape_id": str(shape.shape_id),
                                     "scale": round(scale, 3)}))
    options.append(Remedy(
        "leave", "Leave it as it is",
        "Recorded as a decision. The estimate here is a calculation, not a "
        "render, and a line that just fits can read as one that just does not."))

    return DesignFinding(
        finding_id=_finding_id("fit", ["overflow", s_idx,
                                       str(shape.shape_id)]),
        kind="fit", severity="error" if over > height * 0.35 else "warning",
        headline=f"{_label(shape)!r} has more text than its box holds",
        detail=(f"The text needs about {natural / EMU_IN:.2f}in and the box is "
                f"{height / EMU_IN:.2f}in, so roughly {over / EMU_IN:.2f}in of "
                f"it falls outside. PowerPoint is not set to shrink it, so the "
                f"overflow prints over whatever is below or is clipped at the "
                f"box edge."
                + ("" if scale >= MIN_SHRINK else
                   f" There is too much of it for a type tweak to absorb: "
                   f"fitting it would mean {scale * 100:.0f}% of the current "
                   f"size, so this is a copy-length problem rather than a "
                   f"formatting one.")
                + f" Height is estimated from the resolved type sizes, not "
                  f"from a render."),
        slides=[s_idx],
        evidence={"needs_in": round(natural / EMU_IN, 2),
                  "box_in": round(height / EMU_IN, 2),
                  "over_in": round(over / EMU_IN, 2), "places": 1},
        options=options)


def _escape_finding(s_idx, shape, box, card, card_box, escape, worst,
                    slide_w, slide_h):
    side, amount = worst
    # The move that seats it back inside, on the axis it escaped.
    dx = dy = 0
    if side == "left":
        dx = amount
    elif side == "right":
        dx = -amount
    elif side == "top":
        dy = amount
    else:
        dy = -amount

    options = [Remedy(
        "seat", f"Move it {abs(amount) / EMU_IN:.2f}in back inside the card",
        f"The shortest move that puts it inside {_label(card)!r}. Only this "
        f"shape moves, so if it is one of a row of labels the others stay where "
        f"they are.",
        op="offset", params={"slide_index": s_idx,
                             "shape_id": str(shape.shape_id),
                             "dx": int(dx), "dy": int(dy)})]
    # Narrowing a text box reflows its text, which is why the migration refuses
    # to do it (qc.migrate, "wider than the margins"). Offered here anyway,
    # because on a card the reflow is usually what the designer wanted and the
    # alternative is type crossing a visible edge - but the note says so.
    if side in ("left", "right"):
        options.append(Remedy(
            "narrow", f"Narrow the box by {abs(amount) / EMU_IN:.2f}in to the "
                      f"card's edge",
            "Keeps the box's other edge where it is. The text REFLOWS, so line "
            "breaks change and the block may get taller; check it before you "
            "download.",
            op="resize", params={"slide_index": s_idx,
                                 "shape_id": str(shape.shape_id),
                                 "dw": -int(abs(amount)),
                                 "anchor": "right" if side == "left" else "left"}))
    else:
        grow = abs(amount)
        options.append(Remedy(
            "grow_card", f"Make {_label(card)!r} {grow / EMU_IN:.2f}in taller "
                         f"to hold it",
            "Changes the card, not the text. Right when the card was drawn too "
            "small; wrong when the card's size is the grid, because this one "
            "will no longer match its neighbours.",
            op="resize", params={"slide_index": s_idx,
                                 "shape_id": str(card.shape_id),
                                 "dh": int(grow),
                                 "anchor": "bottom" if side == "top" else "top"}))
    options.append(Remedy(
        "leave", "Leave it crossing the edge",
        "Recorded as a decision. Type deliberately breaking a card's edge is a "
        "real device."))

    where = {"left": "past its left edge", "right": "past its right edge",
             "top": "above its top edge", "bottom": "below its bottom edge"}[side]
    return DesignFinding(
        finding_id=_finding_id("fit", ["escape", s_idx, str(shape.shape_id),
                                       str(card.shape_id), side]),
        kind="fit",
        severity="error" if amount > EMU_IN // 4 else "warning",
        headline=f"{_label(shape)!r} runs outside {_label(card)!r}",
        detail=(f"The text box reaches {amount / EMU_IN:.2f}in {where}, so the "
                f"words cross a visible boundary and read as belonging to "
                f"whatever is on the other side of it. The card itself is "
                f"where the designer put it; it is the text box that does not "
                f"fit inside it."),
        slides=[s_idx],
        evidence={"escape_in": round(amount / EMU_IN, 2), "side": side,
                  "places": 1},
        options=options)


# --- outside the stated frame -------------------------------------------


def _frame_findings(prs) -> list[DesignFinding]:
    """Text sitting outside the frame the master states.

    Nothing is removed for this, here or in the migration, and the reason is
    that the tool cannot tell a leftover from a deliberate one. A page number, a
    source line and a client's own corner badge all live outside the content
    frame on purpose; a stranded eyebrow from the deck's previous design lives
    there by accident, and they look identical from the XML.

    So the answer is to say what is out there and let the person who knows
    decide - which is also why identical furniture is collapsed into ONE finding
    across the whole deck. A numbered badge on twenty-one slides is one
    question, and asking it twenty-one times is how a designer learns to skip
    this page.
    """
    from .stylespec import dominant_master, read_presentation_space

    master = dominant_master(prs)
    space = read_presentation_space(prs, master) if master is not None else None
    if not space or space.get("problem") or not space.get("box_emu"):
        return []
    fl, ft, fr, fb = space["box_emu"]
    slide_w, slide_h = prs.slide_width, prs.slide_height

    groups: dict[tuple, list] = {}
    for s_idx, slide in enumerate(prs.slides):
        for item in placed_shapes(slide):
            shape, box = item.shape, item.box
            if box is None or item.grouped or _is_frame_marker(shape):
                continue
            if getattr(shape, "is_placeholder", False):
                continue  # the master's own furniture, where the master put it
            if not _text_of(shape):
                continue  # a rule or a mark outside the frame is a composition
            if _is_backdrop(box, slide_w, slide_h):
                continue
            inside = _cover(box, (fl, ft, fr, fb))
            if inside > 0.5:
                continue
            if (box[0] >= fl - FRAME_SLACK and box[1] >= ft - FRAME_SLACK
                    and box[2] <= fr + FRAME_SLACK
                    and box[3] <= fb + FRAME_SLACK):
                continue
            key = (shape.name or "", box[0] // POS_BIN, box[1] // POS_BIN,
                   _text_of(shape).isdigit())
            groups.setdefault(key, []).append((s_idx, shape, box, inside))

    findings = []
    for key, members in sorted(groups.items(), key=lambda kv: -len(kv[1])):
        slides = sorted({m[0] for m in members})
        first = members[0]
        box = first[2]
        numeric = key[3]
        what = ("a number" if numeric else f"{_label(first[1])!r}")
        # The move that brings it inside: the shortest translation that seats
        # its own top-left on the frame's nearest corner region.
        dx = max(0, fl - box[0]) or min(0, fr - box[2])
        dy = max(0, ft - box[1]) or min(0, fb - box[3])
        options = []
        if dx or dy:
            options.append(Remedy(
                "inside", f"Move {'each' if len(members) > 1 else 'it'} inside "
                          f"the frame ({dx / EMU_IN:+.2f}in, {dy / EMU_IN:+.2f}in)",
                f"Moves {_plural(len(members), 'shape')} the shortest distance that "
                f"puts them inside the presentation space. They may then "
                f"overlap what is already there - the overlap check below runs "
                f"again after you apply this.",
                op="offset_many",
                params={"dx": int(dx), "dy": int(dy),
                        "targets": [{"slide_index": m[0],
                                     "shape_id": str(m[1].shape_id)}
                                    for m in members]}))
        options.append(Remedy(
            "remove", f"Remove {'them' if len(members) > 1 else 'it'} "
                      f"({_plural(len(members), 'shape')})",
            "Deletes the shape and keeps its own XML, so Undo puts back the "
            "same words in the same box with the same formatting. Use this for "
            "furniture the new master already provides.",
            op="delete",
            params={"targets": [{"slide_index": m[0],
                                 "shape_id": str(m[1].shape_id)}
                                for m in members]}))
        options.append(Remedy(
            "leave", "Leave it outside the frame",
            "Recorded as a decision. Page numbers, source lines and corner "
            "badges belong outside the content frame."))

        findings.append(DesignFinding(
            finding_id=_finding_id("frame", list(key)),
            kind="frame", severity="info",
            headline=f"{what} sits outside the presentation space "
                     f"on {_plural(len(slides), 'slide')}",
            detail=(f"At {box[0] / EMU_IN:.2f}in, {box[1] / EMU_IN:.2f}in, "
                    f"{'wholly' if first[3] == 0 else f'{(1 - first[3]) * 100:.0f}% '}"
                    f"outside the frame the master states, and in no "
                    f"placeholder. The formatting pass left it alone: it only "
                    f"removes text the whole deck treats as a stray "
                    f"(qc.migrate.stray_texts), and repeated furniture is not "
                    f"that. Whether it belongs is a question about the master, "
                    f"not about this slide."),
            slides=slides,
            evidence={"places": len(members), "inside_share": round(first[3], 3),
                      "left_in": round(box[0] / EMU_IN, 2),
                      "top_in": round(box[1] / EMU_IN, 2)},
            options=options))
    return findings


# --- the pass ------------------------------------------------------------


_ORDER = {"error": 0, "warning": 1, "info": 2}


def scan(deck_bytes: bytes, profile_cfg: dict | None = None) -> list[DesignFinding]:
    """Every design finding on this deck, worst first.

    Ordered by severity and then by how many places it touches, because that is
    the order a designer would work in: the text nobody can read, then the
    colour that is wrong on forty shapes, then the badge in the corner.
    """
    # The memo holds elements alive, so it is emptied per scan rather than
    # growing for the life of the process. A concurrent scan clearing it can
    # only cost the other one cache misses, never a wrong answer - the identity
    # check in _is_frame_marker is what guarantees that.
    _MARKER_MEMO.clear()
    prs = Presentation(io.BytesIO(deck_bytes))
    from .stylespec import dominant_master

    master = dominant_master(prs)
    palette = palette_of(profile_cfg, master) if master is not None else {}

    findings: list[DesignFinding] = []
    # Four of the five checks below walk every slide, and nothing here mutates
    # the deck, so the walk is done once per slide and shared (_PLACED_MEMO).
    # The cache closes with this block: no box outlives the scan that measured
    # it, and no Presentation is pinned by it.
    with _placed_cache():
        for step in (lambda: _palette_findings(prs, palette),
                     lambda: _contrast_findings(prs, palette),
                     lambda: _overlap_findings(prs),
                     lambda: _fit_findings(prs),
                     lambda: _frame_findings(prs)):
            try:
                findings.extend(step())
            except Exception as exc:
                # One check failing must not take the other three with it. A
                # page that shows three of four answers and says which one is
                # missing is usable; a 500 is not.
                findings.append(DesignFinding(
                    finding_id=_finding_id("error", [str(exc)]),
                    kind="error", severity="info",
                    headline="One of the design checks could not run",
                    detail=f"{type(exc).__name__}: {exc}. The other checks on "
                           f"this page ran normally.",
                    slides=[], options=[], evidence={"places": 0}))

    findings.sort(key=lambda f: (_ORDER.get(f.severity, 3), -f.places,
                                 f.finding_id))
    return findings


# --- handing the decisions over ------------------------------------------
#
# A designer can ask the tool to answer every card itself (the design page's
# "let the tool decide"). That is not a default creeping back in: a default is
# the tool answering a question nobody asked, and this is one deliberate action,
# taken once, whose every consequence is listed with an Undo beside it.
#
# What it must NOT do is pretend to know things this file has already said it
# cannot know. Those kinds are named here rather than in the page, because the
# reason lives with the check.

UNDECIDABLE_KINDS = {
    "frame": "a page number, a source line and a stranded eyebrow are "
             "identical from the file; whether something belongs outside the "
             "frame is a question about the master, not one the deck answers",
    "error": "a check that could not run has nothing to propose",
}


def auto_choice(finding):
    """The remedy the tool picks when a designer hands the decision over, or
    None when this kind of finding is not the tool's to decide.

    The FIRST option with something to perform, every time. Not a new ranking:
    each check already puts its own recommendation first and says why in the
    note the designer reads ("offered first because it changes the smallest
    thing on the slide that can fix the problem"; "nothing moves, nothing
    changes size on screen, and the type scale is untouched, which is why it is
    first"). A second opinion computed here would be free to disagree with the
    page, and then the tool would be recommending one thing and doing another.

    This is why the fit check's ordering matters as much as its option list: it
    IS the recommendation, and it is what the tool does when a designer hands
    the decision over.
    """
    if finding.kind in UNDECIDABLE_KINDS:
        return None
    return next((o for o in (finding.options or ()) if o.op), None)


def auto_skip_reason(finding) -> str | None:
    """Why the tool left this one for the designer, in the words the page
    shows, or None when it would decide it."""
    if finding.kind in UNDECIDABLE_KINDS:
        return UNDECIDABLE_KINDS[finding.kind]
    if not next((o for o in (finding.options or ()) if o.op), None):
        return "nothing here can be changed automatically"
    return None


def finding_shapes(finding) -> list[tuple]:
    """(slide_index, shape_id) for every shape this finding's remedies would
    touch, deduped, in the order they appear.

    Read off the remedies rather than recorded separately on purpose: what a
    finding is ABOUT is exactly what its fixes would act on, and a second list
    kept alongside them could disagree. The page draws its highlight boxes from
    this, so a box on the render is a promise that clicking a remedy changes
    that shape.
    """
    out: list[tuple] = []
    seen = set()

    def add(slide_index, shape_id):
        if slide_index is None or shape_id is None:
            return
        key = (int(slide_index), str(shape_id))
        if key not in seen:
            seen.add(key)
            out.append(key)

    for option in finding.options or ():
        params = option.params or {}
        for target in params.get("targets") or ():
            add(target.get("slide_index"), target.get("shape_id"))
        add(params.get("slide_index"), params.get("shape_id"))
        if params.get("below") is not None:
            add(params.get("slide_index"), params.get("below"))
    return out


def slide_rects(deck_bytes: bytes, findings: list) -> dict:
    """{slide_index: [rect]} for drawing findings onto a rendered slide.

    Coordinates are fractions of the slide, so the page can lay them over an
    image of any size. Group children are transformed into slide space
    (placed_shapes), which is the whole reason this does not reuse
    render.shape_rects: that one reads a shape's raw offset, and a shape inside
    a group would get a box somewhere else entirely.
    """
    prs = Presentation(io.BytesIO(deck_bytes))
    slide_w, slide_h = prs.slide_width, prs.slide_height
    wanted: dict[int, list] = {}
    for finding in findings:
        for slide_index, shape_id in finding_shapes(finding):
            wanted.setdefault(slide_index, []).append((shape_id, finding))

    out: dict[int, list] = {}
    for slide_index, items in wanted.items():
        if not (0 <= slide_index < len(prs.slides)):
            continue
        boxes = {str(p.shape.shape_id): p.box
                 for p in placed_shapes(prs.slides[slide_index])}
        rects = []
        for shape_id, finding in items:
            box = boxes.get(str(shape_id))
            if box is None:
                continue
            rects.append({
                "x": max(0.0, box[0] / slide_w), "y": max(0.0, box[1] / slide_h),
                "w": min(1.0, (box[2] - box[0]) / slide_w),
                "h": min(1.0, (box[3] - box[1]) / slide_h),
                "severity": finding.severity, "kind": finding.kind,
                "finding_id": finding.finding_id,
                "label": finding.headline,
            })
        out[slide_index] = rects
    return out


def summary(findings: list[DesignFinding]) -> dict:
    out = {"total": len(findings), "by_kind": {}, "by_severity": {}}
    for f in findings:
        out["by_kind"][f.kind] = out["by_kind"].get(f.kind, 0) + 1
        out["by_severity"][f.severity] = out["by_severity"].get(f.severity, 0) + 1
    return out


__all__ = ["DesignFinding", "Remedy", "scan", "summary", "palette_of",
           "contrast_ratio", "luminance", "hex_of", "parse_hex",
           "placed_shapes", "slide_ground", "shape_fill", "chroma",
           "hue_distance", "natural_text_height", "finding_shapes",
           "slide_rects", "auto_choice", "auto_skip_reason",
           "UNDECIDABLE_KINDS"]
