"""Audit engine: one pipeline for audit and (later) fix.

run_audit() opens the deck, runs the mandatory pre-flight scan, builds the
Arabic index (the guard layer every module consults), then runs the selected
detection modules. Each module exposes detect(ctx) -> list[FindingRecord].
Preview/audit is this same pipeline with writes disabled by construction:
v1 modules only ever flag.
"""

import importlib
import json
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .profile import Profile
from .records import MODULES, FindingRecord, make_record
from .util import iter_shapes_deep


@dataclass
class AuditContext:
    prs: object
    profile: Profile
    deck_path: Path
    # (slide_index, shape_id, paragraph_index, run_index) for every Arabic
    # run; run_index None = paragraph-level rtl hit. shape_id is str.
    arabic_runs: set = field(default_factory=set)
    # (slide_index, shape_id) with any Arabic content: the coarse guard.
    arabic_shapes: set = field(default_factory=set)

    def shape_has_arabic(self, slide_index: int, shape_id) -> bool:
        return (slide_index, str(shape_id)) in self.arabic_shapes

    def run_is_arabic(self, slide_index: int, shape_id, para_idx: int, run_idx: int) -> bool:
        return (slide_index, str(shape_id), para_idx, run_idx) in self.arabic_runs


@dataclass
class AuditResult:
    deck_path: str
    profile_id: str
    profile_version: int
    slides: int
    records: list[FindingRecord]
    summary: dict

    def to_manifest(self) -> dict:
        return {
            "deck": self.deck_path,
            "profile_id": self.profile_id,
            "profile_version": self.profile_version,
            "slides": self.slides,
            "summary": self.summary,
            "records": [r.to_dict() for r in self.records],
        }

    def save_manifest(self, path: Path):
        Path(path).write_text(json.dumps(self.to_manifest(), indent=2),
                              encoding="utf-8")


def _build_arabic_index(prs, ctx: AuditContext):
    from spike.arabic import scan_presentation

    for hit in scan_presentation(prs):
        ctx.arabic_shapes.add((hit.slide_index, hit.shape_id))
        if hit.run_index is not None:
            ctx.arabic_runs.add(
                (hit.slide_index, hit.shape_id, hit.paragraph_index, hit.run_index))


_DIAGRAM_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"


def _is_smart_art(shape) -> bool:
    """python-pptx exposes no has_smart_art; SmartArt is a graphic frame whose
    a:graphicData uri is the diagram namespace."""
    el = shape._element
    for gd in el.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}graphicData"):
        if gd.get("uri") == _DIAGRAM_URI:
            return True
    return False


def preflight(prs) -> list[FindingRecord]:
    """Flag content we will never modify and may not perfectly preserve
    (SmartArt, charts, media, OLE) so no deck is silently downgraded."""
    records = []
    for s_idx, slide in enumerate(prs.slides):
        for shape, shape_path in iter_shapes_deep(slide.shapes):
            kind = None
            if _is_smart_art(shape):
                kind = "SmartArt"
            elif getattr(shape, "has_chart", False):
                kind = "chart"
            elif shape.shape_type in (MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.LINKED_PICTURE,
                                      MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
                                      MSO_SHAPE_TYPE.LINKED_OLE_OBJECT):
                kind = str(shape.shape_type)
            if kind:
                records.append(make_record(
                    slide_index=s_idx, shape_id=shape.shape_id, shape_path=shape_path,
                    module="preflight", issue_type="preflight.unmodifiable_content",
                    severity="info", action="flagged", confidence="deterministic",
                    message=f"Contains {kind}: preserved verbatim, not audited or modified.",
                ))
    return records


def _load_module(key: str):
    return importlib.import_module(f"qc.modules.{key}")


def run_audit(deck_path, profile: Profile | str, modules: list[str] | None = None) -> AuditResult:
    deck_path = Path(deck_path)
    if isinstance(profile, str):
        profile = Profile.load(profile)
    selected = list(modules) if modules else list(MODULES)
    unknown = [m for m in selected if m not in MODULES]
    if unknown:
        raise ValueError(f"Unknown modules {unknown}; valid: {list(MODULES)}")

    prs = Presentation(deck_path)
    ctx = AuditContext(prs=prs, profile=profile, deck_path=deck_path)
    _build_arabic_index(prs, ctx)

    records: list[FindingRecord] = list(preflight(prs))
    for key in selected:
        records.extend(_load_module(key).detect(ctx))

    summary = {
        "by_severity": dict(Counter(r.severity for r in records)),
        "by_issue_type": dict(Counter(r.issue_type for r in records)),
        "by_module": dict(Counter(r.module for r in records)),
        "arabic_flagged": sum(1 for r in records if r.arabic_flag),
        "total": len(records),
    }
    return AuditResult(
        deck_path=str(deck_path),
        profile_id=profile.id,
        profile_version=profile.version,
        slides=len(prs.slides),
        records=records,
        summary=summary,
    )
