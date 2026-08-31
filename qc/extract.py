"""Deterministic extraction: the ground truth every judgment pass argues from.

    python -m qc.extract deck.pptx                  > deck.json
    python -m qc.extract master.pptx --layouts      > master_layouts.json
    python -m qc.extract deck.pptx --master m.pptx  > pipeline_input.json

This runs BEFORE any model sees the deck, and it is the only place the numbers
come from. Every hex code, every EMU coordinate, every placeholder index a
later step uses is read out of the file here. A vision pass is handed this JSON
next to the rendered pictures, and what it may answer with is a judgment naming
things by the ids in here - never a measurement of its own. Ask a model for EMU
and you get plausible EMU.

Three modules in this package read style facts, and they answer different
questions. They must not be confused for each other:

    stylespec.py  what does a submitted MASTER declare?    (reads the master)
    bootstrap.py  what conventions does a deck follow?     (infers rules)
    extract.py    what is actually on these slides, and
                  what does this master offer to put them  (states facts, and
                  on?                                       nothing else)

Nothing here infers, thresholds, or judges. It produces no findings, and a
value it cannot read comes back as null rather than as a default that reads
like a fact.

TWO THINGS ARE REPORTED FOR EVERY COLOUR, and that distinction is the whole
point of the module for a palette review: what the file SAYS (`as_written` - an
explicit srgbClr, a schemeClr reference to a theme slot, or "inherited" because
the shape states nothing at all), and what it RESOLVES TO (`hex`, after the
clrMap, the theme's own scheme, and any tint/shade/lumMod on the reference). A
deck that is on-palette through theme references and a deck carrying forty
hand-typed near-navies can resolve to nearly the same picture and are entirely
different files. Only the first field tells them apart, and "which colours are
explicit" is a question a designer asks before they change anything.

THE RESOLVERS ARE THE AUDIT'S OWN (qc.design.shape_fill, slide_ground,
_run_color; spike.resolver.resolve_run), not python-pptx's accessors.
python-pptx answers None for every inherited colour and font, which is most of
them - a master that states its body colour once is the normal case - so a
reader that only looks at the run reports a deck with no colours in it. Sharing
the resolvers also means the JSON a model reasons over and the numbers a
finding is measured against cannot drift apart, which they would the day one of
them was improved.

GEOMETRY IS IN SLIDE COORDINATES, group transforms composed
(qc.design.placed_shapes). A shape's own left/top inside a group is in the
group's child space; handing those over as slide positions is how a card's icon
ends up "at" the wrong side of the slide. Group children are included, because
the icon welded to a card is exactly the shape a layout question is about, and
the group itself is listed too, because a group is what a designer selects.

Shape ids are the same currency the component pass already uses - str(shape_id),
unique within a slide - so "does this id exist" is one rule across every pass
that validates a model's answer (qc.components.inventory).
"""

import argparse
import io
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from spike.color_resolver import resolve_color_element
from spike.ns import find
from spike.resolver import resolve_run, theme_fonts
# _color_child and _run_color are this package's only definitions of "which
# element carries the colour" and "what colour is a run actually drawn in".
# Forking either one to keep this module's imports tidy would be how the
# extractor and the audit start disagreeing about the same deck.
from .design import (_color_child, _run_color, hex_of, placed_shapes,
                     shape_fill, slide_ground)
from .stylespec import alt_text, dominant_master, extract_layouts, extract_theme

SCHEMA_VERSION = 1


# --------------------------------------------------------------- primitives


def emu_to_dict(emu_val):
    """EMU, and the same number in inches.

    Both, because the two readers of this file want different ones: execution
    code works in EMU (it is what the file stores and the only unit that
    round-trips), and a person checking the JSON against what they see in
    PowerPoint thinks in inches."""
    if emu_val is None:
        return None
    return {"emu": int(emu_val), "inches": round(Emu(int(emu_val)).inches, 3)}


def _position(box):
    """A slide-space (left, top, right, bottom) box as the position block.

    Width and height are carried as well as the edges. They are derivable, and
    they are also what a placeholder comparison is actually about, so a
    consumer should not have to subtract to ask "is this the same size as the
    layout's box"."""
    if box is None:
        return None
    left, top, right, bottom = box
    return {"left": emu_to_dict(left), "top": emu_to_dict(top),
            "width": emu_to_dict(right - left),
            "height": emu_to_dict(bottom - top),
            "right": emu_to_dict(right), "bottom": emu_to_dict(bottom)}


_TRANSFORMS = ("tint", "shade", "lumMod", "lumOff", "alpha")


def _as_written(color_el) -> dict | None:
    """How ONE colour element was written, before anything resolves it.

    An explicit srgbClr and a schemeClr pointing at accent1 are the same
    picture and different intentions, and every palette question a designer
    asks starts from which of the two it is."""
    if color_el is None:
        return None
    tag = color_el.tag.split("}")[-1]
    out: dict = {}
    if tag == "srgbClr":
        out["as_written"] = "explicit_rgb"
        out["written_hex"] = (color_el.get("val") or "").upper() or None
    elif tag == "schemeClr":
        out["as_written"] = "theme_color"
        out["theme_slot"] = color_el.get("val")
    else:
        # sysClr, prstClr, scrgbClr, hslClr. Rare in client decks and named
        # rather than lumped into "unknown", because a reader seeing the token
        # can go and look at it.
        out["as_written"] = tag

    mods = {}
    for name in _TRANSFORMS:
        el = find(color_el, f"a:{name}")
        if el is None or not el.get("val"):
            continue
        try:
            mods[name] = round(int(el.get("val")) / 100_000.0, 4)
        except ValueError:
            continue
    if mods:
        # A tint on a theme reference is a deliberate 40% of the brand navy,
        # not a different navy. Recording it is what keeps that from being
        # reported as an off-palette colour.
        out["transforms"] = mods
    return out


def _written_fill(parent) -> dict | None:
    """How a shape's own a:solidFill under `parent` was written, or None when
    it states no solid fill of its own (inherited, or not a flat colour)."""
    solid = find(parent, "a:solidFill")
    if solid is None:
        return None
    return _as_written(_color_child(solid))


# ------------------------------------------------------------------- colour


def _fill(shape, slide, master) -> dict:
    """What this shape is painted, and how the file said so."""
    rgb, paint = shape_fill(shape, slide, master)
    out = {"paint": paint, "hex": hex_of(rgb) if rgb else None,
           "as_written": "inherited"}
    written = _written_fill(find(shape._element, "p:spPr"))
    if written:
        out.update(written)
    return out


def _line(shape, master) -> dict | None:
    """The shape's stroke, when it states one. Rules, dividers and card borders
    are a real part of a client's palette and they are never text or fill, so a
    palette review that skips them misses a whole class of off-brand colour."""
    ln = find(find(shape._element, "p:spPr"), "a:ln")
    if ln is None:
        return None
    if find(ln, "a:noFill") is not None:
        return {"paint": "none"}
    written = _written_fill(ln)
    if written is None:
        return None
    solid = find(ln, "a:solidFill")
    rgb = resolve_color_element(_color_child(solid), master)
    width = ln.get("w")
    out = {"paint": "solid", "hex": hex_of(rgb) if rgb else None}
    out.update(written)
    if width:
        try:
            out["width"] = emu_to_dict(int(width))
        except ValueError:
            pass
    return out


def _background(slide, master) -> dict:
    """The ground this slide is actually painted on, and WHICH LEVEL said so.

    The cascade is real and a designer needs to see it: a slide's own p:bg beats
    the layout's, which beats the master's, and a deck that states none anywhere
    lands on the theme's bg1 slot. "Explicit" means a different thing at each
    level - a hex typed onto one slide is a local override, the same hex on the
    master is the brand - so the level is reported next to the value."""
    rgb, where = slide_ground(slide, master)
    out = {"hex": hex_of(rgb) if rgb else None, "source": where,
           "stated_by": None, "as_written": None}
    for container, level in ((slide, "slide"), (slide.slide_layout, "layout"),
                             (master, "master")):
        element = getattr(container, "_element", None)
        bg = find(find(element, "p:cSld"), "p:bg")
        if bg is None:
            continue
        out["stated_by"] = level
        ref = find(bg, "p:bgRef")
        if ref is not None:
            written = _as_written(_color_child(ref)) or {}
            written["as_written"] = "theme_background"
            written["bg_fill_style_idx"] = ref.get("idx")
            out.update(written)
        else:
            pr = find(bg, "p:bgPr")
            written = _written_fill(pr)
            if written:
                out.update(written)
            elif pr is not None and find(pr, "a:noFill") is None:
                # A picture or a gradient: there is no single value, and saying
                # so is the honest answer. Contrast is not judged over it.
                out["as_written"] = "picture_or_gradient"
        break
    if out["stated_by"] is None:
        out["stated_by"] = "theme"
    return out


# --------------------------------------------------------------------- text


def _paragraphs(shape, slide, prs, master) -> list[dict]:
    """The words, and the font and colour they are ACTUALLY drawn in.

    Every font and colour here is resolved through the same cascade the audit
    walks, and each one carries the level that supplied it (`from`). That field
    is what turns "Arial 12pt" into something a designer can act on: set on the
    run is a local override somebody typed, inherited from the master's
    bodyStyle is the brand doing its job."""
    out = []
    for para in shape.text_frame.paragraphs:
        runs = []
        for run in para.runs:
            if not run.text:
                continue
            font = resolve_run(run, para, shape, slide, prs)
            rgb, source = _run_color(run, para, shape, slide, prs, master)
            colour = {"hex": hex_of(rgb) if rgb else None, "from": source,
                      "as_written": "inherited"}
            written = _written_fill(find(run._r, "a:rPr"))
            if written:
                colour.update(written)
            runs.append({
                "text": run.text,
                "font_name": font.family.value,
                "font_from": font.family.source,
                "complex_script_font": (font.cs_family.value
                                        if font.cs_family else None),
                "size_pt": font.size_pt.value,
                "size_from": font.size_pt.source,
                "bold": bool(font.bold.value),
                "color": colour,
            })
        if not runs and not para.text.strip():
            continue
        out.append({"level": para.level,
                    "alignment": str(para.alignment) if para.alignment else None,
                    "runs": runs})
    return out


# ------------------------------------------------------------------- shapes


def content_type(shape) -> str | None:
    """What the shape IS, one word, resolved in a fixed order.

    A chart and a table are graphic frames that also answer to has_text_frame,
    so the order matters: read the wrong one first and every chart on the deck
    is reported as a text box."""
    try:
        if shape.has_chart:
            return "chart"
        if shape.has_table:
            return "table"
    except Exception:
        pass
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return "group"
    try:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return "text"
    except Exception:
        pass
    return None


def extract_shape(placed, slide, prs, master) -> dict:
    """One shape, as fact. `placed` carries its slide-space box and its place
    in the paint order (qc.design.placed_shapes)."""
    shape = placed.shape
    info = {
        "id": str(shape.shape_id),
        "name": shape.name,
        "shape_type": str(shape.shape_type) if shape.shape_type else None,
        "content": content_type(shape),
        "z": placed.z,
        "in_group": placed.grouped,
        "top_level_index": placed.top,
        "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
        "position": _position(placed.box),
        "rotation": float(getattr(shape, "rotation", 0) or 0),
        "fill": _fill(shape, slide, master),
    }
    line = _line(shape, master)
    if line:
        info["line"] = line

    if info["is_placeholder"]:
        try:
            pf = shape.placeholder_format
            info["placeholder"] = {"idx": pf.idx,
                                   "type": str(pf.type) if pf.type else None}
        except Exception:
            info["placeholder"] = None

    try:
        alt = alt_text(shape)
    except Exception:
        alt = ""
    if alt:
        # Alt text is how this tool's own markers are declared (the
        # presentation-space rectangle), so it is never noise here.
        info["alt_text"] = alt

    try:
        if shape.has_text_frame and shape.text_frame.text.strip():
            info["paragraphs"] = _paragraphs(shape, slide, prs, master)
    except Exception as exc:
        info["text_error"] = f"{type(exc).__name__}: {exc}"
    return info


# --------------------------------------------------------------------- decks


def _open(source) -> tuple[object, str]:
    """A Presentation from a path or from bytes, and a label for the JSON.

    Bytes matter: the web layer holds uploads in memory and never writes the
    client's deck to disk, so an extractor that only took paths could not be
    called from the pipeline it exists for."""
    if isinstance(source, (bytes, bytearray)):
        return Presentation(io.BytesIO(bytes(source))), "(uploaded deck)"
    return Presentation(str(source)), str(source)


def _master_of(slide, fallback):
    """The master THIS slide resolves through, not the deck's dominant one.

    A deck that could not be fully rebuilt carries two masters on purpose
    (qc.applymaster.ApplyResult.stragglers), and resolving a straggler's colours
    against the applied master would report the wrong hex for every shape on
    it."""
    try:
        return slide.slide_layout.slide_master
    except Exception:
        return fallback


def extract_deck(source) -> dict:
    """Every slide, every shape, resolved. The working deck's ground truth."""
    prs, label = _open(source)
    fallback = dominant_master(prs)
    masters = list(prs.slide_masters)
    result = {
        "schema_version": SCHEMA_VERSION,
        "source_file": label,
        "slide_width": emu_to_dict(prs.slide_width),
        "slide_height": emu_to_dict(prs.slide_height),
        "master_count": len(masters),
        "theme": extract_theme(prs),
        "slides": [],
    }

    for idx, slide in enumerate(prs.slides):
        master = _master_of(slide, fallback)
        layout = slide.slide_layout
        try:
            master_index = next((i for i, m in enumerate(masters)
                                 if m.part is master.part), None)
        except Exception:
            master_index = None
        slide_info = {
            "slide_index": idx,
            "slide_id": slide.slide_id,
            "layout_name": layout.name,
            # The archetype token, verbatim. It is the master's own statement of
            # what a layout is for, and it is what layout matching reads second
            # (qc.applymaster.plan_assignments).
            "layout_type": layout._element.get("type"),
            "layout_placeholder_count": len(layout.placeholders),
            "master_index": master_index,
            "background": _background(slide, master),
            "shapes": [extract_shape(placed, slide, prs, master)
                       for placed in placed_shapes(slide)],
        }
        result["slides"].append(slide_info)
    return result


def extract_master_layouts(source) -> dict:
    """What the master OFFERS: one entry per layout, with its archetype and its
    placeholder geometry.

    Delegated to qc.stylespec.extract_layouts rather than re-read here, because
    that reader is already the canonical one and a second representation of a
    layout is a second thing to keep in step. Assets are left out: this view
    exists to be reasoned over next to the rendered pictures, and a base64
    background image in the middle of it helps nobody."""
    prs, label = _open(source)
    masters = list(prs.slide_masters)
    layouts = []
    for m_index, master in enumerate(masters):
        for entry in extract_layouts(master, embed_assets=False):
            entry = dict(entry)
            entry["master_index"] = m_index
            layouts.append(entry)
    return {
        "schema_version": SCHEMA_VERSION,
        "source_file": label,
        "slide_width": emu_to_dict(prs.slide_width),
        "slide_height": emu_to_dict(prs.slide_height),
        "master_count": len(masters),
        "theme": extract_theme(prs),
        "theme_fonts": (theme_fonts(masters[0]) if masters else {}),
        "layouts": layouts,
    }


def extract_pair(deck_source, master_source) -> dict:
    """Both sides of the question, in one document.

    Layout matching and the palette review are both comparisons - what the deck
    does against what the master offers - and a pass handed the two halves
    separately can be handed halves that came from different files."""
    return {"schema_version": SCHEMA_VERSION,
            "deck": extract_deck(deck_source),
            "master": extract_master_layouts(master_source)}


def palette_inventory(deck: dict) -> dict:
    """Every distinct colour in an extracted deck, and where each one came from.

    A roll-up, not a new read: it only counts what extract_deck already stated.
    It is here because "what is the palette of this deck, and which of it was
    typed in by hand" is the question a designer opens the tool with, and
    answering it from the slide list means every caller writing the same loop.

    Backgrounds are counted separately from fills and text, because "which
    colours are explicit for backgrounds" is a different question from which
    colours appear at all."""
    seen: dict[str, dict] = {}

    def note(hex_value, role, as_written, slot, slide_index):
        if not hex_value:
            return
        entry = seen.setdefault(hex_value, {
            "hex": hex_value, "uses": 0, "roles": {}, "written": {},
            "theme_slots": [], "slides": []})
        entry["uses"] += 1
        entry["roles"][role] = entry["roles"].get(role, 0) + 1
        entry["written"][as_written] = entry["written"].get(as_written, 0) + 1
        if slot and slot not in entry["theme_slots"]:
            entry["theme_slots"].append(slot)
        if slide_index not in entry["slides"]:
            entry["slides"].append(slide_index)

    for slide in deck.get("slides", []):
        s_idx = slide.get("slide_index")
        bg = slide.get("background") or {}
        note(bg.get("hex"), f"background ({bg.get('stated_by')})",
             bg.get("as_written"), bg.get("theme_slot"), s_idx)
        for shape in slide.get("shapes", []):
            fill = shape.get("fill") or {}
            note(fill.get("hex"), "fill", fill.get("as_written"),
                 fill.get("theme_slot"), s_idx)
            line = shape.get("line") or {}
            note(line.get("hex"), "line", line.get("as_written"),
                 line.get("theme_slot"), s_idx)
            for para in shape.get("paragraphs", []):
                for run in para.get("runs", []):
                    colour = run.get("color") or {}
                    note(colour.get("hex"), "text", colour.get("as_written"),
                         colour.get("theme_slot"), s_idx)

    colours = sorted(seen.values(), key=lambda c: (-c["uses"], c["hex"]))
    # extract_theme nests the scheme under "colors" alongside the clrMap and
    # the fonts; the roll-up wants the twelve slots and nothing else.
    theme = dict(((deck.get("theme") or {}).get("colors") or {}))
    return {
        "theme_slots": theme,
        "colours": colours,
        "explicit_count": sum(1 for c in colours
                              if c["written"].get("explicit_rgb")),
        "distinct_count": len(colours),
    }


def font_inventory(deck: dict) -> dict:
    """Every distinct typeface in an extracted deck, and where each one comes
    from.

    A roll-up, not a new read: it only counts what extract_deck already stated.
    The SOURCE is the part that matters and the part python-pptx cannot give
    you - "Arial, set on the run" is somebody typing over the brand, and "Arial,
    from the master's bodyStyle" is the brand itself. Both look identical on the
    slide and only one of them is a finding.

    Sizes travel with the family rather than as their own list, because a
    designer asks "what is the body copy" and not "which sizes exist".
    """
    seen: dict[tuple, dict] = {}
    theme = ((deck.get("theme") or {}).get("fonts") or {})

    for slide in deck.get("slides", []):
        s_idx = slide.get("slide_index")
        for shape in slide.get("shapes", []):
            for para in shape.get("paragraphs", []):
                for run in para.get("runs", []):
                    family = run.get("font_name")
                    if not family:
                        continue
                    # The level that supplied it, not the whole cascade path: a
                    # designer acts on "the run" or "the master", not on
                    # "layout placeholder 1 -> theme(+mn-lt)".
                    source = str(run.get("font_from") or "unknown")
                    stated = source.startswith(("run", "paragraph"))
                    key = (family, source)
                    entry = seen.setdefault(key, {
                        "family": family, "from": source,
                        "set_by_hand": stated, "uses": 0, "sizes": {},
                        "slides": [], "complex_script": None})
                    entry["uses"] += 1
                    size = run.get("size_pt")
                    if size is not None:
                        size = round(float(size), 1)
                        entry["sizes"][size] = entry["sizes"].get(size, 0) + 1
                    if s_idx not in entry["slides"]:
                        entry["slides"].append(s_idx)
                    if run.get("complex_script_font"):
                        entry["complex_script"] = run["complex_script_font"]

    fonts = sorted(seen.values(), key=lambda f: (-f["uses"], f["family"]))
    for entry in fonts:
        entry["sizes"] = [{"pt": pt, "uses": n}
                          for pt, n in sorted(entry["sizes"].items(),
                                              key=lambda kv: -kv[1])]
    return {
        "theme_fonts": theme,
        "fonts": fonts,
        "distinct_families": len({f["family"] for f in fonts}),
        "set_by_hand": sum(1 for f in fonts if f["set_by_hand"]),
    }


# ----------------------------------------------------------------------- CLI


def main(argv=None):
    parser = argparse.ArgumentParser(
        description="Deterministic ground-truth extraction from a .pptx.")
    parser.add_argument("pptx_path")
    parser.add_argument("--layouts", action="store_true",
                        help="read the master's layout inventory instead of "
                             "slide content")
    parser.add_argument("--master", metavar="MASTER.pptx",
                        help="also read this master, and emit both halves as "
                             "one pipeline input document")
    parser.add_argument("--palette", action="store_true",
                        help="emit the colour roll-up instead of the full "
                             "structure")
    parser.add_argument("--fonts", action="store_true",
                        help="emit the type roll-up instead of the full "
                             "structure")
    args = parser.parse_args(argv)

    if not Path(args.pptx_path).exists():
        parser.error(f"no such file: {args.pptx_path}")

    if args.layouts:
        data = extract_master_layouts(args.pptx_path)
    elif args.master:
        data = extract_pair(args.pptx_path, args.master)
    else:
        data = extract_deck(args.pptx_path)

    if args.palette:
        data = palette_inventory(data.get("deck", data))
    elif args.fonts:
        data = font_inventory(data.get("deck", data))

    json.dump(data, sys.stdout, indent=2, ensure_ascii=False)
    sys.stdout.write("\n")


if __name__ == "__main__":
    main()
