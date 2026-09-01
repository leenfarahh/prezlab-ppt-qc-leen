"""Fix engine (v1.5 core): apply selected, safe, deterministic fixes.

Rules from the PRD:
- Per-change accept/reject: only explicitly selected record_ids are applied.
- Arabic guard: records with arabic_flag are never fixable here, whatever
  the caller selects.
- Surgical writes: each fix touches only the XML element it names. The
  original bytes are never modified; output is a new file.
- The applied records flip action to "changed"; everything else is untouched.

Fixable issue types (the v1.5 deterministic tier):
  font.family_out_of_set          set the run's latin typeface (needs locator)
  master_slide.placeholder_geometry_off
                                  remove the slide placeholder's xfrm override
                                  so it re-inherits layout/master geometry
  shape_size.size_mismatch        set width/height to the cohort dominant
  header_footer.text_mismatch     set the footer placeholder's text
"""

import io
import re
from dataclasses import dataclass

from pptx import Presentation

from .util import is_backdrop, iter_shapes_deep, rides_with

_SIZE_RE = re.compile(r"^(\d+)x(\d+)$")
_LOC_RE = re.compile(r"^p(\d+)/r(\d+)$")
_TABLE_LOC_RE = re.compile(r"^t(\d+),(\d+)/p(\d+)/r(\d+)$")

FIXABLE_ISSUES = ("font.family_out_of_set", "font.size_off_role",
                  "master_slide.placeholder_geometry_off",
                  "shape_size.size_mismatch", "header_footer.text_mismatch",
                  "margin_alignment.edge_misaligned",
                  "margin_alignment.uneven_spacing",
                  "master_slide.foreign_master",
                  "font.title_autofit_shrunk",
                  "header_footer.fake_slide_number",
                  "header_footer.footer_off_canvas",
                  "margin_alignment.panel_row_misaligned",
                  "margin_alignment.cluster_rhythm",
                  "typography.case_inconsistent",
                  "typography.redundant_size_override",
                  "typography.size_inconsistent",
                  "margin_alignment.content_overflow",
                  "font.cs_typeface_missing",
                  "margin_alignment.recurring_off_position",
                  "margin_alignment.body_band_intrusion",
                  "margin_alignment.body_below_band",
                  "color_palette.off_palette_rgb",
                  "margin_alignment.space_edge_misaligned",
                  "margin_alignment.component_edge_misaligned",
                  # Changes the shape TREE, which is the most invasive thing
                  # this tool does to a slide, so it is never pre-ticked
                  # (_MODEL_JUDGED).
                  "margin_alignment.should_be_grouped")

# Deliberately absent from the list above, and to stay absent:
# margin_alignment.heading_past_margin. A title or standfirst running past a
# margin is a house-style question the client answers, not a defect this tool
# corrects (design lead, 19/08/2026). It is emitted flag-only, with no computed
# target, so adding it here would still fix nothing - but it would put a tick
# box in front of a designer, which is the wrong invitation.

# Fixes that MOVE or RESIZE shapes; each is checked after application and
# reverted if it introduced an overlap that was not there before.
_POSITIONAL_ISSUES = {"margin_alignment.edge_misaligned",
                      "margin_alignment.uneven_spacing",
                      "shape_size.size_mismatch",
                      "master_slide.placeholder_geometry_off",
                      "header_footer.footer_off_canvas",
                      "margin_alignment.panel_row_misaligned",
                      "margin_alignment.cluster_rhythm",
                      "margin_alignment.content_overflow",
                      "margin_alignment.recurring_off_position",
                      "margin_alignment.body_band_intrusion",
                      "margin_alignment.body_below_band",
                      "margin_alignment.space_edge_misaligned",
                      "margin_alignment.component_edge_misaligned"}

# Heuristic detections are still tickable (a designer's tick IS per-change
# validation) but carry lower confidence than the default floor.
_MIN_CONFIDENCE = {
    # medium = a judgment about what belongs together, checked against the
    # geometry but not derived from it. Tickable, never pre-ticked.
    "margin_alignment.should_be_grouped": ("deterministic", "high", "medium"),
    "margin_alignment.edge_misaligned": ("deterministic", "high", "medium"),
    "margin_alignment.uneven_spacing": ("deterministic", "high", "medium", "low"),
    "margin_alignment.cluster_rhythm": ("deterministic", "high", "medium"),
    "typography.size_inconsistent": ("deterministic", "high", "medium"),
    # medium = the run is not a placeholder, so which ROLE it plays was
    # inferred rather than read off the layout. Tickable: the target is the
    # profile's own number either way, and a wrong role is visible at a glance.
    "font.size_off_role": ("deterministic", "high", "medium"),
    # deterministic = clone-master repoint (visual no-op); medium = PowerPoint
    # re-applies the layout via COM behind designer approval
    "master_slide.foreign_master": ("deterministic", "high", "medium"),
    # medium = every column on the slide starts below the guide, so the guide
    # being the intended line is inference rather than evidence. Still tickable:
    # this move is never pre-selected either way (_BLOCK_MOVE_IS_APPROVAL).
    "margin_alignment.body_below_band": ("deterministic", "high", "medium"),
    # medium = the nearest palette colour is ambiguous (deltaE 5-10). Tickable,
    # because a designer's tick IS the per-change validation, and never
    # pre-ticked (_COLOR_TICK_IS_APPROVAL).
    "color_palette.off_palette_rgb": ("deterministic", "high", "medium"),
    # medium by construction: a component review is a judgment the model made
    # about the design, re-verified against the geometry but never promoted
    # to deterministic. Tickable, never pre-ticked.
    "margin_alignment.component_edge_misaligned": ("high", "medium"),
}


# Recolouring is script-neutral: a fill or a run colour changes no glyph and no
# shaping, so an Arabic deck's palette is as fixable as an English one. The
# Arabic guard is there for text that gets re-typed or re-fonted.
_ARABIC_SAFE_RECOLOR = {"color_palette.off_palette_rgb"}

# Geometry-only fixes are script-neutral for the same reason: moving or resizing
# a shape never opens its text, so a misaligned card is fixable in Arabic exactly
# as in English. Everything that edits runs or text keeps the Arabic guard
# (guard relaxation approved 12/08/2026 after an 89-slide Arabic deck came
# back 72% blocked with 120 computed geometry targets refused).
_ARABIC_SAFE_GEOMETRY = {
    "margin_alignment.edge_misaligned",
    "margin_alignment.uneven_spacing",
    "margin_alignment.cluster_rhythm",
    "margin_alignment.panel_row_misaligned",
    "margin_alignment.content_overflow",
    "margin_alignment.recurring_off_position",
    "margin_alignment.body_band_intrusion",
    "margin_alignment.body_below_band",
    "margin_alignment.space_edge_misaligned",
    "margin_alignment.component_edge_misaligned",
    "shape_size.size_mismatch",
    "master_slide.placeholder_geometry_off",
}

# Font substitution on Arabic runs changes shaping, so it is fixable but
# NEVER pre-selected: the designer's tick is the approval (requested
# 12/08/2026: "convert all fonts to Sakkal" on an Arabic deck whose
# learned convention was exactly that).
_ARABIC_TICK_IS_APPROVAL = {
    "font.family_out_of_set",
    "font.cs_typeface_missing",
}


# Fixes that move the WHOLE body of a slide. The tool can be certain the line
# was broken and still not be the one to decide that every element on the slide
# should shift, so the designer's tick is the approval (design decision,
# 20/08/2026).
_BLOCK_MOVE_IS_APPROVAL = {
    "margin_alignment.body_band_intrusion",
    "margin_alignment.body_below_band",
}


# A colour swap the eye can see. Inside the auto-replace band (deltaE <= 5) a
# palette snap is invisible on screen and pre-ticking it is safe; past that the
# colour visibly changes, and whether this deck's off-palette colour is a
# mistake or a decision is not the tool's call - brand and client colours sit
# off-palette on purpose. Fixable either way; pre-ticked only when invisible
# (design lead, 24/08/2026).
_COLOR_TICK_IS_APPROVAL = {"color_palette.off_palette_rgb"}


# Anything a language model judged. The geometry is re-verified by code before
# the record exists, so the NUMBERS are the tool's; what came from the model is the
# claim that these shapes are one component and that this line is the intended
# one. That is a design judgment, and a design judgment is the designer's to
# confirm - so it is offered, never pre-selected, whatever its confidence.
_MODEL_JUDGED = {"margin_alignment.component_edge_misaligned",
                 "margin_alignment.should_be_grouped"}


def tick_reason(record: dict) -> str | None:
    """Why this fix is never pre-selected, in the words the UI shows, or None
    when it may be pre-ticked on the usual evidence."""
    if record["issue_type"] == "margin_alignment.should_be_grouped":
        return ("the model judged these shapes to be one object and this writes "
                "that into the file; grouping is easy to undo and easy to miss, "
                "so ticking it is your approval")
    if record["issue_type"] in _MODEL_JUDGED:
        return ("The model grouped these shapes and chose the line; the geometry "
                "was checked but the judgment is yours to confirm")
    if record["issue_type"] in _BLOCK_MOVE_IS_APPROVAL:
        return ("moves every element on the slide down together: ticking it "
                "is your approval")
    if (record["issue_type"] in _COLOR_TICK_IS_APPROVAL
            and record["confidence"] != "high"):
        return ("the nearest palette colour is a visible change, not a snap: "
                "ticking it is your approval")
    if record["arabic_flag"] \
            and record["issue_type"] in _ARABIC_TICK_IS_APPROVAL:
        return "Arabic font substitution: ticking it is your explicit approval"
    return None


# Why a whole ISSUE TYPE has no automatic fix. Keyed by type, because the answer
# is a property of the check rather than of the slide it fired on.
#
# WRITTEN DOWN BECAUSE THE UI WAS SAYING NOTHING. A row that reads "no automatic
# fix" and stops invites exactly one conclusion - the tool is half-finished -
# and on a slide where three rows say it in a column, that conclusion is
# unavoidable (design lead, 31/08/2026). Every one of these has a reason, and
# most of the reasons are that the tool would have to make a decision that is
# not its to make.
_NO_FIX_REASON = {
    "margin_alignment.outside_safe_zone":
        "the breach is measured, the correction is not: moving the shape in "
        "could push it onto its neighbour, and shrinking it changes the "
        "composition. Nudge it in PowerPoint",
    "margin_alignment.heading_past_margin":
        "a heading running wide is a house-style question for the client, not "
        "a defect to correct",
    "header_footer.missing":
        "adding a placeholder means inserting a shape the slide does not have. "
        "Apply the layout that carries it, or add it on the master",
    "header_footer.position_mismatch":
        "the master states where this belongs and the slide disagrees; moving "
        "it slide-by-slide hides that rather than fixing it",
    "header_footer.font_mismatch":
        "the footer's type comes from the master, so correcting it here would "
        "pin one slide out of step with the rest",
    "font.mixed_weight":
        "which weight was intended is a design decision - the tool can see "
        "that a paragraph mixes them, not which one is right",
    "font.theme_ref_disallowed":
        "the profile asks for stated families and this run resolves through "
        "the theme; changing that is a master decision",
    "color_palette.disallowed_theme_slot":
        "the slot is wrong at the theme level, so every deck on this master "
        "has it. Fixing one shape leaves the cause in place",
    "master_slide.layout_outlier":
        "a layout used once is worth a look, not a change",
    "master_slide.no_usable_master":
        "there is nothing to correct against until the deck has a master",
    "margin_alignment.squeezed_text":
        "the box is too small for its text and both are the designer's: "
        "resizing changes the layout, re-typing changes the copy",
    "margin_alignment.text_overlap":
        "two text blocks overlap and which one should move is a composition "
        "decision",
    "margin_alignment.text_anchor_mismatch":
        "vertical anchoring is a deliberate choice as often as it is a slip",
    "margin_alignment.overlap_check_capped":
        "a note about the run, not a finding on the slide",
    "shape_size.off_grid":
        "the grid is advisory in this profile; snapping to it would resize "
        "shapes nobody asked to resize",
    "preflight.unmodifiable_content":
        "the tool could not open this content, so it cannot change it either",
}


def no_fix_reason(record: dict) -> str:
    """Why this row has no tick, in the words the UI shows. "" when it has one.

    Answers in the order the checks actually run, because a designer reads the
    FIRST reason as the reason: an Arabic block outranks a missing target, which
    outranks the type having no fix at all.
    """
    if record.get("action") == "changed":
        return ""
    if is_fixable(record):
        return ""
    issue = record.get("issue_type", "")
    if issue in FIXABLE_ISSUES:
        # The type has a fix; THIS record cannot use it.
        if record.get("arabic_flag"):
            return ("Arabic content: this correction edits runs, and the tool "
                    "never re-types Arabic without a designer")
        if record.get("action") != "flagged":
            return "already handled during the run"
        allowed = _MIN_CONFIDENCE.get(issue, ("deterministic", "high"))
        if record.get("confidence") not in allowed:
            return (f"read with {record.get('confidence')} confidence, which is "
                    f"below the bar for changing a client's file automatically")
        return ("the check found no safe target to correct this to on this "
                "slide")
    return _NO_FIX_REASON.get(
        issue, "this check reports what it sees and does not compute a "
               "correction")


def needs_explicit_tick(record: dict) -> bool:
    """True for fixes the UI must never pre-select: Arabic font substitutions,
    where ticking is the designer's explicit approval, and whole-slide body
    moves, where the tick is the decision that the slide should shift."""
    return tick_reason(record) is not None


def is_fixable(record: dict) -> bool:
    """A record the UI may offer for apply. Arabic content is fixable for
    pure-geometry issues (the text is never touched) and for font
    substitutions (never pre-selected; the tick is the approval);
    confidence must be deterministic or high; the fix must have an
    explicit target."""
    if record["issue_type"] not in FIXABLE_ISSUES:
        return False
    if record["arabic_flag"] \
            and record["issue_type"] not in _ARABIC_SAFE_GEOMETRY \
            and record["issue_type"] not in _ARABIC_SAFE_RECOLOR \
            and record["issue_type"] not in _ARABIC_TICK_IS_APPROVAL:
        return False
    if record["action"] != "flagged":
        return False
    allowed = _MIN_CONFIDENCE.get(record["issue_type"], ("deterministic", "high"))
    if record["confidence"] not in allowed:
        return False
    if record["issue_type"] == "master_slide.placeholder_geometry_off":
        return True  # target state is "inherit"; no new_value needed
    return record.get("new_value") is not None


@dataclass
class FixOutcome:
    record_id: str
    outcome: str  # "changed" | "skipped"
    reason: str = ""


def _find_shape(slide, shape_id: str):
    for shape, _path in iter_shapes_deep(slide.shapes):
        if str(shape.shape_id) == shape_id:
            return shape
    return None


# The perceptual floor, the same 0.03in qc.copilot and qc.components measure
# against. Stated again rather than imported because importing a judgment pass
# into the fixer would invert the dependency: the fixer is what those passes
# feed, and it must not need them present to run.
ON_THE_LINE_EMU = 28575


def _measured_now(shape, record):
    """The value this record's `old_value` was read from, re-read off the deck
    as it stands now, or None when the record does not name a measurable."""
    prop = record.get("property") or ""
    if prop.endswith(".x"):
        return shape.left
    if prop.endswith(".y"):
        return shape.top
    if prop == "spPr.xfrm.ext":
        return f"{shape.width}x{shape.height}"
    return None


def still_open(record: dict, slide) -> bool:
    """Whether a MODEL-RAISED record still describes this deck.

    A measured record is re-derived by the re-audit after every fix, so it can
    never go stale. A vision record cannot be re-derived without asking the
    model again, and until now that meant it was simply dropped: applying one
    unrelated font fix re-audited the deck, replaced the manifest, and took
    every copilot and component suggestion on the page with it (reproduced
    01/09/2026 - one font fix, five suggestions gone, no message). The designer
    then had to spend another slide's worth of API calls to get back a list
    they had already been shown.

    So they are carried across instead, and each one is re-checked here against
    the geometry it was measured from. Three answers:

      the shape is gone            drop it, there is nothing to move
      it has moved or resized      drop it: the model judged a slide that no
                                   longer exists, and its computed target was
                                   read off the old positions
      it is already on the line    drop it, something else fixed it

    Anything else survives with its target intact, because nothing it depends
    on has changed. Kept deliberately strict: a stale target is a wrong move
    applied with confidence, which is worse than a suggestion a designer has to
    ask for again.
    """
    shape = _find_shape(slide, str(record["shape_id"]))
    if shape is None:
        return False
    current = _measured_now(shape, record)
    if current is None:
        return True                     # nothing to re-measure; the judgment stands
    if str(current) != str(record.get("old_value")):
        return False
    try:
        if abs(int(record["new_value"]) - int(current)) <= ON_THE_LINE_EMU:
            return False
    except (TypeError, ValueError, KeyError):
        pass
    return True


def _slide_rects(slide) -> dict:
    """{shape_id: (container, (l, t, r, b), texty)}. Group children carry
    their group's id as container: their EMU live in group space, so
    rectangles only compare within the same container."""
    rects = {}

    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def walk(shapes, container):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes, str(shape.shape_id))
            l, t = shape.left, shape.top
            w, h = shape.width, shape.height
            if None in (l, t, w, h):
                continue
            texty = bool(getattr(shape, "has_text_frame", False)
                         and shape.text_frame.text.strip())
            rects[str(shape.shape_id)] = (container, (l, t, l + w, t + h),
                                          texty)

    walk(slide.shapes, "")
    return rects


def _overlap_area(a, b) -> int:
    ox = min(a[2], b[2]) - max(a[0], b[0])
    oy = min(a[3], b[3]) - max(a[1], b[1])
    return ox * oy if (ox > 0 and oy > 0) else 0


def _overlaps(a, b, min_cover=0.1) -> bool:
    smaller = min((a[2] - a[0]) * (a[3] - a[1]),
                  (b[2] - b[0]) * (b[3] - b[1]))
    return smaller > 0 and _overlap_area(a, b) / smaller > min_cover


# a text-on-text overlap that GROWS by this share of the smaller shape is
# a deepened collision even though the boxes already touched
_DEEPEN_RATIO = 0.15


def _collision_created(slide, moved_ids: set, before: dict) -> str | None:
    """Non-None when a moved shape now overlaps a same-container shape it
    did not overlap before the fix, or - for two TEXT shapes - overlaps it
    substantially deeper than before. A fix must never make the slide
    worse (real-deck findings: 14/07/2026, a spacing fix walked a logo
    into its neighbor; 12/08/2026, RTL text boxes overlap as rectangles
    even when the text does not, so a fix pushed one column's text into
    another without creating any NEW overlap)."""
    after = _slide_rects(slide)
    for mid in moved_ids:
        if mid not in after or mid not in before:
            continue
        container, rect, texty = after[mid]
        old_rect = before[mid][1]
        for oid, (other_container, other, other_texty) in after.items():
            if oid == mid or other_container != container:
                continue
            other_before = before.get(oid, (other_container, other, False))[1]
            if (_overlaps(rect, other)
                    and not _overlaps(old_rect, other_before)):
                return (f"snapping would push shape {mid} into shape {oid}; "
                        "left for manual fix")
            if texty and other_texty:
                grew = (_overlap_area(rect, other)
                        - _overlap_area(old_rect, other_before))
                smaller = min((rect[2] - rect[0]) * (rect[3] - rect[1]),
                              (other[2] - other[0]) * (other[3] - other[1]))
                if smaller > 0 and grew > _DEEPEN_RATIO * smaller:
                    return (f"snapping would push shape {mid} deeper into "
                            f"text shape {oid}; left for manual fix")
    return None


def _restore_rects(slide, moved_ids: set, before: dict) -> None:
    for shape, _p in iter_shapes_deep(slide.shapes):
        sid = str(shape.shape_id)
        if sid in moved_ids and sid in before:
            l, t, r, b = before[sid][1]
            shape.left, shape.top = l, t
            shape.width, shape.height = r - l, b - t


# 2/3 of a shape inside another = contained
# (keep in sync with qc/modules/margin_alignment.py CONTAIN_MIN)
_CONTAIN_MIN = 0.66


def _slide_size(slide) -> tuple:
    try:
        prs = slide.part.package.main_document_part.presentation
        return prs.slide_width, prs.slide_height
    except Exception:
        return 12192000, 6858000


def _carried_contents(slide, container, along: str | None = None) -> list:
    """The container's COLLECTION: everything that composes with it and must
    travel with it.

    Two kinds of member, one rule ("what a designer had selected when they
    dragged this"):

    - CONTAINED: >= _CONTAIN_MIN of its area inside the container's bounds - a
      panel's labels, photos, icons. Orphaning these aligned the background and
      left the content behind (real-deck finding, 19/07/2026: a photo-grid panel
      snapped left away from its photos).
    - RIDING: welded to it (overlapping) or sitting within a satellite gap
      ACROSS the axis of the move - a corner rule drawn over a photo, a quote
      mark above a paragraph, a chip stacked on its box. Leaving these behind is
      what pulled the brackets off the photos on the leadership-quotes slide
      (design lead, 20/08/2026). See qc.util.rides_with for why the axis
      matters: shapes ALONG the move are its spacing peers, not its collection.

    Satellites are never chained: each is tested against THIS container only, so
    a slide of touching decorative shapes cannot drag half the canvas along.
    Backdrops and pinned page furniture ride nothing."""
    rects = _slide_rects(slide)
    key = str(container.shape_id)
    if key not in rects:
        return []
    sw, sh = _slide_size(slide)
    c_container, c_box, _tx = rects[key]
    cl, ct, cr, cb = c_box
    c_area = (cr - cl) * (cb - ct)
    carried = []
    for shape, _p in iter_shapes_deep(slide.shapes):
        sid = str(shape.shape_id)
        if sid == key or sid not in rects:
            continue
        other_container, box, _tx2 = rects[sid]
        if other_container != c_container:
            continue
        l, t, r, b = box
        area = (r - l) * (b - t)
        if area <= 0:
            continue
        ix = max(0, min(r, cr) - max(l, cl))
        iy = max(0, min(b, cb) - max(t, ct))
        if ix * iy >= _CONTAIN_MIN * c_area:
            # This candidate is the container. Contents ride containers, never
            # the reverse: an icon nudged onto its peers' inset must not drag
            # the panel it sits in, or the inset fix and the panel-row fix stop
            # composing (regression: test_inset_fix_then_panel_move_converge).
            continue
        contained = area < c_area and ix * iy >= _CONTAIN_MIN * area
        member = contained or rides_with(box, c_box, along)
        # Furniture is checked on BOTH branches: a recurring footer logo that
        # happens to sit inside a lifted image's box is still furniture, and it
        # travelled up onto the photo when only the riding branch asked (real-
        # deck finding, 21/07/2026, the 'Strategy&' logo).
        if (member and not is_backdrop(box, sw, sh)
                and not _pinned_furniture(slide, shape, l, t)):
            carried.append(shape)
    return carried


# An alignment record that states the set it was measured against, written by
# the pass that knows it (qc.copilot.synthesize) as "align-x:6,8,10" /
# "align-y:...". A record without one is every record this tool has ever
# written before now, and it behaves exactly as it did.
_ALIGN_LOCATORS = ("align-x:", "align-y:")


def _cluster_ids(record) -> set:
    """The peer set an alignment record was measured against, or empty."""
    loc = record.get("locator") or ""
    if not loc.startswith(_ALIGN_LOCATORS):
        return set()
    return {sid for sid in loc.split(":", 1)[1].split(",") if sid}


def _peer_pinned(slide, record, target_id) -> set:
    """Ids that must NOT ride the shape this record moves: the other members of
    its alignment cluster, and whatever lives inside them.

    A SHAPE BEING HELD TO A LINE CANNOT ALSO BE A SATELLITE OF A SHAPE BEING
    MOVED ONTO THAT LINE. _carried_contents cannot see the difference - it
    reads overlap and adjacency, and for a vertical move it carries what sits
    beside the mover, which on a row of cards is the rest of the row. The fix
    then seats the stray on the line and drags its neighbours, already on that
    line, the same distance off it. Measured on a ten-circle grid (01/09/2026):
    the row's spread went from 0.184in to 0.163in, one circle came onto the line
    and two left it. qc.fixer._fix_space_edge and _fix_component_edge each say
    the same thing about their own membership; this is the third place it is
    true and the first where a model supplied the membership.

    WELDED MEMBERS ARE NOT PINNED. Overlapping shapes are one object in this
    engine's terms (qc.util.rides_with: "composed by definition"), so a badge
    sitting on top of a peer travels with whichever of them moves. Only
    neighbours that merely share the line are held still.

    A peer's own CONTENTS are pinned with it: an icon inside the card beside
    ours sits well within the satellite gap of ours, and carrying it would tear
    it off the card it belongs to. Containment is not chained - contents are
    tested against the peers, never against each other - so one crowded slide
    cannot pin itself solid.
    """
    peers = _cluster_ids(record) - {str(target_id)}
    if not peers:
        return set()
    rects = _slide_rects(slide)
    target = rects.get(str(target_id))
    if target is None:
        return set()
    t_container, t_box, _texty = target

    pinned = set()
    for pid in peers:
        entry = rects.get(pid)
        if entry is None or entry[0] != t_container:
            continue
        if _overlap_area(entry[1], t_box) > 0:
            continue                       # welded to the mover: it travels
        pinned.add(pid)

    held = [(rects[pid][0], rects[pid][1]) for pid in pinned]
    for sid, (container, box, _t) in rects.items():
        if sid in pinned or sid == str(target_id):
            continue
        area = (box[2] - box[0]) * (box[3] - box[1])
        if area <= 0:
            continue
        for p_container, p_box in held:
            if p_container == container and \
                    _overlap_area(box, p_box) >= _CONTAIN_MIN * area:
                pinned.add(sid)
                break
    return pinned


def _translate(shapes, dx: int, dy: int) -> str | None:
    """Move shapes by a delta, materializing complete geometry first."""
    for shape in shapes:
        error = _materialize_xfrm(shape)
        if error:
            return f"shape {shape.shape_id}: {error}"
    for shape in shapes:
        if dx:
            shape.left = shape.left + dx
        if dy:
            shape.top = shape.top + dy
    return None


_FURNITURE_STRIP = 0.88  # the bottom strip where page furniture lives


def _recurring_furniture(slide, shape, left: int, top: int) -> bool:
    """True when the same-named shape sits at (nearly) the same spot on 3+
    slides: the footer logo / source / page-number layer. Furniture never
    rides a lift (real-deck finding, 21/07/2026: the 'Strategy&' logo text
    overlapped a lifted image and travelled up onto the photo; the
    designer keeps furniture pinned)."""
    try:
        prs = slide.part.package.main_document_part.presentation
    except Exception:
        return False
    count = 0
    for other in prs.slides:
        for sh in other.shapes:
            if (sh.name == shape.name and sh.left is not None
                    and abs(sh.left - left) <= 36000
                    and abs(sh.top - top) <= 36000):
                count += 1
                break
        if count >= 3:
            return True
    return False


def _pinned_furniture(slide, shape, left: int, top: int) -> bool:
    """Page furniture that never travels with anything: in the bottom strip AND
    recurring across the deck. Both halves are needed - a source line low on one
    slide is content, and a repeated header chip is not in the strip."""
    try:
        slide_h = slide.part.package.main_document_part.presentation.slide_height
    except Exception:
        slide_h = 6858000
    return (top >= _FURNITURE_STRIP * slide_h
            and _recurring_furniture(slide, shape, left, top))


def _collection_riders(slide, anchor_ids: list, along: str | None = None) -> list:
    """Every satellite of the given anchors: the collection minus the anchors.

    One notion of "what travels with this shape" for the whole engine
    (_carried_contents). Lifting an image without its labels stranded them
    (ground truth, 20/07/2026, the wheel slide's bottom clusters); snapping a
    photo without its corner rule pulled the bracket off it (design lead,
    20/08/2026). Both are the same mistake, so both ask the same question."""
    seen = set(str(a) for a in anchor_ids)
    out = []
    for aid in anchor_ids:
        anchor = _find_shape(slide, str(aid))
        if anchor is None:
            continue
        for shape in _carried_contents(slide, anchor, along):
            sid = str(shape.shape_id)
            if sid not in seen:
                seen.add(sid)
                out.append(shape)
    return out


def _fix_lift(shape, record, slide=None) -> str | None:
    """Move the tail of a same-size cohort back onto its rhythm, satellites
    riding along."""
    if slide is None:
        return "lift needs slide context"
    loc = record.get("locator") or ""
    along_col = loc.startswith("lift-col:")
    ids = loc.split(":", 1)[1].split(",")
    by_id = {str(sh.shape_id): sh for sh, _p in iter_shapes_deep(slide.shapes)}
    anchors = [by_id.get(sid) for sid in ids]
    if any(a is None for a in anchors):
        return "cohort member no longer resolves (deck changed since audit)"
    try:
        delta = int(record["new_value"]) - int(record["old_value"])
    except (TypeError, ValueError):
        return "gap values are not EMU integers"
    moving = anchors + _collection_riders(slide, ids,
                                          "y" if along_col else "x")
    return _translate(moving,
                      0 if along_col else delta,
                      delta if along_col else 0)


def _fix_rescale(shape, record, slide=None) -> str | None:
    """The designer's select-all shrink: scale every content shape by the
    stored factor, anchored at the left margin, one proportional transform.
    Alignments and gaps scale together, so nothing inside the selection can
    collide with anything else inside it."""
    if slide is None:
        return "rescale needs slide context"
    loc = record.get("locator") or ""
    try:
        _tag, permille, anchor, ids = loc.split(":", 3)
        scale = int(permille) / 1000.0
        anchor_l, anchor_t = (int(v) for v in anchor.split(","))
        id_list = ids.split(",")
    except (ValueError, TypeError):
        return "record has no usable scale locator"
    by_id = {str(sh.shape_id): sh for sh, _p in iter_shapes_deep(slide.shapes)}
    shapes = [by_id.get(sid) for sid in id_list]
    if any(s is None for s in shapes):
        return "content member no longer resolves (deck changed since audit)"
    for s in shapes:
        error = _materialize_xfrm(s)
        if error:
            return f"shape {s.shape_id}: {error}"
    bbox_l = min(s.left for s in shapes)
    bbox_t = min(s.top for s in shapes)
    for s in shapes:
        s_left, s_top = s.left, s.top  # read before any write
        s.left = anchor_l + int((s_left - bbox_l) * scale)
        s.top = anchor_t + int((s_top - bbox_t) * scale)
        s.width = int(s.width * scale)
        s.height = int(s.height * scale)
    return None


def _fix_body_band(shape, record, slide=None) -> str | None:
    """Seat the slide's body on the master's body guide: ONE vertical translate
    of every content shape, so the strip the master keeps clear under the header
    is cleared without changing anything inside the body.

    The locator carries both the delta and the full selection ('band:<dy>:<ids>')
    because the move is a slide-level decision, not a per-shape one: the audit
    decided which shapes are the body, and re-deriving that here would let the
    two disagree. Shapes contained in the movers ride for free - they are
    top-level members of the selection themselves, or children of a group whose
    own offset moves them."""
    if slide is None:
        return "band fix needs slide context"
    loc = record.get("locator") or ""
    if not loc.startswith("band:"):
        return "record has no band target"
    try:
        _tag, delta, ids = loc.split(":", 2)
        dy = int(delta)
    except (ValueError, TypeError):
        return "band locator is not a delta and an id list"
    by_id = {str(sh.shape_id): sh for sh, _p in iter_shapes_deep(slide.shapes)}
    id_list = [sid for sid in ids.split(",") if sid]
    moving = [by_id.get(sid) for sid in id_list]
    if not moving or any(s is None for s in moving):
        return "content member no longer resolves (deck changed since audit)"
    # Satellites ride, as with every other move: a corner rule welded to a
    # photo, or a mark the audit kept out of the body because it sits in the
    # header, still belongs to what is moving.
    moving += _collection_riders(slide, id_list, "y")
    return _translate(moving, 0, dy)


def _planned_movers(slide, rec) -> set:
    """Shape ids a positional fix will move: the target (or the row/column
    tail) plus any visually contained shapes the fixers carry along.
    Computed before the fix runs so the collision guard checks and can
    restore the complete set."""
    target = str(rec["shape_id"])
    base = [target]
    loc = rec.get("locator") or ""
    if loc.startswith(("row:", "col:")):
        ids = loc.split(":", 1)[1].split(",")
        if target in ids:
            base = ids[ids.index(target):]
    elif loc.startswith(("dist-row:", "dist-col:")):
        base = loc.split(":", 1)[1].split(",")  # a distribution moves them all
    elif loc.startswith(("lift-row:", "lift-col:")):
        ids = loc.split(":", 1)[1].split(",")
        sats = _collection_riders(slide, ids,
                                  "y" if loc.startswith("lift-col:") else "x")
        return set(ids) | {str(s.shape_id) for s in sats}
    elif loc.startswith("scale:"):
        return set(loc.split(":", 3)[3].split(","))
    elif loc.startswith("band:"):
        ids = loc.split(":", 2)[2].split(",")
        sats = _collection_riders(slide, ids, "y")
        return set(ids) | {str(s.shape_id) for s in sats}
    elif loc.startswith("edge:"):
        # every stray in the stack moves, each with its own contents
        ids = loc.split(":", 3)[2].split(",")
        sats = _collection_riders(slide, ids, "x")
        return set(ids) | {str(s.shape_id) for s in sats}
    moved = set(base)
    # The axis must match the one the FIXER will use, or the guard checks a
    # different set of shapes than the one that moves.
    prop = rec.get("property") or ""
    axis = "y" if prop.endswith(".y") else "x" if prop.endswith(".x") else None
    if rec["issue_type"] in ("margin_alignment.edge_misaligned",
                             "margin_alignment.uneven_spacing",
                             "margin_alignment.panel_row_misaligned",
                             "margin_alignment.recurring_off_position",
                             # a frame snap carries its contents like any other
                             # edge fix; without it here the claim guard cannot
                             # see them and two stacked members each move twice
                             "margin_alignment.space_edge_misaligned"):
        by_id = {str(sh.shape_id): sh
                 for sh, _p in iter_shapes_deep(slide.shapes)}
        for bid in base:
            sh = by_id.get(bid)
            if sh is not None:
                # Pinned peers are subtracted here for the same reason they are
                # in _fix_edge: the guard has to check and restore the set that
                # ACTUALLY moves, and the claim set has to leave the other
                # strays in the cluster free to be fixed in the same round -
                # otherwise the first record claims the whole row and every
                # other stray on it comes back "shares shapes with a fix
                # already applied".
                pinned = _peer_pinned(slide, rec, bid)
                moved.update(str(c.shape_id)
                             for c in _carried_contents(slide, sh, axis)
                             if str(c.shape_id) not in pinned)
    if loc.startswith(("row:", "col:")):
        sats = _collection_riders(slide, sorted(moved),
                                  "x" if loc.startswith("row:") else "y")
        moved.update(str(s.shape_id) for s in sats)
    return moved


def _materialize_xfrm(shape) -> str | None:
    """Write the shape's effective geometry into an explicit, COMPLETE xfrm
    before any positional edit.

    Root cause of a real-deck corruption (14/07/2026): removing a
    placeholder's xfrm (geometry snap) and then setting .left on it creates
    an off-WITHOUT-ext transform; python-pptx reads the missing extent from
    the layout, but PowerPoint renders the placeholder as a degenerate
    sliver (one letter per line). Setting all four values first guarantees
    off+ext always travel together."""
    left, top = shape.left, shape.top
    width, height = shape.width, shape.height
    if left is None or top is None or width is None or height is None:
        return "shape geometry is not fully resolvable"
    shape.left, shape.top, shape.width, shape.height = left, top, width, height
    return None


def _set_cs_typeface(run, typeface: str) -> None:
    """Set the complex-script typeface, respecting rPr child order (a:cs
    must precede a:sym / hyperlink / a:rtl, which Arabic runs often
    carry)."""
    from pptx.oxml.ns import qn

    rpr = run._r.get_or_add_rPr()
    cs = rpr.find(qn("a:cs"))
    if cs is None:
        cs = rpr.makeelement(qn("a:cs"), {})
        after = [qn("a:sym"), qn("a:hlinkClick"), qn("a:hlinkMouseOver"),
                 qn("a:rtl"), qn("a:extLst")]
        anchor = next((el for el in rpr if el.tag in after), None)
        if anchor is not None:
            anchor.addprevious(cs)
        else:
            rpr.append(cs)
    cs.set("typeface", typeface)


def _fix_font_family(shape, record) -> str | None:
    loc = record.get("locator") or ""
    tm = _TABLE_LOC_RE.match(loc)
    if tm:
        if not getattr(shape, "has_table", False):
            return "shape has no table"
        row_i, col_i, p_idx, r_idx = (int(g) for g in tm.groups())
        table = shape.table
        try:
            cell = table.rows[row_i].cells[col_i]
        except IndexError:
            return "table cell no longer resolves (deck changed since audit)"
        paras = cell.text_frame.paragraphs
    else:
        m = _LOC_RE.match(loc)
        if not m:
            return "record has no run locator"
        p_idx, r_idx = int(m.group(1)), int(m.group(2))
        if not getattr(shape, "has_text_frame", False):
            return "shape has no text frame"
        paras = shape.text_frame.paragraphs
    if p_idx >= len(paras) or r_idx >= len(paras[p_idx].runs):
        return "locator no longer resolves (deck changed since audit)"
    run = paras[p_idx].runs[r_idx]
    if (record.get("property") or "") == "rPr.cs.typeface":
        _set_cs_typeface(run, record["new_value"])
    else:
        run.font.name = record["new_value"]
    return None


def _fix_off_palette_color(shape, record) -> str | None:
    """Repaint one off-palette surface with the nearest palette colour.

    Two surfaces, one record shape: `property` says which. "spPr.solidFill" is
    the shape's own fill; "rPr.solidFill" is one run's text colour, addressed by
    the record's locator - a shape-wide repaint would take a deliberate accent
    word with it.

    python-pptx does the writing. Setting a solid fill means removing whatever
    fill was there and inserting a:solidFill at the one position the schema
    allows, and hand-rolled XML that gets that wrong produces a file PowerPoint
    offers to repair rather than open (the same reason qc.remedy defers to it).
    """
    from pptx.dml.color import RGBColor

    raw = str(record.get("new_value") or "").lstrip("#")
    if len(raw) != 6:
        return "new_value is not a six-digit hex colour"
    try:
        rgb = RGBColor.from_string(raw.upper())
    except ValueError:
        return f"'{raw}' is not a colour"

    prop = record.get("property") or ""
    if prop.startswith("rPr"):
        m = _LOC_RE.match(record.get("locator") or "")
        if not m:
            return "record has no run locator"
        if not getattr(shape, "has_text_frame", False):
            return "shape has no text frame"
        paras = shape.text_frame.paragraphs
        p_idx, r_idx = int(m.group(1)), int(m.group(2))
        if p_idx >= len(paras) or r_idx >= len(paras[p_idx].runs):
            return "locator no longer resolves (deck changed since audit)"
        paras[p_idx].runs[r_idx].font.color.rgb = rgb
        return None

    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb
    except (AttributeError, NotImplementedError, TypeError) as exc:
        # A picture, a group or a connector has no fill python-pptx will set.
        return f"this shape's fill cannot be set ({type(exc).__name__})"
    return None


def _fix_placeholder_geometry(shape, record) -> str | None:
    from spike.ns import find

    spPr = find(shape._element, "p:spPr")
    if spPr is None:
        spPr = find(shape._element, "a:spPr")
    xfrm = find(spPr, "a:xfrm") if spPr is not None else None
    if xfrm is None:
        return "no explicit geometry override present"
    spPr.remove(xfrm)  # revert to inherited layout/master geometry
    return None


def _fix_shape_size(shape, record) -> str | None:
    m = _SIZE_RE.match(record.get("new_value") or "")
    if not m:
        return "new_value is not WxH EMU"
    error = _materialize_xfrm(shape)
    if error:
        return error
    shape.width, shape.height = int(m.group(1)), int(m.group(2))
    return None


def _fix_footer_text(shape, record) -> str | None:
    if not getattr(shape, "has_text_frame", False):
        return "shape has no text frame"
    shape.text_frame.text = record["new_value"]
    return None


def _fix_retype_upper(shape, record) -> str | None:
    """Retype every run ALL-CAPS, the designer's own move (never the cap
    attribute: the ground-truth deck has zero cap="all" runs). The original
    text is preserved in the record's old_value."""
    if not getattr(shape, "has_text_frame", False):
        return "shape has no text frame"
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.text = run.text.upper()
    return None


def _fix_strip_size_overrides(shape, record) -> str | None:
    """Remove run-level sz attributes the audit verified as restating the
    inherited size (visual no-op)."""
    from pptx.oxml.ns import qn

    if not getattr(shape, "has_text_frame", False):
        return "shape has no text frame"
    paras = shape.text_frame.paragraphs
    removed = 0
    for loc in (record.get("locator") or "").split(","):
        m = _LOC_RE.match(loc.strip())
        if not m:
            continue
        p_idx, r_idx = int(m.group(1)), int(m.group(2))
        if p_idx >= len(paras) or r_idx >= len(paras[p_idx].runs):
            return "locator no longer resolves (deck changed since audit)"
        rpr = paras[p_idx].runs[r_idx]._r.find(qn("a:rPr"))
        if rpr is not None and rpr.get("sz") is not None:
            rpr.attrib.pop("sz")
            removed += 1
    return None if removed else "no size overrides left to remove"


def _fix_set_run_sizes(shape, record) -> str | None:
    """Set every run to the sibling-majority size (sz is pt*100)."""
    from pptx.util import Pt

    try:
        size = int(record["new_value"]) / 100.0
    except (TypeError, ValueError):
        return "new_value is not an sz integer"
    if not getattr(shape, "has_text_frame", False):
        return "shape has no text frame"
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size)
    return None


def _fix_role_size(shape, record) -> str | None:
    """Set every run in the shape to the role's target size, in POINTS.

    Not _fix_set_run_sizes: that one reads new_value as an OOXML `sz`
    (hundredths of a point) because it comes from a sibling run's own
    attribute, where this reads a profile value that is stated in points.
    Sharing the function would have set a 44pt title to 0.44pt.

    WHAT THIS HARD-CODES. When the size was inherited - from the layout, the
    master or the theme - writing it onto the run pins it, so a later change to
    the master's title size stops reaching this slide. That is the same trade
    the family fix already makes (_fix_font_family), and it is why the record
    keeps its source in the message: a designer reading "source
    layout.placeholder" is being told the value is currently tracking the
    master, and ticking the fix is choosing the deck over the master for this
    run (design lead, 31/08/2026).
    """
    from pptx.util import Pt

    try:
        points = float(record["new_value"])
    except (TypeError, ValueError):
        return "new_value is not a point size"
    if points <= 0:
        return "target size is not a positive point value"
    if not getattr(shape, "has_text_frame", False):
        return "shape has no text frame"
    touched = 0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(points)
            touched += 1
    if not touched:
        return "no runs to resize"
    return None


def _fix_stop_autofit(shape, record) -> str | None:
    """'Stop Fitting Text to This Placeholder': replace the shrink-on-overflow
    autofit with noAutofit so the title renders at its intended size. The
    designer reviews the result; a long title can overflow after this."""
    from pptx.enum.text import MSO_AUTO_SIZE

    if not getattr(shape, "has_text_frame", False):
        return "shape has no text frame"
    shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    return None


def _fix_fake_slide_number(shape, record, slide=None) -> str | None:
    """Replace a literal-number text box with a real slide-number placeholder
    cloned from the layout (or master). This is the designer's own move on
    the ground-truth deck: delete the fake, insert the real placeholder,
    never nudge. The clone's xfrm is stripped so it keeps inheriting the
    layout geometry live."""
    import copy

    from pptx.enum.shapes import PP_PLACEHOLDER
    from spike.ns import find

    if slide is None:
        return "replacement needs slide context"
    source_ph, from_layout = None, False
    try:
        layout = slide.slide_layout
        sources = ((layout, True), (layout.slide_master, False))
    except Exception:
        sources = ()
    for source, is_layout in sources:
        if source is None:
            continue
        for ph in source.placeholders:
            if ph.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                source_ph, from_layout = ph, is_layout
                break
        if source_ph is not None:
            break
    if source_ph is None:
        return "layout defines no slide-number placeholder to inherit"

    already = any(
        p.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER
        for p in slide.placeholders)
    if not already:
        new_sp = copy.deepcopy(source_ph._element)
        # A layout clone inherits geometry live (slide ph idx matches the
        # layout ph). A master clone keeps its explicit xfrm: its idx has no
        # layout twin, so inheritance would be resolver-dependent.
        if from_layout:
            spPr = find(new_sp, "p:spPr")
            xfrm = find(spPr, "a:xfrm") if spPr is not None else None
            if xfrm is not None:
                spPr.remove(xfrm)
        cNvPr = find(new_sp, "p:nvSpPr/p:cNvPr")
        if cNvPr is None:
            return "layout placeholder has no id element to rewrite"
        max_id = 0
        for sh, _p in iter_shapes_deep(slide.shapes):
            max_id = max(max_id, int(sh.shape_id))
        cNvPr.set("id", str(max_id + 1))
        slide.shapes._spTree.append(new_sp)
    # the fake goes last, so a failure above leaves the slide intact
    shape._element.getparent().remove(shape._element)
    return None


def _fix_pin(shape, record, slide=None) -> str | None:
    """Snap a cross-slide anchor (title bar, header chip) back to its
    deck-wide modal position, contents riding."""
    loc = record.get("locator") or ""
    if not loc.startswith("pin:"):
        return "record has no pin target"
    try:
        tx, ty = (int(v) for v in loc[4:].split(","))
    except ValueError:
        return "pin target is not a coordinate pair"
    error = _materialize_xfrm(shape)
    if error:
        return error
    dx, dy = tx - shape.left, ty - shape.top
    carried = _carried_contents(slide, shape) if slide is not None else []
    return _translate([shape] + carried, dx, dy)


def _fix_edge(shape, record, slide=None) -> str | None:
    """Snap the shape's misaligned edge to its alignment-cluster median;
    the record's property suffix names the axis (.x = left, .y = top).
    Shapes visually contained in the moved shape travel with it, so a
    panel's labels and photos stay composed."""
    try:
        value = int(record["new_value"])
    except (TypeError, ValueError):
        return "new_value is not an EMU integer"
    error = _materialize_xfrm(shape)
    if error:
        return error
    vertical = (record.get("property") or "").endswith(".y")
    delta = value - (shape.top if vertical else shape.left)
    carried = []
    if slide is not None:
        # The record's own peers never ride it. Without this the shapes that
        # DEFINE the line get dragged off it by the shape being moved onto it
        # (see _peer_pinned). A record naming no cluster pins nothing, which is
        # the behaviour every measured edge record has always had.
        pinned = _peer_pinned(slide, record, shape.shape_id)
        carried = [c for c in _carried_contents(slide, shape,
                                                "y" if vertical else "x")
                   if str(c.shape_id) not in pinned]
    return _translate([shape] + carried,
                      0 if vertical else delta,
                      delta if vertical else 0)


def _fix_space_edge(shape, record, slide=None) -> str | None:
    """Snap every stacked element that drifted inboard back onto the frame's
    leading edge, each carrying its own contents.

    One fix for the whole stack, because "these all start on one line" is one
    decision. The record's new_value is the FRAME EDGE, not a left coordinate:
    on an RTL slide the shapes align by their right edges, so a stack of
    different widths lands on different lefts and the same edge.
    """
    if slide is None:
        return "frame snap needs slide context"
    loc = record.get("locator") or ""
    if not loc.startswith("edge:"):
        return "record has no frame-edge target"
    try:
        edge = int(record["new_value"])
    except (TypeError, ValueError):
        return "new_value is not an EMU integer"
    parts = loc.split(":", 3)
    if len(parts) < 3:
        return "frame-edge locator is malformed"
    rtl = parts[1] == "r"
    ids = [sid for sid in parts[2].split(",") if sid]
    by_id = {str(sh.shape_id): sh for sh, _p in iter_shapes_deep(slide.shapes)}
    movers = [by_id.get(sid) for sid in ids]
    if not movers or any(m is None for m in movers):
        return "stack member no longer resolves (deck changed since audit)"

    # NOTHING IN THE STACK RIDES. Shapes stacked across the axis of a move are
    # exactly what rides_with treats as a collection, so left alone every member
    # would be carried by its neighbour's translate as well as moved by its own
    # - and the member already ON the frame, which this fix deliberately does
    # not move, would be dragged a full delta to the wrong side of it. The
    # stack's position is entirely decided here: the strays go to the edge, the
    # rest stay, and neither is anybody's satellite.
    pinned = set(parts[3].split(",")) if len(parts) > 3 else set()
    pinned |= {str(m.shape_id) for m in movers}
    for member in movers:
        error = _materialize_xfrm(member)
        if error:
            return f"shape {member.shape_id}: {error}"
        target = (edge - member.width) if rtl else edge
        delta = target - member.left
        if not delta:
            continue
        carried = [c for c in _carried_contents(slide, member, "x")
                   if str(c.shape_id) not in pinned]
        error = _translate([member] + carried, delta, 0)
        if error:
            return error
    return None


def _fix_group(shape, record, slide=None) -> str | None:
    """Wrap a component's shapes in a real group.

    The group's own extent is the members' bounding box, and its CHILD extent is
    set to the same rectangle. That is what keeps the members exactly where they
    are: a child's offset is in the group's child space, so a child space equal
    to the parent space is an identity transform - the same two rectangles
    qc.design._group_transform reads from the other end.

    Inserted at the position of the member that was FIRST in the drawing order,
    so the group sits where its topmost part sat and nothing changes what covers
    what.
    """
    from lxml import etree

    if slide is None:
        return "grouping needs the slide"
    wanted = [i for i in str(record.get("new_value") or "").split(",") if i]
    if len(wanted) < 2:
        return "a group needs at least two shapes"

    spTree = slide.shapes._spTree
    by_id = {str(s.shape_id): s for s in slide.shapes}
    members = [by_id[i] for i in wanted if i in by_id]
    if len(members) != len(wanted):
        # The slide moved on since the review. Grouping what is left would be a
        # guess at what the component was.
        return "some of those shapes are no longer on the slide"

    boxes = [(s.left, s.top, s.width, s.height) for s in members]
    if any(v is None for b in boxes for v in b):
        return "one of those shapes states no size of its own"
    left = min(b[0] for b in boxes)
    top = min(b[1] for b in boxes)
    right = max(b[0] + b[2] for b in boxes)
    bottom = max(b[1] + b[3] for b in boxes)
    if right <= left or bottom <= top:
        return "those shapes have no extent between them"

    order = list(spTree)
    first = min(members, key=lambda s: order.index(s._element))
    next_id = max(int(i) for i in by_id if i.isdigit()) + 1

    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    name = f"Group {next_id}"
    grp = etree.fromstring(
        f'<p:grpSp xmlns:p="{P}" xmlns:a="{A}">'
        f'<p:nvGrpSpPr>'
        f'<p:cNvPr id="{next_id}" name="{name}"/><p:cNvGrpSpPr/><p:nvPr/>'
        f'</p:nvGrpSpPr>'
        f'<p:grpSpPr><a:xfrm>'
        f'<a:off x="{int(left)}" y="{int(top)}"/>'
        f'<a:ext cx="{int(right - left)}" cy="{int(bottom - top)}"/>'
        f'<a:chOff x="{int(left)}" y="{int(top)}"/>'
        f'<a:chExt cx="{int(right - left)}" cy="{int(bottom - top)}"/>'
        f'</a:xfrm></p:grpSpPr></p:grpSp>')

    first._element.addprevious(grp)
    for member in members:
        grp.append(member._element)
    return None


def _fix_component_edge(shape, record, slide=None) -> str | None:
    """Move a whole component onto a stated line.

    The component is named in the locator ("comp:<axis>:<ids>") because the
    entity was decided by the review that raised this, not re-derived here.
    That is the point of the layer: _carried_contents infers what travels with
    a shape from overlap and adjacency, and on a stack of blocks it decides
    each one carries its neighbours - so every member moves once for itself and
    again for the member beside it. A component states its own membership, so
    nothing rides and nothing moves twice.
    """
    if slide is None:
        return "component move needs slide context"
    loc = record.get("locator") or ""
    if not loc.startswith("comp:"):
        return "record names no component"
    parts = loc.split(":", 2)
    if len(parts) < 3:
        return "component locator is malformed"
    axis = parts[1]
    if axis not in ("top", "left", "right"):
        return f"'{axis}' is not an edge this fix understands"
    try:
        target = int(record["new_value"])
    except (TypeError, ValueError):
        return "new_value is not an EMU integer"

    by_id = {str(sh.shape_id): sh for sh, _p in iter_shapes_deep(slide.shapes)}
    members = [by_id.get(sid) for sid in parts[2].split(",") if sid]
    if not members or any(m is None for m in members):
        return "component member no longer resolves (deck changed since audit)"
    for member in members:
        error = _materialize_xfrm(member)
        if error:
            return f"shape {member.shape_id}: {error}"

    # ONE delta for the whole component, measured on its bounding box, so the
    # arrangement inside it is untouched.
    if axis == "top":
        current = min(m.top for m in members)
    elif axis == "left":
        current = min(m.left for m in members)
    else:
        current = max(m.left + m.width for m in members)
    delta = target - current
    if not delta:
        return None
    return _translate(members, 0 if axis == "top" else delta,
                      delta if axis == "top" else 0)


def _fix_distribute(record, slide) -> str | None:
    """Distribute a line of shapes evenly: first and last stay anchored,
    the middles spread so every gap matches (the designer's 'distribute
    horizontally/vertically'). Each moved shape carries its visual
    contents, same as every other positional fix."""
    loc = record.get("locator") or ""
    along_row = loc.startswith("dist-row:")
    ids = loc.split(":", 1)[1].split(",")
    by_id = {str(sh.shape_id): sh for sh, _p in iter_shapes_deep(slide.shapes)}
    line = [by_id.get(sid) for sid in ids]
    if any(sh is None for sh in line):
        return "line member no longer resolves (deck changed since audit)"
    if len(line) < 3:
        return "distribution needs at least three shapes"

    run = (lambda s: s.left) if along_row else (lambda s: s.top)
    size = (lambda s: s.width) if along_row else (lambda s: s.height)
    line = sorted(line, key=run)
    span = run(line[-1]) + size(line[-1]) - run(line[0])
    inner = span - sum(size(s) for s in line)
    if inner < 0:
        return "shapes overlap; nothing sensible to distribute"
    gap = inner // (len(line) - 1)

    moving = []
    cursor = run(line[0])
    targets = []
    for i, member in enumerate(line):
        targets.append(cursor)
        cursor += size(member) + gap
    for member, target in zip(line, targets):
        delta = target - run(member)
        if delta == 0:
            continue
        carried = _carried_contents(slide, member,
                                    "x" if along_row else "y")
        error = _translate([member] + carried,
                           delta if along_row else 0,
                           0 if along_row else delta)
        if error:
            return error
    return None


def _fix_uneven_spacing(shape, record, slide=None) -> str | None:
    """Close the odd gap by translating this shape AND everything after it
    in its row (or below it in its column) by the same delta, so no other
    gap in the line changes."""
    loc = record.get("locator") or ""
    if loc.startswith(("dist-row:", "dist-col:")):
        return _fix_distribute(record, slide) if slide is not None \
            else "distribute needs slide context"
    if loc.startswith("row:"):
        along_row = True
    elif loc.startswith("col:"):
        along_row = False
    else:
        return "record has no row context"
    if slide is None:
        return "row fix needs slide context"
    try:
        delta = int(record["new_value"]) - int(record["old_value"])
    except (TypeError, ValueError):
        return "gap values are not EMU integers"
    line_ids = loc.split(":", 1)[1].split(",")
    by_id = {str(sh.shape_id): sh for sh, _p in iter_shapes_deep(slide.shapes)}
    line = [by_id.get(sid) for sid in line_ids]
    if any(sh is None for sh in line):
        return "row member no longer resolves (deck changed since audit)"
    target = str(record["shape_id"])
    start = line_ids.index(target)
    moving = list(line[start:])
    moving_ids = {str(m.shape_id) for m in moving}
    # each moved shape carries its collection: contents, and the partners
    # stacked across the axis of the move (a chip on its box tears off the pair
    # otherwise - real-deck finding, 12/08/2026, RTL strategy map)
    for sat in _collection_riders(slide, [str(m.shape_id) for m in line[start:]],
                                  "x" if along_row else "y"):
        sid = str(sat.shape_id)
        if sid not in moving_ids:
            moving_ids.add(sid)
            moving.append(sat)
    return _translate(moving,
                      delta if along_row else 0,
                      0 if along_row else delta)


_FIXERS = {
    "font.family_out_of_set": _fix_font_family,
    "master_slide.placeholder_geometry_off": _fix_placeholder_geometry,
    "shape_size.size_mismatch": _fix_shape_size,
    "header_footer.text_mismatch": _fix_footer_text,
    "margin_alignment.edge_misaligned": _fix_edge,
    "margin_alignment.uneven_spacing": _fix_uneven_spacing,
    "font.title_autofit_shrunk": _fix_stop_autofit,
    "header_footer.fake_slide_number": _fix_fake_slide_number,
    "header_footer.footer_off_canvas": _fix_edge,
    "margin_alignment.panel_row_misaligned": _fix_edge,
    "margin_alignment.cluster_rhythm": _fix_lift,
    "typography.case_inconsistent": _fix_retype_upper,
    "typography.redundant_size_override": _fix_strip_size_overrides,
    "typography.size_inconsistent": _fix_set_run_sizes,
    "font.size_off_role": _fix_role_size,
    "margin_alignment.content_overflow": _fix_rescale,
    "font.cs_typeface_missing": _fix_font_family,
    "margin_alignment.recurring_off_position": _fix_pin,
    "margin_alignment.body_band_intrusion": _fix_body_band,
    # the same single translate; only the sign of the delta differs
    "margin_alignment.body_below_band": _fix_body_band,
    "color_palette.off_palette_rgb": _fix_off_palette_color,
    # a single-axis snap to a stated edge: the same fixer the cluster snap
    # uses, because the only difference is where the target came from
    "margin_alignment.space_edge_misaligned": _fix_space_edge,
    "margin_alignment.component_edge_misaligned": _fix_component_edge,
    "margin_alignment.should_be_grouped": _fix_group,
}


@dataclass
class FixResult:
    cleaned_bytes: bytes
    outcomes: list[FixOutcome]
    applied: int
    records: list[dict]  # updated copies; applied ones carry action="changed"


def _apply_unify(deck_bytes: bytes, unify_recs: list[dict],
                 outcomes: list[FixOutcome]) -> tuple[bytes, int]:
    """Master-unification records operate on the package, not on shapes:
    clone repoints (locator 'dedup:<layout part>') run everywhere; layout
    re-application (locator 'com:<layout name>') needs desktop PowerPoint."""
    from .unify import com_unify, dedup

    applied = 0
    repoints: dict[int, tuple[str, str]] = {}   # slide -> (record_id, part)
    assignments: dict[int, tuple[str, str]] = {}  # slide -> (record_id, name)
    for rec in unify_recs:
        loc = rec.get("locator") or ""
        if loc.startswith("dedup:"):
            repoints[rec["slide_index"]] = (rec["record_id"], loc[6:])
        elif loc.startswith("com:"):
            assignments[rec["slide_index"]] = (rec["record_id"], loc[4:])
        else:
            outcomes.append(FixOutcome(rec["record_id"], "skipped",
                                       "record has no unify route"))

    if repoints:
        try:
            deck_bytes = dedup(deck_bytes,
                               {i: part for i, (_r, part) in repoints.items()})
            for _i, (rid, _p) in sorted(repoints.items()):
                outcomes.append(FixOutcome(rid, "changed"))
                applied += 1
        except Exception as exc:
            for _i, (rid, _p) in sorted(repoints.items()):
                outcomes.append(FixOutcome(rid, "skipped", f"dedup failed: {exc}"))

    if assignments:
        new_bytes, errors = com_unify(
            deck_bytes, {i: name for i, (_r, name) in assignments.items()})
        if new_bytes is None:
            reason = next(iter(errors.values()), "PowerPoint unavailable")
            for _i, (rid, _n) in sorted(assignments.items()):
                outcomes.append(FixOutcome(rid, "skipped", reason))
        else:
            deck_bytes = new_bytes
            for i, (rid, _n) in sorted(assignments.items()):
                if i in errors:
                    outcomes.append(FixOutcome(rid, "skipped", errors[i]))
                else:
                    outcomes.append(FixOutcome(rid, "changed"))
                    applied += 1
    return deck_bytes, applied


def apply_fixes(deck_bytes: bytes, records: list[dict], record_ids: set[str]) -> FixResult:
    """Open the deck from bytes, apply the selected fixable records, return
    the cleaned deck as new bytes plus per-record outcomes."""
    prs = Presentation(io.BytesIO(deck_bytes))
    by_id = {r["record_id"]: r for r in records}
    outcomes: list[FixOutcome] = []
    unify_recs: list[dict] = []
    applied = 0

    # Deterministic semantic order: geometry snaps (which REMOVE the xfrm to
    # re-inherit) run before positional edits (which SET xfrm values), so a
    # shape receiving both is first snapped to its layout baseline and then
    # nudged from there. UUID order made this a per-shape coin flip and, with
    # the old fixers, could leave an off-without-ext transform (real-deck
    # corruption, 14/07/2026).
    # Panel-row moves run late: a contained shape's fix (computed relative
    # to its container's audit-time position) must land before the panel
    # moves and carries it, so both fixes compose to the designer's result
    # instead of double-moving the content (ground truth, 20/07/2026: band
    # row + icon inset on the capabilities slide only converge in this
    # order). Off-canvas footer fixes run LAST of all: their target zone is
    # often occupied until a rhythm lift clears it (the wheel slide's
    # bottom clusters sat exactly where the source line belongs).
    # Rescale runs after every in-selection fix (their targets were computed
    # on unscaled geometry; alignment survives an affine transform) but
    # before the footer fix (whose target is an absolute layout baseline).
    # The header-band push shares the rescale's stage for the same reason and
    # one more: both move the whole body, so two of them on one slide are
    # genuinely ambiguous, and sharing a stage makes the claim check skip the
    # second with a reason instead of composing two guesses.
    def _rank(rec) -> int:
        issue = rec.get("issue_type") if rec else None
        return {"master_slide.placeholder_geometry_off": 0,
                "margin_alignment.panel_row_misaligned": 2,
                "margin_alignment.content_overflow": 3,
                "margin_alignment.body_band_intrusion": 3,
                "margin_alignment.body_below_band": 3,
                "header_footer.footer_off_canvas": 4}.get(issue, 1)

    def _order(rid: str):
        return (_rank(by_id.get(rid)), rid)

    # One move per shape per STAGE per round: two same-stage row fixes with
    # different deltas on vertically paired shapes (a chip and its goals
    # box) tear the pair apart, and a shape can even ride one fix as a
    # satellite and then be moved again by its own (real-deck finding,
    # 12/08/2026, RTL strategy map). The first fix claims its movers;
    # same-stage fixes touching any of them wait for the next re-audit
    # round. CROSS-stage double-moves are the designed pipeline (icon
    # inset, then its panel carries it; rescale transforms everything).
    claimed: dict[tuple, set] = {}

    for rid in sorted(record_ids, key=_order):
        rec = by_id.get(rid)
        if rec is None:
            outcomes.append(FixOutcome(rid, "skipped", "unknown record id"))
            continue
        if not is_fixable(rec):
            reason = "Arabic content, manual review" if rec.get("arabic_flag") \
                else "not a fixable record"
            outcomes.append(FixOutcome(rid, "skipped", reason))
            continue
        slide_index = rec["slide_index"]
        if rec["issue_type"] == "master_slide.foreign_master":
            if slide_index >= len(prs.slides):
                outcomes.append(FixOutcome(rid, "skipped", "slide index out of range"))
            else:
                unify_recs.append(rec)  # package-level; applied after this pass
            continue
        if slide_index >= len(prs.slides):
            outcomes.append(FixOutcome(rid, "skipped", "slide index out of range"))
            continue
        shape = _find_shape(prs.slides[slide_index], rec["shape_id"])
        if shape is None:
            outcomes.append(FixOutcome(rid, "skipped", "shape not found"))
            continue
        slide = prs.slides[slide_index]
        positional = rec["issue_type"] in _POSITIONAL_ISSUES
        before = _slide_rects(slide) if positional else None
        # planned BEFORE the fix runs: containment is judged on pre-move
        # geometry, and the guard must be able to restore every moved shape
        movers = _planned_movers(slide, rec) if positional else set()

        if positional:
            taken = claimed.setdefault((slide_index, _rank(rec)), set())
            if movers & taken:
                outcomes.append(FixOutcome(
                    rid, "skipped",
                    "shares shapes with a fix already applied this round; "
                    "re-audit and apply again"))
                continue

        fixer = _FIXERS[rec["issue_type"]]
        if fixer in (_fix_uneven_spacing, _fix_edge, _fix_fake_slide_number,
                     _fix_lift, _fix_rescale, _fix_pin, _fix_body_band,
                     _fix_space_edge, _fix_component_edge, _fix_group):
            error = fixer(shape, rec, slide=slide)
        else:
            error = fixer(shape, rec)

        if not error and positional:
            error = _collision_created(slide, movers, before)
            if error:
                _restore_rects(slide, movers, before)

        if error:
            outcomes.append(FixOutcome(rid, "skipped", error))
        else:
            outcomes.append(FixOutcome(rid, "changed"))
            applied += 1
            if positional:
                claimed.setdefault((slide_index, _rank(rec)),
                                   set()).update(movers)

    out = io.BytesIO()
    prs.save(out)
    cleaned = out.getvalue()

    if unify_recs:
        cleaned, unify_applied = _apply_unify(cleaned, unify_recs, outcomes)
        applied += unify_applied

    changed_ids = {o.record_id for o in outcomes if o.outcome == "changed"}
    updated = []
    for r in records:
        r2 = dict(r)
        if r2["record_id"] in changed_ids:
            r2["action"] = "changed"
        updated.append(r2)
    return FixResult(cleaned_bytes=cleaned, outcomes=outcomes,
                     applied=applied, records=updated)
