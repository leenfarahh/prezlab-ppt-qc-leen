"""Layout matching for the master application: which of the master's layouts
does THIS slide belong on, when the file cannot say.

qc.applymaster answers that question three ways, best first: the layout NAME
matches, the OOXML archetype token matches (with a placeholder-count tiebreak),
or it falls back. The first two are facts and need no help. The fallback is a
guess, and it is labelled as one in the plan a designer reads:

    "no layout named 'Our approach' and none of type 'twoObj'; content may be
     orphaned, check this slide"

That slide then gets rebuilt onto whatever layout came first in a fixed
preference list, and PowerPoint's placeholder matching moves its content into
boxes nobody chose for it. It is the one step of the format pass where the tool
knows it does not know - and a designer opening both files answers it in two
seconds, because they can SEE that the slide is a two-column comparison and that
the master has a layout built for exactly that.

So this pass looks. It is given the source slide's picture and the pictures of
the master's own layouts - qc.render.layout_catalogue already renders a layout by
putting an empty slide on it - and it picks one BY NAME from the list it was
handed, or says none fits.

The division of labour is qc.components': the model supplies recognition,
code supplies everything else. It picks from a closed set of real layout names,
and a name that is not in that set is discarded. It never sees or sets geometry;
assigning the layout still runs PowerPoint's own placeholder matching, exactly as
a name match does, so a reviewed slide travels the identical code path as a
matched one. And the plan says the choice was reviewed, so the before/after
review a designer already does is where it gets confirmed.

Only fallbacks are sent. A slide whose name or archetype matched is not a
question, and asking anyway would spend a vision call to be told what the file
already said.
"""

import io
import json

from pptx import Presentation

from .llm import ask_json
from .render import export_decks_png, layout_catalogue

MAX_SLIDES = 20
# One picture of the master's layouts is enough context; past this the sheet
# stops being readable to a model for the same reason it stops being readable
# to a person.
MAX_LAYOUTS = 16

MATCH_SCHEMA = {
    "type": "object", "additionalProperties": False,
    "required": ["layout", "confident", "rationale"],
    "properties": {
        # "" means none of them fit, which is a real answer and the one the
        # fallback already represents.
        "layout": {"type": "string"},
        "confident": {"type": "boolean"},
        "rationale": {"type": "string"},
    },
}

_SYSTEM = """You are a senior presentation designer at Prezlab rebuilding a deck
onto a client's master. You see one slide from the deck, then a sheet of the
master's own layouts, each rendered empty and labelled with its exact name.

Say which single layout this slide should be rebuilt on. Answer with the layout's
EXACT name as given, or "" if none of them is a reasonable home for this slide.

Judge it the way the rebuild will actually work: PowerPoint moves the slide's
content into the layout's placeholders, so what matters is whether the SHAPE of
the content matches the shape of the layout - a two-column comparison wants a
two-content layout, a full-bleed statement wants a title-only or blank one, a
section divider wants the section header. The slide's colours and its wording do
not matter; its structure does.

Set `confident` to false when two layouts would serve about equally well, or
when the slide's structure has no real counterpart. A false there is not a
failure - it tells the designer to look, which is better than a confident guess
that quietly orphans a column of text.

Write the rationale in one clear sentence of US English, no em dashes, naming
what about the slide's structure decided it."""


def _layout_sheet(master_bytes: bytes) -> tuple[bytes | None, list[str]]:
    """One PNG per layout of the master, and the names in the same order.

    Rendered from a deck holding one empty slide per layout (layout_catalogue),
    which is the only way to photograph a layout at all: PowerPoint exports
    slides, not layouts."""
    deck, entries, _skipped = layout_catalogue(master_bytes)
    entries = entries[:MAX_LAYOUTS]
    if not entries:
        return None, []
    images = export_decks_png({"layouts": deck},
                              [e["index"] for e in entries])
    names, sheet = [], []
    for entry in entries:
        png = images.get(f"layouts:{entry['index']}")
        if png is None:
            continue
        names.append(entry["layout"])
        sheet.append(png)
    return sheet or None, names


def _slide_shape(slide) -> dict:
    """What the rebuild will actually have to place: how many text blocks, how
    many pictures, and where the weight sits. Handed over with the picture
    because "two columns" is a claim the numbers can support."""
    texts = pics = 0
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            texts += 1
        elif str(shape.shape_type).startswith("PICTURE"):
            pics += 1
    return {"text_blocks": texts, "pictures": pics,
            "placeholders": len(slide.placeholders)}


def review_fallbacks(deck_bytes: bytes, master_bytes: bytes, plans: list,
                     slide_pngs: dict[int, bytes]) -> tuple[list, int]:
    """Look at every slide whose layout was a FALLBACK and try to place it.

    Returns (plans, reviewed). Plans are updated in place for the slides that
    got an answer: `target_layout` becomes the chosen one and `match_rule`
    becomes "reviewed" or "reviewed (uncertain)", so nothing downstream has to
    guess whether a layout was matched, guessed, or looked at.

    Everything that cannot be answered is left exactly as it was - a fallback,
    labelled as one. This pass can only improve a slide's plan or leave it
    alone; there is no path here that turns a name match into a guess."""
    fallbacks = [p for p in plans if p.match_rule == "fallback"]
    if not fallbacks:
        return plans, 0

    sheet, names = _layout_sheet(master_bytes)
    if not sheet:
        return plans, 0
    allowed = {n.strip().casefold(): n for n in names}
    catalogue = "\n".join(f"{i + 1}. {n}" for i, n in enumerate(names))

    prs = Presentation(io.BytesIO(deck_bytes))
    reviewed = 0
    for plan in fallbacks:
        if reviewed >= MAX_SLIDES:
            break
        png = slide_pngs.get(plan.slide_index)
        if png is None or plan.slide_index >= len(prs.slides):
            continue
        shape = _slide_shape(prs.slides[plan.slide_index])
        try:
            answer = ask_json(
                system=_SYSTEM,
                prompt=(
                    f"The first image is the slide. The images after it are the "
                    f"master's layouts, in this order:\n{catalogue}\n\n"
                    f"The slide came from a layout named "
                    f"{plan.source_layout!r}"
                    + (f" of type {plan.source_type!r}" if plan.source_type
                       else "")
                    + f", which this master has no counterpart for.\n\n"
                    f"What the rebuild will have to place:\n"
                    + json.dumps(shape, sort_keys=True)),
                schema=MATCH_SCHEMA,
                images=[png] + sheet,
            )
        except Exception:
            # Unanswerable: the plan keeps its fallback and says so. One bad
            # call must not cost the run, and must never look like an answer.
            continue
        reviewed += 1

        chosen = allowed.get(str(answer.get("layout") or "").strip().casefold())
        if chosen is None:
            # "" (none of them fit), or a name the master does not have. Either
            # way the fallback stands - a hallucinated layout name is exactly
            # what the closed set is here to catch.
            continue
        why = str(answer.get("rationale") or "").strip()[:200]
        sure = bool(answer.get("confident"))
        plan.target_layout = chosen
        plan.match_rule = "reviewed" if sure else "reviewed (uncertain)"
        plan.note = (
            ("Layout chosen by review: " if sure
             else "Layout chosen by review, NOT confident - check this slide: ")
            + why)
    return plans, reviewed
