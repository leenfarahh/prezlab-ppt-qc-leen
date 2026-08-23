"""Move a slide's content into the master it was just given.

Applying a layout is only half the job. PowerPoint's CustomLayout assignment
remaps content that lives in PLACEHOLDERS; free-floating shapes are left
exactly where they were. A deck built by an export tool has no placeholders at
all, so the visible result is the master's empty placeholders showing their
prompt text on top of the original content, which is what a designer sees as
"the master applied but nothing moved".

This pass closes that gap, per slide:

    1. title / subtitle text moves INTO the master's placeholders, which then
       style and position it - including where the text sits inside its box.
       The master's own vertical anchor is followed, never overridden: a title
       the master hangs at the bottom of its box is a design decision, and
       hoisting it to the top only moved the empty space to the other side of
       the heading (design lead, 21/08/2026)
    2. everything else translates as one block into the master's content
       region, so relative arrangement is preserved exactly
    3. page furniture the master now supplies is removed from the content
    4. placeholders left empty are removed, so no prompt text remains

Which margin the block is seated against follows the SCRIPT: an Arabic slide
reads right to left, so its block binds to the right margin, an English one to
the left. Judged per slide, because bilingual decks run both.

Three deliberate restraints:

A HEADING that ends up past a margin is reported and left exactly as it is.
Whether a title or standfirst may break the margin frame is the client's house
style, not a defect, so nothing here moves or resizes one to fit; the designer
gets the fact and asks. What is measured is the BOX, never the text drawn
inside it - a long line spilling out of a correctly placed box is a copy-length
conversation, and estimating glyph widths would replace a fact about the deck
with a guess about the reader's fonts.

Text moved into a placeholder is moved as TEXT, not as formatted runs. The
placeholder then styles it from the master, which is the whole point of
"match the master". Run-level overrides (a bold word, a coloured span) are
lost, and every such move is reported so the loss is visible rather than
discovered later.

Remaining content is TRANSLATED, never scaled. Scaling a text box does not
scale its font, so a "helpful" shrink produces overflowing text that looks
fine in the XML and broken on screen. When a block genuinely does not fit, it
is still seated on the line the master draws for the top of its body - the
strip a master reserves under the header stays empty on every slide - and the
overflow at the bottom is reported for a designer to rework. Where the master
draws no such line, the old restraint holds: the block is moved as far as it
can go and no further.
"""

import io
from dataclasses import dataclass

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from spike.ns import find
from .util import (full_height_panel, heading_ids, iter_shapes_deep,
                   max_font_pt, recurring_anchors, slide_is_rtl)

# Placeholder roles this pass will fill from free content.
_TITLE_TYPES = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
_SUBTITLE_TYPES = (PP_PLACEHOLDER.SUBTITLE,)
# Furniture the master provides; a slide's own copy is then duplication.
_FURNITURE_TYPES = (PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER,
                    PP_PLACEHOLDER.DATE)

# A shape must overlap a placeholder by at least this share of its own area
# before overlap counts as evidence it belongs there.
MIN_OVERLAP_SHARE = 0.25
# Any overlap this large with a placeholder that now holds text is a visible
# collision. Far lower than the matching threshold above: deciding a shape
# BELONGS somewhere needs real evidence, whereas noticing two things printing
# on top of each other does not.
MIN_COLLISION_SHARE = 0.08
# Gap left below the lowest header placeholder before content starts.
CONTENT_GAP_EMU = 137160  # 0.15in, the alignment module's intent window
# Bottom strip where page furniture lives.
FURNITURE_BAND = 0.88
# Content this wide is a deliberate full-bleed element, not indented content,
# and the left margin does not apply to it.
FULL_BLEED_SHARE = 0.97
# A shape sharing this much of its own area with a full-bleed element is sitting
# ON it - a logo stamped on a footer band, a mark on a coloured stripe - and
# belongs to it rather than to the text elsewhere on the slide. Not strict
# containment: a logo routinely overhangs the band it sits on by a hair.
ON_A_BLEED_SHARE = 0.9
# A shape smaller than this on BOTH sides is a marker, not content, and does not
# get to say where the content block begins. think-cell parks a 0.0017in square
# named "think-cell data - do not delete" in the top-left corner of every slide
# it has ever touched, and PowerPoint leaves the same stub behind for any
# embedded OLE object; a shape 0.04mm across cannot be the thing a designer
# seated on a margin. BOTH sides have to fail: a hairline rule or a thin divider
# is degenerate on ONE axis and is content.
MIN_ANCHOR_SIDE_EMU = 45720  # 0.05in
# Frame sources the master STATES rather than the tool inferring: a rectangle
# the designer drew and named, or the guides they set. A frame derived from
# placeholder extents is an inference and does not get to bind anything.
# (qc.stylespec.infer_grid)
STATED_FRAME_SOURCES = ("presentation_space", "guides")

# A heading box within this of a margin is ON it. Guides are stored in eighths
# of a point while placeholder extents are typed in inches, so a title a
# designer placed exactly on the guide reads a few hundred EMU past it; that is
# rounding, not a heading breaking the frame.
# (keep in sync with qc/modules/margin_alignment.py HEADING_SLACK_EMU)
HEADING_SLACK_EMU = 36000  # 1mm


@dataclass
class ContentChange:
    slide_index: int      # zero-based
    action: str
    detail: str
    # "info" for routine moves, "alert" for anything a designer must look at.
    # Removals are always alerts: content leaving the deck is the one outcome
    # nobody should have to notice for themselves.
    severity: str = "info"
    # Full text of anything removed, so the report can offer it back verbatim
    # rather than a truncated preview a designer would have to retype.
    removed_text: str | None = None
    # The removed shape's own XML, and a handle for it. Reporting the text was
    # only half an answer: a designer who wants it back had to retype it and
    # place it by eye. With the element kept, putting it back is exact - same
    # words, same box, same formatting - and the decision stays theirs
    # (design lead, 20/08/2026: "add the option to bring back the pieces that
    # were removed").
    removed_xml: str | None = None
    restore_id: str | None = None
    # How to put THIS change back, exactly, and a handle to ask for it by. A
    # list of operations replayed in order (qc.undo); None means there is
    # nothing to revert, which is two different things and the review page says
    # which: a change that only REPORTS (a heading left sitting past a margin
    # moved nothing), or one this pass did not make in the first place (the
    # layout assignment is PowerPoint's own work - see qc.applymaster - and
    # cannot be taken back a slide at a time from here).
    #
    # Every operation stores the STATE, never a delta to re-derive: an offset
    # carries the coordinates the shape had, a removal carries the element
    # itself. Re-deriving would make undo a second guess at the same problem,
    # and a wrong guess here silently damages a deck a designer has approved.
    change_id: str | None = None
    undo: list[dict] | None = None

    def __str__(self) -> str:
        mark = "!! " if self.severity == "alert" else ""
        return f"slide {self.slide_index + 1}: {mark}{self.action} - {self.detail}"


def _box(shape):
    l, t, w, h = shape.left, shape.top, shape.width, shape.height
    if None in (l, t, w, h) or w <= 0 or h <= 0:
        return None
    return (l, t, l + w, t + h)


def _anchor_boxes(boxes) -> list:
    """The boxes allowed to define the content block's edges.

    Degenerate shapes are dropped here and nowhere else: they still TRAVEL with
    the block, because leaving a stub behind while its slide moves serves
    nobody, but they do not get to say where the block BEGINS. A think-cell data
    object sitting at (0.002in, 0.002in) made itself the top-left corner of the
    block on 24 of a 26-slide deck's slides, so seating the block on the
    master's body line moved the real content - which was already on that line -
    1.90in DOWN, and 1.1in to 1.9in of every slide off the bottom of the page
    (real deck, 23/08/2026). The move was reported as an overflow each time,
    which is not the same as not making it."""
    return [b for b in boxes
            if (b[2] - b[0]) >= MIN_ANCHOR_SIDE_EMU
            or (b[3] - b[1]) >= MIN_ANCHOR_SIDE_EMU]


def _overlap_share(shape_box, ph_box) -> float:
    """Overlap area as a share of the SHAPE's area: a small eyebrow fully
    inside the title band scores 1.0, while a full-width banner that merely
    crosses it scores low."""
    if shape_box is None or ph_box is None:
        return 0.0
    x0 = max(shape_box[0], ph_box[0])
    y0 = max(shape_box[1], ph_box[1])
    x1 = min(shape_box[2], ph_box[2])
    y1 = min(shape_box[3], ph_box[3])
    if x1 <= x0 or y1 <= y0:
        return 0.0
    area = (shape_box[2] - shape_box[0]) * (shape_box[3] - shape_box[1])
    return ((x1 - x0) * (y1 - y0)) / area if area else 0.0


# Shared with the audit's heading rule (qc.util), so the two passes cannot
# disagree about which shape on a slide is the title.
_max_font_pt = max_font_pt


def _text_of(shape) -> str:
    return shape.text_frame.text.strip() if getattr(
        shape, "has_text_frame", False) else ""


def _delete(shape):
    el = shape._element
    el.getparent().remove(el)


def _shape_xml(shape) -> str | None:
    """The shape's own element, serialized. Kept so a removal can be undone
    exactly rather than described: same words, same box, same formatting."""
    from lxml import etree

    try:
        return etree.tostring(shape._element, encoding="unicode")
    except Exception:
        return None


def _insert_undo(shape) -> list[dict] | None:
    """The undo for a REMOVAL: put this element back. Called before the delete,
    because after it the element is no longer reachable from the tree."""
    xml = _shape_xml(shape)
    return [{"op": "insert", "xml": xml}] if xml else None


def _offset_undo(shapes) -> list[dict] | None:
    """The undo for a MOVE: the coordinates these shapes hold right now. Called
    before the move, and it records where each one WAS rather than how far it
    is about to travel, so replaying it cannot drift."""
    ops = [{"op": "offset", "shape_id": str(s.shape_id),
            "left": int(s.left), "top": int(s.top)}
           for s in shapes if s.left is not None and s.top is not None]
    return ops or None


def _renumber(element, first_id: int) -> None:
    """Give the restored subtree fresh shape ids starting at first_id.

    Shape ids are unique per slide, and the id the piece had when it was removed
    is not reserved for it. Putting one back with a stale id is a duplicate id,
    which PowerPoint reports as a damaged file rather than as a duplicate (bug
    found 20/08/2026: a repeat restore, e.g. a browser resubmitting the form,
    corrupted the deck)."""
    from pptx.oxml.ns import qn

    for i, cNvPr in enumerate(element.iter(qn("p:cNvPr"))):
        cNvPr.set("id", str(first_id + i))


def _covers(slide, shape) -> list:
    """Names of the shapes a restored piece now prints over, if any.

    Reported, never resolved. A piece goes back at the coordinates it was
    removed from - that is where the designer put it, and it is what makes the
    restore an undo rather than a re-layout (design lead, 21/08/2026: "ticked
    pieces should be put back where they were originally, not into the body
    area"). An earlier version hunted for the nearest clear space instead, and
    landed an eyebrow three inches down the slide, which is not where an eyebrow
    belongs however empty that spot was."""
    box = _box(shape)
    if box is None:
        return []
    names = []
    for other in slide.shapes:
        if other.shape_id == shape.shape_id:
            continue
        other_box = _box(other)
        if other_box is None:
            continue
        if (_overlap_share(box, other_box) >= MIN_COLLISION_SHARE
                or _overlap_share(other_box, box) >= MIN_COLLISION_SHARE):
            names.append(_ph_label(other) if other.is_placeholder
                         else (_text_of(other)[:24] or other.name))
    return names


def restore_shapes(deck_bytes: bytes, items: list) -> tuple[bytes, list]:
    """Put removed shapes back, one per (slide_index, xml) item. Returns the
    new deck bytes and one outcome dict per piece that went back
    ({restore_id, shifted_emu, detail}).

    The counterpart to the sweep above, and the reason the sweep is defensible
    at all: the tool decides what the master has no room for, the designer
    decides whether that decision was right (design lead, 20/08/2026). The
    element goes back LAST among the shapes, so it lands on top in z-order -
    visible, not hidden under whatever took its place.

    Three things separate this from a plain append, and all three were real
    damage (20/08/2026, "bringing back selected pieces screwed up the whole
    presentation content"):

    - p:extLst must be the LAST child of a shape tree. PowerPoint writes one on
      some slides, so appending put the shape after it and the deck opened with
      a repair prompt. python-pptx never appends for this reason either; it
      inserts before p:extLst.
    - the ids are renumbered, so restoring the same piece twice cannot collide
      with itself.

    The piece goes back at its ORIGINAL coordinates. Where the master has since
    filled that spot, the piece prints over it and the report says exactly what
    it covers: the designer asked for this content back, and deciding where it
    belongs on the rebuilt slide is their call, not this pass's.

    The XML is parsed with python-pptx's OWN parser, not lxml's. lxml returns
    generic elements, and a generic element spliced into a python-pptx tree is a
    shape the library can no longer read: touching it raises, and nothing that
    walks the slide afterwards - here, the audit, or the next pass - can be
    trusted to see it correctly."""
    from pptx.oxml import parse_xml

    prs = Presentation(io.BytesIO(deck_bytes))
    outcomes = []
    for item in items:
        index = item.get("slide_index")
        xml = item.get("removed_xml")
        if xml is None or index is None or index >= len(prs.slides):
            continue
        try:
            element = parse_xml(xml)
        except Exception:
            continue
        slide = prs.slides[index]
        first_id = slide.shapes._next_shape_id
        _renumber(element, first_id)
        slide.shapes._spTree.insert_element_before(element, "p:extLst")
        shape = next((s for s in slide.shapes
                      if str(s.shape_id) == str(first_id)), None)
        covers = _covers(slide, shape) if shape is not None else []
        if shape is not None:
            # Findable in PowerPoint's selection pane, which is how a designer
            # gets to a piece that now overlaps something.
            shape.name = f"RESTORED {shape.name}"[:255]
        detail = "back where it was"
        if covers:
            detail += (f", printing over {', '.join(covers[:3])}"
                       + (f" and {len(covers) - 3} more" if len(covers) > 3
                          else ""))
        outcomes.append({
            "restore_id": item.get("restore_id"),
            "slide_index": index,
            "covers": covers,
            "detail": detail,
        })
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue(), outcomes


def _direction_of(paragraph) -> tuple:
    """(rtl, algn) as the paragraph STATES them, either None when it does not."""
    pPr = find(paragraph._p, "a:pPr")
    if pPr is None:
        return None, None
    rtl = pPr.get("rtl")
    return (rtl in ("1", "true") if rtl is not None else None), pPr.get("algn")


def _set_direction(paragraph, rtl: bool, algn: str | None) -> None:
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1" if rtl else "0")
    pPr.set("algn", algn or ("r" if rtl else "l"))


def _set_text(placeholder, text: str, source=None) -> bool:
    """Replace a placeholder's text, keeping paragraph breaks. Returns True when
    any paragraph had to be marked right-to-left.

    Run-level formatting is intentionally not carried across: the placeholder
    must style this from the master. DIRECTION is not formatting, though, and
    that distinction was a real defect (design lead, 21/08/2026): the source
    paragraph said rtl="1" algn="r", the plain-text move dropped both, and the
    Arabic title inherited the English master's left-to-right alignment - so
    every Arabic deck came back with its headings starting on the wrong edge.

    Which way a language reads is a fact about the language, not a house style
    the master gets to overrule. The source paragraph's own direction wins; with
    nothing stated, Arabic text is marked RTL and starts on the right."""
    from spike.arabic import contains_arabic

    tf = placeholder.text_frame
    lines = [ln for ln in text.split("\n")]
    tf.text = lines[0] if lines else ""
    for line in lines[1:]:
        tf.add_paragraph().text = line

    stated = []
    if source is not None and getattr(source, "has_text_frame", False):
        stated = [_direction_of(p) for p in source.text_frame.paragraphs
                  if p.text.strip()]
    marked = False
    for i, para in enumerate(tf.paragraphs):
        src_rtl, src_algn = stated[i] if i < len(stated) else (None, None)
        rtl = src_rtl if src_rtl is not None else contains_arabic(para.text)
        if not rtl:
            continue
        _set_direction(para, True, src_algn)
        marked = True
    return marked


def _placeholders_by_type(container):
    found = {}
    for ph in container.placeholders:
        try:
            found.setdefault(ph.placeholder_format.type, []).append(ph)
        except Exception:
            continue
    return found


def _layout_supplies(slide, types) -> bool:
    """Whether the slide's LAYOUT defines these placeholders.

    Read from the layout, not the slide: PowerPoint does not materialise an
    empty footer or slide-number placeholder on every slide, so asking the
    slide would answer "no furniture" for a master that clearly has some, and
    the deck would keep its hand-drawn duplicates."""
    try:
        by_type = _placeholders_by_type(slide.slide_layout)
    except Exception:
        return False
    return any(by_type.get(t) for t in types)


def _pick(candidates, ph, slide, prs, prefer_largest: bool):
    """The free shape that best belongs in this placeholder.

    Overlap with the master's own placeholder box qualifies a candidate,
    because the master is the statement of where a title goes. Type size then
    decides, since two shapes routinely share the title band (a small-caps
    eyebrow above a heading) and only one of them is the heading.

    When the cascade cannot resolve a point size, box height stands in. That
    fallback is not cosmetic: with every size reading 0 the sort collapsed to
    overlap alone, and a compact eyebrow sitting wholly inside the title band
    beats a wide heading that overhangs it. That is exactly how the eyebrow
    ended up in the title placeholder with the real heading left on top of
    it."""
    scored = []
    for shape in candidates:
        share = _overlap_share(_box(shape), _box(ph))
        if share < MIN_OVERLAP_SHARE:
            continue
        size = _max_font_pt(shape, slide, prs)
        height_pt = (shape.height or 0) / 12700.0
        scored.append({"share": share, "size": size,
                       "height_pt": height_pt, "shape": shape})
    if not scored:
        return None

    # Mixing resolved point sizes with height proxies across candidates would
    # compare different units, so pick one basis for the whole comparison.
    basis = "size" if any(s["size"] for s in scored) else "height_pt"
    if prefer_largest:
        # Box height breaks a size tie before overlap does: an eyebrow and a
        # heading can both resolve to the same inherited size AND both sit
        # wholly inside the title band, and then insertion order was deciding
        # which one became the title.
        scored.sort(key=lambda s: (s[basis], s["height_pt"], s["share"]),
                    reverse=True)
    else:
        scored.sort(key=lambda s: (s["share"], s[basis]), reverse=True)
    return scored[0]["shape"]


def _header_floor(title_ph, sub_ph) -> int:
    """Bottom of the master's header band: nothing above this is content."""
    floor = 0
    for ph in (title_ph, sub_ph):
        b = _box(ph) if ph is not None else None
        if b:
            floor = max(floor, b[3])
    return floor


# A standfirst is a line of text, not a panel. These bound what the subtitle
# fallback will accept, because it CONSUMES the shape it picks: an unbounded
# version swallowed a full-slide diagram into the subtitle placeholder and
# deleted it. Destroying content is far worse than leaving a subtitle empty.
MAX_SUBTITLE_AREA_SHARE = 0.10
MAX_SUBTITLE_CHARS = 200


def _header_sized(box, text: str, header_floor: int, prs) -> bool:
    """Whether a shape is a LINE OF HEADER TEXT rather than body content that
    merely starts high on the slide.

    The single gate for both ranking a title candidate and sweeping unplaced
    text away, deliberately shared: anything eligible for removal must be
    something that could have been a heading in the first place. Testing only
    the top edge deleted a full-slide diagram whose top sat 0.15in above the
    header floor."""
    if box is None:
        return False
    if not text:
        # No text, no "header text". A corner rule, a bracket or a mark drawn in
        # the header band is part of a composition, and the sweep below REMOVES
        # what this returns true for: an empty-text graphic was being deleted
        # and reported as "unplaced text", which is a defect however it is
        # phrased (design lead, 20/08/2026).
        return False
    if box[3] > header_floor + CONTENT_GAP_EMU:
        return False  # extends into the body: not header text
    area = (box[2] - box[0]) * (box[3] - box[1])
    if area > MAX_SUBTITLE_AREA_SHARE * prs.slide_width * prs.slide_height:
        return False
    return len(text) <= MAX_SUBTITLE_CHARS


# A shape sharing this much of its own HEIGHT with a body shape is standing in
# the same row as that shape. A quarter is deliberately low: the whole point is
# a row whose members do not line up exactly, and two labels set 0.25in apart on
# the client's table shared only a third of their height.
ROW_SHARE = 0.25


def _row_bound(box, body_boxes) -> bool:
    """Whether a header-band shape lines up with CONTENT rather than with the
    page's header.

    Vertical overlap only, and no test of horizontal distance. A row's members
    are side by side by definition, so the gap between them says nothing, and
    requiring adjacency would have failed on the very case this is for: the
    client's two column headings sit 0.9in and 2.6in from the left with the
    numbers they label starting at 4.3in."""
    height = box[3] - box[1]
    if height <= 0:
        return False
    return any(min(box[3], other[3]) - max(box[1], other[1])
               >= ROW_SHARE * height for other in body_boxes)


def _band_candidates(slide, prs):
    """(candidates, everything) for one slide: the header-band text that COULD be
    a stray, and every free box on the slide that could vouch for it."""
    boxes = [(s, _box(s)) for s in slide.shapes
             if not getattr(s, "is_placeholder", False)]
    boxes = [(s, b) for s, b in boxes if b is not None]
    by_type = _placeholders_by_type(slide)
    title = next(iter(sum((by_type.get(t, []) for t in _TITLE_TYPES), [])), None)
    sub = next(iter(sum((by_type.get(t, []) for t in _SUBTITLE_TYPES), [])), None)
    floor = _header_floor(title, sub)
    return ([(s, b) for s, b in boxes
             if _header_sized(b, _text_of(s), floor, prs)], boxes)


def stray_texts(prs) -> set:
    """The header-band texts this deck treats as STRAYS, decided once for the
    whole deck. Anything else up there is content that happens to sit high.

    Three things make this a deck-wide question rather than a per-slide one, and
    each of them was a wrong answer first.

    A ROW HAS NO BODY TO LINE UP WITH when the whole row sits above the line. A
    Gantt's eighteen month numbers all ended at 1.27in against a body beginning
    at 1.90in, so asking each one whether it lines up with the BODY answered no
    eighteen times and the sweep took the lot. They line up with each OTHER,
    which is what a row is, so mates are looked for among everything on the
    slide - not only among what the block will move.

    BOILERPLATE IS BOILERPLATE WHEREVER IT LANDS. "To be translated" sat alone
    at the top of most slides and, on the seven where it happened to share a band
    with a numbered badge, looked exactly like a row member. Judged per slide it
    came off five slides and stayed on seven, which is the worst of both
    outcomes. Judged by TEXT across the deck it goes from all of them or none.

    AND A STRAY CANNOT VOUCH FOR ANOTHER STRAY. Two working notes stamped side by
    side at the top of a slide each made the other look like a row, so both
    survived until the set was iterated to a fixpoint.

    On the client's deck (23/08/2026) this removes 25 pieces of text - every one
    of the 21 "To be translated" stamps, three "not comprehensive" notes and one
    orphaned column heading - and keeps the rest, including all eighteen month
    labels and every numbered badge and card heading. Removing everything above
    the line took 60 and gutted the tables.

    Known gap, and it errs the safe way: two strays that appear ONLY side by side
    and nowhere else vouch for each other forever, and both survive. The fixpoint
    can add to the set but never break a mutual cycle. Keeping two notes is a far
    better failure than a member count that would delete a two-label row, so it
    is left alone until a deck actually shows the problem."""
    slides = [_band_candidates(slide, prs) for slide in prs.slides]
    strays: set = set()
    # Bounded: each pass can only ADD to the set, and there are finitely many
    # texts, so this settles. The cap is a backstop, not a strategy.
    for _pass in range(8):
        grew = False
        for candidates, everything in slides:
            for shape, box in candidates:
                text = _text_of(shape)
                if text in strays:
                    continue
                mates = [b for other, b in everything
                         if other is not shape and _text_of(other) not in strays]
                if not _row_bound(box, mates):
                    strays.add(text)
                    grew = True
        if not grew:
            break
    return strays


def _rank_header_text(candidates, header_floor: int, slide, prs) -> list:
    """Header text ordered biggest type first: [title, subtitle, ...].

    The rule is deliberately blunt because that is how a designer reads a
    slide: the largest type is the heading, the next largest is the
    subheading. Ranking once and taking the first two in order beats searching
    per placeholder, which let the two searches disagree about the same shape.

    Ties break by vertical position, higher wins. An eyebrow set at the same
    size as its heading is a real pattern, and on a tie the thing nearer the
    top of the slide is the heading.

    Restricted to the HEADER band rather than the whole slide. Taking the
    globally largest text would let a big figure in a chart become the title,
    and an earlier unbounded version consumed a full-slide diagram into the
    subtitle placeholder and deleted it. The bounds below are what keep a
    panel from being mistaken for a line of text."""
    sw, sh = prs.slide_width, prs.slide_height
    scored = []
    for shape in candidates:
        box = _box(shape)
        if box is None:
            continue
        if not _header_sized(box, _text_of(shape), header_floor, prs):
            continue
        size = _max_font_pt(shape, slide, prs)
        # Box height stands in when the cascade resolves no size, so a set of
        # all-inherited candidates does not collapse to a single rank.
        scale = size or (box[3] - box[1]) / 12700.0
        scored.append((scale, box[1], shape))
    scored.sort(key=lambda s: (-s[0], s[1]))
    return [s[2] for s in scored]


def _header_zone_pick(candidates, header_floor: int, slide, prs):
    """The largest-type text that genuinely LIVES in the header band.

    Containment is tested on the whole box, not just its top edge. A tall
    content block starting high on the slide has its top in the header band
    while occupying most of the canvas, and it is not a subtitle."""
    sw, sh = prs.slide_width, prs.slide_height
    scored = []
    for shape in candidates:
        b = _box(shape)
        if b is None:
            continue
        if b[3] > header_floor + CONTENT_GAP_EMU:
            continue  # extends below the header band: not header content
        area = (b[2] - b[0]) * (b[3] - b[1])
        if area > MAX_SUBTITLE_AREA_SHARE * sw * sh:
            continue
        if len(_text_of(shape)) > MAX_SUBTITLE_CHARS:
            continue
        scored.append((_max_font_pt(shape, slide, prs),
                       (b[3] - b[1]), shape))
    if not scored:
        return None
    scored.sort(key=lambda s: (s[0], s[1]), reverse=True)
    return scored[0][2]


def _ph_label(ph) -> str:
    """A placeholder's role in the words a designer uses, not the enum repr."""
    try:
        kind = ph.placeholder_format.type
    except Exception:
        return "placeholder"
    return {PP_PLACEHOLDER.TITLE: "title",
            PP_PLACEHOLDER.CENTER_TITLE: "title",
            PP_PLACEHOLDER.SUBTITLE: "subtitle",
            PP_PLACEHOLDER.BODY: "body"}.get(kind, str(kind).split()[0].lower())


def _drop_background_override(slide, slide_index) -> list[ContentChange]:
    """Delete the slide's own p:bg so the master's background shows through.

    A slide-level background beats the layout's and the master's, and exported
    decks routinely stamp an explicit white fill on every slide. That is why a
    deck can adopt every other part of a master and still come out on the wrong
    ground colour.

    This is not a colour decision, which is why it belongs here and text colour
    does not: nothing picks a new value, an override is removed so the master's
    own declaration wins. Only done when the layout or master actually declares
    a background, so a slide is never stripped down to nothing."""
    cSld = slide._element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}cSld")
    if cSld is None:
        return []
    own = cSld.findall(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}bg")
    if not own:
        return []

    inherits = False
    for source in (getattr(slide, "slide_layout", None),
                   getattr(getattr(slide, "slide_layout", None),
                           "slide_master", None)):
        if source is None:
            continue
        if find(source._element, "p:cSld/p:bg") is not None:
            inherits = True
            break
    if not inherits:
        return []

    from lxml import etree

    kept = [etree.tostring(bg, encoding="unicode") for bg in own]
    for bg in own:
        cSld.remove(bg)
    return [ContentChange(
        slide_index, "dropped background override",
        "the slide carried its own background, which beats the master's; "
        "removed so the master's background applies",
        undo=[{"op": "bg", "xml": x} for x in kept if x])]


def _remove_duplicates(slide_index, free, filled_phs) -> list[ContentChange]:
    """Drop free shapes repeating text that just went into a placeholder.

    Runs BEFORE the content block moves. Ordering matters: the block move
    relocates the leftover copy out of the placeholder's box, and a check that
    runs afterwards sees no overlap and leaves the deck with the same line
    printed twice, once styled by the master and once not."""
    changes = []
    for ph in filled_phs:
        ph_box, ph_text = _box(ph), _text_of(ph).casefold()
        if ph_box is None or not ph_text:
            continue
        for shape in list(free):
            if _text_of(shape).casefold() != ph_text:
                continue
            if _overlap_share(_box(shape), ph_box) < MIN_COLLISION_SHARE:
                continue
            changes.append(ContentChange(
                slide_index, "removed duplicated text",
                f"{_text_of(shape)[:40]!r} was a second copy of the text now "
                f"in the {_ph_label(ph)} placeholder",
                undo=_insert_undo(shape)))
            free.remove(shape)
            _delete(shape)
    return changes


def _resolve_collisions(slide_index, free, filled_phs, prs,
                        was_clear=None) -> list[ContentChange]:
    """Clear anything this pass would otherwise leave printing on top of a
    placeholder it just filled.

    This exists because the pass CAUSES these collisions: it puts text into a
    placeholder and can leave a sibling shape (an eyebrow, an old heading) in
    the same box. Handing that to a later audit stage would make a designer
    clean up after the tool, and the audit would have to re-derive with
    heuristics what this function knows for certain, having just done it.

    `was_clear` is the set of shape ids that did NOT overlap the header band
    BEFORE this pass moved anything, and it is what holds that sentence to its
    word. A shape already standing in the band arrived that way: this pass did
    not cause it, and clearing it means moving one member of a row on its own -
    a Gantt's twenty-one month labels were pushed 0.93in clear of the table they
    label (design lead, 23/08/2026). Those are reported instead. Omit it and
    every collision is treated as ours, which is the old behaviour.

    Colour is deliberately NOT touched here. Text may now be unreadable on a
    new dark background, but choosing a colour is a design judgment and the
    audit owns contrast with real WCAG thresholds."""
    changes: list[ContentChange] = []
    boxes = [(ph, _box(ph)) for ph in filled_phs]
    boxes = [(ph, b) for ph, b in boxes if b is not None]

    # Shape in the OUTER loop on purpose. Placeholder-first double-counted a
    # shape that overlapped both the title and the subtitle, moving it once per
    # placeholder and logging two entries for one nudge. One shape, one
    # decision, one line in the report.
    for shape in list(free):
        box = _box(shape)
        if box is None:
            continue
        hits = [(ph, b) for ph, b in boxes
                if _overlap_share(box, b) >= MIN_COLLISION_SHARE]
        if not hits:
            continue

        label = _text_of(shape)[:40] or "a shape"
        names = ", ".join(sorted({_ph_label(ph) for ph, _b in hits}))

        if was_clear is not None and str(shape.shape_id) not in was_clear:
            changes.append(ContentChange(
                slide_index, "overlap needs a designer",
                f"{label!r} was already standing in the header band before this "
                f"pass moved anything, so the overlap with the {names} "
                f"placeholder is not this pass's doing. Nothing was moved: "
                f"clearing it would move one member of its row on its own. "
                f"Check it"))
            continue
        # Clear the whole header band, not just the placeholders currently hit.
        # Shifting only past the title dropped the eyebrow straight onto the
        # subtitle, trading one collision for another; header placeholders span
        # the slide width, so the band is the thing to get below.
        band_bottom = max(b[3] for _ph, b in boxes)
        shift = band_bottom + CONTENT_GAP_EMU - box[1]

        # Only TEXT is nudged. This function exists because the pass puts text
        # into a placeholder and can leave a sibling line of text printing on top
        # of it; a graphic has no such problem, and the same reasoning that keeps
        # a corner rule, a bracket or a mark out of the remnant sweep
        # (_header_sized) keeps it out of this move. It is part of a composition,
        # and a header placeholder spans the whole slide width, so a box overlap
        # with one says little about where its words actually are.
        #
        # Nudging graphics broke the arrangement it was meant to protect. A
        # 1.70x0.02in decorative bar sitting at 0.37in was pushed 1.53in down
        # into the body on three slides, and two full-width table rules were
        # pushed 0.63in and 0.40in - different distances, so they collapsed onto
        # each other and onto the row beneath (design lead, 23/08/2026). One
        # block, one move; anything else here is a report.
        if not _text_of(shape):
            changes.append(ContentChange(
                slide_index, "overlap needs a designer",
                f"{label!r} overlaps the {names} placeholder. Nothing was moved: "
                f"it carries no text, so it is part of a composition rather than "
                f"a line competing with the heading, and moving it alone would "
                f"break whatever it was drawn with. A header placeholder spans "
                f"the whole slide width, so the two may not touch on screen at "
                f"all. Check it"))
            continue

        if shift > 0 and box[3] + shift <= prs.slide_height:
            back = _offset_undo([shape])
            shape.top += shift
            changes.append(ContentChange(
                slide_index, "nudged clear of a placeholder",
                f"{label!r} moved down {shift / 914400:.2f}in so it no longer "
                f"prints over the {names} placeholder",
                undo=back))
        else:
            changes.append(ContentChange(
                slide_index, "overlap needs a designer",
                f"{label!r} overlaps the {names} placeholder and cannot clear "
                f"it without running off the canvas"))
    return changes


# Two text boxes sharing this much of the smaller one's area are printing on
# top of each other, not merely adjacent.
MIN_TEXT_OVERLAP_SHARE = 0.30


def _report_text_overlaps(slide_index, free) -> list[ContentChange]:
    """Flag text still printing over other text among the content.

    Reported, never moved. Placeholder collisions are this pass's own doing and
    it fixes them; content-on-content overlap is a layout judgment about the
    designer's own arrangement, and guessing which of two blocks should give
    way is how a tool wrecks a deliberate composition. The audit's alignment
    module owns that call, with the tolerances calibrated for it."""
    text_shapes = [(s, _box(s)) for s in free if _text_of(s)]
    text_shapes = [(s, b) for s, b in text_shapes if b is not None]
    changes = []
    seen = set()
    for i, (shape_a, box_a) in enumerate(text_shapes):
        for shape_b, box_b in text_shapes[i + 1:]:
            share = max(_overlap_share(box_a, box_b), _overlap_share(box_b, box_a))
            if share < MIN_TEXT_OVERLAP_SHARE:
                continue
            key = tuple(sorted((shape_a.shape_id, shape_b.shape_id)))
            if key in seen:
                continue
            seen.add(key)
            changes.append(ContentChange(
                slide_index, "text overlaps text",
                f"{_text_of(shape_a)[:28]!r} and {_text_of(shape_b)[:28]!r} "
                f"share {share * 100:.0f}% of their area; left as they are, "
                f"since choosing which one moves is a layout decision"))
    return changes


def _margin_frame(slide, prs):
    """(left, top, right, bottom, source, body_top) EDGES in slide coordinates
    from the master's own DRAWING GUIDES when the designer drew them, else from
    its placeholder extents. None for any side the master does not state;
    `source` is "guides", "placeholders" or None, and callers that need a STATED
    frame rather than an inferred one check it.

    `body_top` is the sixth edge, and the one a header band makes necessary: the
    guide the master draws to say where CONTENT begins, which is not the top
    margin (that is where the page begins) and not the header placeholder's
    floor (that is where one layout's title box happens to end). None unless
    the master drew it. See qc.stylespec.read_content_band.

    A presentation-space rectangle wins, then guides, then placeholder extents
    (qc.stylespec.infer_grid ranks them and says which it used). The first two
    are stated intentions - one drawn and named, one set as guides - while
    placeholder extents are an inference from where the title happens to sit,
    which left body content indented a few millimetres inside the title on a
    real deck.

    All FOUR sides are read, not just left and right. A master that states a
    bottom margin was stating it about content too, and reading only half the
    frame let a block sit below a line the master had drawn."""
    try:
        from .stylespec import infer_grid

        master = slide.slide_layout.slide_master
        grid = infer_grid(prs, master) or {}
    except Exception:
        return (None, None, None, None, None, None)
    margins = grid.get("margins_emu") or {}
    if margins.get("left") is None or margins.get("right") is None:
        return (None, None, None, None, None, None)
    top, bottom = margins.get("top"), margins.get("bottom")
    return (margins["left"],
            None if top is None else top,
            prs.slide_width - margins["right"],
            None if bottom is None else prs.slide_height - bottom,
            grid.get("source"),
            grid.get("body_top_emu"))


def _content_region(prs, header_phs, footer_phs, frame):
    """Where free content may live, on all four sides: inside the master's
    margin frame, clear of the header placeholders above and the footer band
    below. Each side takes the TIGHTER of the two statements, because both are
    real - guides say where the page's content area is, a header placeholder
    says what already occupies the top of it.

    Only a frame the designer DREW contributes its top and bottom. Left and
    right fall back to placeholder extents because a master's content
    placeholders do bracket the page horizontally, but their vertical extent
    describes one layout's boxes, not the page: on a 4:3 template widened to
    16:9 the inferred bottom lands well above the real one, and a block clamped
    to it stops short of where it was meant to go.

    The top prefers the master's BODY CEILING over its top margin when the
    master drew one. They are different lines and only one of them is about
    content: seeding from the top margin left the header placeholder's floor as
    the only thing holding the body down, so a layout with no subtitle
    placeholder started its body at the page's top margin, inside the strip the
    master reserves under the header."""
    sw, sh = prs.slide_width, prs.slide_height
    left, m_top, right, m_bottom, source, body_top = frame
    stated = source in STATED_FRAME_SOURCES
    if left is None:
        boxes = [b for b in (_box(p) for p in header_phs + footer_phs) if b]
        if boxes:
            left = min(b[0] for b in boxes)
            right = max(b[2] for b in boxes)
        else:
            left, right = 0, sw

    top = m_top if (stated and m_top is not None) else 0
    if stated and body_top is not None:
        top = body_top
    for ph in header_phs:
        b = _box(ph)
        if b:
            top = max(top, b[3] + CONTENT_GAP_EMU)
    bottom = m_bottom if (stated and m_bottom is not None) else sh
    for ph in footer_phs:
        b = _box(ph)
        if b and b[1] > sh * 0.5:  # only furniture genuinely low on the slide
            bottom = min(bottom, b[1] - CONTENT_GAP_EMU)
    return (left, top, right, bottom)


def _heading_margin_notes(slide_index, slide, prs, frame) -> list[ContentChange]:
    """Report a heading whose BOX ends up outside the master's margin frame.

    Reported and never corrected, on any side. Whether a title or standfirst
    may break the margin is the client's house style, not a defect, so this
    pass states the fact and leaves the geometry exactly as the master put it
    (design lead, 19/08/2026). Moving the heading would also be the one move
    that undoes the master this pass has just applied.

    The BOX is what is measured, never the text inside it. A long line spilling
    out of a correctly placed box is a copy-length problem for the designer and
    the client to settle; it is not this pass's to detect, and guessing at glyph
    widths would make it a rendering estimate rather than a fact about the
    deck.

    Only a frame the master STATES counts. Placeholder-derived margins are an
    inference from where the master's own title happens to sit, so measuring a
    heading against them compares a heading to a heading: on a master whose
    placeholders no longer span its slide (a 4:3 template widened to 16:9) every
    slide reports a title breaking a margin nobody ever drew.

    The frame's top is the PAGE's top margin, never a body ceiling: a
    presentation space drawn around the body only states where the body begins,
    and qc.stylespec.infer_grid separates the two before this ever sees them."""
    left, top, right, bottom, source, _body_top = frame
    if left is None or source not in STATED_FRAME_SOURCES:
        return []
    heads = heading_ids(slide, prs)
    changes = []
    for shape in slide.shapes:
        if str(shape.shape_id) not in heads:
            continue
        box = _box(shape)
        if box is None or not _text_of(shape):
            continue
        breaches = [(side, over) for side, over in (
            ("left", left - box[0]),
            ("top", 0 if top is None else top - box[1]),
            ("right", box[2] - right),
            ("bottom", 0 if bottom is None else box[3] - bottom),
        ) if over > HEADING_SLACK_EMU]
        if not breaches:
            continue
        sides = ", ".join(side for side, _over in breaches)
        worst = max(over for _side, over in breaches)
        changes.append(ContentChange(
            slide_index, "heading past the margin",
            f"{_text_of(shape)[:40]!r} sits {worst / 914400:.2f}in past the "
            f"{sides} margin. Nothing was moved or resized: whether a heading "
            f"may break the margin is the client's call, so ask them before "
            f"changing it",
            severity="alert"))
    return changes


def _master_footer_text(slide) -> str | None:
    """The footer string the master actually stamps, if any."""
    for source in (getattr(slide, "slide_layout", None),
                   getattr(getattr(slide, "slide_layout", None),
                           "slide_master", None)):
        if source is None:
            continue
        try:
            for ph in source.placeholders:
                if ph.placeholder_format.type == PP_PLACEHOLDER.FOOTER:
                    text = _text_of(ph)
                    if text:
                        return text
        except Exception:
            continue
    return None


def _is_page_furniture(shape, slide_number: int, sh_height: int,
                       master_footer: str | None) -> str | None:
    """A free text box DUPLICATING furniture the master now provides.

    Deliberately narrow. An earlier version removed any short text low on the
    slide, which deleted source lines and attributions ("Source: internal
    analysis 2026") because they happened to sit in the footer band. Losing a
    designer's content is far worse than leaving a duplicate footer, so only
    two things qualify: a bare page number matching this slide, and text that
    matches the master's own footer string."""
    box = _box(shape)
    if box is None or box[1] < FURNITURE_BAND * sh_height:
        return None
    text = _text_of(shape)
    if not text:
        return None
    if text.isdigit() and (text.lstrip("0") or "0") == str(slide_number):
        return "page number"
    if master_footer and text.casefold() == master_footer.casefold():
        return "footer text"
    return None


def migrate_slide(slide, slide_index: int, prs, anchors=None,
                  strays=None) -> list[ContentChange]:
    """`anchors` is qc.util.recurring_anchors and `strays` is stray_texts, both
    read over the WHOLE deck: one tells page furniture from content low on a
    single slide, the other tells a boilerplate note from content sitting high.
    Neither question can be answered from one slide, and both are passed in
    because computing them per slide would make the pass quadratic. Computed here
    when a caller has only one slide in mind."""
    if anchors is None:
        anchors = recurring_anchors(prs)
    if strays is None:
        strays = stray_texts(prs)
    changes: list[ContentChange] = _drop_background_override(slide, slide_index)
    by_type = _placeholders_by_type(slide)
    title_ph = next(iter(sum((by_type.get(t, []) for t in _TITLE_TYPES), [])), None)
    sub_ph = next(iter(sum((by_type.get(t, []) for t in _SUBTITLE_TYPES), [])), None)
    furniture_phs = sum((by_type.get(t, []) for t in _FURNITURE_TYPES), [])

    # Recorded before anything moves: who was clear of the header band to begin
    # with. Read afterwards it is unanswerable, and it is the difference between
    # a collision this pass caused and one it inherited
    # (_resolve_collisions).
    band_floor_now = _header_floor(title_ph, sub_ph)
    was_clear = {str(s.shape_id) for s in slide.shapes
                 if not getattr(s, "is_placeholder", False)
                 and (_box(s) or (0, band_floor_now + 1, 0, 0))[1] >= band_floor_now}

    free = [s for s in slide.shapes if not getattr(s, "is_placeholder", False)]

    # --- 1. page furniture the master now supplies -----------------------
    if furniture_phs or _layout_supplies(slide, _FURNITURE_TYPES):
        master_footer = _master_footer_text(slide)
        for shape in list(free):
            what = _is_page_furniture(shape, slide_index + 1, prs.slide_height,
                                      master_footer)
            if what:
                changes.append(ContentChange(
                    slide_index, "removed duplicate furniture",
                    f"{what} {_text_of(shape)!r}; the master supplies this",
                    undo=_insert_undo(shape)))
                free.remove(shape)
                _delete(shape)

    # --- 2. title and subtitle into their placeholders -------------------
    # Title candidates are searched at every depth, not just the top level. A
    # designer who grouped the eyebrow with the heading, or a converter that
    # wrapped the header in a group, would otherwise hide the title from this
    # pass entirely: a group carries no text of its own, so the shape that
    # actually holds "Full Lifecycle Management" is never even a candidate.
    # The block move still works on TOP-LEVEL shapes, so pulling one text
    # shape out of a group leaves the rest of that group intact and moving
    # together.
    text_shapes = [s for s, _path in iter_shapes_deep(slide.shapes)
                   if not getattr(s, "is_placeholder", False) and _text_of(s)]
    header_floor = _header_floor(title_ph, sub_ph)
    ranked = _rank_header_text(text_shapes, header_floor, slide, prs)

    placed: set[str] = set()
    for ph, label in ((title_ph, "title"), (sub_ph, "subtitle")):
        if ph is None or _text_of(ph):
            continue
        # Largest type is the title, second largest the subtitle. Ranking once
        # over the whole header and taking the top two in order is what makes
        # that a single rule rather than two competing per-placeholder searches.
        # A duplicated heading is skipped rather than printed twice under two
        # different master styles; the collision pass removes the leftover.
        pick = next((s for s in ranked
                     if _text_of(s).casefold() not in placed), None)
        if pick is None:
            continue
        ranked.remove(pick)
        text = _text_of(pick)
        # Both halves of the move, captured before either happens: the
        # placeholder as it was (empty, and styled by the master) and the shape
        # the text came out of. Undoing one without the other would either lose
        # the wording or print it twice.
        was_empty = _shape_xml(ph)
        came_from = _shape_xml(pick)
        rtl_marked = _set_text(ph, text, source=pick)
        grouped = pick not in free
        back = [op for op in (
            {"op": "replace", "shape_id": str(ph.shape_id), "xml": was_empty}
            if was_empty else None,
            {"op": "insert", "xml": came_from} if came_from else None,
        ) if op]
        changes.append(ContentChange(
            slide_index, f"{label} into placeholder",
            f"{text[:60]!r} now styled by the master; run-level formatting "
            f"dropped" + (" (lifted out of a group)" if grouped else "")
            + ("; kept right-to-left, which the master's own paragraph "
               "direction would have reversed" if rtl_marked else ""),
            undo=back or None))
        placed.add(text.casefold())
        if not grouped:
            free.remove(pick)
        text_shapes.remove(pick)
        _delete(pick)

    # --- 2b. text duplicated by the placeholder fill ---------------------
    filled_now = [ph for ph in (title_ph, sub_ph)
                  if ph is not None and _text_of(ph)]
    if filled_now:
        changes.extend(_remove_duplicates(slide_index, free, filled_now))

    # --- 3. remaining content into the master's content region -----------
    header_phs = [p for p in (title_ph, sub_ph) if p is not None]
    # Read once and passed on: the frame is what both the block move and the
    # heading report measure against, and re-deriving it per caller is how the
    # two would end up disagreeing about where the margins are.
    frame = _margin_frame(slide, prs)
    region = _content_region(prs, header_phs, furniture_phs, frame)

    # Shapes already living in the furniture band are page furniture the deck
    # keeps by choice (a source line, a logo strip). They are NOT part of the
    # content block: dragging them down with it pushed them clean off the
    # canvas in testing, which is worse than leaving them where they were.
    #
    # Page-deep panels are parked for the same reason and it matters more now
    # that the body ceiling binds: a 94%-height image anchored to the bottom
    # edge is a background, not body content, and with it in the block the
    # ceiling move pushed an entire slide off the canvas (real-master check,
    # 20/08/2026). Parked, it stays exactly where the designer put it while the
    # body seats itself on the guide.
    # The bottom strip alone is not enough to call something furniture, and on
    # its own it parked 25 shapes on this deck of which NONE were furniture: the
    # bottom row of a table at 6.61in, a chart's axis labels at 6.80in, a legend
    # at 7.09in, a band of arrows at 6.65in. Each stayed behind while the
    # composition it belongs to moved (design lead, 23/08/2026). Real page
    # furniture on these decks is a PLACEHOLDER the master supplies, so it never
    # reaches this test at all.
    #
    # So both halves are required, exactly as qc.fixer._pinned_furniture already
    # asks it: low on the slide AND recurring at one position across the deck. A
    # source line low on one slide is content; a logo strip on twenty slides is
    # not. The risk this trades against is real - content dragged past the bottom
    # edge - and it is reported as an overflow rather than hidden by leaving the
    # shape behind.
    band_top = FURNITURE_BAND * prs.slide_height

    def _pinned(shape) -> bool:
        box = _box(shape)
        if box is None:
            return False
        if full_height_panel(box[1], box[3] - box[1], prs.slide_height):
            return True
        return (box[1] >= band_top
                and (slide_index, str(shape.shape_id)) in anchors)

    movable = [s for s in free if not _pinned(s)]
    parked = len(free) - len(movable)

    # Header remnants are separated from body content by where each shape
    # STARTED, and they get their own placement. This is the difference between
    # respecting a margin and preserving a distance: with everything in one
    # block, a leftover eyebrow sitting high on the slide became the block's
    # top edge, so aligning that edge to the body margin pushed the real
    # content down by however far the EYEBROW had to travel. The cards landed
    # 1.5in below where the master says the body begins. Body content is now
    # placed against the master's body margin regardless of what sits above it.
    # RIDERS: shapes that cannot be MEASURED but must still travel. A perfectly
    # horizontal or vertical connector has a width or a height of exactly zero,
    # so _box refuses it - correctly, since it has no area to overlap, contain or
    # anchor anything with - and dropping it here meant it never moved at all.
    # A deck of tables came back with every rule and divider stranded where it
    # was while the rows it separates moved beneath it (design lead, 23/08/2026:
    # "boxes and lines that are supposed to be moved as one entity should stay
    # that way"). It is one line of code and it was silent on 14 of 26 slides.
    #
    # They ride and nothing more: no say in where the block starts, no part in
    # the remnant sweep, no vote on the bleed. A line is not content to be
    # judged, it is content to be carried.
    remnants, body, riders = [], [], []
    for shape in movable:
        box = _box(shape)
        if box is None:
            if shape.left is not None and shape.top is not None:
                riders.append(shape)
            continue
        # Two gates, and both are needed. _header_sized asks whether this could
        # be a line of header text at all - a diagram whose top sits just above
        # the floor is body content, and an earlier top-edge-only test deleted
        # one. `strays` then asks whether the DECK treats this text as a stray,
        # which is the question a single slide cannot answer (stray_texts).
        if (_header_sized(box, _text_of(shape), header_floor, prs)
                and _text_of(shape) in strays):
            remnants.append((shape, box))
        else:
            body.append((shape, box))

    # Second look, and the sweep below is why it has to happen: a shape that
    # lines up with CONTENT is content, whatever band its box happens to end in.
    #
    # A table's column headings are the case that proves it. On the client's
    # Gantt slide, "team members" and "months of work" ended at 1.84in and
    # 1.59in - just above the 1.90in header cutoff - while the month numbers
    # labelling the same row ended at 2.01in, just below it. Same row, same 14pt
    # type, drawn as one thing by the designer; the cutoff fell between them and
    # the two words were deleted while the numbers were kept (design lead,
    # 23/08/2026). The cutoff cannot be moved to fix that - wherever it sits,
    # some row straddles it - so the question has to be asked of the content
    # instead of the band.
    #
    # Keeping is the safe direction. A spared shape travels with the block and
    # gets nudged clear of the master's header if it still collides; a swept one
    # is gone, and a designer has to notice and put it back.
    # Whether this pass PLACED anything is deliberately not consulted. Gating
    # the sweep on it was tried and withdrawn (design lead, 23/08/2026): a slide
    # whose placeholders PowerPoint had already filled then kept every stray in
    # the band, and the strays travelled into the body with the block - so a
    # working note reading "To be translated" ended up inside the content area
    # instead of out of the deck. Text above the line where the body begins is
    # unplaced whatever put the heading there, and the client's instruction is
    # plain: remove it, flag it, let me put it back.
    #
    # The cost of that instruction, stated because it is real: on a slide where
    # PowerPoint widened the subtitle placeholder to the layout's full width, the
    # heading it absorbed now prints at the far margin, so a stray restored to
    # its own coordinates can land under it. That is the layout assignment's
    # doing, not the sweep's, and the restore says what it now covers.
    #
    # Remnants stack directly under the header band; the body starts below
    # them. A missing eyebrow slot in the master is why they need somewhere to
    # go at all, and stacking keeps reading order without coupling the body's
    # position to the remnant's original one.
    boxes = [b for _s, b in body]
    # What the block is MEASURED from, which is not the same as what moves:
    # every body shape below travels, and only these say where it starts
    # (_anchor_boxes).
    anchors = _anchor_boxes(boxes)
    body_top_target = region[1]

    # Header furniture does not get to be the body's top edge either. A rule, a
    # bracket or a corner mark drawn above the line where content begins is part
    # of the header's composition and carries no text, so the remnant sweep
    # leaves it alone - correctly, since it is not "unplaced text" - but leaving
    # it in the ANCHOR set made a 0.02in decorative bar sitting at 0.37in the top
    # of the block on the client's deck, and seating THAT on the body line pushed
    # the real content 1.53in down and off the page.
    #
    # Dropped only while something else still reaches the body region. A deck
    # whose whole body sits above the line has to be brought DOWN to it, and
    # excluding every anchor there would leave exactly those slides untouched.
    if body_top_target:
        reaching = [b for b in anchors if b[3] > body_top_target]
        if reaching:
            anchors = reaching

    # Nor does anything down in the bottom strip. It still TRAVELS - that is what
    # keeps a legend with its chart and a table's last row with its table - but a
    # footnote, an axis label or a legend low on the page is not where the body
    # begins or ends. Leaving them in made a slide whose only free content is one
    # footer bar at 7.17in seat that bar on the body line, hoisting it 5.27in to
    # the top of an otherwise empty page.
    #
    # With nothing above the strip there is nothing to seat, and the block is
    # left alone rather than invented a position for.
    anchors = [b for b in anchors if b[1] < band_top]

    # Header text that reached no placeholder is unaccounted for: the master
    # defines no slot for it, and leaving it floating is what produced the
    # overlapping headers and the dead space. It is REMOVED and reported as an
    # alert carrying its full text AND its own XML, so a designer can put back
    # anything that mattered exactly as it was. Removing content silently would
    # be indefensible; keeping it and hoping is what the earlier stacking did,
    # and it collided on full slides.
    for shape, box in remnants:
        text = _text_of(shape)
        where = (f"ended at {box[3] / 914400:.2f}in, above the "
                 f"{body_top_target / 914400:.2f}in where the body begins"
                 if body_top_target else "sat above the body")
        changes.append(ContentChange(
            slide_index, "removed unplaced text",
            f"{text!r} {where}, and the master has no placeholder for it. "
            f"Removed rather than carried into the content: put it back with "
            f"Undo if it belongs in the deck",
            severity="alert", removed_text=text,
            removed_xml=_shape_xml(shape),
            restore_id=f"{slide_index}-{shape.shape_id}",
            undo=_insert_undo(shape)))
        if shape in free:
            free.remove(shape)
        _delete(shape)
    if anchors:
        cl = min(b[0] for b in anchors)
        ct = min(b[1] for b in anchors)
        cr = max(b[2] for b in anchors)
        cb = max(b[3] for b in anchors)

        # All four sides are consulted, and they divide the way a translate
        # forces them to: TOP and LEFT bind (the block's edge is set to the
        # margin), BOTTOM and RIGHT clamp and report. A move cannot seat a
        # block against opposite margins at once - only a resize could, and
        # resizing a text box reflows its text.
        #
        # Signed on purpose: content sitting BELOW the body margin is pulled up
        # to meet it. Only ever moving down could close no gap, which is the
        # whole complaint. Bound only when the master actually STATES where the
        # body begins - a header placeholder's floor or a top guide. With
        # neither, region[1] is zero and a "correction" would yank content to
        # the very top edge of the slide.
        dy = (body_top_target - ct) if region[1] > 0 else 0
        # Where the master DREW its body ceiling, that line binds absolutely:
        # the strip it reserves under the header stays empty on every slide,
        # including the over-full ones. The clamp below used to apply here too,
        # and because it is bounded by the bottom margin it collapsed to zero on
        # exactly the slides that overflow - so the deck came back with the top
        # of the block sitting in the reserved strip on every busy slide, which
        # is the complaint this answers (design decision, 20/08/2026: an honest
        # alert about the bottom beats a silently broken header band).
        #
        # An already-overflowing block may still be pulled UP (that reduces the
        # overflow) but is never pushed further DOWN. max(0, ...) is what makes
        # the clamp one-directional; a plain min() would drag content upward
        # past its margin and back into the header.
        #
        # What decides is whether the master STATES where content begins, not
        # whether it happens to state it as a guide PAIR. Tying the rule to the
        # header band was a bug with a deck-wide symptom (design lead,
        # 21/08/2026): a master stating its frame with a presentation-space
        # rectangle has no band, so the clamp applied, so every slide whose
        # content is too tall kept the position it arrived with - and the deck
        # read as "not following the presentation space" from end to end.
        stated = frame[4] in STATED_FRAME_SOURCES
        if not stated:
            dy = min(dy, max(0, region[3] - cb))
        # Not even the CANVAS outranks a stated frame. Holding a too-tall block
        # back to keep it on the page was tried and refused: the frame's top line
        # binds on every slide or the deck loses the one line that makes its
        # headers read as one, and an overflow that is reported is better than a
        # header band that is quietly broken on the busy slides (design lead,
        # 23/08/2026, re-confirming the 21/08 decision against the alternative).
        # What the overflow costs is stated instead, in inches, below.
        # The START margin BINDS, the same way the body-top margin does: the
        # block's edge is set to it rather than merely kept inside it.
        # Correcting only breaches left content sitting a few millimetres
        # inside the title's own left edge, which reads as misalignment even
        # though no margin was crossed.
        #
        # Which margin is the start depends on the SCRIPT, not on the tool:
        # Arabic reads right to left, so the right margin is where its block
        # begins and binding the left edge leaves the reading edge ragged
        # (design lead, 20/08/2026). Judged per slide, because a bilingual deck
        # runs both (qc.util.slide_is_rtl).
        rtl = slide_is_rtl(slide)
        dx = (region[2] - cr) if rtl else (region[0] - cl)

        # Full-bleed content is exempt from the sideways move, and exempt ALONE.
        # A band or image deliberately running edge to edge is not indented
        # content and dragging it to the margin would destroy the effect - but it
        # is one shape's exemption, not the slide's. Vetoing the whole block on
        # its account left the deck's Arabic cover title sitting in the left half
        # of the page because a 1.28in white footer band happened to run edge to
        # edge underneath it (design lead, 23/08/2026, "the title should be right
        # aligned since it's Arabic"). It was: right-aligned inside a box the
        # width of half the slide, which is not what anyone means.
        #
        # So the block's SIDE edges are measured from the shapes the margin can
        # actually apply to, and the bleed elements keep their x while still
        # travelling in y with everything else. Vertically the block stays one
        # thing; horizontally a bleed element was never part of the arrangement,
        # which is the same reading that parks a page-deep panel above.
        # What "alone" covers: the bleed element AND whatever is sitting on it.
        # A logo stamped on a full-width footer band belongs to the band, not to
        # the text three inches above it, and carrying it along took the client's
        # MWAN mark from the bottom-left corner of its cover to the middle of the
        # strip. Anything mostly inside a bleed element keeps its x with it.
        #
        # The cost, stated: where a bleed element covers the content area itself
        # rather than a strip of it, everything is inside it and the slide gets no
        # sideways seat at all. That is what the old whole-block veto did on every
        # such slide, so it is no loss, and the move report names what was held.
        bleed_width = FULL_BLEED_SHARE * prs.slide_width
        bleed_boxes = [b for b in boxes if (b[2] - b[0]) >= bleed_width]

        def _bleeds(box) -> bool:
            if (box[2] - box[0]) >= bleed_width:
                return True
            return any(_overlap_share(box, b) >= ON_A_BLEED_SHARE
                       for b in bleed_boxes)

        indented = [b for b in anchors if not _bleeds(b)]
        if indented:
            il = min(b[0] for b in indented)
            ir = max(b[2] for b in indented)
            dx = (region[2] - ir) if rtl else (region[0] - il)
        else:
            # Nothing but bleed elements: there is no indented content to seat,
            # and moving them sideways is the one thing this must not do.
            il, ir, dx = cl, cr, 0
        too_wide = (ir - il) > (region[2] - region[0])

        # The MOVE is decided from the anchors; the OVERFLOW is measured from
        # everything that actually travels. They are different sets now that a
        # legend low on the page rides along without saying where the block ends,
        # and reporting only the anchors let that legend leave the page unsaid.
        low = max([b[3] for _s, b in body] or [cb])
        fits = ((low + dy) <= region[3]
                and (il + dx) >= region[0] and (ir + dx) <= region[2])
        movable = [s for s, _b in body]
        edge = "right" if rtl else "left"
        if dx or dy:
            travelling = [s for s, _b in body] + riders
            back = _offset_undo(travelling)
            held = 0
            for shape in travelling:
                box = _box(shape)
                if box is not None and _bleeds(box):
                    held += 1
                else:
                    shape.left += dx
                shape.top += dy
            # Names the frame it was seated on, because "was the presentation
            # space used on this slide?" is otherwise unanswerable from the
            # report, and a stale stored master looks identical to a bug.
            source = {"presentation_space": "the presentation space",
                      "guides": "the master's guides",
                      "placeholders": "the master's placeholder extents "
                                      "(no guides or presentation space to "
                                      "read)"}.get(frame[4], "the master")
            landed = ct + dy
            detail = (f"{len(travelling)} shape(s) shifted {dx / 914400:+.2f}in, "
                      f"{dy / 914400:+.2f}in onto {source}: body now starts at "
                      f"{landed / 914400:.2f}in"
                      + ("" if abs(landed - body_top_target) <= HEADING_SLACK_EMU
                         else f", which is {abs(landed - body_top_target) / 914400:.2f}in "
                              f"off the {body_top_target / 914400:.2f}in the frame states")
                      + f", seated on the {edge} margin"
                      + (" (Arabic reads right to left)" if rtl else ""))
            if held and dx:
                detail += (f"; {held} shape(s) running edge to edge, and what "
                           f"sits on them, kept their own left and right edges "
                           f"so the bleed survives; they moved down with the "
                           f"rest")
            if remnants:
                detail += (f"; {len(remnants)} unplaced header shape(s) removed "
                           f"and listed above")
            if parked:
                detail += (f"; {parked} shape(s) in the footer band or running "
                           f"the depth of the page left in place rather than "
                           f"dragged off the canvas")
            changes.append(ContentChange(
                slide_index, "content block moved", detail, undo=back))
        # Where the master states nothing, the block is only kept INSIDE an
        # inferred frame, never seated on it - so it can end up short of the
        # line, and the report has to say why rather than leave a designer
        # comparing the deck against a rectangle that was never read.
        landed = ct + dy
        if (not stated and body_top_target
                and abs(landed - body_top_target) > HEADING_SLACK_EMU):
            changes.append(ContentChange(
                slide_index, "body not seated on a stated frame",
                f"the body starts at {landed / 914400:.2f}in, not the "
                f"{body_top_target / 914400:.2f}in this pass aimed at. The "
                f"master states no content frame to bind to (read from "
                f"{frame[4] or 'nothing'}), so the block was kept inside the "
                f"inferred margins instead of seated on a line. Draw a "
                f"presentation-space rectangle on the master, or set guides, "
                f"and this becomes exact",
                severity="alert"))
        if too_wide:
            changes.append(ContentChange(
                slide_index, "wider than the margins",
                f"the content block is {(ir - il) / 914400:.2f}in wide against a "
                f"{(region[2] - region[0]) / 914400:.2f}in content area, so it "
                f"cannot sit inside both margins; aligned to the {edge} margin "
                f"and left for a designer (not scaled: narrowing a text box "
                f"reflows its text)",
                severity="alert"))
        if not fits:
            # Says WHERE it does not fit and by how much, because the answer
            # differs: past the bottom margin is a rework conversation, past the
            # slide edge is content that will not print at all.
            over_margin = (low + dy) - region[3]
            over_canvas = (low + dy) - prs.slide_height
            spill = []
            if over_margin > 0:
                spill.append(f"{over_margin / 914400:.2f}in past the bottom "
                             f"margin")
            if over_canvas > 0:
                spill.append(f"{over_canvas / 914400:.2f}in past the slide edge, "
                             f"where it will not print")
            where = ("; it now runs " + " and ".join(spill)) if spill else ""
            held = (" The line the master states for the top of its body was "
                    "held, so the strip under the header is still clear."
                    if stated and dy > 0 else "")
            # Measured, not characterised. "Taller or wider than the content
            # region" reads like a nudge is missing; "7.08in of content against
            # a 4.90in area" is a rework conversation, and it is the same fact
            # (design lead, 23/08/2026, on a slide whose block was half again
            # the height of the frame).
            size = (f"the block is {(low - ct) / 914400:.2f}in tall and "
                    f"{(ir - il) / 914400:.2f}in wide against a "
                    f"{(region[3] - region[1]) / 914400:.2f}in by "
                    f"{(region[2] - region[0]) / 914400:.2f}in content region")
            changes.append(ContentChange(
                slide_index, "content does not fit",
                f"{size}{where}. Left for a designer (not scaled: shrinking a "
                f"text box does not shrink its type).{held}",
                severity="alert" if over_canvas > 0 else "info"))

    # --- 3b. collisions this pass would otherwise ship -------------------
    # Runs AFTER the block move, so it only sees what the block move could not
    # resolve. A deck whose content already reaches the bottom edge cannot
    # shift down at all without leaving the canvas, so the header stays
    # occupied and these are the shapes still standing in it.
    filled = [ph for ph in (title_ph, sub_ph)
              if ph is not None and _text_of(ph)]
    if filled:
        changes.extend(_resolve_collisions(slide_index, free, filled, prs,
                                           was_clear))

    # --- 3c. text-on-text left among the content -------------------------
    changes.extend(_report_text_overlaps(slide_index, free))

    # --- 4. placeholders left empty --------------------------------------
    for ph in list(slide.placeholders):
        try:
            ph_type = ph.placeholder_format.type
        except Exception:
            continue
        if ph_type in _FURNITURE_TYPES:
            continue  # footer/slide-number placeholders are meant to be empty
        if not _text_of(ph):
            changes.append(ContentChange(
                slide_index, "removed empty placeholder",
                f"{ph_type} had nothing to hold; its prompt text would show "
                f"in the editor",
                undo=_insert_undo(ph)))
            _delete(ph)

    # --- 5. headings outside the margin frame ----------------------------
    # Last, so it describes the slide as it will actually ship rather than an
    # intermediate state. Reported only; see _heading_margin_notes.
    changes.extend(_heading_margin_notes(slide_index, slide, prs, frame))
    return changes


def migrate_deck(deck_bytes: bytes) -> tuple[bytes, list[ContentChange]]:
    """Run the content migration over every slide of an already-restyled deck."""
    prs = Presentation(io.BytesIO(deck_bytes))
    changes: list[ContentChange] = []
    # Read once for the deck: what recurs across slides is the only way to tell
    # a footer from a source line that happens to sit low on one slide.
    anchors = recurring_anchors(prs)
    # And which header-band texts this deck treats as strays, before a single
    # slide is touched: judged per slide, the same note came off some slides and
    # stayed on others (stray_texts).
    strays = stray_texts(prs)
    for idx, slide in enumerate(prs.slides):
        try:
            changes.extend(migrate_slide(slide, idx, prs, anchors, strays))
        except Exception as exc:
            changes.append(ContentChange(
                idx, "migration skipped",
                f"{type(exc).__name__}: {exc}; the slide is unchanged"))
    # Stamped here rather than at each change site: the id has to survive a
    # round trip through the review page's form, so it must be assigned once,
    # after the order is final, and never re-derived from the change's own text.
    for n, change in enumerate(changes):
        change.change_id = f"c{n}"
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue(), changes
