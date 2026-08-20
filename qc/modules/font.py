"""Font audit module (v1, audit-only).

Checks every non-empty run against the profile font role rules:
family membership (latin for non-Arabic runs, complex-script for Arabic
runs), size against the role target, mixed bold weight inside a single
paragraph, and disallowed theme font references. Never mutates the deck;
every record's action is "flagged".
"""

from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn

from qc.records import make_record
from qc.util import font_role, iter_shapes_deep
from spike.arabic import cs_typeface
from spike.resolver import resolve_run

MODULE = "font"
ARABIC_NOTE = "Arabic content, manual review"

_TITLE_PH_TYPES = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)


def _pct(raw: str | None) -> float | None:
    """ST_TextFontScalePercentOrPercentString: '80000' (thousandths of a
    percent) or '80%'. Returns percent as float, or None."""
    if raw is None:
        return None
    raw = raw.strip()
    try:
        if raw.endswith("%"):
            return float(raw[:-1])
        return int(raw) / 1000.0
    except ValueError:
        return None


def _autofit_records(ctx) -> list:
    """Titles PowerPoint has shrunk to fit (designer workflow step 3).

    When 'shrink text on overflow' fires, PowerPoint records the applied
    scale as fontScale (and often lnSpcReduction) on the slide shape's own
    bodyPr normAutofit element, so detection is a direct fact read, no
    layout inheritance needed. The fix is the designer's own prescription,
    'Stop Fitting Text to This Placeholder' (noAutofit); the title then
    renders at full size and can overflow, so it is never pre-selected:
    the designer ticks it and judges the before/after preview.
    """
    records = []
    for slide_index, slide in enumerate(ctx.prs.slides):
        for shape in slide.placeholders:
            if shape.placeholder_format.type not in _TITLE_PH_TYPES:
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            if not tf.text.strip():
                continue
            bodyPr = tf._txBody.find(qn("a:bodyPr"))
            fit = bodyPr.find(qn("a:normAutofit")) if bodyPr is not None else None
            if fit is None:
                continue
            scale = _pct(fit.get("fontScale"))
            lnsp = _pct(fit.get("lnSpcReduction"))
            shrunk = (scale is not None and scale < 99.95) or (lnsp or 0) > 0
            if not shrunk:
                continue
            old = f"{scale:g}% scale" if scale is not None else "autofit active"
            if lnsp:
                old += f", line spacing -{lnsp:g}%"
            arabic = ctx.shape_has_arabic(slide_index, str(shape.shape_id))
            msg = (f"Title is auto-shrunk to fit its placeholder ({old}). "
                   "The fix applies 'Stop Fitting Text to This Placeholder' "
                   "so the title renders at its intended size; it may then "
                   "overflow, so review the preview and shorten or resize "
                   "if needed.")
            if arabic:
                msg += f" {ARABIC_NOTE}."
            # Evidence-based severity: a title crushed to 80% or less (or a
            # heavy line-spacing squeeze) is visibly wrong; a mild shrink is
            # a judgment call.
            visibly = (scale is not None and scale <= 80) or (lnsp or 0) >= 20
            records.append(make_record(
                slide_index=slide_index, shape_id=str(shape.shape_id),
                shape_path=shape.name, module=MODULE,
                issue_type="font.title_autofit_shrunk",
                severity="error" if visibly else "warning",
                action="flagged", confidence="high",
                arabic_flag=arabic,
                property="bodyPr.normAutofit",
                old_value=old, new_value="no autofit (100%)",
                profile_rule_id="font.title_no_autofit",
                message=msg,
            ))
    return records


def _table_records(slide_index, shape, shape_path, profile) -> list:
    """Font checks inside TABLE cells (28 tables / 421 text cells escaped
    the audit on a real Arabic deck, 12/08/2026). Table font inheritance
    flows through table styles, not the placeholder chain, so only
    EXPLICIT run-level typefaces are judged: a stray Calibri typed onto a
    run, or an Arabic run whose author set a latin font but no
    complex-script one."""
    from spike.arabic import contains_arabic

    rules = profile.get("font.roles.body") or {}
    allowed_latin = rules.get("latin") or []
    allowed_cs = rules.get("complex_script") or []
    cs_target = allowed_cs[0] if allowed_cs else None
    records = []
    common = dict(slide_index=slide_index, shape_id=str(shape.shape_id),
                  shape_path=shape_path, module=MODULE)
    seen_tc = set()
    for r_idx, row in enumerate(shape.table.rows):
        for c_idx, cell in enumerate(row.cells):
            if id(cell._tc) in seen_tc:
                continue  # merged span: one cell, one set of findings
            seen_tc.add(id(cell._tc))
            for p_idx, para in enumerate(cell.text_frame.paragraphs):
                for run_idx, run in enumerate(para.runs):
                    if not run.text or not run.text.strip():
                        continue
                    rpr = run._r.find(qn("a:rPr"))
                    if rpr is None:
                        continue
                    loc = f"t{r_idx},{c_idx}/p{p_idx}/r{run_idx}"
                    latin_el = rpr.find(qn("a:latin"))
                    fam = (latin_el.get("typeface")
                           if latin_el is not None else None)
                    if contains_arabic(run.text):
                        cs_el = rpr.find(qn("a:cs"))
                        cs = (cs_el.get("typeface")
                              if cs_el is not None else None)
                        note = ((f" The fix sets '{cs_target}'; shaping "
                                 "changes with the font, so it is never "
                                 "pre-selected: ticking it is your approval.")
                                if cs_target else f" {ARABIC_NOTE}.")
                        if cs is not None and allowed_cs and cs not in allowed_cs:
                            records.append(make_record(
                                **common,
                                issue_type="font.family_out_of_set",
                                locator=loc, severity="error",
                                action="flagged", confidence="high",
                                arabic_flag=True,
                                property="rPr.cs.typeface",
                                old_value=cs, new_value=cs_target,
                                profile_rule_id="font.roles.body.complex_script",
                                message=(f"Table cell ({r_idx + 1},{c_idx + 1}): "
                                         f"complex-script typeface '{cs}' not in "
                                         f"allowed set {allowed_cs}.{note}"),
                            ))
                        elif cs is None and fam is not None and cs_target:
                            # styled by hand (explicit latin) but the Arabic
                            # side was forgotten
                            records.append(make_record(
                                **common,
                                issue_type="font.cs_typeface_missing",
                                locator=loc, severity="warning",
                                action="flagged", confidence="high",
                                arabic_flag=True,
                                property="rPr.cs.typeface",
                                old_value=None, new_value=cs_target,
                                profile_rule_id="font.roles.body.complex_script",
                                message=(f"Table cell ({r_idx + 1},{c_idx + 1}): "
                                         f"Arabic run styled with latin "
                                         f"'{fam}' but no complex-script "
                                         f"typeface.{note}"),
                            ))
                    elif (fam and allowed_latin and fam not in allowed_latin
                          and not fam.startswith("+")):
                        records.append(make_record(
                            **common,
                            issue_type="font.family_out_of_set",
                            locator=loc, severity="error",
                            action="flagged", confidence="deterministic",
                            property="rPr.latin.typeface",
                            old_value=fam, new_value=allowed_latin[0],
                            profile_rule_id="font.roles.body.latin",
                            message=(f"Table cell ({r_idx + 1},{c_idx + 1}): "
                                     f"latin family '{fam}' not in allowed "
                                     f"set {allowed_latin}."),
                        ))
    return records


def detect(ctx) -> list:
    records = _autofit_records(ctx)
    profile = ctx.profile
    tolerance = profile.get("font.size_tolerance_pt", 0.5)
    theme_refs_allowed = profile.get("font.theme_font_refs_allowed", True)

    for slide_index, slide in enumerate(ctx.prs.slides):
        for shape, shape_path in iter_shapes_deep(slide.shapes):
            if getattr(shape, "has_table", False):
                records.extend(_table_records(slide_index, shape,
                                              shape_path, profile))
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            role = font_role(shape)
            rules = profile.get(f"font.roles.{role}")
            if not rules:
                continue
            shape_id = str(shape.shape_id)
            common = dict(slide_index=slide_index, shape_id=shape_id,
                          shape_path=shape_path, module=MODULE)

            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                # (resolved bold, run is Arabic) per non-empty run, for the
                # per-paragraph mixed weight check.
                bold_seen: list[tuple[bool, bool]] = []

                for r_idx, run in enumerate(para.runs):
                    if not run.text or not run.text.strip():
                        continue
                    eff = resolve_run(run, para, shape, slide, ctx.prs)
                    is_arabic = ctx.run_is_arabic(slide_index, shape_id, p_idx, r_idx)
                    bold_seen.append((bool(eff.bold.value), is_arabic))
                    # run addressing so the fixer can apply run-level changes
                    loc = f"p{p_idx}/r{r_idx}"

                    if is_arabic:
                        # Arabic guard, scoped (12/08/2026): the latin family
                        # is never audited or substituted on an Arabic run;
                        # the complex-script typeface IS auditable and now
                        # carries a proposed target (first allowed cs font),
                        # but the fix is never pre-selected - Arabic shaping
                        # changes with the font, so the tick is the
                        # designer's explicit approval.
                        cs = cs_typeface(run)
                        allowed_cs = rules.get("complex_script", [])
                        target = allowed_cs[0] if allowed_cs else None
                        note = ((f" The fix sets '{target}'; shaping changes "
                                 "with the font, so it is never pre-selected: "
                                 "ticking it is your approval.")
                                if target else f" {ARABIC_NOTE}.")
                        if cs is None:
                            records.append(make_record(
                                **common,
                                issue_type="font.cs_typeface_missing",
                                locator=loc,
                                severity="warning", action="flagged",
                                confidence="high", arabic_flag=True,
                                property="rPr.cs.typeface",
                                old_value=None, new_value=target,
                                profile_rule_id=f"font.roles.{role}.complex_script",
                                message=(f"Arabic run has no complex-script typeface "
                                         f"(role {role}).{note}"),
                            ))
                        elif cs not in allowed_cs:
                            records.append(make_record(
                                **common,
                                issue_type="font.family_out_of_set",
                                locator=loc,
                                severity="error", action="flagged",
                                confidence="high", arabic_flag=True,
                                property="rPr.cs.typeface",
                                old_value=cs, new_value=target,
                                profile_rule_id=f"font.roles.{role}.complex_script",
                                message=(f"Complex-script typeface '{cs}' not in allowed "
                                         f"set {allowed_cs} for role {role}.{note}"),
                            ))
                    else:
                        family = eff.family.value
                        allowed_latin = rules.get("latin", [])
                        if theme_refs_allowed is False and "theme(" in eff.family.source:
                            records.append(make_record(
                                **common,
                                issue_type="font.theme_ref_disallowed",
                                locator=loc,
                                severity="error", action="flagged",
                                confidence="high",
                                property="rPr.latin.typeface",
                                old_value=family, new_value=None,
                                profile_rule_id="font.theme_font_refs_allowed",
                                message=(f"Font resolves through a theme reference "
                                         f"({eff.family.source}) but the profile "
                                         f"disallows theme font refs."),
                            ))
                        if family not in allowed_latin:
                            # Deterministic only when the run itself carries the
                            # family; inherited values leave room for cascade
                            # interpretation differences.
                            conf = ("deterministic"
                                    if eff.family.source.startswith("run.rPr")
                                    else "high")
                            records.append(make_record(
                                **common,
                                issue_type="font.family_out_of_set",
                                locator=loc,
                                severity="error", action="flagged",
                                confidence=conf,
                                property="rPr.latin.typeface",
                                old_value=family,
                                new_value=allowed_latin[0] if allowed_latin else None,
                                profile_rule_id=f"font.roles.{role}.latin",
                                message=(f"Latin family '{family}' (source "
                                         f"{eff.family.source}) not in allowed set "
                                         f"{allowed_latin} for role {role}."),
                            ))

                    target_size = rules.get("size_pt")
                    size = eff.size_pt.value
                    if (target_size is not None and size is not None
                            and abs(size - target_size) > tolerance):
                        msg = (f"Size {size}pt (source {eff.size_pt.source}) is off the "
                               f"role {role} target {target_size}pt "
                               f"(tolerance {tolerance}pt).")
                        if is_arabic:
                            msg += f" {ARABIC_NOTE}."
                        records.append(make_record(
                            **common,
                            issue_type="font.size_off_role",
                            locator=loc,
                            severity="warning", action="flagged",
                            confidence="high" if shape.is_placeholder else "medium",
                            arabic_flag=is_arabic,
                            property="rPr.sz",
                            old_value=size, new_value=target_size,
                            profile_rule_id=f"font.roles.{role}.size_pt",
                            message=msg,
                        ))

                if len({b for b, _ in bold_seen}) > 1:
                    records.append(make_record(
                        **common,
                        issue_type="font.mixed_weight",
                        severity="info", action="flagged",
                        confidence="high",
                        arabic_flag=any(a for _, a in bold_seen),
                        property="rPr.b",
                        old_value="mixed bold values within paragraph",
                        new_value=None,
                        profile_rule_id=f"font.roles.{role}.allowed_weights",
                        message=(f"Paragraph {p_idx} mixes bold and regular runs "
                                 f"(role {role}); often intentional emphasis, "
                                 "listed for completeness."),
                    ))
    return records
