"""header_footer module: audits the deck's page furniture.

Template checks (require header_footer.template in the profile):
  header_footer.missing            error   deterministic (medium for the
                                           fake-footer textbox heuristic)
  header_footer.text_mismatch      warning deterministic
  header_footer.position_mismatch  warning high
  header_footer.font_mismatch      warning medium

Furniture checks (always on; the deck's own layout/master is the baseline,
no profile needed - ground-truth calibration 20/07/2026: the designer
stamped identical page furniture on all 13 slides, rescuing two page
numbers from the top-right corner and one source line from below the
slide edge):
  header_footer.fake_slide_number  a plain text box carrying the literal
                                   page number (or a slidenum field outside
                                   any placeholder); breaks on reorder and
                                   escapes Insert > Header & Footer. Fixable
                                   when the layout/master defines a real
                                   slide-number placeholder to inherit.
  header_footer.footer_off_canvas  footer-zone text extending below the
                                   slide edge (invisible/clipped when
                                   presented). Fixable by aligning its
                                   bottom to the layout footer baseline.
"""

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn

from qc.records import make_record
from spike.arabic import contains_arabic
from spike.resolver import resolve_run

MODULE = "header_footer"

# fake page numbers live in the slide's periphery strips
_TOP_STRIP = 0.10     # shape fully above 10% of slide height
_BOTTOM_STRIP = 0.88  # shape starts below 88% of slide height
_MAX_FURNITURE_H = 432000   # 12mm: page furniture is small
_OFF_CANVAS_EMU = 18000     # 0.5mm past the slide edge = off canvas

# Property names use the OOXML placeholder type tokens (ph/@type).
_FOOTER_PROP = "ph.footer"
_SLDNUM_PROP = "ph.sldNum"
_DATE_PROP = "ph.dt"


def _find_placeholder(slide, ph_type):
    for ph in slide.placeholders:
        if ph.placeholder_format.type == ph_type:
            return ph
    return None


def _expected_position(position_emu):
    """Accept {'left':..,'top':..} or a (left, top) sequence."""
    if isinstance(position_emu, dict):
        return position_emu.get("left"), position_emu.get("top")
    if isinstance(position_emu, (list, tuple)) and len(position_emu) >= 2:
        return position_emu[0], position_emu[1]
    return None, None


def _find_fake_footer(ctx, slide, footer_text):
    """A top-level text box in the bottom 20% of the slide whose text contains
    the required footer text is likely serving as a hand-drawn footer."""
    threshold = int(ctx.prs.slide_height * 0.8)
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX or not shape.has_text_frame:
            continue
        if shape.top is None or shape.top <= threshold:
            continue
        if footer_text in shape.text_frame.text:
            return shape
    return None


def _missing_record(s_idx, prop, rule_id, what):
    return make_record(
        slide_index=s_idx, shape_id="-", shape_path=None, module=MODULE,
        issue_type="header_footer.missing", property=prop,
        severity="error", action="flagged", confidence="deterministic",
        new_value=what, profile_rule_id=rule_id,
        message=f"Required {what} placeholder is absent on this slide.",
    )


def _footer_font_record(ctx, s_idx, slide, footer_ph, allowed_latin, role):
    """One record per slide when any non-Arabic footer run resolves to a latin
    family outside the profile role list. Arabic runs are never family-checked
    here (complex-script auditing is the font module's concern)."""
    bad_families = []
    shape_arabic = ctx.shape_has_arabic(s_idx, footer_ph.shape_id)
    for p_idx, para in enumerate(footer_ph.text_frame.paragraphs):
        for r_idx, run in enumerate(para.runs):
            if not run.text.strip():
                continue
            if contains_arabic(run.text) or ctx.run_is_arabic(
                    s_idx, footer_ph.shape_id, p_idx, r_idx):
                continue
            eff = resolve_run(run, para, footer_ph, slide, ctx.prs)
            if eff.family.value not in allowed_latin:
                bad_families.append(eff.family.value)
    if not bad_families:
        return None
    return make_record(
        slide_index=s_idx, shape_id=footer_ph.shape_id, shape_path=None,
        module=MODULE, issue_type="header_footer.font_mismatch",
        property="font.latin", severity="warning", action="flagged",
        confidence="medium", old_value=bad_families[0],
        new_value=allowed_latin[0] if allowed_latin else None,
        arabic_flag=shape_arabic,
        profile_rule_id=f"font.roles.{role}.latin",
        message=(f"Footer text resolves to '{bad_families[0]}', expected one "
                 f"of {allowed_latin} for role '{role}'."),
    )


def _inherited_placeholder(slide, ph_type):
    """The layout's placeholder of this type, falling back to the master:
    the geometry a slide-level placeholder would inherit."""
    try:
        layout = slide.slide_layout
        sources = (layout, layout.slide_master)
    except Exception:
        return None
    for source in sources:
        if source is None:
            continue
        try:
            for ph in source.placeholders:
                if ph.placeholder_format.type == ph_type:
                    return ph
        except Exception:
            return None
    return None


def _has_slidenum_field(shape):
    for fld in shape._element.iter(qn("a:fld")):
        if fld.get("type") == "slidenum":
            return True
    return False


def _is_literal_page_number(text: str, number: int) -> bool:
    t = text.strip()
    return t.isdigit() and (t.lstrip("0") or "0") == str(number)


def _furniture_records(ctx):
    """Layout-baseline checks that need no profile template: the deck's own
    master already says where page furniture belongs."""
    records = []
    H = ctx.prs.slide_height
    for s_idx, slide in enumerate(ctx.prs.slides):
        number = s_idx + 1
        for shape in slide.shapes:  # furniture is top-level by construction
            if getattr(shape, "is_placeholder", False):
                continue  # drifted real placeholders are master_slide's beat
            if not getattr(shape, "has_text_frame", False):
                continue
            l, t = shape.left, shape.top
            w, h = shape.width, shape.height
            if None in (l, t, w, h):
                continue
            text = shape.text_frame.text.strip()
            arabic = ctx.shape_has_arabic(s_idx, shape.shape_id)
            in_strip = t > _BOTTOM_STRIP * H or (t + h) < _TOP_STRIP * H
            looks_number = (_is_literal_page_number(text, number)
                            or _has_slidenum_field(shape))

            if in_strip and h <= _MAX_FURNITURE_H and looks_number:
                lph = _inherited_placeholder(slide, PP_PLACEHOLDER.SLIDE_NUMBER)
                target = None
                if lph is not None and lph.left is not None and lph.top is not None:
                    target = f"sldNum placeholder ({lph.left}, {lph.top})"
                msg = ("Page number is a plain text box, not a slide-number "
                       "placeholder: it will not renumber when slides move "
                       "and ignores Insert > Header & Footer. ")
                msg += ("Replace it with a real placeholder inheriting the "
                        "layout position." if target else
                        "The layout defines no slide-number placeholder; "
                        "add one to the master, then re-audit.")
                if arabic:
                    msg += " Arabic content, manual review."
                records.append(make_record(
                    slide_index=s_idx, shape_id=shape.shape_id,
                    shape_path=None, module=MODULE,
                    issue_type="header_footer.fake_slide_number",
                    property=_SLDNUM_PROP,
                    severity="error" if target else "warning",
                    action="flagged", confidence="high",
                    old_value=text or "(slidenum field)", new_value=target,
                    arabic_flag=arabic,
                    profile_rule_id="header_footer.furniture",
                    message=msg,
                ))
            elif text and t > 0.70 * H and (t + h) > H + _OFF_CANVAS_EMU:
                fph = _inherited_placeholder(slide, PP_PLACEHOLDER.FOOTER)
                target = None
                if (fph is not None and fph.top is not None
                        and fph.height is not None):
                    target = fph.top + fph.height - h  # share the baseline
                below = t >= H
                msg = (f"Footer text sits {'entirely below' if below else 'past'} "
                       "the slide edge and will be cut off or invisible when "
                       "presented. ")
                msg += ("Fix aligns its bottom to the layout footer baseline."
                        if target is not None else
                        "Move it above the slide edge (layout has no footer "
                        "placeholder to take a baseline from).")
                if arabic:
                    msg += " Arabic content, manual review."
                records.append(make_record(
                    slide_index=s_idx, shape_id=shape.shape_id,
                    shape_path=None, module=MODULE,
                    issue_type="header_footer.footer_off_canvas",
                    property="spPr.xfrm.off.y",
                    severity="error", action="flagged", confidence="high",
                    old_value=t,
                    new_value=None if target is None else int(target),
                    arabic_flag=arabic,
                    profile_rule_id="header_footer.furniture",
                    message=msg,
                ))
    return records


def detect(ctx):
    profile = ctx.profile
    records = _furniture_records(ctx)
    tpl = profile.get("header_footer.template") or {}
    footer_text = tpl.get("footer_text")
    need_number = bool(tpl.get("slide_number"))
    date_cfg = tpl.get("date") or {}
    need_date = bool(date_cfg.get("enabled"))
    if footer_text is None and not need_number and not need_date:
        return records  # template not enforced; furniture checks still ran

    position_emu = tpl.get("position_emu")
    exp_left, exp_top = _expected_position(position_emu)
    role = tpl.get("font_role", "caption")
    allowed_latin = profile.get(f"font.roles.{role}.latin") or []
    tolerance = profile.get("master_slide.geometry_tolerance_emu", 9525)

    for s_idx, slide in enumerate(ctx.prs.slides):
        if footer_text is not None:
            footer_ph = _find_placeholder(slide, PP_PLACEHOLDER.FOOTER)
            if footer_ph is None:
                fake = _find_fake_footer(ctx, slide, footer_text)
                if fake is not None:
                    records.append(make_record(
                        slide_index=s_idx, shape_id=fake.shape_id,
                        shape_path=None, module=MODULE,
                        issue_type="header_footer.missing",
                        property=_FOOTER_PROP, severity="error",
                        action="flagged", confidence="medium",
                        new_value=footer_text,
                        arabic_flag=ctx.shape_has_arabic(s_idx, fake.shape_id),
                        profile_rule_id="header_footer.template.footer_text",
                        message=("Footer placeholder is absent; a text box at "
                                 "the footer position may be serving as the "
                                 "footer."),
                    ))
                else:
                    records.append(_missing_record(
                        s_idx, _FOOTER_PROP,
                        "header_footer.template.footer_text", "footer"))
            else:
                footer_arabic = ctx.shape_has_arabic(s_idx, footer_ph.shape_id)
                actual_text = footer_ph.text_frame.text
                if actual_text != footer_text:
                    records.append(make_record(
                        slide_index=s_idx, shape_id=footer_ph.shape_id,
                        shape_path=None, module=MODULE,
                        issue_type="header_footer.text_mismatch",
                        property="text", severity="error", action="flagged",
                        confidence="deterministic", old_value=actual_text,
                        new_value=footer_text, arabic_flag=footer_arabic,
                        profile_rule_id="header_footer.template.footer_text",
                        message=(f"Footer text is '{actual_text}', expected "
                                 f"'{footer_text}'."),
                    ))
                # Placeholder geometry can still be None after python-pptx
                # inheritance resolution; only compare axes we can read.
                if (exp_left is not None or exp_top is not None) and \
                        footer_ph.left is not None and footer_ph.top is not None:
                    off_left = (exp_left is not None
                                and abs(footer_ph.left - exp_left) > tolerance)
                    off_top = (exp_top is not None
                               and abs(footer_ph.top - exp_top) > tolerance)
                    if off_left or off_top:
                        msg = (f"Footer placeholder at ({footer_ph.left}, "
                               f"{footer_ph.top}) EMU, expected ({exp_left}, "
                               f"{exp_top}) EMU.")
                        if footer_arabic:
                            # Footer position is direction-sensitive in RTL
                            # decks; never auto-fix without a human look.
                            msg += " Arabic content, manual review."
                        records.append(make_record(
                            slide_index=s_idx, shape_id=footer_ph.shape_id,
                            shape_path=None, module=MODULE,
                            issue_type="header_footer.position_mismatch",
                            property="position", severity="warning",
                            action="flagged", confidence="high",
                            old_value=f"({footer_ph.left}, {footer_ph.top})",
                            new_value=f"({exp_left}, {exp_top})",
                            arabic_flag=footer_arabic,
                            profile_rule_id="header_footer.template.position_emu",
                            message=msg,
                        ))
                font_rec = _footer_font_record(
                    ctx, s_idx, slide, footer_ph, allowed_latin, role)
                if font_rec is not None:
                    records.append(font_rec)

        if need_number and _find_placeholder(
                slide, PP_PLACEHOLDER.SLIDE_NUMBER) is None:
            records.append(_missing_record(
                s_idx, _SLDNUM_PROP,
                "header_footer.template.slide_number", "slide number"))

        if need_date and _find_placeholder(slide, PP_PLACEHOLDER.DATE) is None:
            records.append(_missing_record(
                s_idx, _DATE_PROP, "header_footer.template.date", "date"))

    return records
