"""margin_alignment: safe-zone breaches, edge misalignment, uneven spacing,
squeezed text boxes, and text-on-text overlap.

Real-deck tuning (14/07/2026, consulting-deck feedback): the original pass
only compared shapes whose edges sat within ~0.8mm of each other, checked
left edges only, and never looked inside groups, so the misalignments
designers actually see (3-20mm, tops as often as lefts, inside grouped
cards) sailed through. Now:

- clustering uses an INTENT WINDOW (default 0.15in): shapes whose edges sit
  within it are treated as meant-to-align; deviation from the cluster median
  beyond the profile tolerance flags (and snaps, on approval)
- both axes: left edges of columns, top edges of rows
- spacing is checked along rows AND down stacks; a single odd gap gets the
  translate-the-tail fix, several odd gaps get an honest flag-only record
- group interiors are audited in group coordinate space (unscaled groups
  only: a resized group renders child EMU at a different scale, so
  comparisons there would lie); safe-zone stays top-level, the group's own
  bounding box is what margins see
- a text box squeezed so narrow it wraps a letter per line is flagged
- two text-bearing shapes overlapping is flagged (text on text is almost
  never intent in a consulting deck)

Two standing rules for everything below:

The BOX is what a margin is measured against, never the text drawn inside it
(design lead, 19/08/2026). A text frame's insets, its autofit scale and the
glyph extent of a long line are all rendering, and rendering varies with the
font actually installed; the stored spPr xfrm is the only edge the deck itself
states. So all four sides compare box edges - left/top against the margin,
left+width / top+height against the opposite one - and a line of text spilling
outside a box that sits correctly is not a margin breach here.

The one thing inside a box that IS compared is its vertical anchor, because
that is an attribute the deck states rather than a rendering: a row of aligned
boxes whose text anchors disagree reads as ragged however the fonts resolve.
See _text_anchor_mismatch.

A HEADING past a margin is flagged and never acted on. Whether a title or
standfirst may break the margin is the client's house style, not a defect, so
it gets its own issue type, no computed target, and no place in the rescale
selection. See qc.util.heading_ids.

The HEADER BAND is the exception to "the safe zone governs text only". A master
that draws two horizontal guides under its subtitle is reserving the strip
between them: the subtitle stops above it, the body starts below it, and the
gap stays empty on every slide. That is a statement about the page, not about
copy, so a photo standing in the strip is the same defect as a paragraph and
gets the same finding - one per slide, cleared by one move of the whole body
block. See _band_intrusion and qc.stylespec.read_content_band.
"""

import statistics
from collections import Counter

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from qc.records import make_record
from qc.util import (full_height_panel, heading_ids, is_backdrop,
                     rides_with, slide_is_rtl)

MODULE = "margin_alignment"

ARABIC_NOTE = "Arabic content, manual review"

DEFAULT_INTENT_WINDOW = 137160   # 0.15 inch: "these were meant to line up"
OBVIOUS_EDGE_EMU = 36000         # ~1mm: a deviation the eye actually sees
CONTAIN_MIN = 0.66               # 2/3 of a shape inside another = contained
                                 # (keep in sync with qc/fixer.py)
SQUEEZED_MAX_WIDTH = 548640      # 0.6 inch
SQUEEZED_MIN_ASPECT = 3.0        # height > 3x width reads as a strip
# Boxes within this height ratio of each other are one line of labels rather
# than assorted content that happens to start on the same line, which is what
# makes their text anchors comparable at all.
ANCHOR_H_RATIO = 1.25
_ANCHOR_LABEL = {"t": "top", "ctr": "middle", "b": "bottom",
                 "just": "top", "dist": "top"}
OVERLAP_MIN_RATIO = 0.4          # of either shape's area: substantial cover
# The text-overlap check stays O(n^2)-sane. 80 was set before the pass was
# profiled and it cut the answer rather than the cost: this list is already
# filtered to text-bearing shapes, so 80 of them is a genuinely crowded slide -
# exactly where collisions live - and the pairs are cheap now that the frame
# marker read is memoized. Raised with the design pass's own cap (250), and it
# reports what it dropped instead of going quiet (measured 24/08/2026).
MAX_PAIRWISE_SHAPES = 250

# Panel rows (ground truth, 20/07/2026: 49% of a designer's geometry work
# was moving panels as blocks onto shared row lines; the tool's per-shape
# window never saw a 6.2mm panel offset).
PANEL_WINDOW = 540000            # 15mm: panel rows drift bigger than text
PANEL_MIN_W = 720000             # 20mm wide and
PANEL_MIN_H = 360000             # 10mm tall: a panel is substantial
PANEL_H_RATIO = 1.25             # same row = same height class
PANEL_X_OVERLAP = 0.2            # side-by-side, not stacked

# Cluster rhythm (ground truth, 20/07/2026: the wheel slide's bottom image
# clusters sat 11.6/16.4mm below their column rhythm; the designer lifted
# each cluster - image, labels, underlines - as one axis-locked move).
LIFT_MAX_RATIO = 6               # odd gap beyond this = a deliberate break
LIFT_OBVIOUS_RATIO = 3           # unmistakable rhythm break
COHORT_BIN_EMU = 50800           # ~1.4mm size bins: "the same kind of item"

# Content overflow (ground truth, 20/07/2026: the rankings slide's content
# spanned 0..338.7mm - past BOTH slide edges - and the designer's fix was
# one select-all proportional rescale to 92.8% anchored at the left
# margin, not fifty nudges).
SCALE_MIN = 0.80                 # below this a rescale is too drastic to offer
SCALE_SKIP = 0.97                # above this the per-shape checks do it better
MARGIN_SLACK_EMU = 72000         # 2mm: designers park bleed ON the edge
MIN_BREACHERS = 4                # a LAYOUT is too wide, not one stray shape

# A heading whose box sits this close to a margin is ON it. Guides are stored
# in eighths of a point and placeholder extents are typed in inches, so the two
# disagree by a few hundred EMU on a box a designer placed exactly; reporting
# that as "the title breaks the margin" would spend a designer's attention on
# rounding.
HEADING_SLACK_EMU = 36000        # 1mm

# Same rounding allowance for the body ceiling, and for the same reason: on the
# client master measured, the master's OWN body placeholder starts 0.007in above
# its body guide, so a shape a designer aligned to the master reads as a
# breach without this.
BAND_SLACK_EMU = 36000           # 1mm
# Page furniture lives in the bottom strip and never rides a body move (the
# same strip qc/fixer.py and qc/migrate.py keep pinned).
FURNITURE_STRIP = 0.88
# An intrusion at least this deep is unmistakable; below it the shape is a
# millimetre or two off a line and the finding is a judgment call, so it is
# reported as one (same evidence-based severity as the edge checks).
BAND_OBVIOUS_EMU = 108000        # 3mm

# Cross-slide anchors (title bars, header chips): a straggler within this
# window of the deck-wide modal position snaps back to it; further away it
# reads as a deliberately different layout and is left alone.
PIN_WINDOW = 1080000             # 30mm


def _margin_breaches(left, top, width, height, margins, slide_w, slide_h):
    """[(side, overshoot_emu)] for every margin the shape's BOX crosses.

    All four sides, measured the same way and from the same source of truth:
    the stored xfrm. Text insets and glyph extents are deliberately not
    consulted (see the module docstring) - a box on the line is on the line,
    whatever its text does inside it."""
    return [(side, over) for side, over in (
        ("left", margins.get("left", 0) - left),
        ("top", margins.get("top", 0) - top),
        ("right", (left + width) - (slide_w - margins.get("right", 0))),
        ("bottom", (top + height) - (slide_h - margins.get("bottom", 0))),
    ) if over > 0]


def _body_band(profile, prs):
    """(subtitle_floor, body_top) for the strip a master reserves under its
    header, or (None, None) when nobody states one.

    The profile is asked first: it is the client's stated rule, and it outlives
    the master file it was read from. A profile written before this field
    existed falls back to the guides on the deck's OWN master - the same source
    the profile would have been projected from, and the source the migration
    pass already reads its frame from (qc.migrate._margin_frame). Nothing is
    inferred: no guides, no band, and no slide is measured against a line the
    client never drew.

    The ceiling can also come from a PRESENTATION SPACE rectangle drawn around
    the body. infer_grid already works out which of the two things such a
    rectangle is - a frame around the whole page has a top margin, a frame
    around the body has a line where content begins - and sets body_top_emu
    only in the second case (space_states == "body"). This used to insist the
    source be "guides", so a master that states the body's top as a rectangle
    rather than a guide pair got no ceiling at all, and every slide on it passed
    however far below the line its content began (design lead, 24/08/2026).

    A rectangle gives a ceiling and no floor, and that asymmetry is fine because
    the two directions need different evidence: content that has crept UP needs
    the floor to tell an eyebrow in the header from body content in the strip
    (_band_intrusion requires both), while content sitting BELOW the line needs
    only the line (_band_shortfall). Note what is NOT read here - the safe-zone
    top margin. That is where the PAGE begins, not where the body does; using it
    as a ceiling fires on every normal slide in the deck."""
    band = profile.get("geometry.body_band_emu") or {}
    if band.get("body_top"):
        return (band.get("subtitle_floor"), band["body_top"])
    try:
        from qc.stylespec import dominant_master, infer_grid

        grid = infer_grid(prs, dominant_master(prs)) or {}
    except Exception:
        return (None, None)
    if grid.get("source") not in ("guides", "presentation_space"):
        return (None, None)
    return (grid.get("subtitle_floor_emu"), grid.get("body_top_emu"))


def _space_box(profile, prs):
    """(left, top, right, bottom) of the presentation space in slide EMU, or
    None when nobody states one.

    The frame reaches this stage the same way the band does: the profile first,
    because it is the client's stated rule and outlives the master file it came
    from, then the deck's OWN master. A profile projected from a master carries
    the space as its safe-zone margins (qc.stylespec.infer_grid), which is why
    that key is read here rather than a second one being invented - two places
    to state one rectangle is two places for them to disagree.

    Until now the frame only reached the design pass's outside-the-frame check.
    The audit knew the margins but not that they were a designer's stated
    rectangle, so it had nothing to align TO and fell back to aligning shapes to
    each other (design lead, 24/08/2026)."""
    slide_w, slide_h = prs.slide_width, prs.slide_height
    margins = profile.get("geometry.safe_zone_margins_emu") or {}
    if all(margins.get(side) is not None
           for side in ("left", "top", "right", "bottom")):
        return (margins["left"], margins["top"],
                slide_w - margins["right"], slide_h - margins["bottom"])
    try:
        from qc.stylespec import dominant_master, read_presentation_space

        space = read_presentation_space(prs, dominant_master(prs))
    except Exception:
        return None
    if not space or space.get("problem") or not space.get("box_emu"):
        return None
    return tuple(space["box_emu"])


# Past this far inside the frame, an element is INDENTED rather than misaligned,
# and the tool says nothing. The two failure modes sit at opposite ends of the
# scale, which is why this is a band and not a floor: a stray of a hundredth of
# an inch is a nudge nobody meant, and a quarter inch is a designer's indent - a
# sub-bullet, an inset aside, a hanging quote. A rule with only a lower bound
# reports the deliberate ones too, which is how a check earns its way into being
# ignored. The lower bound is the client's own stated alignment tolerance, so
# "off the line" means what the profile says it means.
SPACE_EDGE_MAX_EMU = 228600      # 0.25in


def _space_edge_stacks(s_idx, body, measurable, space, rtl, edge_tol):
    """Stacked elements not starting on the frame's leading edge.

    "All stacked elements should start at the left/right boundary of the space,
    not just the first element" (design lead, 24/08/2026). The complaint names
    the symptom precisely: with a column of blocks whose first sits on the frame
    and whose rest are indented a few millimetres, _edge_misaligned takes the
    MAJORITY as the line - so the one element that was right gets pulled off the
    frame to join the ones that were wrong.

    The frame only counts as live on a slide where something actually uses it, so
    a stack is judged only when one of its own members sits on the edge. That
    keeps a deliberately inset column - a card's contents, an indented aside -
    out of it: nothing in such a column touches the frame, so nothing in it is
    asked to.

    Leading edge, not left edge: an RTL slide's stacks run from the right, and
    holding Arabic content to a left margin would be the tool imposing a reading
    direction the deck does not have."""
    if space is None or len(measurable) < 2:
        return []
    left, _top, right, _bottom = space
    records = []
    for stack in _x_columns(measurable):
        if len(stack) < 2:
            continue          # one element is not a stack
        def _off(it):
            return (right - (it["l"] + it["w"])) if rtl else (it["l"] - left)

        offs = [_off(it) for it in stack]
        if not any(abs(o) <= edge_tol for o in offs):
            continue          # nothing in this stack uses the frame
        strays = [(it, o) for it, o in zip(stack, offs)
                  if edge_tol < o <= SPACE_EDGE_MAX_EMU]
        if not strays:
            continue
        # ONE record for the stack, not one per stray. "Not just the first
        # element" is a single decision about a column, and three tick boxes
        # that all snap to the same edge would be three fixes claiming each
        # other's shapes - the fix engine's one-move-per-shape guard then
        # applies the first and defers the rest, so a designer ticking all
        # three would watch two of them do nothing.
        #
        # locator "edge:<l|r>:<strays>:<stack>". The side is stated rather than
        # inferred: the fixer has to know whether it is aligning left edges to
        # the frame or right edges to it, and working that out from the
        # coordinates would be a second guess at the slide's reading direction
        # that could disagree with the one made here. The fixer moves the
        # strays; _guide_beats_median needs the whole stack, because the shape
        # the cluster median wants to move is the one this rule leaves alone -
        # the member already sitting on the frame, out-voted by its indented
        # siblings.
        stray_ids = ",".join(str(it["shape"].shape_id) for it, _o in strays)
        stack_ids = ",".join(str(it["shape"].shape_id) for it in stack)
        worst, off = max(strays, key=lambda pair: pair[1])
        # RTL snaps each shape's RIGHT edge to the frame, so shapes of different
        # widths get different lefts; that is the fixer's business, and the
        # record states the frame edge itself.
        target = right if rtl else left
        on_line = len(stack) - len(strays)
        records.append(make_record(
            slide_index=s_idx, shape_id=worst["shape"].shape_id,
            shape_path=None, module=MODULE,
            issue_type="margin_alignment.space_edge_misaligned",
            severity="error" if on_line >= len(strays) else "warning",
            action="flagged", confidence="high",
            locator=f"edge:{'r' if rtl else 'l'}:{stray_ids}:{stack_ids}",
            property="spPr.xfrm.off.x",
            old_value=worst["l"], new_value=target,
            profile_rule_id="geometry.safe_zone_margins_emu",
            arabic_flag=any(it["arabic"] for it, _o in strays),
            message=_geo_msg(
                f"{len(strays)} element(s) stacked here sit up to "
                f"{off / 36000:.1f}mm inside the presentation space's "
                f"{'right' if rtl else 'left'} edge while {on_line} stacked "
                f"with them sit on it; the fix snaps them all out to the "
                f"frame, so the column starts on one line",
                any(it["arabic"] for it, _o in strays)),
        ))
    return records


def _clusters(items, key_index: int, span: int):
    """Greedy 1-D clustering: sort by the keyed value and group items whose
    values all sit within `span` of the cluster minimum."""
    clusters = []
    current = []
    for item in sorted(items, key=lambda t: t[key_index]):
        if current and item[key_index] - current[0][key_index] > span:
            clusters.append(current)
            current = []
        current.append(item)
    if current:
        clusters.append(current)
    return clusters


def _msg(base: str, arabic: bool) -> str:
    return f"{base}; {ARABIC_NOTE}" if arabic else base


def _geo_msg(base: str, arabic: bool) -> str:
    """For fixable GEOMETRY findings: the move never opens the text, so
    Arabic content gets a transparency note, not a manual-review demand."""
    return (f"{base}; contains Arabic text (geometry-only fix, text "
            "untouched)") if arabic else base


def _unscaled_groups(shapes):
    """Groups (recursively) whose child space renders 1:1, i.e. ext ==
    chExt within 1%. A resized group scales child EMU on screen, so
    child-space comparisons there are unreliable and skipped."""
    for shape in shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.GROUP:
            continue
        xfrm = shape._element.find(qn("p:grpSpPr") + "/" + qn("a:xfrm"))
        ok = False
        if xfrm is not None:
            ext = xfrm.find(qn("a:ext"))
            ch = xfrm.find(qn("a:chExt"))
            if ext is not None and ch is not None:
                try:
                    cx, ccx = int(ext.get("cx")), int(ch.get("cx"))
                    cy, ccy = int(ext.get("cy")), int(ch.get("cy"))
                    ok = (ccx and ccy
                          and abs(cx - ccx) <= 0.01 * ccx
                          and abs(cy - ccy) <= 0.01 * ccy)
                except (TypeError, ValueError):
                    ok = False
        if ok:
            yield shape
        # descend regardless: nested unscaled groups are still auditable
        yield from _unscaled_groups(shape.shapes)


MOSAIC_MIN_MEMBERS = 5
MOSAIC_TOUCH_EMU = 2000   # abutting or near-abutting counts as touching


def _mosaic_members(pool) -> set:
    """Indexes of shapes that form a MOSAIC: a connected component of >=5
    mutually touching/abutting shapes (a segments wheel, a stepped podium
    chart, a fan diagram). Such fragments compose ONE visual object whose
    internal geometry follows the diagram, not the layout grid - clustering
    them against outside edges produces nonsense (ground truth, 19/07/2026:
    the diagram slides carried most of the residual alignment noise on a
    designer-finished deck)."""
    n = len(pool)
    parent = list(range(n))

    def find(i):
        while parent[i] != i:
            parent[i] = parent[parent[i]]
            i = parent[i]
        return i

    for i in range(n):
        a = pool[i]
        for j in range(i + 1, n):
            b = pool[j]
            ix = (min(a[1] + a[3], b[1] + b[3]) - max(a[1], b[1])
                  + MOSAIC_TOUCH_EMU)
            iy = (min(a[2] + a[4], b[2] + b[4]) - max(a[2], b[2])
                  + MOSAIC_TOUCH_EMU)
            if ix > 0 and iy > 0:
                ri, rj = find(i), find(j)
                if ri != rj:
                    parent[ri] = rj
    sizes: dict[int, int] = {}
    for i in range(n):
        sizes[find(i)] = sizes.get(find(i), 0) + 1
    return {i for i in range(n) if sizes[find(i)] >= MOSAIC_MIN_MEMBERS}


def _wall_cohort_ids(pool) -> set:
    """shape_ids belonging to a WALL: >=6 same-size shapes (photo grids,
    logo walls). Their spacing varies organically; a computed nudge is a
    guess, so wall spacing findings are advisory (flag-only)."""
    from collections import Counter

    def key(item):
        return (round(item[3] / 25400), round(item[4] / 25400))  # ~2mm bins

    counts = Counter(key(item) for item in pool)
    return {str(item[0].shape_id) for item in pool if counts[key(item)] >= 6}


def _containment(pool, slide_w=0, slide_h=0) -> dict:
    """index -> index of the SMALLEST pool shape this one COMPOSES with: the
    anchor of its collection.

    Two ways to belong, and the same consequence either way:

    - CONTAINED: >= CONTAIN_MIN of its area inside a strictly larger shape. A
      card's icon or label belongs to its card and must be compared to how peers
      sit in THEIR cards, never thrown into a global cluster with unrelated
      edges (real-deck finding, 19/07/2026: icons snapped out of their cards
      onto a neighboring cluster).
    - RIDING: welded to a larger shape or sitting within a satellite gap of it,
      inside its span on the other axis (qc.util.rides_with). A corner rule
      against a photo, a quote mark above a paragraph, a caption under it. These
      used to join the absolute clusters as if they were free shapes, so each
      could be nudged on its own and the collection came apart (design lead,
      20/08/2026, the leadership-quotes slide).

    Backdrops are never anchors: a shape that big is the ground the collection
    sits on, and everything on the slide would "belong" to it."""
    out: dict[int, int] = {}
    for i, a in enumerate(pool):
        a_area = a[3] * a[4]
        if a_area <= 0:
            continue
        a_box = (a[1], a[2], a[1] + a[3], a[2] + a[4])
        best = None
        best_area = None
        for j, b in enumerate(pool):
            if i == j:
                continue
            b_area = b[3] * b[4]
            if b_area <= a_area:
                continue
            b_box = (b[1], b[2], b[1] + b[3], b[2] + b[4])
            if slide_w and is_backdrop(b_box, slide_w, slide_h):
                continue
            ix = max(0, min(a_box[2], b_box[2]) - max(a_box[0], b_box[0]))
            iy = max(0, min(a_box[3], b_box[3]) - max(a_box[1], b_box[1]))
            if ix * iy >= CONTAIN_MIN * a_area or rides_with(a_box, b_box):
                if best is None or b_area < best_area:
                    best, best_area = j, b_area
        if best is not None:
            out[i] = best
    return out


def _placed_pool(ctx, s_idx, shapes):
    """(shape, left, top, width, height, arabic) for unrotated shapes with
    resolved geometry, in the coordinate space the caller iterates."""
    pool = []
    for shape in shapes:
        left, top = shape.left, shape.top
        width, height = shape.width, shape.height
        if left is None or top is None or width is None or height is None:
            continue
        if getattr(shape, "rotation", 0):
            continue
        arabic = ctx.shape_has_arabic(s_idx, shape.shape_id)
        pool.append((shape, left, top, width, height, arabic))
    return pool


def detect(ctx):
    profile = ctx.profile
    margins = profile.get("geometry.safe_zone_margins_emu") or {}
    edge_tol = profile.get("geometry.alignment.edge_tolerance_emu", 9525)
    spacing_tol = profile.get("geometry.alignment.spacing_tolerance_emu", 9525)
    window = profile.get("geometry.alignment.intent_window_emu",
                         DEFAULT_INTENT_WINDOW)
    slide_w = ctx.prs.slide_width
    slide_h = ctx.prs.slide_height

    # Full-bleed exemption (real-deck tuning): a shape covering most of the
    # slide, or spanning nearly its full width/height, is design intent
    # (background image, edge-to-edge band) and legitimately crosses the
    # safe zone. It neither flags nor joins the alignment pools.
    bleed_area = profile.get("geometry.full_bleed_min_area", 0.7)
    span_ratio = profile.get("geometry.full_span_min_ratio", 0.95)
    slide_area = slide_w * slide_h

    from qc.util import recurring_anchors

    anchors = recurring_anchors(ctx.prs)
    band_floor, band_ceiling = _body_band(profile, ctx.prs)
    space = _space_box(profile, ctx.prs)

    records = []
    for s_idx, slide in enumerate(ctx.prs.slides):
        # Resolved once per slide and shared: the breach check and the overflow
        # rescale must agree about which shape is the title, or one of them
        # reports a heading the other has already excused.
        head_ids = heading_ids(slide, ctx.prs)
        placed = []  # top-level, non-bleed pool for alignment + spacing
        # Every top-level shape with resolved geometry, flagged rather than
        # filtered. The header-band rule needs the shapes the pools above drop:
        # a picture (no text), and a rotated element that must still travel with
        # the body block even though its stored box cannot be measured.
        surface = []
        for shape in slide.shapes:
            left, top = shape.left, shape.top
            width, height = shape.width, shape.height
            # Placeholder geometry inheritance can still resolve to None.
            if left is None or top is None or width is None or height is None:
                continue
            arabic = ctx.shape_has_arabic(s_idx, shape.shape_id)
            full_bleed = (width * height >= bleed_area * slide_area
                          or width >= span_ratio * slide_w
                          or height >= span_ratio * slide_h)
            is_head = str(shape.shape_id) in head_ids
            surface.append({
                "shape": shape, "t": top, "h": height, "arabic": arabic,
                # left/width too: the band rules group content into side-by-side
                # columns, and a column is a horizontal span (_x_columns)
                "l": left, "w": width,
                "head": is_head, "bleed": full_bleed,
                "ph": bool(getattr(shape, "is_placeholder", False)),
                "rot": bool(getattr(shape, "rotation", 0)),
            })
            # Computed HERE, before anything else can act on this shape. The
            # cross-slide pin below is a fix, pre-selectable and confident, and
            # a heading the client may want running wide must not be moved by
            # it. A full-bleed or rotated shape is excused as it always was.
            head_note = None
            if (is_head and not full_bleed
                    and not getattr(shape, "rotation", 0)
                    and getattr(shape, "has_text_frame", False)
                    and shape.text_frame.text.strip()):
                head_note = _heading_past_margin(
                    s_idx, shape,
                    _margin_breaches(left, top, width, height, margins,
                                     slide_w, slide_h)
                    + _heading_past_ceiling(top, height, band_floor,
                                            band_ceiling), arabic)

            anchor = anchors.get((s_idx, str(shape.shape_id)))
            if anchor is not None:
                # cross-slide anchors never join local pools; a straggler
                # gets pinned back to the deck-wide modal position instead -
                # unless it is a heading already outside the frame, which this
                # tool reports and never repositions
                if head_note is not None:
                    records.append(head_note)
                    continue
                ml, mt = anchor
                off = max(abs(left - ml), abs(top - mt))
                if edge_tol < off <= PIN_WINDOW:
                    records.append(make_record(
                        slide_index=s_idx, shape_id=shape.shape_id,
                        shape_path=None, module=MODULE,
                        issue_type="margin_alignment.recurring_off_position",
                        severity="error", action="flagged",
                        confidence="high",
                        property="spPr.xfrm.off",
                        locator=f"pin:{ml},{mt}",
                        old_value=f"({left}, {top})",
                        new_value=f"({ml}, {mt})",
                        profile_rule_id="geometry.alignment.edge_tolerance_emu",
                        arabic_flag=arabic,
                        message=_geo_msg(
                            f"this element sits at the same spot on most "
                            f"slides but is {off // 36000}mm off here; the "
                            f"fix snaps it back to the deck-wide position",
                            arabic),
                    ))
                continue
            if full_bleed:
                continue  # full-bleed / full-span: intentional by definition
            if getattr(shape, "rotation", 0):
                # Stored xfrm offsets do not describe the rendered bbox of a
                # rotated shape, so no comparison is trustworthy. Only text
                # shapes get the manual-check note (safe zones govern text).
                if getattr(shape, "has_text_frame", False) \
                        and shape.text_frame.text.strip():
                    records.append(make_record(
                        slide_index=s_idx, shape_id=shape.shape_id, shape_path=None,
                        module=MODULE, issue_type="margin_alignment.outside_safe_zone",
                        severity="warning", action="skipped", confidence="low",
                        property="spPr.xfrm.off",
                        profile_rule_id="geometry.safe_zone_margins_emu",
                        arabic_flag=arabic,
                        message=_msg("rotated shape, stored bounding box "
                                     "unreliable, manual check", arabic),
                    ))
                continue
            placed.append((shape, left, top, width, height, arabic))

            if is_head:
                # One shape, one answer. A heading is never handed on to the
                # plain breach record, which offers no client question and
                # reads as a defect rather than a decision.
                if head_note is not None:
                    records.append(head_note)
                continue

            # Safe zones govern TEXT content: designers keep text on the
            # grid while decorative graphics roam past it by design (ground
            # truth, 19/07/2026: 82% of graphic "breaches" on a finished
            # deck were intentional).
            has_text = (getattr(shape, "has_text_frame", False)
                        and shape.text_frame.text.strip())
            if not has_text:
                continue
            breached = _margin_breaches(left, top, width, height, margins,
                                        slide_w, slide_h)
            if not breached:
                continue
            records.append(make_record(
                slide_index=s_idx, shape_id=shape.shape_id, shape_path=None,
                module=MODULE, issue_type="margin_alignment.outside_safe_zone",
                severity="warning", action="flagged", confidence="deterministic",
                property="spPr.xfrm.off",
                profile_rule_id="geometry.safe_zone_margins_emu",
                arabic_flag=arabic,
                message=_msg("shape breaches safe zone edges: "
                             f"{', '.join(side for side, _o in breached)}",
                             arabic),
            ))

        # pools: the slide surface plus each unscaled group's interior
        pools = [("", placed)]
        for group in _unscaled_groups(slide.shapes):
            pool = _placed_pool(ctx, s_idx, group.shapes)
            if len(pool) >= 2:
                pools.append((f"group {group.shape_id}", pool))

        for scope, pool in pools:
            # Contained shapes (a card's icon, label, body) are judged by
            # how they sit WITHIN their container, relative to peers in
            # theirs; only free shapes and the containers themselves join
            # the absolute clusters.
            cont = _containment(pool, slide_w, slide_h)
            free_all = [item for i, item in enumerate(pool) if i not in cont]
            mosaic = _mosaic_members(free_all)  # indexes into free_all
            free = [item for i, item in enumerate(free_all) if i not in mosaic]
            contained = [(item, pool[ci]) for i, item in enumerate(pool)
                         for ci in (cont.get(i),) if ci is not None]
            panel_recs = _panel_row_misaligned(s_idx, pool, cont, edge_tol,
                                               scope)
            records.extend(panel_recs)
            # a panel already judged at row level must not also be nudged by
            # the per-shape pass (one shape, one move)
            suppress = {(r.shape_id, r.property) for r in panel_recs}

            def _fresh(recs):
                return [r for r in recs
                        if (r.shape_id, r.property) not in suppress]

            records.extend(_fresh(_edge_misaligned(s_idx, free, "left",
                                                   edge_tol, window, scope)))
            records.extend(_fresh(_edge_misaligned(s_idx, free, "top",
                                                   edge_tol, window, scope)))
            records.extend(_relative_edge_misaligned(
                s_idx, contained, "left", edge_tol, window, scope))
            records.extend(_relative_edge_misaligned(
                s_idx, contained, "top", edge_tol, window, scope))
            # rhythm runs on the pre-mosaic pool: a same-size cohort is a
            # stronger structural claim than mosaic adjacency (a connector
            # brushing one image must not break its column's rhythm line),
            # and the wall rule already keeps repeated diagram fragments out
            rhythm = (_cohort_rhythm(s_idx, free_all, "col", window,
                                     spacing_tol, scope)
                      + _cohort_rhythm(s_idx, free_all, "row", window,
                                       spacing_tol, scope))
            records.extend(rhythm)
            # a shape already owed a rhythm lift must not also get the
            # weaker band-line spacing nudge (one shape, one move)
            lifted = {(r.shape_id, r.locator.split(":", 1)[0][5:])
                      for r in rhythm}
            for direction in ("row", "col"):
                records.extend(
                    r for r in _uneven_spacing(s_idx, free, direction,
                                               window, spacing_tol, scope)
                    if (r.shape_id, direction) not in lifted)
            records.extend(_squeezed_text(s_idx, pool, scope))
            records.extend(_text_overlap(s_idx, pool, scope))
            records.extend(_text_anchor_mismatch(s_idx, free, window, scope))

        records.extend(_content_overflow(s_idx, placed, margins,
                                         slide_w, slide_h, head_ids,
                                         slide_is_rtl(slide)))
        records.extend(_band_intrusion(s_idx, surface, band_floor,
                                       band_ceiling, slide_h))
        # Content sitting BELOW the line needs only the line. A master that
        # states its body's top as a presentation-space rectangle rather than a
        # guide pair gives a ceiling and no floor, and that is enough here
        # (see _body_band) - which is what makes "adjacent elements start at the
        # top of the space" hold on those masters.
        if band_ceiling and not band_floor:
            body, measurable = _body_members(surface, slide_h)
            records.extend(_band_shortfall(s_idx, body, measurable,
                                           band_ceiling, slide_h))
        body, measurable = _body_members(surface, slide_h)
        records.extend(_space_edge_stacks(s_idx, body, measurable, space,
                                          slide_is_rtl(slide), edge_tol))

    return _guide_beats_median(records)


def _guide_beats_median(records: list) -> list:
    """Drop top-edge cluster fixes for shapes a body-guide fix already claims.

    The two rules genuinely disagree, and one of them is right. On the slide
    that prompted the guide rule, _edge_misaligned proposed pulling the photo
    DOWN 0.25in onto the cards' line (they are the majority, so they set the
    median) while the guide rule proposes lifting both UP onto the line the
    master states. Left in, a designer ticking both gets the photo moved twice
    from two targets computed off the same starting position, and the second
    move lands somewhere neither rule asked for.

    The guide wins because it is stated rather than inferred: a median is the
    tool reading intent off where things happen to sit, and a guide is the
    designer having said where they go. Only the TOP axis is affected - a
    left-edge cluster and a body guide are answers to different questions."""
    claimed_y, claimed_x = set(), set()
    for r in records:
        if r.issue_type == "margin_alignment.body_below_band":
            loc = r.locator or ""
            if loc.startswith("band:"):
                claimed_y.update(loc.split(":", 2)[2].split(","))
        elif r.issue_type == "margin_alignment.space_edge_misaligned":
            # The frame's leading edge, same argument on the other axis - and the
            # claim covers the WHOLE stack, because the shape the median fix
            # wants to move is precisely the one this rule leaves alone: the
            # member already sitting on the frame, which the majority of
            # indented siblings out-votes.
            loc = r.locator or ""
            if loc.startswith("edge:"):
                claimed_x.update(loc.split(":", 3)[3].split(","))
            else:
                claimed_x.add(str(r.shape_id))
    if not (claimed_y or claimed_x):
        return records

    def _superseded(r) -> bool:
        if r.issue_type != "margin_alignment.edge_misaligned":
            return False
        prop = r.property or ""
        sid = str(r.shape_id)
        return ((prop.endswith(".y") and sid in claimed_y)
                or (prop.endswith(".x") and sid in claimed_x))

    return [r for r in records if not _superseded(r)]


def _heading_past_ceiling(top, height, floor, ceiling) -> list:
    """[("body ceiling", overshoot)] for a header that has grown through the
    strip the master reserves and into the body area.

    Measured from ABOVE only: the box has to start in the header (above the
    subtitle's floor) and end below the body's ceiling. A title deliberately set
    lower down - a section divider, a hero slide - starts below the floor and is
    not a header that overflowed, so it says nothing here.

    Reported through _heading_past_margin, which means reported and never acted
    on: a standfirst that runs long is a copy decision, and the master's own
    subtitle box already overhangs its floor by 2mm, so the floor itself is not
    a line worth measuring a box against. The ceiling is: cross it and the body
    has lost the line it starts on."""
    if not floor or not ceiling or ceiling <= floor:
        return []
    if top >= floor:
        return []
    over = (top + height) - ceiling
    return [("body ceiling", over)] if over > 0 else []


def _body_members(surface, slide_h, floor=None, ceiling=None):
    """(body, measurable) - what counts as this slide's body content, and which
    of it can be measured.

    ONE definition, shared by every rule that holds the body to a stated line.
    Three rules ask the question now (the band intrusion, the shortfall under
    the ceiling, the frame's leading edge) and they have to agree: a shape one
    of them calls furniture and another calls body would be moved by one and
    left by the other.

    Out by construction: headings (this tool never repositions a title),
    placeholders (their geometry is master_slide's business, measured against the
    layout rather than a guide), full-bleed and page-deep elements
    (qc.util.full_height_panel), and the bottom furniture strip.

    HEADER-ZONE shapes are out too, but only when a band is stated - a box that
    starts above the subtitle floor and stops inside the strip is descender
    slack, and counting it as body made its top, up in the header, the line the
    whole block was measured from (a 36mm move asked for on the client's own
    sample deck). Without a floor there is no header zone to define, so the
    exclusion cannot apply.

    Rotated shapes cannot be MEASURED - their stored box is not their rendered
    one - but they stay in `body` because they travel with a translation, which
    does not care which way a shape faces."""
    def _header_zone(it) -> bool:
        return (floor is not None and ceiling is not None
                and it["t"] < floor
                and it["t"] + it["h"] <= ceiling + BAND_SLACK_EMU)

    body = [it for it in surface
            if not it["head"] and not it["ph"] and not it["bleed"]
            and not _header_zone(it)
            and not full_height_panel(it["t"], it["h"], slide_h)
            and it["t"] < FURNITURE_STRIP * slide_h]
    return body, [it for it in body if not it["rot"]]


def _band_intrusion(s_idx, surface, floor, ceiling, slide_h):
    """Content standing in the strip the master keeps clear under its header.

    The gap between the subtitle's floor and the body's ceiling is white space
    by declaration - the master draws a guide on each side of it - and it is
    what makes the headers read as one line down a deck. A deck that fills it
    loses that line on the slides that do (client convention, 20/08/2026).

    Wider than the safe-zone rule in two deliberate ways. It counts PICTURES
    and graphics, not just text: the strip is about the page's rhythm, so a
    photo standing in it breaks exactly what a paragraph would, and on the deck
    that prompted this rule every intruder was a photo, which the text-only
    safe-zone check never looked at. And it is judged per SLIDE with ONE move
    for the whole body: pushing the intruders down by themselves would drop
    each photo onto its own caption.

    Out of the move by construction: headings (this tool never repositions a
    title), placeholders (their geometry is master_slide's business, measured
    against the layout rather than a guide), full-bleed and page-deep elements
    (qc.util.full_height_panel), the bottom furniture strip, and HEADER-ZONE
    shapes - a box that starts above the floor
    and stops inside the strip. That last exclusion is not a nicety: a text box
    carries descender slack below its last line, the master's own subtitle
    placeholder overhangs its floor by 2.4mm because of it, and counting such a
    box as body content made its top - up in the header - the line the whole
    block was measured from, which asked for a 36mm move on the client's own
    sample deck. A header that grows all the way THROUGH the strip is a
    different thing and has its own finding (_heading_past_ceiling).

    Rotated shapes cannot be MEASURED - their stored box is not their rendered
    one - but they do travel with the block, because a translation does not
    care which way a shape faces.

    Both guides are required. With only a ceiling stated there is no way to
    tell an eyebrow sitting legitimately in the header from body content that
    has crept up into it, and flagging the eyebrow would be noise on every
    slide."""
    if not floor or not ceiling or ceiling <= floor:
        return []

    body, measurable = _body_members(surface, slide_h, floor, ceiling)
    intruders = [it for it in measurable if it["t"] < ceiling - BAND_SLACK_EMU]
    if not intruders:
        return _band_shortfall(s_idx, body, measurable, ceiling, slide_h)
    dy = ceiling - min(it["t"] for it in intruders)
    lowest = max(it["t"] + it["h"] for it in body)
    # A block that cannot come down without leaving the canvas is reported and
    # left alone. The migration pass seats the block on the ceiling anyway
    # because it is rebuilding the slide onto the master and says so loudly;
    # here the deck is finished, and a tick should never push a designer's
    # content off the page.
    fits = lowest + dy <= slide_h
    ids = ",".join(str(it["shape"].shape_id) for it in body)
    arabic = any(it["arabic"] for it in body)
    highest = min(intruders, key=lambda it: it["t"])
    strip = f"{floor / 914400:.2f}in-{ceiling / 914400:.2f}in"
    deep = f"{dy / 36000:.1f}mm"
    if fits:
        msg = (f"{len(intruders)} element(s) stand {deep} into the strip the "
               f"master keeps clear under the header ({strip}); the fix moves "
               f"the whole body block down onto the body guide, so the "
               f"arrangement inside it does not change")
    else:
        msg = (f"{len(intruders)} element(s) stand {deep} into the strip the "
               f"master keeps clear under the header ({strip}), and the body "
               f"cannot drop that far without pushing content off the slide; "
               f"nothing is moved, the slide needs a rework")
    return [make_record(
        slide_index=s_idx, shape_id=highest["shape"].shape_id, shape_path=None,
        module=MODULE, issue_type="margin_alignment.body_band_intrusion",
        severity="error" if dy >= BAND_OBVIOUS_EMU else "warning",
        action="flagged", confidence="high",
        locator=f"band:{dy}:{ids}" if fits else None,
        property="spPr.xfrm.off.y",
        old_value=highest["t"], new_value=dy if fits else None,
        profile_rule_id="geometry.body_band_emu",
        arabic_flag=arabic,
        message=_geo_msg(msg, arabic),
    )]


# How far below the body guide a column has to start before the gap is a mistake
# rather than a composition. 6mm is about the smallest a designer reads as "not
# on the line" when the column beside it does sit on the line.
BAND_SHORTFALL_EMU = 216000        # 6mm
# A slide needs more than one column for this to mean anything: with a single
# column there is no line to be off, only a composition that begins where it
# begins.
BAND_MIN_COLUMNS = 2


def _x_columns(items):
    """Body content grouped into side-by-side columns by horizontal overlap.

    Columns are what the body guide is about. A second ROW legitimately starts
    lower than the first and shares its column's horizontal span, so grouping on
    x and judging only each column's TOPMOST shape is what separates "this
    column missed the line" from "this is the next row down"."""
    ordered = sorted(items, key=lambda it: it["l"])
    columns: list[list] = []
    edge = None
    for it in ordered:
        if columns and edge is not None and it["l"] <= edge:
            columns[-1].append(it)
            edge = max(edge, it["l"] + it["w"])
        else:
            columns.append([it])
            edge = it["l"] + it["w"]
    return columns


def _band_shortfall(s_idx, body, measurable, ceiling, slide_h):
    """Body columns starting BELOW the guide the master says the body starts on.

    The mirror of the intrusion above, and the half that was missing. The strip
    rule was written for content that had crept UP into the header's clear band,
    so it only ever asked shapes to come DOWN; a column that began well below
    the guide was silently clean (design lead, 24/08/2026: "the picture on the
    left and the elements on the right should both be aligned to the top of the
    presentation space but only the picture was aligned").

    Silence there is what produced the reported result, because the cluster
    rule is still running: with a photo at 2.30in and a row of cards at 2.55in,
    _edge_misaligned aligns the photo TO THE CARDS - they are the majority line -
    and the slide then passes with its whole body below where the master says the
    body begins. Aligned to each other, and to nothing the master states.

    PER COLUMN, not per block and not per shape. Per block preserves the very
    stagger being complained about: lifting photo-and-cards together by the
    photo's shortfall puts the photo on the guide and leaves the cards below it,
    which is exactly the screenshot. Per shape would flatten a deliberate
    stagger inside a column into a row. Per column is the designer's own
    unit - "these two columns both start at the top" - and it is the one the
    guide is drawn for.

    A column that already sits on the guide is the strongest evidence the guide
    is live on this slide, so the columns that missed it are reported as errors;
    with every column low it is a warning, because a body that begins low all
    the way across may be the composition.
    """
    if len(measurable) < 2:
        return []
    columns = _x_columns(measurable)
    if len(columns) < BAND_MIN_COLUMNS:
        return []
    tops = [min(it["t"] for it in col) for col in columns]
    anchored = any(abs(t - ceiling) <= BAND_SLACK_EMU for t in tops)

    # Columns that start on the SAME low line are one decision, not one each: a
    # row of three cards side by side is three columns geometrically and one
    # problem to a designer, and three tick boxes that all move by 0.65in is the
    # per-occurrence noise this module exists to avoid.
    groups: dict[int, list] = {}
    for col, top in zip(columns, tops):
        short = top - ceiling
        if short < BAND_SHORTFALL_EMU:
            continue
        groups.setdefault(short // BAND_SLACK_EMU, []).extend(col)

    records = []
    for members in groups.values():
        top = min(it["t"] for it in members)
        short = top - ceiling
        highest = min(members, key=lambda it: it["t"])
        # Everything in the group travels, so the arrangement inside it survives
        # the lift; _fix_body_band takes the delta and the member ids.
        ids = ",".join(str(it["shape"].shape_id) for it in members)
        arabic = any(it["arabic"] for it in members)
        n_cols = sum(1 for col, t in zip(columns, tops)
                     if t - ceiling >= BAND_SHORTFALL_EMU
                     and (t - ceiling) // BAND_SLACK_EMU
                     == short // BAND_SLACK_EMU)
        what = f"{n_cols} columns" if n_cols > 1 else "this column"
        records.append(make_record(
            slide_index=s_idx, shape_id=highest["shape"].shape_id,
            shape_path=None, module=MODULE,
            issue_type="margin_alignment.body_below_band",
            severity="error" if anchored else "warning",
            action="flagged", confidence="high" if anchored else "medium",
            locator=f"band:{-short}:{ids}",
            property="spPr.xfrm.off.y",
            old_value=top, new_value=ceiling,
            profile_rule_id="geometry.body_band_emu",
            arabic_flag=arabic,
            message=_geo_msg(
                f"{what} start{'' if n_cols > 1 else 's'} "
                f"{short / 36000:.1f}mm below the guide the master says the "
                f"body starts on ({ceiling / 914400:.2f}in)"
                + (", while another column on this slide sits on it"
                   if anchored else
                   ", and so does every other column on this slide")
                + f"; the fix lifts {len(members)} element(s) onto the guide "
                  f"together, so the arrangement inside them does not change",
                arabic),
        ))
    return records


def _heading_past_margin(s_idx, shape, breached, arabic):
    """A title or standfirst whose BOX crosses a margin.

    Flag-only by construction: no new_value, no locator, and the issue type is
    absent from qc.fixer.FIXABLE_ISSUES, so there is nothing for the UI to tick
    and nothing for the fix engine to apply. A heading that breaks the margin
    is a house-style question - some clients want the title to run wide, some
    want it held to the frame - and answering it by nudging the box is how a
    tool overrules a brand (design lead, 19/08/2026).

    None when the overshoot is inside HEADING_SLACK_EMU: the box is on the
    line, and so is the answer."""
    breached = [(side, over) for side, over in breached
                if over > HEADING_SLACK_EMU]
    if not breached:
        return None
    # A margin side is named as a margin; the body ceiling is a line of its own
    # and reads wrong called one ("past the body ceiling margin").
    sides = ", ".join(f"{side} margin" if side in
                      ("left", "top", "right", "bottom") else side
                      for side, _over in breached)
    worst = max(over for _side, over in breached)
    return make_record(
        slide_index=s_idx, shape_id=shape.shape_id, shape_path=None,
        module=MODULE, issue_type="margin_alignment.heading_past_margin",
        severity="warning", action="flagged", confidence="deterministic",
        property="spPr.xfrm.off",
        old_value=f"{shape.left}, {shape.top}, {shape.width}x{shape.height}",
        new_value=None,
        profile_rule_id="geometry.safe_zone_margins_emu",
        arabic_flag=arabic,
        message=_msg(
            f"the heading's text box runs {worst / 36000:.0f}mm past the "
            f"{sides} margin. Nothing is moved or resized: whether a heading "
            f"may break the margin is the client's house style, so ask them "
            f"before changing it", arabic),
    )


def _content_overflow(s_idx, placed, margins, slide_w, slide_h, head_ids=(),
                      rtl=False):
    """A LAYOUT too wide for the margin frame: many shapes cross it on both
    sides, and the designer's fix is one proportional select-all rescale
    anchored at the margin the content starts from - the left one, or the RIGHT
    one on an Arabic slide - never fifty per-shape nudges (ground truth: 92.8%
    on the rankings slide). Deliberately quiet otherwise: a
    single bleed sidebar or edge stamp is design intent, one stray footer
    is the off-canvas check's beat, and shapes parked off the slide (whose
    center is outside) are junk for manual cleanup, not scale drivers.

    Headings are out of the selection entirely, not merely spared the move: a
    title running wide must neither count as evidence that the LAYOUT is too
    wide, nor be scaled down with it. Otherwise the one element this tool is
    forbidden to touch would be the element that triggers the biggest change
    it can make."""
    ml = margins.get("left", 0)
    mr = margins.get("right", 0)
    mt = margins.get("top", 0)
    mb = margins.get("bottom", 0)
    def _bleed_panel(it):
        """Edge-anchored and spanning most of that axis: a full-height
        sidebar or full-width band bleeding off the canvas by design (the
        designer's own deliverable keeps a 94%-height sidebar at x=0)."""
        l_, t_, w_, h_ = it[1], it[2], it[3], it[4]
        touches_x = l_ <= MARGIN_SLACK_EMU or l_ + w_ >= slide_w - MARGIN_SLACK_EMU
        touches_y = t_ <= MARGIN_SLACK_EMU or t_ + h_ >= slide_h - MARGIN_SLACK_EMU
        return ((touches_x and h_ >= 0.85 * slide_h)
                or (touches_y and w_ >= 0.85 * slide_w))

    content = [it for it in placed
               if not getattr(it[0], "is_placeholder", False)
               and str(it[0].shape_id) not in head_ids
               and not _bleed_panel(it)
               and 0 <= it[1] + it[3] // 2 <= slide_w
               and 0 <= it[2] + it[4] // 2 <= slide_h]
    if len(content) < 3:
        return []
    # A breacher is a SUBSTANTIAL shape crossing the frame by a visible
    # amount: sub-centimeter artifacts, corner ornaments, and shapes a
    # couple of millimeters over the line say nothing about layout width.
    deep = 2 * MARGIN_SLACK_EMU
    breachers = [it for it in content
                 if it[3] >= 360000 and it[4] >= 180000
                 and (it[1] < ml - deep
                      or it[1] + it[3] > slide_w - mr + deep)]
    left_breach = any(it[1] < ml - deep for it in breachers)
    right_breach = any(it[1] + it[3] > slide_w - mr + deep
                       for it in breachers)
    if len(breachers) < MIN_BREACHERS or not (left_breach and right_breach):
        return []

    l = min(it[1] for it in content)
    t = min(it[2] for it in content)
    r = max(it[1] + it[3] for it in content)
    b = max(it[2] + it[4] for it in content)
    usable_w = slide_w - ml - mr
    if usable_w <= 0 or r <= l or b <= t:
        return []
    scale = usable_w / (r - l)   # width-driven, like the designer's move
    if scale >= SCALE_SKIP:
        return []
    permille = int(scale * 1000)  # floor: never overshoot the frame
    scale = permille / 1000.0
    arabic = any(it[5] for it in content)
    ids = ",".join(str(it[0].shape_id) for it in content)
    # The rescale is anchored on the margin the content STARTS from, which
    # mirrors under RTL: an Arabic layout shrunk against the left margin drifts
    # away from its reading edge (design lead, 20/08/2026).
    anchor_l = (slide_w - mr - int((r - l) * scale)) if rtl else ml
    anchor_t = max(0, min(t, slide_h - mb - int((b - t) * scale)))
    fixable = scale >= SCALE_MIN
    pct = f"{scale * 100:.1f}%"
    msg = (f"the layout is wider than the margin frame: {len(breachers)} "
           f"shapes cross it on both sides "
           f"({(r - l) / 36000:.0f}mm of content against a "
           f"{usable_w / 36000:.0f}mm frame). ")
    msg += (f"The fix rescales all {len(content)} content shapes to {pct} "
            f"in one proportional move anchored at the "
            f"{'right' if rtl else 'left'} margin, placeholders untouched"
            if fixable else
            f"Fitting would need {pct}, too drastic to automate; rework the "
            "layout by hand")
    return [make_record(
        slide_index=s_idx, shape_id=content[0][0].shape_id, shape_path=None,
        module=MODULE, issue_type="margin_alignment.content_overflow",
        severity="warning", action="flagged", confidence="high",
        locator=(f"scale:{permille}:{anchor_l},{anchor_t}"
                 f":{ids}") if fixable else None,
        property="spPr.xfrm",
        old_value=f"({l}, {t})-({r}, {b})",
        new_value=permille if fixable else None,
        profile_rule_id="geometry.safe_zone_margins_emu",
        arabic_flag=arabic,
        message=_geo_msg(msg, arabic),
    )]


_AXIS = {
    # key index into (shape, left, top, w, h, arabic); property axis suffix
    "left": (1, "spPr.xfrm.off.x"),
    "top": (2, "spPr.xfrm.off.y"),
}


def _largest_agreeing(vals: list, tol: int) -> list:
    """The largest subset of values that all sit within tol of one value."""
    best = []
    for v in vals:
        group = [u for u in vals if abs(u - v) <= tol]
        if len(group) > len(best):
            best = group
    return best


def _panel_row_misaligned(s_idx, pool, cont, edge_tol, scope=""):
    """Sibling PANELS (substantial containers, side by side, same height
    class) whose row is broken: one panel sits off the top (or bottom) line
    the others share. The fix moves the panel AND everything inside it -
    designers move panels as blocks, never as lone rectangles (ground
    truth, 20/07/2026: a 69-shape Consultants panel moved 6.2mm as one
    selection; the per-shape window never saw it)."""
    records = []
    where = f" (inside {scope})" if scope else ""
    panels = []
    for ci in sorted(set(cont.values())):
        item = pool[ci]
        if item[3] >= PANEL_MIN_W and item[4] >= PANEL_MIN_H:
            panels.append(item)
    if len(panels) < 3:
        return records

    for row in _clusters(panels, 2, PANEL_WINDOW):
        if len(row) < 3:
            continue
        heights = [it[4] for it in row]
        if max(heights) > PANEL_H_RATIO * min(heights):
            continue  # different height classes: not one row
        side_by_side = all(
            (min(a[1] + a[3], b[1] + b[3]) - max(a[1], b[1]))
            <= PANEL_X_OVERLAP * min(a[3], b[3])
            for i, a in enumerate(row) for b in row[i + 1:])
        if not side_by_side:
            continue

        tops = [it[2] for it in row]
        bottoms = [it[2] + it[4] for it in row]
        top_grp = _largest_agreeing(tops, edge_tol)
        bot_grp = _largest_agreeing(bottoms, edge_tol)
        by_bottom = len(bot_grp) > len(top_grp)
        group = bot_grp if by_bottom else top_grp
        if len(group) < 2 or len(group) == len(row):
            continue
        line = int(statistics.median(group))
        deviants = [it for it in row
                    if abs((it[2] + it[4] if by_bottom else it[2]) - line)
                    > edge_tol]
        if not deviants or len(deviants) > len(row) / 3:
            continue
        for item in deviants:
            shape, arabic = item[0], item[5]
            target = line - item[4] if by_bottom else line
            off = abs(item[2] - target)
            records.append(make_record(
                slide_index=s_idx, shape_id=shape.shape_id, shape_path=None,
                module=MODULE,
                issue_type="margin_alignment.panel_row_misaligned",
                severity="error", action="flagged", confidence="high",
                property="spPr.xfrm.off.y",
                old_value=item[2], new_value=target,
                profile_rule_id="geometry.alignment.edge_tolerance_emu",
                arabic_flag=arabic,
                message=_geo_msg(
                    f"panel sits {off} EMU off the {'bottom' if by_bottom else 'top'} "
                    f"line its {len(row) - len(deviants)} sibling panels "
                    f"share; the fix moves the panel and everything inside "
                    f"it as one block{where}", arabic),
            ))
    return records


def _edge_misaligned(s_idx, placed, axis, edge_tol, window, scope=""):
    records = []
    key, prop = _AXIS[axis]
    where = f" (inside {scope})" if scope else ""
    # Candidate cluster: >= 3 edges within the intent window of each other
    # (near-aligned intent) whose overall spread exceeds tolerance.
    for cluster in _clusters(placed, key, window):
        if len(cluster) < 3:
            continue
        edges = [item[key] for item in cluster]
        if max(edges) - min(edges) <= edge_tol:
            continue
        median = statistics.median(edges)
        deviants = [item for item in cluster
                    if abs(item[key] - median) > edge_tol]
        # Alignment intent is only inferable from a clear aligned majority.
        # In an organically ragged pool (photo grids, logo walls) half the
        # members "deviate" from an arbitrary median: ambiguous, emit nothing
        # (real-deck tuning, 14/07/2026).
        if len(deviants) > len(cluster) / 3:
            continue
        for item in deviants:
            shape, arabic = item[0], item[5]
            edge = item[key]
            off = abs(edge - int(median))
            # Evidence-based severity: one shape VISIBLY off a tight line-up
            # of several others is confidently wrong (error); a sub-visual
            # nudge or a small cluster stays a judgment call (warning).
            obvious = len(cluster) >= 4 and off >= OBVIOUS_EDGE_EMU
            records.append(make_record(
                slide_index=s_idx, shape_id=shape.shape_id, shape_path=None,
                module=MODULE, issue_type="margin_alignment.edge_misaligned",
                severity="error" if obvious else "warning",
                action="flagged",
                confidence="high" if obvious else "medium",
                property=prop,
                old_value=edge, new_value=int(median),
                profile_rule_id="geometry.alignment.edge_tolerance_emu",
                arabic_flag=arabic,
                message=_geo_msg(f"{axis} edge off cluster median by "
                             f"{off} EMU "
                             f"(median {int(median)} EMU){where}", arabic),
            ))
    return records


def _relative_edge_misaligned(s_idx, contained, axis, edge_tol, window,
                              scope=""):
    """Alignment for shapes living inside containers, in container-relative
    coordinates: an icon 0.2in from its card's left edge is compared to how
    its peers sit in THEIR (similar-size) cards. The snap target composes
    the shape with its own container, not with a global cluster."""
    records = []
    key, prop = _AXIS[axis]
    size_key = 3 if axis == "left" else 4  # container width / height
    where = f" (inside {scope})" if scope else ""
    if len(contained) < 3:
        return records

    # only compare across similar-size containers (same family of cards)
    triples = [(item, cont, item[key] - cont[key], cont[size_key])
               for item, cont in contained]
    size_span = max(1, int(0.08 * statistics.median(t[3] for t in triples)))
    for size_cohort in _clusters(triples, 3, size_span):
        if len(size_cohort) < 3:
            continue
        for cluster in _clusters(size_cohort, 2, window):
            if len(cluster) < 3:
                continue
            rels = [t[2] for t in cluster]
            if max(rels) - min(rels) <= edge_tol:
                continue
            median = statistics.median(rels)
            deviants = [t for t in cluster if abs(t[2] - median) > edge_tol]
            if len(deviants) > len(cluster) / 3:
                continue  # no clear shared inset to hold anyone to
            for item, cont, rel, _sz in deviants:
                shape, arabic = item[0], item[5]
                off = abs(rel - int(median))
                obvious = len(cluster) >= 4 and off >= OBVIOUS_EDGE_EMU
                records.append(make_record(
                    slide_index=s_idx, shape_id=shape.shape_id,
                    shape_path=None, module=MODULE,
                    issue_type="margin_alignment.edge_misaligned",
                    severity="error" if obvious else "warning",
                    action="flagged",
                    confidence="high" if obvious else "medium",
                    property=prop,
                    old_value=item[key],
                    new_value=int(cont[key] + median),
                    profile_rule_id="geometry.alignment.edge_tolerance_emu",
                    arabic_flag=arabic,
                    message=_geo_msg(
                        f"{axis} edge sits {off} EMU off the shared inset "
                        f"its {len(cluster) - len(deviants)} peers use inside "
                        f"their cards; the snap re-composes it with its own "
                        f"card{where}", arabic),
                ))
    return records


def _cohort_rhythm(s_idx, placed, direction, window, spacing_tol, scope=""):
    """Rhythm check over cohorts of SAME-SIZE shapes (an image column, a
    card row): when every gap matches except one, the tail of the line is
    lifted back onto the rhythm. The fix moves each tail shape WITH its
    satellites (the labels and underlines riding beside it), which is how
    the designer executes the same move. Stronger evidence than band-line
    spacing: same-size repetition IS the rhythm claim."""
    records = []
    band_key, run_key, size_key = (1, 2, 4) if direction == "col" else (2, 1, 3)
    where = f" (inside {scope})" if scope else ""
    label = "vertical" if direction == "col" else "horizontal"

    cohorts: dict = {}
    for item in placed:
        key = (round(item[3] / COHORT_BIN_EMU), round(item[4] / COHORT_BIN_EMU))
        cohorts.setdefault(key, []).append(item)

    candidates = []
    for members in cohorts.values():
        if len(members) < 3:
            continue
        if len(members) >= 6:
            # a photo/logo WALL: its spacing breathes organically and stays
            # advisory-only (see _wall_cohort_ids); rhythm lifts are for
            # small cohorts where the beat is unmistakable
            continue
        for line in _clusters(members, band_key, window):
            if len(line) < 3:
                continue
            line = sorted(line, key=lambda t: t[run_key])
            gaps = [line[i + 1][run_key]
                    - (line[i][run_key] + line[i][size_key])
                    for i in range(len(line) - 1)]
            if any(g <= 0 for g in gaps):
                continue
            median = statistics.median_low(gaps)
            if median <= 0:
                continue
            # A rhythm claim needs a TIGHT sequence: the normal gap smaller
            # than the items themselves (a 22mm image column breathing 3mm).
            # Sparse stacks (small labels far apart) claim nothing.
            if median >= min(it[size_key] for it in line):
                continue
            odd = [(i, g) for i, g in enumerate(gaps)
                   if abs(g - median) > spacing_tol]
            if len(odd) != 1:
                continue
            i, gap = odd[0]
            if gap >= LIFT_MAX_RATIO * median:
                continue  # a gap this large reads as a section break
            candidates.append(dict(line=line, tail=line[i + 1:],
                                   gap=gap, median=median))

    def _rides(b_item, a_item) -> bool:
        """b (smaller) rides a: b's center on the lift axis sits inside
        a's span, within a satellite gap on the cross axis (keep in sync
        with qc/fixer.py _lift_satellites)."""
        if direction == "col":
            center = b_item[2] + b_item[4] // 2
            in_span = a_item[2] <= center <= a_item[2] + a_item[4]
            gap_x = max(0, max(b_item[1], a_item[1])
                        - min(b_item[1] + b_item[3], a_item[1] + a_item[3]))
        else:
            center = b_item[1] + b_item[3] // 2
            in_span = a_item[1] <= center <= a_item[1] + a_item[3]
            gap_x = max(0, max(b_item[2], a_item[2])
                        - min(b_item[2] + b_item[4], a_item[2] + a_item[4]))
        return in_span and gap_x <= 360000

    def _area(cand):
        return statistics.median(it[3] * it[4] for it in cand["tail"])

    # A smaller cohort (labels) riding a bigger one (images) shares its
    # rhythm break; the big cohort's lift carries it as a satellite, so a
    # second lift would double-move it.
    keep = []
    for b in candidates:
        ridden = any(a is not b and _area(a) > _area(b)
                     and all(any(_rides(bi, ai) for ai in a["tail"])
                             for bi in b["tail"])
                     for a in candidates)
        if not ridden:
            keep.append(b)

    for cand in keep:
        line, tail = cand["line"], cand["tail"]
        gap, median = cand["gap"], cand["median"]
        tail_ids = ",".join(str(m[0].shape_id) for m in tail)
        anchor, arabic = tail[0][0], any(m[5] for m in line)
        obvious = len(line) >= 4 and gap >= LIFT_OBVIOUS_RATIO * median
        records.append(make_record(
            slide_index=s_idx, shape_id=anchor.shape_id, shape_path=None,
            module=MODULE, issue_type="margin_alignment.cluster_rhythm",
            severity="error" if obvious else "warning",
            action="flagged",
            confidence="high" if obvious else "medium",
            locator=f"lift-{direction}:{tail_ids}",
            property="spPr.xfrm.off",
            old_value=gap, new_value=median,
            profile_rule_id="geometry.alignment.spacing_tolerance_emu",
            arabic_flag=arabic,
            message=_geo_msg(
                f"a line of {len(line)} same-size items keeps a "
                f"{median} EMU {label} rhythm, but this one sits after a "
                f"{gap} EMU gap; the fix lifts it (and everything after "
                f"it) back onto the rhythm, labels riding along{where}",
                arabic),
        ))
    return records


def _uneven_spacing(s_idx, placed, direction, window, spacing_tol, scope=""):
    """direction 'row': shapes sharing a top band, gaps run left-to-right.
    direction 'col': shapes sharing a left band, gaps run top-to-bottom."""
    records = []
    band_key, run_key, size_key = (2, 1, 3) if direction == "row" else (1, 2, 4)
    where = f" (inside {scope})" if scope else ""
    label = ("horizontal" if direction == "row" else "vertical")

    for line in _clusters(placed, band_key, window):
        if len(line) < 3:
            continue
        line = sorted(line, key=lambda t: t[run_key])
        gaps = [line[i + 1][run_key] - (line[i][run_key] + line[i][size_key])
                for i in range(len(line) - 1)]
        # Overlapping shapes: not a spacing pattern we can reason about.
        if any(gap <= 0 for gap in gaps):
            continue
        # median_low anchors on an actual gap so a single outlier among the
        # neighbors reads as the odd one out rather than skewing the center.
        median = statistics.median_low(gaps)
        if median <= 0:
            continue
        # A gap >= 5x the median reads as an intentional visual break, which
        # makes the whole line ambiguous: emit nothing.
        if any(gap >= 5 * median for gap in gaps):
            continue
        odd = [(i, gap) for i, gap in enumerate(gaps)
               if abs(gap - median) > spacing_tol]
        if not odd:
            continue
        line_ids = ",".join(str(m[0].shape_id) for m in line)
        walls = _wall_cohort_ids(placed)
        in_wall = sum(1 for m in line if str(m[0].shape_id) in walls)
        if len(odd) == 1 and in_wall >= len(line) - 1:
            # photo/logo wall rows breathe organically; a computed nudge is
            # a guess there, so the finding is advisory only
            i, gap = odd[0]
            shape, arabic = line[i + 1][0], line[i + 1][5]
            records.append(make_record(
                slide_index=s_idx, shape_id=shape.shape_id, shape_path=None,
                module=MODULE, issue_type="margin_alignment.uneven_spacing",
                severity="warning", action="flagged", confidence="low",
                property="spPr.xfrm.off",
                old_value=gap, new_value=None,
                profile_rule_id="geometry.alignment.spacing_tolerance_emu",
                arabic_flag=arabic,
                message=_geo_msg(f"uneven {label} spacing in a grid of same-size "
                             f"items: gap of {gap} EMU vs median {median} "
                             f"EMU; adjust by eye if it bothers{where}",
                             arabic),
            ))
            continue
        if len(odd) == 1:
            # single odd gap: fixable by translating the tail of the line
            i, gap = odd[0]
            shape, arabic = line[i + 1][0], line[i + 1][5]
            records.append(make_record(
                slide_index=s_idx, shape_id=shape.shape_id, shape_path=None,
                module=MODULE, issue_type="margin_alignment.uneven_spacing",
                severity="warning", action="flagged", confidence="low",
                locator=f"{direction}:{line_ids}",
                property="spPr.xfrm.off",
                old_value=gap, new_value=median,
                profile_rule_id="geometry.alignment.spacing_tolerance_emu",
                arabic_flag=arabic,
                message=_geo_msg(f"uneven {label} spacing: gap of {gap} EMU before "
                             f"this shape vs median gap {median} EMU{where}",
                             arabic),
            ))
        else:
            # several odd gaps: the honest fix is a full even distribution
            # (first and last anchored, middles spread) - offered as a
            # tickable suggestion, never pre-selected
            shape, arabic = line[0][0], any(m[5] for m in line)
            gaps_txt = ", ".join(str(g) for g in gaps)
            equal_gap = sum(gaps) // len(gaps)
            records.append(make_record(
                slide_index=s_idx, shape_id=shape.shape_id, shape_path=None,
                module=MODULE, issue_type="margin_alignment.uneven_spacing",
                severity="warning", action="flagged", confidence="low",
                locator=f"dist-{direction}:{line_ids}",
                property="spPr.xfrm.off",
                old_value=gaps_txt, new_value=equal_gap,
                profile_rule_id="geometry.alignment.spacing_tolerance_emu",
                arabic_flag=arabic,
                message=_geo_msg(f"irregular {label} spacing across {len(line)} "
                             f"shapes (gaps {gaps_txt} EMU); the fix "
                             f"distributes them evenly (~{equal_gap} EMU "
                             f"gaps), first and last stay put{where}",
                             arabic),
            ))
    return records


def _anchor_of(shape):
    """The vertical anchor a text box STATES, or None when its anchor cannot
    honestly be compared with its neighbours'.

    Three exclusions, each because the comparison would otherwise be a guess:

    A PLACEHOLDER's anchor comes from the master, and the master's own vertical
    anchor is followed rather than overridden anywhere in this tool - a title
    the master hangs at the bottom of its box is a design decision
    (qc.migrate). Reading it here would mean reporting the master as a defect.

    A box with a:spAutoFit has a height that TRACKS its text, so the anchor
    draws nothing: the box hugs the copy either way, and every such row would
    report a difference nobody can see.

    An empty box has no text for an anchor to place."""
    if getattr(shape, "is_placeholder", False):
        return None
    if not getattr(shape, "has_text_frame", False):
        return None
    if not shape.text_frame.text.strip():
        return None
    bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
    if bodyPr is None or bodyPr.find(qn("a:spAutoFit")) is not None:
        return None
    # OOXML's default when a:bodyPr states no anchor.
    return bodyPr.get("anchor") or "t"


def _text_anchor_mismatch(s_idx, placed, window, scope=""):
    """Boxes that line up, holding text that does not.

    The one alignment defect the box-edge rules cannot see. A row of equal-size
    labels whose boxes share a top edge to the EMU still reads as ragged when
    one of them anchors its text to the middle or the bottom of the box: the
    boxes are aligned and the words are not (design lead's caution: align the
    text box AND the text inside it). Nothing here measures a glyph - the
    anchor is an attribute the deck states, so the disagreement is a fact and
    not a rendering estimate, which is what keeps this inside the module's
    standing rule about boxes versus text.

    Flag-only, and warning severity, for the same reason a heading past a
    margin is: a deliberately bottom-anchored caption under a row of top-
    anchored ones is a real design, so the tool states the difference and lets
    the designer settle it. Absent from qc.fixer.FIXABLE_ISSUES, so there is
    nothing to tick.

    Reported only where the row's intent is unambiguous - three or more boxes,
    one height class, and a clear majority anchor - the same evidence bar the
    edge-cluster rules use."""
    records = []
    where = f" (inside {scope})" if scope else ""
    items = []
    for shape, left, top, width, height, arabic in placed:
        anchor = _anchor_of(shape)
        if anchor is not None:
            items.append((shape, left, top, width, height, arabic, anchor))

    for row in _clusters(items, 2, window):
        if len(row) < 3:
            continue
        heights = [it[4] for it in row]
        if min(heights) <= 0 or max(heights) / min(heights) > ANCHOR_H_RATIO:
            continue    # different-size boxes are not one line of labels
        counts = Counter(it[6] for it in row)
        if len(counts) < 2:
            continue
        majority, agreed = counts.most_common(1)[0]
        deviants = [it for it in row if it[6] != majority]
        if len(deviants) > len(row) / 3:
            continue    # no clear intent to deviate FROM
        for shape, _l, _t, _w, _h, arabic, anchor in deviants:
            records.append(make_record(
                slide_index=s_idx, shape_id=shape.shape_id, shape_path=None,
                module=MODULE,
                issue_type="margin_alignment.text_anchor_mismatch",
                severity="warning", action="flagged", confidence="medium",
                property="bodyPr.anchor",
                old_value=anchor, new_value=majority,
                profile_rule_id="geometry.alignment.edge_tolerance_emu",
                arabic_flag=arabic,
                message=_geo_msg(
                    f"this box lines up with {agreed} others on the same row, "
                    f"but its text sits at the "
                    f"{_ANCHOR_LABEL.get(anchor, anchor)} of the box while "
                    f"theirs sits at the "
                    f"{_ANCHOR_LABEL.get(majority, majority)}, so the words do "
                    f"not line up even though the boxes do. Nothing is changed: "
                    f"a deliberately different anchor is a design decision"
                    f"{where}", arabic),
            ))
    return records


def _squeezed_text(s_idx, placed, scope=""):
    """A text box so narrow its text wraps a letter or two per line: reads
    as a vertical strip of characters (real-deck feedback)."""
    records = []
    where = f" (inside {scope})" if scope else ""
    for shape, _left, _top, width, height, arabic in placed:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if len(text) < 12:
            continue
        if width < SQUEEZED_MAX_WIDTH and height > SQUEEZED_MIN_ASPECT * width:
            records.append(make_record(
                slide_index=s_idx, shape_id=shape.shape_id, shape_path=None,
                module=MODULE, issue_type="margin_alignment.squeezed_text",
                severity="error", action="flagged", confidence="high",
                property="spPr.xfrm.ext",
                old_value=f"{width}x{height}", new_value=None,
                profile_rule_id="geometry.alignment.edge_tolerance_emu",
                arabic_flag=arabic,
                message=_msg("text box is squeezed into a narrow strip "
                             f"({width} EMU wide for {len(text)} characters); "
                             f"its text likely wraps a letter per line, widen "
                             f"or rewrap it{where}", arabic),
            ))
    return records


def _text_overlap(s_idx, placed, scope=""):
    """Two text-bearing shapes overlapping: near-certain collision in a
    consulting deck (real-deck feedback: broken boxes sitting on titles)."""
    records = []
    where = f" (inside {scope})" if scope else ""
    texty = [item for item in placed
             if getattr(item[0], "has_text_frame", False)
             and item[0].text_frame.text.strip()]
    capped = max(0, len(texty) - MAX_PAIRWISE_SHAPES)
    if capped:
        texty = texty[:MAX_PAIRWISE_SHAPES]
        records.append(make_record(
            slide_index=s_idx, shape_id="-", shape_path=None, module=MODULE,
            issue_type="margin_alignment.overlap_check_capped",
            severity="info", action="flagged", confidence="high",
            property=None, old_value=len(texty) + capped,
            new_value=MAX_PAIRWISE_SHAPES,
            profile_rule_id=None, arabic_flag=False,
            message=(f"{len(texty) + capped} text elements on this slide is "
                     f"past the {MAX_PAIRWISE_SHAPES} this check compares pair "
                     f"by pair, so {capped} were left out and the overlap list "
                     f"for this slide is incomplete{where}"),
        ))

    # Grid siblings: >= 4 same-size text shapes (photo labels, logo captions)
    # sit tightly and their padded boxes brush each other by design; a pair of
    # them overlapping is layout rhythm, not collision (real-deck tuning).
    def _size_key(item):
        return (round(item[3] / 25400), round(item[4] / 25400))  # ~2mm bins

    sizes: dict = {}
    for item in texty:
        sizes.setdefault(_size_key(item), []).append(item)
    grid_ids = {str(item[0].shape_id)
                for members in sizes.values() if len(members) >= 4
                for item in members}

    hits: dict = {}   # anchor shape_id -> (item, [other ids], arabic)
    for i in range(len(texty)):
        for j in range(i + 1, len(texty)):
            a, b = texty[i], texty[j]
            if (str(a[0].shape_id) in grid_ids
                    and str(b[0].shape_id) in grid_ids):
                continue
            ox = min(a[1] + a[3], b[1] + b[3]) - max(a[1], b[1])
            oy = min(a[2] + a[4], b[2] + b[4]) - max(a[2], b[2])
            if ox <= 0 or oy <= 0:
                continue
            areas = (a[3] * a[4], b[3] * b[4])
            if min(areas) <= 0:
                continue
            cover = (ox * oy) / min(areas)
            # substantial cover of EITHER shape: tightly packed label grids
            # brush against each other without covering anything
            if cover < OVERLAP_MIN_RATIO:
                continue
            # near-full containment is deliberate layering (a label placed
            # on a text-bearing card), not a collision
            if cover >= 0.9:
                continue
            anchor = b  # later in z-order: the one sitting on top
            entry = hits.setdefault(str(anchor[0].shape_id),
                                    (anchor, [], False))
            entry[1].append(str(a[0].shape_id))
            hits[str(anchor[0].shape_id)] = (entry[0], entry[1],
                                             entry[2] or a[5] or b[5])
    for _sid, (item, others, arabic) in hits.items():
        shape = item[0]
        records.append(make_record(
            slide_index=s_idx, shape_id=shape.shape_id, shape_path=None,
            module=MODULE, issue_type="margin_alignment.text_overlap",
            severity="error", action="flagged", confidence="high",
            property="spPr.xfrm.off",
            old_value=f"overlaps shapes {', '.join(others)}", new_value=None,
            profile_rule_id="geometry.alignment.edge_tolerance_emu",
            arabic_flag=arabic,
            message=_msg(f"text box covers {int(OVERLAP_MIN_RATIO * 100)}%+ "
                         f"of {len(others)} other text box(es) (shapes "
                         f"{', '.join(others)}); separate or merge "
                         f"them{where}", arabic),
        ))
    return records
