"""shape_size module: deck-wide autoshape size consistency (audit-only, v1).

Exact-match tier only: shapes are grouped into cohorts by their preset
geometry (a:prstGeom @prst). For each cohort large enough to be meaningful,
the dominant size is the per-axis median; members deviating beyond the
profile tolerance are flagged. Fuzzy cross-preset similarity is v2.

Also emits off_grid records when the profile column grid is enabled.
"""

from collections import defaultdict
from statistics import median

from pptx.enum.shapes import MSO_SHAPE_TYPE

from qc.records import make_record
from qc.util import iter_shapes_deep
from spike.ns import find

MODULE = "shape_size"


def _collect_autoshapes(ctx):
    """All autoshapes with a preset geometry and explicit size, deck-wide.

    Returns list of (slide_index, shape, shape_path, prst, width, height).
    Placeholder-inherited geometry can resolve to None even on real shapes,
    so both axes must be explicitly present.
    """
    collected = []
    for s_idx, slide in enumerate(ctx.prs.slides):
        for shape, shape_path in iter_shapes_deep(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                continue
            prst_geom = find(shape._element, "p:spPr/a:prstGeom")
            if prst_geom is None:
                continue
            prst = prst_geom.get("prst")
            if not prst:
                continue
            width, height = shape.width, shape.height
            if width is None or height is None:
                continue
            collected.append((s_idx, shape, shape_path, prst, int(width), int(height)))
    return collected


def _near_clusters(members, tolerance: int, near_ratio: float):
    """Split a same-preset cohort into clusters of NEAR-DUPLICATE sizes.

    Real-deck tuning: a deck legitimately uses the same preset at many
    intentional sizes (cards, stats, dividers). The fixable signal is
    "almost identical, probably meant to be identical", so only sizes
    within near_ratio of each other form a cluster; anything further apart
    is intentional variety and never compared."""
    def near(a, b):
        return (abs(a[4] - b[4]) <= max(tolerance, near_ratio * max(a[4], b[4]))
                and abs(a[5] - b[5]) <= max(tolerance, near_ratio * max(a[5], b[5])))

    clusters = []
    for m in sorted(members, key=lambda t: t[4] * t[5]):
        for cluster in clusters:
            if near(m, cluster[0]):
                cluster.append(m)
                break
        else:
            clusters.append([m])
    return clusters


def _detect_size_mismatch(ctx, shapes) -> list:
    cfg = ctx.profile.module_config(MODULE)
    tolerance = cfg.get("size_tolerance_emu", 9525)
    min_cohort = cfg.get("min_cohort_size", 3)
    near_ratio = cfg.get("near_miss_ratio", 0.08)

    cohorts = defaultdict(list)
    for entry in shapes:
        cohorts[entry[3]].append(entry)

    records = []
    for prst, cohort in cohorts.items():
        for members in _near_clusters(cohort, tolerance, near_ratio):
            if len(members) < min_cohort:
                continue
            dom_cx = int(round(median(m[4] for m in members)))
            dom_cy = int(round(median(m[5] for m in members)))
            for s_idx, shape, shape_path, _, width, height in members:
                if abs(width - dom_cx) <= tolerance and abs(height - dom_cy) <= tolerance:
                    continue
                records.append(make_record(
                    slide_index=s_idx,
                    shape_id=shape.shape_id,
                    shape_path=shape_path,
                    module=MODULE,
                    issue_type="shape_size.size_mismatch",
                    severity="warning",
                    action="flagged",
                    confidence="high",
                    property="spPr.xfrm.ext",
                    old_value=f"{width}x{height}",
                    new_value=f"{dom_cx}x{dom_cy}",
                    arabic_flag=ctx.shape_has_arabic(s_idx, shape.shape_id),
                    profile_rule_id="shape_size.size_tolerance_emu",
                    message=(f"Shape size {width}x{height} deviates from dominant "
                             f"{dom_cx}x{dom_cy} EMU for near-duplicate cluster "
                             f"of {len(members)} '{prst}' shapes."),
                ))
    return records


def _detect_off_grid(ctx, shapes) -> list:
    if not ctx.profile.get("geometry.grid.enabled", False):
        return []
    columns = ctx.profile.get("geometry.grid.columns", 12)
    margin_left = ctx.profile.get("geometry.safe_zone_margins_emu.left", 0)
    margin_right = ctx.profile.get("geometry.safe_zone_margins_emu.right", 0)
    tolerance = ctx.profile.module_config(MODULE).get("size_tolerance_emu", 9525)

    usable = ctx.prs.slide_width - margin_left - margin_right
    if usable <= 0 or not columns:
        return []
    col_width = usable / columns

    records = []
    for s_idx, shape, shape_path, prst, width, height in shapes:
        # Group children carry offsets in the group child coordinate space,
        # not slide space, so grid checks would be meaningless there.
        if shape_path is not None:
            continue
        left = shape.left
        if left is None:
            continue
        offset = left - margin_left
        nearest_col = round(offset / col_width)
        boundary = int(round(nearest_col * col_width))
        if abs(offset - boundary) <= tolerance:
            continue
        records.append(make_record(
            slide_index=s_idx,
            shape_id=shape.shape_id,
            shape_path=shape_path,
            module=MODULE,
            issue_type="shape_size.off_grid",
            severity="warning",
            action="flagged",
            confidence="low",
            property="spPr.xfrm.off.x",
            old_value=str(left),
            new_value=str(margin_left + boundary),
            arabic_flag=ctx.shape_has_arabic(s_idx, shape.shape_id),
            profile_rule_id="geometry.grid",
            message=(f"Shape left edge is {abs(offset - boundary)} EMU off the "
                     f"nearest of {columns} column boundaries."),
        ))
    return records


def detect(ctx) -> list:
    shapes = _collect_autoshapes(ctx)
    records = _detect_size_mismatch(ctx, shapes)
    records.extend(_detect_off_grid(ctx, shapes))
    return records
