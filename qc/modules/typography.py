"""typography module: the deck's text SYSTEM.

Ground truth (designer deliverable, 20/07/2026): of the designer's ~96
text changes, ~93 were rule-shaped - 40 labels uppercased as COMPLETE
sibling sets (never partial, always by literal retype), 11 size overrides
deleted where they merely restated the inherited size, and 35 sibling
rescales. Font families and colors are the font/color modules' beats;
this module owns case, redundant overrides, and sibling size agreement.

Detections:
  typography.case_inconsistent       a label whose sibling set (same-size
                                     label shapes) or the profile's label
                                     case convention says ALL-CAPS while it
                                     is not. Fix retypes the text upper
                                     (the designer's own move; the original
                                     text is preserved in old_value).
  typography.redundant_size_override run-level sz that equals the size the
                                     run would inherit anyway: deleting it
                                     is a visual no-op that stops the deck
                                     fighting its master.
  typography.size_inconsistent       one sibling label breaks the single
                                     font size the rest of its set shares.
"""

from pptx.oxml.ns import qn

from qc.records import make_record
from spike.resolver import resolve_run

MODULE = "typography"

MAX_LABEL_WORDS = 6
LABEL_MIN_SZ = 800          # OOXML sz units (pt * 100): labels live at
LABEL_MAX_SZ = 1600         # 8-16pt (measured on the ground-truth deck)
SIBLING_BIN_EMU = 72000     # ~2mm size bins define "the same kind of label"
MIN_SET = 3
_SENTENCE_CHARS = ".?!;"


def _run_sizes(shape) -> list:
    """Explicit run sz values (pt*100) across the shape, in order."""
    out = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            rpr = run._r.find(qn("a:rPr"))
            sz = rpr.get("sz") if rpr is not None else None
            if sz is not None:
                out.append(int(sz))
    return out


def label_text(shape) -> str | None:
    """The shape's text when it walks and talks like a LABEL: short,
    non-placeholder, no sentence punctuation, explicitly sized 8-16pt.
    (Structural signature measured on the ground-truth deck: 39/40 of the
    designer's uppercased shapes match it.)"""
    if getattr(shape, "is_placeholder", False):
        return None
    if not getattr(shape, "has_text_frame", False):
        return None
    text = shape.text_frame.text.strip()
    if not text or len(text.split()) > MAX_LABEL_WORDS:
        return None
    if any(ch in text for ch in _SENTENCE_CHARS):
        return None
    sizes = _run_sizes(shape)
    if not sizes or any(s < LABEL_MIN_SZ or s > LABEL_MAX_SZ for s in sizes):
        return None
    return text


def _is_upper(text: str) -> bool:
    return text == text.upper() and text != text.lower()


def _sibling_sets(labelled) -> list:
    sets: dict = {}
    for shape, text in labelled:
        key = (round(shape.width / SIBLING_BIN_EMU),
               round(shape.height / SIBLING_BIN_EMU))
        sets.setdefault(key, []).append((shape, text))
    return [members for members in sets.values() if len(members) >= MIN_SET]


def _case_records(ctx, s_idx, labelled, convention):
    records = []
    flagged = set()

    def emit(shape, text, why):
        if shape.shape_id in flagged:
            return
        flagged.add(shape.shape_id)
        records.append(make_record(
            slide_index=s_idx, shape_id=shape.shape_id, shape_path=None,
            module=MODULE, issue_type="typography.case_inconsistent",
            severity="warning", action="flagged", confidence="high",
            property="text", old_value=text, new_value=text.upper(),
            profile_rule_id="typography.label_case",
            message=f"label case breaks the convention ({why}); the fix "
                    "retypes it ALL-CAPS, the original text stays in the "
                    "record",
        ))

    for members in _sibling_sets(labelled):
        caps = [m for m in members if _is_upper(m[1])]
        lowers = [m for m in members if not _is_upper(m[1])]
        if caps and lowers and len(caps) >= len(lowers):
            for shape, text in lowers:
                emit(shape, text,
                     f"{len(caps)} of its {len(members)} sibling labels "
                     "are ALL-CAPS")
    if convention == "upper":
        for shape, text in labelled:
            if not _is_upper(text):
                emit(shape, text, "profile label case is 'upper'")
    return records


def _redundant_size_records(ctx, s_idx, slide):
    """Runs whose explicit sz restates what they would inherit anyway.
    Verified by resolving the run with the attribute lifted; anything whose
    inherited chain bottoms out at the hard default is left alone."""
    records = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        redundant = []
        for p_idx, para in enumerate(shape.text_frame.paragraphs):
            for r_idx, run in enumerate(para.runs):
                rpr = run._r.find(qn("a:rPr"))
                sz = rpr.get("sz") if rpr is not None else None
                if sz is None:
                    continue
                rpr.attrib.pop("sz")
                try:
                    eff = resolve_run(run, para, shape, slide, ctx.prs)
                finally:
                    rpr.set("sz", sz)
                if (eff.size_pt.source != "hard-default"
                        and abs(eff.size_pt.value - int(sz) / 100.0) < 0.01):
                    redundant.append(f"p{p_idx}/r{r_idx}")
        if redundant:
            arabic = ctx.shape_has_arabic(s_idx, shape.shape_id)
            records.append(make_record(
                slide_index=s_idx, shape_id=shape.shape_id, shape_path=None,
                module=MODULE,
                issue_type="typography.redundant_size_override",
                severity="warning", action="flagged",
                confidence="deterministic",
                property="rPr.sz", locator=",".join(redundant),
                old_value=f"{len(redundant)} run(s)", new_value="inherit",
                arabic_flag=arabic,
                profile_rule_id="typography.redundant_size_override",
                message=f"{len(redundant)} run(s) restate the exact size "
                        "they inherit from the layout/master; removing the "
                        "override changes nothing visually and stops the "
                        "deck fighting its template",
            ))
    return records


def _size_records(ctx, s_idx, labelled):
    """One sibling label breaking the single size its set shares."""
    records = []
    for members in _sibling_sets(labelled):
        sized = []
        for shape, text in members:
            sizes = set(_run_sizes(shape))
            if len(sizes) == 1:
                sized.append((shape, text, sizes.pop()))
        if len(sized) < MIN_SET:
            continue
        from collections import Counter

        counts = Counter(s for _sh, _t, s in sized)
        majority, n = counts.most_common(1)[0]
        if n < 2 * len(sized) / 3:
            continue
        for shape, text, size in sized:
            if size == majority:
                continue
            records.append(make_record(
                slide_index=s_idx, shape_id=shape.shape_id, shape_path=None,
                module=MODULE, issue_type="typography.size_inconsistent",
                severity="warning", action="flagged", confidence="medium",
                property="rPr.sz", old_value=size, new_value=majority,
                arabic_flag=ctx.shape_has_arabic(s_idx, shape.shape_id),
                profile_rule_id="typography.size_tolerance",
                message=f"label is {size / 100:g}pt while {n} of its "
                        f"{len(sized)} sibling labels share "
                        f"{majority / 100:g}pt",
            ))
    return records


def detect(ctx):
    convention = ctx.profile.get("typography.label_case")
    records = []
    for s_idx, slide in enumerate(ctx.prs.slides):
        labelled = []
        for shape in slide.shapes:
            if ctx.shape_has_arabic(s_idx, shape.shape_id):
                continue  # Arabic has no letter case; sizes stay manual
            text = label_text(shape)
            if text is not None:
                labelled.append((shape, text))
        records.extend(_case_records(ctx, s_idx, labelled, convention))
        records.extend(_size_records(ctx, s_idx, labelled))
        records.extend(_redundant_size_records(ctx, s_idx, slide))
    return records


def learn_label_case(prs) -> str | None:
    """Bootstrap: the deck's own label-case convention, when it clearly has
    one (>=60% of at least 10 label-class shapes are ALL-CAPS)."""
    total = caps = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            text = label_text(shape)
            if text is not None:
                total += 1
                if _is_upper(text):
                    caps += 1
    if total >= 10 and caps / total >= 0.6:
        return "upper"
    return None
