"""Apply the master, then audit what it left.

Applying a master and auditing a deck were two doors, and a designer holding a
client deck needed both of them in that order, every time. The order is not a
preference either: half of what an audit finds on a raw client deck is about to
be rewritten by the master, and half of what the master leaves behind is
invisible until it has been applied. Auditing first reports defects that the
rebuild deletes; auditing second reports the deck a designer is actually going
to send (design lead, 27/08/2026).

Two calls, with a decision between them:

    plan()      reads the master, works out which of its layouts every slide
                belongs on, and collects the slides it could not place with
                confidence (qc.applymaster.plan_assignments, qc.layoutpick).
                Nothing is rebuilt. Cheap, deterministic, no PowerPoint.

    -- the designer picks a layout for each of those slides --

    run()       rebuilds every slide on its layout, writes the presentation
                space, migrates the content, then audits the result
                (qc.applymaster, qc.pspace, qc.migrate, qc.engine).

THE PAUSE IS THE POINT (design lead, 31/08/2026). This used to be one call, and
the slides nothing matched were placed mid-run by a vision model. That put a
network round trip inside the one operation that rewrites a designer's file, and
it answered a question the designer was better placed to answer anyway - they
can see the slide is a two-column comparison and that the master has a layout
for exactly that. They just had nowhere to say so. Now they do, and applying a
master is end-to-end deterministic: same deck, same master, same picks, same
file.

NOTHING HERE RENDERS HTML AND NOTHING HERE HOLDS STATE. It takes bytes and
returns a Plan or a Prep, so the pipeline can be run and asserted without a web
server, and so the routes above it are bookkeeping rather than a second
implementation.

WHAT FAILS, AND WHAT THAT COSTS. A step that produces the deliverable raises;
everything else degrades to a sentence. Losing the coverage report means the
rebuilt deck still stands, and losing the audit means the rebuilt deck is still
a rebuilt deck. A pass that took the whole run down with it would be the worst
of both: no file, and no report either.
"""

import io
import os
import tempfile
from dataclasses import dataclass, field
from pathlib import Path


class PrepError(RuntimeError):
    """The run cannot produce a deck, and the designer has to do something
    about it.

    Carries the HTTP status the route should answer with, because "that deck
    has no slides" and "PowerPoint is not available on this host" are a bad
    request and a missing dependency, and a caller that flattened both to 400
    would be telling a designer to fix a file that is fine.
    """

    def __init__(self, message: str, status: int = 400):
        super().__init__(message)
        self.status = status


@dataclass
class Prep:
    """Everything one run produced, in the order the page reads it."""

    filename: str
    source: bytes                         # the upload, exactly as it arrived
    deck: bytes | None = None             # after the master and the migration

    # --- the master half
    plans: list = field(default_factory=list)
    errors: dict = field(default_factory=dict)
    applied: int = 0
    masters: int = 1
    stragglers: list = field(default_factory=list)
    space_notes: list = field(default_factory=list)
    changes: list = field(default_factory=list)
    match_note: str = ""

    # --- what the master could not build
    coverage: object | None = None
    suggestions: list = field(default_factory=list)
    pictures: dict = field(default_factory=dict)
    suggest_note: str = ""

    # --- what is left on the slides
    manifest: dict | None = None
    audit_note: str = ""

    @property
    def slides(self) -> int:
        return len(self.plans)

    @property
    def findings(self) -> int:
        """Open audit findings on the rebuilt deck, ignoring the engine's own
        preflight rows - those are notes about the run, not defects a designer
        fixes."""
        if not self.manifest:
            return 0
        return sum(1 for r in self.manifest.get("records") or []
                   if r.get("module") != "preflight")


# --- reading the master ---------------------------------------------------


def read_master(master_bytes: bytes):
    """(presentation, its dominant master, its layouts, its frame, its size).

    The layouts come out of the FILE rather than out of an archived spec: the
    master is the source of truth and it cannot drift from itself, where a spec
    read last month can (qc.web._formattable_profiles says the same thing about
    a stored master).
    """
    from pptx import Presentation

    from .stylespec import dominant_master, extract_layouts, infer_grid

    try:
        prs = Presentation(io.BytesIO(master_bytes))
    except Exception as exc:
        raise PrepError(f"Could not read that master: {type(exc).__name__}: "
                        f"{exc}", 422) from exc
    target = dominant_master(prs)
    if target is None:
        raise PrepError("That master file has no slide master, so there is "
                        "nothing to apply.")
    layouts = extract_layouts(target, embed_assets=False)
    space = (infer_grid(prs, target) or {}).get("presentation_space")
    return prs, target, layouts, space, (prs.slide_width, prs.slide_height)


# --- the layout decision --------------------------------------------------


@dataclass
class Plan:
    """A deck read against a master, with nothing rebuilt yet.

    Held between the two halves of the run: the upload produces one of these,
    the designer answers the choices on it, and run() consumes it. It carries
    the SOURCE BYTES rather than an open Presentation because a python-pptx
    object cannot outlive the request that opened it, and re-reading a file
    already in memory costs nothing next to the rebuild that follows.
    """

    filename: str
    source: bytes
    layouts: list = field(default_factory=list)
    plans: list = field(default_factory=list)        # SlidePlan, all slides
    choices: list = field(default_factory=list)      # Choice, EVERY slide
    space: object | None = None
    master_size: tuple = ()

    @property
    def slides(self) -> int:
        return len(self.plans)

    @property
    def undecided(self) -> int:
        """Slides that are a QUESTION, not slides on the page.

        Those were the same number until the layout step started showing the
        whole deck (qc.layoutpick.choices, 02/09/2026). Counting the page would
        report every matched slide as a decision nobody made, and the run's own
        note and the coverage report are both written off this."""
        from .layoutpick import undecided as count_undecided

        return count_undecided(self.choices)


def plan(data: bytes, filename: str, master_bytes: bytes) -> Plan:
    """Work out where every slide would go, and which of those are guesses.

    The cheap half of the run and the only half a designer waits on before
    seeing anything: no PowerPoint, no renderer, no network. It reads two files
    that are already in memory and does arithmetic over them, so a deck that
    turns out to be wrong is rejected in a moment rather than after a rebuild.
    """
    from pptx import Presentation

    from .applymaster import plan_assignments
    from .layoutpick import choices as layout_choices

    _, _, layouts, space, master_size = read_master(master_bytes)
    try:
        deck_prs = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise PrepError(f"Could not read that deck: {type(exc).__name__}: "
                        f"{exc}", 422) from exc

    plans = plan_assignments(deck_prs, layouts)
    if not plans:
        raise PrepError("That deck has no slides to work on.")

    # Degrades to "nothing to decide" rather than raising. A deck whose layout
    # choices cannot be computed is still a deck that rebuilds on the targets
    # plan_assignments picked, which is exactly what happened before this step
    # existed - so the cost of losing it is the decision, not the deliverable.
    try:
        choices = layout_choices(deck_prs, layouts, plans)
    except Exception:
        choices = []

    return Plan(filename=filename, source=data, layouts=layouts, plans=plans,
                choices=choices, space=space, master_size=master_size)


# --- asked for, not run: what the master is missing -----------------------


def propose_layouts(prep: Prep, master_bytes: bytes) -> Prep:
    """A layout to add per gap, drawn on the master's own frame. In place.

    The prescription half. "Eleven slides fell back" is a sentence a designer
    has to interpret; "add a two-column comparison layout, and here is where its
    boxes land on your margins" is one they can act on (qc.layoutsuggest).

    A BUTTON, NOT A STEP (31/08/2026). This ran inside the rebuild until the
    model came out of the master application. It is the last pass in the prepare
    flow that asks a model anything, and leaving it in the run would have meant
    the deterministic half still waiting on a network call to finish - so it
    moved to the results page, where a designer presses it when the gaps turn
    out to be worth acting on. Most runs never need it.

    Nothing is built and nothing is written to the client's master: a layout
    carries type styles, guides and brand furniture that are a designer's to
    add. This draws a picture of one.
    """
    from pptx import Presentation

    coverage = prep.coverage
    if coverage is None:
        prep.suggest_note = ("There is no coverage report for this run, so "
                             "there are no gaps to propose against.")
        return prep
    wanted = (list(coverage.gaps or [])
              + list(getattr(coverage, "misfit_clusters", None) or []))
    if not wanted:
        prep.suggest_note = ("Every slide found a layout in this master, so "
                             "there is nothing to propose.")
        return prep

    try:
        from .layoutsuggest import suggest, wireframe

        _, target, layouts, space, _ = read_master(master_bytes)
        deck_prs = Presentation(io.BytesIO(prep.deck or prep.source))
        suggestions, asked, unreachable = suggest(coverage, deck_prs, layouts,
                                                  space, target)
        pictures = {i: wireframe(s, deck_prs)
                    for i, s in enumerate(suggestions)}
    except Exception as exc:
        prep.suggest_note = (f"The layouts could not be proposed "
                             f"({type(exc).__name__}). The gaps are read from "
                             f"the plans and are unaffected.")
        return prep

    if unreachable and not suggestions:
        # NOT the same sentence as a rejected proposal. Nothing was proposed
        # because nothing was asked, and a designer told their proposal was
        # discarded would go looking at the gaps for a fault that is not there.
        prep.suggest_note = (f"The layouts could not be proposed: the model "
                             f"could not be reached ({unreachable}). The gaps "
                             f"are read from the plans and are unaffected.")
        return prep
    if asked and not suggestions:
        prep.suggest_note = ("No layout could be proposed for these groups. "
                             "The gaps stand on their own; a proposal that did "
                             "not answer the group it was asked about is "
                             "discarded rather than shown.")
        return prep

    prep.suggestions, prep.pictures, prep.suggest_note = suggestions, pictures, ""
    return prep


# --- the rebuild ----------------------------------------------------------


def build(prepared: Plan, master_bytes: bytes, *, match_note: str = "") -> Prep:
    """Roll up the gaps and rebuild. Returns a Prep with `deck` set, or raises
    PrepError.

    Takes the plans as DECIDED rather than re-deriving them, so the layouts a
    designer approved on the previous page are the layouts that get applied.
    Re-planning here would silently discard every pick.

    The gaps are computed from the plans that are ABOUT TO BE APPLIED, before
    the rebuild rather than after it, so the report and the file cannot
    disagree about which slide went where.
    """
    from pptx import Presentation

    from .applymaster import apply_master

    data, layouts, plans = prepared.source, prepared.layouts, prepared.plans
    space, master_size = prepared.space, prepared.master_size
    prep = Prep(filename=prepared.filename, source=data, plans=plans,
                match_note=match_note)

    # A failure here costs the panel it fills, never the run: a designer with a
    # rebuilt deck and no coverage report is better off than one with neither.
    try:
        from .layoutgap import report as layout_coverage

        deck_prs = Presentation(io.BytesIO(data))
        prep.coverage = layout_coverage(deck_prs, layouts, plans)
    except Exception:
        prep.coverage = None

    result = apply_master(data, master_bytes, plans)
    if result.fatal:
        raise PrepError(result.fatal, 503)
    prep.plans = result.plans
    prep.errors = result.errors
    prep.applied = result.applied
    prep.masters = result.masters
    prep.stragglers = result.stragglers

    # The frame goes in BEFORE the content moves, so the migration seats every
    # slide's body on the frame it reads from that slide's own master rather
    # than on whatever the original design implied (qc.pspace).
    from .pspace import ensure_presentation_space

    try:
        deck, prep.space_notes = ensure_presentation_space(
            result.deck,
            fallback_box=(space or {}).get("box_emu")
            if not (space or {}).get("problem") else None,
            fallback_size=master_size)
    except Exception as exc:
        deck = result.deck
        prep.space_notes = [f"The presentation space could not be written into "
                            f"the deck ({type(exc).__name__}: {exc}). The "
                            f"layouts were still applied; the marker is missing "
                            f"from the masters."]

    # Applying the layout is only half of it. PowerPoint remaps placeholder
    # content, but a deck of free-floating shapes keeps them exactly where they
    # were, leaving the master's empty placeholders on top of the old content.
    from .migrate import migrate_deck

    try:
        prep.deck, prep.changes = migrate_deck(deck)
    except Exception as exc:
        prep.deck = deck
        prep.changes = [type("C", (), {
            "slide_index": 0, "action": "migration skipped",
            "detail": f"{type(exc).__name__}: {exc}; layouts applied, content "
                      f"left where it was"})()]
    return prep


# --- step 5: the audit ----------------------------------------------------


def audit(prep: Prep, profile_obj) -> Prep:
    """Read the REBUILT deck against the profile, in place on the same Prep.

    The rebuilt one, and that is the whole reason this pass sits after the
    other four: an audit of the upload reports margins the master is about to
    reset and fonts it is about to replace, and a designer working through that
    list is fixing a file that no longer exists.

    A failure leaves `manifest` None and says why. The deck is already built by
    then, so downloading it is unaffected.
    """
    from .engine import run_audit

    if prep.deck is None:
        prep.audit_note = ("There is no rebuilt deck to audit, so the slides "
                           "were not read.")
        return prep

    fd, name = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)          # Windows cannot delete a file while this is open
    tmp = Path(name)
    try:
        tmp.write_bytes(prep.deck)
        result = run_audit(tmp, profile_obj)
    except Exception as exc:
        prep.audit_note = (f"The rebuilt deck could not be audited "
                           f"({type(exc).__name__}: {exc}), so the slide "
                           f"findings below are missing. The rebuild itself is "
                           f"unaffected and the file is ready to download.")
        return prep
    finally:
        tmp.unlink(missing_ok=True)   # PRD: uploads auto-delete after processing

    manifest = result.to_manifest()
    manifest["deck"] = prep.filename   # a temp path is meaningless to a designer
    prep.manifest = manifest
    return prep


def run(prepared: Plan, master_bytes: bytes, profile_obj, *,
        match_note: str = "") -> Prep:
    """The second half: rebuild the deck on the decided layouts, then audit it.

    `prepared` comes from plan() with the designer's picks already written onto
    its plans (qc.layoutpick.apply_picks). Calling this with an untouched Plan
    is legitimate and means "use every suggestion as it stands"."""
    return audit(build(prepared, master_bytes, match_note=match_note),
                 profile_obj)


# --- what the page leads with ---------------------------------------------


def headline(prep: Prep) -> str:
    """One sentence saying what happened and what is left, because the numbers
    a designer wants are not on the same list.

    Read off the Prep rather than passed in, so the sentence at the top of the
    page and the cards underneath it cannot become different claims."""
    bits = [f"Rebuilt {prep.applied} of {prep.slides} slide"
            f"{'s' if prep.slides != 1 else ''} on the master"]
    if prep.errors:
        bits.append(f"{len(prep.errors)} could not be rebuilt and were left "
                    f"exactly as they were")
    cov = prep.coverage
    unplaced = getattr(cov, "unplaced", 0) if cov is not None else 0
    if unplaced:
        n_gaps = len(getattr(cov, "gaps", None) or [])
        bits.append(f"{unplaced} had no layout in it, which is {n_gaps} thing"
                    f"{'s' if n_gaps != 1 else ''} the master is missing")
    if prep.manifest is None:
        bits.append("and the audit did not run")
    elif prep.findings:
        bits.append(f"and the rebuilt deck has {prep.findings} finding"
                    f"{'s' if prep.findings != 1 else ''} left on it")
    else:
        bits.append("and the rebuilt deck is clean against the profile")
    return ". ".join(bits) + "."
