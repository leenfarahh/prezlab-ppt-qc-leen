"""Put the presentation-space marker into a formatted deck.

Reading the frame is Stage 1's job (qc.stylespec.read_presentation_space). This
is the other half: after a master has been applied, the OUTPUT deck has to carry
the marker itself, on every slide master it ended up with.

Why it is not already there. Applying a master copies the design, so a marker
the designer drew travels with it - but only to where they drew it. On the
client master it lives on ONE layout ("USE THIS"), so a slide sitting on any
other layout inherits no frame at all, and a slide that could not be rebuilt
keeps the deck's original design alive, which never had one (see
qc.applymaster.ApplyResult.stragglers). A designer opening that deck and
reaching for ToolsToo's align-to-presentation-space finds nothing to align to,
on exactly the slides that most need it.

The marker goes on the slide MASTER. That is where a frame governing the whole
deck belongs - every slide inherits from it whatever layout it is on - and it is
the first place the read looks, so a re-audit of the output finds the same frame
the format was seated against rather than a weaker guess from the guides.

The geometry is read from the OUTPUT deck wherever the applied design already
states it, not from the submitted master file. PowerPoint resizes a loaded
design to the deck's slide size, so the file it just wrote is the only place
where the frame's real numbers are; taking them from the master would plant a
16:9 rectangle on a 4:3 deck and call it the frame. Only when the output states
none anywhere does the submitted master's box get used, scaled, and the scaling
reported.

The rectangle written is invisible by construction - no fill, no line - because
a marker that prints appears on every slide of the delivered deck.
"""

import io

from pptx import Presentation
from pptx.oxml.ns import qn

from .stylespec import _space_from, read_presentation_space

# What ToolsToo stamps, and therefore what this writes: the add-in reads its
# presentation space off the alt text, so a marker carrying only a matching
# NAME would be found by this tool and not by the designer's own.
PS_ALT = "ToolsToo_PS"
PS_NAME = "Presentation space"


def _insert_marker(container, box_emu) -> None:
    """Draw the marker rectangle into a master or layout.

    python-pptx cannot add a shape to a master, so the element goes in
    directly. Placement before p:extLst is what keeps the part schema-valid:
    spTree children are ordered, and PowerPoint repairs a file that gets this
    wrong rather than opening it."""
    from pptx.oxml.shapes.autoshape import CT_Shape

    left, top, right, bottom = box_emu
    ids = [int(el.get("id")) for el in
           container.shapes._spTree.iter(qn("p:cNvPr"))
           if (el.get("id") or "").isdigit()]
    sp = CT_Shape.new_autoshape_sp(max(ids or [1]) + 1, PS_NAME, "rect",
                                   left, top, right - left, bottom - top)
    container.shapes._spTree.insert_element_before(sp, "p:extLst")
    shape = next(s for s in container.shapes if s._element is sp)
    # Invisible, and stated as such. A rectangle's fill and line otherwise come
    # from the shape STYLE, so "no fill" has to be written rather than left out.
    shape.fill.background()
    shape.line.fill.background()
    cNvPr = sp.find(qn("p:nvSpPr")).find(qn("p:cNvPr"))
    cNvPr.set("descr", PS_ALT)
    # Drawn by a person, not inherited: without this PowerPoint treats the
    # shape as part of the layout's own furniture in some views.
    sp.find(qn("p:nvSpPr")).find(qn("p:nvPr")).set("userDrawn", "1")


def _scale(box_emu, from_size, to_size):
    """A box read at one slide size, restated at another."""
    fw, fh = from_size
    tw, th = to_size
    if not fw or not fh:
        return box_emu
    return [int(round(box_emu[0] * tw / fw)), int(round(box_emu[1] * th / fh)),
            int(round(box_emu[2] * tw / fw)), int(round(box_emu[3] * th / fh))]


def frame_in(deck_bytes: bytes):
    """(box_emu, where) for the frame a deck already states, or (None, None).

    Every master is asked, not just the dominant one, because the marker's
    whole job here is to serve the masters that do not have it - including the
    original design a straggler slide kept alive."""
    prs = Presentation(io.BytesIO(deck_bytes))
    for master in prs.slide_masters:
        space = read_presentation_space(prs, master)
        if space and not space.get("problem"):
            return space["box_emu"], space.get("source")
    return None, None


def stamp_master(master_bytes: bytes, box_emu) -> tuple[bytes, str]:
    """A COPY of this master carrying the presentation-space marker.

    The other half of ensure_presentation_space. That one writes the marker into
    a deck this tool just built, where the frame is already known; this writes it
    into the MASTER, which is where the frame should have been stated in the
    first place and where stating it ends the guessing for every deck formatted
    against it afterwards.

    A COPY, and never in place. A client's master is not a file this tool edits
    on a hunch - the same rule qc.layoutsuggest keeps by proposing layouts rather
    than building them. What comes back is bytes for a designer to look at, open
    in PowerPoint, and decide to keep. Nothing is written to the template store
    by this function.

    Why it is safe to write when a whole LAYOUT is not: a presentation space is
    an invisible rectangle that states a decision the designer has already made.
    It carries no type styles, no brand furniture and no design opinion, it
    prints nothing, and deleting it returns the file exactly to where it was.

    Raises ValueError when the box is not a usable rectangle inside the slide,
    because a frame stamped outside the canvas would be read back as the frame
    and seat every deck against it.
    """
    prs = Presentation(io.BytesIO(master_bytes))
    sw, sh = prs.slide_width, prs.slide_height
    try:
        left, top, right, bottom = (int(v) for v in box_emu)
    except (TypeError, ValueError) as exc:
        raise ValueError("the presentation space needs four numbers: left, "
                         "top, right and bottom") from exc
    if right <= left or bottom <= top:
        raise ValueError("the presentation space has no area: its right edge "
                         "must be past its left, and its bottom past its top")
    if left < 0 or top < 0 or right > sw or bottom > sh:
        raise ValueError(
            f"the presentation space falls outside the slide "
            f"({sw / 914400:.2f}in by {sh / 914400:.2f}in). Every edge has to "
            f"sit on the canvas, or every deck formatted against this master "
            f"would be seated on a frame that is not on the page")

    masters = list(prs.slide_masters)
    already = [m for m in masters if _space_from(m, sw, sh) is not None]
    stamped = 0
    for master in masters:
        if _space_from(master, sw, sh) is not None:
            continue
        _insert_marker(master, (left, top, right, bottom))
        stamped += 1

    if not stamped:
        return master_bytes, (
            f"Every slide master in this file already states a presentation "
            f"space, so nothing was added. Delete the existing rectangle first "
            f"if you meant to replace it.")

    out = io.BytesIO()
    prs.save(out)
    note = (f"Presentation space stamped onto {stamped} slide master(s): an "
            f"invisible rectangle with the alt text {PS_ALT}, "
            f"{(right - left) / 914400:.2f}in by {(bottom - top) / 914400:.2f}in "
            f"at {left / 914400:.2f}in from the left and "
            f"{top / 914400:.2f}in from the top. Open it in PowerPoint to check "
            f"it, then use it in place of the original: every deck formatted "
            f"against it is seated on a stated frame rather than on an "
            f"inference, and ToolsToo can align to it.")
    if already:
        note += (f" {len(already)} master(s) already had one and were left "
                 f"alone.")
    return out.getvalue(), note


def ensure_presentation_space(deck_bytes: bytes, fallback_box=None,
                              fallback_size=None) -> tuple[bytes, list[str]]:
    """Give every slide master in the deck a presentation-space marker.

    Returns the deck and one note per thing a designer would want to know:
    what was added and where the numbers came from. A deck that already states
    the frame on every master is returned untouched with a note saying so, so
    "was the presentation space carried over?" is always answerable from the
    result page rather than by opening master view and hunting.

    `fallback_box`/`fallback_size` are the submitted master's frame and slide
    size, used only when the output states no frame anywhere."""
    prs = Presentation(io.BytesIO(deck_bytes))
    sw, sh = prs.slide_width, prs.slide_height
    notes: list[str] = []

    box, where = None, None
    for master in prs.slide_masters:
        space = read_presentation_space(prs, master)
        if space and not space.get("problem"):
            box, where = space["box_emu"], space.get("source")
            break

    if box is None:
        if not fallback_box:
            return deck_bytes, ["This master states no presentation space, so "
                                "none was added to the deck. Draw the rectangle "
                                "on the master and resubmit it to get one."]
        box = fallback_box
        where = "the master file"
        if fallback_size and tuple(fallback_size) != (sw, sh):
            box = _scale(box, fallback_size, (sw, sh))
            notes.append(
                f"The master is {fallback_size[0] / 914400:.2f}in wide and this "
                f"deck is {sw / 914400:.2f}in, so the presentation space was "
                f"scaled to fit. Check it against the deck before relying on it.")

    added = 0
    for master in prs.slide_masters:
        # Its OWN shapes, not what its layouts carry: a marker on one layout
        # serves that layout only, which is the gap this closes.
        if _space_from(master, sw, sh) is not None:
            continue
        _insert_marker(master, box)
        added += 1

    if not added:
        notes.append("Every slide master in this deck already carries the "
                     "presentation space; nothing was added.")
        return deck_bytes, notes

    out = io.BytesIO()
    prs.save(out)
    notes.insert(0, (
        f"Presentation space added to {added} slide master(s) as an invisible "
        f"rectangle with the alt text {PS_ALT}, "
        f"{(box[2] - box[0]) / 914400:.2f}in by {(box[3] - box[1]) / 914400:.2f}in "
        f"at {box[0] / 914400:.2f}in from the left and "
        f"{box[1] / 914400:.2f}in from the top, read from {where}. Every slide "
        f"now inherits it whatever layout it is on, and ToolsToo can align to "
        f"it."))
    return out.getvalue(), notes
