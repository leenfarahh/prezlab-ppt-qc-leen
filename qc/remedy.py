"""Apply the remedy a designer picked, and record exactly how to take it back.

qc.design states the options; this performs the one that was chosen. The split
is not tidiness - it is what keeps the page honest. Detection can run on a deck
nobody has touched, as often as it likes, and produce the same findings; nothing
here runs unless a designer picked something.

Two rules govern everything in this file, and both are inherited from qc.undo
because a second undo mechanism with different guarantees would be worse than
none:

THE UNDO IS CAPTURED BEFORE THE CHANGE, AS STATE. A recolour stores the shape's
whole element as it stood, not "it used to be #1F3864"; a move stores the
coordinates the shape held, not the distance it travelled; a reorder stores the
index it sat at. Re-deriving any of those at undo time means computing the same
thing twice and hoping the two agree, and when they disagree the deck a designer
approved is quietly wrong.

AND EVERY OPERATION SPEAKS qc.undo's VOCABULARY. The ops written here are
replayed by qc.undo.apply_undo, the same function the format review page has
been using - so undo is not a feature of this page, it is the same machinery
reaching one step further. That is why a colour change stores a whole-element
`replace` rather than an inverse recolour: `replace` already exists, is already
exact, and is already tested.

python-pptx's own fill and font APIs do the writing wherever they can. Setting a
solid fill means removing whatever fill was there and inserting a:solidFill at
the one position the schema allows, and hand-rolled XML that gets that wrong
produces a file PowerPoint offers to repair rather than open.
"""

import io
from dataclasses import dataclass, field

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu

from spike.ns import find

from .design import parse_hex
from .util import iter_shapes_deep


@dataclass
class Applied:
    """One remedy, performed. `undo` is what qc.undo replays to take it back."""
    finding_id: str
    remedy_id: str
    kind: str
    headline: str
    label: str
    done: bool
    detail: str
    # The slides this decision was about, copied off the finding. Kept because
    # the page shows a decision on the slide it belongs to, and a "leave it"
    # touches no shape at all - so `touched` cannot answer "which slide was
    # this?" for exactly the decisions a designer is most likely to revisit.
    slides: list[int] = field(default_factory=list)
    undo: list[dict] = field(default_factory=list)
    # (slide_index, shape_id) for everything this remedy wrote to. Two remedies
    # that touched one shape cannot be undone independently - the first one's
    # stored element predates the second one's change, so replaying it alone
    # would silently erase work the designer also approved. The page uses this
    # to say so, and to take both back together (see followers).
    touched: list[tuple] = field(default_factory=list)


def _by_id(slide) -> dict:
    return {str(s.shape_id): s for s, _p in iter_shapes_deep(slide.shapes)}


def _xml(shape) -> str | None:
    try:
        return etree.tostring(shape._element, encoding="unicode")
    except Exception:
        return None


def _replace_undo(shape, label: str) -> dict | None:
    """The undo for ANY change inside one shape: the element as it stands.

    One op covers recolouring a fill, a line, one run or every run, and it
    covers them exactly. The alternative - an inverse op per surface - is four
    code paths that each have to know how to put back a colour that may have
    been inherited rather than stated, and "it had no colour of its own" is the
    case they would all get wrong."""
    blob = _xml(shape)
    if blob is None:
        return None
    return {"op": "replace", "shape_id": str(shape.shape_id), "xml": blob,
            "label": label}


def _top_level(shape, spTree):
    """The element in the shape tree whose subtree holds this shape - the thing
    a designer selects, and the only thing whose position in the drawing order
    means anything. A shape inside a group has no z-order of its own relative to
    the slide; its group does."""
    element = shape._element
    while element is not None and element.getparent() is not spTree:
        element = element.getparent()
    return element


# --- the operations -------------------------------------------------------


def _do_set_color(prs, params, notes, touched, undo):
    rgb = parse_hex(params.get("hex"))
    if rgb is None:
        return "that is not a colour"
    color = RGBColor(*rgb)
    seen_shapes = set()
    for target in params.get("targets") or []:
        idx = target.get("slide_index")
        if idx is None or not (0 <= idx < len(prs.slides)):
            continue
        slide = prs.slides[idx]
        shape = _by_id(slide).get(str(target.get("shape_id")))
        if shape is None:
            continue
        key = (idx, str(target.get("shape_id")))
        if key not in seen_shapes:
            op = _replace_undo(shape, target.get("surface") or "shape")
            if op is not None:
                undo.append(dict(op, slide_index=idx))
            seen_shapes.add(key)
            touched.append(key)

        surface = target.get("surface") or "fill"
        try:
            if surface == "fill":
                shape.fill.solid()
                shape.fill.fore_color.rgb = color
                notes.append(f"{shape.name!r} filled #{params['hex']}")
            elif surface == "line":
                shape.line.color.rgb = color
                notes.append(f"{shape.name!r} outlined #{params['hex']}")
            else:
                n = _recolor_runs(shape, target.get("locator"), color)
                if n:
                    notes.append(f"{n} run(s) set to #{params['hex']}")
        except Exception as exc:
            notes.append(f"failed on {shape.name!r}: {type(exc).__name__}: {exc}")
    return None


def _recolor_runs(shape, locator, color) -> int:
    """Recolour one run, or every run when no locator is given.

    A locator is the normal case for contrast: only the runs that actually fail
    are named, so a deliberate accent word inside the same box keeps its colour.
    """
    if not getattr(shape, "has_text_frame", False):
        return 0
    want = None
    if locator and locator.startswith("p"):
        try:
            p_part, r_part = locator.split("/")
            want = (int(p_part[1:]), int(r_part[1:]))
        except (ValueError, IndexError):
            want = None
    n = 0
    for p_idx, para in enumerate(shape.text_frame.paragraphs):
        for r_idx, run in enumerate(para.runs):
            if want is not None and (p_idx, r_idx) != want:
                continue
            run.font.color.rgb = color
            n += 1
    return n


def _do_offset(prs, params, notes, touched, undo):
    targets = params.get("targets")
    if targets is None:
        targets = [{"slide_index": params.get("slide_index"),
                    "shape_id": params.get("shape_id")}]
    dx, dy = int(params.get("dx") or 0), int(params.get("dy") or 0)
    for target in targets:
        idx = target.get("slide_index")
        if idx is None or not (0 <= idx < len(prs.slides)):
            continue
        shape = _by_id(prs.slides[idx]).get(str(target.get("shape_id")))
        if shape is None or shape.left is None or shape.top is None:
            continue
        undo.append({"op": "offset", "slide_index": idx,
                     "shape_id": str(shape.shape_id),
                     "left": int(shape.left), "top": int(shape.top)})
        touched.append((idx, str(shape.shape_id)))
        shape.left = int(shape.left) + dx
        shape.top = int(shape.top) + dy
        notes.append(f"{shape.name!r} moved {dx / 914400:+.2f}in, "
                     f"{dy / 914400:+.2f}in")
    return None


def _do_zorder(prs, params, notes, touched, undo):
    idx = params.get("slide_index")
    if idx is None or not (0 <= idx < len(prs.slides)):
        return "that slide is no longer in the deck"
    slide = prs.slides[idx]
    lookup = _by_id(slide)
    mover = lookup.get(str(params.get("shape_id")))
    anchor = lookup.get(str(params.get("below")))
    if mover is None or anchor is None:
        return "one of the two shapes is no longer on the slide"
    spTree = slide.shapes._spTree
    mover_el = _top_level(mover, spTree)
    anchor_el = _top_level(anchor, spTree)
    if mover_el is None or anchor_el is None or mover_el is anchor_el:
        return "those two shapes are in the same group, so neither is in " \
               "front of the other"
    children = list(spTree)
    undo.append({"op": "zorder", "slide_index": idx,
                 "shape_id": str(mover.shape_id),
                 "index": children.index(mover_el)})
    touched.append((idx, str(mover.shape_id)))
    spTree.remove(mover_el)
    anchor_el.addprevious(mover_el)
    notes.append(f"{mover.name!r} moved behind {anchor.name!r}")
    return None


def _do_delete(prs, params, notes, touched, undo):
    for target in params.get("targets") or []:
        idx = target.get("slide_index")
        if idx is None or not (0 <= idx < len(prs.slides)):
            continue
        slide = prs.slides[idx]
        shape = _by_id(slide).get(str(target.get("shape_id")))
        if shape is None:
            continue
        element = shape._element
        parent = element.getparent()
        if parent is None:
            continue
        blob = _xml(shape)
        if blob is None:
            notes.append(f"{shape.name!r} could not be copied, so it was left "
                         f"in place rather than removed without a way back")
            continue
        undo.append({"op": "insert", "slide_index": idx, "xml": blob,
                     "index": list(parent).index(element)})
        touched.append((idx, str(shape.shape_id)))
        label = ""
        try:
            label = shape.text_frame.text.strip()[:24]
        except Exception:
            pass
        parent.remove(element)
        notes.append(f"removed {label or shape.name!r} from slide {idx + 1}")
    return None


def _do_resize(prs, params, notes, touched, undo):
    """Grow or shrink one shape, holding the edge the caller names.

    `anchor` is which edge stays put: growing a card downward (anchor "top") and
    growing it upward (anchor "bottom") are different fixes for different
    slides, and a resize that silently picks one is a resize that moves things
    the designer did not ask to move."""
    idx = params.get("slide_index")
    if idx is None or not (0 <= idx < len(prs.slides)):
        return "that slide is no longer in the deck"
    shape = _by_id(prs.slides[idx]).get(str(params.get("shape_id")))
    if shape is None:
        return "that shape is no longer on the slide"
    if None in (shape.left, shape.top, shape.width, shape.height):
        return "that shape states no size of its own, so it cannot be resized"
    dw, dh = int(params.get("dw") or 0), int(params.get("dh") or 0)
    if not dw and not dh:
        return "nothing to resize"
    if shape.width + dw <= 0 or shape.height + dh <= 0:
        return "that would leave the shape with no size"

    op = _replace_undo(shape, "size")
    if op is None:
        return "that shape could not be copied, so nothing was changed"
    undo.append(dict(op, slide_index=idx))
    touched.append((idx, str(shape.shape_id)))

    anchor = params.get("anchor") or ("top" if dh else "left")
    if anchor == "right":
        shape.left = int(shape.left) - dw
    if anchor == "bottom":
        shape.top = int(shape.top) - dh
    shape.width = int(shape.width) + dw
    shape.height = int(shape.height) + dh
    notes.append(f"{shape.name!r} resized {dw / 914400:+.2f}in by "
                 f"{dh / 914400:+.2f}in, holding its {anchor} edge")
    return None


def _do_scale_text(prs, params, notes, touched, undo):
    """Multiply every run's size in one shape.

    Resolved sizes, not stated ones: most runs in a real deck state no size at
    all and inherit it, so scaling only what is written down would leave an
    inherited 18pt paragraph at 18pt and shrink the one run beside it - which
    looks like a bug rather than a fix."""
    from pptx.util import Pt

    from spike.resolver import resolve_run

    idx = params.get("slide_index")
    if idx is None or not (0 <= idx < len(prs.slides)):
        return "that slide is no longer in the deck"
    slide = prs.slides[idx]
    shape = _by_id(slide).get(str(params.get("shape_id")))
    if shape is None or not getattr(shape, "has_text_frame", False):
        return "that text box is no longer on the slide"
    try:
        scale = float(params.get("scale") or 0)
    except (TypeError, ValueError):
        scale = 0.0
    if not 0.1 < scale < 3.0:
        return "that is not a usable scale"

    op = _replace_undo(shape, "type sizes")
    if op is None:
        return "that shape could not be copied, so nothing was changed"
    undo.append(dict(op, slide_index=idx))
    touched.append((idx, str(shape.shape_id)))

    changed = 0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue
            try:
                size = float(resolve_run(run, para, shape, slide,
                                         prs).size_pt.value)
            except Exception:
                continue
            # A floor, not a rounding: type below 8pt is not small, it is
            # unreadable, and shrinking to fit is not worth trading one
            # legibility defect for another.
            run.font.size = Pt(max(8.0, round(size * scale, 1)))
            changed += 1
    if not changed:
        undo.clear()
        touched.clear()
        return "no run in that box has a size this could scale"
    notes.append(f"{changed} run(s) in {shape.name!r} scaled to "
                 f"{scale * 100:.0f}%")
    return None


def _do_autofit(prs, params, notes, touched, undo):
    """Set shrink-on-overflow on one text box, so PowerPoint keeps it fitting.

    The one fix on this page that keeps working after the designer edits the
    copy: it states the intent (this text must fit) rather than a size that
    happens to satisfy it today."""
    from pptx.enum.text import MSO_AUTO_SIZE

    idx = params.get("slide_index")
    if idx is None or not (0 <= idx < len(prs.slides)):
        return "that slide is no longer in the deck"
    shape = _by_id(prs.slides[idx]).get(str(params.get("shape_id")))
    if shape is None or not getattr(shape, "has_text_frame", False):
        return "that text box is no longer on the slide"
    op = _replace_undo(shape, "autofit")
    if op is None:
        return "that shape could not be copied, so nothing was changed"
    undo.append(dict(op, slide_index=idx))
    touched.append((idx, str(shape.shape_id)))
    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    notes.append(f"{shape.name!r} set to shrink its text on overflow")
    return None


def _do_set_insets(prs, params, notes, touched, undo):
    """Set one text box's internal margins.

    The cheapest fix for text that does not fit: padding is invisible on the
    slide, so returning it to the default gives the words room without moving
    the shape or touching the type scale. Every edge is written explicitly
    rather than cleared, because clearing an attribute leaves the box inheriting
    from a placeholder that may state its own padding, and "the default" then
    means something different per slide.
    """
    idx = params.get("slide_index")
    if idx is None or not (0 <= idx < len(prs.slides)):
        return "that slide is no longer in the deck"
    shape = _by_id(prs.slides[idx]).get(str(params.get("shape_id")))
    if shape is None or not getattr(shape, "has_text_frame", False):
        return "that text box is no longer on the slide"

    edges = {}
    for edge in ("left", "right", "top", "bottom"):
        value = params.get(edge)
        if value is None:
            continue
        value = int(value)
        if value < 0:
            return "a negative margin is not a margin"
        edges[edge] = value
    if not edges:
        return "no margins to set"

    op = _replace_undo(shape, "margins")
    if op is None:
        return "that shape could not be copied, so nothing was changed"
    undo.append(dict(op, slide_index=idx))
    touched.append((idx, str(shape.shape_id)))

    frame = shape.text_frame
    for edge, value in edges.items():
        setattr(frame, f"margin_{edge}", Emu(value))
    notes.append(f"{shape.name!r} internal margins set to "
                 + ", ".join(f"{e} {v / 914400:.2f}in"
                             for e, v in sorted(edges.items())))
    return None


def _do_front(prs, params, notes, touched, undo):
    """Bring one shape to the front of the drawing order."""
    idx = params.get("slide_index")
    if idx is None or not (0 <= idx < len(prs.slides)):
        return "that slide is no longer in the deck"
    slide = prs.slides[idx]
    shape = _by_id(slide).get(str(params.get("shape_id")))
    if shape is None:
        return "that shape is no longer on the slide"
    spTree = slide.shapes._spTree
    element = _top_level(shape, spTree)
    if element is None:
        return "that shape is inside a group, so it has no order of its own"
    children = list(spTree)
    undo.append({"op": "zorder", "slide_index": idx,
                 "shape_id": str(shape.shape_id),
                 "index": children.index(element)})
    touched.append((idx, str(shape.shape_id)))
    spTree.remove(element)
    # Before p:extLst, which must stay last or PowerPoint repairs the file.
    spTree.insert_element_before(element, "p:extLst")
    notes.append(f"{shape.name!r} brought to the front")
    return None


def _do_set_theme_color(prs, params, notes, touched, undo):
    """Repoint a surface at a THEME SLOT instead of a literal hex.

    Strictly better than snapping to the same colour as a hex, when the palette
    value is a theme slot: the surface then moves with the theme, which is what
    a theme is for, and the next rebrand does not have to find it again."""
    from pptx.oxml.ns import qn

    slot = str(params.get("slot") or "")
    if not slot:
        return "no theme slot was named"
    seen = set()
    for target in params.get("targets") or []:
        idx = target.get("slide_index")
        if idx is None or not (0 <= idx < len(prs.slides)):
            continue
        slide = prs.slides[idx]
        shape = _by_id(slide).get(str(target.get("shape_id")))
        if shape is None:
            continue
        key = (idx, str(target.get("shape_id")))
        if key not in seen:
            op = _replace_undo(shape, target.get("surface") or "shape")
            if op is not None:
                undo.append(dict(op, slide_index=idx))
            seen.add(key)
            touched.append(key)
        parents = []
        if (target.get("surface") or "fill") == "fill":
            shape.fill.solid()          # ensure a solidFill exists to repoint
            parents.append(find(shape._element, "p:spPr"))
        elif target.get("surface") == "line":
            parents.append(find(find(shape._element, "p:spPr"), "a:ln"))
        else:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        parents.append(find(run._r, "a:rPr"))
        done = 0
        for parent in parents:
            solid = find(parent, "a:solidFill") if parent is not None else None
            if solid is None:
                continue
            for child in list(solid):
                solid.remove(child)
            scheme = solid.makeelement(qn("a:schemeClr"), {"val": slot})
            solid.append(scheme)
            done += 1
        if done:
            notes.append(f"{shape.name!r} repointed at the theme's {slot}")
    return None


_OPS = {"set_color": _do_set_color, "offset": _do_offset,
        "offset_many": _do_offset, "zorder": _do_zorder, "delete": _do_delete,
        "resize": _do_resize, "scale_text": _do_scale_text,
        "set_insets": _do_set_insets,
        "autofit": _do_autofit, "front": _do_front,
        "set_theme_color": _do_set_theme_color}


# --- the pass ------------------------------------------------------------


def apply(deck_bytes: bytes, picks: list) -> tuple[bytes, list[Applied]]:
    """Perform each (finding, remedy) pick against the deck, in the order given.

    picks: [(DesignFinding, Remedy)]. A remedy with no `op` - every finding's
    "leave it" - is recorded as a decision and changes nothing, which is the
    point of offering it: "I looked at this and it is fine" is an answer, and a
    page that cannot record it asks the same question again next week.

    Returns the new deck and one Applied per pick, in the order given. A pick
    that could not be performed comes back with done=False and a detail saying
    why, never a success the deck does not show.
    """
    prs = Presentation(io.BytesIO(deck_bytes))
    out: list[Applied] = []
    for finding, remedy in picks:
        notes: list[str] = []
        touched: list[tuple] = []
        undo: list[dict] = []
        if not remedy.op:
            out.append(Applied(
                finding.finding_id, remedy.remedy_id, finding.kind,
                finding.headline, remedy.label, True,
                "left as it is, on purpose", list(finding.slides), [], []))
            continue
        handler = _OPS.get(remedy.op)
        if handler is None:
            out.append(Applied(
                finding.finding_id, remedy.remedy_id, finding.kind,
                finding.headline, remedy.label, False,
                f"this build cannot perform '{remedy.op}'",
                list(finding.slides), [], []))
            continue
        try:
            refusal = handler(prs, remedy.params, notes, touched, undo)
        except Exception as exc:
            refusal = f"{type(exc).__name__}: {exc}"
        failed = [n for n in notes if n.startswith("failed")]
        done = refusal is None and bool(undo) and len(failed) < len(notes or [1])
        detail = refusal or ("; ".join(notes[:4]) if notes else
                             "nothing on the deck matched this any more")
        out.append(Applied(finding.finding_id, remedy.remedy_id, finding.kind,
                           finding.headline, remedy.label, done, detail,
                           list(finding.slides), undo if done else [],
                           touched if done else []))
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), out


def followers(applied: list, finding_id: str) -> list:
    """The applied remedies that have to come back with this one: itself, and
    every LATER one that wrote to a shape it wrote to.

    The same correctness requirement qc.undo.followers exists for, on a
    different key. There it is the slide, because the migration works down a
    slide in order and each step is computed on what the last one left. Here it
    is the shape, because these remedies are independent by construction -
    recolouring a navy on slide 3 says nothing about an overlap on slide 11 -
    EXCEPT when two of them touched the same shape. Then the first one's stored
    element predates the second one's change, and replaying it alone would put
    back a version of the shape from before work the designer also approved,
    with no sign that anything was lost.
    """
    index = next((i for i, a in enumerate(applied)
                  if a.finding_id == finding_id), None)
    if index is None:
        return []
    chain = [applied[index]]
    reach = set(applied[index].touched)
    for later in applied[index + 1:]:
        if reach & set(later.touched):
            chain.append(later)
            reach |= set(later.touched)
    return chain


def undo_items(chain: list) -> list[dict]:
    """The items qc.undo.apply_undo wants, from a chain of Applied.

    Each op carries the slide it belongs to (written by the handlers above), so
    one Applied that touched three slides - a palette snap usually does - splits
    into one item per slide rather than being handed to apply_undo with a single
    slide index that is wrong for two thirds of it.
    """
    items = []
    for entry in chain:
        by_slide: dict[int, list] = {}
        for op in entry.undo:
            by_slide.setdefault(op.get("slide_index"), []).append(op)
        for slide_index, ops in by_slide.items():
            items.append({"change_id": entry.finding_id,
                          "slide_index": slide_index,
                          "action": entry.label, "ops": ops})
    return items


__all__ = ["Applied", "apply", "followers", "undo_items"]
