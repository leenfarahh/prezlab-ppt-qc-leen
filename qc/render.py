"""Slide rendering via desktop PowerPoint COM, for the before/after diff.

This is the LOCAL PILOT renderer: it drives the PowerPoint installed on this
machine in an interactive user session, which Microsoft supports (the
unsupported scenario is unattended servers; the production tier uses
Microsoft Graph per the PRD, pending U5). Fidelity is exact by definition,
including Arabic shaping, because it IS PowerPoint.

Renders are serialized behind a lock: PowerPoint is single-instance and COM
objects must stay on the thread that created them.
"""

import io
import os
import shutil
import subprocess
import tempfile
import threading
from pathlib import Path

from pptx import Presentation

from .config import RENDERER
from .util import iter_shapes_deep

_RENDER_LOCK = threading.Lock()
RENDER_WIDTH = 1360  # px; height follows the deck's aspect ratio


def _libreoffice_available() -> bool:
    return bool(shutil.which("soffice") or shutil.which("libreoffice"))


def export_decks_png(decks: dict[str, bytes], slide_indices: list[int],
                     width: int = RENDER_WIDTH) -> dict[str, bytes]:
    """Render the given zero-based slides of each named deck to PNG bytes.

    decks: {"before": pptx_bytes, "after": pptx_bytes}
    Returns {"before:12": png_bytes, "after:12": ...}. Raises RuntimeError
    with a readable message when no renderer is available.

    Renderer selection (config.RENDERER): "com" forces PowerPoint (Windows),
    "libreoffice" forces the Linux/cloud path, "none" disables rendering,
    "auto" (default) uses COM when it works and falls back to LibreOffice.
    """
    if RENDERER == "none":
        raise RuntimeError("rendering disabled (QC_RENDERER=none)")
    if RENDERER == "libreoffice":
        return _export_libreoffice(decks, slide_indices, width)
    try:
        return _export_com(decks, slide_indices, width)
    except RuntimeError:
        if RENDERER == "com":
            raise
        if _libreoffice_available():
            return _export_libreoffice(decks, slide_indices, width)
        raise


def _contiguous(indices: list[int]) -> list[tuple[int, int]]:
    """`indices` as (first, last) runs, sorted and deduped. One pdftoppm call
    covers a run, so a whole-deck render is one call rather than two hundred."""
    out: list[tuple[int, int]] = []
    for idx in sorted(set(indices)):
        if out and idx == out[-1][1] + 1:
            out[-1] = (out[-1][0], idx)
        else:
            out.append((idx, idx))
    return out


def _pages_to_png(pdftoppm: str, pdf, name: str, slide_indices: list[int],
                  tmp, width: int) -> dict[str, bytes]:
    """The wanted pages of `pdf`, rendered to PNG.

    ONE PROCESS PER RUN OF CONSECUTIVE PAGES, not one per page. pdftoppm parses
    the whole PDF on startup, so rendering a 200-slide deck a page at a time
    launched 200 processes and re-parsed the document 200 times - and the whole
    deck is exactly what _ensure_thumbs asks for. A contiguous range is the
    common case (a window of slides, or all of them), so this is usually a
    single call.

    Page numbers come back from pdftoppm as a fixed-width suffix it chooses
    itself, so the files are matched by sorted order within the run rather than
    by a name guessed here.
    """
    out: dict[str, bytes] = {}
    for first, last in _contiguous(slide_indices):
        stem = tmp / f"{name}-{first}"
        args = [pdftoppm, "-png", "-scale-to-x", str(width), "-scale-to-y", "-1",
                "-f", str(first + 1), "-l", str(last + 1)]
        if first == last:
            args.append("-singlefile")
        subprocess.run(args + [str(pdf), str(stem)],
                       check=True, capture_output=True, timeout=600)
        if first == last:
            png = stem.with_suffix(".png")
            if png.exists():
                out[f"{name}:{first}"] = png.read_bytes()
            continue
        # Multi-page: one file per page, suffixed -1, -01, -001 ... depending on
        # the page count. Sorted order IS page order for a fixed-width suffix.
        pages = sorted(tmp.glob(f"{stem.name}-*.png"))
        for offset, png in enumerate(pages):
            if first + offset <= last:
                out[f"{name}:{first + offset}"] = png.read_bytes()
    return out


def _export_libreoffice(decks: dict[str, bytes], slide_indices: list[int],
                        width: int = RENDER_WIDTH) -> dict[str, bytes]:
    """Render via LibreOffice headless (pptx -> pdf) then poppler (pdf page ->
    png). The cloud/Linux path. Fidelity differs from PowerPoint, especially
    Arabic shaping, so this is labeled a demo-grade renderer in the UI."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice or not pdftoppm:
        raise RuntimeError("LibreOffice/poppler not installed on this host")

    out: dict[str, bytes] = {}
    with _RENDER_LOCK:
        tmp = Path(tempfile.mkdtemp(prefix="qc-lo-"))
        try:
            for name, data in decks.items():
                src = tmp / f"{name}.pptx"
                src.write_bytes(data)
                # isolate the LO user profile so concurrent runs do not clash
                env = dict(os.environ, HOME=str(tmp))
                subprocess.run(
                    [soffice, "--headless", "--norestore",
                     f"-env:UserInstallation=file://{tmp / 'louser'}",
                     "--convert-to", "pdf", "--outdir", str(tmp), str(src)],
                    check=True, capture_output=True, timeout=300, env=env)
                pdf = tmp / f"{name}.pdf"
                if not pdf.exists():
                    raise RuntimeError(f"LibreOffice did not produce a PDF for {name}")
                out.update(_pages_to_png(pdftoppm, pdf, name, slide_indices,
                                         tmp, width))
        finally:
            shutil.rmtree(tmp, ignore_errors=True)
    return out


def _export_com(decks: dict[str, bytes], slide_indices: list[int],
                width: int = RENDER_WIDTH, attempts: int = 2) -> dict[str, bytes]:
    """Render through desktop PowerPoint, with one retry on a fresh instance.

    Presentations.Open fails intermittently - roughly one run in three when the
    suite starts and quits PowerPoint many times - with a bare "Failed." and no
    reason. The instance is wedged by then, so retrying against it is pointless;
    _export_com_once tears its own instance down completely (qc.unify.force_quit)
    and the second attempt starts clean, which is what actually recovers.

    A DISPATCH refusal is not retried. That is the host saying PowerPoint cannot
    be automated at all, and the advice it carries is the answer."""
    last = None
    for _attempt in range(max(1, attempts)):
        try:
            return _export_com_once(decks, slide_indices, width)
        except RuntimeError:
            raise
        except Exception as exc:
            last = exc
    from .unify import com_failure_advice

    raise RuntimeError(com_failure_advice(last))


def _export_com_once(decks: dict[str, bytes], slide_indices: list[int],
                     width: int = RENDER_WIDTH) -> dict[str, bytes]:
    import pythoncom
    import win32com.client

    from .unify import sweep_automation

    out: dict[str, bytes] = {}
    with _RENDER_LOCK:
        pythoncom.CoInitialize()
        tmp_dir = Path(tempfile.mkdtemp(prefix="qc-render-"))
        app = None
        # Clear leftovers from an interrupted run FIRST. DispatchEx does not
        # start a second PowerPoint - it attaches to whatever /AUTOMATION
        # instance is already running - so a wedged leftover is not something
        # this run can work around, and a snapshot taken after the dispatch
        # could never see it. What the sweep hands back is what survived, which
        # IS the teardown's snapshot, so this is one process listing rather than
        # two (qc.unify.sweep_automation, qc.unify.force_quit).
        started = sweep_automation()
        try:
            try:
                app = win32com.client.DispatchEx("PowerPoint.Application")
                app.DisplayAlerts = 1  # ppAlertsNone
            except Exception as exc:
                from .unify import com_failure_advice

                raise RuntimeError(com_failure_advice(exc)) from exc

            for name, data in decks.items():
                deck_path = tmp_dir / f"{name}.pptx"
                deck_path.write_bytes(data)
                pres = app.Presentations.Open(str(deck_path), -1, 0, 0)
                try:
                    aspect = pres.PageSetup.SlideHeight / pres.PageSetup.SlideWidth
                    height = int(width * aspect)
                    for idx in slide_indices:
                        if idx >= pres.Slides.Count:
                            continue
                        png_path = tmp_dir / f"{name}-{idx}.png"
                        pres.Slides(idx + 1).Export(str(png_path), "PNG",
                                                    width, height)
                        out[f"{name}:{idx}"] = png_path.read_bytes()
                finally:
                    try:
                        pres.Close()
                    except Exception:
                        pass  # a wedged presentation is force_quit's problem
        finally:
            if app is not None:
                # Quit, and make sure. A Quit that quietly fails leaves a
                # windowless POWERPNT.EXE behind, and every later render on the
                # host fails against it and leaks another (qc.unify.force_quit).
                from .unify import force_quit

                force_quit(app, started)
            # release the COM proxy BEFORE CoUninitialize, or the teardown
            # trips RPC_E_DISCONNECTED (0x80010108) while PowerPoint exits
            app = None
            import gc

            gc.collect()
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass
            for f in tmp_dir.glob("*"):
                try:
                    f.unlink()
                except OSError:
                    pass
            try:
                tmp_dir.rmdir()
            except OSError:
                pass
    return out


# Enough to see a master's whole vocabulary without asking PowerPoint to render
# a hundred slides for one page view. A master carrying more than this is
# truncated and the page says so; a silent cut reads as "the master has fewer
# layouts than it does".
#
# THE ONLY CAP. qc.layoutmatch used to keep a second, smaller one of its own
# (that module is gone; the cap it argued with is not),
# so the sheet shown to the vision pass was 16 of the 24 rendered here - and a
# slide whose real home was layout 20 came back "no layout in this master fits",
# which qc.layoutgap reports as a MISSING LAYOUT and the page states as "checked
# against every layout in the master". Two numbers for one idea turned a
# truncation into a finding about the client's master (30/08/2026).
MAX_LAYOUTS = 24


def layout_catalogue(deck_bytes: bytes) -> tuple[bytes, list[dict], int]:
    """A deck holding ONE blank slide per layout, and what each one is.

    This is how a layout gets rendered at all: PowerPoint exports SLIDES, not
    layouts, so the only way to photograph a layout is to put an empty slide on
    it. Adding those slides to a copy whose own slides have been dropped keeps
    the render honest - what you see is the layout's own furniture, placeholders
    and background, with nothing of the deck's content on top of it - and keeps
    it quick, because PowerPoint then opens a file with fifteen slides instead
    of the deck's two hundred.

    Every master is walked, not just the dominant one. A deck that could not
    rebuild every slide keeps its ORIGINAL master alive alongside the applied
    one (qc.applymaster.ApplyResult.stragglers), and which layouts came from
    which master is exactly what the review is for.

    Returns (deck_bytes, entries, skipped) where each entry is
    {"index", "master", "layout"} - index being the slide index to render - and
    `skipped` counts the layouts past MAX_LAYOUTS, so the page can say it
    truncated instead of implying the master is smaller than it is.
    """
    from pptx.oxml.ns import qn

    prs = Presentation(io.BytesIO(deck_bytes))
    id_list = prs.slides._sldIdLst
    for sldId in list(id_list):
        rId = sldId.get(qn("r:id"))
        id_list.remove(sldId)
        if rId:
            try:
                prs.part.drop_rel(rId)
            except KeyError:
                pass

    entries: list[dict] = []
    skipped = 0
    for m_index, master in enumerate(prs.slide_masters):
        for layout in master.slide_layouts:
            if len(entries) >= MAX_LAYOUTS:
                skipped += 1
                continue
            prs.slides.add_slide(layout)
            entries.append({"index": len(entries), "master": m_index,
                            "layout": layout.name or f"Layout {len(entries) + 1}"})
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), entries, skipped


def layout_previews(before: bytes, after: bytes) -> dict:
    """Every layout of the submitted deck, and every layout of the rebuilt one.

    Returns {"before": [entry...], "after": [entry...], "images": {key: png},
    "error": str|None} with image keys "before:<i>" / "after:<i>" matching each
    entry's index, i.e. the same shape build_diff returns, so the page renders
    both the same way.

    A RENDER failure does not lose the entries. Which layouts a deck arrived with
    and which it has now is read out of the file by python-pptx and needs no
    PowerPoint at all; letting an export failure take the whole answer down left
    the page saying "No layouts to show" about a master carrying twelve of them
    (design lead, 23/08/2026). The pictures are the illustration, not the
    finding."""
    before_deck, before_entries, cut_b = layout_catalogue(before)
    after_deck, after_entries, cut_a = layout_catalogue(after)
    images: dict[str, bytes] = {}
    error = None
    try:
        if before_entries:
            images.update(export_decks_png({"before": before_deck},
                                           [e["index"] for e in before_entries]))
        if after_entries:
            images.update(export_decks_png({"after": after_deck},
                                           [e["index"] for e in after_entries]))
    except Exception as exc:
        error = f"{type(exc).__name__}: {exc}"
    return {"before": before_entries, "after": after_entries, "images": images,
            "truncated": bool(cut_b or cut_a), "error": error}


def slide_previews(before: bytes, after: bytes, indices: list[int]) -> dict:
    """The deck itself, before and after, for the given slides.

    Separate from build_diff because there is nothing to highlight here: the
    format pass rewrites a whole slide onto a new layout rather than fixing
    named shapes, so a rectangle around "the changed shape" would be a rectangle
    around the slide. The comparison IS the finding."""
    if not indices:
        return {"indices": [], "images": {}, "error": None}
    try:
        images = export_decks_png({"before": before, "after": after}, indices)
    except Exception as exc:
        return {"indices": list(indices), "images": {},
                "error": f"{type(exc).__name__}: {exc}"}
    return {"indices": list(indices), "images": images, "error": None}


def _fraction_box(left, top, width, height, slide_w, slide_h) -> dict:
    """A shape's box as fractions of the slide, CLIPPED to the slide.

    The two callers below clamped the origin and the size independently -
    max(0, left/w) with min(1, width/w) - which is right only while the shape is
    on the canvas. A shape hanging off the left edge had its origin pulled to 0
    and kept its full width, so the highlight came out shifted: it covered the
    wrong part of the slide, on exactly the shapes an audit flags for being
    outside the frame.

    Clipping the far edge and deriving the size from the clipped corners keeps
    the rectangle over the part of the shape a designer can actually see.
    """
    x0 = max(0.0, min(1.0, left / slide_w))
    y0 = max(0.0, min(1.0, top / slide_h))
    x1 = max(0.0, min(1.0, (left + width) / slide_w))
    y1 = max(0.0, min(1.0, (top + height) / slide_h))
    return {"x": x0, "y": y0, "w": max(0.0, x1 - x0), "h": max(0.0, y1 - y0)}


def shape_rects(deck_bytes: bytes, wanted: dict[int, list[dict]]) -> dict[int, list[dict]]:
    """Normalized highlight rectangles for changed shapes, per slide.

    wanted: {slide_index: [applied record dicts]} from the fix pass.
    Returns {slide_index: [{"x","y","w","h" (0-1 floats), "label"}]}.
    Rectangles come from the deck's own geometry at read time, so they are
    correct on BOTH sides of a geometry fix (before shows the override
    position, after shows the inherited one)."""
    prs = Presentation(io.BytesIO(deck_bytes))
    slide_w, slide_h = prs.slide_width, prs.slide_height
    out: dict[int, list[dict]] = {}

    for slide_idx, records in wanted.items():
        if slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]
        by_id = {str(sh.shape_id): sh for sh, _p in iter_shapes_deep(slide.shapes)}
        rects = []
        for rec in records:
            shape = by_id.get(str(rec["shape_id"]))
            if shape is None:
                continue
            left, top = shape.left, shape.top
            width, height = shape.width, shape.height
            if None in (left, top, width, height):
                continue
            rects.append({
                **_fraction_box(left, top, width, height, slide_w, slide_h),
                "label": rec["issue_type"],
            })
        out[slide_idx] = rects
    return out


def build_diff(before: bytes, after: bytes, applied_records: list[dict]) -> dict:
    """Render + overlay data for EVERY slide in the deck, changed or not.

    It used to build the list from the applied records alone, so a slide
    nothing was fixed on simply was not there. A designer counting slides down
    the review found 1, 2, 3, 5 and had no way to tell "slide 4 was fine" from
    "slide 4 was skipped" - and those are opposite facts about a deck about to
    go to a client (design lead, 31/08/2026).

    So every slide is shown and the untouched ones carry `changes: 0`, which is
    the page's cue to say so out loud. The cost is real and worth naming: this
    exports both decks at every slide rather than at the handful that changed,
    so a long deck takes proportionally longer on the first visit. The result
    is cached on the job, so it is paid once.

    Returns {"slides": [{"index", "changes", "labels", "before_rects",
    "after_rects"}], "images": {"before:i": png, ...}}."""
    wanted: dict[int, list[dict]] = {}
    for rec in applied_records:
        wanted.setdefault(rec["slide_index"], []).append(rec)

    try:
        total = len(Presentation(io.BytesIO(after)).slides)
    except Exception:
        # A deck that cannot be reopened is not a reason to lose the review of
        # the slides that were fixed; fall back to what the records name.
        total = (max(wanted) + 1) if wanted else 0
    indices = list(range(total)) or sorted(wanted)
    if not indices:
        return {"slides": [], "images": {}}

    images = export_decks_png({"before": before, "after": after}, indices)
    # Rects are still only asked for where something changed: an untouched
    # slide has nothing to outline, and highlighting it would be a lie.
    before_rects = shape_rects(before, wanted)
    after_rects = shape_rects(after, wanted)

    slides = []
    for idx in indices:
        records = wanted.get(idx) or []
        slides.append({
            "index": idx,
            "changes": len(records),
            "labels": sorted({r["issue_type"] for r in records}),
            "before_rects": before_rects.get(idx, []),
            "after_rects": after_rects.get(idx, []),
        })
    return {"slides": slides, "images": images}


_SEV_WORST = {"error": 0, "warning": 1, "info": 2}


def pin_numbers(records: list[dict]) -> dict[str, int]:
    """record_id -> pin number: shape-bound findings are numbered per slide
    in list order, one pin per flagged shape (findings sharing a shape share
    its pin). The SAME rule numbers the rects in audit_rects and the badge
    in the report row, so the two always agree. Slide-level records get no
    pin; the report marks them 'whole slide'."""
    pins: dict[str, int] = {}
    per_slide: dict[int, dict[str, int]] = {}
    for rec in records:
        sid = str(rec.get("shape_id") or "-")
        if sid == "-" or rec.get("action") == "changed" \
                or rec.get("module") == "preflight":
            continue
        slide_pins = per_slide.setdefault(rec["slide_index"], {})
        if sid not in slide_pins:
            slide_pins[sid] = len(slide_pins) + 1
        pins[rec["record_id"]] = slide_pins[sid]
    return pins


def audit_rects(deck_bytes: bytes, records: list[dict]) -> dict[int, list[dict]]:
    """Highlight rectangles for the audit views: one rect per flagged shape
    per slide, deduped across records (worst severity wins, issue labels
    joined), numbered to match the pin badges in the findings list.
    Slide-level records (shape_id '-') carry no rect; they appear only in
    the findings list."""
    prs = Presentation(io.BytesIO(deck_bytes))
    slide_w, slide_h = prs.slide_width, prs.slide_height

    by_slide: dict[int, dict[str, dict]] = {}
    for rec in records:
        sid = str(rec.get("shape_id") or "-")
        if sid == "-" or rec.get("action") == "changed" \
                or rec.get("module") == "preflight":
            continue
        slide_slots = by_slide.setdefault(rec["slide_index"], {})
        if sid not in slide_slots:
            slide_slots[sid] = {"pin": len(slide_slots) + 1, "labels": [],
                                "severity": "info", "arabic": False,
                                "record_ids": []}
        slot = slide_slots[sid]
        if rec["issue_type"] not in slot["labels"]:
            slot["labels"].append(rec["issue_type"])
        if _SEV_WORST.get(rec["severity"], 3) < _SEV_WORST.get(slot["severity"], 3):
            slot["severity"] = rec["severity"]
        slot["arabic"] = slot["arabic"] or bool(rec.get("arabic_flag"))
        slot["record_ids"].append(rec["record_id"])

    out: dict[int, list[dict]] = {}
    for slide_idx, shapes in by_slide.items():
        if slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]
        by_id = {str(sh.shape_id): sh for sh, _p in iter_shapes_deep(slide.shapes)}
        rects = []
        for sid, slot in shapes.items():
            shape = by_id.get(sid)
            if shape is None:
                continue
            left, top = shape.left, shape.top
            width, height = shape.width, shape.height
            if None in (left, top, width, height):
                continue
            rects.append({
                "pin": slot["pin"],
                **_fraction_box(left, top, width, height, slide_w, slide_h),
                "label": " · ".join(slot["labels"]),
                "severity": slot["severity"],
                "arabic": slot["arabic"],
                "record_ids": slot["record_ids"],
            })
        out[slide_idx] = rects
    return out
