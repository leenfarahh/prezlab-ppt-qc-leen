"""Stage 1: read a submitted master's DESIGN SURFACE into a Style Spec.

    python -m qc.stylespec "master.pptx" [--json spec.json]

This module reads the slideMaster, its slideLayouts, and the theme part. It
never looks at slide content. That is the whole point: a designer submits a
finished master, which may carry one sample slide or none at all, and the
question is what the master DECLARES, not what any slides happen to do.

Contrast with qc/bootstrap.py, which surveys what slides actually do in order
to infer rules from a finished deck. Both produce style facts; they answer
different questions and must not be confused:

    bootstrap.py   "what conventions does this deck follow?"    (reads slides)
    stylespec.py   "what does this master define?"              (reads master)

What is read, and from where:

    theme       theme1.xml         scheme colors by role, clrMap, major/minor
                                   fonts per script
    layouts     each slideLayout   name, archetype type, placeholder geometry
                                   per placeholder, marked explicit vs
                                   inherited from the master
    master      slideMaster        background, footer/slide-number/date
                                   furniture, brand marks
    grid        slideMaster        drawing guides, and the margin/column
                                   convention they imply

The output is a self-contained JSON document. Nothing downstream re-reads the
original master file, so a spec can be archived and replayed against a future
deck for the same client without resubmitting a master.

Verified against desktop PowerPoint (18/08/2026), because the guide format is
not in the ECMA base spec: PowerPoint stores per-master drawing guides in the
slideMaster's p:extLst as p15:sldGuideLst/p15:guide, positions in EIGHTHS OF
A POINT measured from the top-left slide edge. A missing pos attribute means
0; orient="horz" is a horizontal guide and a missing orient means vertical.

Known gap: legacy p:guideLst guides in ppt/viewProps.xml are NOT read. That
format could not be produced by the PowerPoint on this machine, so its units
are unverified, and emitting a silently wrong margin is worse than emitting
none. Decks whose guides live only there report grid.source = null.
"""

import argparse
import base64
import json
import statistics
from collections import Counter, defaultdict
from datetime import datetime, timezone
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from spike.color_resolver import (clr_map, color_scheme, resolve_color_element,
                                  resolve_solid_fill)
from spike.ns import find
from spike.resolver import _theme_element, theme_fonts
from .util import iter_shapes_deep

SPEC_VERSION = 1

P15 = "http://schemas.microsoft.com/office/powerpoint/2012/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
# Guides are stored in eighths of a point; a point is 12700 EMU.
GUIDE_UNIT_EMU = 12700 / 8

# Page furniture sits in the margins by design and must never teach the grid
# where the content area begins.
FURNITURE_PH = frozenset({"ftr", "sldNum", "dt"})


# ------------------------------------------------------------------- master


def dominant_master(prs):
    """The master most slides live on; the deck's first master when no slide
    resolves to one.

    The no-slides path is a normal case here, not an edge case: a submitted
    MASTER SLIDE file may carry a single sample slide or none at all, and its
    design surface still has to be readable."""
    masters = list(prs.slide_masters)
    if not masters:
        return None
    counts = Counter()
    for slide in prs.slides:
        try:
            part = slide.slide_layout.slide_master.part
        except Exception:
            continue
        for i, m in enumerate(masters):
            if m.part is part:
                counts[i] += 1
                break
    if not counts:
        return masters[0]
    return masters[counts.most_common(1)[0][0]]


def _hex(rgb) -> str:
    return "%02X%02X%02X" % tuple(rgb)


# -------------------------------------------------------------------- theme

_SCRIPTS = (("latin", "lt"), ("east_asian", "ea"), ("complex_script", "cs"))


def extract_theme(prs) -> dict:
    """What the theme part DECLARES: scheme colors by role, the master's color
    map (bg1 -> lt1 and friends), and the major/minor font per script.

    Read from theme1.xml, never inferred from usage. A deck whose every shape
    hardcodes literal RGB still has a theme, and restyling needs the roles,
    not the observed values. Empty dict when the deck has no master."""
    master = dominant_master(prs)
    if master is None:
        return {}
    fonts = theme_fonts(master)
    return {
        "colors": {slot: _hex(rgb) for slot, rgb in color_scheme(master).items()},
        "color_map": dict(clr_map(master)),
        "fonts": {
            "major": {script: fonts.get(f"+mj-{key}") for script, key in _SCRIPTS},
            "minor": {script: fonts.get(f"+mn-{key}") for script, key in _SCRIPTS},
        },
    }


# ------------------------------------------------------------- brand marks

# A brand mark is page furniture, not artwork: anything covering more than
# this share of the canvas is a background or a hero image.
MAX_LOGO_AREA_SHARE = 0.12
MAX_LOGO_EDGE_SHARE = 0.55
# A picture on no master and no layout has to recur across this many slides
# before it reads as a hand-stamped logo rather than as content.
MIN_LOGO_SLIDE_STAMPS = 3
# Occurrences further apart than the perceptual floor (~0.8mm, the same floor
# the alignment module uses) mean the mark has no single home position, so
# nothing downstream should propagate one.
LOGO_POSITION_TOLERANCE_EMU = 28575


def _image_key(shape):
    """sha1 of an embedded picture's bytes: identity that survives renaming,
    resizing, and re-cropping. None for anything that is not an embedded
    picture, including linked pictures, which carry no blob."""
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return None
    try:
        return shape.image.sha1
    except Exception:
        return None


def _furniture_sized(shape, sw, sh) -> bool:
    l, t, w, h = shape.left, shape.top, shape.width, shape.height
    if None in (l, t, w, h) or w <= 0 or h <= 0:
        return False
    if w * h > MAX_LOGO_AREA_SHARE * sw * sh:
        return False
    return w <= MAX_LOGO_EDGE_SHARE * sw and h <= MAX_LOGO_EDGE_SHARE * sh


def find_brand_marks(prs, include_slides: bool = True) -> list[dict]:
    """Embedded pictures that behave like brand furniture, strongest evidence
    first.

    Evidence ranks by WHERE a mark lives, because location is what makes it
    brand rather than content: a picture on the slide master is furniture by
    construction; a picture on layouts is furniture per layout; a picture only
    ever stamped on slides is a logo the master should have carried and does
    not, which is worth surfacing precisely because it is a defect.

    include_slides=False restricts the scan to the design surface, which is
    what Stage 1 wants when the submission is a master rather than a deck.

    v1 reads embedded pictures only. A logo drawn as grouped vector shapes or
    as a freeform is a known gap: identifying it needs shape-tree
    fingerprinting, not an image hash."""
    sw, sh = prs.slide_width, prs.slide_height
    master = dominant_master(prs)
    if master is None:
        return []

    marks = defaultdict(lambda: {"on_master": False, "layouts": [],
                                 "slide_count": 0, "boxes": [],
                                 "design_boxes": []})

    def record(shape, bucket, layout_name=None):
        key = _image_key(shape)
        if key is None or not _furniture_sized(shape, sw, sh):
            return
        m = marks[key]
        box = (shape.left, shape.top, shape.width, shape.height)
        m["boxes"].append(box)
        if bucket == "master":
            m["on_master"] = True
            m["design_boxes"].append(box)
        elif bucket == "layout":
            if layout_name not in m["layouts"]:
                m["layouts"].append(layout_name)
            m["design_boxes"].append(box)
        else:
            m["slide_count"] += 1

    for shape, _path in iter_shapes_deep(master.shapes):
        record(shape, "master")
    for layout in master.slide_layouts:
        for shape, _path in iter_shapes_deep(layout.shapes):
            record(shape, "layout", layout.name)
    if include_slides:
        for slide in prs.slides:
            for shape, _path in iter_shapes_deep(slide.shapes):
                record(shape, "slide")

    out = []
    for key, m in marks.items():
        if not m["on_master"] and not m["layouts"] \
                and m["slide_count"] < MIN_LOGO_SLIDE_STAMPS:
            continue
        # Position comes from the design surface when the mark lives there;
        # slide stamps only ever speak for themselves.
        boxes = m["design_boxes"] or m["boxes"]
        lefts = [b[0] for b in boxes]
        tops = [b[1] for b in boxes]
        out.append({
            "image_sha1": key,
            "scope": ("master" if m["on_master"]
                      else "layouts" if m["layouts"] else "slides"),
            "layouts": sorted(m["layouts"]),
            "slide_count": m["slide_count"],
            "position_emu": {
                "left": int(statistics.median(lefts)),
                "top": int(statistics.median(tops)),
                "width": int(statistics.median(b[2] for b in boxes)),
                "height": int(statistics.median(b[3] for b in boxes)),
            },
            "position_varies": (
                max(lefts) - min(lefts) > LOGO_POSITION_TOLERANCE_EMU
                or max(tops) - min(tops) > LOGO_POSITION_TOLERANCE_EMU),
        })
    out.sort(key=lambda m: (m["scope"] == "master", len(m["layouts"]),
                            m["slide_count"]), reverse=True)
    return out


def extract_brand(prs, include_slides: bool = True) -> dict:
    """The deck's logo plus the runners-up. `logo` is the pick so consumers
    read one field; `logo_alternates` keeps what was rejected visible to the
    design lead reviewing the spec."""
    marks = find_brand_marks(prs, include_slides=include_slides)
    return {"logo": marks[0] if marks else None, "logo_alternates": marks[1:]}


# ------------------------------------------------- placeholders and layouts


def _ph_token(shape) -> str | None:
    """The raw OOXML p:ph/@type token ('ctrTitle', 'body', 'ftr', 'sldNum').

    The raw token, not the python-pptx enum, because Stage 2 has to WRITE
    these back into generated layouts. OOXML omits the attribute for body
    placeholders, so an absent type means 'body'."""
    ph = find(shape._element, ".//p:nvSpPr/p:nvPr/p:ph")
    if ph is None:
        ph = find(shape._element, ".//p:nvPicPr/p:nvPr/p:ph")
    if ph is None:
        ph = find(shape._element, ".//p:nvGraphicFramePr/p:nvPr/p:ph")
    if ph is None:
        return None
    return ph.get("type", "body")


def _has_own_xfrm(shape) -> bool:
    """True when this shape carries its own a:xfrm rather than inheriting one.

    python-pptx resolves layout placeholder geometry against the master
    silently, so .left on a layout placeholder can be a master value. Stage 2
    needs to know which layouts actually pin their own geometry and which
    ride the master, because only the second kind follows a master edit."""
    return find(shape._element, "p:spPr/a:xfrm") is not None


def _ph_size_pt(shape):
    """The type size this placeholder DECLARES for its first outline level, or
    None when it inherits the master's.

    Read because layouts disagree: the client master measured on 20/08/2026
    declares 24, 25, 28 and 32pt for the title across its twelve layouts, so
    which title size a slide gets depends on which layout it landed on. That is
    a fact about the master, and a design lead can only fix it once - in the
    master - if the read shows it."""
    el = find(shape._element, "p:txBody/a:lstStyle/a:lvl1pPr/a:defRPr")
    if el is None or not el.get("sz"):
        return None
    try:
        return int(el.get("sz")) / 100.0
    except (TypeError, ValueError):
        return None


def _ph_autofit(shape) -> str | None:
    """"normAutofit" (shrink text on overflow), "noAutofit", or None. A title on
    shrink-to-fit is the second reason two slides show different title sizes:
    the longer one is scaled down at render time."""
    bodyPr = find(shape._element, "p:txBody/a:bodyPr")
    if bodyPr is None:
        return None
    for tag in ("a:normAutofit", "a:noAutofit", "a:spAutoFit"):
        if find(bodyPr, tag) is not None:
            return tag.split(":")[1]
    return None


def _placeholder_entry(shape) -> dict:
    return {
        "idx": shape.placeholder_format.idx,
        "type": _ph_token(shape),
        "name": shape.name,
        "position_emu": {"left": shape.left, "top": shape.top,
                         "width": shape.width, "height": shape.height},
        "geometry_source": "explicit" if _has_own_xfrm(shape) else "inherited",
        "size_pt": _ph_size_pt(shape),
        "autofit": _ph_autofit(shape),
    }


def _placeholders(container) -> list[dict]:
    out = []
    for ph in container.placeholders:
        try:
            out.append(_placeholder_entry(ph))
        except Exception:
            continue
    return sorted(out, key=lambda p: (p["idx"] is None, p["idx"]))


# A background image is part of the brand, so the spec carries its BYTES, not
# just a reference: the contract is that nothing downstream re-reads the master
# file. Past this cap the bytes are left out rather than bloating every spec,
# and the omission is stated in the spec so no consumer assumes it has them.
MAX_EMBEDDED_ASSET_BYTES = 4 * 1024 * 1024


def _pct(el, attr: str):
    """An OOXML percentage attribute (thousandths of a percent) as a float."""
    if el is None or el.get(attr) is None:
        return None
    return int(el.get(attr)) / 1000.0


def _rect_pct(el) -> dict | None:
    """a:fillRect / a:srcRect insets, as percentages per side."""
    if el is None:
        return None
    sides = {side: _pct(el, side) for side in ("l", "t", "r", "b")}
    if not any(v for v in sides.values()):
        return None  # an empty fillRect means "no inset", not "unknown"
    return sides


def _image_asset(container, rId: str, embed: bool) -> dict | None:
    """Identity, dimensions, and optionally the bytes of a related image."""
    if not rId:
        return None
    try:
        part = container.part.related_part(rId)
        img = part.image
    except Exception:
        # A linked (not embedded) picture, or a broken relationship. Say so
        # rather than pretending there is no image at all.
        return {"sha1": None, "unavailable": "image is linked or missing"}
    blob = img.blob
    asset = {
        "sha1": img.sha1,
        "format": img.ext,
        "content_type": img.content_type,
        "px": {"width": img.size[0], "height": img.size[1]},
        "dpi": list(img.dpi) if img.dpi else None,
        "bytes": len(blob),
        "data_base64": None,
        "embed_skipped": None,
    }
    if not embed:
        asset["embed_skipped"] = "asset embedding disabled for this extraction"
    elif len(blob) > MAX_EMBEDDED_ASSET_BYTES:
        asset["embed_skipped"] = (
            f"{len(blob) / 1e6:.1f} MB exceeds the "
            f"{MAX_EMBEDDED_ASSET_BYTES / 1e6:.0f} MB embed cap")
    else:
        asset["data_base64"] = base64.b64encode(blob).decode("ascii")
    return asset


def _blip_fill(blip_fill, container, embed: bool) -> dict:
    """How a picture fill renders, which is as much a design decision as the
    picture itself: a stretched photo, a tiled texture, and a 12%-alpha
    watermark are three different backgrounds sharing one image."""
    blip = find(blip_fill, "a:blip")
    rId = blip.get(f"{{{R_NS}}}embed") if blip is not None else None

    tile = find(blip_fill, "a:tile")
    stretch = find(blip_fill, "a:stretch")
    recolor = [tag.split(":")[1] for tag in
               ("a:duotone", "a:grayscl", "a:clrChange", "a:biLevel")
               if find(blip, tag) is not None]

    fill = {
        "mode": "tile" if tile is not None else "stretch" if stretch is not None
                else "unspecified",
        "stretch_insets_pct": _rect_pct(find(stretch, "a:fillRect")),
        "tile": None,
        "crop_pct": _rect_pct(find(blip_fill, "a:srcRect")),
        "alpha_pct": _pct(find(blip, "a:alphaModFix"), "amt"),
        "recolor": recolor,
    }
    if tile is not None:
        fill["tile"] = {
            "offset_emu": {"x": tile.get("tx"), "y": tile.get("ty")},
            "scale_pct": {"x": _pct(tile, "sx"), "y": _pct(tile, "sy")},
            "align": tile.get("algn"), "flip": tile.get("flip"),
        }
    return {"kind": "image", "fill": fill,
            "image": _image_asset(container, rId, embed)}


def _background(container, master, embed_assets: bool = True) -> dict | None:
    """The cSld background: a solid colour resolved to hex, a picture fill
    captured in full, a theme fill reference, or the kind of fill when no
    single colour exists."""
    bg = find(container._element, "p:cSld/p:bg")
    if bg is None:
        return None
    bgPr = find(bg, "p:bgPr")
    if bgPr is not None:
        rgb = resolve_solid_fill(bgPr, master)
        if rgb:
            return {"kind": "solid", "hex": _hex(rgb)}
        blip_fill = find(bgPr, "a:blipFill")
        if blip_fill is not None:
            return _blip_fill(blip_fill, container, embed_assets)
        for tag, kind in (("a:gradFill", "gradient"), ("a:pattFill", "pattern"),
                          ("a:noFill", "none")):
            if find(bgPr, tag) is not None:
                return {"kind": kind}
        return {"kind": "other"}
    bgRef = find(bg, "p:bgRef")
    if bgRef is not None:
        # bgRef carries its colour element directly, not wrapped in solidFill.
        # The colour is only half the story: idx points into the theme's
        # bgFillStyleLst, and that style can itself be a picture or gradient,
        # in which case the colour is just the phClr substituted into it.
        rgb = None
        for tag in ("a:schemeClr", "a:srgbClr"):
            el = find(bgRef, tag)
            if el is not None:
                rgb = resolve_color_element(el, master)
                break
        return {"kind": "theme_ref", "idx": bgRef.get("idx"),
                "hex": _hex(rgb) if rgb else None,
                "theme_fill_kind": _theme_bg_fill_kind(master, bgRef.get("idx"))}
    return None


def _theme_bg_fill_kind(master, idx) -> str | None:
    """What the theme's bgFillStyleLst entry behind a p:bgRef actually is.

    idx 1001+ indexes bgFillStyleLst (1-based); 1..999 indexes fillStyleLst;
    0 and 1000 mean no fill. Reported so a picture or gradient background
    inherited from the THEME is not silently flattened to its phClr colour."""
    try:
        i = int(idx)
    except (TypeError, ValueError):
        return None
    if i in (0, 1000):
        return "none"
    theme = _theme_element(master)
    lst_tag = "a:bgFillStyleLst" if i > 1000 else "a:fillStyleLst"
    offset = 1001 if i > 1000 else 1
    lst = find(theme, f".//a:fmtScheme/{lst_tag}")
    if lst is None:
        return None
    children = list(lst)
    pos = i - offset
    if not 0 <= pos < len(children):
        return None
    local = etree.QName(children[pos]).localname
    return {"solidFill": "solid", "blipFill": "image", "gradFill": "gradient",
            "pattFill": "pattern", "noFill": "none"}.get(local, local)


def extract_layouts(master, embed_assets: bool = True) -> list[dict]:
    """One entry per slideLayout: name, archetype, and placeholder geometry.

    `type` is the layout's own OOXML archetype token (title, obj, twoObj,
    secHead, blank, ...). That token is how Stage 2 matches a designer's
    layouts against the template bank, so it is recorded verbatim rather than
    re-derived from placeholder shapes."""
    out = []
    for i, layout in enumerate(master.slide_layouts):
        out.append({
            "index": i,
            "name": layout.name,
            "type": layout._element.get("type"),
            "placeholders": _placeholders(layout),
            "background": _background(layout, master, embed_assets),
            "shape_count": len(layout.shapes),
        })
    return out


# ------------------------------------------------------- master furniture


def _furniture(master, token: str) -> dict | None:
    """The master's footer / slide-number / date placeholder, if it has one.

    Slide-number and date placeholders hold an a:fld, not literal text, and
    python-pptx renders that field as private-use marker characters
    ('\\ue000#\\ue001'). Reporting those as `text` would put mojibake in the
    spec and invite Stage 2 to stamp it onto a layout, so a field is reported
    as a field and `text` stays null."""
    for ph in master.placeholders:
        if _ph_token(ph) != token:
            continue
        entry = {"present": True, "field": None, "text": None,
                 "position_emu": {"left": ph.left, "top": ph.top,
                                  "width": ph.width, "height": ph.height}}
        if getattr(ph, "has_text_frame", False):
            fld = find(ph._element, ".//a:fld")
            if fld is not None:
                entry["field"] = fld.get("type")
            else:
                entry["text"] = ph.text_frame.text.strip() or None
        return entry
    return {"present": False}


_TXSTYLE_ROLES = (("title", "p:titleStyle"), ("body", "p:bodyStyle"),
                  ("other", "p:otherStyle"))


def extract_text_styles(master) -> dict:
    """The master's p:txStyles at outline level 1: the type size and typeface
    the master DECLARES for titles, body copy, and everything else.

    This is a declaration, not a survey. It is where a title size target comes
    from when the submission is a master with no slides to measure. Theme font
    references (+mj-lt and friends) are resolved to real family names so the
    spec never hands a downstream stage a token it has to resolve itself."""
    fonts = theme_fonts(master)
    txStyles = find(master._element, "p:txStyles")
    out = {}
    for role, tag in _TXSTYLE_ROLES:
        defRPr = find(find(find(txStyles, tag), "a:lvl1pPr"), "a:defRPr")
        if defRPr is None:
            out[role] = None
            continue

        def typeface(child):
            el = find(defRPr, child)
            face = el.get("typeface") if el is not None else None
            return fonts.get(face, face) if face else None

        bold = defRPr.get("b")
        out[role] = {
            "size_pt": int(defRPr.get("sz")) / 100.0 if defRPr.get("sz") else None,
            "latin": typeface("a:latin"),
            "complex_script": typeface("a:cs"),
            "bold": bold in ("1", "true") if bold is not None else None,
        }
    return out


def extract_master(master, embed_assets: bool = True) -> dict:
    """Master-level fixed elements: background, page furniture, the declared
    text styles, and the placeholder geometry every layout inherits from."""
    return {
        "name": master.name if hasattr(master, "name") else None,
        "background": _background(master, master, embed_assets),
        "placeholders": _placeholders(master),
        "text_styles": extract_text_styles(master),
        "footer": _furniture(master, "ftr"),
        "slide_number": _furniture(master, "sldNum"),
        "date": _furniture(master, "dt"),
    }


# --------------------------------------------------------------- grid/guides


def read_guides(master) -> dict:
    """Drawing guides from the master's p15:sldGuideLst, in EMU.

    Format verified against desktop PowerPoint (18/08/2026): positions are in
    eighths of a point from the top-left slide edge, a missing pos means 0,
    orient="horz" is horizontal and a missing orient means vertical."""
    vertical, horizontal = [], []
    for g in master._element.iter(f"{{{P15}}}guide"):
        pos = int(g.get("pos", "0")) * GUIDE_UNIT_EMU
        if g.get("orient") == "horz":
            horizontal.append(int(round(pos)))
        else:
            vertical.append(int(round(pos)))
    return {"vertical_emu": sorted(vertical), "horizontal_emu": sorted(horizontal)}


# ---------------------------------------------------- presentation space
#
# The frame a designer DRAWS on the master to say where content may live: a
# rectangle named "presentation space". It outranks everything else here,
# because everything else is a reading of intent and this IS the intent.
#
# Asked for by the design lead (21/08/2026): "some cases have multiple margins,
# so presentation space is safer". A master can carry an outer page margin, a
# column grid, a band under the header and a bleed line, and picking which of
# those a block should seat against is a guess the tool keeps getting wrong. One
# named rectangle ends the guessing.

PRESENTATION_SPACE_NAMES = ("presentationspace", "presentationarea",
                            "contentspace", "contentarea")
# A frame drawn a hair off the canvas edge is rounding, not a mistake.
EDGE_SLACK = 36000              # 1mm


def _norm_name(name: str) -> str:
    return "".join(ch for ch in (name or "").lower() if ch.isalnum())


def _space_from(container, sw, sh):
    """The presentation-space rectangle among a container's own shapes."""
    for shape in container.shapes:
        if _norm_name(shape.name) not in PRESENTATION_SPACE_NAMES:
            continue
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
        if None in (l, t, w, h) or w <= 0 or h <= 0:
            continue
        # A marker is meant to be invisible, and a rectangle a designer just
        # drew is not: its fill and line come from the shape STYLE (p:style
        # fillRef / lnRef), with nothing in spPr to find. So invisibility has to
        # be stated - an explicit noFill on both - and anything else prints on
        # every slide of every deck, which the read says out loud.
        spPr = find(shape._element, "p:spPr")
        fill_off = find(spPr, "a:noFill") is not None
        ln = find(spPr, "a:ln")
        line_off = ln is not None and find(ln, "a:noFill") is not None
        # A frame has to be ON the page. A rectangle hanging off the canvas -
        # left over from a resized master, or drawn on the wrong slide size -
        # would hand every downstream pass a negative margin, so it is read,
        # reported and NOT used.
        problem = None
        if l < -EDGE_SLACK or t < -EDGE_SLACK \
                or l + w > sw + EDGE_SLACK or t + h > sh + EDGE_SLACK:
            problem = ("the rectangle is not inside the slide, so it cannot be "
                       "a content frame; check it against this master's slide "
                       "size")
        return {"left": l, "top": t, "right": sw - (l + w),
                "bottom": sh - (t + h), "box_emu": [l, t, l + w, t + h],
                "prints": not (fill_off and line_off), "name": shape.name,
                "problem": problem}
    return None


def read_presentation_space(prs, master) -> dict | None:
    """The designer's declared content frame, as margins from each edge, or
    None when the master does not carry one.

    Looked for on the MASTER first, because a frame that governs the deck
    belongs where the deck's every slide inherits from. A marker found only on a
    layout is still read - a designer trying this out will put it wherever seems
    natural - and `source` says where it came from so the read can ask for it to
    be moved."""
    if master is None:
        return None
    sw, sh = prs.slide_width, prs.slide_height
    space = _space_from(master, sw, sh)
    if space:
        space["source"] = "master"
        return space
    for layout in master.slide_layouts:
        space = _space_from(layout, sw, sh)
        if space:
            space["source"] = f"layout '{layout.name}'"
            return space
    return None


# A guide within this of the canvas centre line is a CENTRE guide: masters
# carry one on each axis as a placement aid, and it says nothing about where
# content begins.
CENTER_GUIDE_TOL_EMU = 9525


def read_content_band(horizontal: list, slide_h: int) -> tuple:
    """(subtitle_floor_emu, body_top_emu): the pair of horizontal guides a
    master draws under its subtitle to reserve a strip of white space between
    the header and the body.

    The convention, confirmed against the client masters (20/08/2026): the
    outermost horizontal guides are the page's top and bottom margins, and
    between them sit two more - the floor the SUBTITLE may not cross, and the
    ceiling the BODY may not cross. The strip between those two stays empty on
    every slide, which is what makes a deck's headers read as one line down the
    deck. On the master measured, the master's own body placeholder starts
    within 0.01in of the second guide, which is the master stating the same
    thing twice.

    Read only when the master states it UNAMBIGUOUSLY. Two interior guides in
    the top half are the band; one is a stated body ceiling with no reserved
    strip; anything else (none, or three and more) is not a band this can name,
    and guessing which pair a designer meant would seat content on a line
    nobody drew. Centre guides are dropped first: a guide on the canvas centre
    line is a placement aid present on most masters, and it would otherwise
    read as a body ceiling half way down the slide."""
    if len(horizontal) < 3:
        return (None, None)
    interior = [g for g in horizontal[1:-1]
                if g < slide_h / 2
                and abs(g - slide_h // 2) > CENTER_GUIDE_TOL_EMU]
    if len(interior) == 2:
        return (interior[0], interior[1])
    if len(interior) == 1:
        return (None, interior[0])
    return (None, None)


def _infer_columns(vertical, tolerance_emu: int = 9525):
    """(columns, gutter_emu) when vertical guides form an even column grid.

    A designer's column grid marks both edges of every column, so the gaps
    alternate column, gutter, column, gutter. Anything that does not alternate
    evenly is not a column grid and gets no guess: a wrong column count is
    worse for Stage 2 than an absent one."""
    if len(vertical) < 4 or len(vertical) % 2:
        return None, None
    widths = [vertical[i + 1] - vertical[i] for i in range(0, len(vertical) - 1, 2)]
    gutters = [vertical[i + 1] - vertical[i] for i in range(1, len(vertical) - 1, 2)]
    if not gutters or min(widths) <= 0 or min(gutters) <= 0:
        return None, None
    if max(widths) - min(widths) > tolerance_emu:
        return None, None
    if max(gutters) - min(gutters) > tolerance_emu:
        return None, None
    return len(widths), int(statistics.median(gutters))


def _content_extent(container):
    """Bounding box of the container's CONTENT placeholders. Page furniture
    is excluded: footers and slide numbers live in the margins by design and
    would drag the content area out to the slide edges."""
    boxes = []
    for ph in container.placeholders:
        if _ph_token(ph) in FURNITURE_PH:
            continue
        l, t, w, h = ph.left, ph.top, ph.width, ph.height
        if None in (l, t, w, h) or w <= 0 or h <= 0:
            continue
        boxes.append((l, t, w, h))
    if not boxes:
        return None
    return (min(b[0] for b in boxes), min(b[1] for b in boxes),
            max(b[0] + b[2] for b in boxes), max(b[1] + b[3] for b in boxes))


def infer_grid(prs, master) -> dict:
    """The margin and column convention the master implies.

    Three sources, best first, and `source` records which was used so a reviewer
    knows how much to trust the numbers:

    1. the PRESENTATION SPACE, a rectangle the designer drew and named. Not an
       inference at all, which is why it wins: a master can carry an outer page
       margin, a column grid, a header band and a bleed line, and choosing among
       them is the guess this ends (design lead, 21/08/2026).
    2. drawing GUIDES, a stated intention but one the tool still has to
       interpret - which pair is the frame, which is the band, which is the
       centre line.
    3. the master's own content PLACEHOLDERS, an inference from where its title
       happens to sit.

    The header band always comes from the guides: it is a pair of lines inside
    the frame, and a single rectangle cannot state it."""
    sw, sh = prs.slide_width, prs.slide_height
    guides = read_guides(master)
    v, h = guides["vertical_emu"], guides["horizontal_emu"]
    space = read_presentation_space(prs, master)
    out = {"guides": guides, "margins_emu": None, "columns": None,
           "gutter_emu": None, "source": None,
           "subtitle_floor_emu": None, "body_top_emu": None,
           "presentation_space": space}

    # The band is read whenever the guides state one, whatever the frame source.
    if len(h) >= 3:
        out["subtitle_floor_emu"], out["body_top_emu"] = read_content_band(h, sh)

    if space is not None and not space.get("problem"):
        out["margins_emu"] = {side: space[side]
                              for side in ("left", "top", "right", "bottom")}
        out["columns"], out["gutter_emu"] = _infer_columns(v)
        out["source"] = "presentation_space"
        return out

    if len(v) >= 2 and len(h) >= 2:
        out["margins_emu"] = {"left": v[0], "right": sw - v[-1],
                              "top": h[0], "bottom": sh - h[-1]}
        out["columns"], out["gutter_emu"] = _infer_columns(v)
        # The top margin is where the PAGE begins; the body ceiling is where
        # CONTENT begins, and on these masters they are two different lines
        # with the header band between them (see read_content_band).
        out["subtitle_floor_emu"], out["body_top_emu"] = read_content_band(h, sh)
        out["source"] = "guides"
        return out

    box = _content_extent(master)
    if box is not None:
        left, top, right, bottom = box
        out["margins_emu"] = {"left": left, "top": top,
                              "right": sw - right, "bottom": sh - bottom}
        out["source"] = "placeholders"
    return out


# ------------------------------------------------------------------- output


def extract_style_spec(prs, source: str | None = None,
                       embed_assets: bool = True) -> dict:
    """The canonical Stage 1 artifact. Self-contained by design: nothing
    downstream re-reads the master file."""
    master = dominant_master(prs)
    spec = {
        "spec_version": SPEC_VERSION,
        "meta": {
            "source_file": source,
            "extracted_at": datetime.now(timezone.utc).isoformat(timespec="seconds"),
            "slide_size_emu": {"width": prs.slide_width, "height": prs.slide_height},
            "master_count": len(list(prs.slide_masters)),
            "slide_count": len(prs.slides),
        },
        "theme": extract_theme(prs),
        "master": (extract_master(master, embed_assets)
                   if master is not None else None),
        "layouts": (extract_layouts(master, embed_assets)
                    if master is not None else []),
        # A master submission has no slides to hand-stamp anything on, so the
        # design surface is the whole story.
        "brand": extract_brand(prs, include_slides=False),
        "grid": infer_grid(prs, master) if master is not None else None,
    }
    return spec


# ----------------------------------------------------- spec -> audit profile

# Palette slots worth enforcing. Hyperlink colors are deliberately excluded:
# they are link states, not brand palette, and flagging text for using them
# would be noise.
_PALETTE_SLOTS = ("dk1", "lt1", "dk2", "lt2", "accent1", "accent2", "accent3",
                  "accent4", "accent5", "accent6")

# Perceptual-floor tolerances, the same ones bootstrap uses. Professionally
# formatted decks carry sub-millimetre offsets no eye can see; flagging them
# is pure noise (ground-truth calibration, 19/07/2026).
VISUAL_TOLERANCE_EMU = 28575          # 0.79mm
INTENT_WINDOW_EMU = 137160            # 0.15in


def spec_to_profile(spec: dict, profile_id: str, name: str) -> dict:
    """Project a Style Spec onto a formatting profile the audit engine reads.

    The spec is canonical; the profile is a VIEW of it shaped for the existing
    modules. Nothing here re-reads the master, which is the decoupling the
    pipeline depends on: an archived spec can be replayed against a future
    deck for the same client without the master ever being resubmitted.

    Mapping choices, deliberately conservative:
    - Heading roles take the theme's MAJOR font, body roles the MINOR font,
      because that is what the theme means by major and minor.
    - Size targets come from the master's declared txStyles, and only for
      title and subtitle. Body copy legitimately varies across a real deck, so
      a single body target would flood the report (same reasoning as
      bootstrap, which reaches it from the opposite direction).
    - The palette is the theme's own slots, each carrying its theme_ref, so a
      later theme edit moves the palette with it.
    - Layouts are allow-listed by name: a master's layouts are the layouts its
      decks are supposed to use.
    """
    theme = spec.get("theme") or {}
    fonts = theme.get("fonts") or {}
    major = fonts.get("major") or {}
    minor = fonts.get("minor") or {}
    grid = spec.get("grid") or {}
    master = spec.get("master") or {}
    styles = master.get("text_styles") or {}
    size = spec.get("meta", {}).get("slide_size_emu", {})
    sw = size.get("width") or 12192000
    sh = size.get("height") or 6858000

    def families(source, fallback):
        latin = [f for f in (source.get("latin"), fallback.get("latin")) if f]
        cs = [f for f in (source.get("complex_script"),
                          fallback.get("complex_script")) if f]
        # dict.fromkeys keeps first-seen order while dropping duplicates
        return list(dict.fromkeys(latin)), list(dict.fromkeys(cs))

    title_style = styles.get("title") or {}
    body_style = styles.get("body") or {}
    roles = {}
    for role, theme_font, declared in (("title", major, title_style),
                                       ("subtitle", major, title_style),
                                       ("body", minor, body_style),
                                       ("caption", minor, body_style)):
        latin, cs = families(declared, theme_font)
        entry = {"latin": latin, "complex_script": cs,
                 "allowed_weights": ["regular", "bold"]}
        if role in ("title", "subtitle") and declared.get("size_pt"):
            entry["size_pt"] = declared["size_pt"]
        roles[role] = entry

    colors = theme.get("colors") or {}
    named = [{"name": slot, "hex": colors[slot], "theme_ref": slot,
              "allowed_tints": [], "allowed_shades": []}
             for slot in _PALETTE_SLOTS if slot in colors]

    margins = grid.get("margins_emu") or {
        "left": round(sw * 0.0375), "right": round(sw * 0.0375),
        "top": round(sh * 0.040), "bottom": round(sh * 0.053),
    }
    # Carried only when the master DREW it. There is no sensible default for
    # where a body begins: a guessed ceiling would be a line the client never
    # drew, and every slide would be measured against it.
    band = None
    if grid.get("body_top_emu"):
        band = {"subtitle_floor": grid.get("subtitle_floor_emu"),
                "body_top": grid["body_top_emu"]}
    footer = master.get("footer") or {}
    sldnum = master.get("slide_number") or {}

    return {
        "id": profile_id,
        "name": name,
        "client_scope": None, "project_scope": None,
        "is_default": False, "version": 1,
        "owner": (f"from master {spec.get('meta', {}).get('source_file')}, "
                  f"review with design lead"),
        "config": {
            "theme": theme,
            "brand": spec.get("brand") or {},
            "style_spec_source": {
                "spec_version": spec.get("spec_version"),
                "extracted_at": spec.get("meta", {}).get("extracted_at"),
                "grid_source": grid.get("source"),
            },
            "master_slide": {
                "enforce_existing_only": True, "pinned_layout_id": None,
                "layout_allowlist": [lay["name"] for lay in spec.get("layouts", [])],
                "geometry_tolerance_emu": 9525,
            },
            "font": {"roles": roles, "theme_font_refs_allowed": True,
                     "size_tolerance_pt": 0.5},
            "color_palette": {
                "theme_color_slots": list(_PALETTE_SLOTS),
                "named_colors": named,
                "on_palette_mode": "by_name",
                "match_tolerance_deltaE": 2.0,
                "auto_replace_max_deltaE": 5.0,
                "ambiguity_band_deltaE": 10.0,
            },
            "geometry": {
                "safe_zone_margins_emu": margins,
                "body_band_emu": band,
                "grid": {"columns": grid.get("columns") or 12,
                         "gutter_emu": grid.get("gutter_emu") or 0,
                         "enabled": bool(grid.get("columns"))},
                "alignment": {"edge_tolerance_emu": VISUAL_TOLERANCE_EMU,
                              "center_tolerance_emu": VISUAL_TOLERANCE_EMU,
                              "spacing_tolerance_emu": VISUAL_TOLERANCE_EMU,
                              "intent_window_emu": INTENT_WINDOW_EMU},
            },
            "shape_size": {"size_tolerance_emu": 9525, "min_cohort_size": 3,
                           "preserve_picture_aspect": True,
                           "dominant_size_strategy": "median"},
            "header_footer": {"template": {
                "footer_text": footer.get("text"),
                "slide_number": bool(sldnum.get("present")),
                "date": {"enabled": False, "format": "DD/MM/YYYY"},
                "position_emu": None, "font_role": "caption"}},
            "typography": {"label_case": None},
        },
    }


def main(argv=None):
    ap = argparse.ArgumentParser(
        description="Extract a Style Spec from a submitted master .pptx.")
    ap.add_argument("master", help="Path to the .pptx carrying the master")
    ap.add_argument("--json", dest="json_out", default=None,
                    help="Write the full Style Spec JSON here")
    args = ap.parse_args(argv)

    path = Path(args.master)
    spec = extract_style_spec(Presentation(path), source=path.name)

    theme = spec["theme"]
    print(f"\nStyle Spec: {path.name}  |  spec v{spec['spec_version']}")
    print(f"  canvas: {spec['meta']['slide_size_emu']['width']} x "
          f"{spec['meta']['slide_size_emu']['height']} EMU  |  "
          f"masters: {spec['meta']['master_count']}  |  "
          f"layouts: {len(spec['layouts'])}")
    if theme.get("colors"):
        accents = [f"{k}={v}" for k, v in theme["colors"].items()
                   if k.startswith("accent")]
        print(f"  theme accents: {' '.join(accents)}")
        mj, mn = theme["fonts"]["major"], theme["fonts"]["minor"]
        print(f"  theme fonts: major={mj['latin']}/{mj['complex_script']}  "
              f"minor={mn['latin']}/{mn['complex_script']}")

    grid = spec["grid"] or {}
    if grid.get("source"):
        m = grid["margins_emu"]
        cols = (f", {grid['columns']} columns / {grid['gutter_emu']} EMU gutter"
                if grid["columns"] else "")
        print(f"  grid ({grid['source']}): L{m['left']} R{m['right']} "
              f"T{m['top']} B{m['bottom']}{cols}")
        g = grid["guides"]
        print(f"  guides: {len(g['vertical_emu'])} vertical, "
              f"{len(g['horizontal_emu'])} horizontal")
        if grid.get("body_top_emu"):
            floor = grid.get("subtitle_floor_emu")
            band = (f", reserved strip {(grid['body_top_emu'] - floor) / 36000:.0f}mm "
                    f"under the subtitle floor at {floor / 914400:.2f}in"
                    if floor else " (no reserved strip stated)")
            print(f"  body begins: {grid['body_top_emu'] / 914400:.2f}in{band}")
    else:
        print("  grid: no guides and no master placeholders to infer from")

    mst = spec["master"] or {}
    print(f"  background: {mst.get('background')}")
    for key in ("footer", "slide_number", "date"):
        f = mst.get(key) or {}
        if f.get("present"):
            extra = f" text={f['text']!r}" if f.get("text") else ""
            print(f"  {key}: present{extra}")
    logo = spec["brand"]["logo"]
    print(f"  logo: {logo['image_sha1'][:12]} on {logo['scope']}" if logo
          else "  logo: none on the design surface")

    print("\n  layouts:")
    for lay in spec["layouts"]:
        explicit = sum(1 for p in lay["placeholders"]
                       if p["geometry_source"] == "explicit")
        print(f"    [{lay['index']:>2}] {lay['name'][:34]:34} "
              f"type={str(lay['type']):9} "
              f"ph={len(lay['placeholders'])} ({explicit} explicit)")

    if args.json_out:
        Path(args.json_out).write_text(
            json.dumps(spec, indent=2, ensure_ascii=False), encoding="utf-8")
        print(f"\nStyle Spec written: {args.json_out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
