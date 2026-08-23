"""Put one reported change back, exactly.

The migration pass is a series of decisions a designer signs off, and a review
that can only say "no" to the whole deck is not a review. This is the other
half: every change qc.migrate makes carries the STATE it changed
(ContentChange.undo), and this replays that state onto the finished deck.

Three properties are the whole design, and each was chosen against an
alternative that looks simpler and is wrong:

STATE, NEVER A DELTA. An offset operation carries the coordinates the shape
held, not the distance it travelled. Re-deriving "so move it back up 1.90in"
means computing the same thing twice and hoping the two agree; when they do not,
the deck a designer approved is quietly wrong. With the coordinates stored, undo
is an assignment.

LAST CHANGE UNDONE FIRST. Changes overlap: the block move and the collision
nudge both touch the same shape, the nudge running on the position the move left
it in. Replaying them front to back puts the shape back to its pre-move
coordinates and then immediately forward again to its pre-nudge ones, which are
post-move - so the shape ends up exactly where it started and the undo silently
does nothing (real deck, 23/08/2026: 5 of 721 shapes came back wrong, all of
them nudged after being moved). Reversing the order is what makes undo mean
undo. Each is then retired by the caller, so a resubmitted form cannot
double-apply one.

WHAT CANNOT BE UNDONE SAYS SO. The layout assignment is PowerPoint's own work
(qc.applymaster drives desktop PowerPoint to do it), so it cannot be taken back
one slide at a time from here - the honest answer is to say that on the page and
offer the original upload, not to fake a reversal by moving shapes around.
"""

import io

from pptx import Presentation

from .migrate import _covers, _renumber, _text_of


def _by_id(slide):
    """Every shape on the slide by id, GROUP CHILDREN INCLUDED.

    Top-level only was enough while undo served the migration, which never
    reaches inside a group. The design pass does - a run recoloured for contrast
    is often inside a badge or a card - and a shape this cannot find is an undo
    that silently reports success and changes nothing, which is the one outcome
    qc.design must never produce.
    """
    from .util import iter_shapes_deep

    return {str(s.shape_id): s for s, _path in iter_shapes_deep(slide.shapes)}


def _op_offset(slide, op) -> str | None:
    """Put one shape back at the coordinates it held."""
    shape = _by_id(slide).get(str(op.get("shape_id")))
    if shape is None:
        return None
    was = (shape.left, shape.top)
    shape.left, shape.top = int(op["left"]), int(op["top"])
    if was == (shape.left, shape.top):
        return None
    return (f"{(_text_of(shape)[:20] or shape.name)!r} back to "
            f"{op['left'] / 914400:.2f}in, {op['top'] / 914400:.2f}in")


def _op_replace(slide, op) -> str | None:
    """Swap a shape's element for the one stored, in place.

    In place matters: a placeholder put back by delete-then-insert lands at the
    END of the shape tree, so it prints over the content instead of under it,
    and PowerPoint no longer treats it as the layout's placeholder in the same
    order. Replacing the element keeps its position in the tree and therefore
    its z-order.

    Parsed with python-pptx's own parser rather than lxml's, for the reason
    qc.migrate.restore_shapes gives: a generic lxml element spliced into a
    python-pptx tree is a shape the library can no longer read."""
    from pptx.oxml import parse_xml

    shape = _by_id(slide).get(str(op.get("shape_id")))
    if shape is None or not op.get("xml"):
        return None
    try:
        element = parse_xml(op["xml"])
    except Exception:
        return None
    old = shape._element
    old.addprevious(element)
    old.getparent().remove(old)
    return f"{op.get('label') or 'placeholder'} back as it was"


def _op_insert(slide, op) -> str | None:
    """Put a removed element back, on top of the slide.

    The same three hardenings as qc.migrate.restore_shapes, and for the same
    reasons: insert BEFORE p:extLst (a shape tree's extLst must stay last or
    PowerPoint offers to repair the file), renumber the subtree (a stale shape
    id is a duplicate id, which reads as damage rather than as a duplicate), and
    report what the returned piece now prints over rather than hunting for a
    clear spot - where the designer put it is where it goes back.

    An optional `index` puts it back at the position in the shape tree it was
    removed from, so it returns to its own place in the drawing order rather
    than on top of everything. The migration's removals omit it and keep the
    original behaviour: a swept header remnant coming back over the content is
    exactly the alert that says the sweep was wrong. A design-pass deletion is
    the opposite case - the designer chose to remove a badge and is putting it
    back where it was, not raising it above the deck."""
    from pptx.oxml import parse_xml

    if not op.get("xml"):
        return None
    try:
        element = parse_xml(op["xml"])
    except Exception:
        return None
    first_id = slide.shapes._next_shape_id
    _renumber(element, first_id)
    spTree = slide.shapes._spTree
    index = op.get("index")
    if index is not None and 0 <= int(index) <= len(spTree):
        spTree.insert(int(index), element)
    else:
        spTree.insert_element_before(element, "p:extLst")
    shape = next((s for s in slide.shapes if str(s.shape_id) == str(first_id)),
                 None)
    if shape is None:
        return "put back"
    covers = _covers(slide, shape)
    label = _text_of(shape)[:24] or shape.name
    if covers:
        return (f"{label!r} put back where it was, printing over "
                f"{', '.join(covers[:3])}"
                + (f" and {len(covers) - 3} more" if len(covers) > 3 else ""))
    return f"{label!r} put back where it was"


def _op_bg(slide, op) -> str | None:
    """Put the slide's own p:bg back as the first child of p:cSld.

    Position is not cosmetic here: the schema orders cSld's children and p:bg
    is the first of them, so appending it produces a file PowerPoint repairs."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import qn

    if not op.get("xml"):
        return None
    cSld = slide._element.find(qn("p:cSld"))
    if cSld is None:
        return None
    try:
        element = parse_xml(op["xml"])
    except Exception:
        return None
    for existing in cSld.findall(qn("p:bg")):
        cSld.remove(existing)
    cSld.insert(0, element)
    return "the slide's own background is back, so it beats the master's again"


def _op_zorder(slide, op) -> str | None:
    """Put one shape back at the position in the drawing order it held.

    Stored as the raw child index of the shape tree, not as "above shape 7".
    A relative record has to be re-derived against a tree that may have moved
    on since, and the re-derivation is a second guess at the thing already
    known; the index is what the shape actually had. Reordering never changes
    how many children the tree has, so the index a reorder was recorded against
    stays valid for it - and apply_undo replays last-first, so two reorders on
    one slide unwind in the order that makes each of them true.
    """
    shape = _by_id(slide).get(str(op.get("shape_id")))
    if shape is None or op.get("index") is None:
        return None
    element = shape._element
    parent = element.getparent()
    if parent is None:
        return None
    index = int(op["index"])
    if not (0 <= index <= len(parent)):
        return None
    if list(parent).index(element) == index:
        return None
    parent.remove(element)
    parent.insert(index, element)
    return f"{(_text_of(shape)[:20] or shape.name)!r} back to its old " \
           f"position in the drawing order"


_OPS = {"offset": _op_offset, "replace": _op_replace, "insert": _op_insert,
        "bg": _op_bg, "zorder": _op_zorder}


def apply_undo(deck_bytes: bytes, items: list) -> tuple[bytes, list[dict]]:
    """Replay the stored state for each item onto the deck.

    items: [{"change_id", "slide_index", "action", "ops": [...]}] - normally
    read straight off the ContentChange objects the format run produced, IN THE
    ORDER THE RUN REPORTED THEM. They are replayed back to front (see above);
    outcomes come back in the order they were given, so a caller can zip them
    against its own list.

    Returns the new deck and one outcome per item: {change_id, done, detail}.
    `done` is False for an item nothing could be applied from - a shape the
    deck no longer has, XML that will not parse - and the detail says so
    rather than reporting a success the deck does not show. A review that
    claims to have undone something it did not is worse than one that refuses.
    """
    prs = Presentation(io.BytesIO(deck_bytes))
    outcomes = []
    for item in reversed(list(items)):
        index = item.get("slide_index")
        ops = item.get("ops") or []
        if index is None or not (0 <= index < len(prs.slides)):
            outcomes.append({"change_id": item.get("change_id"), "done": False,
                             "detail": "that slide is no longer in the deck"})
            continue
        slide = prs.slides[index]
        notes = []
        for op in ops:
            handler = _OPS.get(op.get("op"))
            if handler is None:
                continue
            try:
                note = handler(slide, op)
            except Exception as exc:
                note = None
                notes.append(f"failed: {type(exc).__name__}: {exc}")
            if note:
                notes.append(note)
        done = bool(notes) and not all(n.startswith("failed") for n in notes)
        if not notes:
            notes = ["nothing to change: the deck no longer holds what this "
                     "change touched"]
        outcomes.append({"change_id": item.get("change_id"),
                         "slide_index": index, "done": done,
                         "detail": "; ".join(notes[:4])})
    out = io.BytesIO()
    prs.save(out)
    outcomes.reverse()  # back into the caller's order
    return out.getvalue(), outcomes


def undoable(change) -> bool:
    """Whether the review page may offer an Undo for this change."""
    return bool(getattr(change, "undo", None))


def followers(changes: list, change_id: str) -> list:
    """The changes that have to come back with this one: itself and every later
    undoable change ON THE SAME SLIDE.

    Not an extra helpfulness - a correctness requirement. The pass works down a
    slide in order and every step is computed on the state the steps before it
    left, so undoing one in the middle produces a slide the pass never made. The
    client's Gantt slide is the plain version of it: two column headings were
    swept, and the block was then seated on the frame WITHOUT them. Putting the
    headings back on their own returned them to their original coordinates while
    the table stayed 0.55in lower, so the headings landed on the eyebrow instead
    of on their columns (design lead, 23/08/2026, "the kept text should also go
    back to its place instead of overlapping them").

    Other slides are untouched. Changes on one slide say nothing about another,
    and taking a whole deck back is what the original upload is for.
    """
    index = next((i for i, c in enumerate(changes)
                  if getattr(c, "change_id", None) == change_id), None)
    if index is None:
        return []
    slide = changes[index].slide_index
    return [c for c in changes[index:]
            if c.slide_index == slide and getattr(c, "undo", None)]


def expand(changes: list, wanted_ids) -> list:
    """Every change the requested ones drag with them, in the order the run
    reported them (which is the order apply_undo expects)."""
    ids = set()
    for change_id in wanted_ids or ():
        ids.update(c.change_id for c in followers(changes, change_id))
    return [c for c in changes if getattr(c, "change_id", None) in ids
            and getattr(c, "undo", None)]


__all__ = ["apply_undo", "expand", "followers", "undoable"]
