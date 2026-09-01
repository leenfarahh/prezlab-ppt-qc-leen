"""Shared traversal helpers for audit modules."""

import statistics
from collections import Counter, defaultdict

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

TITLE_TYPES = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
HEADING_TYPES = TITLE_TYPES + (PP_PLACEHOLDER.SUBTITLE,)


def recurring_anchors(prs, min_slides: int = 3, pos_bin: int = 36000,
                      size_bin: int = 72000) -> dict:
    """Cross-slide ANCHORS: top-level shapes whose (name, size class)
    recur at one modal position on min_slides+ slides - title bars, header
    chips, logos. Returns {(slide_index, shape_id_str): (modal_l, modal_t)}
    for every occurrence of each anchored key.

    Real-deck finding (12/08/2026, RTL strategy deck): local cluster fixes
    nudged per-slide copies of the SAME title bar by different amounts, so
    titles that used to sit identically across slides came back scattered.
    Anchored shapes are excluded from local alignment pools and instead
    snapped to their deck-wide modal position."""
    occ = defaultdict(list)
    for s_idx, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
            if None in (l, t, w, h):
                continue
            key = (sh.name, round(w / size_bin), round(h / size_bin))
            occ[key].append((s_idx, str(sh.shape_id), l, t))
    out = {}
    for key, items in occ.items():
        if len({i[0] for i in items}) < min_slides:
            continue
        bins = Counter((round(i[2] / pos_bin), round(i[3] / pos_bin))
                       for i in items)
        modal_bin, n = bins.most_common(1)[0]
        if n < min_slides:
            continue  # recurs but wanders: no position to hold anyone to
        in_bin = [i for i in items
                  if (round(i[2] / pos_bin), round(i[3] / pos_bin)) == modal_bin]
        ml = int(statistics.median(i[2] for i in in_bin))
        mt = int(statistics.median(i[3] for i in in_bin))
        for s_idx, sid, _l, _t in items:
            out[(s_idx, sid)] = (ml, mt)
    return out


# --------------------------------------------------------- page-deep panels
#
# Asked by both the audit and the migration, so they cannot disagree about what
# counts as body content. A background image, a full-height sidebar or a bleed
# band runs the depth of the page by design: it must not set where the body
# begins, and it must not be dragged when the body moves (real-master finding,
# 20/08/2026 - a 94%-height image starting on the top margin asked for a 36mm
# body move in the audit and pushed a whole slide off the canvas in the
# migration).

FULL_SPAN_SHARE = 0.85       # spans most of the canvas height
EDGE_SLACK_EMU = 72000       # 2mm: designers park bleed ON the edge


def full_height_panel(top, height, slide_h) -> bool:
    """True for a shape anchored to the top or bottom edge that spans most of
    the canvas height. Edge-anchored AND page-deep, because either alone is
    ordinary: a card column can be tall, and a chip can sit on the edge."""
    if top is None or height is None:
        return False
    touches = (top <= EDGE_SLACK_EMU
               or top + height >= slide_h - EDGE_SLACK_EMU)
    return touches and height >= FULL_SPAN_SHARE * slide_h


# ------------------------------------------------------------- collections
#
# A designer does not place a photo, its corner rule, its quote mark, its
# paragraph and its caption. They place a COLLECTION, and the spacing inside it
# IS the composition (design lead, 20/08/2026, comparing a before/after of the
# leadership-quotes slide: the tool had moved members of each collection
# separately and the brackets no longer met their photos).
#
# So a collection is an ANCHOR plus its SATELLITES: the substantial shape, and
# the smaller shapes that ride it. Satellites attach to an anchor, never to each
# other, which is what keeps a slide of touching decorative shapes from
# collapsing into one 114-shape "collection" (measured on the client's own
# sample deck: pure proximity chaining does exactly that).
#
# Both the audit and the migration ask this question, and a positional fix that
# answered it differently from the check that raised it is how a composition
# comes apart.

SATELLITE_GAP_EMU = 360000     # 10mm: a caption, rule or mark beside its anchor

# How much of a satellite has to sit ON its anchor before the two are one
# object. Half, because the cases this rule exists for - a rule across a photo,
# a chip on its box, a label on its card - are at or near the whole of the
# satellite, and the case it must exclude - two text boxes whose empty margins
# lap over each other - is a fraction of one.
WELD_MIN_SHARE = 0.5
BACKDROP_SLIDE_SHARE = 0.35    # covering this much canvas: a backdrop, not a
                               # partner, and it rides nothing


def rides_with(sat, anchor, along: str | None = None,
               gap: int = SATELLITE_GAP_EMU) -> bool:
    """True when `sat` composes with `anchor` and must travel with it. Boxes
    are (l, t, r, b).

    Overlapping shapes are welded: a corner rule drawn across a photo is part
    of it whatever else is true. Otherwise membership is judged ACROSS the axis
    of the move, which is the distinction the whole engine turns on:

        along="x" (a horizontal move) carries what is stacked with the anchor -
        the caption under it, the chip above it - because those share its column
        and a horizontal move without them tears the pair apart (real-deck
        finding, 12/08/2026, RTL strategy map).

        along="y" (a vertical move) carries what sits beside it - the label to
        its right, the underline - because those share its row.

    The other direction is deliberately NOT carried: shapes along the axis of
    the move are its spacing PEERS, and dragging them is how a rhythm fix
    cancels itself (the image column would move as a lump and the odd gap would
    survive). `along=None` carries both, for a fix that moves on both axes."""
    sl, st, sr, sb = sat
    al, at, ar, ab = anchor
    if sr <= sl or sb <= st or ar <= al or ab <= at:
        return False
    ox = min(sr, ar) - max(sl, al)
    oy = min(sb, ab) - max(st, at)
    if ox > 0 and oy > 0:
        # OVERLAPPING IS NOT THE SAME AS TOUCHING. Any overlap at all used to
        # weld the two, and on a deck of text boxes that is most of the slide:
        # a box is routinely much bigger than the words in it, so a full-width
        # title's box laps over the header below it while nothing visible
        # touches. The engine then treated the title as a passenger of the
        # header, and a fix asked to nudge the header 0.3in moved the title
        # too - into the next column's header, which the collision guard
        # correctly refused, so the whole fix was dropped and the page said
        # "Applied 0 fixes" (reproduced 01/09/2026 on a two-column slide).
        #
        # A weld is when the satellite LIVES on the anchor: a corner rule drawn
        # across a photo, a chip on its box, a label on its card - the overlap
        # is most of the satellite. A graze is measured against the satellite's
        # own area rather than the anchor's because the question is whether
        # THIS shape would be torn by being left behind, and a shape with a
        # corner over the mover would not.
        if ox * oy >= WELD_MIN_SHARE * (sr - sl) * (sb - st):
            return True                  # composed by definition
        # Otherwise it is only a graze, and it still gets the ordinary
        # stacked/beside test below: a caption whose box laps its image is
        # carried for sharing its column, not for touching it.
    cx, cy = (sl + sr) // 2, (st + sb) // 2
    stacked = al <= cx <= ar and -oy <= gap     # same column, small gap
    beside = at <= cy <= ab and -ox <= gap      # same row, small gap
    if along == "x":
        return stacked
    if along == "y":
        return beside
    return stacked or beside


def is_backdrop(box, slide_w, slide_h) -> bool:
    """A shape big enough to be the ground the content sits on rather than a
    member of any collection. It overlaps everything, so without this a fix on
    one small label would drag the background image with it."""
    l, t, r, b = box
    if r <= l or b <= t:
        return False
    return ((r - l) * (b - t) >= BACKDROP_SLIDE_SHARE * slide_w * slide_h
            or full_height_panel(t, b - t, slide_h))


# ------------------------------------------------------------ RTL slides
#
# Margins mirror under RTL: the RIGHT margin is where an Arabic block starts,
# the way the left margin is for English. Binding the left edge on an Arabic
# slide leaves the reading edge ragged, which is the first thing a reader of
# the language sees (design lead, 20/08/2026).

RTL_TEXT_SHARE = 0.5           # more Arabic than latin letters: an RTL slide


def slide_is_rtl(slide) -> bool:
    """True when this slide's text reads right-to-left.

    Judged per SLIDE, not per deck: a bilingual deck runs English and Arabic
    slides side by side, and the deck-level answer would mirror the wrong ones.
    An explicit rtl paragraph is taken at its word; otherwise the script that
    carries more of the slide's letters wins, so one English product name on an
    Arabic slide does not flip it back."""
    from spike.arabic import ARABIC_RANGES, paragraph_is_rtl

    arabic = latin = 0
    for shape, _path in iter_shapes_deep(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        for para in shape.text_frame.paragraphs:
            if paragraph_is_rtl(para):
                return True
            for ch in para.text:
                code = ord(ch)
                if any(lo <= code <= hi for lo, hi in ARABIC_RANGES):
                    arabic += 1
                elif ch.isalpha():
                    latin += 1
    return arabic > 0 and arabic >= RTL_TEXT_SHARE * (arabic + latin)


def iter_shapes_deep(shapes, path: str = ""):
    """Yield (shape, shape_path) for every shape, descending into groups.
    shape_path is the group ancestry ('12/3' = shape 3 inside group 12),
    None for top-level shapes (Appendix A.2)."""
    for shape in shapes:
        shape_path = path or None
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield shape, shape_path
            child_path = f"{path}/{shape.shape_id}" if path else str(shape.shape_id)
            yield from iter_shapes_deep(shape.shapes, child_path)
        else:
            yield shape, shape_path


# --------------------------------------------------------------- headings
#
# A heading that runs past a margin is never corrected by this tool: whether a
# title may break the margin is a house-style decision the client owns, so the
# only right move is to flag it and let the designer ask. Both the audit and
# the migration need the same answer to "is this shape a heading", and they
# must not disagree, so the question is answered once here.

# Bounds for inferring a heading on a deck that carries no placeholders (every
# export-tool deck). Deliberately tight: a wrong "this is the title" quietly
# changes which findings fire, so a heading has to look like one.
HEADING_BAND = 0.25             # top quarter of the slide
HEADING_MAX_AREA_SHARE = 0.10   # a line of text, not a panel
HEADING_MAX_CHARS = 200
HEADING_MIN_LEAD = 1.25         # visibly bigger than the text it heads
HEADING_MIN_TEXT_SHAPES = 3     # hierarchy needs something to be the top of


def max_font_pt(shape, slide, prs) -> float:
    """Largest effective point size in a shape, resolved through the OOXML
    cascade (python-pptx returns None for inherited sizes)."""
    from spike.resolver import resolve_run

    best = 0.0
    if not getattr(shape, "has_text_frame", False):
        return best
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue
            try:
                size = resolve_run(run, para, shape, slide, prs).size_pt.value
            except Exception:
                size = None
            if size:
                best = max(best, size)
    return best


def _heading_candidates(slide, prs) -> list:
    """(scale, top, shape) for shapes that could be a line of header text.

    Point sizes and box heights are never mixed inside one comparison: with
    every size inherited the cascade resolves nothing, and comparing a resolved
    28pt against a 0.6in proxy ranks by unit rather than by size."""
    cands = []
    slide_area = prs.slide_width * prs.slide_height
    for shape in slide.shapes:  # top level only; a heading is not group content
        if getattr(shape, "is_placeholder", False):
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if not text or len(text) > HEADING_MAX_CHARS:
            continue
        top, height, width = shape.top, shape.height, shape.width
        if None in (top, height, width):
            continue
        if width * height > HEADING_MAX_AREA_SHARE * slide_area:
            continue
        cands.append([max_font_pt(shape, slide, prs), top,
                      height / 12700.0, shape])
    if not cands:
        return []
    if any(c[0] for c in cands):
        # Once ANY size resolves, the ones that did not are dropped rather than
        # ranked at zero: an unresolvable shape would otherwise sit at the
        # bottom of the pool and drag the median with it, so the shape above it
        # would clear the "visibly bigger" bar by default.
        return [(c[0], c[1], c[3]) for c in cands if c[0]]
    return [(c[2], c[1], c[3]) for c in cands]


def heading_ids(slide, prs) -> set:
    """shape_ids (as str) of the shapes carrying this slide's title and
    subtitle.

    Placeholders answer this outright, and when the slide has them that is the
    only answer taken: a placeholder role is the deck's own statement of what a
    shape is, not an inference (the same restraint `font_role` keeps).

    An export-tool deck has no placeholders at all, so a bounded fallback
    applies the designer's own rule - the largest type on the slide is the
    heading - but only where there is real hierarchy to read: at least three
    lines of text, the candidate visibly bigger than the rest, sitting in the
    top band, and shaped like a line rather than a panel. A lone text box is
    not a heading just because nothing else competes with it."""
    ids = set()
    try:
        for ph in slide.placeholders:
            if ph.placeholder_format.type in HEADING_TYPES:
                ids.add(str(ph.shape_id))
    except Exception:
        pass
    if ids:
        return ids

    cands = _heading_candidates(slide, prs)
    if len(cands) < HEADING_MIN_TEXT_SHAPES:
        return ids
    ranked = sorted(cands, key=lambda c: (-c[0], c[1]))
    title = ranked[0]
    rest = [c[0] for c in ranked[1:]]
    if title[1] >= HEADING_BAND * prs.slide_height:
        return ids  # the biggest type is not at the top: no title to read
    if title[0] < HEADING_MIN_LEAD * statistics.median(rest):
        return ids  # everything is the same size: no hierarchy stated
    ids.add(str(title[2].shape_id))
    for scale, top, shape in ranked[1:]:
        # the standfirst: the next size down, still in the band, under the title
        if (scale < title[0] and top >= title[1]
                and top < HEADING_BAND * prs.slide_height):
            ids.add(str(shape.shape_id))
            break
    return ids


def font_role(shape) -> str:
    """Map a shape to a profile font role. v1 inference is deliberately
    conservative: title/subtitle placeholders map to their roles, everything
    else is body. Caption is never inferred (PRD: hierarchy inference beyond
    placeholders is v2)."""
    if shape.is_placeholder:
        t = shape.placeholder_format.type
        if t in TITLE_TYPES:
            return "title"
        if t == PP_PLACEHOLDER.SUBTITLE:
            return "subtitle"
    return "body"
