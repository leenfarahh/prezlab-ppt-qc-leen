"""Local pilot web UI for Flow B (audit only).

    python -m qc.web            # http://127.0.0.1:8000

Deliberately a LOCAL PILOT, not the production deployment: binds loopback by
default, no auth (PRD forbids network exposure without the Entra sign-in
gate), uploads are audited from a temp file and deleted immediately after
processing, manifests are held in memory only. The production web tier (job
queue, storage, SSO) is PRD Section 6 and comes after the residency ruling.
"""

import io
import json
import os
import tempfile
import threading
import uuid
import zipfile
from collections import OrderedDict
from contextlib import asynccontextmanager
from pathlib import Path

from fastapi import FastAPI, File, Form, Request, UploadFile
from fastapi.responses import (HTMLResponse, JSONResponse, RedirectResponse,
                               Response)
from fastapi.staticfiles import StaticFiles

from .engine import run_audit
from .profile import PROFILES_DIR, Profile
from .records import MODULES
from .ui import render_index, render_report

from .config import DEMO_BANNER
from .config import MAX_UPLOAD_MB as _MAX_UPLOAD_MB

MAX_UPLOAD_BYTES = _MAX_UPLOAD_MB * 1024 * 1024  # configurable (QC_MAX_UPLOAD_MB)
# A .pptx is a zip; XML compresses ~100x, so a small upload can expand to GBs
# in RAM when parsed (review finding: 0.26MB deck -> ~206MB peak). Bound the
# declared uncompressed size and the compression ratio before parsing. The
# absolute uncompressed cap scales with the upload cap; the ratio check is the
# real zip-bomb defense regardless of absolute size.
MAX_UNCOMPRESSED_BYTES = MAX_UPLOAD_BYTES * 6
MAX_COMPRESSION_RATIO = 120.0
RATIO_CHECK_FLOOR = 64 * 1024 * 1024  # small decks legitimately compress well
# What to tell a designer when no model key is set. Named once: the provider is
# configurable, so hard-coding one vendor's variable in a message is how the
# component review ended up asking for a key it does not use.
_MODEL_KEY_NOTE = "set GEMINI_API_KEY in the .env file at the project root and restart."

MAX_STORED_MANIFESTS = 50
# Deck bytes are retained in memory for the newest jobs only, so "apply
# selected fixes" can run without re-upload. Older jobs keep their manifest
# but lose the deck (fix option expires). Nothing is written to disk.
MAX_DECKS_IN_MEMORY = 5

# Every route outside the public set requires a signed-in session when set
# (PRD: no network-reachable endpoint without an auth gate). Driven by env
# (cloud) or the --lan flag (LAN pilot).
from .config import AI_ENABLED
from .config import AUTH_REQUIRED as _ENV_AUTH_REQUIRED

AUTH_REQUIRED = _ENV_AUTH_REQUIRED
_PUBLIC_PATHS = ("/signin", "/signout", "/health", "/me")
_PUBLIC_PREFIXES = ("/static/",)

def _bootstrap_admin():
    """Create the first-run admin when the user table is empty.

    Auth is mandatory on a cloud deploy, so a fresh one with no users would lock
    everybody out including the person who deployed it. No effect once any user
    exists."""
    from .config import BOOTSTRAP_ADMIN
    from .store import add_user, list_users

    if BOOTSTRAP_ADMIN and not list_users():
        add_user(BOOTSTRAP_ADMIN, "admin")


@asynccontextmanager
async def _lifespan(_app):
    """Startup and shutdown. `@app.on_event("startup")` is deprecated in this
    FastAPI and warns twice on every import; a lifespan handler is the
    supported form and is what the next major version will require."""
    _bootstrap_admin()
    yield


app = FastAPI(title="Prezlab PPT QC", docs_url=None, redoc_url=None,
              lifespan=_lifespan)
app.state.auth_required = _ENV_AUTH_REQUIRED


@app.middleware("http")
async def require_signin(request: Request, call_next):
    # app.state has one identity even under `python -m` (which loads this
    # module twice, as __main__ and qc.web); the module global alone would
    # miss that. The global is still honored for the test monkeypatch path.
    active = getattr(request.app.state, "auth_required", False) or AUTH_REQUIRED
    if active:
        path = request.url.path
        public = path in _PUBLIC_PATHS or any(
            path.startswith(p) for p in _PUBLIC_PREFIXES)
        if not public:
            from .auth import current_user

            if current_user(request) is None:
                if request.method == "GET":
                    from fastapi.responses import RedirectResponse

                    return RedirectResponse("/signin", status_code=303)
                return JSONResponse({"error": "sign in required"}, status_code=401)
    return await call_next(request)
app.mount("/static", StaticFiles(directory=str(Path(__file__).parent / "static")),
          name="static")
# job_id -> {"manifest": dict, "deck": bytes|None, "cleaned": bytes|None,
#            "filename": str, "profile": str}
_jobs: OrderedDict[str, dict] = OrderedDict()
_jobs_lock = threading.Lock()


from .web_admin import register as _register_admin

_register_admin(app)


def _profiles_meta() -> list[dict]:
    """[{id, name}] for the profile picker, read from the profile JSONs."""
    out = []
    for p in sorted(PROFILES_DIR.glob("*.json")):
        try:
            data = json.loads(p.read_text(encoding="utf-8"))
            out.append({"id": p.stem, "name": data.get("name", p.stem)})
        except Exception:
            out.append({"id": p.stem, "name": p.stem})
    return out


@app.middleware("http")
async def reject_oversized_bodies(request: Request, call_next):
    """Reject oversized uploads from the declared Content-Length BEFORE
    Starlette spools the whole body to disk (review finding: the in-handler
    cap only fires after the full body has already been received)."""
    if request.method == "POST":
        declared = request.headers.get("content-length")
        if declared and declared.isdigit() and int(declared) > MAX_UPLOAD_BYTES + 1_048_576:
            mb = int(declared) // 1_048_576
            msg = (f"That upload is about {mb} MB, over this instance's "
                   f"{_MAX_UPLOAD_MB} MB cap.")
            if DEMO_BANNER:
                msg += (" Large client decks belong on the office LAN "
                        "instance, which accepts up to 250 MB and keeps "
                        "files on-site.")
            return HTMLResponse(
                render_index(_pickable_profiles(), MODULES, msg),
                status_code=413)
    return await call_next(request)


def _zip_bomb_reason(data: bytes) -> str | None:
    """Non-None when the archive declares an implausible expansion. A broken
    zip returns None here; python-pptx raises and the 422 path handles it."""
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            total = sum(i.file_size for i in z.infolist())
    except zipfile.BadZipFile:
        return None
    if total > MAX_UNCOMPRESSED_BYTES:
        return f"declares {total / 1e6:.0f} MB uncompressed (limit {MAX_UNCOMPRESSED_BYTES / 1e6:.0f} MB)"
    if total > RATIO_CHECK_FLOOR and total / max(len(data), 1) > MAX_COMPRESSION_RATIO:
        return "implausible compression ratio for a presentation"
    return None


def _profiles() -> list[str]:
    return sorted(p.stem for p in PROFILES_DIR.glob("*.json"))


# Audit a deck against rules derived from the deck itself (Sense 1): no brand
# needed, flags slides that break from the deck's own dominant conventions.
SELF_PROFILE = "__self__"
_SELF_META = {"id": SELF_PROFILE, "label": "Match the deck itself",
              "name": "No brand needed; flags slides that break from the "
                      "deck's own conventions"}

# Audit a deck against a master submitted in the same request: Stage 1 reads
# the master's design surface, and the resulting Style Spec becomes the rules
# for this one job. Nothing is written to disk, so a designer with no admin
# rights can do it; saving it as a reusable profile is a separate, gated step
# on the /master page.
MASTER_PROFILE = "__master__"
_MASTER_META = {"id": MASTER_PROFILE, "label": "Match a master I'll upload",
                "name": "Reads the master's theme, layouts, and grid, then "
                        "checks the deck against them"}

# Neither sentinel corresponds to a profile file, so anything that tunes or
# reloads a SAVED profile has to refuse them.
_EPHEMERAL = (SELF_PROFILE, MASTER_PROFILE)


def _ephemeral(profile_key: str) -> bool:
    return profile_key in _EPHEMERAL


# Model-backed features are refused at the route, not just hidden in the
# UI: hiding a button is a presentation choice, and QC_AI=0 has to mean no
# request leaves this machine even for someone POSTing the endpoint directly.
_AI_OFF = {"error": "AI features are disabled on this instance (QC_AI=0). "
                    "No request is sent to the model API."}


def _ai_disabled_response():
    return JSONResponse(_AI_OFF, status_code=503) if not AI_ENABLED else None


def _pickable_profiles() -> list[dict]:
    return [_SELF_META, _MASTER_META] + _profiles_meta()


def _attachment(filename: str) -> dict:
    """Content-Disposition that survives non-ASCII names (RFC 5987): an
    ASCII fallback plus filename* in UTF-8. HTTP headers are latin-1, so a
    bare Arabic filename crashes the response (real-deck finding,
    12/08/2026: every download/export died on an Arabic deck name)."""
    from urllib.parse import quote

    fallback = "".join(c if 31 < ord(c) < 127 and c not in '"\\' else "_"
                       for c in filename).strip() or "download"
    return {"Content-Disposition":
            f'attachment; filename="{fallback}"; '
            f"filename*=UTF-8''{quote(filename, safe='')}"}


def _resolve_profile(profile: str, data: bytes) -> Profile:
    """A Profile object for the chosen id, or one bootstrapped from the deck
    itself when the self-consistency option is picked."""
    if profile == SELF_PROFILE:
        from pptx import Presentation

        from .bootstrap import build_profile

        prs = Presentation(io.BytesIO(data))
        return Profile(build_profile(prs, "self", "This deck (self-consistency)"))
    return Profile.load(profile)


@app.get("/", response_class=HTMLResponse)
def index():
    return render_index(_pickable_profiles(), MODULES)


@app.get("/health")
def health():
    return {"status": "ok"}


# ------------------------------------------------------- Stage 1: the master
# Reading a master is a different job from auditing a deck: the input is the
# design surface, the output is a Style Spec, and the spec (not the file) is
# what everything downstream consumes. Specs are held in memory beside jobs,
# under the same rule as decks: nothing is written to disk.
MAX_SPECS_IN_MEMORY = 20
_specs: OrderedDict[str, dict] = OrderedDict()
_specs_lock = threading.Lock()


def _remember_spec(spec: dict, master_bytes: bytes | None = None) -> str:
    """Hold the spec and the master it came from. The master is kept because a
    spec DESCRIBES a design system and cannot be applied as one: restyling a
    slide needs real slideLayout parts for PowerPoint to match placeholders
    against. Saving the spec as a profile hands these bytes to the template
    store; until then they live in memory like everything else here."""
    spec_id = uuid.uuid4().hex
    with _specs_lock:
        _specs[spec_id] = {"spec": spec, "master": master_bytes}
        while len(_specs) > MAX_SPECS_IN_MEMORY:
            _specs.popitem(last=False)
    return spec_id


def _get_spec(spec_id: str) -> dict | None:
    with _specs_lock:
        held = _specs.get(spec_id)
    return held["spec"] if held else None


def _get_spec_master(spec_id: str) -> bytes | None:
    with _specs_lock:
        held = _specs.get(spec_id)
    return held["master"] if held else None


def _prep_intake_page(request: Request, message: str = "", status: int = 200,
                      saved: str = "", spec: dict | None = None,
                      spec_id: str = "", spec_message: str = "",
                      replaced: bool = False):
    """The ONE prepare intake: read a master, save it as a profile, apply it.

    Every step of Prepare a deck answers on this page, including the ones that
    used to have pages of their own (reading a master, checking coverage,
    applying a master). `spec` is how the read master gets back here instead of
    onto a page of its own."""
    from .ui_prep import render_prep_intake
    from .unify import com_available
    from .web_admin import _editor

    can_look, why = _can_look()
    return HTMLResponse(render_prep_intake(
        _formattable_profiles(), message=message, com_ready=com_available(),
        look=can_look, look_note=why, saved=saved,
        can_save=_editor(request) is not None, spec=spec, spec_id=spec_id,
        spec_message=spec_message, replaced=replaced), status_code=status)


@app.post("/master", response_class=HTMLResponse)
def master_read(request: Request, master: UploadFile = File(...)):
    """Step 1: read the master, and answer on the page that asked.

    A form endpoint, not a page. There is no GET /master any more: the master
    is dropped on Prepare a deck and the spec renders under the form."""
    from pptx import Presentation

    from .stylespec import extract_style_spec

    filename = master.filename or "master.pptx"
    if not filename.lower().endswith(".pptx"):
        return _prep_intake_page(request, "Only .pptx files are accepted.", 400)

    data = master.file.read(MAX_UPLOAD_BYTES + 1)
    if len(data) > MAX_UPLOAD_BYTES:
        return _prep_intake_page(
            request, f"File exceeds the {_MAX_UPLOAD_MB} MB cap.", 413)
    bomb = _zip_bomb_reason(data)
    if bomb:
        return _prep_intake_page(request, f"File rejected: {bomb}.", 413)

    try:
        spec = extract_style_spec(Presentation(io.BytesIO(data)), source=filename)
    except Exception as exc:
        return _prep_intake_page(
            request,
            f"Could not read that master: {type(exc).__name__}: {exc}", 422)

    return _prep_intake_page(request, spec=spec,
                             spec_id=_remember_spec(spec, data))


@app.post("/spec/{spec_id}/pspace")
def spec_stamp_pspace(spec_id: str, left: str = Form(...), top: str = Form(...),
                      right: str = Form(...), bottom: str = Form(...)):
    """Hand back a COPY of this master carrying a presentation-space rectangle.

    The frame a master states is the single thing that ends the guessing
    downstream: without one, the content area is inferred from where the
    master's own placeholders happen to sit, and every deck formatted against it
    is seated on that inference (qc.stylespec.infer_grid). This writes the
    designer's own numbers in, so the next run reads a stated frame.

    A DOWNLOAD, not an edit. The stored master is untouched and nothing is
    written to the template store: the designer opens the copy, checks it, and
    re-uploads it if they want it. A client's master is not a file this tool
    changes on its own (qc.pspace.stamp_master).
    """
    from .pspace import stamp_master

    data = _get_spec_master(spec_id)
    if data is None:
        return JSONResponse(
            {"error": "that master is no longer held in memory; read it again"},
            status_code=410)
    try:
        box = tuple(int(round(float(v) * 914400))
                    for v in (left, top, right, bottom))
    except (TypeError, ValueError):
        return JSONResponse(
            {"error": "the four edges have to be numbers, in inches"},
            status_code=400)
    try:
        out, note = stamp_master(data, box)
    except ValueError as exc:
        return JSONResponse({"error": str(exc)}, status_code=400)
    except Exception as exc:
        return JSONResponse(
            {"error": f"that master could not be stamped: "
                      f"{type(exc).__name__}: {exc}"}, status_code=422)

    spec = _get_spec(spec_id) or {}
    source = (spec.get("meta") or {}).get("source_file") or "master"
    stem = Path(source).stem or "master"
    return Response(
        content=out,
        media_type="application/vnd.openxmlformats-officedocument."
                   "presentationml.presentation",
        headers={**_attachment(f"{stem}-presentation-space.pptx"),
                 # The note is what the designer needs and a file download has
                 # nowhere to put it, so it rides on a header the page reads.
                 "X-QC-Note": note[:900]})


@app.get("/spec/{spec_id}.json")
def spec_json(spec_id: str):
    """The Style Spec itself: the canonical artifact, downloadable so it can be
    archived and replayed against a future deck without the master."""
    spec = _get_spec(spec_id)
    if spec is None:
        return JSONResponse({"error": "spec expired or unknown"}, status_code=404)
    source = (spec.get("meta") or {}).get("source_file") or "master"
    stem = Path(source).stem or "master"
    return JSONResponse(spec, headers={
        "Content-Disposition": f'attachment; filename="{stem}-stylespec.json"'})


@app.post("/spec/{spec_id}/profile", response_class=HTMLResponse)
def spec_to_profile_route(request: Request, spec_id: str, name: str = Form(""),
                          target: str = Form("")):
    """Turn the master just read into a profile: a new one, or an existing one
    pointed at this file.

    `target` is the second door and it is the one that stops the profile list
    filling up with near-duplicates. Revising a master and reading it again is
    the normal case, and creating "Client X 2" for it means half the team
    prepares decks on last month's file."""
    from datetime import date

    from fastapi.responses import RedirectResponse

    from .profile import PROFILES_DIR
    from .stylespec import spec_to_profile
    from .web_admin import (_editor, _slugify, _unique_pid,
                            replace_profile_master)

    # Every refusal comes back to Prepare a deck rather than to an error page,
    # and the two that still have a spec bring it with them: a designer who
    # mistyped the name should not have to read the master again to retry.
    spec = _get_spec(spec_id)
    if spec is None:
        return _prep_intake_page(
            request, "That spec has expired. Specs are held in memory only, "
                     "so read the master again.", 404)

    editor = _editor(request)
    if editor is None:
        return _prep_intake_page(
            request, "Saving a profile needs a lead or admin. Sign in, then "
                     "save it; the master stays read below.", 403,
            spec=spec, spec_id=spec_id)

    if target:
        # The master bytes, not the spec: a profile has to carry a real .pptx
        # for PowerPoint to match placeholders against, and a spec describes a
        # design system rather than being one.
        master_bytes = _get_spec_master(spec_id)
        if not master_bytes:
            return _prep_intake_page(
                request, "That master file is no longer held in memory, so "
                         "there is nothing to replace it with. Read it again.",
                410)
        if target not in _profiles():
            return _prep_intake_page(request, "Unknown profile.", 404,
                                     spec=spec, spec_id=spec_id)
        source = (spec.get("meta") or {}).get("source_file") or "master.pptx"
        saved, note = replace_profile_master(target, master_bytes, source,
                                             editor["name"])
        if not saved:
            return _prep_intake_page(request, "", 422, spec=spec,
                                     spec_id=spec_id, spec_message=note)
        # Redirect rather than render: this POST is not idempotent (it bumps the
        # profile's version), and a designer who refreshes the result should not
        # walk the version number up one press at a time.
        return RedirectResponse(f"/prep?saved={target}&replaced=1",
                                status_code=303)

    slug = _slugify(name)
    if not slug:
        return _prep_intake_page(
            request, "", 400, spec=spec, spec_id=spec_id,
            spec_message="Give the profile a name with letters or numbers.")

    pid = _unique_pid(slug)
    profile = spec_to_profile(spec, pid, name.strip())
    profile["owner"] = (f"{editor['name']} (from master "
                        f"{(spec.get('meta') or {}).get('source_file')}, "
                        f"{date.today().strftime('%d/%m/%Y')})")
    (PROFILES_DIR / f"{pid}.json").write_text(
        json.dumps(profile, indent=2, ensure_ascii=False), encoding="utf-8")

    # The rules alone cannot restyle a slide; keep the master file so this
    # profile can be applied, not just audited against.
    master_bytes = _get_spec_master(spec_id)
    if master_bytes:
        from .templates import save_master

        save_master(pid, master_bytes)
    return RedirectResponse(f"/prep?saved={pid}", status_code=303)


def _profile_from_master(master: UploadFile | None):
    """(Profile, spec, error_message) for the upload-a-master option.

    The master is read for its design surface, projected onto a profile, and
    then dropped: only the derived rules survive into the job. That is the
    Style Spec contract holding at the web tier, not just in the library."""
    from pptx import Presentation

    from .stylespec import extract_style_spec, spec_to_profile

    if master is None or not master.filename:
        return None, None, ("Pick the master .pptx to check this deck against, "
                            "or choose a different rule source.")
    if not master.filename.lower().endswith(".pptx"):
        return None, None, "The master must be a .pptx file."

    raw = master.file.read(MAX_UPLOAD_BYTES + 1)
    if len(raw) > MAX_UPLOAD_BYTES:
        return None, None, f"The master exceeds the {_MAX_UPLOAD_MB} MB cap."
    bomb = _zip_bomb_reason(raw)
    if bomb:
        return None, None, f"The master was rejected: {bomb}."

    try:
        spec = extract_style_spec(Presentation(io.BytesIO(raw)),
                                  source=master.filename)
    except Exception as exc:
        return None, None, (f"Could not read that master: "
                            f"{type(exc).__name__}: {exc}")

    stem = Path(master.filename).stem or "master"
    # A readable id, because it lands in the report header and in history.
    return (Profile(spec_to_profile(spec, f"master:{stem}",
                                    f"Master: {master.filename}")),
            spec, None)


# ----------------------------------------------- Stage 3: apply the master
# Auditing reads a deck; this REWRITES it, so it gets its own page rather than
# hiding behind the audit form. Results are held in memory like every other
# job: nothing is written to disk except the profile's stored master.
MAX_FORMAT_JOBS = 10
_format_jobs: OrderedDict[str, dict] = OrderedDict()
_format_lock = threading.Lock()

# Decks read against a master and waiting on their layout decision. Distinct
# from a format job, and deliberately so: a Plan has no rebuilt deck, no
# manifest and nothing to download, and putting one in _format_jobs would make
# every page that reads a job have to ask whether this one is real yet.
#
# HALF THE CEILING OF A FORMAT JOB, because a Plan is not cheap and it is not
# finished work. It holds the upload, the master, and the PNGs of every slide
# and layout on its page - most of what a job holds, for something nobody has
# pressed Apply on yet. Ten of these alongside ten format jobs is two client
# decks in memory per run. They expire oldest first, and a designer who leaves
# the layout page open past that gets "start again" rather than a stack trace.
MAX_PLANS = 5
_plans: OrderedDict[str, dict] = OrderedDict()
_plans_lock = threading.Lock()


def _remember_plan_stage(plan, profile: str, profile_obj,
                         master_bytes: bytes) -> str:
    """Hold a planned deck between step 1 and step 2, and return its id.

    The MASTER BYTES are held with it. Re-reading the profile's stored master
    on submit would be one fewer thing in memory and would also mean a master
    replaced between the two presses gets applied to layouts chosen against the
    old one - which is a corrupted rebuild that nothing on the page would
    explain."""
    plan_id = uuid.uuid4().hex
    with _plans_lock:
        _plans[plan_id] = {"plan": plan, "profile": profile,
                           "profile_obj": profile_obj, "master": master_bytes,
                           "shots": {}, "layout_thumbs": {},
                           "render_note": ""}
        while len(_plans) > MAX_PLANS:
            _plans.popitem(last=False)
    return plan_id


def _get_plan(plan_id: str) -> dict | None:
    with _plans_lock:
        return _plans.get(plan_id)


def _drop_plan(plan_id: str) -> None:
    """Once the rebuild has run, the Plan is dead weight: its plans live on the
    job and its source bytes are held there too."""
    with _plans_lock:
        _plans.pop(plan_id, None)


def _formattable_profiles() -> list[dict]:
    """Every saved profile, flagged with whether it carries a master. Profiles
    without one are still listed by the page so the absence is explained
    rather than looking like a missing option.

    Each one also reports WHICH master it would apply: when the stored file was
    stored, and what frame that file states. A profile whose stored copy
    predates the designer's latest master formats decks on the old one, and the
    only symptom is something missing from the output - a presentation-space
    rectangle, in the case that made this necessary (design lead, 21/08/2026).
    Read at the point of use, because that is where the surprise happens."""
    from datetime import datetime

    from .stylespec import dominant_master, extract_layouts, infer_grid
    from .templates import has_master, load_master, master_info

    out = []
    for meta in _profiles_meta():
        pid = meta["id"]
        n_layouts = 0
        frame = stored = None
        if has_master(pid):
            try:
                from pptx import Presentation

                prs = Presentation(io.BytesIO(load_master(pid)))
                master = dominant_master(prs)
                n_layouts = len(extract_layouts(master, embed_assets=False)) \
                    if master is not None else 0
                if master is not None:
                    frame = infer_grid(prs, master).get("source")
                info = master_info(pid)
                if info:
                    stored = datetime.fromtimestamp(
                        info["modified"]).strftime("%d/%m/%Y")
            except Exception:
                n_layouts = 0
        out.append({"id": pid, "name": meta["name"],
                    "has_master": has_master(pid), "layouts": n_layouts,
                    "frame": frame, "master_stored": stored})
    return out


# ------------------------------------------- pre-flight: layout coverage
# The same question the format result answers, asked BEFORE anything is
# rewritten. Applying a master is a rewrite that needs desktop PowerPoint and
# takes about a second a slide; learning afterwards that the master had no
# layout for a third of the deck means running the whole pass twice. This route
# writes nothing, stores nothing, and returns no file.


def _can_look() -> tuple[bool, str]:
    """Whether the slides can be looked at on this host, and why not when they
    cannot.

    Four separate reasons, each with its own sentence, because a designer is
    entitled to know which of them they got: "switched off here" and "no key"
    are somebody's decision, and the other two are a broken setting.

    WHAT IS LOST IS THE JUDGMENT, NOT THE REPORT. Every count on every page is
    read from the two files and stands without a model - the coverage, the gaps,
    the audit, and the layout choices. What a host without this cannot do is ask
    what a designer would ADJUST (qc.copilot, qc.components) or propose a layout
    for a gap (qc.layoutsuggest)."""
    from .config import RENDERER
    from .llm import api_configured, configuration_note

    # Before the key check, because this one is true WITH a valid key: the model
    # id names a model the endpoint will not answer on, which fails as a 404
    # inside every pass rather than as a setting anyone can see.
    misconfigured = configuration_note()
    if misconfigured and AI_ENABLED:
        return False, (f"The model settings on this host are wrong, so the "
                       f"slides cannot be looked at. {misconfigured}")
    if not AI_ENABLED:
        return False, ("Model-backed features are switched off on this host, so "
                       "the slides are not looked at. Everything else on this "
                       "page is read from the files and is unaffected.")
    if RENDERER == "none":
        return False, ("Rendering is disabled on this host, and looking at a "
                       "slide means rendering it first. Everything else on this "
                       "page is read from the files and is unaffected.")
    if not api_configured():
        return False, ("No model key is configured on this host, so the slides "
                       "are not looked at. Set GEMINI_API_KEY in the .env file "
                       "at the project root and restart. Everything else on "
                       "this page is read from the files and is unaffected.")
    return True, ""



# How many checked decks are held for a re-check. The loop this exists for is
# short - look at the gaps, build the layout, check again - so a handful is
# plenty, and the deck bytes are the expensive part to keep.
MAX_CHECK_JOBS = 8
_check_jobs: OrderedDict[str, dict] = OrderedDict()


def _run_check(request: Request, data: bytes, filename: str,
               master_bytes: bytes, master_name: str, profile: str,
               can_look: bool, why: str, check_id: str | None = None):
    """The coverage check for a REVISED master, against the deck already
    prepared under this id.

    There is no longer a pre-flight door onto this: coverage is reported on the
    Prepare a deck result as a matter of course, and the only reason to run it
    again is that the designer has built the missing layout and wants to know
    whether the gap closed. Failures therefore go back to Prepare a deck, which
    is the page they came from.
    """
    from pptx import Presentation

    from .applymaster import plan_assignments
    from .layoutgap import misfits as find_misfits
    from .layoutgap import report as layout_coverage
    from .stylespec import dominant_master, extract_layouts, infer_grid
    from .ui_check import render_check_result

    def _fail(msg, code=400):
        return _prep_intake_page(request, msg, code)

    try:
        deck_prs = Presentation(io.BytesIO(data))
        master_prs = Presentation(io.BytesIO(master_bytes))
        target = dominant_master(master_prs)
        if target is None:
            return _fail("That master file has no slide master.")
        # The same layout list and the same planner the format pass uses, so a
        # slide reported here as having no home is the same slide that would
        # fall back there. A second reading could differ, and a pre-flight that
        # disagrees with the run it precedes is worse than no pre-flight.
        target_layouts = extract_layouts(target, embed_assets=False)
        plans = plan_assignments(deck_prs, target_layouts)
    except Exception as exc:
        return _fail(f"Could not read that deck: {type(exc).__name__}: {exc}", 422)

    if not plans:
        return _fail("That deck has no slides to check.")

    # Two kinds of slide need a decision: the ones nothing matched, and the ones
    # that matched by NAME and whose content does not fit the boxes that layout
    # offers. The second kind is invisible to matching by construction, since a
    # name match is a claim about intent (qc.layoutgap.misfits).
    #
    # NEITHER IS GUESSED AT HERE ANY MORE (31/08/2026). This used to render the
    # slides and ask a vision model to place them, which is the pass that became
    # the layout step on Prepare a deck. A re-check is asking one question -
    # "did the layout I just built close the gap?" - and the honest answer to it
    # is arithmetic over the revised master, not a second opinion that can come
    # back differently from the one the rebuild will use.
    look_note = ""
    fallbacks = [p.slide_index for p in plans if p.match_rule == "fallback"]
    try:
        suspect = [m.slide_index
                   for m in find_misfits(deck_prs, target_layouts, plans)]
    except Exception:
        suspect = []
    worth_looking = fallbacks + [i for i in suspect if i not in fallbacks]
    if worth_looking:
        look_note = (f"{len(worth_looking)} slide"
                     f"{'s' if len(worth_looking) != 1 else ''} still need a "
                     f"layout decision against this master. Prepare the deck "
                     f"to choose one for each of them.")

    coverage = layout_coverage(deck_prs, target_layouts, plans)

    # The proposal half. Only the re-check runs it here: this is where a
    # designer decides what to add to a master, and the format result is read
    # after the rebuild, when the answer is too late to act on cheaply. One call
    # per group and no render, since a layout that does not exist yet cannot be
    # photographed.
    wanted = list(coverage.gaps or []) + list(coverage.misfit_clusters or [])
    suggestions, pictures, suggest_note = [], {}, ""
    if wanted and can_look:
        try:
            from .layoutsuggest import suggest, wireframe

            space = (infer_grid(master_prs, target) or {}).get(
                "presentation_space")
            suggestions, asked, unreachable = suggest(
                coverage, deck_prs, target_layouts, space, target)
            pictures = {i: wireframe(s, deck_prs)
                        for i, s in enumerate(suggestions)}
            if unreachable and not suggestions:
                # NOT the same sentence as a rejected proposal: nothing was
                # proposed because nothing was asked.
                suggest_note = (
                    f"The layouts could not be proposed: the model could not "
                    f"be reached ({unreachable}). The findings above are read "
                    f"from the plans and are unaffected.")
            elif asked and not suggestions:
                suggest_note = (
                    "No layout could be proposed for these groups. The findings "
                    "above stand on their own; a proposal that did not answer "
                    "the group it was asked about is discarded rather than "
                    "shown.")
        except Exception as exc:
            suggest_note = (f"The layouts could not be proposed "
                            f"({type(exc).__name__}). The coverage above is "
                            f"unaffected.")
    elif wanted:
        suggest_note = why or (
            "Proposing a layout needs a model, and none is configured here. "
            "The findings above are read from the file and stand without one.")

    # Held for the re-check, not stored: the loop this page exists for is
    # propose, build, check again, and asking a designer to re-upload the deck
    # every time they revise a layout is the friction that stops them checking.
    if check_id is None:
        check_id = uuid.uuid4().hex
    with _jobs_lock:
        _check_jobs[check_id] = {"deck": data, "filename": filename,
                                 "profile": profile}
        _check_jobs.move_to_end(check_id)
        while len(_check_jobs) > MAX_CHECK_JOBS:
            _check_jobs.popitem(last=False)

    profile_meta = next((p for p in _profiles_meta() if p["id"] == profile), None)
    return HTMLResponse(render_check_result(
        deck_name=filename,
        profile_name=(master_name or (profile_meta or {}).get("name", profile)),
        profile_id=profile, coverage=coverage, look_note=look_note,
        suggestions=suggestions, pictures=pictures,
        suggest_note=suggest_note, check_id=check_id,
        master_name=master_name))


def _read_master(master, profile: str):
    """(bytes, name, error) for whichever master this check is against.

    An uploaded file wins over the profile's stored one: it is the file the
    designer just changed, and checking the old copy would answer a question
    nobody asked. It is read for this check and dropped - promoting it to the
    profile is a separate decision, taken in step 1 of Prepare a deck, and
    doing it silently here would change what every future run applies.
    """
    from .templates import load_master

    if master is not None and getattr(master, "filename", ""):
        if not master.filename.lower().endswith(".pptx"):
            return None, "", "The master must be a .pptx file."
        blob = master.file.read(MAX_UPLOAD_BYTES + 1)
        if len(blob) > MAX_UPLOAD_BYTES:
            return None, "", f"The master exceeds the {_MAX_UPLOAD_MB} MB cap."
        bomb = _zip_bomb_reason(blob)
        if bomb:
            return None, "", f"The master was rejected: {bomb}."
        return blob, master.filename, None

    if profile:
        if profile not in _profiles():
            return None, "", "Unknown profile."
        blob = load_master(profile)
        if not blob:
            return None, "", ("That profile carries no master file, so there is "
                              "nothing to check against. Save a master as a "
                              "profile on Prepare a deck first.")
        return blob, "", None
    return None, "", ("Pick a profile to check against, or upload the master "
                      "file directly.")


@app.post("/check/{check_id}/again", response_class=HTMLResponse)
def check_again(request: Request, check_id: str, master: UploadFile = File(...)):
    """Check the same deck against the master the designer has just revised.

    The other half of a proposal. Suggesting a layout and stopping there leaves
    a designer to work out for themselves whether what they built closed the gap,
    and the answer to that is this page run again (design lead, 26/08/2026).

    The deck is the one already checked, held in memory, so a designer revising a
    master three times uploads three masters rather than three masters and three
    decks.
    """
    can_look, why = _can_look()
    with _jobs_lock:
        job = _check_jobs.get(check_id)
    if job is None:
        return _prep_intake_page(
            request, "That deck is no longer held in memory, so the revised "
                     "master has nothing to be checked against. Prepare the "
                     "deck again.", 404)

    master_bytes, master_name, error = _read_master(master, "")
    if error:
        return _prep_intake_page(request, error, 400)

    return _run_check(request, job["deck"], job["filename"], master_bytes,
                      master_name, job.get("profile") or "", can_look,
                      why, check_id=check_id)


def _perform_removals(job, remove_ids) -> int:
    """Take out the pieces named on this job, and record what came back.

    ONE implementation, because there are now two doors onto it: the tick list
    on the review page and a sentence in the ask box. Two copies of this
    bookkeeping is how the removed list and the change list start disagreeing
    about what is in the deck, and the Undo button reads both.

    Returns how many pieces were taken out; the reason for a failure lands on
    job["remove_error"], which is what the page renders.
    """
    from .migrate import apply_removals

    # Anything already taken out is skipped rather than attempted twice: a
    # browser resubmitting the POST would otherwise report a failure for a
    # piece that is correctly gone.
    done = set(job.get("removed") or [])
    wanted = set(remove_ids or []) - done
    ops = [c.remove_op for c in job.get("changes") or []
           if getattr(c, "remove_id", None) in wanted
           and getattr(c, "remove_op", None)]
    if not ops:
        return 0
    try:
        deck, performed = apply_removals(job["deck"], ops)
    except Exception as exc:
        job["remove_error"] = f"{type(exc).__name__}: {exc}"
        return 0
    with _format_lock:
        job["deck"] = deck
        # New bytes, so every cached picture of the old ones is wrong. A prep
        # job is also an audit job, so its design renders and findings go too.
        _drop_review_renders(job)
        if job.get("manifest") is not None:
            _invalidate_renders(job)
        changes = list(job.get("changes") or [])
        # Stamped past the existing ids so an undo of a removal cannot collide
        # with an undo of a migration change.
        for n, change in enumerate(performed, start=len(changes)):
            change.change_id = f"c{n}"
        changes.extend(performed)
        job["changes"] = changes
        job["removed"] = sorted(done | wanted)
    # Verify after write, the same as every other path that changes the deck.
    # A prep job carries an audit of these bytes, and a count taken before a
    # shape left the slide is a description of a file that no longer exists.
    if job.get("manifest") is not None:
        _reaudit_in_place(job)
    return len(performed)


@app.post("/format/{job_id}/remove", response_class=HTMLResponse)
def format_remove(job_id: str, remove_ids: list[str] = Form(None)):
    """Take out the pieces a designer ticked, and nothing else.

    The mirror image of the route below, and it replaced it as the normal path
    (design lead, 26/08/2026): the migration used to remove first and offer the
    pieces back, so a designer's first sight of the deck was already missing
    things. Now it proposes, and this performs.

    Each removal comes back as an ordinary change carrying its own undo, so one
    a designer regrets is taken back by the same Undo button as every other
    change rather than by a second mechanism.
    """
    with _format_lock:
        job = _format_jobs.get(job_id)
    # A job with no prep is not a job this can answer on: the page it returns is
    # drawn from the run, and every run comes through Prepare a deck now.
    if job is None or job.get("deck") is None or job.get("prep") is None:
        return JSONResponse({"error": "unknown or expired job"}, status_code=404)

    _perform_removals(job, remove_ids)

    # Back to the page the tick list is on. Every format job is a prepared
    # deck now that Prepare a deck is the only door onto the master pass, so
    # there is one page to come back to and it is the one that sent us here.
    return _prep_page(job_id, job, error=job.get("remove_error")
                      or job.get("restore_error"))


@app.post("/format/{job_id}/restore", response_class=HTMLResponse)
def format_restore(job_id: str, restore_ids: list[str] = Form(None)):
    """Put selected removed pieces back into the rebuilt deck.

    The migration removes header text the master has no placeholder for, and
    says so loudly. This is the other half of that promise: the designer, not
    the tool, decides whether the removal was right, and the piece goes back
    exactly as it was rather than being retyped from the report."""
    from .migrate import restore_shapes

    with _format_lock:
        job = _format_jobs.get(job_id)
    # A job with no prep is not a job this can answer on: the page it returns is
    # drawn from the run, and every run comes through Prepare a deck now.
    if job is None or job.get("deck") is None or job.get("prep") is None:
        return JSONResponse({"error": "unknown or expired job"}, status_code=404)

    # Anything already back is skipped rather than inserted twice: a browser
    # resubmitting this POST (refresh, back button) would otherwise put a second
    # copy of the same shape on the slide.
    wanted = set(restore_ids or []) - set(job.get("restored") or [])
    items = [{"slide_index": c.slide_index,
              "removed_xml": getattr(c, "removed_xml", None),
              "restore_id": getattr(c, "restore_id", None)}
             for c in job.get("changes") or []
             if getattr(c, "restore_id", None) in wanted]
    if items:
        try:
            deck, outcomes = restore_shapes(job["deck"], items)
        except Exception as exc:
            job["restore_error"] = f"{type(exc).__name__}: {exc}"
        else:
            with _format_lock:
                job["deck"] = deck
                # New bytes, so every picture of the old ones is wrong - the same
                # reason an undo drops them. This was missed: restoring a piece
                # here and then opening the review tab showed the deck without
                # it.
                _drop_review_renders(job)
                # Kept per piece, not as a bare list of ids: a piece nudged clear
                # of the master's header is back in a different place than it
                # left, and a designer has to be told which.
                notes = dict(job.get("restored_notes") or {})
                for o in outcomes:
                    if o.get("restore_id"):
                        notes[o["restore_id"]] = o["detail"]
                job["restored_notes"] = notes
                job["restored"] = sorted(notes)

    # Back to the page the tick list is on. Every format job is a prepared
    # deck now that Prepare a deck is the only door onto the master pass, so
    # there is one page to come back to and it is the one that sent us here.
    return _prep_page(job_id, job, error=job.get("remove_error")
                      or job.get("restore_error"))


# How many slides the deck view renders on ONE PAGE. Every render is a
# PowerPoint export, so an uncapped page on a 200-slide deck is a several-minute
# wait. This is a page size, not a limit on what can be reviewed: the page pages
# through the rest, because a review that silently ends at slide 20 of a 26-slide
# deck reads as six slides the tool declined to show (design lead, 23/08/2026).
REVIEW_PAGE_SIZE = 20


def _reviewable(job) -> list[int]:
    """Every slide worth reviewing: the ones this run changed, plus any that
    failed. A slide the content pass did not touch has nothing to review and
    would only push the ones that do further down the page."""
    changed = {c.slide_index for c in (job.get("changes") or [])}
    changed |= set((job.get("errors") or {}).keys())
    return sorted(changed)


def _review_slides(job, page: int = 0) -> list[int]:
    """One page of reviewable slides."""
    every = _reviewable(job)
    start = max(0, page) * REVIEW_PAGE_SIZE
    return every[start:start + REVIEW_PAGE_SIZE]


def _layout_use(plans, key: str) -> dict:
    """{layout name: slides on it} from the assignment plans, for the badges on
    the master view. A layout nothing sits on is worth seeing as unused."""
    counts: dict[str, int] = {}
    for p in plans or []:
        name = getattr(p, key, None)
        if name:
            counts[name] = counts.get(name, 0) + 1
    return counts


def _url_keys(images: dict, prefix: str) -> dict:
    """Re-key renders as the page asks for them: "before:7" -> "slide_before_7".

    One namespace, URL-safe, and prefixed per view. Without the prefix a
    12-layout master and a 12-slide deck collide on "before:3" and the master
    view starts serving slide pictures - which is a wrong answer that looks like
    a right one, the worst kind for a review page."""
    return {f"{prefix}_{k.replace(':', '_', 1)}": png
            for k, png in (images or {}).items()}


def _ensure_review(job, view: str, page: int = 0) -> str | None:
    """Render what this view needs, once per page, and cache it on the job.
    Returns a readable reason when nothing could be rendered - the page then
    shows the change list and the Undo buttons without pictures, which is the
    part a designer cannot do without."""
    from .render import layout_previews, slide_previews

    key = "review_master" if view == "master" else f"review_deck_{page}"
    if job.get(key) is not None:
        return job.get(key + "_error")
    if job.get("source") is None:
        job[key] = {}
        job[key + "_error"] = ("the uploaded deck is no longer in memory, so "
                              "there is nothing to compare against")
        return job[key + "_error"]
    try:
        if view == "master":
            out = layout_previews(job["source"], job["deck"])
            out["images"] = _url_keys(out.get("images"), "layout")
        else:
            out = slide_previews(job["source"], job["deck"],
                                 _review_slides(job, page))
            out["images"] = _url_keys(out.get("images"), "slide")
    except Exception as exc:  # reading the DECK failed: there is no answer at all
        job[key] = {}
        job[key + "_error"] = f"{type(exc).__name__}: {exc}."
        return job[key + "_error"]
    # A render failure is kept WITH the result, not instead of it. Which layouts
    # the deck arrived with and which it has now needs no PowerPoint; discarding
    # the whole answer made the page say "No layouts to show" about a master
    # carrying twelve (design lead, 23/08/2026).
    job[key] = out
    job[key + "_error"] = out.get("error")
    return out.get("error")


def _drop_review_renders(job) -> None:
    """Forget every cached render. Called after an undo: the deck has changed,
    and a stale "after" picture beside a row marked undone is the one thing this
    page must never show.

    The render TAG goes with them, and it is the other half of that sentence:
    clearing the server's copy only makes the page ask again, and the browser
    was answering out of its own cache under an unchanged URL
    (_render_tag)."""
    for key in [k for k in list(job) if k.startswith(("review_master",
                                                      "review_deck"))]:
        job[key] = None
    job["render_tag"] = None


def _review_page(job_id: str, job, view: str, undo_error: str | None = None,
                 page: int = 0):
    from .ui_review import render_review

    error = _ensure_review(job, view, page)
    key = "review_master" if view == "master" else f"review_deck_{page}"
    previews = job.get(key) or {}
    every = _reviewable(job)
    profile = job.get("profile")
    meta = next((p for p in _profiles_meta() if p["id"] == profile), None)
    return HTMLResponse(render_review(
        job_tabs=_tabs_for(job_id, job, "review"),
        deck_name=job["filename"],
        profile_name=(meta or {}).get("name", profile or ""),
        job_id=job_id, view=view, previews=previews,
        changes=job.get("changes") or [], plans=job.get("plans") or [],
        errors=job.get("errors") or {},
        undone=set(job.get("undone") or []),
        notes=job.get("undo_notes") or {},
        shown=_review_slides(job, page),
        reviewable=len(every),
        page=page, page_size=REVIEW_PAGE_SIZE,
        total_slides=len(job.get("plans") or []),
        masters=job.get("masters") or 1,
        used_before=_layout_use(job.get("plans"), "source_layout"),
        used_after=_layout_use(job.get("plans"), "target_layout"),
        truncated=bool(previews.get("truncated")),
        img_tag=_render_tag(job),
        render_error=error, undo_error=undo_error))


@app.get("/format/{job_id}/review", response_class=HTMLResponse)
def format_review(job_id: str, view: str = "master", page: int = 0):
    """Before and after for one format run, as the template and as the deck."""
    with _format_lock:
        job = _format_jobs.get(job_id)
    if job is None or job.get("deck") is None:
        return HTMLResponse(
            "<p>Unknown or expired job. Format the deck again.</p>",
            status_code=404)
    return _review_page(job_id, job, "deck" if view == "deck" else "master",
                        page=max(0, page))


@app.post("/format/{job_id}/undo", response_class=HTMLResponse)
def format_undo(job_id: str, change_ids: list[str] = Form(None),
                page: int = Form(0)):
    """Take one reported change back, exactly as it was.

    Anything already undone is skipped rather than replayed: a browser
    resubmitting this POST would otherwise insert a second copy of a restored
    shape, or reset a shape a later undo has since moved."""
    from .undo import apply_undo, expand

    with _format_lock:
        job = _format_jobs.get(job_id)
    if job is None or job.get("deck") is None:
        return HTMLResponse(
            "<p>Unknown or expired job. Format the deck again.</p>",
            status_code=404)

    # Expanded before it is filtered: undoing a change means undoing what came
    # after it on that slide too, or the slide lands in a state the run never
    # produced (qc.undo.followers).
    changes = job.get("changes") or []
    wanted = set(change_ids or []) - set(job.get("undone") or [])
    already = set(job.get("undone") or [])
    items = [{"change_id": c.change_id, "slide_index": c.slide_index,
              "action": c.action, "ops": c.undo}
             for c in expand(changes, wanted)
             if c.change_id not in already]
    error = None
    if items:
        try:
            deck, outcomes = apply_undo(job["deck"], items)
        except Exception as exc:
            error = f"{type(exc).__name__}: {exc}"
        else:
            with _format_lock:
                job["deck"] = deck
                notes = dict(job.get("undo_notes") or {})
                for o in outcomes:
                    if o.get("change_id"):
                        notes[o["change_id"]] = o["detail"]
                job["undo_notes"] = notes
                job["undone"] = sorted(
                    set(job.get("undone") or [])
                    | {o["change_id"] for o in outcomes
                       if o.get("done") and o.get("change_id")})
                _drop_review_renders(job)
    return _review_page(job_id, job, "deck", undo_error=error,
                        page=max(0, page))


@app.get("/review-img/{job_id}/{key}.png")
def review_img(job_id: str, key: str):
    from fastapi.responses import Response

    with _format_lock:
        job = _format_jobs.get(job_id)
    if job is None:
        return JSONResponse({"error": "unknown or expired job"}, status_code=404)
    # Every review bucket is searched by PREFIX, not by an exact name: the deck
    # view caches per page ("review_deck_0", "review_deck_1"), and naming the two
    # buckets literally here meant page 0's images were cached under a key this
    # route never looked in, so the deck view rendered 18 pictures and served
    # none of them. The image keys themselves are already unique across views
    # (_url_keys), so a prefix scan cannot return the wrong picture.
    png = None
    for bucket, value in job.items():
        if not str(bucket).startswith("review_") or not isinstance(value, dict):
            continue
        png = (value.get("images") or {}).get(key)
        if png is not None:
            break
    if png is None:
        return JSONResponse({"error": "no such render"}, status_code=404)
    return Response(content=png, media_type="image/png",
                    headers={"Cache-Control": "private, max-age=600"})


@app.get("/format/{job_id}/download")
def format_download(job_id: str):
    from fastapi.responses import Response

    with _format_lock:
        job = _format_jobs.get(job_id)
    if job is None or job["deck"] is None:
        return JSONResponse({"error": "unknown or expired job"}, status_code=404)
    stem = Path(job["filename"]).stem or "deck"
    return Response(
        job["deck"],
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers=_attachment(f"{stem} (master applied).pptx"))


# ------------------------------------------- prepare: the master, then the audit
#
# One upload, one page, in the order the work actually happens. Applying a
# master and auditing a deck were two routes and two pages, and a designer
# needed both of them in that order every time - so the pipeline is qc.prep and
# this is bookkeeping around it.


def _register_prep(job_id: str, prep, profile: str, profile_obj) -> dict:
    """One job, in both registries, under one id.

    A prepared deck IS a format job and an audit job: it has plans and a
    before/after review, and it has a manifest and design cards. Putting ONE
    dict in both registries is what lets every page that already exists - the
    review, the design cards, the checklist, the download, the ask box - work
    on it without a second copy of any of them.

    One dict rather than two, and that is the part that matters. Two would each
    hold their own `deck`, and the first fix applied from the design page would
    leave the download serving the file as it stood before it. One page has to
    mean one file (design lead, 27/08/2026).

    A run whose audit failed is registered as a format job only. The design
    page reads the manifest for the slide count on its first line, so a job
    with no manifest is not an audit job however much it looks like one, and
    half-registering it would trade a missing link for a stack trace.
    """
    job = {
        # --- the format half
        "deck": prep.deck, "source": prep.source, "filename": prep.filename,
        "profile": profile, "plans": prep.plans, "errors": prep.errors,
        "applied": prep.applied, "masters": prep.masters,
        "stragglers": prep.stragglers, "space_notes": prep.space_notes,
        "coverage": prep.coverage, "changes": prep.changes,
        "restored": [], "removed": [], "undone": [], "undo_notes": {},
        # --- the audit half
        "manifest": prep.manifest, "cleaned": None,
        "profile_obj": profile_obj, "master_spec": None,
        # --- what makes it one job rather than two
        "prep": prep,
    }
    with _format_lock:
        _format_jobs[job_id] = job
        while len(_format_jobs) > MAX_FORMAT_JOBS:
            _format_jobs.popitem(last=False)

    # The UPLOAD, held under the same id, so "add the missing layout and check
    # again" works from this page too. The gaps panel is the same one the
    # prepared deck draws, and its re-check form needs a deck to re-read; the
    # upload is the right one, because the coverage was computed against it
    # before anything was rebuilt.
    with _jobs_lock:
        _check_jobs[job_id] = {"deck": prep.source, "filename": prep.filename,
                               "profile": profile}
        _check_jobs.move_to_end(job_id)
        while len(_check_jobs) > MAX_CHECK_JOBS:
            _check_jobs.popitem(last=False)

    if prep.manifest is None:
        return job
    with _jobs_lock:
        _jobs[job_id] = job
        _expire_old_decks()
    return job


def _prep_layout_review(job_id: str, job) -> None:
    """The visual layout pass, on the REBUILT deck, as part of the prepare run.

    Part of the run rather than a button, because a button on a page two clicks
    down is a feature that does not run. The geometric alignment pass answers
    "were these meant to line up?" from proximity alone and therefore has to be
    cautious: it will not compare shapes more than 0.15in apart, it drops any
    cluster smaller than three, and it diverts repeated structures and touching
    clusters out of its pools entirely. Those are the right calls WITHOUT a
    model. With one, a designer's eye supplies the intent that proximity was
    standing in for, and the code still supplies every number (qc.copilot).

    Runs AFTER the audit and on the rebuilt deck, for the same reason the audit
    does: half of what is wrong on the raw file is about to be rewritten by the
    master, and judging the slides before that is judging a file nobody will
    send.

    Degrades to a sentence, never raises. Losing this pass costs the alignment
    judgments; it must not cost the rebuilt deck or the audit that came with it.
    The thumbnails it renders are cached on the job, so the design QC page that
    a designer opens next does not render them a second time.
    """
    job["layout_note"] = ""
    job["layout_ok"] = True
    if job.get("manifest") is None:
        return  # the audit did not run; the page already says so, once

    can_look, why = _can_look()
    if not can_look:
        job["layout_note"] = why
        job["layout_ok"] = False
        return

    from .copilot import run_copilot

    try:
        _ensure_thumbs(job_id, job)
        new_records, reviewed = run_copilot(job["deck"], job["thumbs"],
                                            job["manifest"])
    except Exception as exc:
        job["layout_note"] = (f"The slides could not be looked at for "
                              f"alignment: {type(exc).__name__}: {exc}. "
                              f"Everything else on this page still stands.")
        job["layout_ok"] = False
        return

    _merge_records(job, new_records)
    n = len(new_records)
    job["layout_note"] = (
        f"The visual model looked at {reviewed} slide"
        f"{'s' if reviewed != 1 else ''} and found {n} alignment "
        f"{'thing' if n == 1 else 'things'} a designer would adjust"
        + (", listed with the rest below and tickable like any other finding."
           if n else
           ". Nothing on the slides it saw reads as out of line."))


def _prep_page(job_id: str, job, banner: str = "", error: str | None = None):
    """The prepared deck's one page, drawn from the job rather than from the
    run that produced it. Every count on it therefore reflects the deck as it
    now stands, including after a fix applied from the ask box."""
    from .prep import headline
    from .ui_prep import render_prep_result

    prep = job["prep"]
    prep.deck = job.get("deck")
    prep.manifest = job.get("manifest")

    design_open = design_error = None
    auto: dict = {}
    per_slide: dict = {}
    if job.get("manifest") is not None:
        design_open = design_count(job)
        design_error = job.get("design_error")
        auto = _auto_plan(job, 0)
        answered = {a.finding_id for a in job.get("design_applied") or []}
        every = [f for f in (job.get("design") or [])
                 if f.finding_id not in answered]
        records = [r for r in (job["manifest"].get("records") or [])
                   if r.get("module") != "preflight"]
        per_slide = _design_severity_map(every, records, answered)

    meta = next((p for p in _profiles_meta() if p["id"] == job.get("profile")),
                None)
    chat_ready, chat_note = _chat_available()
    return HTMLResponse(render_prep_result(
        prep=prep, job_id=job_id, tabs=_tabs_for(job_id, job, "overview"),
        profile_name=(meta or {}).get("name") or job.get("profile") or "",
        headline=headline(prep), auto=auto, design_open=design_open,
        design_error=design_error, per_slide=per_slide,
        banner=banner, error=error, chat=chat_ready, chat_note=chat_note,
        # What the remove and restore buttons on this page have already done.
        # These used to be hardcoded empty, so a piece put back went back into
        # the deck and the page still offered to put it back.
        restored=job.get("restored") or [],
        restored_notes=job.get("restored_notes") or {},
        restore_error=job.get("restore_error"),
        removed=job.get("removed") or [],
        remove_error=job.get("remove_error"),
        layout_note=job.get("layout_note") or "",
        layout_ok=job.get("layout_ok", True)))


@app.get("/prep", response_class=HTMLResponse)
def prep_intake(request: Request, saved: str = "", replaced: str = ""):
    return _prep_intake_page(request, saved=saved, replaced=bool(replaced))


@app.get("/prep/{job_id}", response_class=HTMLResponse)
def prep_view(job_id: str):
    """The prepared deck's page, for a run that already happened.

    Its own GET because the page is now the doorway to four others, and a page
    a designer cannot get back to is a page they will not leave (the audit
    report learned the same lesson)."""
    with _format_lock:
        job = _format_jobs.get(job_id)
    if job is None or job.get("prep") is None:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES,
                                         "Unknown or expired job."),
                            status_code=404)
    return _prep_page(job_id, job)


@app.post("/prep/{job_id}/propose", response_class=HTMLResponse)
def prep_propose(job_id: str):
    """Draw a layout for each gap the master could not fill.

    A BUTTON RATHER THAN A STEP (31/08/2026). This ran inside every prepare run
    until the model came out of the master application; it is the last pass in
    the flow that asks a model anything, and leaving it in the run meant the
    deterministic half still finished by waiting on a network call. Most runs
    have no gaps worth acting on, and the ones that do are a conversation with
    whoever owns the master rather than something to fix this afternoon.
    """
    from .prep import propose_layouts
    from .templates import load_master

    with _format_lock:
        job = _format_jobs.get(job_id)
    if job is None or job.get("prep") is None:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES,
                                         "Unknown or expired job."),
                            status_code=404)

    off = _ai_disabled_response()
    if off is not None:
        return off
    can_look, why = _can_look()
    if not can_look:
        return _prep_page(job_id, job, error=why)

    master = load_master(job.get("profile") or "")
    if not master:
        return _prep_page(job_id, job, error=(
            "That profile no longer carries a master file, so a proposal "
            "cannot be drawn on its frame."))
    try:
        propose_layouts(job["prep"], master)
    except Exception as exc:
        return _prep_page(job_id, job, error=(
            f"The layouts could not be proposed ({type(exc).__name__}: {exc}). "
            f"The gaps below are read from the plans and are unaffected."))
    return _prep_page(job_id, job)


@app.post("/prep", response_class=HTMLResponse)
def prep_deck(request: Request, deck: UploadFile = File(...),
              profile: str = Form(...), look: str = Form(None)):
    """Step 1 -> step 2. Reads the deck against the master and STOPS.

    Nothing is rebuilt here. What used to be one press is now two, with the
    layout decisions in between, so the expensive irreversible half only runs
    against layouts a designer has seen (qc.prep.plan)."""
    from .prep import PrepError
    from .prep import plan as plan_prep
    from .templates import load_master

    def _fail(msg, code=400):
        return _prep_intake_page(request, msg, code)

    filename = deck.filename or "deck.pptx"
    if not filename.lower().endswith(".pptx"):
        return _fail("Only .pptx files are accepted.")
    if profile not in _profiles():
        return _fail("Unknown profile.")
    master_bytes = load_master(profile)
    if not master_bytes:
        return _fail("That profile carries no master file, so there is nothing "
                     "to apply. Save a master as a profile in step 1 first.")

    data = deck.file.read(MAX_UPLOAD_BYTES + 1)
    if len(data) > MAX_UPLOAD_BYTES:
        return _fail(f"File exceeds the {_MAX_UPLOAD_MB} MB cap.", 413)
    bomb = _zip_bomb_reason(data)
    if bomb:
        return _fail(f"File rejected: {bomb}.", 413)

    try:
        profile_obj = _resolve_profile(profile, data)
    except Exception as exc:
        return _fail(f"That profile could not be loaded: "
                     f"{type(exc).__name__}: {exc}", 422)

    try:
        prepared = plan_prep(data, filename, master_bytes)
    except PrepError as exc:
        return _fail(str(exc), exc.status)
    except Exception as exc:
        return _fail(f"That deck could not be read against the master: "
                     f"{type(exc).__name__}: {exc}", 422)

    plan_id = _remember_plan_stage(prepared, profile, profile_obj, master_bytes)
    # The checkbox is remembered here rather than acted on: it decides whether
    # the REBUILT deck gets looked at for alignment, which happens after the
    # rebuild, two presses from now.
    with _plans_lock:
        _plans[plan_id]["look"] = bool(look)
    return _layouts_page(plan_id)


# ---------------------------------------------- Step 2: choosing the layouts
# The one page between an upload and a rewritten file. Only the slides the file
# could not place with confidence appear on it; everything else is a count.


def _layout_shots(plan_id: str, held: dict) -> None:
    """Render the uncertain slides and the master's layouts, once per plan.

    Two renders, both capped by what the page shows: the slides in question
    (never the whole deck) and the master's layout catalogue (which is one
    small file whatever the deck's length). Cached on the held plan because a
    designer who changes one radio and reloads must not pay for them twice.

    DEGRADES TO WORDS. A host with no renderer still gets the page - layout
    names, what each one holds, and what the slide is asking for are all read
    from the files - and it says so rather than showing broken images.
    """
    if held.get("shots") or held.get("render_note"):
        return
    plan = held["plan"]
    wanted = [c.slide_index for c in plan.choices]
    if not wanted:
        return
    from .render import export_decks_png, layout_catalogue

    try:
        pngs = export_decks_png({"s": plan.source}, wanted)
        held["shots"] = {int(k.split(":", 1)[1]): v for k, v in pngs.items()}
    except Exception as exc:
        held["render_note"] = (
            f"The slides could not be rendered on this machine "
            f"({type(exc).__name__}), so the pictures are missing. What each "
            f"slide holds and what each layout offers are read from the files "
            f"and are unaffected.")
        return
    try:
        catalogue, entries, _ = layout_catalogue(held["master"])
        shots = export_decks_png({"l": catalogue},
                                 [e["index"] for e in entries])
        by_index = {int(k.split(":", 1)[1]): v for k, v in shots.items()}
        held["layout_thumbs"] = {e["layout"]: by_index[e["index"]]
                                 for e in entries if e["index"] in by_index}
    except Exception:
        # The slide pictures are the ones that matter; a missing layout preview
        # leaves a named option with a description, which is still a choice a
        # designer can make.
        held["layout_thumbs"] = {}


def _layouts_page(plan_id: str, message: str = "", status: int = 200):
    held = _get_plan(plan_id)
    if held is None:
        return HTMLResponse(
            render_index(_pickable_profiles(), MODULES,
                         "That layout review has expired. Upload the deck "
                         "again to start over."), status_code=404)
    from .ui_layouts import render_layouts, render_nothing_to_choose

    plan = held["plan"]
    meta = next((p for p in _profiles_meta() if p["id"] == held["profile"]),
                None)
    profile_name = (meta or {}).get("name") or held["profile"]

    if not plan.choices:
        return HTMLResponse(render_nothing_to_choose(
            deck_name=plan.filename, profile_name=profile_name,
            plan_id=plan_id, slides=plan.slides), status_code=status)

    _layout_shots(plan_id, held)
    return HTMLResponse(render_layouts(
        deck_name=plan.filename, profile_name=profile_name, plan_id=plan_id,
        choices=plan.choices,
        layout_names=[l["name"] for l in plan.layouts if l.get("name")],
        matched=plan.slides - plan.undecided,
        slide_shots={i: f"/plan-img/{plan_id}/slide-{i}.png"
                     for i in held.get("shots", {})},
        layout_thumbs={n: f"/plan-img/{plan_id}/layout-{i}.png"
                       for i, n in enumerate(held.get("layout_thumbs", {}))},
        render_note=held.get("render_note", ""), message=message),
        status_code=status)


@app.get("/prep/{plan_id}/layouts", response_class=HTMLResponse)
def prep_layouts(plan_id: str):
    return _layouts_page(plan_id)


@app.get("/plan-img/{plan_id}/{key}.png")
def plan_img(plan_id: str, key: str):
    """A rendered slide or layout from a plan awaiting its decision.

    Served rather than inlined as data URIs: a twenty-slide review with
    twenty-four layout previews is megabytes of base64 in one document, and the
    browser cannot cache any of it across the reload a designer does after
    changing a pick."""
    held = _get_plan(plan_id)
    if held is None:
        return Response(status_code=404)
    png = None
    if key.startswith("slide-"):
        try:
            png = held.get("shots", {}).get(int(key[6:]))
        except ValueError:
            png = None
    elif key.startswith("layout-"):
        try:
            png = list(held.get("layout_thumbs", {}).values())[int(key[7:])]
        except (ValueError, IndexError):
            png = None
    if png is None:
        return Response(status_code=404)
    return Response(png, media_type="image/png",
                    headers={"Cache-Control": "private, max-age=600"})


@app.post("/prep/{plan_id}/layouts", response_class=HTMLResponse)
async def prep_apply_layouts(request: Request, plan_id: str):
    """Step 2 -> step 3. The picks go onto the plans, then the master is
    applied and the rebuilt deck is audited.

    The form is read raw rather than through typed parameters because the field
    names carry the slide index (pick_7, other_7) and FastAPI cannot declare a
    parameter per slide of a deck it has not seen."""
    from .prep import PrepError
    from .prep import run as run_prep

    held = _get_plan(plan_id)
    if held is None:
        return HTMLResponse(
            render_index(_pickable_profiles(), MODULES,
                         "That layout review has expired. Upload the deck "
                         "again to start over."), status_code=404)

    from .layoutpick import LEAVE, apply_picks
    from .layoutpick import note as layout_note

    plan = held["plan"]
    form = await request.form()
    picks = {}
    for choice in plan.choices:
        idx = choice.slide_index
        # The select wins over the radio when it names something: it is only
        # reachable by opening it and choosing, where the radio carries a
        # default nobody had to touch.
        other = (form.get(f"other_{idx}") or "").strip()
        wanted = other or (form.get(f"pick_{idx}") or "").strip()
        if wanted:
            picks[idx] = wanted

    moved = apply_picks(plan.plans, picks, plan.layouts)
    refused = sum(1 for v in picks.values() if v == LEAVE)
    match_note = layout_note(len(plan.choices), len(picks) - refused, moved,
                             refused)

    try:
        prep = run_prep(plan, held["master"], held["profile_obj"],
                        match_note=match_note)
    except PrepError as exc:
        return _layouts_page(plan_id, str(exc), exc.status)
    except Exception as exc:
        return _layouts_page(plan_id, f"That deck could not be prepared: "
                                      f"{type(exc).__name__}: {exc}", 422)

    job_id = uuid.uuid4().hex
    job = _register_prep(job_id, prep, held["profile"], held["profile_obj"])

    # The checkbox from step 1. A designer who asked for the slides to be looked
    # at meant the REBUILT slides: judging the upload is judging a file nobody
    # will send. Runs before the audit is recorded, because it adds records to
    # the manifest and the archived copy should be the one the designer saw.
    if held.get("look"):
        _prep_layout_review(job_id, job)

    if prep.manifest is not None:
        from .auth import current_user
        from .store import record_audit

        user = current_user(request)
        record_audit(prep.manifest, user["name"] if user else "anonymous")
    _drop_plan(plan_id)
    return _prep_page(job_id, job)


@app.post("/audit", response_class=HTMLResponse)
def audit(request: Request, deck: UploadFile = File(...),
          profile: str = Form("prezlab_en"),
          master: UploadFile = File(None),
          modules: list[str] = Form(None)):
    filename = deck.filename or "upload.pptx"
    if not filename.lower().endswith(".pptx"):
        return HTMLResponse(render_index(_pickable_profiles(), MODULES, "Only .pptx files are accepted."), status_code=400)
    if not _ephemeral(profile) and profile not in _profiles():
        return HTMLResponse(render_index(_pickable_profiles(), MODULES, "Unknown profile."), status_code=400)
    selected = [m for m in (modules or []) if m in MODULES] or None

    master_profile, master_spec = None, None
    if profile == MASTER_PROFILE:
        master_profile, master_spec, err = _profile_from_master(master)
        if err:
            return HTMLResponse(
                render_index(_pickable_profiles(), MODULES, err), status_code=400)

    data = deck.file.read(MAX_UPLOAD_BYTES + 1)
    if len(data) > MAX_UPLOAD_BYTES:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES, f"File exceeds the {_MAX_UPLOAD_MB} MB cap."), status_code=413)
    bomb = _zip_bomb_reason(data)
    if bomb:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES, f"File rejected: {bomb}."), status_code=413)

    fd, tmp_name = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)  # Windows cannot delete a file while this handle is open
    tmp = Path(tmp_name)
    try:
        tmp.write_bytes(data)
        try:
            profile_obj = master_profile or _resolve_profile(profile, data)
            result = run_audit(tmp, profile_obj, selected)
        except Exception as exc:
            return HTMLResponse(
                render_index(_pickable_profiles(), MODULES, f"Could not audit this file: {type(exc).__name__}: {exc}"),
                status_code=422)
    finally:
        tmp.unlink(missing_ok=True)  # PRD: uploads auto-delete after processing

    manifest = result.to_manifest()
    manifest["deck"] = filename  # temp path is meaningless to the user
    job_id = uuid.uuid4().hex
    with _jobs_lock:
        _jobs[job_id] = {"manifest": manifest, "deck": data,
                         "cleaned": None, "filename": filename,
                         "profile": profile, "profile_obj": profile_obj,
                         # kept so the report can offer the spec it used; the
                         # master's own bytes are already gone
                         "master_spec": master_spec}
        _expire_old_decks()
    from .auth import current_user
    from .promotion import promoted_issue_types
    from .store import comment_counts, record_audit

    user = current_user(request)
    record_audit(manifest, user["name"] if user else "anonymous")
    # Straight to the slides. The audit report is a list of 2,000 occurrences
    # with no picture; a designer's first question is "what is wrong with THIS
    # slide", and answering it used to cost them a page they had to leave
    # (design lead, 24/08/2026). The report is still there - Design QC links
    # back to it, and it is where the exports, the triage buttons and the
    # comments live - it is just no longer the doorway.
    return RedirectResponse(f"/design/{job_id}", status_code=303)


def _tabs_for(job_id: str, job: dict | None, active: str) -> str:
    """The tab strip for whichever views this job actually has.

    Availability is computed from the job rather than assumed, because the six
    views are not all reachable for every run: an upload that was only audited
    was never rebuilt, so there is no before/after; a prepare run whose audit
    failed has a deck and no findings. A tab that 404s is worse than an absent
    one - it reads as the tool losing the page.
    """
    from .ui import job_tabs

    if job is None:
        return ""
    available, counts = set(), {}
    manifest = job.get("manifest")
    if job.get("prep") is not None:
        available.add("overview")
    if job.get("plans"):
        available.add("review")
    if job.get("deck") is not None:
        available.add("checklist")
    if manifest is not None:
        available |= {"findings", "design"}
        counts["findings"] = sum(
            1 for r in (manifest.get("records") or [])
            if r.get("module") != "preflight")
        open_design = design_count(job)
        if open_design:
            counts["design"] = open_design
    if job.get("applied_records") and job.get("prev_deck") is not None:
        available.add("changes")
    # A strip with one tab on it is a label, not navigation.
    return job_tabs(job_id, active, available, counts) \
        if len(available) > 1 else ""


def _job(job_id: str) -> dict | None:
    return _jobs.get(job_id)


def _apply_audit_fixes(job, selected: set):
    """Apply the selected audit fixes to the job's deck and verify by
    re-auditing. Returns (FixResult, the total before, a staleness note or None).

    ONE path, because there are now two doors onto it: the report's Apply button
    and the tick list on the design page. Two copies of this bookkeeping is how
    the two pages start disagreeing about what the deck currently is - one of
    them forgetting to drop a cached render, say, and showing a picture of the
    deck as it stood before the fix.
    """
    from .fixer import apply_fixes

    before_total = job["manifest"]["summary"]["total"]
    fx = apply_fixes(job["deck"], job["manifest"]["records"], selected)
    changed_ids = {o.record_id for o in fx.outcomes if o.outcome == "changed"}
    applied = [r for r in job["manifest"]["records"]
               if r["record_id"] in changed_ids]
    with _jobs_lock:
        job["prev_deck"] = job["deck"]   # pre-fix bytes for the visual diff
        job["applied_records"] = job.get("applied_records", []) + applied
        job["deck"] = fx.cleaned_bytes   # further fixes build on the cleaned deck
        job["cleaned"] = fx.cleaned_bytes
        # every render and every finding derived from the old bytes, dropped:
        # a fix that moved a shape may have created or cleared an overlap
        _invalidate_renders(job)
    # Verify-after-write: re-audit the cleaned bytes so both pages reflect the
    # actual new state of the deck, not an assumption.
    return fx, before_total, _reaudit_in_place(job)


@app.post("/apply", response_class=HTMLResponse)
def apply(request: Request, job_id: str = Form(...),
          record_ids: list[str] = Form(None)):
    job = _job(job_id)
    if job is None:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES,
                                         "Unknown or expired job."), status_code=404)
    if job["deck"] is None:
        return HTMLResponse(render_report(
            job["manifest"], job_id, can_fix=False,
            tabs=_tabs_for(job_id, job, "findings"),
            banner="The deck is no longer held in memory (newer audits replaced it). "
                   "Re-upload to apply fixes."), status_code=410)
    selected = set(record_ids or [])
    if not selected:
        return render_report(job["manifest"], job_id, can_fix=True,
                             banner="No fixes selected.",
                             assist=AI_ENABLED and not _ephemeral(job["profile"]))

    fx, before_total, stale = _apply_audit_fixes(job, selected)
    skipped = [o for o in fx.outcomes if o.outcome == "skipped"]
    manifest = job["manifest"]
    after_total = manifest["summary"]["total"]
    note = (f"Applied {fx.applied} fix{'es' if fx.applied != 1 else ''}. "
            f"Re-audit of the cleaned deck: {after_total} findings remain "
            f"(was {before_total}).")
    if skipped:
        note += f" Skipped {len(skipped)}."
    if stale:
        note += f" {stale}"
    from .auth import current_user
    from .promotion import promoted_issue_types
    from .store import comment_counts, record_audit

    user = current_user(request)
    record_audit(manifest, user["name"] if user else "anonymous", kind="fix")
    return render_report(manifest, job_id, can_fix=True, banner=note,
                         tabs=_tabs_for(job_id, job, "findings"),
                         has_cleaned=True, diff_href=f"/diff/{job_id}",
                         promoted=promoted_issue_types(),
                         comments=comment_counts(job["filename"]),
                         design=design_count(job),
                         assist=AI_ENABLED and not _ephemeral(job["profile"]))


# ------------------------------------------------------------- design QC
#
# The audit's own modules answer "does this deck match the profile". These are
# the questions that only have answers once a master is on the deck - a palette
# meeting the deck's own colours, text meeting a new background, the master's
# furniture meeting content that was already there - and each of them has more
# than one right fix, so the designer picks and the pick is reversible. Detection
# is qc.design, applying a pick is qc.remedy, taking it back is qc.undo (the same
# machinery the format review page uses, reaching one step further).


def _palette_cfg(job) -> dict:
    profile = job.get("profile_obj")
    try:
        return profile.module_config("color_palette") if profile else {}
    except Exception:
        return {}


def _ensure_design(job) -> str | None:
    """Run the design pass once per state of the deck, and cache it on the job.

    Cleared (design=None) after every apply and every undo, because the findings
    ARE a description of the current bytes: a card offering to fix something the
    designer just fixed is the one thing this page must not show."""
    if job.get("design") is not None:
        return job.get("design_error")
    if job.get("deck") is None:
        job["design"] = []
        return None
    from .design import scan

    try:
        job["design"] = scan(job["deck"], _palette_cfg(job))
        job["design_error"] = None
    except Exception as exc:
        job["design"] = []
        job["design_error"] = (f"The design checks could not run on this deck: "
                               f"{type(exc).__name__}: {exc}")
    return job.get("design_error")


def design_count(job) -> int | None:
    """How many open design decisions this deck has, for the badge on the audit
    report. None when the pass could not run - the link still appears, because a
    missing link reads as a missing feature."""
    if _ensure_design(job):
        return None
    answered = {a.finding_id for a in job.get("design_applied") or []}
    return sum(1 for f in job.get("design") or []
               if f.finding_id not in answered)


def _reaudit_in_place(job) -> str | None:
    """Re-run the audit over the deck as the design pass has left it, keeping
    the SAME job.

    In place, and that is the whole point: /reaudit mints a new job id, which
    would strand every design decision recorded against this one. The rule the
    audit flow already follows - verify after write, never assume - applies
    just as much to a change made from this page."""
    if job.get("deck") is None:
        return None
    fd, tmp_name = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    tmp = Path(tmp_name)
    try:
        tmp.write_bytes(job["deck"])
        profile = job.get("profile_obj") or Profile.load(job["profile"])
        result = run_audit(tmp, profile)
    except Exception as exc:
        return (f"The deck was changed, but it could not be re-audited "
                f"({type(exc).__name__}: {exc}), so the audit report's counts "
                f"are from before these changes.")
    finally:
        tmp.unlink(missing_ok=True)
    manifest = result.to_manifest()
    manifest["deck"] = job["filename"]
    job["manifest"] = manifest
    return None


# How many slides are rendered per PowerPoint call. Flipping one slide at a time
# would mean one COM round trip per click (a second or more each); rendering the
# whole deck up front would mean minutes of wait on a 200-slide deck before the
# first slide appears. A window is neither: the first click pays for eight
# slides, the next seven are instant.
DESIGN_WINDOW = 8


def _ensure_design_shot(job, index: int) -> str | None:
    """Render the window of slides containing `index`, once. Returns a readable
    reason when nothing could be rendered - the page then shows every finding
    and every remedy without a picture, which is the part a designer cannot do
    without."""
    shots = job.setdefault("design_shots", {})
    if index in shots:
        return None
    if job.get("deck") is None:
        return "the deck is no longer held in memory."
    start = (index // DESIGN_WINDOW) * DESIGN_WINDOW
    total = job["manifest"]["slides"]
    wanted = [i for i in range(start, min(total, start + DESIGN_WINDOW))
              if i not in shots]
    if not wanted:
        return None
    from .render import export_decks_png

    lock = _render_lock_for(job, "design_shot_lock")
    with lock:
        if index in shots:
            return None
        try:
            images = export_decks_png({"deck": job["deck"]}, wanted,
                                      width=THUMB_WIDTH)
        except Exception as exc:
            job["design_shot_error"] = f"{type(exc).__name__}: {exc}."
            return job["design_shot_error"]
        for key, png in (images or {}).items():
            shots[int(key.split(":", 1)[1])] = png
    if index not in shots:
        return ("PowerPoint returned no image for this slide. Rendering needs "
                "desktop PowerPoint or LibreOffice on this machine.")
    return None


def _design_severity_map(findings, records, answered) -> dict:
    """{slide_index: {severity: n}} for the dots on the slide strip, from both
    passes: a slide the designer should stop at is one with anything on it, and
    which check found it is not the question the strip answers."""
    out: dict[int, dict] = {}
    for f in findings:
        if f.finding_id in answered:
            continue
        for index in f.slides:
            slot = out.setdefault(index, {})
            slot[f.severity] = slot.get(f.severity, 0) + 1
    for r in records or []:
        if r.get("module") == "preflight":
            continue
        slot = out.setdefault(r["slide_index"], {})
        slot[r["severity"]] = slot.get(r["severity"], 0) + 1
    return out


def _design_page(job_id: str, job, banner: str = "", error: str | None = None,
                 view: str = "slide", current: int = 0):
    from .design import slide_rects
    from .promotion import promoted_issue_types
    from .ui_design import render_design

    err = _ensure_design(job) or error
    meta = next((p for p in _profiles_meta() if p["id"] == job.get("profile")),
                None)
    applied = job.get("design_applied") or []
    answered = {a.finding_id for a in applied}
    every = [f for f in (job.get("design") or [])
             if f.finding_id not in answered]
    total = job["manifest"]["slides"]
    current = max(0, min(current, max(0, total - 1)))
    records = [r for r in job["manifest"]["records"]
               if r.get("module") != "preflight"]

    # A finding that touches ONE slide belongs to that slide's page. One that
    # spans several is a single decision about the deck and cannot honestly be
    # asked on any one of them, so it goes to the deck view (and is counted
    # there, on the tab, so it is not lost).
    on_slide = [f for f in every if f.slides == [current]]
    deck_wide = [f for f in every if len(f.slides) != 1]

    shot_error = None
    rects: list = []
    if view != "deck":
        shot_error = _ensure_design_shot(job, current)
        if shot_error is None and job.get("deck") is not None:
            try:
                rects = slide_rects(job["deck"], on_slide).get(current, [])
            except Exception:
                rects = []

    chat_ok, chat_why = _chat_available()
    # A prepared deck came from one page carrying the coverage, the gaps and
    # these findings together. Sending a designer to the audit report from here
    # would strand the half of the answer that is about the master.
    back = ((f"/prep/{job_id}", "Back to the prepared deck")
            if job.get("prep") is not None else None)
    return HTMLResponse(render_design(
        job_tabs=_tabs_for(job_id, job, "design"),
        chat=chat_ok, chat_note=chat_why, back=back,
        deck_name=job["filename"],
        profile_name=(meta or {}).get("name") or job.get("profile") or "",
        job_id=job_id, view=view, current=current, total_slides=total,
        findings=on_slide, deck_findings=deck_wide, applied=applied,
        audit_records=[r for r in records if r["slide_index"] == current],
        rects=rects,
        per_slide=_design_severity_map(every, records, answered),
        banner=banner, error=err, render_error=shot_error,
        has_deck=job.get("deck") is not None,
        can_fix=job.get("deck") is not None,
        promoted=promoted_issue_types(),
        shot_tag=_render_tag(job),
        auto=_auto_plan(job, current)))


def _render_tag(job) -> str:
    """A short token for the deck bytes the pictures on a page were rendered
    FROM, to be spent in the image URL.

    Dropping the server's cached renders is only HALF of an update. The browser
    already holds the old PNG, cached for ten minutes under a URL that depends
    on nothing but the job and the slide number, so a designer who applied a
    decision was handed back the same URL and shown the deck as it stood before
    they touched it: the row said applied, the change really had been made, and
    the slide did not move (design lead, 26/08/2026). A URL that changes with
    the bytes cannot serve a stale picture, and the long cache stays useful for
    the pictures that genuinely have not changed.

    Hashed rather than counted. A counter is a second record of the same fact
    and it drifts the first time a route changes the deck and forgets to bump
    it; a digest of the bytes cannot disagree with the bytes. Computed once per
    change rather than once per page, on the same discipline the renders
    already follow - every place that assigns new deck bytes clears this too.
    """
    tag = job.get("render_tag")
    if tag is None:
        import hashlib

        deck = job.get("deck")
        tag = (hashlib.blake2s(deck, digest_size=6).hexdigest()
               if deck else "none")
        job["render_tag"] = tag
    return tag


def _expire_old_decks() -> None:
    """Trim the job registry, and drop the deck bytes of all but the newest few.

    Manifests stay - a designer can still read the report - but the bytes and
    everything derived from them go, which is the truth: the fix option really
    has expired. The dict is shared with the format registry, so this expires
    the download too.

    ONE COPY. This was written out twice, verbatim, in the two places that
    register a job, and the copies had already drifted from the rest of the
    codebase: both dropped the deck, the diff, the thumbnails and the rects, and
    neither dropped `design_shots` - a dict of rendered PNGs, by far the largest
    thing hanging off a job. An evicted deck was freeing its bytes and keeping
    its pictures (30/08/2026).

    Callers hold _jobs_lock.
    """
    while len(_jobs) > MAX_STORED_MANIFESTS:
        _jobs.popitem(last=False)
    with_deck = [k for k, v in _jobs.items() if v.get("deck") is not None]
    for key in with_deck[:-MAX_DECKS_IN_MEMORY]:
        job = _jobs[key]
        job["deck"] = None
        job.pop("prev_deck", None)
        # Everything else derived from those bytes, by the same rule that drops
        # it when a fix lands, so a new derived cache cannot be forgotten here.
        _invalidate_renders(job)


def _invalidate_renders(job) -> None:
    """Everything derived from the deck bytes, dropped. A cached thumbnail beside
    a row marked applied is a picture of the deck before the fix."""
    job["diff"] = None
    job["thumbs"] = None
    job["rects"] = None
    job["design"] = None
    job["design_shots"] = {}
    job["render_tag"] = None
    job.pop("design_shot_error", None)
    # The extraction is derived from the same bytes and goes stale the same way:
    # a recolour changes the palette roll-up, and a checklist showing the old one
    # beside a row marked applied is the same lie as a stale thumbnail.
    job.pop("extracted", None)


def _extracted(job) -> dict | None:
    """This job's deck as ground truth (qc.extract), read once and kept.

    Two pages want it - the colour and type checklist, and the ask box's palette
    facts - and both read it fresh on EVERY request. That is a full parse plus a
    colour and font resolution of every run on every slide (about 1.4s for a
    200-slide deck), repeated for a document that cannot change between requests
    unless a fix lands, and a fix already calls _invalidate_renders.

    Returns None when the deck bytes are gone, which the callers already handle:
    the checklist says so and the ask box drops the palette from its facts.
    """
    if job.get("deck") is None:
        return None
    cached = job.get("extracted")
    if cached is not None:
        return cached
    from .extract import extract_deck

    job["extracted"] = extract_deck(job["deck"])
    return job["extracted"]


@app.get("/audit/{job_id}", response_class=HTMLResponse)
def audit_view(request: Request, job_id: str):
    """The audit report for a job that already ran.

    The report used to exist only as the response to the POST that produced it,
    so navigating away from it - to the design page, say - was one-way. A page a
    designer cannot get back to is a page they will not leave."""
    from .promotion import promoted_issue_types
    from .store import comment_counts

    job = _job(job_id)
    if job is None:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES,
                                         "Unknown or expired job."),
                            status_code=404)
    return render_report(job["manifest"], job_id,
                         tabs=_tabs_for(job_id, job, "findings"),
                         can_fix=job.get("deck") is not None,
                         has_cleaned=job.get("cleaned") is not None,
                         promoted=promoted_issue_types(),
                         comments=comment_counts(job["filename"]),
                         design=design_count(job),
                         assist=AI_ENABLED and not _ephemeral(job["profile"]))


@app.get("/design/{job_id}", response_class=HTMLResponse)
def design_page(job_id: str, view: str = "slide", n: int = 0):
    job = _job(job_id)
    if job is None:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES,
                                         "Unknown or expired job."),
                            status_code=404)
    return _design_page(job_id, job,
                        view="deck" if view == "deck" else "slide",
                        current=max(0, n))


@app.get("/design-img/{job_id}/{idx}.png")
def design_img(job_id: str, idx: int):
    """One rendered slide for the design page.

    Its own route and its own cache rather than /thumb: that one is filled by
    _ensure_thumbs, which renders the WHOLE deck and whose callers (the visual
    PDF, the preview overlay) rely on it being complete. This page fills a
    window at a time, and a half-filled thumbs cache would make those callers
    think they had every slide."""
    from fastapi.responses import Response

    job = _job(job_id)
    if job is None:
        return JSONResponse({"error": "unknown or expired job"}, status_code=404)
    error = _ensure_design_shot(job, idx)
    png = (job.get("design_shots") or {}).get(idx)
    if png is None:
        return JSONResponse({"error": error or "no such slide"}, status_code=503)
    return Response(content=png, media_type="image/png",
                    headers={"Cache-Control": "private, max-age=600"})


def _perform_picks(job, picks):
    """Perform (finding, remedy) pairs, record them, re-audit. Returns
    (applied, staleness note or None).

    Shared by the card-by-card apply and by "let the tool decide", because the
    difference between those two is only who chose the remedy. A second copy of
    this would be a second set of rules about what an applied decision is, and
    Undo reads that record."""
    from .remedy import apply as apply_remedies

    deck, applied = apply_remedies(job["deck"], picks)
    acted = [a for a in applied if a.done and a.undo]
    with _jobs_lock:
        if acted:
            job["prev_deck"] = job["deck"]   # pre-change bytes, for the diff
            job["deck"] = deck
            job["cleaned"] = deck
        job["design_applied"] = (job.get("design_applied") or []) + applied
        _invalidate_renders(job)
    return applied, (_reaudit_in_place(job) if acted else None)


def _picks_note(applied: list) -> str:
    acted = [a for a in applied if a.done and a.undo]
    left = [a for a in applied if a.done and not a.undo]
    failed = [a for a in applied if not a.done]
    note = []
    if acted:
        note.append(f"Applied {len(acted)} change"
                    f"{'s' if len(acted) != 1 else ''}")
    if left:
        note.append(f"recorded {len(left)} as deliberate")
    if failed:
        note.append(f"{len(failed)} could not be applied and say why below")
    return "; ".join(note)


@app.post("/design/{job_id}/apply", response_class=HTMLResponse)
async def design_apply(request: Request, job_id: str):
    """Perform the remedy picked on each card, and nothing else.

    The radio group per finding is named pick_<finding_id>, so the form is read
    rather than declared: the set of findings is not known until the deck is
    scanned, and a fixed signature would have to be a list of opaque pairs."""
    job = _job(job_id)
    if job is None:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES,
                                         "Unknown or expired job."),
                            status_code=404)
    if job.get("deck") is None:
        return _design_page(job_id, job)

    _ensure_design(job)
    by_id = {f.finding_id: f for f in job.get("design") or []}
    already = {a.finding_id for a in job.get("design_applied") or []}
    form = await request.form()
    try:
        back_to = max(0, int(form.get("n") or 0))
    except (TypeError, ValueError):
        back_to = 0
    picks = []
    for key, value in form.multi_items():
        if not str(key).startswith("pick_"):
            continue
        finding = by_id.get(str(key)[5:])
        if finding is None or finding.finding_id in already:
            continue
        remedy = next((o for o in finding.options
                       if o.remedy_id == str(value)), None)
        if remedy is not None:
            picks.append((finding, remedy))
    if not picks:
        return _design_page(job_id, job, current=back_to,
                            banner="No card had an answer picked, so nothing "
                                   "was changed.")

    try:
        applied, stale = _perform_picks(job, picks)
    except Exception as exc:
        return _design_page(job_id, job, current=back_to,
                            error=f"Nothing was changed: {type(exc).__name__}: "
                                  f"{exc}")
    return _design_page(job_id, job, current=back_to,
                        banner=_picks_note(applied) + ".", error=stale)


# ------------------------------------------------- the audit's own fixes, here
#
# The rows under "Also on this slide" were read-only for one release, on the
# argument that duplicating the report's tick box would give a designer two
# buttons for one action. The argument was wrong in the only way that matters:
# it is not two actions, it is one action reachable from the place the problem
# is visible. A designer looking at slide 7, reading "Calibri is not in the
# allowed set", was being told to go to another page and find the same row
# (design lead, 24/08/2026).
#
# So the tick is here too, and it is THE SAME tick: same qc.fixer, same
# selection rules, same verify-after-write (_apply_audit_fixes). What is not
# duplicated is the engine.


@app.post("/design/{job_id}/fix", response_class=HTMLResponse)
def design_fix(job_id: str, record_ids: list[str] = Form(None),
               n: int = Form(0)):
    """Apply the audit fixes ticked on the design page, and come back to the
    slide they were ticked on."""
    job = _job(job_id)
    if job is None:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES,
                                         "Unknown or expired job."),
                            status_code=404)
    back_to = max(0, n or 0)
    if job.get("deck") is None:
        return _design_page(job_id, job, current=back_to)
    selected = set(record_ids or [])
    if not selected:
        return _design_page(job_id, job, current=back_to,
                            banner="No fix was ticked, so nothing was changed.")

    fx, before_total, stale = _apply_audit_fixes(job, selected)
    skipped = [o for o in fx.outcomes if o.outcome == "skipped"]
    note = (f"Applied {fx.applied} fix{'es' if fx.applied != 1 else ''}. "
            f"Re-audit of the deck: "
            f"{job['manifest']['summary']['total']} findings remain "
            f"(was {before_total}).")
    if skipped:
        note += (f" Skipped {len(skipped)}: "
                 + "; ".join(sorted({o.reason for o in skipped if o.reason}))
                 + ".")
    return _design_page(job_id, job, current=back_to, banner=note, error=stale)


def _auto_targets(job, scope: str, current: int) -> tuple[list, list]:
    """What "let the tool decide" would look at: (fixable audit records, open
    design findings), for one slide or for the whole deck.

    The same function counts the work for the button and selects it for the
    route, so the number a designer reads before pressing is the number that
    gets acted on."""
    from .fixer import is_fixable

    deck_wide = scope == "deck"
    records = [r for r in job["manifest"]["records"]
               if r.get("module") != "preflight" and is_fixable(r)
               and (deck_wide or r["slide_index"] == current)]
    answered = {a.finding_id for a in job.get("design_applied") or []}
    findings = [f for f in (job.get("design") or [])
                if f.finding_id not in answered
                and (deck_wide or f.slides == [current])]
    return records, findings


def _auto_plan(job, current: int) -> dict:
    """The counts behind the two buttons on the hand-it-over card, and the
    checks' own words for what they would not answer."""
    from .design import auto_choice, auto_skip_reason
    from .fixer import needs_explicit_tick

    if job.get("deck") is None:
        return {}
    _ensure_design(job)
    plan = {}
    for scope in ("slide", "deck"):
        records, findings = _auto_targets(job, scope, current)
        held = [r for r in records if needs_explicit_tick(r)]
        decides = [f for f in findings if auto_choice(f) is not None]
        left = [f for f in findings if auto_choice(f) is None]
        plan[scope] = {
            "fixes": len(records) - len(held),
            "held": len(held),
            "picks": len(decides),
            "left": len(left),
            "reasons": sorted({auto_skip_reason(f) for f in left} - {None}),
        }
    return plan


def _auto_apply(job, scope: str, current: int, holds_too: bool) -> dict:
    """Apply everything the tool has an answer for, on one slide or the deck.

    Two passes, in this order and for this reason: the audit's own fixes are
    deterministic conformance to the profile, and the design decisions are
    judgments about how the deck then LOOKS. Making the judgments first would
    mean judging a slide the tool is about to change - a contrast reading taken
    before a font substitution, an overlap measured before a shape is snapped
    back onto its margin. So the deck is brought into line first, re-audited,
    and only then is it looked at.

    Returns {note, stale, error, stage, untouched, held}. `stage` says how far
    it got, because a failure in pass two leaves a deck that pass one already
    changed and re-audited, and telling a designer "nothing happened" then
    would be false. The button on the design page and the sentence in the ask
    box both come through here, so they cannot mean different things.
    """
    from .design import auto_choice
    from .fixer import needs_explicit_tick

    out = {"note": [], "stale": None, "error": None, "stage": "none",
           "untouched": [], "held": []}

    # --- pass one: the audit's own fixes
    records, _ = _auto_targets(job, scope, current)
    out["held"] = [r for r in records if needs_explicit_tick(r)]
    chosen = {r["record_id"] for r in records
              if holds_too or not needs_explicit_tick(r)}
    if chosen:
        try:
            fx, _before, out["stale"] = _apply_audit_fixes(job, chosen)
        except Exception as exc:
            out["error"] = (f"Nothing was changed: the audit fixes could not "
                            f"be applied ({type(exc).__name__}: {exc}), so the "
                            f"design decisions were not attempted either.")
            return out
        out["stage"] = "fixes"
        out["note"].append(f"applied {fx.applied} audit fix"
                           f"{'es' if fx.applied != 1 else ''}")
        skipped = len([o for o in fx.outcomes if o.outcome == "skipped"])
        if skipped:
            out["note"].append(f"skipped {skipped} the fixer could not perform")

    # --- pass two: the design decisions, on the deck pass one left behind
    _ensure_design(job)
    _records, findings = _auto_targets(job, scope, current)
    picks = []
    for finding in findings:
        remedy = auto_choice(finding)
        if remedy is None:
            out["untouched"].append(finding)
        else:
            picks.append((finding, remedy))
    if picks:
        try:
            applied, pick_stale = _perform_picks(job, picks)
        except Exception as exc:
            # Pass one already landed and is already re-audited; saying so is
            # the difference between a designer who knows what state the deck
            # is in and one who presses the button again.
            out["error"] = (f"The design decisions could not be performed: "
                            f"{type(exc).__name__}: {exc}")
            return out
        done = _picks_note(applied)
        if done:
            out["note"].append(done[0].lower() + done[1:])
        out["stale"] = out["stale"] or pick_stale
    out["stage"] = "all"
    return out


@app.post("/design/{job_id}/auto", response_class=HTMLResponse)
def design_auto(job_id: str, scope: str = Form("slide"), n: int = Form(0),
                include_holds: str = Form(None)):
    """Let the tool answer everything it has an answer for.

    Every decision made here lands in the same list as a hand-picked one, with
    the same Undo. That is the whole safety story: this is a starting point a
    designer corrects, not a commitment they have to accept.
    """
    from .design import auto_skip_reason

    job = _job(job_id)
    if job is None:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES,
                                         "Unknown or expired job."),
                            status_code=404)
    back_to = max(0, n or 0)
    scope = "deck" if scope == "deck" else "slide"
    if job.get("deck") is None:
        return _design_page(job_id, job, current=back_to)
    holds_too = bool(include_holds)

    result = _auto_apply(job, scope, back_to, holds_too)
    note, stale = result["note"], result["stale"]
    untouched, held = result["untouched"], result["held"]
    if result["error"]:
        return _design_page(
            job_id, job, current=back_to,
            banner=(("Applied the audit fixes only: " + "; ".join(note) + ".")
                    if note and result["stage"] == "fixes" else ""),
            error=result["error"])

    where = "the whole deck" if scope == "deck" else f"slide {back_to + 1}"
    if not note:
        banner = (f"Nothing on {where} was the tool's to decide, so nothing "
                  f"changed.")
    else:
        banner = f"Decided {where}: " + "; ".join(note) + "."
    reasons = sorted({auto_skip_reason(f) for f in untouched} - {None})
    if untouched:
        banner += (f" {len(untouched)} left for you, because "
                   + "; and ".join(reasons) + ".")
    if held and not holds_too:
        banner += (f" {len(held)} fix{'es' if len(held) != 1 else ''} still "
                   f"ask{'' if len(held) != 1 else 's'} for your explicit "
                   f"approval and {'were' if len(held) != 1 else 'was'} not "
                   f"applied.")
    return _design_page(job_id, job, current=back_to, banner=banner,
                        error=stale)


def _perform_undo(job, wanted: set) -> dict:
    """Take the named decisions back, exactly, and say what came with them.

    Returns {chain, dragged, put_back, error}. One implementation for the same
    reason the apply path has one: the Undo button and the ask box are two
    doors onto a replay that must never happen twice, and a second copy of the
    chain arithmetic is how that starts.
    """
    from .remedy import followers, undo_items
    from .undo import apply_undo

    applied = list(job.get("design_applied") or [])

    # Chains are unioned before anything is replayed: two requested decisions
    # can drag the same third one, and undoing it twice would put back a state
    # from before the first replay.
    chain_ids: set = set()
    for finding_id in wanted:
        chain_ids.update(a.finding_id for a in followers(applied, finding_id))
    chain = [a for a in applied if a.finding_id in chain_ids]
    dragged = [a for a in chain if a.finding_id not in wanted]

    error = None
    outcomes = []
    items = undo_items([a for a in chain if a.undo])
    if items and job.get("deck") is not None:
        try:
            deck, outcomes = apply_undo(job["deck"], items)
        except Exception as exc:
            error = (f"The undo failed and nothing was changed: "
                     f"{type(exc).__name__}: {exc}")
        else:
            with _jobs_lock:
                job["deck"] = deck
                job["cleaned"] = deck
                _invalidate_renders(job)
    elif items:
        error = ("The deck is no longer in memory, so these changes cannot be "
                 "reversed here. The decision has been cleared from the list.")

    if error is None or not items:
        with _jobs_lock:
            job["design_applied"] = [a for a in applied
                                     if a.finding_id not in chain_ids]
        if items:
            _reaudit_in_place(job)

    return {"chain": chain, "dragged": dragged, "error": error,
            "put_back": sum(1 for o in outcomes if o.get("done"))}


def _undo_note(result: dict) -> str:
    chain, dragged = result["chain"], result["dragged"]
    put_back = result["put_back"]
    note = [f"{len(chain)} decision{'s' if len(chain) != 1 else ''} reopened"]
    if put_back:
        note.append(f"{put_back} change{'s' if put_back != 1 else ''} put back "
                    f"exactly as {'they were' if put_back != 1 else 'it was'}")
    if dragged:
        note.append(f"including {len(dragged)} that touched the same shape and "
                    f"could not come back on its own")
    return "; ".join(note) + "."


@app.post("/design/{job_id}/undo", response_class=HTMLResponse)
def design_undo(job_id: str, finding_ids: list[str] = Form(None),
                n: int = Form(None)):
    """Take one decision back, exactly, and say what came with it."""
    job = _job(job_id)
    if job is None:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES,
                                         "Unknown or expired job."),
                            status_code=404)
    # Which view to answer on: the slide the button was pressed on, or the
    # deck-wide tab when it carried no slide.
    view = "slide" if n is not None else "deck"
    back_to = max(0, n or 0)
    wanted = set(finding_ids or [])
    if not wanted:
        return _design_page(job_id, job, view=view, current=back_to)

    result = _perform_undo(job, wanted)
    # Back to the view the Undo button was on. Returning a designer to slide 1
    # after they pressed a button on slide 7 is the small rudeness that makes a
    # tool tiring to use.
    return _design_page(job_id, job, view=view, current=back_to,
                        banner=_undo_note(result), error=result["error"])


def _render_lock_for(job, name: str) -> threading.Lock:
    """The lock guarding one job's renders, kept ON the job.

    It used to live in a module-level dict keyed by job id and, in one case, by
    `id(job)`. Both were wrong in the same two ways: nothing ever removed an
    entry, so the dict grew for the life of the process and pinned a Lock per
    job ever seen; and id() is reused by CPython the moment an evicted job is
    collected, so a new job could be handed the lock of a dead one and two
    requests could render into each other.

    On the job, it is created once, shared by every request for that job, and
    collected with it. setdefault is the atomic step, exactly as before.
    """
    return job.setdefault(name, threading.Lock())

THUMB_WIDTH = 1100


def _ensure_thumbs(job_id: str, job: dict) -> None:
    """Render every slide of the job's current deck once (PowerPoint COM),
    cache in memory. Raises RuntimeError when rendering is unavailable."""
    if job.get("thumbs"):
        return
    if job.get("deck") is None:
        raise RuntimeError("deck bytes no longer in memory; re-upload to preview")
    from .render import export_decks_png

    lock = _render_lock_for(job, "thumbs_lock")
    with lock:
        if job.get("thumbs"):
            return
        n = job["manifest"]["slides"]
        images = export_decks_png({"deck": job["deck"]}, list(range(n)),
                                  width=THUMB_WIDTH)
        job["thumbs"] = {int(k.split(":", 1)[1]): v for k, v in images.items()}


def _ensure_rects(job: dict) -> dict:
    if job.get("rects") is None:
        from .render import audit_rects

        flagged = [r for r in job["manifest"]["records"]
                   if r["module"] != "preflight"]
        job["rects"] = audit_rects(job["deck"], flagged)
    return job["rects"]


# --------------------------------------------------- ask about this deck
# One door onto what the passes recorded, because a designer holding a client
# deck does not arrive knowing which of the five pages answers their question.
# It answers and it navigates; it cannot act (qc.chat).


def _chat_available() -> tuple[bool, str]:
    """Whether a question can be answered on this host, and why not when it
    cannot. Same three reasons as the coverage review, and worth the same
    sentence each: a text box that silently does nothing is worse than an
    absence with a reason beside it."""
    from .llm import api_configured

    if not AI_ENABLED:
        return False, ("Model-backed features are switched off on this host, so "
                       "there is nothing to ask. The reports below are the same "
                       "facts it would have read.")
    if not api_configured():
        return False, ("No model key on this host, so questions cannot be "
                       "answered and the visual passes are off too. Set "
                       "GEMINI_API_KEY in the .env file at the project root and "
                       "restart. The reports below are the same facts it would "
                       "have read.")
    return True, ""


@app.get("/checklist/{job_id}", response_class=HTMLResponse)
def checklist(job_id: str):
    """What this deck is made of: every colour and every typeface, with the level
    each one comes from.

    Read-only and read fresh out of the deck's own bytes. It serves both job
    kinds off one route because the question does not depend on which pass ran -
    a designer wants the palette whether they audited the deck or rebuilt it."""
    from .extract import font_inventory, palette_inventory
    from .ui_checklist import render_checklist

    job = _job(job_id)
    back = f"/design/{job_id}"
    if job is None:
        with _format_lock:
            job = _format_jobs.get(job_id)
        back = f"/format/{job_id}/review?view=deck"
    if job is not None and job.get("prep") is not None:
        back = f"/prep/{job_id}"
    if job is None:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES,
                                        "Unknown or expired job."),
                            status_code=404)
    if job.get("deck") is None:
        return HTMLResponse(render_index(
            _pickable_profiles(), MODULES,
            "The deck is no longer held in memory, so its colours and type "
            "cannot be read. Re-upload it."), status_code=410)

    try:
        deck = _extracted(job)
        palette, fonts = palette_inventory(deck), font_inventory(deck)
    except Exception as exc:
        return HTMLResponse(render_index(
            _pickable_profiles(), MODULES,
            f"That deck's colours could not be read: {type(exc).__name__}: "
            f"{exc}"), status_code=422)

    return HTMLResponse(render_checklist(
        tabs=_tabs_for(job_id, job, "checklist"),
        deck_name=job.get("filename") or "deck.pptx", job_id=job_id,
        back=back, palette=palette, fonts=fonts))


def _chat_job(job_id: str) -> tuple[dict | None, str]:
    """(job, kind) for whichever registry holds it.

    A prep job is in both, and it answers as "prep": it has plans AND records,
    so answering it as one or the other would leave half of what a designer can
    see on the page out of the fact sheet."""
    job = _job(job_id)
    if job is not None:
        return job, ("prep" if job.get("prep") is not None else "audit")
    with _format_lock:
        job = _format_jobs.get(job_id)
    if job is None:
        return None, ""
    return job, ("prep" if job.get("prep") is not None else "format")


# How many proposed plans one job keeps. A designer asks two or three things
# before pressing anything; beyond that the older ones describe a deck that has
# since moved, and a plan resolved against a deck that no longer exists is the
# one thing this must not perform.
MAX_PENDING_PLANS = 8


def _remember_plan(job, plan) -> str:
    """Hold a proposed plan until it is confirmed, and hand back its handle."""
    pending = job.get("chat_plans")
    if pending is None:
        pending = job["chat_plans"] = OrderedDict()
    token = uuid.uuid4().hex[:12]
    pending[token] = plan
    while len(pending) > MAX_PENDING_PLANS:
        pending.popitem(last=False)
    return token


@app.post("/chat/{job_id}")
async def chat(request: Request, job_id: str):
    """Answer one question about one job, and propose what it asked for.

    THIS ROUTE CHANGES NOTHING. An action comes back as a plan and a handle;
    performing it is the route below, which the designer reaches by pressing a
    button that says what would happen."""
    off = _ai_disabled_response()
    if off is not None:
        return off

    available, why = _chat_available()
    if not available:
        return JSONResponse({"error": why}, status_code=503)

    from .chat import ask

    try:
        payload = await request.json()
    except Exception:
        payload = {}
    question = str((payload or {}).get("q") or "")

    job, kind = _chat_job(job_id)
    if job is None:
        return JSONResponse({"error": "That job is no longer held in memory. "
                                      "Re-run the deck and ask again."},
                            status_code=404)

    try:
        out = ask(job, kind, job_id, question)
    except Exception as exc:
        # A model that cannot be reached is not an answer of "nothing found".
        #
        # WITH the reason, not just the exception class. "That could not be
        # answered (RuntimeError)" is what this said for a fortnight while
        # every call was in fact dying on a closed HTTP client, and the type
        # name alone was not enough for anyone to guess that (design lead,
        # 27/08/2026 - "why doesn't the chatbot work?"). This is a local pilot
        # tool: the person reading the box is the person who can fix the cause.
        reason = " ".join(str(exc).split())[:300]
        return JSONResponse(
            {"error": f"That could not be answered. {type(exc).__name__}: "
                      f"{reason or 'no detail given'}. The reports on this "
                      f"page are unaffected."},
            status_code=503)

    plan = out.pop("plan", None)
    if plan is not None:
        out["plan"] = {"token": _remember_plan(job, plan),
                       "summary": plan.summary, "changes": plan.changes,
                       "name": plan.name}
    return JSONResponse(out)


def _perform_plan(job, plan) -> tuple[str, bool]:
    """Perform one confirmed plan. Returns (what happened, redraw the page).

    Every branch calls the SAME function the button on the page calls, and none
    of them re-implements a fix, a re-audit or an undo chain. That is what makes
    the ask box a new door onto the tool rather than a second tool: there is no
    change reachable from here that is not in the Undo list, and none that skips
    the verify-after-write.
    """
    if plan.name == "fix_findings":
        fx, before, stale = _apply_audit_fixes(job, set(plan.record_ids))
        after = job["manifest"]["summary"]["total"]
        note = (f"Applied {fx.applied} fix{'es' if fx.applied != 1 else ''}. "
                f"Re-audit of the deck: {after} findings remain (was {before}).")
        skipped = [o for o in fx.outcomes if o.outcome == "skipped"]
        if skipped:
            note += f" Skipped {len(skipped)}."
        return note + (f" {stale}" if stale else ""), True

    if plan.name == "decide":
        result = _auto_apply(job, plan.scope, plan.slide or 0,
                             plan.include_holds)
        if result["error"] and result["stage"] == "none":
            return result["error"], False
        where = ("the whole deck" if plan.slide is None
                 else f"slide {plan.slide + 1}")
        if result["note"]:
            note = f"Decided {where}: " + "; ".join(result["note"]) + "."
        else:
            note = (f"Nothing on {where} turned out to be the tool's to "
                    f"decide, so nothing changed.")
        if result["error"]:
            note += " " + result["error"]
        if result["untouched"]:
            note += f" {len(result['untouched'])} left for you."
        return note + (f" {result['stale']}" if result["stale"] else ""), True

    if plan.name == "take_remedy":
        applied, stale = _perform_picks(job, plan.picks)
        note = _picks_note(applied) or "Nothing was changed"
        if not note.endswith("."):
            note += "."
        return note + (f" {stale}" if stale else ""), True

    if plan.name == "undo":
        result = _perform_undo(job, set(plan.finding_ids))
        if result["error"]:
            return result["error"], False
        return _undo_note(result), True

    if plan.name == "recheck":
        stale = _reaudit_in_place(job)
        if stale:
            return stale, False
        total = job["manifest"]["summary"]["total"]
        return (f"Read the deck again: {total} finding"
                f"{'s' if total != 1 else ''} on it as it now stands."), True

    if plan.name == "remove_pieces":
        removed = _perform_removals(job, set(plan.remove_ids))
        if not removed:
            return (job.get("remove_error")
                    or "Nothing was removed."), False
        return (f"Removed {removed} piece{'s' if removed != 1 else ''}. Each "
                f"one is in the change list with its own Undo."), True

    return "That plan is not one this build can perform.", False


@app.post("/chat/{job_id}/do")
async def chat_do(request: Request, job_id: str):
    """Perform one plan the ask box proposed, once.

    THE TOKEN IS SPENT ON USE. A plan describes a deck in a particular state -
    these record ids, this many findings - and performing it twice would apply
    it to a deck the first pass already changed. So it is popped before
    anything runs, and a second press is told the plan is gone rather than
    quietly doing it again.
    """
    off = _ai_disabled_response()
    if off is not None:
        return off

    try:
        payload = await request.json()
    except Exception:
        payload = {}
    token = str((payload or {}).get("token") or "")

    job, _kind = _chat_job(job_id)
    if job is None:
        return JSONResponse({"error": "That job is no longer held in memory, "
                                      "so nothing was changed."},
                            status_code=404)
    plan = (job.get("chat_plans") or {}).pop(token, None)
    if plan is None:
        return JSONResponse(
            {"error": "That plan is no longer on offer, so nothing was "
                      "changed. Ask again and the tool will re-read the deck "
                      "as it now stands."}, status_code=409)
    if plan.changes and job.get("deck") is None:
        return JSONResponse(
            {"error": "The deck is no longer held in memory, so nothing could "
                      "be changed."}, status_code=410)

    try:
        note, redraw = _perform_plan(job, plan)
    except Exception as exc:
        return JSONResponse(
            {"error": f"That could not be performed ({type(exc).__name__}: "
                      f"{exc}). Open the page it belongs to and check the "
                      f"deck before pressing anything else."}, status_code=500)
    return JSONResponse({"done": True, "note": note, "reload": redraw})


@app.get("/signin", response_class=HTMLResponse)
def signin_page():
    from .store import list_users
    from .ui import render_signin

    return render_signin(list_users())


@app.post("/signin", response_class=HTMLResponse)
def signin(request: Request, name: str = Form(...), pin: str = Form(...)):
    from fastapi.responses import RedirectResponse

    from .store import (create_session, get_user, has_pin, list_users,
                        set_pin, verify_pin)
    from .ui import render_signin

    if get_user(name) is None:
        return HTMLResponse(render_signin(list_users(), "Unknown name."),
                            status_code=404)
    if not has_pin(name):
        try:
            set_pin(name, pin)  # first sign-in sets the PIN
        except ValueError as exc:
            return HTMLResponse(render_signin(list_users(), str(exc)),
                                status_code=400)
    elif not verify_pin(name, pin):
        return HTMLResponse(render_signin(list_users(), "Wrong PIN."),
                            status_code=401)
    resp = RedirectResponse("/", status_code=303)
    resp.set_cookie("qc_session", create_session(name), max_age=30 * 24 * 3600,
                    httponly=True)
    resp.delete_cookie("qc_user")
    return resp


@app.post("/signout")
def signout(request: Request):
    from .store import delete_session

    delete_session(request.cookies.get("qc_session", ""))
    resp = JSONResponse({"ok": True})
    resp.delete_cookie("qc_session")
    resp.delete_cookie("qc_user")
    return resp


@app.get("/me")
def me(request: Request):
    from .auth import current_user
    from .store import list_users

    user = current_user(request)
    if user:
        user = {k: v for k, v in user.items() if k != "pin_hash"}
    return JSONResponse({"user": user, "users": [
        {k: v for k, v in u.items() if k != "pin_hash"} for u in list_users()]})


@app.post("/whoami")
async def whoami(request: Request):
    from .store import get_user

    data = await request.json()
    user = get_user(data.get("name", ""))
    if user is None:
        return JSONResponse({"error": "unknown user"}, status_code=404)
    resp = JSONResponse({"ok": True, "user": user})
    resp.set_cookie("qc_user", user["name"], max_age=90 * 24 * 3600)
    return resp


@app.get("/comments")
def get_comments(deck: str, slide: int):
    from .store import comments_for

    return JSONResponse({"comments": comments_for(deck, slide)})


@app.post("/comments")
async def post_comment(request: Request):
    from .store import add_comment, get_user

    from .auth import current_user

    user = current_user(request)
    if user is None:
        return JSONResponse({"error": "sign in first (top right)"},
                            status_code=401)
    author = user["name"]
    data = await request.json()
    try:
        comment = add_comment(deck=data["deck"], slide_index=int(data["slide_index"]),
                              author=author, text=data.get("text", ""),
                              record_id=data.get("record_id"))
    except (KeyError, ValueError) as exc:
        return JSONResponse({"error": str(exc)}, status_code=400)
    return JSONResponse({"ok": True, "comment": comment})


@app.post("/triage")
async def triage(request: Request):
    from .triage import log_triage

    data = await request.json()
    job = _job(data.get("job_id", ""))
    state = data.get("state", "")
    if job is None or state not in ("confirmed", "false_positive", "cleared"):
        return JSONResponse({"error": "bad request"}, status_code=400)
    record = next((r for r in job["manifest"]["records"]
                   if r["record_id"] == data.get("record_id")), None)
    if record is None:
        return JSONResponse({"error": "unknown record"}, status_code=404)

    states = job.setdefault("triage", {})
    if state == "cleared":
        states.pop(record["record_id"], None)
    else:
        states[record["record_id"]] = state
    log_triage(record, state, job["filename"], job["manifest"]["profile_id"])
    counts = {"confirmed": sum(1 for v in states.values() if v == "confirmed"),
              "false_positive": sum(1 for v in states.values() if v == "false_positive")}
    return JSONResponse({"ok": True, "counts": counts})


@app.post("/assist/{job_id}")
def assist_questions(job_id: str):
    """Build clarifying questions from this job's findings. The issued
    question set (with actions) is held server-side; /assist/apply accepts
    only question ids, never client-supplied action payloads."""
    from .assist import aggregate, generate_questions

    off = _ai_disabled_response()
    if off is not None:
        return off

    job = _job(job_id)
    if job is None:
        return JSONResponse({"error": "unknown or expired job id"}, status_code=404)
    if _ephemeral(job["profile"]):
        which = ("self-consistency audit" if job["profile"] == SELF_PROFILE
                 else "master you uploaded")
        extra = ("" if job["profile"] == SELF_PROFILE else
                 " Save the master as a profile from the Read a master page "
                 "first, then audit against it.")
        return JSONResponse(
            {"error": f"The assistant tunes a saved profile; the {which} has "
                      f"none.{extra}"}, status_code=400)
    profile_obj = job.get("profile_obj")
    agg = aggregate(job["manifest"], job.get("deck"),
                    profile_obj.config if profile_obj is not None else None)
    questions, source = generate_questions(agg)
    job["assist"] = {q["id"]: q for q in questions}
    public = [{k: q[k] for k in ("id", "question", "rationale", "impact")}
              for q in questions]
    return {"questions": public, "source": source,
            "profile": job["profile"]}


@app.post("/assist/{job_id}/apply")
def assist_apply(request: Request, job_id: str,
                 accepted: list[str] = Form(None)):
    from .assist import apply_actions
    from .web_admin import _editor

    off = _ai_disabled_response()
    if off is not None:
        return off

    job = _job(job_id)
    if job is None:
        return JSONResponse({"error": "unknown or expired job id"}, status_code=404)
    editor = _editor(request)
    if editor is None:
        return JSONResponse(
            {"error": "Profile changes need a lead or admin; sign in with "
                      "such a role first."}, status_code=403)
    issued = job.get("assist") or {}
    actions = [issued[qid]["action"] for qid in (accepted or [])
               if qid in issued]
    if not actions:
        return JSONResponse({"error": "no accepted questions"}, status_code=400)
    result = apply_actions(job["profile"], actions, editor["name"])
    return {"profile": job["profile"], **result}


@app.post("/reaudit/{job_id}", response_class=HTMLResponse)
def reaudit(request: Request, job_id: str):
    """Re-run the audit on the retained deck bytes with the profile as it
    is NOW (e.g. right after assistant answers updated it)."""
    job = _job(job_id)
    if job is None:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES,
                                         "Unknown or expired job."), status_code=404)
    if job["deck"] is None:
        return HTMLResponse(render_index(
            _pickable_profiles(), MODULES,
            "The deck is no longer held in memory; re-upload it."), status_code=410)

    fd, tmp_name = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    tmp = Path(tmp_name)
    try:
        tmp.write_bytes(job["deck"])
        # A master-derived profile cannot be re-resolved: the master's bytes
        # were dropped after extraction, by design. Reuse the profile the job
        # already carries. Saved profiles still re-resolve, which is the whole
        # point of this route (it picks up assistant edits).
        profile_obj = (job["profile_obj"] if job["profile"] == MASTER_PROFILE
                       else _resolve_profile(job["profile"], job["deck"]))
        result = run_audit(tmp, profile_obj)
    finally:
        tmp.unlink(missing_ok=True)

    manifest = result.to_manifest()
    manifest["deck"] = job["filename"]
    new_id = uuid.uuid4().hex
    with _jobs_lock:
        _jobs[new_id] = {"manifest": manifest, "deck": job["deck"],
                         "cleaned": None, "filename": job["filename"],
                         "profile": job["profile"], "profile_obj": profile_obj,
                         "master_spec": job.get("master_spec")}
    from .auth import current_user
    from .promotion import promoted_issue_types
    from .store import comment_counts, record_audit

    user = current_user(request)
    record_audit(manifest, user["name"] if user else "anonymous")
    return render_report(manifest, new_id, can_fix=True,
                         tabs=_tabs_for(new_id, _job(new_id), "findings"),
                         promoted=promoted_issue_types(),
                         comments=comment_counts(job["filename"]),
                         assist=AI_ENABLED and not _ephemeral(job["profile"]),
                         banner="Re-audited with the updated profile.")


@app.post("/copilot/{job_id}", response_class=HTMLResponse)
def copilot(request: Request, job_id: str):
    """Design copilot: render the slides, let the vision pass propose layout
    actions, verify them in code, and merge the survivors into the report
    as ordinary tickable suggestions."""
    # qc.llm's check, not qc.assist's. Both now resolve to the same key, but the
    # seam is still the right one to ask: qc.llm is where the model lives, and
    # asking a pass about another pass's credentials is how this button came to
    # refuse on a correctly configured host in the first place.
    from .llm import api_configured
    off = _ai_disabled_response()
    if off is not None:
        return off

    from .copilot import run_copilot

    job = _job(job_id)
    if job is None:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES,
                                         "Unknown or expired job."), status_code=404)
    if job["deck"] is None:
        return HTMLResponse(render_index(
            _pickable_profiles(), MODULES,
            "The deck is no longer held in memory; re-upload it."), status_code=410)

    def _report(banner):
        # Back to the page the button was on. A prepared deck reaching the
        # audit report here would lose the coverage and the gaps, which are
        # half of what that page is for.
        if job.get("prep") is not None:
            return _prep_page(job_id, job, banner=banner)

        from .promotion import promoted_issue_types
        from .store import comment_counts

        return render_report(job["manifest"], job_id, can_fix=True,
                             tabs=_tabs_for(job_id, job, "findings"),
                             banner=banner, promoted=promoted_issue_types(),
                             comments=comment_counts(job["filename"]),
                             assist=AI_ENABLED and not _ephemeral(job["profile"]))

    if not api_configured():
        return _report("Design copilot needs an API key for the configured "
                       "model provider on the server.")
    try:
        _ensure_thumbs(job_id, job)
    except RuntimeError as exc:
        return _report(f"Design copilot needs slide renders: {exc}")

    new_records, reviewed = run_copilot(job["deck"], job["thumbs"],
                                        job["manifest"])
    _merge_records(job, new_records)
    return _report(f"Design copilot reviewed {reviewed} slide"
                   f"{'s' if reviewed != 1 else ''} and added "
                   f"{len(new_records)} suggestion"
                   f"{'s' if len(new_records) != 1 else ''} "
                   "(tickable, never pre-selected).")


def _merge_records(job, new_records: list[dict]) -> None:
    """Fold reviewed suggestions into the manifest and re-count it.

    Shared by the copilot and the component review: two passes appending to one
    manifest with two copies of the summary arithmetic is how the counts on the
    report start disagreeing with the rows under them."""
    from collections import Counter

    if not new_records:
        return
    manifest = job["manifest"]
    manifest["records"].extend(new_records)
    records = manifest["records"]
    manifest["summary"] = {
        "by_severity": dict(Counter(r["severity"] for r in records)),
        "by_issue_type": dict(Counter(r["issue_type"] for r in records)),
        "by_module": dict(Counter(r["module"] for r in records)),
        "arabic_flagged": sum(1 for r in records if r["arabic_flag"]),
        "total": len(records),
    }
    job["rects"] = None  # pins are renumbered on next preview


@app.post("/components/{job_id}", response_class=HTMLResponse)
def components(request: Request, job_id: str):
    """Component review: the visual model names what the things on each slide
    ARE and which line they belong on; code measures, computes the targets and
    emits ordinary tickable records.

    Its own button rather than part of the audit, for the same two reasons the
    copilot has one: it costs a vision call per slide, and it sends slide
    IMAGES to the API. Neither belongs on the path a designer takes by
    default. It is offered on BOTH result pages, because it answers the
    question a designer asks on either one.

    Gated on qc.llm, not on qc.assist. This asked for a second vendor's key and
    package while run_components has gone through qc.llm with every other
    judgment pass since 30/08/2026 - so the button refused with a message naming
    a key the pass does not use and would never have needed (31/08/2026)."""
    from .llm import api_configured, configuration_note

    off = _ai_disabled_response()
    if off is not None:
        return off

    from .components import run_components

    job = _job(job_id)
    if job is None:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES,
                                         "Unknown or expired job."),
                            status_code=404)
    if job["deck"] is None:
        return HTMLResponse(render_index(
            _pickable_profiles(), MODULES,
            "The deck is no longer held in memory; re-upload it."),
            status_code=410)

    def _report(banner):
        # Back to the page the button was on. A prepared deck reaching the
        # audit report here would lose the coverage and the gaps, which are
        # half of what that page is for.
        if job.get("prep") is not None:
            return _prep_page(job_id, job, banner=banner)

        from .promotion import promoted_issue_types
        from .store import comment_counts

        return render_report(job["manifest"], job_id, can_fix=True,
                             tabs=_tabs_for(job_id, job, "findings"),
                             banner=banner, promoted=promoted_issue_types(),
                             comments=comment_counts(job["filename"]),
                             assist=AI_ENABLED and not _ephemeral(job["profile"]))

    if not api_configured():
        return _report(f"Component review needs a model key on the server: "
                       f"{configuration_note() or _MODEL_KEY_NOTE}")
    try:
        _ensure_thumbs(job_id, job)
    except RuntimeError as exc:
        return _report(f"Component review needs slide renders: {exc}")

    # The frame is handed in rather than re-read: the audit already resolved it
    # from the profile or the deck's own master, and a second read could pick a
    # different rectangle than the records were measured against.
    from pptx import Presentation

    from .modules.margin_alignment import _space_box

    try:
        profile = job.get("profile_obj") or Profile.load(job["profile"])
        space = _space_box(profile, Presentation(io.BytesIO(job["deck"])))
    except Exception:
        space = None

    new_records, reviewed = run_components(job["deck"], job["thumbs"],
                                           job["manifest"], space)
    _merge_records(job, new_records)
    frame = "against the master's frame" if space is not None \
        else "with no frame stated, so only component-to-component lines"
    return _report(f"Component review looked at {reviewed} slide"
                   f"{'s' if reviewed != 1 else ''} {frame} and added "
                   f"{len(new_records)} suggestion"
                   f"{'s' if len(new_records) != 1 else ''} "
                   "(tickable, never pre-selected).")


@app.get("/slide/{job_id}/{idx}")
def slide_preview(job_id: str, idx: int):
    job = _job(job_id)
    if job is None:
        return JSONResponse({"error": "unknown or expired job id"}, status_code=404)
    try:
        _ensure_thumbs(job_id, job)
    except RuntimeError as exc:
        return JSONResponse({"error": str(exc)}, status_code=503)
    if idx not in job["thumbs"]:
        return JSONResponse({"error": "no such slide"}, status_code=404)
    rects = _ensure_rects(job).get(idx, [])
    return JSONResponse({"png": f"/thumb/{job_id}/{idx}.png", "rects": rects})


@app.get("/thumb/{job_id}/{idx}.png")
def thumb_png(job_id: str, idx: int):
    from fastapi.responses import Response

    job = _job(job_id)
    png = ((job or {}).get("thumbs") or {}).get(idx)
    if png is None:
        return JSONResponse({"error": "no such render"}, status_code=404)
    return Response(content=png, media_type="image/png",
                    headers={"Cache-Control": "private, max-age=600"})


@app.get("/diff/{job_id}.pdf")
def diff_pdf(job_id: str):
    from fastapi.responses import Response

    from .render import build_diff
    from .report import render_diff_pdf

    job = _job(job_id)
    if job is None:
        return JSONResponse({"error": "unknown or expired job id"}, status_code=404)
    if not job.get("applied_records") or job.get("prev_deck") is None:
        return JSONResponse({"error": "no applied fixes to compare yet"},
                            status_code=404)
    if job.get("diff") is None:
        try:
            job["diff"] = build_diff(job["prev_deck"], job["cleaned"],
                                     job["applied_records"])
        except RuntimeError as exc:
            return JSONResponse({"error": str(exc)}, status_code=503)
    pdf = render_diff_pdf(job["filename"], job["diff"])
    name = Path(job["filename"]).stem
    return Response(content=pdf, media_type="application/pdf",
                    headers=_attachment(f"review-{name}.pdf"))


@app.get("/diff/{job_id}", response_class=HTMLResponse)
def diff(job_id: str):
    from .render import build_diff
    from .ui import render_diff

    job = _job(job_id)
    if job is None:
        return HTMLResponse(render_index(_pickable_profiles(), MODULES,
                                         "Unknown or expired job."), status_code=404)
    if not job.get("applied_records") or job.get("prev_deck") is None:
        return render_diff(job["filename"], job_id, None,
                           tabs=_tabs_for(job_id, job, "changes"),
                           error="No applied fixes to compare yet. Apply fixes first.")
    if job.get("diff") is None:
        try:
            # first visit renders via desktop PowerPoint; cached afterwards
            job["diff"] = build_diff(job["prev_deck"], job["cleaned"],
                                     job["applied_records"])
        except RuntimeError as exc:
            return render_diff(job["filename"], job_id, None,
                           tabs=_tabs_for(job_id, job, "changes"), error=str(exc))
    return render_diff(job["filename"], job_id, job["diff"])


@app.get("/render/{job_id}/{key}.png")
def render_png(job_id: str, key: str):
    from fastapi.responses import Response

    job = _job(job_id)
    images = (job or {}).get("diff", {}) or {}
    png = images.get("images", {}).get(key.replace("-", ":", 1))
    if png is None:
        return JSONResponse({"error": "no such render"}, status_code=404)
    return Response(content=png, media_type="image/png",
                    headers={"Cache-Control": "private, max-age=3600"})


@app.get("/annotated/{job_id}")
def annotated(job_id: str):
    from fastapi.responses import Response

    from .annotate import build_annotated
    from .store import comments_for

    job = _job(job_id)
    if job is None or job.get("deck") is None:
        return JSONResponse({"error": "deck no longer in memory; re-upload"},
                            status_code=404)
    out = build_annotated(job["deck"], job["manifest"],
                          comments_for(job["filename"]))
    name = Path(job["filename"]).stem
    return Response(
        content=out,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers=_attachment(f"{name}.annotated.pptx"))


@app.get("/download/{job_id}")
def download(job_id: str):
    from fastapi.responses import Response

    job = _job(job_id)
    if job is None or job["cleaned"] is None:
        return JSONResponse({"error": "no cleaned deck for this job"}, status_code=404)
    name = Path(job["filename"]).stem
    return Response(
        content=job["cleaned"],
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers=_attachment(f"{name}.cleaned.pptx"))


@app.get("/report/{job_id}.pdf")
def report_pdf(job_id: str):
    from fastapi.responses import Response

    from .report import render_pdf, render_visual_audit_pdf

    job = _job(job_id)
    if job is None:
        return JSONResponse({"error": "unknown or expired job id"}, status_code=404)
    content = None
    if job.get("deck") is not None:
        try:
            # visual report: rendered slides with the findings highlighted
            _ensure_thumbs(job_id, job)
            content = render_visual_audit_pdf(job["manifest"], job["thumbs"],
                                              _ensure_rects(job))
        except RuntimeError:
            content = None  # rendering unavailable: text-only fallback
    if content is None:
        content = render_pdf(job["manifest"])
    return Response(content=content, media_type="application/pdf",
                    headers=_attachment(f"audit-{Path(job['filename']).stem}.pdf"))


@app.get("/report/{job_id}.csv")
def report_csv(job_id: str):
    from fastapi.responses import Response

    from .report import render_csv

    job = _job(job_id)
    if job is None:
        return JSONResponse({"error": "unknown or expired job id"}, status_code=404)
    return Response(content=render_csv(job["manifest"]), media_type="text/csv",
                    headers=_attachment(f"audit-{Path(job['filename']).stem}.csv"))


@app.get("/manifest/{job_id}")
def manifest(job_id: str):
    # single atomic lookup: eviction by a concurrent /audit must yield a
    # clean 404, not a KeyError 500 (review finding)
    job = _jobs.get(job_id)
    if job is None:
        return JSONResponse({"error": "unknown or expired job id"}, status_code=404)
    return JSONResponse(job["manifest"])


def main():
    import argparse
    import socket

    import uvicorn

    from . import auth, web

    ap = argparse.ArgumentParser(description="Prezlab PPT QC web app.")
    ap.add_argument("--host", default=None,
                    help="Bind address (default 127.0.0.1, or 0.0.0.0 with --lan).")
    ap.add_argument("--port", type=int, default=8000)
    ap.add_argument("--lan", action="store_true",
                    help="LAN pilot deployment: bind all interfaces, REQUIRE "
                         "sign-in on every route, and disable the legacy "
                         "name-cookie identity.")
    ap.add_argument("--reload", action="store_true",
                    help="Restart on source changes. Use while developing: "
                         "without it a running server keeps serving the code "
                         "it was started with, and edits appear to do nothing.")
    args = ap.parse_args()

    host = args.host or ("0.0.0.0" if args.lan else "127.0.0.1")
    if args.lan:
        app.state.auth_required = True  # single identity across -m double-import
        web.AUTH_REQUIRED = True
        auth.STRICT_SESSIONS = True
        try:
            probe = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
            probe.connect(("10.255.255.255", 1))
            ip = probe.getsockname()[0]
            probe.close()
        except OSError:
            ip = "this-machine"
        print(f"LAN mode: sign-in enforced on every route.")
        print(f"Team URL: http://{ip}:{args.port}")
    elif host != "127.0.0.1":
        print("WARNING: binding beyond loopback WITHOUT --lan leaves routes "
              "open to anonymous use. Use --lan for network exposure.")
    # Printed unconditionally: log_level="warning" suppresses uvicorn's own
    # startup line, so without this the server looks hung, and there is no
    # visible marker of WHEN it started or whether it predates a code change.
    print(f"Serving http://{'localhost' if host == '127.0.0.1' else host}:"
          f"{args.port}   (Ctrl+C to stop)")
    if args.reload:
        # Reload needs an import string, not the app object, so uvicorn can
        # re-import the module in a fresh worker.
        uvicorn.run("qc.web:app", host=host, port=args.port,
                    log_level="warning", reload=True)
    else:
        print("Code changes need a restart; start with --reload to pick them "
              "up automatically.")
        uvicorn.run(app, host=host, port=args.port, log_level="warning")


if __name__ == "__main__":
    main()
