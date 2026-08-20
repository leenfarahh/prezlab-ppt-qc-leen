"""Generate the synthetic spike corpus into fixtures/.

Each deck plants known conditions and writes a .truth.json next to it so the
experiments can be scored. The synthetic corpus proves mechanics; PRD 11.1
still requires real Prezlab decks (incl. SmartArt, which python-pptx cannot
create) before spike sign-off.
"""

import base64
import io
import json
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Emu, Inches, Pt

from .ns import find, qn

FIXTURES = Path(__file__).resolve().parent.parent / "fixtures"

_PNG_1PX = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


def _new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _set_master_style(master, style_tag: str, sz: int, latin: str):
    """Plant a known value in master p:txStyles/<style>/a:lvl1pPr/a:defRPr."""
    txStyles = find(master.element, "p:txStyles")
    style = find(txStyles, f"p:{style_tag}")
    lvl = find(style, "a:lvl1pPr")
    if lvl is None:
        lvl = style.makeelement(qn("a:lvl1pPr"), {})
        style.insert(0, lvl)
    defRPr = find(lvl, "a:defRPr")
    if defRPr is None:
        defRPr = lvl.makeelement(qn("a:defRPr"), {})
        lvl.append(defRPr)
    defRPr.set("sz", str(sz))
    latin_el = find(defRPr, "a:latin")
    if latin_el is None:
        latin_el = defRPr.makeelement(qn("a:latin"), {})
        defRPr.append(latin_el)
    latin_el.set("typeface", latin)


def _set_shape_lststyle_size(shape, sz: int):
    """Plant a size in the shape's own txBody lstStyle (cascade layer 3)."""
    txBody = shape.text_frame._txBody
    lst = find(txBody, "a:lstStyle")
    if lst is None:
        lst = txBody.makeelement(qn("a:lstStyle"), {})
        # schema order: bodyPr, lstStyle, p*
        txBody.insert(1, lst)
    lvl = lst.makeelement(qn("a:lvl1pPr"), {})
    lst.append(lvl)
    defRPr = lvl.makeelement(qn("a:defRPr"), {})
    defRPr.set("sz", str(sz))
    lvl.append(defRPr)


def _set_run_theme_font_ref(run, ref: str):
    """Point the run's a:latin at a theme font reference like +mn-lt."""
    rPr = run._r.get_or_add_rPr()
    latin = find(rPr, "a:latin")
    if latin is None:
        latin = rPr.makeelement(qn("a:latin"), {})
        rPr.append(latin)
    latin.set("typeface", ref)


# --- decks ------------------------------------------------------------------


def make_clean(path: Path) -> dict:
    prs = _new_prs()
    master = prs.slide_masters[0]
    _set_master_style(master, "titleStyle", 4000, "Georgia")
    _set_master_style(master, "bodyStyle", 1800, "Trebuchet MS")

    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Formatting QC Spike"
    s1.placeholders[1].text = "Synthetic corpus baseline"

    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Inherited title"
    body = s2.placeholders[1]
    body.text_frame.text = "Inherited body run"
    p2 = body.text_frame.add_paragraph()
    r_explicit = p2.add_run()
    r_explicit.text = "Explicit 20pt run"
    r_explicit.font.size = Pt(20)
    p3 = body.text_frame.add_paragraph()
    r_theme = p3.add_run()
    r_theme.text = "Theme minor font run"
    _set_run_theme_font_ref(r_theme, "+mn-lt")

    tb_styled = s2.shapes.add_textbox(Inches(8), Inches(1), Inches(4), Inches(1))
    _set_shape_lststyle_size(tb_styled, 1500)
    tb_styled.text_frame.text = "lstStyle 15pt textbox"

    tb_bare = s2.shapes.add_textbox(Inches(8), Inches(2.5), Inches(4), Inches(1))
    tb_bare.text_frame.text = "Bare textbox (cascade tail)"

    prs.save(path)
    return {
        "planted": {
            "master_title": {"size_pt": 40.0, "family": "Georgia"},
            "master_body": {"size_pt": 18.0, "family": "Trebuchet MS"},
            "explicit_run": {"slide": 1, "text": "Explicit 20pt run", "size_pt": 20.0},
            "theme_ref_run": {"slide": 1, "text": "Theme minor font run", "ref": "+mn-lt"},
            "lststyle_textbox": {"slide": 1, "text": "lstStyle 15pt textbox", "size_pt": 15.0},
            "bare_textbox": {"slide": 1, "text": "Bare textbox (cascade tail)"},
        }
    }


def make_bilingual(path: Path) -> dict:
    prs = _new_prs()
    truth_hits = []

    s = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    s.shapes.title.text = "Bilingual fixture"

    tb = s.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf = tb.text_frame
    # para 0: mixed EN + AR runs, no rtl attribute (unicode-only detection)
    tf.text = "Quarterly Review "
    r_ar = tf.paragraphs[0].add_run()
    r_ar.text = "مراجعة ربع سنوية"
    r_ar.font.name = "Inter"  # latin only; cs deliberately missing
    truth_hits.append({"slide_index": 0, "paragraph_index": 0, "run_index": 1, "reason": "unicode"})

    # para 1: full AR with rtl='1' and a cs typeface set
    p1 = tf.add_paragraph()
    r_full = p1.add_run()
    r_full.text = "الملخص التنفيذي"
    p1._p.get_or_add_pPr().set("rtl", "1")
    from .arabic import set_cs_typeface

    set_cs_typeface(r_full, "Dubai")
    truth_hits.append({"slide_index": 0, "paragraph_index": 1, "run_index": 0, "reason": "unicode+rtl_attr"})

    # para 2: rtl attribute with Latin-only text (rtl_attr-only hit)
    p2 = tf.add_paragraph()
    r_lat = p2.add_run()
    r_lat.text = "Latin text in an RTL paragraph"
    p2._p.get_or_add_pPr().set("rtl", "1")
    truth_hits.append({"slide_index": 0, "paragraph_index": 2, "run_index": None, "reason": "rtl_attr"})

    # Arabic inside a group (group traversal is a known miss risk)
    grouped = {"attempted": True, "created": False}
    try:
        grp = s.shapes.add_group_shape()
        gtb = grp.shapes.add_textbox(Inches(9.5), Inches(2), Inches(3), Inches(1))
        gtb.text_frame.text = "تقرير"
        grouped["created"] = True
        truth_hits.append({"slide_index": 0, "in_group": True, "reason": "unicode"})
    except Exception as exc:  # keep fixture usable even if API differs
        grouped["error"] = str(exc)

    prs.save(path)
    return {
        "expected_hits": truth_hits,
        "expected_hit_count": len(truth_hits),
        "cs_missing_run": {"paragraph_index": 0, "run_index": 1},
        "cs_set_run": {"paragraph_index": 1, "run_index": 0, "typeface": "Dubai"},
        "group": grouped,
    }


def make_theme_colors(path: Path) -> dict:
    prs = _new_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    from pptx.enum.shapes import MSO_SHAPE

    def rect(x_in):
        return s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_in), Inches(1), Inches(2), Inches(1))

    sh_theme = rect(0.5)
    sh_theme.fill.solid()
    sh_theme.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1

    sh_bright = rect(3.0)
    sh_bright.fill.solid()
    sh_bright.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    sh_bright.fill.fore_color.brightness = 0.4  # writes lumMod/lumOff

    sh_tint = rect(5.5)
    sh_tint.fill.solid()
    sh_tint.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)
    srgb = find(find(sh_tint.fill._xPr, "a:solidFill"), "a:srgbClr")
    tint = srgb.makeelement(qn("a:tint"), {})
    tint.set("val", "50000")
    srgb.append(tint)

    sh_lit_on = rect(8.0)  # literal hex, exactly on an (arbitrary) palette color
    sh_lit_on.fill.solid()
    sh_lit_on.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    sh_lit_off = rect(10.5)  # near-miss off-palette literal
    sh_lit_off.fill.solid()
    sh_lit_off.fill.fore_color.rgb = RGBColor(0x21, 0x4F, 0x7A)

    sh_grad = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3), Inches(2), Inches(1))
    sh_grad.fill.gradient()  # no single color: resolver must return None

    prs.save(path)
    return {
        "shapes": {
            "theme_plain": {"shape_id": sh_theme.shape_id, "scheme": "accent1", "transform": None},
            "theme_bright": {"shape_id": sh_bright.shape_id, "scheme": "accent1", "transform": "lum+0.4"},
            "literal_tint": {"shape_id": sh_tint.shape_id, "base_hex": "FF0000", "tint": 0.5},
            "literal_on": {"shape_id": sh_lit_on.shape_id, "hex": "1F4E79"},
            "literal_off_near": {"shape_id": sh_lit_off.shape_id, "hex": "214F7A", "near": "1F4E79"},
            "gradient": {"shape_id": sh_grad.shape_id, "expect": None},
        },
        "palette": {"prezlab_navy": "1F4E79", "prezlab_red": "C00000", "white": "FFFFFF"},
    }


def make_mixed_layouts(path: Path) -> dict:
    prs = _new_prs()
    for i in range(7):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = f"Content slide {i + 1}"
    for i in range(2):
        s = prs.slides.add_slide(prs.slide_layouts[0])
        s.shapes.title.text = f"Divider {i + 1}"
    prs.slides.add_slide(prs.slide_layouts[6])  # blank outlier

    # Displace one content slide's title placeholder (explicit xfrm override)
    moved = prs.slides[2].shapes.title
    orig = {"left": moved.left, "top": moved.top}
    moved.left = moved.left + Inches(2)
    moved.top = moved.top + Inches(1)

    prs.save(path)
    return {
        "dominant_layout": prs.slide_layouts[1].name,
        "outlier_slide_indices": [7, 8, 9],
        "moved_placeholder": {"slide_index": 2, "ph_idx": moved.placeholder_format.idx,
                              "original_emu": orig},
    }


def make_heavy(path: Path) -> dict:
    """Chart (embedded xlsx part), picture, and notes: round-trip surface for U2.
    SmartArt cannot be generated by python-pptx; real decks must cover it."""
    prs = _new_prs()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Heavy assets"
    data = CategoryChartData()
    data.categories = ["Q1", "Q2", "Q3"]
    data.add_series("Revenue", (1.2, 2.3, 3.1))
    s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
                       Inches(5), Inches(4), data)
    s.shapes.add_picture(io.BytesIO(_PNG_1PX), Inches(7), Inches(1.5), Inches(2), Inches(2))
    s.notes_slide.notes_text_frame.text = "Speaker note for round-trip check."
    prs.save(path)
    return {"parts_expected": ["chart", "embedded xlsx", "png media", "notes slide"]}


def make_large(path: Path, n_slides: int = 200) -> dict:
    prs = _new_prs()
    from pptx.enum.shapes import MSO_SHAPE

    for i in range(n_slides):
        s = prs.slides.add_slide(prs.slide_layouts[1 if i % 3 else 5])
        if s.shapes.title is not None:
            s.shapes.title.text = f"Slide {i + 1}: section heading"
        tb = s.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
        tf = tb.text_frame
        tf.text = f"Body copy line one for slide {i + 1}."
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = "Second paragraph with an explicit size."
        r.font.size = Pt(14)
        sh1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(2), Inches(2), Inches(1))
        sh1.fill.solid()
        sh1.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2
        sh2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(2), Inches(1.5), Inches(1))
        sh2.fill.solid()
        sh2.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    prs.save(path)
    return {"slide_count": n_slides}


def main():
    FIXTURES.mkdir(exist_ok=True)
    jobs = {
        "clean": make_clean,
        "bilingual_ar": make_bilingual,
        "theme_colors": make_theme_colors,
        "mixed_layouts": make_mixed_layouts,
        "heavy": make_heavy,
        "large_200": make_large,
    }
    for name, fn in jobs.items():
        pptx_path = FIXTURES / f"{name}.pptx"
        truth = fn(pptx_path)
        (FIXTURES / f"{name}.truth.json").write_text(
            json.dumps(truth, indent=2), encoding="utf-8")
        print(f"built {pptx_path.name}  ({pptx_path.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
