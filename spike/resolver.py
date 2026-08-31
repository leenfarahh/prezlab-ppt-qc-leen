"""U4 (part 1): effective font resolution.

python-pptx returns None for inherited font properties (run.font.size etc.),
so every audit module needs the effective value resolved through the OOXML
cascade (PRD Section 6.5, verified finding):

    run a:rPr
      -> paragraph a:pPr/a:defRPr
      -> shape txBody a:lstStyle (per paragraph level)
      -> layout placeholder lstStyle (matched by ph idx)
      -> master placeholder lstStyle (matched by ph type category)
      -> master p:txStyles (titleStyle / bodyStyle / otherStyle)
      -> presentation p:defaultTextStyle
      -> hard default (18 pt, theme minor latin)

Theme font references (+mj-lt / +mn-lt / +mj-cs / +mn-cs) resolve against the
master's theme part fontScheme.
"""

from dataclasses import dataclass

from lxml import etree
from pptx.enum.shapes import PP_PLACEHOLDER

from .ns import find, qn

_TITLE_TYPES = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)


@dataclass
class Resolved:
    value: object
    source: str  # which cascade level supplied the value


@dataclass
class EffectiveFont:
    family: Resolved
    cs_family: Resolved | None
    size_pt: Resolved
    bold: Resolved


# --- theme fonts -----------------------------------------------------------


def _theme_element(master):
    """The a:theme root related to a slide master (theme1.xml)."""
    for rel in master.part.rels.values():
        if rel.reltype.endswith("/theme"):
            part = rel.target_part
            el = getattr(part, "_element", None)
            return el if el is not None else etree.fromstring(part.blob)
    return None


# A master's theme fonts do not change while a deck is being read, and
# resolve_run asks for them ONCE PER RUN - 3,432 times on a 26-slide deck, each
# one walking the master part's relationships to find theme1.xml (0.4s of an
# 8.1s design scan, measured 24/08/2026). Memoized on the master's part, with
# the part compared by identity on read so a recycled id cannot answer for a
# different master (the same rule as qc.design._MARKER_MEMO).
_THEME_FONT_MEMO: dict[int, tuple] = {}


def theme_fonts(master) -> dict:
    """{'+mj-lt': 'Playfair Display', '+mn-lt': 'Inter', '+mj-cs': ..., '+mn-cs': ...}

    East Asian faces (+mj-ea / +mn-ea) are read too. Prezlab decks do not use
    them, but the theme declares all three scripts and the extracted style
    spec reports what the theme actually says rather than a subset of it."""
    part = getattr(master, "part", None)
    hit = _THEME_FONT_MEMO.get(id(part)) if part is not None else None
    if hit is not None and hit[0] is part:
        return hit[1]
    theme = _theme_element(master)
    out = {}
    scheme = find(theme, ".//a:fontScheme")
    for tag, prefix in (("a:majorFont", "+mj"), ("a:minorFont", "+mn")):
        group = find(scheme, tag)
        if group is None:
            continue
        for child, suffix in (("a:latin", "lt"), ("a:ea", "ea"), ("a:cs", "cs")):
            el = find(group, child)
            if el is not None and el.get("typeface"):
                out[f"{prefix}-{suffix}"] = el.get("typeface")
    if part is not None:
        _THEME_FONT_MEMO[id(part)] = (part, out)
    return out


# --- cascade walk ----------------------------------------------------------


def _props_from_rpr(rpr) -> dict:
    """Extract family / cs_family / size / bold from an rPr-shaped element
    (a:rPr or a:defRPr). Absent attributes stay absent, never defaulted."""
    if rpr is None:
        return {}
    out = {}
    latin = find(rpr, "a:latin")
    if latin is not None and latin.get("typeface"):
        out["family"] = latin.get("typeface")
    cs = find(rpr, "a:cs")
    if cs is not None and cs.get("typeface"):
        out["cs_family"] = cs.get("typeface")
    if rpr.get("sz"):
        out["size_pt"] = int(rpr.get("sz")) / 100.0
    if rpr.get("b") is not None:
        out["bold"] = rpr.get("b") in ("1", "true")
    return out


def _lstStyle_defRPr(txBody, level: int):
    """defRPr from a txBody's lstStyle at the given paragraph level (0-based)."""
    lst = find(txBody, "a:lstStyle")
    if lst is None:
        return None
    return find(find(lst, f"a:lvl{level + 1}pPr"), "a:defRPr")


# Both lookups below are a LINEAR SCAN of a placeholder collection, and
# iterating one runs an lxml xpath per shape it passes (CT_Shape.ph). rpr_layers
# does both on EVERY RUN - so a deck's fonts and colours are resolved by
# re-deriving the same two small maps thousands of times.
#
# Measured on fixtures/large_200.pptx (30/08/2026): after the equivalent fix in
# qc.design._dimensions, these were still 70,423 xpath calls and about a third
# of a design scan. The layout and the master do not gain or lose placeholders
# while a deck is being read, so the maps are built once each.
#
# Memoized on the container, identity-checked on read - the same rule as
# _THEME_FONT_MEMO above: CPython reuses id() once an object is collected, and a
# bare id key would hand one layout's placeholders to another.
_LAYOUT_PH_MEMO: dict[int, tuple] = {}
_MASTER_PH_MEMO: dict[int, tuple] = {}


def layout_ph_index(layout) -> dict:
    """{ph idx: placeholder} for one layout.

    Public because qc.design resolves inherited GEOMETRY through the same map
    and a second copy of it would be a second thing to keep in step."""
    hit = _LAYOUT_PH_MEMO.get(id(layout))
    if hit is not None and hit[0] is layout:
        return hit[1]
    index = {}
    for ph in layout.placeholders:
        try:
            index[ph.element.ph_idx] = ph
        except (ValueError, AttributeError):
            continue
    _LAYOUT_PH_MEMO[id(layout)] = (layout, index)
    return index


def _master_ph_pair(master) -> tuple:
    """(title placeholder, body placeholder) for one master, first of each.

    The same two the scan below would have found, in the same order - the first
    title-ish placeholder and the first BODY one."""
    hit = _MASTER_PH_MEMO.get(id(master))
    if hit is not None and hit[0] is master:
        return hit[1]
    title = body = None
    for ph in master.placeholders:
        try:
            kind = ph.placeholder_format.type
        except (ValueError, AttributeError):
            continue
        if title is None and kind in _TITLE_TYPES:
            title = ph
        if body is None and kind == PP_PLACEHOLDER.BODY:
            body = ph
        if title is not None and body is not None:
            break
    _MASTER_PH_MEMO[id(master)] = (master, (title, body))
    return title, body


def _layout_placeholder(shape, slide):
    if not shape.is_placeholder:
        return None
    return layout_ph_index(slide.slide_layout).get(
        shape.placeholder_format.idx)


def _master_placeholder(shape, master):
    if not shape.is_placeholder:
        return None
    title, body = _master_ph_pair(master)
    return title if shape.placeholder_format.type in _TITLE_TYPES else body


def _master_txstyle_defRPr(shape, master, level: int):
    """p:txStyles on the master: titleStyle for title placeholders, bodyStyle
    for body placeholders, otherStyle for everything else (incl. free shapes)."""
    txStyles = find(master.element, "p:txStyles")
    if txStyles is None:
        return None
    if shape.is_placeholder and shape.placeholder_format.type in _TITLE_TYPES:
        style = find(txStyles, "p:titleStyle")
    elif shape.is_placeholder:
        style = find(txStyles, "p:bodyStyle")
    else:
        style = find(txStyles, "p:otherStyle")
    return find(find(style, f"a:lvl{level + 1}pPr"), "a:defRPr")


def _presentation_default_defRPr(prs, level: int):
    pres_el = getattr(prs.part, "_element", None)
    dts = find(pres_el, "p:defaultTextStyle")
    return find(find(dts, f"a:lvl{level + 1}pPr"), "a:defRPr")


def rpr_layers(run, paragraph, shape, slide, prs) -> list[tuple[str, object]]:
    """(source name, rPr-shaped element or None) for every level of the cascade,
    highest priority first.

    Split out because two different questions walk the same ladder: which font
    and size a run ends up with (resolve_run, below) and what COLOUR it ends up
    with (qc.design, for the contrast check). A second copy of this walk would
    drift from this one, and then the tool's idea of which run is 12pt would not
    match its idea of what colour that same run is - two answers about one run,
    which is worse than either being wrong on its own.
    """
    master = slide.slide_layout.slide_master
    level = paragraph.level
    layers: list[tuple[str, object]] = [
        ("run.rPr", find(run._r, "a:rPr")),
        ("paragraph.defRPr", find(find(paragraph._p, "a:pPr"), "a:defRPr")),
        ("shape.lstStyle", _lstStyle_defRPr(shape.text_frame._txBody, level)),
    ]
    lp = _layout_placeholder(shape, slide)
    if lp is not None and lp.has_text_frame:
        layers.append(("layout.placeholder",
                       _lstStyle_defRPr(lp.text_frame._txBody, level)))
    mp = _master_placeholder(shape, master)
    if mp is not None and mp.has_text_frame:
        layers.append(("master.placeholder",
                       _lstStyle_defRPr(mp.text_frame._txBody, level)))
    layers.append(("master.txStyles",
                   _master_txstyle_defRPr(shape, master, level)))
    layers.append(("presentation.default",
                   _presentation_default_defRPr(prs, level)))
    return layers


def resolve_run(run, paragraph, shape, slide, prs) -> EffectiveFont:
    """Walk the cascade; first level that defines a property wins."""
    master = slide.slide_layout.slide_master

    resolved: dict[str, Resolved] = {}
    for source, rpr in rpr_layers(run, paragraph, shape, slide, prs):
        for key, val in _props_from_rpr(rpr).items():
            if key not in resolved:
                resolved[key] = Resolved(val, source)

    fonts = theme_fonts(master)
    for key in ("family", "cs_family"):
        r = resolved.get(key)
        if r is not None and isinstance(r.value, str) and r.value.startswith("+"):
            mapped = fonts.get(r.value)
            if mapped:
                resolved[key] = Resolved(mapped, f"{r.source} -> theme({r.value})")

    return EffectiveFont(
        family=resolved.get("family", Resolved(fonts.get("+mn-lt", "Calibri"), "hard-default")),
        cs_family=resolved.get("cs_family"),
        size_pt=resolved.get("size_pt", Resolved(18.0, "hard-default")),
        bold=resolved.get("bold", Resolved(False, "hard-default")),
    )
