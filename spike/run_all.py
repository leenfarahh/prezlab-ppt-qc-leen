"""Run the U1-U4 spike experiments against the synthetic corpus and print a
findings memo. Also writes out/spike-memo.md.

U5 (Microsoft Graph rendering + data-residency ruling) is external and is
reported as BLOCKED-ON-INPUTS here.
"""

import json
from pathlib import Path

from pptx import Presentation

from . import u1_master, u2_roundtrip, u3_perf
from .arabic import cs_typeface, scan_presentation
from .color_resolver import ciede2000, color_scheme, nearest_palette_match, resolve_solid_fill
from .ns import find
from .resolver import resolve_run, theme_fonts

ROOT = Path(__file__).resolve().parent.parent
FIXTURES = ROOT / "fixtures"
OUT = ROOT / "out"


def _truth(name: str) -> dict:
    return json.loads((FIXTURES / f"{name}.truth.json").read_text(encoding="utf-8"))


def _find_run(prs, slide_idx: int, text: str):
    slide = prs.slides[slide_idx]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text == text:
                    return run, para, shape, slide
    raise LookupError(text)


def check_u4_fonts() -> dict:
    prs = Presentation(FIXTURES / "clean.pptx")
    truth = _truth("clean")["planted"]
    checks = {}

    run, para, shape, slide = _find_run(prs, 1, truth["explicit_run"]["text"])
    eff = resolve_run(run, para, shape, slide, prs)
    checks["explicit_run_size"] = {
        "expected": truth["explicit_run"]["size_pt"], "got": eff.size_pt.value,
        "source": eff.size_pt.source,
        "pass": eff.size_pt.value == truth["explicit_run"]["size_pt"] and eff.size_pt.source == "run.rPr",
    }

    run, para, shape, slide = _find_run(prs, 1, truth["theme_ref_run"]["text"])
    eff = resolve_run(run, para, shape, slide, prs)
    minor = theme_fonts(slide.slide_layout.slide_master).get("+mn-lt")
    checks["theme_ref_resolution"] = {
        "expected": minor, "got": eff.family.value, "source": eff.family.source,
        "pass": eff.family.value == minor and "theme(+mn-lt)" in eff.family.source,
    }

    run, para, shape, slide = _find_run(prs, 1, truth["lststyle_textbox"]["text"])
    eff = resolve_run(run, para, shape, slide, prs)
    checks["shape_lststyle_size"] = {
        "expected": truth["lststyle_textbox"]["size_pt"], "got": eff.size_pt.value,
        "source": eff.size_pt.source,
        "pass": eff.size_pt.value == truth["lststyle_textbox"]["size_pt"]
                and eff.size_pt.source == "shape.lstStyle",
    }

    run, para, shape, slide = _find_run(prs, 1, truth["bare_textbox"]["text"])
    eff = resolve_run(run, para, shape, slide, prs)
    checks["bare_textbox_cascade_tail"] = {
        "got": eff.size_pt.value, "source": eff.size_pt.source,
        "pass": eff.size_pt.source != "hard-default",  # cascade found a real definition
        "note": "value comes from master otherStyle / presentation default; hand-check vs PowerPoint",
    }

    # Inherited title: planted master titleStyle 40pt Georgia. Template layers
    # may legitimately override; report source so the memo can say which won.
    run, para, shape, slide = _find_run(prs, 1, "Inherited title")
    eff = resolve_run(run, para, shape, slide, prs)
    checks["inherited_title"] = {
        "planted_master_value": truth["master_title"]["size_pt"],
        "got": eff.size_pt.value, "source": eff.size_pt.source,
        "family_got": eff.family.value, "family_source": eff.family.source,
        "pass": eff.size_pt.source != "hard-default",
        "note": "pass = resolved from a real cascade layer (no None-leak); exact layer reported",
    }

    ok = all(c["pass"] for c in checks.values())
    return {"pass": ok, "checks": checks}


def check_u4_arabic() -> dict:
    prs = Presentation(FIXTURES / "bilingual_ar.pptx")
    truth = _truth("bilingual_ar")
    hits = scan_presentation(prs)
    found = len(hits)
    expected = truth["expected_hit_count"]

    # cs typeface assertions
    tf = None
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame and "Quarterly Review" in shape.text_frame.text:
            tf = shape.text_frame
    cs_missing = cs_typeface(tf.paragraphs[0].runs[1]) is None
    cs_set = cs_typeface(tf.paragraphs[1].runs[0]) == truth["cs_set_run"]["typeface"]

    group_covered = (not truth["group"]["created"]) or any(
        h.reason == "unicode" and h.sample == "تقرير" for h in hits)

    return {
        "pass": found >= expected and cs_missing and cs_set and group_covered,
        "expected_hits": expected, "found_hits": found,
        "false_negatives": max(0, expected - found),
        "cs_missing_detected": cs_missing,
        "cs_set_detected": cs_set,
        "group_traversal_covered": group_covered,
        "hits": [vars(h) for h in hits],
    }


def check_u4_color() -> dict:
    prs = Presentation(FIXTURES / "theme_colors.pptx")
    truth = _truth("theme_colors")
    slide = prs.slides[0]
    master = slide.slide_layout.slide_master
    scheme = color_scheme(master)
    by_id = {sh.shape_id: sh for sh in slide.shapes}
    t = truth["shapes"]

    def rgb_of(shape):
        return resolve_solid_fill(find(shape._element, "p:spPr"), master)

    plain = rgb_of(by_id[t["theme_plain"]["shape_id"]])
    bright = rgb_of(by_id[t["theme_bright"]["shape_id"]])
    tinted = rgb_of(by_id[t["literal_tint"]["shape_id"]])
    lit_on = rgb_of(by_id[t["literal_on"]["shape_id"]])
    grad = rgb_of(by_id[t["gradient"]["shape_id"]])

    def lum(rgb):
        return sum(rgb)

    palette = {k: tuple(int(v[i:i + 2], 16) for i in (0, 2, 4))
               for k, v in truth["palette"].items()}
    name_on, de_on = nearest_palette_match(lit_on, palette)
    lit_off = rgb_of(by_id[t["literal_off_near"]["shape_id"]])
    name_off, de_off = nearest_palette_match(lit_off, palette)

    checks = {
        "schemeClr_resolves_to_theme_rgb": {"got": plain, "expected": scheme.get("accent1"),
                                            "pass": plain == scheme.get("accent1")},
        "lumMod_lumOff_applied": {"base": plain, "got": bright,
                                  "pass": bright is not None and bright != plain and lum(bright) > lum(plain),
                                  "note": "lighter than base as planted (+0.4 brightness); exact value pending PowerPoint hand-check"},
        "tint_applied_to_literal": {"base": "FF0000", "got": tinted,
                                    "pass": tinted is not None and tinted != (255, 0, 0) and lum(tinted) > lum((255, 0, 0))},
        "gradient_returns_none": {"got": grad, "pass": grad is None},
        "ciede2000_exact_match_zero": {"deltaE": round(de_on, 4), "matched": name_on,
                                       "pass": name_on == "prezlab_navy" and de_on < 1e-9},
        "ciede2000_near_miss_small": {"deltaE": round(de_off, 2), "matched": name_off,
                                      "pass": name_off == "prezlab_navy" and 0 < de_off < 5},
    }
    return {"pass": all(c["pass"] for c in checks.values()), "checks": checks}


def main():
    OUT.mkdir(exist_ok=True)
    results = {}

    results["U1_master"] = u1_master.run(FIXTURES / "mixed_layouts.pptx", OUT)
    truth1 = _truth("mixed_layouts")
    u1 = results["U1_master"]
    u1_pass = (
        u1["census"]["dominant"] == truth1["dominant_layout"]
        and u1["census"]["outlier_slide_indices"] == truth1["outlier_slide_indices"]
        and u1["deviations_found"] == 1
        and u1["deviations_detail"][0]["slide_index"] == truth1["moved_placeholder"]["slide_index"]
        and u1["residual_after_enforce"] == 0
    )
    u1["pass"] = u1_pass

    results["U2_roundtrip_heavy"] = u2_roundtrip.run(FIXTURES / "heavy.pptx", OUT)
    results["U2_roundtrip_clean"] = u2_roundtrip.run(FIXTURES / "clean.pptx", OUT)
    for key in ("U2_roundtrip_heavy", "U2_roundtrip_clean"):
        r = results[key]
        r["pass"] = (r["roundtrip"]["binary_parts_preserved"]
                     and not r["roundtrip"]["removed"]
                     and r["outputs_parse_ok"]
                     and r["surgical_edit"]["edit_applied"])

    results["U3_perf"] = u3_perf.run(FIXTURES / "large_200.pptx", OUT)
    results["U3_perf"]["pass"] = results["U3_perf"]["seconds"]["total"] < 120.0

    results["U4_fonts"] = check_u4_fonts()
    results["U4_arabic"] = check_u4_arabic()
    results["U4_color"] = check_u4_color()

    results["U5_graph_renderer"] = {
        "pass": None,
        "status": "BLOCKED-ON-INPUTS",
        "needs": ["M365 tenant access for Graph driveItem ?format=pdf test",
                  "Operations/IT-security ruling on decks transiting OneDrive/SharePoint (KSA/UAE residency)"],
    }

    (OUT / "spike-results.json").write_text(json.dumps(results, indent=2, default=str),
                                            encoding="utf-8")

    lines = ["# Spike findings memo (synthetic corpus)", ""]
    for key, res in results.items():
        status = ("PASS" if res.get("pass") else
                  "BLOCKED" if res.get("pass") is None else "FAIL")
        lines.append(f"## {key}: {status}")
        if key == "U3_perf":
            lines.append(f"- {res['slides']} slides, {res['runs_resolved']} runs resolved, "
                         f"total {res['seconds']['total']}s ({res['ms_per_slide']} ms/slide)")
        lines.append("")
    memo = "\n".join(lines)
    (OUT / "spike-memo.md").write_text(memo, encoding="utf-8")

    print(json.dumps({k: (v.get("pass") if isinstance(v, dict) else v)
                      for k, v in results.items()}, indent=2))
    print(f"\nFull results: {OUT / 'spike-results.json'}")
    return results


if __name__ == "__main__":
    main()
