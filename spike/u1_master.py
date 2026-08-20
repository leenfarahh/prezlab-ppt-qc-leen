"""U1: master/layout enforcement without corruption.

Scope per PRD (verified finding): python-pptx cannot re-apply a layout to an
existing populated slide, so v1 enforcement means (a) layout census: dominant
layout + outlier flagging, and (b) the deterministic geometry nudge: a slide
placeholder whose explicit a:xfrm overrides its layout placeholder gets the
override REMOVED, so geometry reverts to inherited layout values (equivalent
to PowerPoint's Reset Layout for position/size).
"""

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation

from .ns import find, qn


@dataclass
class GeometryDeviation:
    slide_index: int
    ph_idx: int
    slide_xfrm: dict
    baseline_xfrm: dict


def layout_census(prs) -> dict:
    counts: dict[str, list[int]] = {}
    for i, slide in enumerate(prs.slides):
        counts.setdefault(slide.slide_layout.name, []).append(i)
    dominant = max(counts, key=lambda k: len(counts[k]))
    outliers = sorted(i for name, idxs in counts.items() if name != dominant for i in idxs)
    return {"counts": {k: len(v) for k, v in counts.items()},
            "dominant": dominant, "outlier_slide_indices": outliers}


def _xfrm_dict(sp_el) -> dict | None:
    """Explicit geometry as {'x','y','cx','cy'} with absent components as
    None. Spike finding: an xfrm may carry a:off WITHOUT a:ext (python-pptx
    writes exactly that when only .left/.top are set), so geometry must be
    read and compared per-component, never as an all-or-nothing tuple."""
    spPr = find(sp_el, "p:spPr")
    if spPr is None:
        spPr = find(sp_el, "a:spPr")
    xfrm = find(spPr, "a:xfrm")
    if xfrm is None:
        return None
    off, ext = find(xfrm, "a:off"), find(xfrm, "a:ext")
    return {
        "x": int(off.get("x")) if off is not None else None,
        "y": int(off.get("y")) if off is not None else None,
        "cx": int(ext.get("cx")) if ext is not None else None,
        "cy": int(ext.get("cy")) if ext is not None else None,
    }


def _effective_baseline_xfrm(slide, ph) -> dict:
    """The geometry a placeholder would inherit with no explicit xfrm,
    resolved per-component through layout ph (same idx) then master ph (same
    type category). Spike finding: comparing against the layout alone misses
    every placeholder whose layout itself inherits from the master."""
    from pptx.enum.shapes import PP_PLACEHOLDER

    title_types = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    idx = ph.placeholder_format.idx

    chain = []
    lp = None
    for cand in slide.slide_layout.placeholders:
        if cand.placeholder_format.idx == idx:
            lp = cand
            break
    if lp is not None:
        chain.append(_xfrm_dict(lp._element))

    ref = lp if lp is not None else ph
    ref_type = ref.placeholder_format.type
    want_title = ref_type in title_types
    for mp in slide.slide_layout.slide_master.placeholders:
        mp_type = mp.placeholder_format.type
        if mp_type == ref_type or (want_title and mp_type in title_types):
            chain.append(_xfrm_dict(mp._element))
            break

    baseline = {"x": None, "y": None, "cx": None, "cy": None}
    for level in chain:
        if level is None:
            continue
        for key in baseline:
            if baseline[key] is None and level[key] is not None:
                baseline[key] = level[key]
    return baseline


def find_geometry_deviations(prs, tolerance_emu: int = 9525) -> list[GeometryDeviation]:
    """Slide placeholders whose explicit xfrm differs from their inherited
    baseline by more than tolerance (default 9525 EMU = 1px at 96dpi, the
    PRD's geometry_tolerance_emu default). Only components the slide
    actually overrides are compared."""
    out = []
    for s_idx, slide in enumerate(prs.slides):
        for ph in slide.placeholders:
            s_x = _xfrm_dict(ph._element)
            if s_x is None:
                continue  # no explicit override => inherits, nothing to check
            base = _effective_baseline_xfrm(slide, ph)
            deviates = any(
                s_x[k] is not None and base[k] is not None
                and abs(s_x[k] - base[k]) > tolerance_emu
                for k in ("x", "y", "cx", "cy")
            )
            if deviates:
                out.append(GeometryDeviation(s_idx, ph.placeholder_format.idx, s_x, base))
    return out


def enforce_placeholder_geometry(prs, deviations: list[GeometryDeviation]) -> int:
    """Remove the explicit a:xfrm from deviating slide placeholders so they
    inherit layout geometry. Touches ONLY the xfrm element (surgical edit)."""
    fixed = 0
    for dev in deviations:
        slide = prs.slides[dev.slide_index]
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != dev.ph_idx:
                continue
            spPr = find(ph._element, "p:spPr")
            if spPr is None:
                spPr = find(ph._element, "a:spPr")
            xfrm = find(spPr, "a:xfrm")
            if xfrm is not None:
                spPr.remove(xfrm)
                fixed += 1
    return fixed


def run(fixture: Path, out_dir: Path) -> dict:
    prs = Presentation(fixture)
    census = layout_census(prs)
    deviations = find_geometry_deviations(prs)

    fixed = enforce_placeholder_geometry(prs, deviations)
    out_path = out_dir / (fixture.stem + ".enforced.pptx")
    prs.save(out_path)

    # Verify: reopen, deviations must be gone and the deck must parse cleanly.
    reopened = Presentation(out_path)
    residual = find_geometry_deviations(reopened)
    return {
        "census": census,
        "deviations_found": len(deviations),
        "deviations_detail": [vars(d) for d in deviations],
        "fixes_applied": fixed,
        "residual_after_enforce": len(residual),
        "enforced_output": str(out_path),
        "reopened_ok": True,
    }
