"""U2: safe write fidelity (round-trip preservation).

Two experiments per deck:
1. Plain round-trip: open + save with no edits. Compare package part
   inventories and bytes. Binary parts (media, embedded xlsx) must survive
   byte-identical; XML parts may re-serialize but must still parse.
2. Surgical edit: change ONE run's font on slide 1, save. The changed-part
   set should be minimal (the slide XML, not the whole package).

"Opens without repair prompt in PowerPoint desktop / M365 web" cannot be
automated here: outputs are written to out/ for that manual check.
"""

import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation


def _parts(path: Path) -> dict[str, bytes]:
    with zipfile.ZipFile(path) as z:
        return {n: z.read(n) for n in z.namelist()}


def _is_xml(name: str) -> bool:
    return name.endswith((".xml", ".rels"))


def _diff(before: dict, after: dict) -> dict:
    added = sorted(set(after) - set(before))
    removed = sorted(set(before) - set(after))
    changed_xml, changed_binary = [], []
    for name in sorted(set(before) & set(after)):
        if before[name] != after[name]:
            (changed_xml if _is_xml(name) else changed_binary).append(name)
    return {"added": added, "removed": removed,
            "changed_xml": changed_xml, "changed_binary": changed_binary}


def _all_xml_parses(path: Path) -> bool:
    with zipfile.ZipFile(path) as z:
        for name in z.namelist():
            if _is_xml(name):
                etree.fromstring(z.read(name))
    return True


def run(fixture: Path, out_dir: Path) -> dict:
    before = _parts(fixture)

    # 1. plain round-trip
    rt_path = out_dir / (fixture.stem + ".roundtrip.pptx")
    Presentation(fixture).save(rt_path)
    rt_diff = _diff(before, _parts(rt_path))

    # 2. surgical edit: one run's font family on the first slide
    prs = Presentation(fixture)
    edited = False
    for shape in prs.slides[0].shapes:
        if edited or not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if para.runs:
                para.runs[0].font.name = "Arial"
                edited = True
                break
    ed_path = out_dir / (fixture.stem + ".edited.pptx")
    prs.save(ed_path)
    ed_diff = _diff(before, _parts(ed_path))

    return {
        "roundtrip": {
            **rt_diff,
            "binary_parts_preserved": not rt_diff["changed_binary"] and not rt_diff["removed"],
            "output": str(rt_path),
        },
        "surgical_edit": {
            **ed_diff,
            "edit_applied": edited,
            "output": str(ed_path),
        },
        "outputs_parse_ok": _all_xml_parses(rt_path) and _all_xml_parses(ed_path),
        "manual_check_required": "Open both outputs in PowerPoint desktop and M365 web: zero repair prompts is the pass bar.",
    }
