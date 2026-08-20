"""U3: processing time on the 200-slide deck.

Times the operations v1 actually performs per audit: package open, full walk
with effective-font resolution on every run + solid-fill color resolution on
every shape + Arabic scan, and open+save (the write path). Rendering is out
of scope here (v1 report PDF needs no renderer; Graph timing is U5).
"""

import time
from pathlib import Path

from pptx import Presentation

from .arabic import scan_presentation
from .color_resolver import resolve_solid_fill
from .ns import find
from .resolver import resolve_run


def _walk_and_resolve(prs) -> dict:
    runs = shapes = fills = 0
    master_cache = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            shapes += 1
            spPr = find(shape._element, "p:spPr")
            if spPr is not None:
                master = slide.slide_layout.slide_master
                if resolve_solid_fill(spPr, master) is not None:
                    fills += 1
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    resolve_run(run, para, shape, slide, prs)
                    runs += 1
    return {"shapes": shapes, "runs_resolved": runs, "solid_fills_resolved": fills}


def run(fixture: Path, out_dir: Path) -> dict:
    t0 = time.perf_counter()
    prs = Presentation(fixture)
    t_open = time.perf_counter() - t0

    t0 = time.perf_counter()
    walk_stats = _walk_and_resolve(prs)
    t_walk = time.perf_counter() - t0

    t0 = time.perf_counter()
    hits = scan_presentation(prs)
    t_arabic = time.perf_counter() - t0

    t0 = time.perf_counter()
    prs.save(out_dir / (fixture.stem + ".perf-save.pptx"))
    t_save = time.perf_counter() - t0

    n = len(prs.slides)
    total = t_open + t_walk + t_arabic + t_save
    return {
        "slides": n,
        **walk_stats,
        "arabic_hits": len(hits),
        "seconds": {
            "open": round(t_open, 3),
            "resolve_walk": round(t_walk, 3),
            "arabic_scan": round(t_arabic, 3),
            "save": round(t_save, 3),
            "total": round(total, 3),
        },
        "ms_per_slide": round(total / n * 1000, 1),
    }
