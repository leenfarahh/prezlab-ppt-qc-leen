"""A slide made of drawn shapes is not an empty slide, and it is not a cover.

"The cover being defaulted when there are shapes that are not seen as content
is wrong - shapes are also types of content even when no text is present"
(design lead, 02/09/2026).

One reading of a slide, made in one function, wrong in one way, showing up in
three places:

  qc.extract.content_type returned None for any shape that was not a picture, a
  chart, a table, a group, or a box with words in it. None means "not content"
  to everything downstream. On the corpus deck that is 150 of 754 top-level
  shapes - the AUTO_SHAPEs, FREEFORMs, LINEs and embedded objects a designed
  deck is actually built from.

  qc.layoutgap.signature therefore reported a slide of shape-drawn cards as
  "no content blocks", so the coverage report clustered a deck of diagrams as a
  deck of empty slides.

  qc.layoutpick.rank ranked the master's layouts against that reading, and
  qc.layoutgap.fits abstains for a layout with no content boxes
  (MIN_BOXES_TO_JUDGE) - which the ranking read as "it fits". A Cover, the one
  layout in every master with nowhere to put anything, was therefore the only
  "fitting" candidate for a slide full of content, and fit sorts above score.

  qc.applymaster.plan_assignments fell back to target_layouts[0] when no
  archetype matched, and the first layout in a master is usually the cover.

Text is one KIND of content. It is not the definition of it.
"""

import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from qc.applymaster import _most_content, plan_assignments
from qc.extract import content_type
from qc.layoutgap import describe, signature
from qc.layoutpick import rank


def _prs():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    return prs


def _reopen(prs):
    buf = io.BytesIO()
    prs.save(buf)
    return Presentation(io.BytesIO(buf.getvalue()))


def _cards(slide, n=3):
    """What a Prezlab slide is mostly made of: drawn cards, no text in them."""
    for i in range(n):
        slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.8 + 4.2 * i), Inches(2.2),
                               Inches(3.6), Inches(3.0))


# ------------------------------------------------------------- what a shape IS


def test_a_wordless_card_is_content():
    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _cards(slide, 1)
    shape = _reopen(prs).slides[0].shapes[0]
    assert content_type(shape) == "shape"


def test_a_rule_and_a_freeform_are_content_too():
    """The two that made up most of the 150. A rule is furniture only when it
    is small enough to be furniture, and that is an AREA question the signature
    already asks - not a question about whether anyone typed in it."""
    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_connector(1, Inches(1), Inches(2), Inches(6), Inches(5))
    builder = slide.shapes.build_freeform(Inches(1), Inches(1))
    builder.add_line_segments([(Inches(4), Inches(1)), (Inches(4), Inches(4))])
    builder.convert_to_shape()

    kinds = {content_type(s) for s in _reopen(prs).slides[0].shapes}
    assert kinds == {"shape"}, kinds


def test_a_card_with_words_in_it_is_still_a_text_block():
    """The order matters the other way too. An autoshape holding a label is a
    text block that happens to have a fill, and calling it a shape would take
    every labelled card off the text count the layout rules read."""
    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1),
                                  Inches(2), Inches(3), Inches(2))
    card.text_frame.text = "Phase one"
    assert content_type(_reopen(prs).slides[0].shapes[0]) == "text"


def test_an_empty_text_box_is_still_not_content():
    """None keeps its meaning: a box that paints nothing is invisible on the
    rendered slide, and counting it would inflate every signature."""
    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
    assert content_type(_reopen(prs).slides[0].shapes[0]) is None


# ------------------------------------------------- what the slide asks for


def test_a_slide_of_drawn_cards_is_not_an_empty_slide():
    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _cards(slide, 3)
    deck = _reopen(prs)

    sig = signature(deck.slides[0], deck.slide_width, deck.slide_height)

    assert sig["blocks"] == 3, sig
    assert sig["shapes"] == 3
    assert sig["columns"] == 3
    assert "no content blocks" not in describe(sig)
    assert "3 drawn shapes" in describe(sig)


def test_a_genuinely_empty_slide_still_says_so():
    """The floor did not move. A slide with nothing on it reads as nothing."""
    prs = _prs()
    prs.slides.add_slide(prs.slide_layouts[6])
    deck = _reopen(prs)
    sig = signature(deck.slides[0], deck.slide_width, deck.slide_height)
    assert sig["blocks"] == 0
    assert describe(sig) == "no title, no content blocks"


def test_a_full_bleed_panel_is_still_a_ground():
    """A drawn shape covering the slide is the ground it sits on, and the area
    rule that said so still runs before the kind is counted."""
    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                           prs.slide_width, prs.slide_height)
    deck = _reopen(prs)
    sig = signature(deck.slides[0], deck.slide_width, deck.slide_height)
    assert sig["blocks"] == 0, "a backdrop was counted as a block of content"


# ------------------------------------------------------------- the cover


def _master_with_a_cover():
    """A cover with nothing but a title, and one content layout. Both carry the
    archetype tokens PowerPoint writes for its own layouts."""
    master = Presentation()
    from qc.stylespec import dominant_master, extract_layouts

    entries = extract_layouts(dominant_master(master), embed_assets=False)
    cover = next(e for e in entries if e["type"] == "title")
    body = next(e for e in entries if e["type"] == "obj")
    return cover, body


def test_a_layout_with_nowhere_to_put_content_does_not_fit_a_slide_with_some():
    """The abstention that became a recommendation. qc.layoutgap.fits returns
    True for a bodyless layout because it declines to judge it; the ranking
    sorts fitting layouts above everything else whatever the score, so the
    Cover led the dropdown for every slide in the deck."""
    cover, body = _master_with_a_cover()
    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _cards(slide, 3)
    deck = _reopen(prs)

    ranked, sig = rank(deck.slides[0], [cover, body],
                       deck.slide_width, deck.slide_height)

    by_name = {c.name: c for c in ranked}
    assert by_name[cover["name"]].fits is False
    assert "no content boxes" in by_name[cover["name"]].why
    assert ranked[0].name != cover["name"], "the cover is still the suggestion"


def test_the_cover_still_fits_a_slide_that_really_is_empty():
    """The guard is about content, not about covers. A slide with nothing on it
    and a layout with nowhere to put anything do agree."""
    cover, body = _master_with_a_cover()
    prs = _prs()
    prs.slides.add_slide(prs.slide_layouts[6])
    deck = _reopen(prs)

    ranked, _sig = rank(deck.slides[0], [cover, body],
                        deck.slide_width, deck.slide_height)

    assert {c.name for c in ranked if c.fits} == {cover["name"], body["name"]}


def test_the_cover_stays_in_the_dropdown():
    """Ranked last, never removed. A designer who wants it picks it - the
    ranking is a suggestion and every layout stays reachable."""
    cover, body = _master_with_a_cover()
    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _cards(slide, 3)
    deck = _reopen(prs)

    ranked, _sig = rank(deck.slides[0], [cover, body],
                        deck.slide_width, deck.slide_height)
    assert cover["name"] in {c.name for c in ranked}


def test_the_last_resort_fallback_is_not_the_first_layout_in_the_file():
    """A master built by a designer carries no archetype tokens - PowerPoint
    only writes those for its own built-in layouts - so plan_assignments ran
    out of rules and took target_layouts[0]. Masters are written front to back
    and the cover is at the front."""
    cover, body = _master_with_a_cover()
    # what a designer's master looks like to this code: named, not typed
    cover = {**cover, "type": None, "name": "Cover"}
    body = {**body, "type": None, "name": "One Column"}

    assert _most_content([cover, body])["name"] == "One Column"

    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _cards(slide, 3)
    plans = plan_assignments(_reopen(prs), [cover, body])

    assert plans[0].match_rule == "fallback"
    assert plans[0].target_layout == "One Column", \
        "a slide of content was sent to the cover"
