"""What the assistant may be asked to DO, and what happens when the ask does
not resolve.

The whole safety argument for an acting chatbot is in this file, and it is four
claims:

  - THE VOCABULARY IS THIS JOB'S OWN IDS. A model is never shown a finding that
    is not open or an issue type this deck does not have, so the common failure
    is a refusal, not a change nobody meant.
  - AN UNRESOLVED REQUEST IS REFUSED WITH WHAT WAS ACTUALLY THERE, never
    performed as an empty success and never widened to something adjacent.
  - THE SUMMARY IS BUILT FROM WHAT RESOLVED. "Six fixes on slide 7" is a count
    of six real record ids, so a designer confirming it is confirming a fact.
  - AND THE PLAN IS A PROPOSAL. Planning changes nothing, the token that
    performs it is spent on use, and every branch of the performing goes
    through the same helper the button on the page goes through.

The model is stubbed throughout (qc.chat.ask_json is the seam).
"""

import io

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Emu

from qc import web
from qc.actions import Refused, plan, vocabulary
from qc.design import DesignFinding, Remedy
from qc.migrate import ContentChange
from qc.remedy import Applied

IN = 914400


# ------------------------------------------------------------------ fixtures


def _deck_bytes() -> bytes:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    for _ in range(4):
        prs.slides.add_slide(prs.slide_layouts[6])
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _record(record_id, slide_index, issue_type="font.family_out_of_set",
            **over):
    base = {"record_id": record_id, "slide_index": slide_index,
            "module": issue_type.split(".")[0], "issue_type": issue_type,
            "severity": "warning", "action": "flagged",
            "confidence": "deterministic", "arabic_flag": False,
            "new_value": "Poppins", "message": "not in the allowed set"}
    base.update(over)
    return base


def _finding(finding_id="f1", slides=(0,)):
    return DesignFinding(
        finding_id=finding_id, kind="palette", severity="warning",
        headline="#203965 is Brand Navy spelled differently",
        detail="two spellings of one colour", slides=list(slides),
        evidence={"places": 2},
        options=[Remedy("snap", "Use Brand Navy (#1F3864) everywhere", "note",
                        op="set_color"),
                 Remedy("leave", "Leave it", "note")])


def _job(records=None, findings=None, applied=None, changes=None,
         deck=True):
    if records is None:
        records = [_record("r1", 6), _record("r2", 6),
                   _record("r3", 8, "margin_alignment.edge_misaligned")]
    return {
        "filename": "client deck.pptx",
        "profile": "prezlab_en",
        "deck": _deck_bytes() if deck else None,
        "manifest": {"slides": 12, "summary": {"total": len(records)},
                     "records": list(records) + [
                         {"module": "preflight", "issue_type": "preflight.skip",
                          "severity": "info", "slide_index": 0,
                          "record_id": "p1", "action": "flagged",
                          "confidence": "high", "arabic_flag": False}]},
        "design": list(findings if findings is not None else [_finding()]),
        "design_applied": list(applied or []),
        "changes": list(changes or []),
        "removed": [],
    }


# ------------------------------------------------------- the vocabulary given


def test_the_vocabulary_names_only_ids_this_job_has():
    vocab = vocabulary(_job())
    types = {e["issue_type"] for e in vocab["issue_types_you_may_name"]}
    assert types == {"font.family_out_of_set",
                     "margin_alignment.edge_misaligned"}
    assert [f["finding"] for f in vocab["open_design_findings"]] == ["f1"]
    assert [r["remedy"] for r in vocab["open_design_findings"][0]["remedies"]] \
        == ["snap", "leave"]
    # The preflight row is a note about the run, not a defect, so it is not
    # something anyone may ask to have fixed.
    assert "preflight.skip" not in types


def test_an_answered_finding_leaves_the_vocabulary():
    answered = Applied(finding_id="f1", remedy_id="snap", kind="palette",
                       headline="h", label="l", done=True, detail="d",
                       undo=[{"op": "noop"}])
    vocab = vocabulary(_job(applied=[answered]))
    assert vocab["open_design_findings"] == []
    assert [d["finding"] for d in vocab["decisions_already_applied"]] == ["f1"]


def test_an_action_with_nothing_to_act_on_is_not_offered():
    """An action that can only refuse is a trap: the model picks it because it
    reads as the right verb."""
    vocab = vocabulary(_job(records=[], findings=[]))
    assert "fix_findings" not in vocab["actions"]
    assert "undo" not in vocab["actions"]
    assert "remove_pieces" not in vocab["actions"]


def test_nothing_is_offered_once_the_deck_is_gone():
    vocab = vocabulary(_job(deck=False))
    assert vocab["actions"] == {}
    assert vocab["can_change_the_deck"] is False


def test_recheck_is_not_offered_on_a_run_that_was_never_audited():
    job = _job()
    job.pop("manifest")
    job["plans"] = []
    assert "recheck" not in vocabulary(job)["actions"]


# ----------------------------------------------------- resolving, and refusing


def test_a_slide_outside_the_deck_is_refused_not_clamped():
    """Clamping would turn 'fix slide 40' on a 12-slide deck into a silent
    change to slide 12."""
    with pytest.raises(Refused) as exc:
        plan(_job(), {"name": "fix_findings", "slide": 40})
    assert "no slide 40" in str(exc.value)
    assert "12" in str(exc.value)


def test_a_slide_with_nothing_fixable_is_refused_with_the_slides_that_do():
    with pytest.raises(Refused) as exc:
        plan(_job(), {"name": "fix_findings", "slide": 2})
    assert "no fixable finding on slide 2" in str(exc.value)
    assert "7" in str(exc.value) and "9" in str(exc.value)


def test_an_issue_type_not_on_that_slide_is_refused_with_what_is():
    with pytest.raises(Refused) as exc:
        plan(_job(), {"name": "fix_findings", "slide": 7,
                      "issue_types": ["typography.case_inconsistent"]})
    assert "font.family_out_of_set" in str(exc.value)


def test_the_summary_counts_exactly_what_resolved():
    out = plan(_job(), {"name": "fix_findings", "slide": 7})
    assert out.record_ids == ["r1", "r2"]
    assert "Apply 2 audit fixes on slide 7" in out.summary
    assert "2 font.family_out_of_set" in out.summary
    assert out.changes is True


def test_naming_no_slide_takes_the_whole_deck():
    out = plan(_job(), {"name": "fix_findings"})
    assert sorted(out.record_ids) == ["r1", "r2", "r3"]
    assert "the whole deck" in out.summary


def test_fixes_held_for_explicit_approval_are_excluded_and_said_so():
    """Arabic font substitutions and whole-slide body moves are never applied
    on the strength of a sentence."""
    held = _record("r9", 6, arabic_flag=True)
    out = plan(_job(records=[_record("r1", 6), held]),
               {"name": "fix_findings", "slide": 7})
    assert out.record_ids == ["r1"]
    assert "held for your explicit approval" in out.summary


def test_a_slide_of_nothing_but_held_fixes_refuses_and_says_how_to_proceed():
    held = _record("r9", 6, arabic_flag=True)
    with pytest.raises(Refused) as exc:
        plan(_job(records=[held]), {"name": "fix_findings", "slide": 7})
    assert "explicit approval" in str(exc.value)


def test_including_the_holds_is_something_that_has_to_be_asked_for():
    held = _record("r9", 6, arabic_flag=True)
    out = plan(_job(records=[held]),
               {"name": "fix_findings", "slide": 7, "include_holds": True})
    assert out.record_ids == ["r9"]


# ------------------------------------------------------------ design remedies


def test_a_remedy_that_is_not_one_of_the_options_is_refused():
    with pytest.raises(Refused) as exc:
        plan(_job(), {"name": "take_remedy", "finding": "f1",
                      "remedy": "recolour-everything"})
    assert "Use Brand Navy" in str(exc.value)


def test_a_finding_that_is_not_open_is_refused():
    with pytest.raises(Refused) as exc:
        plan(_job(), {"name": "take_remedy", "finding": "nope",
                      "remedy": "snap"})
    assert "not open on this deck" in str(exc.value)


def test_a_resolved_remedy_carries_the_real_finding_and_remedy_objects():
    out = plan(_job(), {"name": "take_remedy", "finding": "f1",
                        "remedy": "snap"})
    finding, remedy = out.picks[0]
    assert finding.finding_id == "f1"
    assert remedy.remedy_id == "snap"
    assert "Use Brand Navy" in out.summary


# ------------------------------------------------------------------ handing over


def test_decide_counts_both_passes_and_says_what_is_left():
    undecidable = DesignFinding(
        finding_id="f2", kind="fit", severity="warning", headline="h",
        detail="d", slides=[0], options=[Remedy("a", "A", "n")])
    out = plan(_job(findings=[_finding(), undecidable]), {"name": "decide"})
    assert out.name == "decide"
    assert len(out.record_ids) == 3
    assert "3 audit fixes" in out.summary
    assert "1 design decision" in out.summary


def test_decide_on_a_quiet_slide_is_refused_rather_than_performed_empty():
    with pytest.raises(Refused) as exc:
        plan(_job(records=[], findings=[]), {"name": "decide", "slide": 2})
    assert "Nothing on slide 2" in str(exc.value)


# ------------------------------------------------------------------------ undo


def test_undo_with_no_id_takes_the_most_recent_decision():
    """'Undo that' means the last thing that happened. Guessing wider would
    reverse work nobody asked about."""
    first = Applied(finding_id="a", remedy_id="r", kind="k", headline="h",
                    label="l", done=True, detail="d", undo=[{"op": "x"}])
    last = Applied(finding_id="b", remedy_id="r", kind="k", headline="h",
                   label="l", done=True, detail="d", undo=[{"op": "x"}])
    out = plan(_job(applied=[first, last]), {"name": "undo"})
    assert out.finding_ids == ["b"]


def test_undo_of_something_never_applied_is_refused():
    applied = Applied(finding_id="a", remedy_id="r", kind="k", headline="h",
                      label="l", done=True, detail="d", undo=[{"op": "x"}])
    with pytest.raises(Refused) as exc:
        plan(_job(applied=[applied]), {"name": "undo", "findings": ["zzz"]})
    assert "not in the applied list" in str(exc.value)


def test_undo_on_a_deck_nothing_was_applied_to_is_refused():
    with pytest.raises(Refused) as exc:
        plan(_job(), {"name": "undo"})
    assert "nothing to take back" in str(exc.value)


# -------------------------------------------------------------------- removals


def _proposal(remove_id, slide_index):
    return ContentChange(
        slide_index=slide_index, action="duplicate of the master's footer",
        detail="the master already carries this", severity="alert",
        remove_op={"kind": "shape", "slide_index": slide_index,
                   "shape_id": "5"},
        remove_id=remove_id)


def test_removals_resolve_to_the_pieces_actually_waiting():
    job = _job(changes=[_proposal("x1", 1), _proposal("x2", 3)])
    out = plan(job, {"name": "remove_pieces", "slide": 2})
    assert out.remove_ids == ["x1"]
    assert "slide 2" in out.summary


def test_a_piece_already_removed_is_not_offered_again():
    job = _job(changes=[_proposal("x1", 1)])
    job["removed"] = ["x1"]
    with pytest.raises(Refused) as exc:
        plan(job, {"name": "remove_pieces"})
    assert "Nothing is waiting to be removed" in str(exc.value)


# ------------------------------------------------------------------- the door


def test_an_unknown_action_name_is_refused_with_the_list():
    with pytest.raises(Refused) as exc:
        plan(_job(), {"name": "delete_the_deck"})
    assert "not something this assistant can do" in str(exc.value)


def test_nothing_can_be_planned_once_the_deck_is_gone():
    with pytest.raises(Refused) as exc:
        plan(_job(deck=False), {"name": "fix_findings"})
    assert "no longer held in memory" in str(exc.value)


# ------------------------------------------------- through the ask box and back


@pytest.fixture()
def client(monkeypatch):
    monkeypatch.setattr(web, "AUTH_REQUIRED", False)
    monkeypatch.setattr(web.app.state, "auth_required", False)
    monkeypatch.setattr(web, "AI_ENABLED", True)
    monkeypatch.setattr("qc.web._chat_available", lambda: (True, ""))
    return TestClient(web.app)


def _answer(action=None, text="Six font fixes on slide 7."):
    out = {"answer": text, "links": [], "colours": [], "slides": [],
           "answerable": True}
    if action is not None:
        out["action"] = action
    return out


@pytest.fixture()
def registered(monkeypatch):
    job = _job()
    web._jobs["chatjob1"] = job
    yield job
    web._jobs.pop("chatjob1", None)


def test_asking_for_a_fix_comes_back_as_a_plan_and_changes_nothing(
        client, registered, monkeypatch):
    monkeypatch.setattr(
        "qc.chat.ask_json",
        lambda **k: _answer({"name": "fix_findings", "slide": 7}))
    before = registered["deck"]
    reply = client.post("/chat/chatjob1", json={"q": "fix the fonts on slide 7"})
    assert reply.status_code == 200
    body = reply.json()
    assert body["plan"]["summary"].startswith("Apply 2 audit fixes on slide 7")
    assert body["plan"]["changes"] is True
    assert body["plan"]["token"]
    assert registered["deck"] is before        # nothing has happened


def test_a_question_comes_back_with_no_plan(client, registered, monkeypatch):
    monkeypatch.setattr("qc.chat.ask_json",
                        lambda **k: _answer(None, "Two navies, one hand-typed."))
    body = client.post("/chat/chatjob1", json={"q": "which navy"}).json()
    assert "plan" not in body


def test_an_unresolvable_ask_still_answers_and_says_what_was_there(
        client, registered, monkeypatch):
    monkeypatch.setattr(
        "qc.chat.ask_json",
        lambda **k: _answer({"name": "fix_findings", "slide": 2}))
    body = client.post("/chat/chatjob1", json={"q": "fix slide 2"}).json()
    assert "plan" not in body
    assert "no fixable finding on slide 2" in body["refusal"]


def test_the_token_is_spent_on_use(client, registered, monkeypatch):
    """A plan describes a deck in a particular state. Performing it twice
    applies it to a deck the first pass already changed."""
    monkeypatch.setattr(
        "qc.chat.ask_json",
        lambda **k: _answer({"name": "fix_findings", "slide": 7}))
    token = client.post("/chat/chatjob1",
                        json={"q": "fix slide 7"}).json()["plan"]["token"]

    performed = []

    def _fake_perform(job, plan_obj):
        performed.append(plan_obj)
        return "Applied 2 fixes.", True

    monkeypatch.setattr("qc.web._perform_plan", _fake_perform)
    first = client.post("/chat/chatjob1/do", json={"token": token})
    assert first.status_code == 200
    assert first.json()["reload"] is True
    second = client.post("/chat/chatjob1/do", json={"token": token})
    assert second.status_code == 409
    assert "no longer on offer" in second.json()["error"]
    assert len(performed) == 1


def test_performing_an_unknown_token_changes_nothing(client, registered):
    reply = client.post("/chat/chatjob1/do", json={"token": "made-up"})
    assert reply.status_code == 409


def test_a_plan_goes_through_the_same_helper_the_button_does(
        client, registered, monkeypatch):
    """The safety story in one assertion: there is no fix reachable from the
    ask box that is not the one the page applies, with its re-audit and its
    place in the Undo list."""
    monkeypatch.setattr(
        "qc.chat.ask_json",
        lambda **k: _answer({"name": "fix_findings", "slide": 7}))
    token = client.post("/chat/chatjob1",
                        json={"q": "fix slide 7"}).json()["plan"]["token"]

    seen = {}

    class _Fx:
        applied = 2
        outcomes = []

    def _apply(job, selected):
        seen["selected"] = selected
        job["manifest"]["summary"]["total"] = 1
        return _Fx(), 3, None

    monkeypatch.setattr("qc.web._apply_audit_fixes", _apply)
    body = client.post("/chat/chatjob1/do", json={"token": token}).json()
    assert seen["selected"] == {"r1", "r2"}
    assert "Applied 2 fixes" in body["note"]
    assert "1 findings remain (was 3)" in body["note"]


def test_ai_off_refuses_the_do_route_outright(client, registered, monkeypatch):
    monkeypatch.setattr(web, "AI_ENABLED", False)
    reply = client.post("/chat/chatjob1/do", json={"token": "anything"})
    assert reply.status_code == 503
