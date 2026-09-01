"""The model names the line; the fix must not drag the line off itself.

The gap between detection and execution, found on a real deck (01/09/2026). The
model looked at a five-across grid of goal circles, said the third one sits low,
and named all five as the row it should line up with. Every part of that was
right. Then the fix ran:

    before   2.020  2.020  2.183  2.020  1.999      one circle low
    after    2.020  1.856  2.020  1.856  1.999      two circles high

The stray landed on the line and two of the shapes DEFINING that line were
dragged the same distance off it. Row spread went from 0.184in to 0.163in, so
the numbers said "improved" and the slide was visibly worse.

The cause is not in the model, the prompt or the precision gate. qc.fixer infers
what must travel with a moving shape from overlap and adjacency, and for a
vertical move that means "whatever sits beside it in the same row within 10mm" -
which on a row of cards is the rest of the row (qc.util.rides_with). Geometry
cannot tell a satellite from a peer. The model can, and did, and the answer was
being thrown away at the record boundary: the alignment record named one shape
and forgot the set it was measured against.

So the record now carries its cluster ("align-y:6,8,10,12,14") and the fix pins
those peers against riding, exactly as qc.fixer._fix_space_edge pins a frame
stack and _fix_component_edge pins a component's members. Third time the same
lesson: a shape being held to a line cannot also be a satellite of a shape being
moved onto that line.
"""

import io

from pptx import Presentation
from pptx.util import Emu

import qc.copilot as copilot
from qc.fixer import apply_fixes

IN = 914400
TOL = 28575          # copilot.TOL_EMU: 0.03in, the perceptual floor


def _row(tops, left=1.0, gap=0.25, size=(2.08, 2.10)):
    """A row of same-size cards at the given tops (inches), gap between them
    smaller than the 10mm satellite window - which is what an ordinary row of
    cards looks like and what makes every neighbour a candidate rider."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w, h = size
    cards = []
    for i, top in enumerate(tops):
        card = slide.shapes.add_textbox(
            Emu(int((left + i * (w + gap)) * IN)), Emu(int(top * IN)),
            Emu(int(w * IN)), Emu(int(h * IN)))
        card.text_frame.text = f"Goal {i + 1}"
        cards.append(card)
    buf = io.BytesIO()
    prs.save(buf)
    return prs, slide, cards, buf.getvalue()


def _observation(cards):
    return [{"action": "align_top",
             "shape_ids": [str(c.shape_id) for c in cards],
             "rationale": "One circle sits lower than the rest of the row."}]


def _tops(deck_bytes):
    slide = Presentation(io.BytesIO(deck_bytes)).slides[0]
    return {str(s.shape_id): s.top for s in slide.shapes}


def test_the_record_names_the_set_it_was_measured_against():
    """Without this the fix has no way to know a peer from a satellite."""
    _prs, slide, cards, _data = _row([2.02, 2.02, 2.18, 2.02, 2.00])

    recs = copilot.synthesize(slide, 0, _observation(cards), existing=[])

    assert len(recs) == 1
    assert recs[0]["locator"] == "align-y:" + ",".join(
        sorted((str(c.shape_id) for c in cards), key=int))


def test_only_the_stray_moves_and_the_row_comes_level():
    _prs, slide, cards, data = _row([2.02, 2.02, 2.18, 2.02, 2.00])
    stray = str(cards[2].shape_id)

    recs = copilot.synthesize(slide, 0, _observation(cards), existing=[])
    result = apply_fixes(data, recs, {r["record_id"] for r in recs})
    assert result.applied == 1, [o.reason for o in result.outcomes]

    before, after = _tops(data), _tops(result.cleaned_bytes)
    moved = {sid for sid in before if before[sid] != after[sid]}
    assert moved == {stray}, \
        f"the fix moved shapes that were already on the line: {moved - {stray}}"

    tops = [after[str(c.shape_id)] for c in cards]
    assert max(tops) - min(tops) <= TOL, \
        f"row is still not level: {[t / IN for t in tops]}"


def test_the_peers_that_defined_the_line_are_where_they_were():
    """The specific regression: a shape already on the line must not be dragged
    off it by the shape moving onto it."""
    _prs, slide, cards, data = _row([2.02, 2.02, 2.18, 2.02, 2.00])

    recs = copilot.synthesize(slide, 0, _observation(cards), existing=[])
    after = _tops(apply_fixes(data, recs,
                              {r["record_id"] for r in recs}).cleaned_bytes)

    for card in (cards[1], cards[3]):
        assert after[str(card.shape_id)] == card.top, \
            "a neighbour of the stray was carried along with it"


def test_two_strays_in_one_row_are_both_fixed_in_one_round():
    """Both records name the same cluster, so with the peers pinned neither
    claims the other's shapes and the round does not have to be repeated. Before
    the cluster was carried, the first fix claimed the whole row and the second
    came back "shares shapes with a fix already applied this round"."""
    _prs, slide, cards, data = _row([2.02, 2.20, 2.02, 1.84, 2.02])

    recs = copilot.synthesize(slide, 0, _observation(cards), existing=[])
    assert len(recs) == 2, [r["shape_id"] for r in recs]

    result = apply_fixes(data, recs, {r["record_id"] for r in recs})
    assert result.applied == 2, [o.reason for o in result.outcomes]

    after = _tops(result.cleaned_bytes)
    tops = [after[str(c.shape_id)] for c in cards]
    assert max(tops) - min(tops) <= TOL, \
        f"row is still not level: {[t / IN for t in tops]}"


def test_a_label_inside_the_stray_travels_with_it():
    """Pinning the peers must not pin the mover's OWN contents: a card's label
    is part of the card and moving one without the other tears it apart."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cards, w, h = [], 2.08, 2.10
    for i, top in enumerate((2.02, 2.02, 2.18)):
        card = slide.shapes.add_textbox(
            Emu(int((1.0 + i * (w + 0.25)) * IN)), Emu(int(top * IN)),
            Emu(int(w * IN)), Emu(int(h * IN)))
        card.text_frame.text = f"Goal {i + 1}"
        cards.append(card)
    # A chip well inside the low card, which is the third one.
    chip = slide.shapes.add_textbox(Emu(int((1.0 + 2 * (w + 0.25) + 0.5) * IN)),
                                   Emu(int(2.6 * IN)), Emu(int(0.6 * IN)),
                                   Emu(int(0.5 * IN)))
    chip.text_frame.text = "03"
    buf = io.BytesIO()
    prs.save(buf)
    data = buf.getvalue()

    recs = copilot.synthesize(slide, 0, _observation(cards), existing=[])
    result = apply_fixes(data, recs, {r["record_id"] for r in recs})
    assert result.applied == 1, [o.reason for o in result.outcomes]

    after = _tops(result.cleaned_bytes)
    delta = after[str(cards[2].shape_id)] - cards[2].top
    assert delta != 0
    assert after[str(chip.shape_id)] - chip.top == delta, \
        "the chip inside the moved card stayed behind"


def test_a_record_without_a_cluster_behaves_exactly_as_before():
    """Every measured edge record states no cluster. Those must keep the old
    behaviour - the fix carries what adjacency says it carries - or this change
    would silently alter forty audits' worth of fixes."""
    from qc.fixer import _cluster_ids, _peer_pinned

    _prs, slide, cards, _data = _row([2.02, 2.02, 2.18, 2.02, 2.00])
    plain = {"locator": None, "shape_id": str(cards[2].shape_id)}

    assert _cluster_ids(plain) == set()
    assert _peer_pinned(slide, plain, cards[2].shape_id) == set()
