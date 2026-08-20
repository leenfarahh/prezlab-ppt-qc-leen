"""Annotated .pptx: findings and comments written into speaker notes."""

import io

from pptx import Presentation

from qc.annotate import MARKER, build_annotated
from qc.engine import run_audit


def _manifest(fixtures_dir):
    return run_audit(fixtures_dir / "mixed_layouts.pptx", "prezlab_en").to_manifest()


def _notes(deck_bytes: bytes, idx: int) -> str:
    prs = Presentation(io.BytesIO(deck_bytes))
    slide = prs.slides[idx]
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text


def test_flagged_slide_gets_marker_and_findings(fixtures_dir):
    manifest = _manifest(fixtures_dir)
    deck = (fixtures_dir / "mixed_layouts.pptx").read_bytes()
    out = build_annotated(deck, manifest, [])
    flagged = manifest["records"][0]["slide_index"] if manifest["records"] else 0
    flagged = next(r["slide_index"] for r in manifest["records"]
                   if r["module"] != "preflight")
    notes = _notes(out, flagged)
    assert MARKER in notes
    assert "font.family_out_of_set" in notes or "[" in notes
    assert deck == (fixtures_dir / "mixed_layouts.pptx").read_bytes()  # input untouched


def test_clean_slide_untouched(fixtures_dir):
    manifest = _manifest(fixtures_dir)
    deck = (fixtures_dir / "clean.pptx").read_bytes()
    clean_manifest = run_audit(fixtures_dir / "clean.pptx", "prezlab_en",
                               modules=["header_footer"]).to_manifest()
    out = build_annotated(deck, clean_manifest, [])
    prs = Presentation(io.BytesIO(out))
    assert not any(s.has_notes_slide and MARKER in s.notes_slide.notes_text_frame.text
                   for s in prs.slides)


def test_existing_notes_preserved(fixtures_dir):
    manifest = _manifest(fixtures_dir)
    src = Presentation(fixtures_dir / "mixed_layouts.pptx")
    flagged = next(r["slide_index"] for r in manifest["records"]
                   if r["module"] != "preflight")
    src.slides[flagged].notes_slide.notes_text_frame.text = "Keep this speaker note."
    buf = io.BytesIO()
    src.save(buf)
    out = build_annotated(buf.getvalue(), manifest, [])
    notes = _notes(out, flagged)
    assert "Keep this speaker note." in notes and MARKER in notes


def test_idempotent_reannotation(fixtures_dir):
    manifest = _manifest(fixtures_dir)
    deck = (fixtures_dir / "mixed_layouts.pptx").read_bytes()
    once = build_annotated(deck, manifest, [])
    twice = build_annotated(once, manifest, [])
    flagged = next(r["slide_index"] for r in manifest["records"]
                   if r["module"] != "preflight")
    assert _notes(twice, flagged).count(MARKER) == 1


def test_comments_appear_with_author(fixtures_dir):
    manifest = _manifest(fixtures_dir)
    deck = (fixtures_dir / "mixed_layouts.pptx").read_bytes()
    flagged = next(r["slide_index"] for r in manifest["records"]
                   if r["module"] != "preflight")
    out = build_annotated(deck, manifest, [
        {"slide_index": flagged, "author": "Sanad", "text": "check the title here"}])
    notes = _notes(out, flagged)
    assert "Comments:" in notes and "- Sanad: check the title here" in notes


def test_output_is_valid_presentation(fixtures_dir):
    manifest = _manifest(fixtures_dir)
    deck = (fixtures_dir / "mixed_layouts.pptx").read_bytes()
    out = build_annotated(deck, manifest, [])
    assert len(Presentation(io.BytesIO(out)).slides) == 10
