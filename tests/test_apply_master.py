"""Applying a master to a deck: layout planning, the template store, and the
per-slide copy/apply/delete loop.

The loop itself needs desktop PowerPoint, so those tests skip without it. The
planning and storage layers are pure and always run, which is deliberate: the
part most likely to be wrong (which layout a slide lands on) must be testable
on any machine.
"""

import io
from pathlib import Path

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from qc import web
from qc.applymaster import apply_master, plan_assignments
from qc.stylespec import dominant_master, extract_layouts
from qc.templates import (delete_master, has_master, load_master, master_info,
                          save_master, templates_dir)
from qc.unify import com_available

needs_powerpoint = pytest.mark.skipif(
    not com_available(), reason="applying a master needs desktop PowerPoint")


def _targets(prs=None) -> list[dict]:
    prs = prs or Presentation()
    return extract_layouts(dominant_master(prs), embed_assets=False)


def _deck(layout_indexes=(0, 1, 2)) -> Presentation:
    d = Presentation()
    for i in layout_indexes:
        d.slides.add_slide(d.slide_layouts[i])
    return d


def _bytes(prs) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ------------------------------------------------------------------ planning


def test_layout_matched_by_name_when_both_sides_agree():
    plans = plan_assignments(_deck(), _targets())
    assert [p.match_rule for p in plans] == ["name"] * 3
    assert plans[0].target_layout == "Title Slide"
    assert plans[2].target_layout == "Section Header"


def test_archetype_carries_a_renamed_master():
    """A designer who renamed every layout still meant secHead to mean
    section header; the OOXML archetype is the master's own statement."""
    renamed = [dict(l, name="BRAND " + l["name"]) for l in _targets()]
    plans = plan_assignments(_deck(), renamed)

    assert [p.match_rule for p in plans] == ["archetype"] * 3
    assert plans[2].target_layout == "BRAND Section Header"


def test_fallback_is_labelled_and_explained():
    """A guess must never read as a match: the rule and the reason both land
    in the plan so the report can flag the slide for a human."""
    only_blank = [l for l in _targets() if l["type"] == "blank"]
    plans = plan_assignments(_deck(), only_blank)

    assert {p.match_rule for p in plans} == {"fallback"}
    assert all(p.target_layout == "Blank" for p in plans)
    assert "no layout named" in plans[0].note
    assert "orphaned" in plans[0].note


def test_ambiguous_archetype_prefers_the_closest_placeholder_count():
    """Two layouts share an archetype; a slide should land on the one whose
    shape matches, not simply the first."""
    targets = _targets()
    obj = next(l for l in targets if l["type"] == "obj")
    lean = dict(obj, name="Lean", placeholders=obj["placeholders"][:1])
    rich = dict(obj, name="Rich",
                placeholders=obj["placeholders"] + obj["placeholders"])
    renamed = [l for l in targets if l["type"] != "obj"] + [lean, rich]
    for l in renamed:
        l["name"] = "X " + l["name"] if not l["name"].startswith(("Lean", "Rich")) \
            else l["name"]

    plan = plan_assignments(_deck([1]), renamed)[0]
    assert plan.match_rule == "archetype"
    assert plan.target_layout in ("Lean", "Rich")
    assert "share this type" in plan.note


def test_planning_reads_every_slide_even_with_no_usable_target():
    plans = plan_assignments(_deck(), [])
    assert len(plans) == 3
    assert all(p.target_layout is None and p.match_rule == "none" for p in plans)


# ------------------------------------------------------------ template store


def test_master_round_trips_through_the_store():
    blob = _bytes(Presentation())
    save_master("client_x", blob)

    assert has_master("client_x")
    assert load_master("client_x") == blob
    info = master_info("client_x")
    assert info["bytes"] == len(blob) and len(info["sha1"]) == 40
    assert delete_master("client_x") and not has_master("client_x")


def test_store_honours_the_test_data_dir(tmp_path):
    """Templates must land in the per-test tmp dir like every other artifact.

    Asserted against tmp_path explicitly, not against templates_dir(): the
    first version of this test compared the store to itself, so it passed
    while every test template was being written into the real data/
    directory."""
    save_master("scoped", _bytes(Presentation()))

    assert templates_dir() == tmp_path / "templates"
    assert (tmp_path / "templates" / "scoped.pptx").exists()

    from qc.store import DATA_DIR as REAL_DATA_DIR  # unpatched module constant
    assert not (Path(REAL_DATA_DIR).parent / "data" / "templates"
                / "scoped.pptx").exists()


def test_ids_with_a_colon_are_stored_safely():
    """master:foo is a legal profile id here but an illegal Windows filename."""
    save_master("master:foo", _bytes(Presentation()))
    assert has_master("master:foo")
    assert load_master("master:foo")


def test_path_traversal_ids_are_refused():
    assert save_master("../escape", b"x") is None
    assert load_master("../escape") is None
    assert not has_master("../../etc/passwd")


# ---------------------------------------------------------------- the engine


def test_apply_refuses_clearly_without_powerpoint(monkeypatch):
    """No silent no-op: a machine that cannot do this has to say so."""
    monkeypatch.setattr("qc.applymaster.com_available", lambda: False)
    result = apply_master(_bytes(_deck()), _bytes(Presentation()),
                          plan_assignments(_deck(), _targets()))

    assert result.deck is None
    assert "desktop PowerPoint" in result.fatal
    assert result.applied == 0


@needs_powerpoint
def test_slide_count_never_changes_and_content_survives():
    """The copy/apply/delete ordering exists so the deck is never doubled;
    if that broke, the slide count is where it would show."""
    deck = _deck([1, 1, 1])
    for i, slide in enumerate(deck.slides):
        slide.shapes.title.text = f"Heading {i + 1}"
        slide.placeholders[1].text = f"Body {i + 1}"
    deck_bytes = _bytes(deck)

    master = Presentation()
    targets = _targets(master)
    result = apply_master(deck_bytes, _bytes(master),
                          plan_assignments(deck, targets))

    assert result.fatal is None and not result.errors
    assert result.applied == 3
    out = Presentation(io.BytesIO(result.deck))
    assert len(out.slides) == 3
    for i, slide in enumerate(out.slides):
        text = " ".join(sh.text_frame.text for sh in slide.shapes
                        if sh.has_text_frame)
        assert f"Heading {i + 1}" in text
        assert f"Body {i + 1}" in text


@needs_powerpoint
def test_each_slide_lands_on_its_planned_layout():
    deck = _deck([0, 1, 2])
    master = Presentation()
    targets = _targets(master)
    plans = plan_assignments(deck, targets)

    result = apply_master(_bytes(deck), _bytes(master), plans)
    out = Presentation(io.BytesIO(result.deck))
    assert [s.slide_layout.name for s in out.slides] == \
           [p.target_layout for p in plans]


# ----------------------------------------------------------------- the page


def _client(monkeypatch):
    monkeypatch.setattr(web, "AUTH_REQUIRED", False)
    web.app.state.auth_required = False
    return TestClient(web.app)


def test_page_explains_profiles_that_cannot_be_applied(monkeypatch):
    """The seeded profiles carry no master. Hiding them would read as a bug;
    the page lists them with the reason instead."""
    r = _client(monkeypatch).get("/format")
    assert r.status_code == 200
    assert "carry no master file" in r.text
    assert "prezlab" in r.text.lower()


def test_profile_without_a_master_is_refused(monkeypatch):
    client = _client(monkeypatch)
    r = client.post("/format", data={"profile": "prezlab_en"},
                    files={"deck": ("d.pptx", _bytes(_deck()), "app/x")})
    assert r.status_code == 400
    assert "carries no master" in r.text


def test_unknown_profile_is_refused(monkeypatch):
    client = _client(monkeypatch)
    r = client.post("/format", data={"profile": "nope"},
                    files={"deck": ("d.pptx", _bytes(_deck()), "app/x")})
    assert r.status_code == 400
    assert "Unknown profile" in r.text


def test_non_pptx_is_refused(monkeypatch):
    client = _client(monkeypatch)
    r = client.post("/format", data={"profile": "prezlab_en"},
                    files={"deck": ("notes.txt", b"x", "text/plain")})
    assert r.status_code == 400
    assert "Only .pptx" in r.text


def test_download_of_an_unknown_job_is_a_clean_404(monkeypatch):
    r = _client(monkeypatch).get("/format/deadbeef/download")
    assert r.status_code == 404


# ------------------------------------------- the master that stays behind
#
# A slide that cannot be rebuilt keeps the deck's ORIGINAL design alive to serve
# it, so the output carries two masters and PowerPoint's master view lists the
# original FIRST. A designer opens that view, sees a master with none of the new
# guides or furniture, and reads it as "the master was not copied" - when it
# was, onto the other one (design lead, 21/08/2026).


@needs_powerpoint
def test_a_slide_that_cannot_be_rebuilt_leaves_a_second_master_and_says_so():
    from qc.applymaster import SlidePlan

    deck = _deck([1, 1, 1])
    for i, slide in enumerate(deck.slides):
        slide.shapes.title.text = f"Heading {i + 1}"
    master = Presentation()
    plans = plan_assignments(deck, _targets(master))
    plans[1] = SlidePlan(slide_index=1, source_layout=plans[1].source_layout,
                         source_type=plans[1].source_type,
                         target_layout="No Such Layout", match_rule="none",
                         note="the master defines no usable layout")

    result = apply_master(_bytes(deck), _bytes(master), plans)
    assert result.applied == 2
    assert result.stragglers == [1], result.stragglers
    assert result.masters == 2, "the original master has to stay for that slide"

    out = Presentation(io.BytesIO(result.deck))
    assert len(list(out.slide_masters)) == 2

    from qc.ui_format import render_format_result

    html = render_format_result(
        deck_name="d.pptx", profile_name="p", job_id="j", plans=result.plans,
        errors=result.errors, applied=result.applied, content_changes=[],
        masters=result.masters, stragglers=result.stragglers)
    assert "carries 2 slide masters" in html
    assert "Slide(s) 2" in html
    assert "lists the ORIGINAL master first" in html


def test_a_clean_run_says_nothing_about_masters():
    """The note is for the case that misleads, not a permanent disclaimer."""
    from qc.ui_format import render_format_result

    html = render_format_result(
        deck_name="d.pptx", profile_name="p", job_id="j", plans=[], errors={},
        applied=3, content_changes=[], masters=1, stragglers=[])
    assert "slide masters" not in html


# --------------------------------------------- putting removals back
#
# The migration removes header text the master has no placeholder for and says
# so loudly. These cover the other half of that promise: the designer decides
# whether the removal was right, and the piece goes back exactly as it was
# rather than being retyped off the report (design lead, 20/08/2026).


def _job_with_a_removal(client):
    """A finished format job whose change list carries one removed shape."""
    from qc.migrate import ContentChange, migrate_deck
    from qc.util import iter_shapes_deep

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    box = slide.shapes.add_textbox(914400, 457200, 2743200, 457200)
    box.text_frame.text = "EYEBROW"
    xml = None
    for shape, _p in iter_shapes_deep(slide.shapes):
        if shape.has_text_frame and shape.text_frame.text == "EYEBROW":
            from lxml import etree

            xml = etree.tostring(shape._element, encoding="unicode")
            shape._element.getparent().remove(shape._element)
    assert xml, "the fixture must capture the removed element"

    job_id = "restorejob"
    web._format_jobs[job_id] = {
        "deck": _bytes(deck), "filename": "d.pptx", "profile": "prezlab_en",
        "plans": [], "errors": {}, "applied": 1, "restored": [],
        "changes": [ContentChange(0, "removed unplaced text",
                                  "'EYEBROW' had no slot", severity="alert",
                                  removed_text="EYEBROW", removed_xml=xml,
                                  restore_id="0-9")],
    }
    return job_id


def test_a_removed_piece_can_be_put_back_from_the_result_page(monkeypatch):
    client = _client(monkeypatch)
    job_id = _job_with_a_removal(client)
    try:
        r = client.post(f"/format/{job_id}/restore",
                        data={"restore_ids": ["0-9"]})
        assert r.status_code == 200
        assert "back in the deck" in r.text

        dl = client.get(f"/format/{job_id}/download")
        out = Presentation(io.BytesIO(dl.content)).slides[0]
        assert any(s.has_text_frame and s.text_frame.text == "EYEBROW"
                   for s in out.shapes)
    finally:
        web._format_jobs.pop(job_id, None)


def test_restoring_nothing_leaves_the_deck_alone(monkeypatch):
    client = _client(monkeypatch)
    job_id = _job_with_a_removal(client)
    before = web._format_jobs[job_id]["deck"]
    try:
        r = client.post(f"/format/{job_id}/restore", data={})
        assert r.status_code == 200
        assert web._format_jobs[job_id]["deck"] == before
        # and the offer is still on the page rather than reported as done
        assert "Put the ticked pieces back" in r.text
    finally:
        web._format_jobs.pop(job_id, None)


def test_restore_on_an_unknown_job_is_a_clean_404(monkeypatch):
    r = _client(monkeypatch).post("/format/deadbeef/restore", data={})
    assert r.status_code == 404


def test_resubmitting_the_restore_does_not_double_insert(monkeypatch):
    """A refresh or a back button resends the POST. The piece is already back,
    and a second copy of the same shape is not what anyone asked for."""
    client = _client(monkeypatch)
    job_id = _job_with_a_removal(client)
    try:
        client.post(f"/format/{job_id}/restore", data={"restore_ids": ["0-9"]})
        client.post(f"/format/{job_id}/restore", data={"restore_ids": ["0-9"]})
        out = Presentation(io.BytesIO(
            client.get(f"/format/{job_id}/download").content)).slides[0]
        copies = [s for s in out.shapes
                  if s.has_text_frame and s.text_frame.text == "EYEBROW"]
        assert len(copies) == 1
        ids = [s.shape_id for s in out.shapes]
        assert len(ids) == len(set(ids))
    finally:
        web._format_jobs.pop(job_id, None)


def test_saving_a_profile_from_a_master_stores_the_master(monkeypatch):
    """This is what makes a profile applicable at all: the rules alone cannot
    restyle a slide."""
    from qc.profile import PROFILES_DIR
    from qc.store import add_user

    client = _client(monkeypatch)
    add_user("Lead", "lead")
    client.post("/whoami", json={"name": "Lead"})

    master_bytes = _bytes(Presentation())
    client.post("/master", files={"master": ("brand.pptx", master_bytes, "app/x")})
    spec_id = next(iter(web._specs))
    r = client.post(f"/spec/{spec_id}/profile", data={"name": "Applied Client"},
                    follow_redirects=False)
    assert r.status_code == 303
    pid = r.headers["location"].split("/")[2]
    try:
        assert has_master(pid)
        assert load_master(pid) == master_bytes
    finally:
        (PROFILES_DIR / f"{pid}.json").unlink(missing_ok=True)


@needs_powerpoint
def test_end_to_end_through_the_page(monkeypatch):
    from qc.profile import PROFILES_DIR
    from qc.store import add_user

    client = _client(monkeypatch)
    add_user("Lead", "lead")
    client.post("/whoami", json={"name": "Lead"})
    client.post("/master", files={"master": ("brand.pptx",
                                             _bytes(Presentation()), "app/x")})
    spec_id = next(iter(web._specs))
    pid = client.post(f"/spec/{spec_id}/profile", data={"name": "E2E"},
                      follow_redirects=False).headers["location"].split("/")[2]
    try:
        deck = _deck([1, 1])
        r = client.post("/format", data={"profile": pid},
                        files={"deck": ("rough.pptx", _bytes(deck), "app/x")})
        assert r.status_code == 200
        assert "Rebuilt <b>2</b> of <b>2</b> slides" in r.text

        job_id = next(reversed(web._format_jobs))
        dl = client.get(f"/format/{job_id}/download")
        assert dl.status_code == 200
        assert "master applied" in dl.headers["content-disposition"]
        assert len(Presentation(io.BytesIO(dl.content)).slides) == 2
    finally:
        (PROFILES_DIR / f"{pid}.json").unlink(missing_ok=True)


@needs_powerpoint
def test_the_decks_own_master_is_replaced_not_merged():
    """A deck arriving with its own branded master must leave with only the
    applied one. Designs.Load adds the new master alongside the old; the old
    is only removed because every slide moved off it, so this asserts the
    outcome rather than the mechanism."""
    import os
    import tempfile

    import win32com.client
    from qc.unify import com_available  # noqa: F401  (skip guard already ran)

    app = win32com.client.DispatchEx("PowerPoint.Application")
    path = os.path.join(tempfile.gettempdir(), "own_master_fixture.pptx")
    try:
        prs = app.Presentations.Add(WithWindow=0)
        master = prs.SlideMaster
        for i in range(1, master.CustomLayouts.Count + 1):
            master.CustomLayouts(i).Name = "OLDBRAND " + master.CustomLayouts(i).Name
        for i in range(2):
            prs.Slides.AddSlide(i + 1, master.CustomLayouts(2))
        if os.path.exists(path):
            os.remove(path)
        prs.SaveAs(path)
        prs.Close()
    finally:
        try:
            app.Quit()
        except Exception:
            pass

    with open(path, "rb") as fh:
        deck_bytes = fh.read()
    os.remove(path)

    deck = Presentation(io.BytesIO(deck_bytes))
    assert any("OLDBRAND" in l.name
               for l in deck.slide_masters[0].slide_layouts)

    new_master = Presentation()
    result = apply_master(deck_bytes, _bytes(new_master),
                          plan_assignments(deck, _targets(new_master)))
    assert result.fatal is None and not result.errors

    out = Presentation(io.BytesIO(result.deck))
    assert len(out.slide_masters) == 1, "the old master must not survive"
    assert not any("OLDBRAND" in l.name
                   for m in out.slide_masters for l in m.slide_layouts)
