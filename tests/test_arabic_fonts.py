"""Arabic font substitution: fixable with a proposed target (first allowed
complex-script font), never pre-selected - the tick is the designer's
approval (requested 12/08/2026: convert an Arabic deck's stray fonts to
its own Sakkal Majalla convention)."""

import io

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

import qc.modules.font as font_mod
from qc.fixer import apply_fixes, is_fixable, needs_explicit_tick
from qc.profile import Profile
from tests.conftest import save_and_ctx

IN = 914400

PROFILE = Profile({"id": "t", "config": {"font": {"roles": {"body": {
    "latin": ["Arial"], "complex_script": ["Sakkal Majalla", "Dubai"],
    "allowed_weights": ["regular", "bold"]}}}}})


def _arabic_box(slide, cs=None, rtl_el=False):
    tb = slide.shapes.add_textbox(Emu(IN), Emu(IN), Emu(3000000), Emu(500000))
    tb.text_frame.text = "دراسة الجدوى"
    run = tb.text_frame.paragraphs[0].runs[0]
    rpr = run._r.get_or_add_rPr()
    if cs:
        el = rpr.makeelement(qn("a:cs"), {"typeface": cs})
        rpr.append(el)
    if rtl_el:
        rpr.append(rpr.makeelement(qn("a:rtl"), {}))
    return tb, run


def test_out_of_set_cs_typeface_proposes_target_and_fixes(make_prs, tmp_path):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb, run = _arabic_box(slide, cs="Akhbar MT")
    ctx = save_and_ctx(prs, tmp_path, PROFILE)

    recs = [r.to_dict() for r in font_mod.detect(ctx)
            if r.issue_type == "font.family_out_of_set"
            and r.shape_id == str(tb.shape_id)]
    assert len(recs) == 1
    rec = recs[0]
    assert rec["property"] == "rPr.cs.typeface"
    assert rec["old_value"] == "Akhbar MT"
    assert rec["new_value"] == "Sakkal Majalla"
    assert rec["arabic_flag"] is True
    assert is_fixable(rec)
    assert needs_explicit_tick(rec)

    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {rec["record_id"]})
    assert result.applied == 1
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    fixed = next(s for s in out.shapes if str(s.shape_id) == rec["shape_id"])
    r = fixed.text_frame.paragraphs[0].runs[0]
    cs = r._r.find(qn("a:rPr")).find(qn("a:cs"))
    assert cs.get("typeface") == "Sakkal Majalla"
    assert r.text == "دراسة الجدوى"  # text untouched


def test_missing_cs_typeface_fix_respects_rpr_order(make_prs, tmp_path):
    """Inserting a:cs into an rPr that carries a:rtl must place it BEFORE
    rtl (schema order), or PowerPoint flags the file as corrupt."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb, run = _arabic_box(slide, cs=None, rtl_el=True)
    ctx = save_and_ctx(prs, tmp_path, PROFILE)

    recs = [r.to_dict() for r in font_mod.detect(ctx)
            if r.issue_type == "font.cs_typeface_missing"
            and r.shape_id == str(tb.shape_id)]
    assert len(recs) == 1
    assert recs[0]["new_value"] == "Sakkal Majalla"
    assert is_fixable(recs[0])

    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {recs[0]["record_id"]})
    assert result.applied == 1
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    fixed = next(s for s in out.shapes
                 if str(s.shape_id) == recs[0]["shape_id"])
    rpr = fixed.text_frame.paragraphs[0].runs[0]._r.find(qn("a:rPr"))
    tags = [el.tag for el in rpr]
    assert rpr.find(qn("a:cs")).get("typeface") == "Sakkal Majalla"
    assert tags.index(qn("a:cs")) < tags.index(qn("a:rtl"))


def test_arabic_font_fix_is_never_preselected(make_prs, tmp_path):
    from qc.ui import render_report

    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb, _run = _arabic_box(slide, cs="Akhbar MT")
    ctx = save_and_ctx(prs, tmp_path, PROFILE)
    recs = [r.to_dict() for r in font_mod.detect(ctx)
            if r.issue_type == "font.family_out_of_set"]
    manifest = {"deck": "t.pptx", "profile_id": "t", "profile_version": 1,
                "slides": 1, "summary": {"total": len(recs), "errors": 1,
                                         "warnings": 0, "info": 0,
                                         "arabic": 1},
                "records": recs}
    html = render_report(manifest, job_id="j1", can_fix=True)
    i = html.index(f'value="{recs[0]["record_id"]}" form="applyform"')
    tag = html[html.rindex("<input", 0, i):html.index(">", i) + 1]
    assert " checked" not in tag
    assert "explicit approval" in tag
