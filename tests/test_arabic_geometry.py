"""The Arabic guard, scoped (approved 12/08/2026): pure-geometry fixes are
script-neutral and stay fixable on Arabic shapes; anything touching runs or
text keeps the guard. Trigger: an 89-slide Arabic deck came back 72%
blocked, with 120 computed geometry targets refused."""

import io

from pptx import Presentation
from pptx.util import Emu

from qc.fixer import apply_fixes, is_fixable
from tests.conftest import save_and_ctx

IN = 914400


def _rec(issue_type, arabic=True, **kw):
    base = dict(record_id="r1", job_id=None, slide_index=0, shape_id="9",
                shape_path=None, module=issue_type.split(".")[0],
                issue_type=issue_type, property="spPr.xfrm.off.x",
                old_value="1", new_value="2", severity="warning",
                action="flagged", confidence="high", arabic_flag=arabic,
                profile_rule_id=None, message="m", locator=None,
                created_at="")
    base.update(kw)
    return base


def test_geometry_fixes_stay_fixable_on_arabic_shapes():
    for issue in ("margin_alignment.edge_misaligned",
                  "shape_size.size_mismatch",
                  "margin_alignment.panel_row_misaligned",
                  "master_slide.placeholder_geometry_off"):
        assert is_fixable(_rec(issue, new_value="100x100"
                               if "size" in issue else "2")), issue


def test_text_editing_fixes_keep_the_guard():
    """Font substitution moved to tick-to-approve on 12/08/2026 (see
    tests/test_arabic_fonts.py); everything that edits text stays
    guarded."""
    for issue, extra in (
            ("typography.case_inconsistent", {}),
            ("typography.redundant_size_override",
             {"confidence": "deterministic", "locator": "p0/r0"}),
            ("header_footer.text_mismatch", {}),
            ("header_footer.fake_slide_number", {})):
        rec = _rec(issue, **extra)
        assert not is_fixable(rec), issue
        assert is_fixable({**rec, "arabic_flag": False}), issue


def test_arabic_cluster_moves_but_text_is_untouched(make_prs, en_profile,
                                                    tmp_path):
    """End to end: an Arabic label knocked off its cluster gets snapped
    back, and every character survives byte-identically."""
    import qc.modules.margin_alignment as ma

    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    texts = ("التخطيط العمراني", "دراسات الجدوى", "استراتيجية الأصول",
             "النموذج التشغيلي")
    boxes = []
    for i, t in enumerate(texts):
        left = IN if i < 3 else IN + 120000  # the last one is off the line
        tb = slide.shapes.add_textbox(Emu(left), Emu(IN + i * 700000),
                                      Emu(2500000), Emu(400000))
        tb.text_frame.text = t
        boxes.append(tb)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    recs = [r.to_dict() for r in ma.detect(ctx)
            if r.issue_type == "margin_alignment.edge_misaligned"
            and r.shape_id == str(boxes[3].shape_id)]
    assert len(recs) == 1
    rec = recs[0]
    assert rec["arabic_flag"] is True
    assert "manual review" not in rec["message"]
    assert "text untouched" in rec["message"]
    assert is_fixable(rec)

    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {rec["record_id"]})
    assert result.applied == 1
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    by_id = {str(s.shape_id): s for s in out.shapes}
    assert by_id[rec["shape_id"]].left == IN          # snapped to the line
    for tb, original in zip(boxes, texts):            # every glyph intact
        assert by_id[str(tb.shape_id)].text_frame.text == original
