"""Arabic decks: mirrored margins, titles that start at the top margin,
collections that stay composed, and removals a designer can undo.

Four asks from one review of a leadership-quotes slide (design lead,
20/08/2026), each with the same shape of answer: the tool must respect what the
page already states, and where it decides something for the designer, the
designer must be able to decide back.
"""

import copy
import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Pt

import qc.modules.font as font_module
import qc.modules.margin_alignment as ma
from qc.fixer import apply_fixes, is_fixable
from qc.migrate import migrate_deck, restore_shapes
from qc.profile import Profile
from qc.stylespec import dominant_master, infer_grid

from tests.conftest import save_and_ctx

IN = 914400
MM = 36000
BLANK = 6

_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
_P15 = "http://schemas.microsoft.com/office/powerpoint/2012/main"

AR_TITLE = "وعدنا لكم"
AR_BODY = "يطيب لنا العمل والتعاون مع المركز الوطني لإدارة النفايات"


def _plant_guides(master, vertical_in=(0.60, 12.73), horizontal_in=(0.42, 7.05)):
    """Guides as desktop PowerPoint stores them: eighths of a point from the
    top-left edge, orient omitted for vertical."""
    import itertools

    from lxml import etree

    ext_lst = etree.SubElement(master._element, f"{_P}extLst")
    ext = etree.SubElement(ext_lst, f"{_P}ext")
    ext.set("uri", "{GUIDES}")
    lst = etree.SubElement(ext, f"{{{_P15}}}sldGuideLst")
    gid = itertools.count(1)
    for pos, horz in ([(v, False) for v in vertical_in]
                      + [(h, True) for h in horizontal_in]):
        g = etree.SubElement(lst, f"{{{_P15}}}guide")
        g.set("id", str(next(gid)))
        if horz:
            g.set("orient", "horz")
        g.set("pos", str(int(pos * 72 * 8)))


_TITLE_BOX = (0.48, 0.42, 12.40, 0.92)
_SUBTITLE_BOX = (0.48, 1.40, 12.40, 0.35)


def _header_deck(*, arabic: bool, card_left=1.0, card_width=3.5,
                 anchor=MSO_ANCHOR.BOTTOM):
    """A guided master, a bottom-anchored title, and three cards indented from
    both margins so the binding edge is visible in the result."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    _plant_guides(prs.slide_masters[0])
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)
            ph.text_frame.vertical_anchor = anchor

    def tb(x, y, w, h, text, size=None):
        shape = slide.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                                         Emu(int(w * IN)), Emu(int(h * IN)))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = text
        if size:
            run.font.size = Pt(size)
        return shape

    tb(0.6, 0.55, 9.0, 0.6, AR_TITLE if arabic else "The Heading", 28)
    tb(0.6, 1.20, 9.0, 0.35,
       "سنطبق ممارسات" if arabic else "The standfirst", 16)
    for i in range(3):
        tb(card_left + i * (card_width + 0.3), 2.2, card_width, 2.0,
           (AR_BODY if arabic else f"Card {i + 1}"), 12)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _margins(deck_bytes):
    prs = Presentation(io.BytesIO(deck_bytes))
    grid = infer_grid(prs, dominant_master(prs))
    m = grid["margins_emu"]
    return m["left"], prs.slide_width - m["right"]


def _card_bounds(deck_bytes):
    slide = Presentation(io.BytesIO(deck_bytes)).slides[0]
    boxes = [s for s in slide.shapes
             if not s.is_placeholder and s.has_text_frame
             and s.text_frame.text.strip()]
    return (min(s.left for s in boxes),
            max(s.left + s.width for s in boxes))


# ------------------------------------------------- 1. mirrored margins


def test_an_arabic_block_binds_to_the_right_margin():
    """Arabic reads right to left, so the right margin is where its block
    starts. Binding the left edge leaves the reading edge ragged, which is the
    first thing a reader of the language sees."""
    source = _header_deck(arabic=True)
    _left, right = _margins(source)
    out, changes = migrate_deck(source)
    got_left, got_right = _card_bounds(out)

    assert got_right == right, "the block should sit on the right guide"
    assert got_left > _margins(source)[0], "and not be dragged to the left one"
    moved = [c for c in changes if c.action == "content block moved"]
    assert moved and "right margin" in moved[0].detail
    assert "right to left" in moved[0].detail


def test_an_english_block_still_binds_to_the_left_margin():
    """The mirror is per SLIDE and driven by the script, so nothing changes for
    a latin deck."""
    source = _header_deck(arabic=False)
    left, _right = _margins(source)
    out, changes = migrate_deck(source)

    assert _card_bounds(out)[0] == left
    moved = [c for c in changes if c.action == "content block moved"]
    assert moved and "left margin" in moved[0].detail


def _draw_space(master, left_in, top_in, w_in, h_in):
    """The rectangle a designer draws on the master to state where content may
    live. python-pptx cannot add shapes to a master, so it is planted."""
    from pptx.oxml.shapes.autoshape import CT_Shape

    sp = CT_Shape.new_autoshape_sp(950, "Presentation space", "rect",
                                   int(left_in * IN), int(top_in * IN),
                                   int(w_in * IN), int(h_in * IN))
    master.shapes._spTree.insert_element_before(sp, "p:extLst")
    shape = next(s for s in master.shapes if s.name == "Presentation space")
    shape.fill.background()
    shape.line.fill.background()


def test_the_block_binds_to_the_presentation_space_not_the_guides():
    """A master can carry several sets of margins, so a rectangle the designer
    draws and names is the safer statement of where content lives (design lead,
    21/08/2026). It outranks the guides, and the block seats on IT."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    master = prs.slide_masters[0]
    _plant_guides(master)                      # margins at 0.60in / 12.73in
    _draw_space(master, 1.20, 0.40, 10.5, 6.4)  # the designer's real frame
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)

    def tb(x, y, w, h, text, size):
        shape = slide.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                                         Emu(int(w * IN)), Emu(int(h * IN)))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(size)

    tb(0.6, 0.55, 9.0, 0.6, "The Heading", 28)
    tb(0.6, 1.20, 9.0, 0.35, "The standfirst", 16)
    for i in range(3):
        tb(2.0 + i * 3.0, 2.4, 2.5, 2.0, f"Card {i + 1}", 12)
    buf = io.BytesIO()
    prs.save(buf)
    source = buf.getvalue()

    out, changes = migrate_deck(source)
    left, _right = _card_bounds(out)
    assert left == int(1.20 * IN), "seated on the presentation space, not 0.60in"
    assert [c for c in changes if c.action == "content block moved"]


def test_a_stated_frame_binds_even_when_the_block_overflows():
    """The root cause of "the whole presentation is shifted, not following the
    presentation space" (design lead, 21/08/2026): "the top binds" was tied to
    the guide-pair BAND, and a master stating its frame with a rectangle has no
    band - so the bottom clamp applied, and every slide whose content is too
    tall kept the position it arrived with. Deck-wide, and nothing to do with
    Arabic."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    _draw_space(prs.slide_masters[0], 0.48, 1.90, 12.4, 4.9)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Heading"
    slide.placeholders[1].text = "x"
    # a body block from 0.60in to 7.10in: taller than the 4.9in frame
    tall = slide.shapes.add_textbox(Emu(int(1 * IN)), Emu(int(0.60 * IN)),
                                    Emu(int(8 * IN)), Emu(int(6.5 * IN)))
    tall.text_frame.paragraphs[0].add_run().text = "a very tall block"
    buf = io.BytesIO()
    prs.save(buf)

    out, changes = migrate_deck(buf.getvalue())
    slide = Presentation(io.BytesIO(out)).slides[0]
    block = next(s for s in slide.shapes
                 if s.has_text_frame
                 and s.text_frame.text == "a very tall block")

    assert block.top == Emu(int(1.90 * IN)), "seated on the presentation space"
    moved = [c for c in changes if c.action == "content block moved"]
    assert moved and "body now starts at 1.90in" in moved[0].detail
    # and the overflow it creates is stated, not hidden
    assert [c for c in changes if c.action == "content does not fit"]


def test_an_unstated_frame_keeps_the_clamp_and_says_so():
    """The clamp survives where nothing is stated: with no rectangle and no
    guides, the frame is an inference from placeholder extents and does not get
    to force an overflow. The block then lands short of the line, and the report
    has to say why rather than leave a designer comparing the deck against a
    frame that was never read."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Heading"
    slide.placeholders[1].text = "x"
    # tall enough that the bottom clamp bites, but not a page-deep panel
    tall = slide.shapes.add_textbox(Emu(int(1 * IN)), Emu(int(0.60 * IN)),
                                    Emu(int(8 * IN)), Emu(int(6.0 * IN)))
    tall.text_frame.paragraphs[0].add_run().text = "a very tall block"
    buf = io.BytesIO()
    prs.save(buf)

    out, changes = migrate_deck(buf.getvalue())
    slide = Presentation(io.BytesIO(out)).slides[0]
    block = next(s for s in slide.shapes
                 if s.has_text_frame
                 and s.text_frame.text == "a very tall block")

    assert block.top + block.height <= prs.slide_height, \
        "an inferred frame must not push content off the canvas"
    stuck = [c for c in changes
             if c.action == "body not seated on a stated frame"]
    assert stuck, "the report has to explain the short landing"
    assert "states no content frame" in stuck[0].detail
    assert "presentation-space rectangle" in stuck[0].detail


# ------------------------------------------------- 1b. reading direction


def _titled_deck(*, arabic=True, state_direction=True):
    """A deck whose header text sits in free boxes, the way a converted deck
    arrives, with the direction its own paragraphs state."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)

    def tb(y, text, size, rtl=False):
        shape = slide.shapes.add_textbox(Emu(int(0.6 * IN)), Emu(int(y * IN)),
                                         Emu(int(9 * IN)), Emu(int(0.6 * IN)))
        para = shape.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        if rtl:
            pPr = para._p.get_or_add_pPr()
            pPr.set("rtl", "1")
            pPr.set("algn", "r")
        return shape

    tb(0.55, AR_TITLE if arabic else "Our promise", 28,
       rtl=arabic and state_direction)
    tb(1.20, "سنطبق ممارسات" if arabic else "The standfirst", 16,
       rtl=arabic and state_direction)
    for i in range(3):
        tb(2.5 + i * 0.7, AR_BODY if arabic else "Body copy", 12)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _directions(deck_bytes):
    """(rtl, algn) of every header placeholder paragraph in the output."""
    from pptx.oxml.ns import qn

    slide = Presentation(io.BytesIO(deck_bytes)).slides[0]
    out = []
    for ph in slide.placeholders:
        if not str(ph.placeholder_format.type).startswith(
                ("TITLE", "CENTER_TITLE", "SUBTITLE")):
            continue
        for para in ph.text_frame.paragraphs:
            pPr = para._p.find(qn("a:pPr"))
            out.append((pPr.get("rtl") if pPr is not None else None,
                        pPr.get("algn") if pPr is not None else None))
    return out


def test_arabic_text_keeps_reading_right_to_left_in_the_placeholder():
    """The bug behind "the title is not right to left (still english)": the move
    into a placeholder carried the words and dropped the DIRECTION, so Arabic
    inherited the English master's left-to-right alignment. Which way a language
    reads is not a house style the master gets to overrule."""
    out, changes = migrate_deck(_titled_deck())

    assert _directions(out) == [("1", "r"), ("1", "r")]
    placed = [c for c in changes if "into placeholder" in c.action]
    assert placed and "kept right-to-left" in placed[0].detail


def test_arabic_is_marked_rtl_even_when_the_deck_never_said_so():
    """Converted decks routinely carry Arabic in paragraphs with no direction at
    all. Inheriting an English master then reverses the reading edge, so the
    script itself is the fallback answer."""
    out, _changes = migrate_deck(_titled_deck(state_direction=False))
    assert _directions(out) == [("1", "r"), ("1", "r")]


def test_english_headers_are_left_exactly_as_they_were():
    out, changes = migrate_deck(_titled_deck(arabic=False))
    assert _directions(out) == [(None, None), (None, None)]
    assert not [c for c in changes if "right-to-left" in c.detail]


# --------------------------------------- 2. the master's anchor is followed


def test_the_masters_own_vertical_anchor_is_left_alone():
    """A master that hangs its title at the bottom of the box is stating a
    design: the heading pairs with the subtitle under it. Hoisting the text to
    the top of the box put the heading on the top margin and moved the empty
    space to the other side of it, which reads worse (design lead, 21/08/2026,
    reversing the 20/08 decision after seeing both).

    Where a slide's own title BOX sits somewhere other than the master says,
    that is geometry and master_slide.placeholder_geometry_off owns it."""
    out, changes = migrate_deck(_header_deck(arabic=True,
                                             anchor=MSO_ANCHOR.BOTTOM))
    slide = Presentation(io.BytesIO(out)).slides[0]
    title = next(ph for ph in slide.placeholders
                 if str(ph.placeholder_format.type).startswith(
                     ("TITLE", "CENTER_TITLE")))

    assert title.text_frame.vertical_anchor == MSO_ANCHOR.BOTTOM
    assert title.top == Emu(int(_TITLE_BOX[1] * IN)), "the box must not move"
    assert not [c for c in changes if "anchor" in c.detail.lower()]


def test_nothing_reports_a_title_anchor_any_more(make_prs, en_profile,
                                                 tmp_path):
    """The audit rule went with the override: an anchor the master states is
    not a defect, so there is nothing to flag and nothing to tick."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
    title = slide.shapes.title
    title.text_frame.text = "Annual review"
    title.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    assert not [r for r in font_module.detect(ctx)
                if "anchor" in r.issue_type]


# ----------------------------------------------------- 3. collections


def _quote_column(prs, off_left=None):
    """Three photo collections stacked in a column - photo, a corner rule drawn
    across it, a caption under it - with one photo's left edge off the line."""
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    members = []
    for i in range(3):
        left = (off_left if (i == 1 and off_left is not None) else 20) * MM
        top = (20 + i * 45) * MM
        photo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(left)),
                                       Emu(int(top)), Emu(40 * MM),
                                       Emu(30 * MM))
        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Emu(int(left - 3 * MM)),
                                      Emu(int(top + 6 * MM)), Emu(10 * MM),
                                      Emu(2 * MM))
        caption = slide.shapes.add_textbox(Emu(int(left)),
                                           Emu(int(top + 32 * MM)),
                                           Emu(40 * MM), Emu(6 * MM))
        caption.text_frame.text = f"د. شهاب البرعي {i}"
        members.append((photo, rule, caption))
    return slide, members


def test_a_satellite_is_not_judged_on_its_own(make_prs, en_profile, tmp_path):
    """The corner rules and captions ride their photos, so they never join the
    absolute clusters: only the photo column is judged, and one finding is
    raised for the collection rather than three for its parts."""
    prs = make_prs()
    _slide, members = _quote_column(prs, off_left=22)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    recs = [r.to_dict() for r in ma.detect(ctx)
            if r.issue_type == "margin_alignment.edge_misaligned"]
    riders = {str(m[1].shape_id) for m in members} | \
             {str(m[2].shape_id) for m in members}
    assert recs, "the photo off the column line should still be flagged"
    assert not [r for r in recs if r["shape_id"] in riders], \
        "a satellite must not be nudged on its own"


def test_a_fix_moves_the_whole_collection(make_prs, en_profile, tmp_path):
    """The composition inside a collection is the design: whatever moves, its
    corner rule and caption move by the same delta."""
    prs = make_prs()
    _slide, members = _quote_column(prs, off_left=22)
    photo, rule, caption = members[1]
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    recs = [r.to_dict() for r in ma.detect(ctx)
            if r.issue_type == "margin_alignment.edge_misaligned"
            and r.shape_id == str(photo.shape_id)]
    assert len(recs) == 1 and is_fixable(recs[0])
    before = {s.shape_id: s.left for s in (photo, rule, caption)}

    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {recs[0]["record_id"]})
    assert result.applied == 1, [o.reason for o in result.outcomes]
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    after = {s.shape_id: s.left for s in out.shapes}
    delta = after[photo.shape_id] - before[photo.shape_id]

    assert delta, "the photo should have moved onto the column line"
    assert after[rule.shape_id] == before[rule.shape_id] + delta, \
        "the corner rule was left behind: the collection came apart"
    assert after[caption.shape_id] == before[caption.shape_id] + delta


# -------------------------------------------------- 4. bringing it back


def test_a_graphic_in_the_header_band_is_never_swept_as_unplaced_text():
    """The sweep removes header TEXT the master has no slot for. A corner rule
    or mark carries no text and is part of a composition; deleting it and
    calling it 'unplaced text' was a defect however it was phrased."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    _plant_guides(prs.slide_masters[0])
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)
    mark = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(0.6 * IN)),
                                  Emu(int(0.5 * IN)), Emu(int(0.4 * IN)),
                                  Emu(int(0.3 * IN)))
    mark.name = "CornerRule"
    tb = slide.shapes.add_textbox(Emu(int(0.6 * IN)), Emu(int(2.5 * IN)),
                                 Emu(int(4 * IN)), Emu(int(2 * IN)))
    tb.text_frame.paragraphs[0].add_run().text = AR_BODY
    buf = io.BytesIO()
    prs.save(buf)

    out, changes = migrate_deck(buf.getvalue())
    names = [s.name for s in Presentation(io.BytesIO(out)).slides[0].shapes]
    assert "CornerRule" in names
    assert not [c for c in changes
                if c.action == "removed unplaced text" and not c.removed_text]


AR_EYEBROW = "قطاع الاستدامة"


def _deck_with_an_unplaceable_header_line():
    """Three lines of header text over a master with two header slots: the
    heading and the standfirst are placed, and the eyebrow has nowhere to go."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    _plant_guides(prs.slide_masters[0])
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)

    def tb(x, y, w, h, text, size):
        shape = slide.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                                         Emu(int(w * IN)), Emu(int(h * IN)))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(size)
        return shape

    tb(0.6, 0.45, 3.0, 0.28, AR_EYEBROW, 11)
    tb(0.6, 0.60, 9.0, 0.60, AR_TITLE, 28)
    tb(0.6, 1.25, 9.0, 0.35, "سنطبق ممارسات", 16)
    for i in range(3):
        tb(0.6 + i * 4, 2.5, 3.5, 2.0, AR_BODY, 12)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _texts(deck_bytes):
    slide = Presentation(io.BytesIO(deck_bytes)).slides[0]
    return [s.text_frame.text for s in slide.shapes
            if getattr(s, "has_text_frame", False)]


def test_a_removed_piece_carries_its_own_xml_and_comes_back_whole():
    """Listing the text was half an answer: a designer had to retype it and
    place it by eye. The element is kept, so the piece itself comes back whole -
    same words, same box size, same formatting. Only its vertical position may
    change, and only to a slot where it covers nothing."""
    source = _deck_with_an_unplaceable_header_line()
    before = Presentation(io.BytesIO(source)).slides[0]
    original = next(s for s in before.shapes
                    if s.has_text_frame and s.text_frame.text == AR_EYEBROW)

    out, changes = migrate_deck(source)
    removed = [c for c in changes if c.removed_text == AR_EYEBROW]
    assert removed, "the eyebrow has no slot in the master and is swept"
    assert AR_EYEBROW not in _texts(out)
    assert removed[0].restore_id and removed[0].removed_xml

    back, outcomes = restore_shapes(out, [{"slide_index": c.slide_index,
                                           "removed_xml": c.removed_xml,
                                           "restore_id": c.restore_id}
                                          for c in removed])
    assert [o["restore_id"] for o in outcomes] == [removed[0].restore_id]
    slide = Presentation(io.BytesIO(back)).slides[0]
    got = next(s for s in slide.shapes
               if s.has_text_frame and s.text_frame.text == AR_EYEBROW)
    assert (got.left, got.width, got.height) == \
           (original.left, original.width, original.height)
    assert got.name.startswith("RESTORED"), "findable in the selection pane"
    assert outcomes[0]["detail"]


def _deck_with_an_extension_list():
    """A slide whose shape tree ends with p:extLst, which is what PowerPoint
    writes on any slide it has touched - and applying a master hands every
    slide to PowerPoint."""
    from lxml import etree

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    box = slide.shapes.add_textbox(Emu(IN), Emu(IN), Emu(IN), Emu(IN // 2))
    box.text_frame.text = "PIECE"
    xml = etree.tostring(box._element, encoding="unicode")
    box._element.getparent().remove(box._element)
    ext = etree.SubElement(slide.shapes._spTree, f"{_P}extLst")
    etree.SubElement(ext, f"{_P}ext").set("uri", "{TEST}")
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), xml


def _child_tags(deck_bytes):
    slide = Presentation(io.BytesIO(deck_bytes)).slides[0]
    return [k.tag.split("}")[-1] for k in slide.shapes._spTree]


def test_a_restored_piece_goes_in_before_the_extension_list():
    """p:extLst must be the LAST child of a shape tree. Appending after it
    produced a deck PowerPoint offered to repair, which is how a restore took
    a whole presentation with it (bug, 20/08/2026)."""
    deck, xml = _deck_with_an_extension_list()
    back, outcomes = restore_shapes(deck, [{"slide_index": 0,
                                            "removed_xml": xml,
                                            "restore_id": "0-2"}])

    assert [o["restore_id"] for o in outcomes] == ["0-2"]
    tags = _child_tags(back)
    assert tags[-1] == "extLst", tags
    assert "sp" in tags


def test_restoring_the_same_piece_twice_cannot_collide():
    """Shape ids are unique per slide and a removed piece's id is not reserved
    for it, so the restored subtree is renumbered. A resubmitted form used to
    leave two shapes sharing one id, which PowerPoint reads as damage."""
    deck, xml = _deck_with_an_extension_list()
    item = {"slide_index": 0, "removed_xml": xml, "restore_id": "0-2"}
    once, _ = restore_shapes(deck, [item])
    twice, _ = restore_shapes(once, [item])

    ids = [s.shape_id for s in Presentation(io.BytesIO(twice)).slides[0].shapes]
    assert len(ids) == len(set(ids)), ids
    assert _child_tags(twice)[-1] == "extLst"


def test_restoring_nothing_changes_nothing():
    out, _changes = migrate_deck(_deck_with_an_unplaceable_header_line())
    back, outcomes = restore_shapes(out, [])
    assert outcomes == []
    assert _texts(back) == _texts(out)


def test_a_restored_piece_goes_back_where_it_was_and_says_what_it_covers():
    """Its own position is the point: a restore is an undo, not a re-layout
    (design lead, 21/08/2026). Where the master has since filled that spot the
    piece prints over it, and the report names what it covers rather than
    hunting for empty space three inches down the slide."""
    source = _deck_with_an_unplaceable_header_line()
    out, changes = migrate_deck(source)
    removed = [c for c in changes if c.removed_xml]
    back, outcomes = restore_shapes(out, [{"slide_index": c.slide_index,
                                           "removed_xml": c.removed_xml,
                                           "restore_id": c.restore_id}
                                          for c in removed])
    slide = Presentation(io.BytesIO(back)).slides[0]
    piece = next(s for s in slide.shapes if s.name.startswith("RESTORED"))
    original = next(s for s in Presentation(io.BytesIO(source)).slides[0].shapes
                    if s.has_text_frame and s.text_frame.text == AR_EYEBROW)

    assert (piece.left, piece.top) == (original.left, original.top)
    # the master's title now occupies that spot, and the report says so
    assert outcomes[0]["covers"], "the overlap must be reported, not hidden"
    assert "printing over" in outcomes[0]["detail"]
