"""Assistant triage: aggregation, question generation (offline fallback and
API-failure fallback), action validation, profile updates, and the web routes.
All tests are hermetic - the Claude API is never called."""

import io
import json
import shutil

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Emu

import qc.assist as assist
from qc.web import app
from tests.conftest import job_id_of

client = TestClient(app)


def _rec(issue_type, old_value=None, rule=None, slide=0, action="flagged"):
    return {"record_id": "x", "issue_type": issue_type, "old_value": old_value,
            "profile_rule_id": rule, "slide_index": slide, "action": action}


def _manifest(records):
    return {"records": records, "summary": {}}


PROFILE_CFG = {
    "color_palette": {"named_colors": [{"name": "navy", "hex": "1F4E79"}]},
    "geometry": {"safe_zone_margins_emu": {"left": 457200, "top": 274638,
                                           "right": 457200, "bottom": 365125}},
}


def test_aggregate_groups_and_thresholds():
    records = (
        [_rec("color_palette.off_palette_rgb", "B02E27", slide=i % 2)
         for i in range(4)]
        + [_rec("color_palette.off_palette_rgb", "AAAAAA")] * 2   # below min
        + [_rec("font.family_out_of_set", "Comic Sans MS",
                rule="font.roles.title.latin")] * 3
        + [_rec("font.family_out_of_set", "Dubai",
                rule="font.roles.body.complex_script")] * 2      # below min
        + [_rec("font.family_out_of_set", "NoRule")] * 5          # no role
    )
    agg = assist.aggregate(_manifest(records), None, PROFILE_CFG)
    assert agg["colors"] == [{"hex": "B02E27", "count": 4, "slides": 2}]
    assert agg["fonts"] == [{"family": "Comic Sans MS", "role": "title",
                             "script": "latin", "count": 3}]
    assert agg["margins"] is None
    assert agg["profile"]["palette"] == ["1F4E79"]


def test_margin_proposal_from_deck_edges():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(25):  # content consistently at ~0.3in margins
        tb = slide.shapes.add_textbox(Emu(274320), Emu(274320 + i * 150000),
                                      Emu(1000000), Emu(100000))
        tb.text_frame.text = "content"  # margins are learned from text
    buf = io.BytesIO()
    prs.save(buf)
    records = [_rec("margin_alignment.outside_safe_zone")] * 12
    agg = assist.aggregate(_manifest(records), buf.getvalue(), PROFILE_CFG)
    assert agg["margins"] is not None
    assert agg["margins"]["breaches"] == 12
    assert agg["margins"]["proposed"]["left"] <= 274320


def test_fallback_questions_offline(monkeypatch):
    monkeypatch.setattr(assist, "api_configured", lambda: False)
    records = ([_rec("color_palette.off_palette_rgb", "B02E27")] * 4
               + [_rec("font.family_out_of_set", "Comic Sans MS",
                       rule="font.roles.title.latin")] * 3)
    agg = assist.aggregate(_manifest(records), None, PROFILE_CFG)
    questions, source = assist.generate_questions(agg)
    assert source.startswith("fallback")
    assert len(questions) == 2
    for q in questions:
        assert q["id"] and q["question"] and q["impact"]
    kinds = {q["action"]["type"] for q in questions}
    assert kinds == {"add_color", "add_font"}


def test_api_failure_falls_back(monkeypatch):
    monkeypatch.setattr(assist, "api_configured", lambda: True)

    def boom(agg):
        raise RuntimeError("no network")

    monkeypatch.setattr(assist, "_ask_claude", boom)
    records = [_rec("color_palette.off_palette_rgb", "B02E27")] * 4
    agg = assist.aggregate(_manifest(records), None, PROFILE_CFG)
    questions, source = assist.generate_questions(agg)
    assert questions and source == "fallback (RuntimeError)"


def test_validation_rejects_invented_values():
    agg = {"colors": [{"hex": "B02E27", "count": 4, "slides": 2}],
           "fonts": [], "margins": {"breaches": 12,
                                    "proposed": {"left": 1, "top": 2,
                                                 "right": 3, "bottom": 4}}}
    # a hex the deck never used is dropped
    assert assist._validate(
        {"action": {"type": "add_color", "hex": "123456", "name": "x"}},
        agg) is None
    # margin numbers are always replaced with the locally computed proposal
    q = assist._validate(
        {"action": {"type": "set_margins", "left": 999999, "top": 0,
                    "right": 0, "bottom": 0}}, agg)
    assert q["action"]["left"] == 1 and q["action"]["bottom"] == 4
    # unknown action types are dropped
    assert assist._validate({"action": {"type": "delete_everything"}},
                            agg) is None


@pytest.fixture()
def profiles_dir(monkeypatch, tmp_path):
    import qc.profile as profile_mod

    src = profile_mod.PROFILES_DIR / "prezlab_en.json"
    pdir = tmp_path / "profiles"
    pdir.mkdir()
    shutil.copy(src, pdir / "prezlab_en.json")
    monkeypatch.setattr(profile_mod, "PROFILES_DIR", pdir)
    return pdir


def test_apply_actions_updates_profile(profiles_dir):
    before = json.loads((profiles_dir / "prezlab_en.json").read_text("utf-8"))
    result = assist.apply_actions("prezlab_en", [
        {"type": "add_color", "hex": "B02E27", "name": "deep red"},
        {"type": "add_font", "family": "Figtree", "role": "title",
         "script": "latin"},
        {"type": "set_margins", "left": 274320, "top": 274320,
         "right": 274320, "bottom": 274320},
    ], editor="Sanad")
    assert result["version"] == before["version"] + 1
    assert len(result["applied"]) == 3

    after = json.loads((profiles_dir / "prezlab_en.json").read_text("utf-8"))
    hexes = [c["hex"] for c in
             after["config"]["color_palette"]["named_colors"]]
    assert "B02E27" in hexes
    assert "Figtree" in after["config"]["font"]["roles"]["title"]["latin"]
    assert after["config"]["geometry"]["safe_zone_margins_emu"]["left"] == 274320
    assert "assistant" in after["owner"]


def _audit_job(fixtures_dir, profile="prezlab_en"):
    deck = (fixtures_dir / "mixed_layouts.pptx").read_bytes()
    r = client.post("/audit",
                    files={"deck": ("mixed_layouts.pptx", deck,
                                    "application/octet-stream")},
                    data={"profile": profile})
    assert r.status_code == 200
    job = job_id_of(r)
    return job, client.get(f"/audit/{job}").text


def test_assist_routes_end_to_end(fixtures_dir, profiles_dir, monkeypatch):
    monkeypatch.setattr(assist, "api_configured", lambda: False)
    from qc.store import add_user

    add_user("Razan", "lead")

    job_id, report_html = _audit_job(fixtures_dir)
    assert 'id="askassist"' in report_html  # panel offered on the report

    r = client.post(f"/assist/{job_id}")
    assert r.status_code == 200
    data = r.json()
    assert data["source"].startswith("fallback")
    assert data["profile"] == "prezlab_en"
    ids = [q["id"] for q in data["questions"]]
    assert ids  # mixed_layouts has recurring off-set fonts
    assert all("action" not in q for q in data["questions"])  # ids only

    # unauthenticated apply is rejected; the profile is untouched
    r = client.post(f"/assist/{job_id}/apply", data={"accepted": ids})
    assert r.status_code == 403

    lead = TestClient(app, cookies={"qc_user": "Razan"})
    r = lead.post(f"/assist/{job_id}/apply", data={"accepted": ids})
    assert r.status_code == 200, r.text
    body = r.json()
    assert body["applied"] and body["version"] >= 2

    # re-audit runs against the updated profile and renders a report
    r = lead.post(f"/reaudit/{job_id}")
    assert r.status_code == 200
    assert "Re-audited with the updated profile" in r.text


def test_assist_rejects_self_profile(fixtures_dir):
    job_id, report_html = _audit_job(fixtures_dir, profile="__self__")
    assert 'id="askassist"' not in report_html  # panel not offered
    r = client.post(f"/assist/{job_id}")
    assert r.status_code == 400


def test_assist_unknown_job_404():
    assert client.post("/assist/deadbeef").status_code == 404
