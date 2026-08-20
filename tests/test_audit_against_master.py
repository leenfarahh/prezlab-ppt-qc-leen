"""Auditing a deck against a master submitted in the same request.

This is the plan's flow with Stage 2 still absent: the master's declared
system becomes the rules, and the content deck is checked against them.
"""

import io

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

from qc import web
from qc.stylespec import extract_style_spec

IN = 914400


def _bytes(prs) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _master() -> bytes:
    return _bytes(Presentation())


def _deck_with_off_theme_colour() -> bytes:
    """A deck whose shapes sit far off the master's theme palette, so the
    audit has something real to say."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for _ in range(3):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(IN), Emu(IN),
                                    Emu(IN), Emu(IN))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor.from_string("FF00AA")
    tb = slide.shapes.add_textbox(Emu(IN), Emu(3 * IN), Emu(4 * IN), Emu(IN))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "Off-brand body copy"
    run.font.name = "Comic Sans MS"
    run.font.size = Pt(18)
    return _bytes(prs)


def _client(monkeypatch):
    monkeypatch.setattr(web, "AUTH_REQUIRED", False)
    web.app.state.auth_required = False
    return TestClient(web.app)


# ------------------------------------------------------------------- picker


def test_master_option_offered_on_the_audit_page(monkeypatch):
    r = _client(monkeypatch).get("/")
    assert r.status_code == 200
    assert "__master__" in r.text
    assert "Match a master I&#x27;ll upload" in r.text or "Match a master" in r.text
    # The extra drop zone ships hidden and is revealed by the radio.
    assert 'id="mdrop"' in r.text
    assert "syncRuleSource" in r.text


# --------------------------------------------------------------- happy path


def test_deck_audited_against_an_uploaded_master(monkeypatch):
    client = _client(monkeypatch)
    r = client.post("/audit", data={"profile": "__master__"}, files={
        "deck": ("rough.pptx", _deck_with_off_theme_colour(), "app/x"),
        "master": ("brand.pptx", _master(), "app/x"),
    })
    assert r.status_code == 200
    # The report names the master the rules came from, not a saved profile.
    assert "master:brand" in r.text
    assert "rough.pptx" in r.text

    job = web._jobs[next(reversed(web._jobs))]
    assert job["profile"] == "__master__"
    assert job["profile_obj"].id == "master:brand"
    # The spec is kept so the report can offer it; the master bytes are not.
    assert job["master_spec"]["theme"]["colors"]["accent1"]
    assert "master" not in job


def test_findings_come_from_the_master_rules(monkeypatch):
    """The audit has to actually use the derived palette and fonts, not fall
    back to a default profile."""
    client = _client(monkeypatch)
    r = client.post("/audit", data={"profile": "__master__"}, files={
        "deck": ("rough.pptx", _deck_with_off_theme_colour(), "app/x"),
        "master": ("brand.pptx", _master(), "app/x"),
    })
    job = web._jobs[next(reversed(web._jobs))]
    issues = job["manifest"]["summary"]["by_issue_type"]

    assert any(k.startswith("color_palette.") for k in issues)
    assert any(k.startswith("font.") for k in issues)

    palette = job["profile_obj"].get("color_palette.named_colors")
    theme = extract_style_spec(Presentation())["theme"]["colors"]
    assert {c["hex"] for c in palette} <= set(theme.values())


# ---------------------------------------------------------------- bad input


def test_master_option_without_a_master_file_asks_for_one(monkeypatch):
    client = _client(monkeypatch)
    r = client.post("/audit", data={"profile": "__master__"},
                    files={"deck": ("rough.pptx", _master(), "app/x")})
    assert r.status_code == 400
    assert "Pick the master" in r.text


def test_non_pptx_master_is_refused(monkeypatch):
    client = _client(monkeypatch)
    r = client.post("/audit", data={"profile": "__master__"}, files={
        "deck": ("rough.pptx", _master(), "app/x"),
        "master": ("brand.txt", b"not a deck", "text/plain"),
    })
    assert r.status_code == 400
    assert "must be a .pptx" in r.text


def test_unreadable_master_fails_with_a_message(monkeypatch):
    client = _client(monkeypatch)
    r = client.post("/audit", data={"profile": "__master__"}, files={
        "deck": ("rough.pptx", _master(), "app/x"),
        "master": ("brand.pptx", b"PK\x03\x04garbage", "app/x"),
    })
    assert r.status_code == 400
    assert "Could not read that master" in r.text


def test_a_master_file_is_ignored_when_another_rule_source_is_picked(monkeypatch):
    """Leaving a stale file in the hidden input must not change the rules."""
    client = _client(monkeypatch)
    r = client.post("/audit", data={"profile": "__self__"}, files={
        "deck": ("rough.pptx", _deck_with_off_theme_colour(), "app/x"),
        "master": ("brand.pptx", _master(), "app/x"),
    })
    assert r.status_code == 200
    job = web._jobs[next(reversed(web._jobs))]
    assert job["profile"] == "__self__"
    assert job["master_spec"] is None


# -------------------------------------------------- downstream consequences


def test_reaudit_reuses_the_derived_profile(monkeypatch):
    """The master's bytes are gone by then, so re-resolving would fail."""
    client = _client(monkeypatch)
    client.post("/audit", data={"profile": "__master__"}, files={
        "deck": ("rough.pptx", _deck_with_off_theme_colour(), "app/x"),
        "master": ("brand.pptx", _master(), "app/x"),
    })
    job_id = next(reversed(web._jobs))

    r = client.post(f"/reaudit/{job_id}", follow_redirects=False)
    assert r.status_code == 200
    new_job = web._jobs[next(reversed(web._jobs))]
    assert new_job["profile_obj"].id == "master:brand"
    assert new_job["master_spec"] is not None


def test_assistant_refuses_a_master_derived_profile(monkeypatch):
    """The assistant tunes a profile file; this one was never saved."""
    client = _client(monkeypatch)
    client.post("/audit", data={"profile": "__master__"}, files={
        "deck": ("rough.pptx", _deck_with_off_theme_colour(), "app/x"),
        "master": ("brand.pptx", _master(), "app/x"),
    })
    job_id = next(reversed(web._jobs))

    r = client.post(f"/assist/{job_id}")
    assert r.status_code == 400
    body = r.json()["error"]
    assert "master you uploaded" in body
    assert "Read a master" in body
