"""Title autofit shrink: detection of PowerPoint-shrunk titles and the
'Stop Fitting Text to This Placeholder' fix (designer workflow step 3)."""

import io

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from qc.engine import run_audit
from qc.fixer import apply_fixes, is_fixable
from qc.modules.font import _pct

ISSUE = "font.title_autofit_shrunk"


def _deck_with_title(text: str, font_scale: str | None = None,
                     ln_spc: str | None = None, on_body: bool = False) -> bytes:
    """One-slide deck; optionally plant a recorded shrink on the title (or on
    the body placeholder, to prove scoping)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = text
    slide.placeholders[1].text_frame.text = "body text"
    target = slide.placeholders[1] if on_body else slide.shapes.title
    if font_scale is not None or ln_spc is not None:
        bodyPr = target.text_frame._txBody.find(qn("a:bodyPr"))
        attrs = {}
        if font_scale is not None:
            attrs["fontScale"] = font_scale
        if ln_spc is not None:
            attrs["lnSpcReduction"] = ln_spc
        bodyPr.insert(0, bodyPr.makeelement(qn("a:normAutofit"), attrs))
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _audit(deck_bytes: bytes, tmp_path, profile="prezlab_en"):
    p = tmp_path / "deck.pptx"
    p.write_bytes(deck_bytes)
    result = run_audit(p, profile, modules=["font"])
    return [r.to_dict() for r in result.records if r.issue_type == ISSUE]


def test_pct_parses_both_st_forms():
    assert _pct("80000") == 80.0
    assert _pct("80%") == 80.0
    assert _pct("92.5%") == 92.5
    assert _pct(None) is None
    assert _pct("garbage") is None


def test_shrunk_title_detected(tmp_path):
    recs = _audit(_deck_with_title("A long shrunk title",
                                   font_scale="80000", ln_spc="10000"), tmp_path)
    assert len(recs) == 1
    r = recs[0]
    assert r["old_value"] == "80% scale, line spacing -10%"
    assert r["new_value"] == "no autofit (100%)"
    assert r["severity"] == "error"  # crushed to 80%: visibly wrong
    assert r["confidence"] == "high"  # tickable, never pre-selected
    assert not r["arabic_flag"]
    assert is_fixable(r)


def test_unshrunk_title_not_flagged(tmp_path):
    # autofit enabled but PowerPoint recorded no shrink
    assert _audit(_deck_with_title("Short", font_scale=None), tmp_path) == []
    assert _audit(_deck_with_title("Short", font_scale="100000"), tmp_path) == []


def test_shrunk_body_placeholder_out_of_scope(tmp_path):
    # step 3 targets titles; a shrunk body placeholder is not this finding
    assert _audit(_deck_with_title("Title", font_scale="70000",
                                   on_body=True), tmp_path) == []


def test_arabic_title_keeps_guard(tmp_path):
    recs = _audit(_deck_with_title("عنوان طويل جدا يحتاج تصغير",
                                   font_scale="75000"), tmp_path,
                  profile="prezlab_bilingual")
    assert len(recs) == 1
    assert recs[0]["arabic_flag"]
    assert not is_fixable(recs[0])


def test_stop_fitting_fix_end_to_end(tmp_path):
    deck = _deck_with_title("A long shrunk title", font_scale="80000")
    p = tmp_path / "deck.pptx"
    p.write_bytes(deck)
    result = run_audit(p, "prezlab_en", modules=["font"])
    records = [r.to_dict() for r in result.records]
    ids = {r["record_id"] for r in records if r["issue_type"] == ISSUE}
    assert ids

    fix = apply_fixes(deck, records, ids)
    assert [o.outcome for o in fix.outcomes] == ["changed"]

    prs = Presentation(io.BytesIO(fix.cleaned_bytes))
    bodyPr = prs.slides[0].shapes.title.text_frame._txBody.find(qn("a:bodyPr"))
    assert bodyPr.find(qn("a:normAutofit")) is None
    assert bodyPr.find(qn("a:noAutofit")) is not None
    assert prs.slides[0].shapes.title.text == "A long shrunk title"

    # verify-after-write: re-audit is clean of the issue
    cleaned = tmp_path / "cleaned.pptx"
    cleaned.write_bytes(fix.cleaned_bytes)
    again = run_audit(cleaned, "prezlab_en", modules=["font"])
    assert ISSUE not in again.summary["by_issue_type"]


def test_mild_shrink_stays_warning(tmp_path):
    recs = _audit(_deck_with_title("Slightly long title",
                                   font_scale="92500"), tmp_path)
    assert len(recs) == 1
    assert recs[0]["severity"] == "warning"
