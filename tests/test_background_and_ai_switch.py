"""Picture backgrounds captured in full, and the QC_AI master switch.

A background image is part of the brand, so the spec has to carry the image
itself, not a reference into a master file that Stage 2 is not allowed to
reopen. These tests plant the exact XML PowerPoint writes for a picture
background, captured from a COM-authored probe on 18/08/2026:

    <p:bg><p:bgPr><a:blipFill><a:blip r:embed="rIdN"/>
      <a:stretch><a:fillRect/></a:stretch></a:blipFill>
      <a:effectLst/></p:bgPr></p:bg>
"""

import base64
import io
import struct
import zlib

import pytest
from fastapi.testclient import TestClient
from lxml import etree
from pptx import Presentation

from qc.stylespec import (MAX_EMBEDDED_ASSET_BYTES, extract_style_spec,
                          spec_to_profile)
from qc.ui_master import spec_review
from tests.conftest import job_id_of

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _png(w=8, h=4, rgb=(30, 46, 97)) -> io.BytesIO:
    def chunk(tag, data):
        return (struct.pack(">I", len(data)) + tag + data
                + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF))

    ihdr = struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0)
    raw = b"".join(bytes([0]) + bytes(rgb) * w for _ in range(h))
    return io.BytesIO(b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr)
                      + chunk(b"IDAT", zlib.compress(raw)) + chunk(b"IEND", b""))


def _set_picture_background(container, png, *, alpha=None, tile=False, crop=None):
    """Replace the container's p:bg with a picture fill, as PowerPoint does."""
    _part, rId = container.part.get_or_add_image_part(png)
    cSld = container._element.find(f"{{{P_NS}}}cSld")
    for old in cSld.findall(f"{{{P_NS}}}bg"):
        cSld.remove(old)

    bg = etree.Element(f"{{{P_NS}}}bg")
    bgPr = etree.SubElement(bg, f"{{{P_NS}}}bgPr")
    blip_fill = etree.SubElement(bgPr, f"{{{A_NS}}}blipFill")
    blip = etree.SubElement(blip_fill, f"{{{A_NS}}}blip")
    blip.set(f"{{{R_NS}}}embed", rId)
    if alpha is not None:
        etree.SubElement(blip, f"{{{A_NS}}}alphaModFix").set("amt", str(int(alpha * 1000)))
    if crop:
        src = etree.SubElement(blip_fill, f"{{{A_NS}}}srcRect")
        for side, pct in crop.items():
            src.set(side, str(int(pct * 1000)))
    if tile:
        t = etree.SubElement(blip_fill, f"{{{A_NS}}}tile")
        t.set("tx", "0"); t.set("ty", "0")
        t.set("sx", "50000"); t.set("sy", "50000")
        t.set("algn", "tl")
    else:
        etree.SubElement(etree.SubElement(blip_fill, f"{{{A_NS}}}stretch"),
                         f"{{{A_NS}}}fillRect")
    # p:bg must be the FIRST child of p:cSld per the schema.
    cSld.insert(0, bg)
    return rId


# --------------------------------------------------------------- extraction


def test_picture_background_carries_the_image_itself():
    prs = Presentation()
    _set_picture_background(prs.slide_masters[0], _png())
    bg = extract_style_spec(prs)["master"]["background"]

    assert bg["kind"] == "image"
    img = bg["image"]
    assert img["sha1"] and img["format"] == "png"
    assert img["px"] == {"width": 8, "height": 4}
    assert img["bytes"] > 0
    # The bytes are IN the spec: Stage 2 never reopens the master.
    assert base64.b64decode(img["data_base64"])[:8] == b"\x89PNG\r\n\x1a\n"
    assert img["embed_skipped"] is None


def test_stretch_versus_tile_is_recorded():
    """A stretched photo and a tiled texture are different backgrounds even
    when they share one image."""
    stretched = Presentation()
    _set_picture_background(stretched.slide_masters[0], _png())
    assert extract_style_spec(stretched)["master"]["background"]["fill"]["mode"] \
        == "stretch"

    tiled = Presentation()
    _set_picture_background(tiled.slide_masters[0], _png(), tile=True)
    fill = extract_style_spec(tiled)["master"]["background"]["fill"]
    assert fill["mode"] == "tile"
    assert fill["tile"]["scale_pct"] == {"x": 50.0, "y": 50.0}
    assert fill["tile"]["align"] == "tl"


def test_watermark_alpha_and_crop_are_recorded():
    """A 12%-opacity watermark is not the same design as a full-strength
    image, and losing that detail would restyle a layout wrongly."""
    prs = Presentation()
    _set_picture_background(prs.slide_masters[0], _png(), alpha=12,
                            crop={"l": 5, "r": 5})
    fill = extract_style_spec(prs)["master"]["background"]["fill"]

    assert fill["alpha_pct"] == 12.0
    assert fill["crop_pct"]["l"] == 5.0
    assert fill["crop_pct"]["r"] == 5.0


def test_layout_background_is_captured_separately_from_the_master():
    """A dark image on the title layout over a light master is a real pattern;
    the layout entry has to carry its own background."""
    prs = Presentation()
    layout = list(prs.slide_masters[0].slide_layouts)[0]
    _set_picture_background(layout, _png())
    spec = extract_style_spec(prs)

    assert spec["layouts"][0]["background"]["kind"] == "image"
    # Layouts that declare nothing inherit, and say so by carrying None.
    assert spec["layouts"][1]["background"] is None


def test_embedding_can_be_switched_off():
    prs = Presentation()
    _set_picture_background(prs.slide_masters[0], _png())
    img = extract_style_spec(prs, embed_assets=False)["master"]["background"]["image"]

    assert img["sha1"]          # identity survives
    assert img["data_base64"] is None
    assert "disabled" in img["embed_skipped"]


def test_oversized_asset_is_skipped_and_says_why(monkeypatch):
    """Silently dropping the bytes would let a consumer assume it has them."""
    monkeypatch.setattr("qc.stylespec.MAX_EMBEDDED_ASSET_BYTES", 10)
    prs = Presentation()
    _set_picture_background(prs.slide_masters[0], _png())
    img = extract_style_spec(prs)["master"]["background"]["image"]

    assert img["data_base64"] is None
    assert "exceeds" in img["embed_skipped"]
    assert img["sha1"]


def test_theme_fill_kind_is_reported_for_a_bgref():
    """The default template uses p:bgRef; the colour alone would hide the fact
    that the theme's fill style could be a picture or gradient."""
    bg = extract_style_spec(Presentation())["master"]["background"]
    assert bg["kind"] == "theme_ref"
    assert bg["theme_fill_kind"] == "solid"


def test_spec_with_an_image_background_still_round_trips_json():
    import json

    prs = Presentation()
    _set_picture_background(prs.slide_masters[0], _png())
    spec = extract_style_spec(prs)
    assert json.loads(json.dumps(spec)) == spec


# ------------------------------------------------------------------ the page


def test_page_previews_the_background_image():
    prs = Presentation()
    _set_picture_background(prs.slide_masters[0], _png(), alpha=40)
    spec = extract_style_spec(prs)
    html = spec_review(spec, "sid")

    assert "picture background" in html
    assert "40% opacity" in html
    # Rendered as a real preview, so a designer can confirm the right image.
    assert "src=\"data:image/png;base64," in html


def test_page_marks_a_layout_that_overrides_the_background():
    prs = Presentation()
    layout = list(prs.slide_masters[0].slide_layouts)[0]
    _set_picture_background(layout, _png())
    html = spec_review(extract_style_spec(prs), "sid")

    assert "<th>Background</th>" in html
    assert "inherits" in html          # the layouts that declare nothing


# ------------------------------------------------------------- the AI switch


@pytest.fixture()
def ai_off(monkeypatch):
    from qc import web

    monkeypatch.setattr(web, "AI_ENABLED", False)
    monkeypatch.setattr(web, "AUTH_REQUIRED", False)
    web.app.state.auth_required = False
    return TestClient(web.app)


def _job_with_findings(client):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    from qc import web

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for _ in range(3):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(914400), Emu(914400),
                                    Emu(914400), Emu(914400))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor.from_string("FF00AA")
    buf = io.BytesIO(); prs.save(buf)
    client.post("/audit", data={"profile": "prezlab_en"},
                files={"deck": ("d.pptx", buf.getvalue(), "app/x")})
    return next(reversed(web._jobs))


def test_assistant_route_refuses_when_ai_is_off(ai_off):
    job_id = _job_with_findings(ai_off)
    r = ai_off.post(f"/assist/{job_id}")
    assert r.status_code == 503
    assert "QC_AI=0" in r.json()["error"]


def test_assist_apply_refuses_when_ai_is_off(ai_off):
    job_id = _job_with_findings(ai_off)
    r = ai_off.post(f"/assist/{job_id}/apply", data={"accepted": ["x"]})
    assert r.status_code == 503


def test_copilot_route_refuses_when_ai_is_off(ai_off):
    job_id = _job_with_findings(ai_off)
    r = ai_off.post(f"/copilot/{job_id}", follow_redirects=False)
    assert r.status_code == 503


# The panel is identified by its container and the copilot form's action.
# Button LABELS are not usable markers: the report's JavaScript carries the
# assistant's button text as a string literal whether or not the panel renders.
_PANEL = 'class="card assist'
_COPILOT_FORM = 'action="/copilot/'


def test_report_hides_the_ai_panel_when_ai_is_off(ai_off):
    """Hiding the panel is not the security guarantee (the routes are), but a
    disabled feature must not advertise itself either."""
    prs = Presentation()
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for _ in range(3):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(914400), Emu(914400),
                                    Emu(914400), Emu(914400))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor.from_string("FF00AA")
    buf = io.BytesIO(); prs.save(buf)

    r = ai_off.post("/audit", data={"profile": "prezlab_en"},
                    files={"deck": ("d.pptx", buf.getvalue(), "app/x")})
    assert r.status_code == 200
    # the panel lives on the report, which the upload no longer lands on
    report = ai_off.get(f"/audit/{job_id_of(r)}").text
    assert _PANEL not in report
    assert _COPILOT_FORM not in report


def test_the_switch_is_what_hides_the_panel(monkeypatch):
    """Same deck, AI on: the panel returns. Otherwise the test above could be
    passing because of something unrelated to the switch."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    from qc import web

    monkeypatch.setattr(web, "AI_ENABLED", True)
    monkeypatch.setattr(web, "AUTH_REQUIRED", False)
    web.app.state.auth_required = False

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for _ in range(3):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(914400), Emu(914400),
                                    Emu(914400), Emu(914400))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor.from_string("FF00AA")
    buf = io.BytesIO(); prs.save(buf)

    client = TestClient(web.app)
    r = client.post("/audit", data={"profile": "prezlab_en"},
                    files={"deck": ("d.pptx", buf.getvalue(), "app/x")})
    assert r.status_code == 200
    report = client.get(f"/audit/{job_id_of(r)}").text
    assert _PANEL in report
    assert _COPILOT_FORM in report


def test_the_switch_defaults_on_so_other_instances_are_unchanged(monkeypatch):
    """QC_AI is opt-OUT: an instance that sets nothing keeps its AI features,
    so disabling them on this machine changes nothing for anyone else. The
    parsing is tested directly rather than by reloading qc.config, because a
    reload re-reads this machine's .env and would assert its way to the local
    answer instead of the default one."""
    from qc.config import opt_out_flag

    monkeypatch.delenv("QC_AI", raising=False)
    assert opt_out_flag("QC_AI") is True

    for off in ("0", "false", "no", "off", "OFF", " 0 "):
        monkeypatch.setenv("QC_AI", off)
        assert opt_out_flag("QC_AI") is False, off

    for on in ("1", "true", "yes"):
        monkeypatch.setenv("QC_AI", on)
        assert opt_out_flag("QC_AI") is True, on
