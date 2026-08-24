"""The reserved header band: the strip a master draws under its subtitle,
between the floor the subtitle may not cross and the ceiling the body may not
cross, which stays empty on every slide.

Built from the real client master (PIF, read 20/08/2026): horizontal guides at
0.45in (top margin), 1.65in (subtitle floor), 1.90in (body ceiling), 3.75in
(canvas centre) and 6.80in (bottom margin), with the master's own body
placeholder starting at 1.896in - the master stating the same line twice. The
deck that prompted the rule had photos standing at 1.5in on slide after slide,
which the text-only safe-zone check never looked at.
"""

import copy
import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

import qc.modules.margin_alignment as ma
from qc.fixer import apply_fixes, is_fixable
from qc.profile import Profile
from tests.conftest import save_and_ctx

IN = 914400
MM = 36000
BLANK = 6

FLOOR = int(1.65 * IN)
CEILING = int(1.90 * IN)

_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
_P15 = "http://schemas.microsoft.com/office/powerpoint/2012/main"


def _band_profile(en_profile, floor=FLOOR, ceiling=CEILING):
    """The English profile with a stated band, so the tests exercise the
    profile path rather than the master fallback."""
    data = copy.deepcopy(en_profile.raw)
    data["config"]["geometry"]["body_band_emu"] = (
        {"subtitle_floor": floor, "body_top": ceiling}
        if ceiling else None)
    return Profile(data)


def _plant_guides(master, horizontal_in=()):
    """Guides as desktop PowerPoint stores them: eighths of a point from the
    top-left edge (python-pptx cannot draw them, so the captured format is
    planted directly). Vertical guides come along because infer_grid only
    reads a frame when both axes state one."""
    import itertools

    from lxml import etree

    ext_lst = etree.SubElement(master._element, f"{_P}extLst")
    ext = etree.SubElement(ext_lst, f"{_P}ext")
    ext.set("uri", "{GUIDES}")
    lst = etree.SubElement(ext, f"{{{_P15}}}sldGuideLst")
    gid = itertools.count(1)
    for pos, horz in ([(0.50, False), (12.83, False)]
                      + [(h, True) for h in horizontal_in]):
        g = etree.SubElement(lst, f"{{{_P15}}}guide")
        g.set("id", str(next(gid)))
        if horz:
            g.set("orient", "horz")
        g.set("pos", str(int(pos * 72 * 8)))


def _photo(slide, left, top, w=2 * IN, h=1.6 * IN):
    return slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(left)),
                                  Emu(int(top)), Emu(int(w)), Emu(int(h)))


def _text(slide, left, top, w, h, text):
    tb = slide.shapes.add_textbox(Emu(int(left)), Emu(int(top)),
                                  Emu(int(w)), Emu(int(h)))
    tb.text_frame.text = text
    return tb


def _band_recs(ctx):
    return [r.to_dict() for r in ma.detect(ctx)
            if r.issue_type == "margin_alignment.body_band_intrusion"]


def _photo_slide(prs, photo_top=1.5 * IN):
    """The shape of the client slide: three photos in a row with a caption
    under each, the photos standing in the reserved band."""
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    photos, captions = [], []
    for i in range(3):
        left = (0.6 + i * 2.4) * IN
        photos.append(_photo(slide, left, photo_top))
        captions.append(_text(slide, left, photo_top + 1.7 * IN,
                              2 * IN, 0.8 * IN, f"Caption {i + 1}"))
    return slide, photos, captions


# ------------------------------------------------------------- detection


def test_photos_standing_in_the_band_are_flagged_once_per_slide(
        make_prs, en_profile, tmp_path):
    """Pictures count. The strip is about the page, not about copy, so a photo
    in it is the same defect as a paragraph - and the safe-zone rule, which
    only ever looked at text, saw none of these."""
    prs = make_prs()
    _slide, photos, _caps = _photo_slide(prs)
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))

    recs = _band_recs(ctx)
    assert len(recs) == 1, "one finding per slide, not one per intruder"
    rec = recs[0]
    assert rec["severity"] == "error"
    assert rec["confidence"] == "high"
    assert rec["profile_rule_id"] == "geometry.body_band_emu"
    assert rec["shape_id"] == str(photos[0].shape_id)
    assert int(rec["new_value"]) == CEILING - int(1.5 * IN)
    assert is_fixable(rec)
    assert "3 element(s)" in rec["message"]


def test_content_clear_of_the_band_is_not_flagged(make_prs, en_profile,
                                                  tmp_path):
    prs = make_prs()
    _photo_slide(prs, photo_top=CEILING)
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))

    assert _band_recs(ctx) == []


def test_a_shape_a_millimetre_over_the_ceiling_is_rounding_not_a_breach(
        make_prs, en_profile, tmp_path):
    """The client master's own body placeholder starts 0.007in above its body
    guide, so a shape aligned to the master reads as a breach without slack."""
    prs = make_prs()
    _photo_slide(prs, photo_top=CEILING - 20000)
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))

    assert _band_recs(ctx) == []


def test_header_furniture_above_the_floor_is_left_alone(make_prs, en_profile,
                                                        tmp_path):
    """An eyebrow living in the header band is where it belongs. It is not an
    intruder, and it must not ride the body move down onto the title either."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    eyebrow = _text(slide, 0.6 * IN, 0.5 * IN, 3 * IN, 0.3 * IN, "SECTION")
    _photo(slide, 0.6 * IN, 1.5 * IN)
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))

    recs = _band_recs(ctx)
    assert len(recs) == 1
    assert str(eyebrow.shape_id) not in recs[0]["locator"].split(":")[2]


def test_a_header_box_that_stops_inside_the_strip_is_slack_not_an_intruder(
        make_prs, en_profile, tmp_path):
    """Real-master regression: a text box carries descender slack under its
    last line - the client master's own subtitle placeholder overhangs its
    floor by 2.4mm - and counting such a box as body content made its top, up
    in the header, the line the whole block was measured from. That asked for a
    36mm move on the designer's own sample deck."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    _text(slide, 0.6 * IN, 1.2 * IN, 5 * IN, 0.55 * IN, "The standfirst")
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))

    assert _band_recs(ctx) == [], "the overhang is box slack, not a filled band"


def test_a_full_height_image_is_not_body_content_that_crept_up(
        make_prs, en_profile, tmp_path):
    """Real-master regression: a 94%-height image running from the top margin
    to the bottom edge is design, and pulling it into the body block asked for
    a 36mm move that could never be applied."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    _photo(slide, 0.5 * IN, 0.45 * IN, w=4 * IN, h=7 * IN)
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))

    assert _band_recs(ctx) == []


def test_a_millimetre_or_two_into_the_strip_is_a_warning_not_an_error(
        make_prs, en_profile, tmp_path):
    """Evidence-based severity, as everywhere else in this module: a shape 2mm
    off the guide is a judgment call and must not arrive pre-ticked, while a
    photo standing 10mm into a 6mm strip is confidently wrong."""
    prs = make_prs()
    _photo_slide(prs, photo_top=CEILING - 2 * MM)
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))

    rec = _band_recs(ctx)[0]
    assert rec["severity"] == "warning"
    assert is_fixable(rec), "still tickable: the designer decides"


def test_a_body_move_is_never_pre_ticked(make_prs, en_profile, tmp_path):
    """Moving every element on a slide is a design decision even when the
    broken line is a fact, so the tick is the approval."""
    from qc.fixer import needs_explicit_tick, tick_reason

    prs = make_prs()
    _photo_slide(prs)
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))

    rec = _band_recs(ctx)[0]
    assert rec["severity"] == "error"       # would otherwise be pre-selected
    assert needs_explicit_tick(rec)
    assert "your approval" in tick_reason(rec)


def test_no_band_stated_means_no_finding(make_prs, en_profile, tmp_path):
    """No guides, no band: nothing is measured against a line the client never
    drew, whatever sits high on the slide."""
    prs = make_prs()
    _photo_slide(prs)
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile, ceiling=None))

    assert _band_recs(ctx) == []


def test_a_ceiling_without_a_floor_states_no_strip(make_prs, en_profile,
                                                   tmp_path):
    """One interior guide gives a body ceiling but no reserved strip, and
    without the strip there is no way to tell header furniture from body
    content that crept up into it."""
    prs = make_prs()
    _photo_slide(prs)
    ctx = save_and_ctx(prs, tmp_path,
                       _band_profile(en_profile, floor=None, ceiling=CEILING))

    assert _band_recs(ctx) == []


def test_the_band_is_read_from_the_decks_own_master_when_the_profile_is_silent(
        make_prs, en_profile, tmp_path):
    """A profile written before the band existed still gets the rule: the
    guides on the deck's own master are the same source the profile would have
    been projected from."""
    prs = make_prs()
    _plant_guides(prs.slide_masters[0],
                  horizontal_in=(0.45, 1.65, 1.90, 3.75, 6.80))
    _photo_slide(prs)
    ctx = save_and_ctx(prs, tmp_path, en_profile)   # no band in the profile

    recs = _band_recs(ctx)
    assert len(recs) == 1
    # guides quantise to eighths of a point, so the target is the guide's
    # stored value rather than the inch it was typed as
    assert abs(int(recs[0]["new_value"]) - (CEILING - int(1.5 * IN))) < 2000


def test_a_centre_guide_is_never_read_as_a_body_ceiling(make_prs, en_profile,
                                                        tmp_path):
    """Masters carry a centre guide as a placement aid; read as a ceiling it
    would put the body's start half way down the slide."""
    prs = make_prs()
    _plant_guides(prs.slide_masters[0], horizontal_in=(0.45, 3.75, 6.80))
    _photo_slide(prs)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    assert _band_recs(ctx) == []


def test_a_heading_growing_through_the_ceiling_is_reported_never_moved(
        make_prs, en_profile, tmp_path):
    """A standfirst that has grown into the body area is a copy conversation:
    it gets the heading finding, which carries no target and no tick."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    # three text shapes with real hierarchy, so heading_ids can read a title
    _text(slide, 0.6 * IN, 0.4 * IN, 9 * IN, 0.6 * IN, "The Heading")
    # kept under a tenth of the canvas: a box bigger than that is a panel, and
    # qc.util.heading_ids will not read a panel as a line of header text
    tall = _text(slide, 0.6 * IN, 1.0 * IN, 5 * IN, 1.4 * IN,
                 "A standfirst that has grown down through the reserved strip")
    for i in range(3):
        _text(slide, (0.6 + i * 3) * IN, 3 * IN, 2.5 * IN, 1 * IN, f"Body {i}")
    for para in slide.shapes[0].text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Emu(28 * 12700)
    for para in tall.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Emu(14 * 12700)
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))

    recs = [r.to_dict() for r in ma.detect(ctx)]
    heads = [r for r in recs
             if r["issue_type"] == "margin_alignment.heading_past_margin"]
    assert any("body ceiling" in r["message"] for r in heads)
    assert all(not is_fixable(r) for r in heads)
    # and the standfirst is not in any body-block move either
    for r in recs:
        if r["issue_type"] == "margin_alignment.body_band_intrusion":
            assert str(tall.shape_id) not in (r["locator"] or "")


# ------------------------------------------------------------------- fix


def test_the_fix_moves_the_whole_body_block_and_keeps_its_arrangement(
        make_prs, en_profile, tmp_path):
    """One move for the body, not one per intruder: pushing the photos down
    alone would drop each one onto its own caption."""
    prs = make_prs()
    _slide, photos, captions = _photo_slide(prs)
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))
    recs = _band_recs(ctx)
    rec = recs[0]
    dy = int(rec["new_value"])
    before = {s.shape_id: s.top for s in (photos + captions)}

    result = apply_fixes(ctx.deck_path.read_bytes(), recs, {rec["record_id"]})
    assert result.applied == 1, [o.reason for o in result.outcomes]
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]

    after = {s.shape_id: s.top for s in out.shapes}
    assert all(after[sid] == top + dy for sid, top in before.items())
    tops = sorted(after[p.shape_id] for p in photos)
    assert tops[0] == CEILING, "the block now starts on the body guide"


def test_a_block_that_cannot_come_down_is_reported_not_moved(
        make_prs, en_profile, tmp_path):
    """The audit never pushes a designer's content off the page. The migration
    pass seats the block on the ceiling anyway, because it is rebuilding the
    slide onto the master and says so loudly; a tick on a finished deck is a
    different promise."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    _photo(slide, 0.6 * IN, 1.4 * IN, w=2 * IN, h=1.6 * IN)
    # reaches to 7.2in: the 0.5in the block owes the ceiling would take it past
    # the bottom edge of the canvas
    _text(slide, 0.6 * IN, 3.1 * IN, 4 * IN, 4.1 * IN, "A very long list")
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))

    recs = _band_recs(ctx)
    assert len(recs) == 1
    assert recs[0]["locator"] is None
    assert recs[0]["new_value"] is None
    assert not is_fixable(recs[0])
    assert "needs a rework" in recs[0]["message"]


def test_footer_furniture_never_rides_the_body_move(make_prs, en_profile,
                                                    tmp_path):
    """A source line or page number in the bottom strip stays pinned; dragging
    it with the body is how furniture ends up off the canvas."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    _photo(slide, 0.6 * IN, 1.5 * IN)
    footer = _text(slide, 0.6 * IN, 6.9 * IN, 3 * IN, 0.25 * IN,
                   "Source: internal analysis")
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))

    recs = _band_recs(ctx)
    assert len(recs) == 1
    assert str(footer.shape_id) not in recs[0]["locator"].split(":")[2]


def test_an_arabic_deck_still_gets_the_move(make_prs, bilingual_profile,
                                            tmp_path):
    """Geometry is script-neutral: translating a block never opens its text,
    so an Arabic deck is fixable here exactly as an English one."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    _photo(slide, 0.6 * IN, 1.5 * IN)
    _text(slide, 0.6 * IN, 3.3 * IN, 4 * IN, 0.5 * IN, "وعدنا لكم")
    ctx = save_and_ctx(prs, tmp_path, _band_profile(bilingual_profile))

    recs = _band_recs(ctx)
    assert len(recs) == 1
    assert recs[0]["arabic_flag"] is True
    assert is_fixable(recs[0])
    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {recs[0]["record_id"]})
    assert result.applied == 1, [o.reason for o in result.outcomes]


# ------------------------------- the other direction: a body that starts LOW
#
# The strip rule was written for content that had crept UP into the header's
# clear band, so it only ever asked shapes to come DOWN. A column that began
# well BELOW the guide was silently clean - and that silence is what produced
# the reported result, because the cluster rule is still running: with a photo
# at 2.30in and a row of cards at 2.55in, _edge_misaligned aligns the photo TO
# THE CARDS (they are the majority, so they set the median) and the slide passes
# with its whole body below the line the master states (design lead, 24/08/2026:
# "the picture on the left and the elements on the right should both be aligned
# to the top of the presentation space but only the picture was aligned").


def _low_recs(ctx):
    return [r.to_dict() for r in ma.detect(ctx)
            if r.issue_type == "margin_alignment.body_below_band"]


def _two_column_slide(prs, picture_top, cards_top):
    """The reported shape: a picture down the left, a row of cards to the
    right of it, each column starting on its own line."""
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    picture = _photo(slide, 0.6 * IN, picture_top, 3.2 * IN, 3.4 * IN)
    cards = [_photo(slide, (4.3 + i * 2.7) * IN, cards_top, 2.4 * IN, 1.0 * IN)
             for i in range(3)]
    return slide, picture, cards


def test_a_column_starting_below_the_guide_is_flagged(make_prs, en_profile,
                                                      tmp_path):
    """The reported case exactly: the picture is on the guide, the cards are
    0.65in under it, and that used to be a clean slide."""
    prs = make_prs()
    _slide, _pic, cards = _two_column_slide(prs, CEILING,
                                            CEILING + int(0.65 * IN))
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))

    recs = _low_recs(ctx)
    assert len(recs) == 1, f"expected one decision, got {len(recs)}"
    rec = recs[0]
    assert int(rec["new_value"]) == CEILING, "the target must be the guide"
    assert rec["severity"] == "error", \
        "a column sitting on the guide is evidence the guide is live here"
    ids = set(rec["locator"].split(":", 2)[2].split(","))
    assert ids == {str(c.shape_id) for c in cards}, \
        "the lift must carry the low column and nothing that was on the line"


def test_three_cards_on_one_low_line_are_one_decision(make_prs, en_profile,
                                                      tmp_path):
    """Three columns geometrically, one problem to a designer. Three tick boxes
    that all move by the same 0.65in is the per-occurrence noise this module
    exists to avoid."""
    prs = make_prs()
    _two_column_slide(prs, CEILING, CEILING + int(0.65 * IN))
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))
    assert len(_low_recs(ctx)) == 1


def test_the_guide_beats_the_cluster_median(make_prs, en_profile, tmp_path):
    """Both columns low. _edge_misaligned wants the picture pulled DOWN onto
    the cards' line; the guide rule wants both lifted UP onto the master's.
    Ticking both would move the picture twice from two targets computed off the
    same start, so the inferred median gives way to the stated guide."""
    prs = make_prs()
    _slide, pic, _cards = _two_column_slide(
        prs, CEILING + int(0.40 * IN), CEILING + int(0.65 * IN))
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))
    every = [r.to_dict() for r in ma.detect(ctx)]

    tops = [r for r in every
            if r["issue_type"] == "margin_alignment.edge_misaligned"
            and (r["property"] or "").endswith(".y")
            and str(r["shape_id"]) == str(pic.shape_id)]
    assert not tops, "the median fix survived and now fights the guide fix"
    low = [r for r in every
           if r["issue_type"] == "margin_alignment.body_below_band"]
    assert len(low) == 2, "each column starts on its own line, so each lifts"
    assert all(int(r["new_value"]) == CEILING for r in low)


def test_every_column_low_is_a_warning_not_an_error(make_prs, en_profile,
                                                    tmp_path):
    """With nothing on the guide, the guide being the intended line is
    inference rather than evidence - a body that begins low all the way across
    may be the composition."""
    prs = make_prs()
    _two_column_slide(prs, CEILING + int(0.40 * IN), CEILING + int(0.65 * IN))
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))
    assert {r["severity"] for r in _low_recs(ctx)} == {"warning"}


def test_a_second_row_is_not_a_column_that_missed_the_line(make_prs,
                                                           en_profile,
                                                           tmp_path):
    """A row below the first shares its columns' horizontal span, so only each
    column's TOPMOST shape is judged. Otherwise every two-row layout in the
    deck becomes a finding."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    for i in range(3):
        left = (0.6 + i * 2.4) * IN
        _photo(slide, left, CEILING, 2 * IN, 1.2 * IN)                  # on line
        _photo(slide, left, CEILING + int(2.0 * IN), 2 * IN, 1.2 * IN)  # row two
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))
    assert _low_recs(ctx) == []


def test_a_single_column_is_a_composition_not_a_misalignment(make_prs,
                                                             en_profile,
                                                             tmp_path):
    """With one column there is no line to be off, only content that begins
    where it begins."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    _photo(slide, 3.0 * IN, CEILING + int(1.2 * IN), 4 * IN, 2 * IN)
    _photo(slide, 3.2 * IN, CEILING + int(3.4 * IN), 3 * IN, 1 * IN)
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))
    assert _low_recs(ctx) == []


def test_a_few_millimetres_low_is_not_a_finding(make_prs, en_profile,
                                                tmp_path):
    """6mm is about the smallest gap a designer reads as off the line; under
    that, a rule that moves every element in a column must not fire."""
    prs = make_prs()
    _two_column_slide(prs, CEILING, CEILING + 4 * MM)
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))
    assert _low_recs(ctx) == []


def test_no_band_stated_means_no_low_finding(make_prs, en_profile, tmp_path):
    prs = make_prs()
    _two_column_slide(prs, CEILING, CEILING + int(0.65 * IN))
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile, ceiling=None))
    assert _low_recs(ctx) == []


def test_the_lift_puts_the_column_on_the_guide_and_keeps_its_arrangement(
        make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide, _pic, cards = _two_column_slide(prs, CEILING,
                                           CEILING + int(0.65 * IN))
    # a caption under each card, to prove the column travels together
    caps = [_text(slide, (4.4 + i * 2.7) * IN, CEILING + int(1.8 * IN),
                  2.2 * IN, 0.6 * IN, f"Caption {i + 1}") for i in range(3)]
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))
    recs = [r.to_dict() for r in ma.detect(ctx)]
    low = [r for r in recs
           if r["issue_type"] == "margin_alignment.body_below_band"]
    assert low and is_fixable(low[0])

    gaps_before = [caps[i].top - cards[i].top for i in range(3)]
    out = apply_fixes(ctx.deck_path.read_bytes(), recs, {low[0]["record_id"]})
    assert out.applied == 1, [o.reason for o in out.outcomes]
    after = Presentation(io.BytesIO(out.cleaned_bytes)).slides[0]
    by_id = {str(s.shape_id): s for s in after.shapes}
    assert by_id[str(cards[0].shape_id)].top == CEILING, "not on the guide"
    for i in range(3):
        assert (by_id[str(caps[i].shape_id)].top
                - by_id[str(cards[i].shape_id)].top) == gaps_before[i], \
            "the arrangement inside the column changed"


def test_the_lift_is_never_pre_ticked(make_prs, en_profile, tmp_path):
    """It moves every element in a column. Whether the slide should shift is
    the designer's call, so the tick is the approval."""
    from qc.fixer import needs_explicit_tick, tick_reason

    prs = make_prs()
    _two_column_slide(prs, CEILING, CEILING + int(0.65 * IN))
    ctx = save_and_ctx(prs, tmp_path, _band_profile(en_profile))
    rec = _low_recs(ctx)[0]
    assert needs_explicit_tick(rec)
    assert "approval" in tick_reason(rec)


def test_a_body_frame_states_the_ceiling_even_with_no_guides(tmp_path):
    """The presentation space, carried to this stage.

    infer_grid already works out which of two things a drawn rectangle is - a
    frame around the page has a top margin, a frame around the BODY has a line
    where content begins - and sets body_top_emu only in the second case. This
    module used to insist the source be "guides", so a master stating the body's
    top as a rectangle got no ceiling at all and every slide on it passed
    however far below the line its content began (design lead, 24/08/2026).

    A rectangle gives a ceiling and no floor, and the two directions need
    different evidence: creeping UP needs the floor to tell an eyebrow in the
    header from body content in the strip, sitting BELOW needs only the line."""
    from pptx.util import Emu

    from qc.profile import Profile
    from qc.pspace import ensure_presentation_space
    from tests.test_presentation_space import (_bytes,
                                               _deck_with_marker_on_a_layout)

    raw, _notes = ensure_presentation_space(
        _bytes(_deck_with_marker_on_a_layout()))
    prs = Presentation(io.BytesIO(raw))
    profile = Profile({"id": "t", "config": {"geometry": {
        # page margins only: no band stated anywhere in the profile
        "safe_zone_margins_emu": {"left": int(0.5 * IN), "top": int(0.45 * IN),
                                  "right": int(0.5 * IN),
                                  "bottom": int(0.7 * IN)},
        "alignment": {"edge_tolerance_emu": 9144,
                      "intent_window_emu": 228600}}}})

    floor, ceiling = ma._body_band(profile, prs)
    assert floor is None, "a rectangle cannot state the subtitle floor"
    assert ceiling == 2 * IN, "the rectangle's top IS where the body begins"

    # and it is used: one column on that line, one below it
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])

    def box(x, y, w, h):
        return slide.shapes.add_shape(1, Emu(int(x * IN)), Emu(int(y * IN)),
                                      Emu(int(w * IN)), Emu(int(h * IN)))

    box(0.6, 2.0, 3.0, 3.0)
    for i in range(3):
        box(4.0 + i * 2.7, 2.65, 2.4, 1.0)
    ctx = save_and_ctx(prs, tmp_path, profile)
    low = [r.to_dict() for r in ma.detect(ctx)
           if r.issue_type == "margin_alignment.body_below_band"]
    assert len(low) == 1
    assert int(low[0]["new_value"]) == 2 * IN


def test_the_page_top_margin_is_never_read_as_the_body_top(tmp_path,
                                                           en_profile):
    """The mistake this guards: the safe-zone top is where the PAGE begins, and
    on these masters the body begins well below it with the header between. Used
    as a ceiling it fires on every normal slide in the deck."""
    prs = Presentation()
    profile = _band_profile(en_profile, ceiling=None)
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    _photo(slide, 0.6 * IN, 2.4 * IN)          # well below any page margin
    _photo(slide, 4.0 * IN, 2.4 * IN)
    ctx = save_and_ctx(prs, tmp_path, profile)
    assert _low_recs(ctx) == [], \
        "the page margin was read as the line the body starts on"
