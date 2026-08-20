"""Bootstrap: margins learned from the deck's own content edges, the
perceptual-floor alignment tolerances (ground-truth calibration round,
19/07/2026), and the extracted style spec (theme + brand)."""

import io
import itertools
import struct
import zlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.shapes.picture import CT_Picture
from pptx.util import Emu

from qc.bootstrap import (INTENT_WINDOW_EMU, MARGIN_BUFFER_EMU,
                          MIN_COLOR_COUNT, MIN_LOGO_SLIDE_STAMPS,
                          VISUAL_EDGE_TOLERANCE_EMU, build_profile,
                          extract_theme, find_brand_marks, learn_margins)

IN = 914400
HALF = IN // 2


def _deck(n_shapes=30, left=274320, top=274320, right=274320, bottom=274320,
          full_bleed=False):
    """Shapes laid out to respect exact margins on every side."""
    prs = Presentation()
    sw, sh = prs.slide_width, prs.slide_height
    per_slide = 10
    for s in range((n_shapes + per_slide - 1) // per_slide):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if full_bleed:
            slide.shapes.add_textbox(Emu(0), Emu(0), sw, sh)
        for i in range(per_slide):
            w, h = 1200000, 300000
            x = left if i % 2 == 0 else sw - right - w
            if i == per_slide - 1:  # one shape per slide ends AT the margin
                y = sh - bottom - h
            else:
                y = top + i * ((sh - top - bottom - h) // per_slide)
            tb = slide.shapes.add_textbox(Emu(x), Emu(int(y)), Emu(w), Emu(h))
            tb.text_frame.text = "content"  # margins are learned from TEXT
    return prs


def test_margins_learned_from_content_edges():
    prs = _deck()
    m = learn_margins(prs)
    assert m is not None
    expected = 274320 - MARGIN_BUFFER_EMU  # buffered below observed p5
    assert m["left"] == expected
    assert m["right"] == expected
    assert m["top"] == expected
    assert m["bottom"] == expected


def test_full_bleed_shapes_do_not_drag_margins_to_zero():
    m = learn_margins(_deck(full_bleed=True))
    assert m is not None
    assert m["left"] == 274320 - MARGIN_BUFFER_EMU  # 0-margin bg exempt


def test_sparse_deck_returns_none_and_profile_falls_back():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Emu(IN), Emu(IN), Emu(IN), Emu(IN))
    tb.text_frame.text = "lonely"
    assert learn_margins(prs) is None

    profile = build_profile(prs, "t", "t")
    margins = profile["config"]["geometry"]["safe_zone_margins_emu"]
    assert margins["left"] == round(prs.slide_width * 0.0375)  # canvas ratio


def test_build_profile_uses_learned_margins_and_visual_tolerances():
    prs = _deck()
    profile = build_profile(prs, "t", "t")
    geo = profile["config"]["geometry"]
    assert geo["safe_zone_margins_emu"]["left"] == 274320 - MARGIN_BUFFER_EMU
    assert geo["alignment"]["edge_tolerance_emu"] == VISUAL_EDGE_TOLERANCE_EMU
    assert geo["alignment"]["intent_window_emu"] == INTENT_WINDOW_EMU


# --------------------------------------------------------------- style spec

_SHAPE_IDS = itertools.count(500)


def _png(rgb=(0, 0, 0)) -> io.BytesIO:
    """A valid 1x1 PNG. Distinct rgb values give distinct sha1s, which is how
    brand-mark identity is established."""

    def chunk(tag, data):
        return (struct.pack(">I", len(data)) + tag + data
                + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF))

    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)  # 1x1, 8-bit truecolor
    raw = bytes([0]) + bytes(rgb)                        # filter byte + pixel
    return io.BytesIO(b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr)
                      + chunk(b"IDAT", zlib.compress(raw)) + chunk(b"IEND", b""))


def _stamp(container, png, left=IN, top=IN, w=HALF, h=HALF, name="brandmark"):
    """Put a picture on a slide MASTER or LAYOUT.

    python-pptx exposes add_picture on slides only (MasterShapes and
    LayoutShapes are read-only collections), so the p:pic element and its
    image relationship are built directly. This is the shape a designer's
    finished master actually carries, which is the case that matters most."""
    _image_part, rId = container.part.get_or_add_image_part(png)
    pic = CT_Picture.new_pic(next(_SHAPE_IDS), name, name, rId, left, top, w, h)
    container._element.spTree.append(pic)
    return pic


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _filled(slide, hexval):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(IN), Emu(IN),
                                Emu(HALF), Emu(HALF))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string(hexval)
    return sh


def test_theme_reports_what_the_theme_declares_not_what_slides_use(make_prs):
    """The point of the theme block: a deck whose shapes all hardcode literal
    RGB still has a theme, and the roles are what restyling needs."""
    prs = make_prs()
    _filled(_blank(prs), "123456")

    theme = extract_theme(prs)
    assert {"dk1", "lt1", "accent1", "accent6"} <= set(theme["colors"])
    assert theme["colors"]["accent1"] != "123456"
    assert len(theme["colors"]["accent1"]) == 6
    assert theme["color_map"]["bg1"] == "lt1"
    assert theme["fonts"]["minor"]["latin"]
    assert theme["fonts"]["major"]["latin"]


def test_theme_readable_from_a_file_with_no_slides():
    """The plan's actual submission case: a designer sends the master, not a
    deck. There is nothing to survey and the theme still has to come out."""
    prs = Presentation()
    assert len(prs.slides) == 0
    theme = extract_theme(prs)
    assert theme["colors"]["accent1"]
    assert theme["fonts"]["major"]["latin"]


def test_named_colors_carry_theme_refs_for_exact_slot_matches(make_prs):
    prs = make_prs()
    accent1 = extract_theme(prs)["colors"]["accent1"]
    slide = _blank(prs)
    for hexval in (accent1, "AB12CD"):
        for _ in range(MIN_COLOR_COUNT):
            _filled(slide, hexval)

    named = build_profile(prs, "t", "t")["config"]["color_palette"]["named_colors"]
    ref_by_hex = {c["hex"]: c["theme_ref"] for c in named}
    assert ref_by_hex[accent1] == "accent1"
    assert ref_by_hex["AB12CD"] is None


def test_master_logo_outranks_a_hand_stamped_one(make_prs):
    prs = make_prs()
    _stamp(prs.slide_masters[0], _png((10, 20, 30)))
    for _ in range(4):
        _blank(prs).shapes.add_picture(_png((200, 100, 50)), Emu(IN), Emu(IN),
                                       Emu(HALF), Emu(HALF))

    marks = find_brand_marks(prs)
    assert [m["scope"] for m in marks] == ["master", "slides"]
    assert marks[1]["slide_count"] == 4


def test_full_bleed_picture_is_not_brand_furniture(make_prs):
    prs = make_prs()
    for _ in range(4):
        _blank(prs).shapes.add_picture(_png((7, 7, 7)), Emu(0), Emu(0),
                                       prs.slide_width, prs.slide_height)
    assert find_brand_marks(prs) == []


def test_picture_below_the_stamp_threshold_is_content_not_a_logo(make_prs):
    prs = make_prs()
    for _ in range(MIN_LOGO_SLIDE_STAMPS - 1):
        _blank(prs).shapes.add_picture(_png((9, 9, 9)), Emu(IN), Emu(IN),
                                       Emu(HALF), Emu(HALF))
    assert find_brand_marks(prs) == []


def test_logo_at_different_spots_on_each_layout_is_flagged_as_varying(make_prs):
    prs = make_prs()
    layouts = list(prs.slide_masters[0].slide_layouts)
    _stamp(layouts[0], _png((1, 2, 3)), left=IN, top=IN)
    _stamp(layouts[1], _png((1, 2, 3)), left=5 * IN, top=IN)

    logo = find_brand_marks(prs)[0]
    assert logo["scope"] == "layouts"
    assert len(logo["layouts"]) == 2
    assert logo["position_varies"]


def test_logo_stamped_at_one_spot_has_a_home_position(make_prs):
    prs = make_prs()
    layouts = list(prs.slide_masters[0].slide_layouts)
    for layout in layouts[:2]:
        _stamp(layout, _png((1, 2, 3)), left=3 * IN, top=2 * IN)

    logo = find_brand_marks(prs)[0]
    assert not logo["position_varies"]
    assert logo["position_emu"] == {"left": 3 * IN, "top": 2 * IN,
                                    "width": HALF, "height": HALF}


def test_build_profile_carries_theme_and_brand(make_prs):
    prs = make_prs()
    _stamp(prs.slide_masters[0], _png((4, 5, 6)))
    cfg = build_profile(prs, "t", "t")["config"]
    assert cfg["theme"]["colors"]["accent1"]
    assert cfg["brand"]["logo"]["scope"] == "master"
    assert cfg["brand"]["logo_alternates"] == []
