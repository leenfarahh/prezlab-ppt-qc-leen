"""The ask box: one door onto what the passes recorded.

It answers and it navigates. It never acts (design lead, 26/08/2026), and that
is the first thing these tests hold: there is no path from a question to the
deck, and a link it produces is always a GET to a page that exists.

The other three are the ones that make an answer worth reading:

  - the facts come from the records, so nothing it says can be newer or
    different from what the pages show;
  - a link names a KIND and a slide from those facts, never a URL, so it cannot
    point at another job or at anything that writes;
  - and a claim the facts do not support is DISCARDED rather than shown with a
    hedge, because a designer cannot tell a plausible invention from a fact.

The model is stubbed throughout (qc.chat.ask_json is the seam).
"""

import io
import json

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

import qc.chat as chat
from qc import web
from qc.applymaster import SlidePlan
from qc.design import DesignFinding, Remedy
from qc.layoutgap import Coverage, Gap
from qc.migrate import ContentChange

IN = 914400


# ------------------------------------------------------------------ fixtures


def _deck_bytes(colours=("1F3864", "203965")) -> bytes:
    """A deck with two hand-typed near-navies, so the palette facts have
    something real in them."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i, hexval in enumerate(colours):
        box = slide.shapes.add_textbox(Emu(IN + i * 3 * IN), Emu(IN),
                                       Emu(2 * IN), Emu(IN))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "Confidential client wording"
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor.from_string(hexval)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _audit_job():
    finding = DesignFinding(
        finding_id="f1", kind="palette", severity="warning",
        headline="#203965 is Brand Navy spelled differently",
        detail="two spellings of one color", slides=[0],
        evidence={"places": 2},
        options=[Remedy("snap", "Use Brand Navy (#1F3864) everywhere", "note",
                        op="set_color"),
                 Remedy("leave", "Leave it", "note")])
    return {
        "filename": "client deck.pptx",
        "profile": "prezlab_en",
        "deck": _deck_bytes(),
        "manifest": {"slides": 4, "summary": {},
                     "records": [
                         {"module": "font", "issue_type": "font_family",
                          "severity": "warning", "slide_index": 1},
                         {"module": "font", "issue_type": "font_family",
                          "severity": "warning", "slide_index": 2},
                         {"module": "preflight", "issue_type": "skipped",
                          "severity": "info", "slide_index": 0},
                     ]},
        "design": [finding],
        "design_applied": [],
    }


def _format_job():
    cov = Coverage(slides=6, by_rule={"name": 2, "fallback": 4},
                   gaps=[Gap(label="a title over 2 columns of 2 text blocks",
                             slides=[2, 3, 4, 5], source_layouts=["Blank"],
                             signature={}, closest="Title Only",
                             closest_note="offers 1 content box in 1 column.",
                             refused=4, reviewed=4, asked=1)],
                   used_layouts={"Title Only": 4}, unused_layouts=["Comparison"],
                   reviewed=4, not_reviewed=0, review_ran=True)
    proposal = ContentChange(
        1, "duplicated text left in place",
        "'Our approach' is a second copy of the text now in the title "
        "placeholder", severity="alert", removed_text="Our approach",
        remove_op={"kind": "shape", "slide_index": 1, "shape_id": "7"})
    proposal.remove_id = "r0"
    return {
        "filename": "client deck.pptx",
        "profile": "pif",
        "deck": _deck_bytes(),
        "plans": [SlidePlan(i, "Blank", None, "Title Only",
                            "name" if i < 2 else "fallback")
                  for i in range(6)],
        "errors": {}, "masters": 1, "coverage": cov,
        "changes": [proposal], "removed": [],
    }


@pytest.fixture()
def stub(monkeypatch):
    """A recording stub for the one call this module makes."""
    seen = []

    def _answer(reply):
        def _ask(**kwargs):
            seen.append(kwargs)
            return reply

        monkeypatch.setattr(chat, "ask_json", _ask)
        return seen

    return _answer


_OK = {"answer": "The palette is the thing to look at first.", "links": [],
       "colours": [], "slides": [], "answerable": True}


# --------------------------------------------------------------- the facts


def test_the_facts_come_from_the_records_not_from_the_deck(stub):
    """Nothing it says can be newer or different from what the pages show,
    because it is reading the records the pages render."""
    seen = stub(_OK)
    chat.ask(_audit_job(), "audit", "j1", "what should I fix first?")
    sheet = json.loads(seen[0]["prompt"].split("Fact sheet:\n", 1)[1]
                       .split("\n\nQuestion:")[0])

    assert sheet["kind"] == "audit" and sheet["slides"] == 4
    assert sheet["design_decisions"][0]["id"] == "f1"
    assert sheet["design_decisions"][0]["options"], \
        "the ways out are facts too: they are what the designer will tick"
    # rolled up by kind, not one row per record
    assert sheet["audit_findings_by_kind"][0]["count"] == 2
    assert not [e for e in sheet["audit_findings_by_kind"]
                if e["module"] == "preflight"], "preflight is not a finding"


def test_the_words_on_the_slides_are_never_sent(stub):
    """Same rule as qc.assist. A deck's copy is the part a client would mind."""
    seen = stub(_OK)
    chat.ask(_audit_job(), "audit", "j1", "what does slide 1 say?")
    prompt = seen[0]["prompt"]
    assert "Confidential client wording" not in prompt
    assert "what_you_cannot_see" in prompt, (
        "and the model is told what it cannot see, so it says so instead of "
        "inventing an answer about the copy")


def test_the_palette_is_read_from_the_deck_and_says_what_was_typed(stub):
    seen = stub(_OK)
    chat.ask(_audit_job(), "audit", "j1", "which colours are hand-typed?")
    sheet = json.loads(seen[0]["prompt"].split("Fact sheet:\n", 1)[1]
                       .split("\n\nQuestion:")[0])
    palette = sheet["palette"]
    assert palette["hand_typed_colours"] >= 2
    assert {"1F3864", "203965"} <= {c["hex"] for c in palette["colours"]}


def test_a_deck_no_longer_in_memory_says_so_rather_than_guessing(stub):
    seen = stub(_OK)
    job = _audit_job()
    job["deck"] = None
    chat.ask(job, "audit", "j1", "which colours are hand-typed?")
    sheet = json.loads(seen[0]["prompt"].split("Fact sheet:\n", 1)[1]
                       .split("\n\nQuestion:")[0])
    assert "no longer held in memory" in sheet["palette"]


def test_the_coverage_report_is_handed_over_whole(stub):
    seen = stub(_OK)
    chat.ask(_format_job(), "format", "j2", "is the master missing a layout?")
    sheet = json.loads(seen[0]["prompt"].split("Fact sheet:\n", 1)[1]
                       .split("\n\nQuestion:")[0])
    gap = sheet["coverage"]["gaps"][0]
    assert gap["how_many"] == 4 and gap["closest_layout"] == "Title Only"
    assert gap["checked_and_refused"] == 4
    assert sheet["coverage"]["unused_layouts"] == ["Comparison"]


def test_what_is_waiting_for_a_tick_is_a_fact_too(stub):
    """The proposals are the actionable half of a format run, and a designer
    asking "what is left to do" is asking about these."""
    seen = stub(_OK)
    chat.ask(_format_job(), "format", "j2", "what is left to decide?")
    sheet = json.loads(seen[0]["prompt"].split("Fact sheet:\n", 1)[1]
                       .split("\n\nQuestion:")[0])
    waiting = sheet["removals_waiting_for_a_tick"]
    assert len(waiting) == 1 and waiting[0]["slide"] == 2


# ---------------------------------------------------------------- the links


def test_a_link_is_built_from_a_kind_and_a_slide_never_from_a_url(stub):
    stub({**_OK, "links": [{"kind": "slide_card", "slide": 3,
                            "label": "Open slide 3"}]})
    out = chat.ask(_audit_job(), "audit", "j1", "where do I fix it?")
    assert out["links"] == [{"label": "Open slide 3",
                             "href": "/design/j1?n=2"}]


def test_a_kind_this_job_does_not_have_is_dropped(stub):
    """A format job has no design cards and an audit job has no before/after
    review. A link to the wrong one is a dead end, and a dead end in an answer
    costs more trust than a missing link."""
    stub({**_OK, "links": [{"kind": "review_deck", "label": "Review"},
                           {"kind": "audit_report", "label": "The report"}]})
    out = chat.ask(_audit_job(), "audit", "j1", "where?")
    assert out["links"] == [{"label": "The report", "href": "/audit/j1"}]


def test_an_invented_link_kind_is_dropped(stub):
    stub({**_OK, "links": [{"kind": "apply_all_fixes", "label": "Fix it"},
                           {"kind": "/format/j1/remove", "label": "Remove"}]})
    out = chat.ask(_audit_job(), "audit", "j1", "just fix it")
    assert out["links"] == [], "the closed set is what keeps this read-only"


def test_a_slide_link_past_the_end_of_the_deck_is_dropped(stub):
    stub({**_OK, "links": [{"kind": "slide_card", "slide": 40,
                            "label": "Open slide 40"}]})
    out = chat.ask(_audit_job(), "audit", "j1", "where?")
    assert out["links"] == []


def test_at_most_three_links_come_back(stub):
    stub({**_OK, "links": [{"kind": "slide_card", "slide": n,
                            "label": f"Slide {n}"} for n in (1, 2, 3, 4)]})
    out = chat.ask(_audit_job(), "audit", "j1", "where?")
    assert len(out["links"]) == 3


# ------------------------------------------------------- checking the answer


def test_an_answer_naming_a_colour_the_deck_does_not_use_is_discarded(stub):
    """The failure this check exists for: a confident sentence about #FF0000 on
    a deck that has no red in it. Plausible, wrong, and indistinguishable from a
    fact to the person reading it."""
    stub({**_OK, "answer": "The odd one out is #FF0000, used on two shapes.",
          "colours": ["FF0000"]})
    out = chat.ask(_audit_job(), "audit", "j1", "which colour is off?")
    assert out["checked"] is False
    assert out["answerable"] is False
    assert "#FF0000" in out["answer"] and "discarded" in out["answer"]
    assert out["links"] == []


def test_an_answer_naming_a_slide_past_the_end_is_discarded(stub):
    stub({**_OK, "answer": "Look at slide 19 first.", "slides": [19]})
    out = chat.ask(_audit_job(), "audit", "j1", "where do I start?")
    assert out["checked"] is False
    assert "slide 19" in out["answer"] and "deck of 4" in out["answer"]


def test_a_colour_in_the_prose_is_checked_even_when_it_is_not_declared(stub):
    """Either list alone can be got round by accident. A hex in the sentence is
    a claim a designer reads as a fact whether the model declared it or not."""
    stub({**_OK, "answer": "#00FF00 appears twice.", "colours": []})
    out = chat.ask(_audit_job(), "audit", "j1", "which colour?")
    assert out["checked"] is False


def test_a_colour_the_deck_really_uses_passes(stub):
    stub({**_OK, "answer": "#203965 is a second spelling of #1F3864.",
          "colours": ["203965", "1F3864"], "slides": [1]})
    out = chat.ask(_audit_job(), "audit", "j1", "which colours clash?")
    assert out["checked"] is True
    assert out["answer"].startswith("#203965")


def test_colours_are_not_checked_when_the_palette_could_not_be_read(stub):
    """With the deck's bytes gone there is nothing to check against, and
    inventing a complaint is as bad as missing one."""
    stub({**_OK, "answer": "#FF0000 is off-palette.", "colours": ["FF0000"]})
    job = _audit_job()
    job["deck"] = None
    out = chat.ask(job, "audit", "j1", "which colour is off?")
    assert out["checked"] is True


def test_an_empty_question_is_not_a_call(stub):
    seen = stub(_OK)
    out = chat.ask(_audit_job(), "audit", "j1", "   ")
    assert seen == [], "nothing to ask, nothing sent"
    assert out["answerable"] is False


def test_the_answer_schema_names_things_and_never_describes_a_change():
    """The assistant acts (design lead, 27/08/2026), and the guarantee moved
    rather than went away: the reply can ASK for something from a closed list
    of names and ids, and there is still no field anywhere in it for an
    operation, a shape id, a coordinate or a value to write. It chooses from a
    menu; qc.actions resolves the choice against the job's real records."""
    props = set(chat.ANSWER_SCHEMA["properties"])
    assert props == {"answer", "links", "colours", "slides", "answerable",
                     "action"}
    assert chat.ANSWER_SCHEMA["additionalProperties"] is False
    link = chat.ANSWER_SCHEMA["properties"]["links"]["items"]
    assert set(link["properties"]) == {"kind", "label", "slide"}
    assert link["additionalProperties"] is False

    action = chat.ANSWER_SCHEMA["properties"]["action"]
    assert action["additionalProperties"] is False
    assert set(action["properties"]) == {
        "name", "issue_types", "slide", "finding", "remedy", "findings",
        "removals", "include_holds"}
    assert not {"op", "shape", "shape_id", "value", "colour", "color", "emu",
                "left", "top", "width", "height"} & set(action["properties"])


def test_asking_is_still_not_doing(stub):
    """The confirmation gate did not move. chat.ask resolves an action into a
    PLAN and hands it back; nothing in this module can perform one."""
    stub({**_OK, "action": {"name": "fix_findings"}})
    job = _audit_job()
    before = job["deck"]
    out = chat.ask(job, "audit", "j1", "fix what you can")
    assert job["deck"] is before
    assert "plan" in out
    assert not hasattr(chat, "perform")


def test_no_link_template_writes_anything():
    """Every route it can point at is a GET that renders a page. A POST here
    would be a way around the confirmation gate."""
    for kind, targets in chat._LINKS.items():
        for name, (url, why) in targets.items():
            assert url.startswith("/"), (kind, name)
            assert "remove" not in url and "restore" not in url
            assert "apply" not in url and "fix" not in url and "undo" not in url
            assert why, "a designer has to be told where a link goes"


# ---------------------------------------------------------------- the route


def _client(monkeypatch):
    monkeypatch.setattr(web, "AUTH_REQUIRED", False)
    web.app.state.auth_required = False
    return TestClient(web.app)


def test_the_route_answers_for_a_format_job(monkeypatch):
    monkeypatch.setattr(web, "AI_ENABLED", True, raising=False)
    monkeypatch.setattr("qc.llm.api_configured", lambda: True)
    monkeypatch.setattr(chat, "ask_json",
                        lambda **k: {**_OK, "answer": "Four slides want a "
                                                      "two-column layout.",
                                     "links": [{"kind": "review_deck",
                                                "label": "Review the deck"}]})
    client = _client(monkeypatch)
    web._format_jobs["chatjob"] = _format_job()
    try:
        r = client.post("/chat/chatjob", json={"q": "what is missing?"})
        assert r.status_code == 200
        body = r.json()
        assert "two-column" in body["answer"]
        assert body["links"] == [{"label": "Review the deck",
                                  "href": "/format/chatjob/review?view=deck"}]
    finally:
        web._format_jobs.pop("chatjob", None)


def test_the_route_refuses_when_no_model_is_configured(monkeypatch):
    monkeypatch.setattr(web, "AI_ENABLED", True, raising=False)
    monkeypatch.setattr("qc.llm.api_configured", lambda: False)
    client = _client(monkeypatch)
    web._format_jobs["chatjob2"] = _format_job()
    try:
        r = client.post("/chat/chatjob2", json={"q": "what is missing?"})
        assert r.status_code == 503
        assert "No model key" in r.json()["error"]
    finally:
        web._format_jobs.pop("chatjob2", None)


def test_the_route_refuses_when_ai_is_switched_off(monkeypatch):
    monkeypatch.setattr(web, "AI_ENABLED", False, raising=False)
    r = _client(monkeypatch).post("/chat/anything", json={"q": "hello"})
    assert r.status_code == 503


def test_an_unknown_job_is_a_clean_404(monkeypatch):
    monkeypatch.setattr(web, "AI_ENABLED", True, raising=False)
    monkeypatch.setattr("qc.llm.api_configured", lambda: True)
    r = _client(monkeypatch).post("/chat/deadbeef", json={"q": "hello"})
    assert r.status_code == 404
    assert "no longer held in memory" in r.json()["error"]


def test_a_model_failure_is_not_an_answer_of_nothing_found(monkeypatch):
    """"There is nothing wrong with your deck" and "the model could not be
    reached" mean opposite things to a designer."""
    monkeypatch.setattr(web, "AI_ENABLED", True, raising=False)
    monkeypatch.setattr("qc.llm.api_configured", lambda: True)

    def _boom(**kwargs):
        raise RuntimeError("unreachable")

    monkeypatch.setattr(chat, "ask_json", _boom)
    client = _client(monkeypatch)
    web._format_jobs["chatjob3"] = _format_job()
    try:
        r = client.post("/chat/chatjob3", json={"q": "what is missing?"})
        assert r.status_code == 503
        assert "could not be answered" in r.json()["error"]
    finally:
        web._format_jobs.pop("chatjob3", None)


# ------------------------------------------------ master and audit in one box
#
# The point of the note (design lead, 26/08/2026): a designer's question does not
# know which of the five pages they are on. "Is the master missing a layout?"
# used to be answerable only from a format job, so asking it on the design page
# got "not available" about something the two files state between them. It is
# planned here instead, deterministically, with no render and no model.


def _spec_missing_two_column_layouts():
    from qc.stylespec import extract_style_spec

    spec = extract_style_spec(Presentation(), source="m.pptx",
                              embed_assets=False)
    spec["layouts"] = [l for l in spec["layouts"]
                       if l["type"] in ("titleOnly", "secHead")]
    return spec


def _two_column_deck(slides=3) -> bytes:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    for _ in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        head = slide.shapes.add_textbox(Emu(IN), Emu(IN // 2), Emu(10 * IN),
                                        Emu(IN))
        head.text_frame.text = "Heading"
        for col in (1, 7):
            box = slide.shapes.add_textbox(Emu(col * IN), Emu(2 * IN),
                                           Emu(5 * IN), Emu(3 * IN))
            box.text_frame.text = "Point one"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_the_master_question_is_answerable_from_an_audit_job(stub):
    """No coverage was computed for this job and none needs to have been: the
    layouts and the planner are both deterministic, so the answer is planned on
    the way to answering the question."""
    seen = stub(_OK)
    job = _audit_job()
    job["deck"] = _two_column_deck(3)
    job["master_spec"] = _spec_missing_two_column_layouts()
    job["manifest"]["slides"] = 3

    chat.ask(job, "audit", "j1", "is this master missing a layout?")
    sheet = json.loads(seen[0]["prompt"].split("Fact sheet:\n", 1)[1]
                       .split("\n\nQuestion:")[0])
    cov = sheet["coverage"]
    assert cov["no_layout"] == 3
    assert "2 columns" in cov["gaps"][0]["wants"]
    assert cov["slides_were_looked_at"] is False, \
        "nothing was rendered, so the facts must not imply anything was seen"


def test_a_run_with_no_master_says_there_is_nothing_to_compare(stub):
    """The honest answer, and a different one from "the master is fine"."""
    seen = stub(_OK)
    job = _audit_job()
    job["profile"] = None
    chat.ask(job, "audit", "j1", "is this master missing a layout?")
    sheet = json.loads(seen[0]["prompt"].split("Fact sheet:\n", 1)[1]
                       .split("\n\nQuestion:")[0])
    assert "not available" in sheet["coverage"]
    assert "not judged against a master" in sheet["coverage"]


def test_a_format_job_keeps_the_coverage_it_already_computed(stub):
    """Planning it twice could disagree with the run it describes. The one the
    format pass produced is the one that was applied."""
    seen = stub(_OK)
    chat.ask(_format_job(), "format", "j2", "what is missing?")
    sheet = json.loads(seen[0]["prompt"].split("Fact sheet:\n", 1)[1]
                       .split("\n\nQuestion:")[0])
    assert sheet["coverage"]["gaps"][0]["checked_and_refused"] == 4, \
        "the review's own evidence has to survive, and a fresh plan has none"
