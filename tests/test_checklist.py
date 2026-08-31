"""The colour and type checklist: what the deck is made of.

Not an audit. Nothing here is a finding, nothing has a tick, and nothing changes
the deck - it is the page a designer keeps open beside PowerPoint. What the tests
hold is the one thing it knows that PowerPoint will not show: the LEVEL every
colour and typeface comes from. On screen an explicit hex and a theme reference
are identical; in the file they are the difference between a deck that survives a
rebrand and one that gets hunted through by hand.
"""

import io

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from lxml import etree

from qc import web
from qc.extract import extract_deck, font_inventory, palette_inventory
from qc.ui_checklist import render_checklist

IN = 914400
A = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ------------------------------------------------------------------ fixtures


def _deck_bytes() -> bytes:
    """One typed navy, one shape on accent1, and type set two ways: a run that
    states Arial and a run that inherits the theme's own."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    themed = slide.shapes.add_shape(1, Emu(IN), Emu(IN), Emu(2 * IN), Emu(IN))
    spPr = themed._element.spPr
    for tag in ("a:solidFill", "a:noFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    ln = spPr.find(qn("a:ln"))
    fill = etree.fromstring(
        f'<a:solidFill xmlns:a="{A}"><a:schemeClr val="accent1"/></a:solidFill>')
    if ln is not None:
        ln.addprevious(fill)
    else:
        spPr.append(fill)

    typed = slide.shapes.add_shape(1, Emu(4 * IN), Emu(IN), Emu(2 * IN), Emu(IN))
    typed.fill.solid()
    typed.fill.fore_color.rgb = RGBColor(0x1F, 0x38, 0x64)

    stated = slide.shapes.add_textbox(Emu(IN), Emu(3 * IN), Emu(4 * IN),
                                      Emu(IN // 2))
    run = stated.text_frame.paragraphs[0].add_run()
    run.text = "Set on the run"
    run.font.name = "Arial"
    run.font.size = Pt(11)

    inherited = slide.shapes.add_textbox(Emu(6 * IN), Emu(3 * IN), Emu(4 * IN),
                                        Emu(IN // 2))
    run = inherited.text_frame.paragraphs[0].add_run()
    run.text = "Nothing stated here"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _rolled():
    deck = extract_deck(_deck_bytes())
    return palette_inventory(deck), font_inventory(deck)


# --------------------------------------------------------------- the type roll


def test_type_set_on_a_run_is_told_apart_from_type_that_is_inherited():
    """The distinction the whole page exists for. Both draw the same on screen;
    only one of them survives being restyled onto a new master."""
    _palette, fonts = _rolled()
    by_hand = [f for f in fonts["fonts"] if f["set_by_hand"]]
    inherited = [f for f in fonts["fonts"] if not f["set_by_hand"]]

    assert by_hand and inherited, fonts["fonts"]
    assert any(f["family"] == "Arial" for f in by_hand)
    assert all(f["from"].startswith(("run", "paragraph")) for f in by_hand)
    assert fonts["set_by_hand"] == len(by_hand)


def test_an_inherited_typeface_still_has_a_name():
    """python-pptx answers None for it, which is the trap: a checklist reading
    only what the run states would report a deck with no typography in it."""
    _palette, fonts = _rolled()
    for entry in fonts["fonts"]:
        assert entry["family"], entry
        assert entry["from"], "and it says which level supplied it"


def test_the_theme_fonts_are_stated_separately():
    """What the brand SAYS, as against what the slides do."""
    _palette, fonts = _rolled()
    theme = fonts["theme_fonts"]
    assert theme["major"]["latin"] and theme["minor"]["latin"]


def test_sizes_travel_with_the_family():
    """A designer asks what the body copy is, not which sizes exist."""
    _palette, fonts = _rolled()
    arial = next(f for f in fonts["fonts"] if f["family"] == "Arial"
                 and f["set_by_hand"])
    assert arial["sizes"][0]["pt"] == 11.0


# ------------------------------------------------------------- the page itself


def test_the_page_says_how_many_colours_were_typed_by_hand():
    palette, fonts = _rolled()
    html = render_checklist(deck_name="d.pptx", job_id="j1", back="/design/j1",
                            palette=palette, fonts=fonts)
    assert "typed in by hand" in html
    assert "a rebrand will not reach" in html, \
        "the consequence is the point, not the count"


def test_the_page_shows_the_theme_slots_with_swatches():
    palette, fonts = _rolled()
    html = render_checklist(deck_name="d.pptx", job_id="j1", back="/design/j1",
                            palette=palette, fonts=fonts)
    assert "accent1" in html
    assert "background:#" in html, "a palette without swatches is a table of hex"


def test_the_page_separates_what_the_slides_sit_on():
    """"Which colours are explicit for backgrounds" was asked separately from
    "which colours exist", because a hex on one slide's background beats the
    master and that is a different problem."""
    palette, fonts = _rolled()
    html = render_checklist(deck_name="d.pptx", job_id="j1", back="/design/j1",
                            palette=palette, fonts=fonts)
    assert "What the slides sit on" in html
    assert "background (master)" in html, "the level has to be named"
    assert "inherited" in html, (
        "this deck states no background of its own, which is the state to want "
        "and has to be visible as such rather than as an absence")


def test_the_page_offers_no_way_to_change_anything():
    """It is a checklist. A tick here would make it a fifth place where a deck
    can be edited, and the confirmation cards are the only place that should
    be."""
    palette, fonts = _rolled()
    html = render_checklist(deck_name="d.pptx", job_id="j1", back="/design/j1",
                            palette=palette, fonts=fonts)
    assert "<form" not in html, "a form here would make this a fifth editor"
    assert "<input" not in html and "<button" not in html
    assert "nothing here changes the deck" in html


# ---------------------------------------------------------------- the route


def _client(monkeypatch):
    monkeypatch.setattr(web, "AUTH_REQUIRED", False)
    web.app.state.auth_required = False
    return TestClient(web.app)


def test_the_route_serves_a_format_job(monkeypatch):
    client = _client(monkeypatch)
    web._format_jobs["chkjob"] = {"deck": _deck_bytes(),
                                  "filename": "client deck.pptx"}
    try:
        r = client.get("/checklist/chkjob")
        assert r.status_code == 200
        assert "client deck.pptx" in r.text
        assert "typed in by hand" in r.text
        # back to where the designer came from, not to the audit page
        assert "/format/chkjob/review" in r.text
    finally:
        web._format_jobs.pop("chkjob", None)


def test_a_deck_no_longer_in_memory_says_so(monkeypatch):
    client = _client(monkeypatch)
    web._format_jobs["chkjob2"] = {"deck": None, "filename": "d.pptx"}
    try:
        r = client.get("/checklist/chkjob2")
        assert r.status_code == 410
        assert "no longer held in memory" in r.text
    finally:
        web._format_jobs.pop("chkjob2", None)


def test_an_unknown_job_is_a_clean_404(monkeypatch):
    r = _client(monkeypatch).get("/checklist/deadbeef")
    assert r.status_code == 404


def test_the_chat_can_point_at_it():
    """The gap this closed: the chat could talk about the palette and had
    nowhere to send anybody."""
    from qc.chat import _LINKS

    for kind in ("audit", "format"):
        assert "checklist" in _LINKS[kind]
        url, why = _LINKS[kind]["checklist"]
        assert url == "/checklist/{job}" and why


# ------------------------------------------------ the extraction is cached
#
# Both the checklist page and the ask box's palette facts re-read the deck from
# bytes on EVERY request - a full parse plus a colour and font resolution of
# every run on every slide, about 1.4s for a 200-slide deck, for a document that
# cannot change between requests unless a fix lands (30/08/2026). Cached on the
# job; the danger is a cache that outlives the bytes it described.


def test_the_extraction_is_read_once_per_deck():
    from qc import web

    job = {"deck": _deck_bytes(), "filename": "d.pptx"}
    first = web._extracted(job)
    second = web._extracted(job)
    assert first is second, "the second read is the first read"
    assert job.get("extracted") is first


def test_a_fix_drops_the_cached_extraction():
    """A recolour changes the palette roll-up. A checklist showing the old one
    beside a row marked applied is the same lie as a stale thumbnail, which is
    why this rides on the existing render invalidation rather than a second
    rule that could be forgotten."""
    from qc import web

    job = {"deck": _deck_bytes(), "filename": "d.pptx"}
    web._extracted(job)
    assert job.get("extracted") is not None

    web._invalidate_renders(job)
    assert job.get("extracted") is None, (
        "the extraction is derived from the deck bytes and goes stale with "
        "every other derived thing")


def test_no_deck_means_no_extraction_rather_than_an_error():
    from qc import web

    assert web._extracted({"deck": None}) is None
