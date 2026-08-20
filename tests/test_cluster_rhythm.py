"""Cluster rhythm: a same-size cohort (image column) with ONE broken gap
gets its tail lifted back onto the rhythm, satellites (labels, underlines)
riding along. Modeled on the ground-truth wheel slide: gaps 3/3/3/14.6mm,
designer lifted the bottom cluster 11.6mm as one move."""

import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

import qc.modules.margin_alignment as ma
from qc.fixer import apply_fixes, is_fixable
from tests.conftest import save_and_ctx

IN = 914400
MM = 36000


def _col_with_broken_rhythm(prs, n=4, gap=3 * MM, odd_gap=12 * MM,
                            w=32 * MM, h=22 * MM, left=None):
    # default: distinct left per slide, so multi-slide test decks do not
    # accidentally form cross-slide recurring anchors
    if left is None:
        left = (12 + 7 * len(prs.slides._sldIdLst)) * MM
    """n same-size 'images' stacked with even gaps except the last one,
    each with a label textbox beside it."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    top = 48 * MM
    shapes = []
    for i in range(n):
        if i == n - 1:
            top += odd_gap - gap  # break the rhythm before the last one
        pic = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left),
                                     Emu(top), Emu(w), Emu(h))
        label = slide.shapes.add_textbox(Emu(left + w + int(2.8 * MM)),
                                         Emu(top + 2 * MM),
                                         Emu(43 * MM), Emu(5 * MM))
        label.text_frame.text = f"label {i}"
        shapes.append((pic, label))
        top += h + gap
    return slide, shapes


def _rhythm(ctx):
    return [r.to_dict() for r in ma.detect(ctx)
            if r.issue_type == "margin_alignment.cluster_rhythm"]


def test_broken_rhythm_lifts_tail_with_satellites(make_prs, en_profile,
                                                  tmp_path):
    prs = make_prs()
    slide, shapes = _col_with_broken_rhythm(prs)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    recs = _rhythm(ctx)
    assert len(recs) == 1
    rec = recs[0]
    last_pic, last_label = shapes[-1]
    assert rec["shape_id"] == str(last_pic.shape_id)
    assert rec["severity"] == "error"        # 4 members, 4x the median gap
    assert rec["locator"] == f"lift-col:{last_pic.shape_id}"
    assert is_fixable(rec)

    pic_before, label_before = last_pic.top, last_label.top
    delta = int(rec["new_value"]) - int(rec["old_value"])
    assert delta == -(12 - 3) * MM

    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {rec["record_id"]})
    assert result.applied == 1
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    by_id = {s.shape_id: s for s in out.shapes}
    assert by_id[last_pic.shape_id].top == pic_before + delta
    # the label traveled with its image: composition preserved
    assert by_id[last_label.shape_id].top == label_before + delta
    # and the rhythm is restored: a re-audit stays quiet
    ctx2 = save_and_ctx(Presentation(io.BytesIO(result.cleaned_bytes)),
                        tmp_path, en_profile, name="clean.pptx")
    assert _rhythm(ctx2) == []


def test_rhythm_guards(make_prs, en_profile, tmp_path):
    prs = make_prs()
    # two odd gaps: ambiguous, no record
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    top = 48 * MM
    for gap in (3 * MM, 9 * MM, 14 * MM):
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(12 * MM), Emu(top),
                               Emu(32 * MM), Emu(22 * MM))
        top += 22 * MM + gap
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(12 * MM), Emu(top),
                           Emu(32 * MM), Emu(22 * MM))
    # a gap >= 6x median reads as a section break: no record
    _col_with_broken_rhythm(prs, odd_gap=20 * MM)
    # only two members: no rhythm to claim
    _col_with_broken_rhythm(prs, n=2)
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert _rhythm(ctx) == []


def test_rhythm_suppresses_bandline_duplicate(make_prs, en_profile, tmp_path):
    """The same broken column must not ALSO get the weaker band-line
    uneven_spacing nudge for the same shape."""
    prs = make_prs()
    slide, shapes = _col_with_broken_rhythm(prs)
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    last_pic = shapes[-1][0]
    dupes = [r for r in ma.detect(ctx)
             if r.issue_type == "margin_alignment.uneven_spacing"
             and r.shape_id == str(last_pic.shape_id)
             and (r.locator or "").startswith("col:")]
    assert dupes == []


def test_recurring_furniture_never_rides_a_lift(make_prs, en_profile,
                                                tmp_path):
    """A footer logo that overlaps the lifted image (same name and spot on
    3+ slides) must stay pinned; a unique label beside the image rides.
    Real-deck finding, 21/07/2026: 'Strategy&' travelled up onto a photo."""
    prs = make_prs()
    H = prs.slide_height
    # image column ending deep in the slide, rhythm broken before the last
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tops = (75, 100, 125, 162)  # gaps 3/3/15mm around 22mm images
    for i, top in enumerate(tops):
        pic = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(12 * MM),
                                     Emu(top * MM), Emu(32 * MM), Emu(22 * MM))
        last = pic
    label = slide.shapes.add_textbox(Emu(47 * MM), Emu(171 * MM),
                                     Emu(43 * MM), Emu(5 * MM))
    label.text_frame.text = "unique label"
    logo = slide.shapes.add_textbox(Emu(12 * MM), Emu(181 * MM),
                                    Emu(12 * MM), Emu(3 * MM))
    logo.text_frame.text = "logo"
    logo.name = "FooterLogo"
    # the same logo on two more slides makes it recurring furniture
    for _ in range(2):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        lg = s.shapes.add_textbox(Emu(12 * MM), Emu(181 * MM),
                                  Emu(12 * MM), Emu(3 * MM))
        lg.text_frame.text = "logo"
        lg.name = "FooterLogo"
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    recs = [r for r in _rhythm(ctx) if r["slide_index"] == 0]
    assert len(recs) == 1
    logo_before, label_before = logo.top, label.top
    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {recs[0]["record_id"]})
    assert result.applied == 1
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    by_id = {s.shape_id: s for s in out.shapes}
    delta = int(recs[0]["new_value"]) - int(recs[0]["old_value"])
    assert by_id[label.shape_id].top == label_before + delta  # rides
    assert by_id[logo.shape_id].top == logo_before            # pinned


def test_lift_unblocks_off_canvas_footer(make_prs, en_profile, tmp_path):
    """The wheel-slide scenario end to end: the too-low image cluster
    occupies the footer zone, so the off-canvas source line can only move
    up after the lift. Applying both in one run must fix both (lift runs
    before footer by fix ordering)."""
    import qc.modules.header_footer as hf

    prs = make_prs()
    H = prs.slide_height
    slide, shapes = _col_with_broken_rhythm(prs, odd_gap=14 * MM)
    # source line fully below the slide edge, overlapping where the lifted
    # cluster's old position was
    src = slide.shapes.add_textbox(Emu(12 * MM), Emu(H), Emu(150 * MM),
                                   Emu(int(7.6 * MM)))
    src.text_frame.text = "Source: analysis"
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    recs = _rhythm(ctx) + [r.to_dict() for r in hf.detect(ctx)
                           if r.issue_type == "header_footer.footer_off_canvas"]
    assert len(recs) == 2
    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {r["record_id"] for r in recs})
    assert result.applied == 2, [o.reason for o in result.outcomes]
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    moved_src = next(s for s in out.shapes
                     if getattr(s, "has_text_frame", False)
                     and s.text_frame.text.startswith("Source"))
    assert moved_src.top + moved_src.height <= H
