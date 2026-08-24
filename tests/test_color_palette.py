"""Tests for the color_palette detection module (planted-violation decks)."""

import copy

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from qc.modules.color_palette import detect
from qc.profile import Profile
from spike.color_resolver import ciede2000
from tests.conftest import save_and_ctx

NAVY = (0x1F, 0x4E, 0x79)


def _hx(hexstr):
    return tuple(int(hexstr[i:i + 2], 16) for i in (0, 2, 4))


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_filled_shape(slide, hexstr, left_in=1.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_in),
                                Inches(1), Inches(2), Inches(1))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string(hexstr)
    return sh


def test_on_palette_fills_yield_zero(make_prs, en_profile, tmp_path):
    """Exact palette hex and a deltaE ~0.4 near-hit are both compliant."""
    prs = make_prs()
    slide = _blank_slide(prs)
    _add_filled_shape(slide, "1F4E79", left_in=1)
    assert ciede2000(_hx("214F7A"), NAVY) <= 2.0
    _add_filled_shape(slide, "214F7A", left_in=4)
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert detect(ctx) == []


def test_near_miss_fill_warns_with_replacement(make_prs, en_profile, tmp_path):
    """deltaE in (2, 5]: warning, high confidence, new_value = nearest hex."""
    de = ciede2000(_hx("2A5A8A"), NAVY)
    assert 2.0 < de <= 5.0
    prs = make_prs()
    _add_filled_shape(_blank_slide(prs), "2A5A8A")
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = detect(ctx)
    assert len(recs) == 1
    r = recs[0]
    assert r.issue_type == "color_palette.off_palette_rgb"
    assert r.severity == "warning"
    assert r.confidence == "high"
    assert r.property == "spPr.solidFill"
    assert r.old_value == "2A5A8A"
    assert r.new_value == "1F4E79"
    assert "prezlab_navy" in r.message
    assert r.action == "flagged"
    assert r.slide_index == 0
    assert r.profile_rule_id == "color_palette.named_colors"


def test_ambiguity_band_fill_names_the_target_but_asks_first(make_prs,
                                                            en_profile,
                                                            tmp_path):
    """deltaE in (5, 10]: warning, medium confidence, and the nearest palette
    colour NAMED as the target.

    It used to carry none, so the row read "no automatic fix" for a colour the
    tool could name and reach (design lead, 24/08/2026). What changes past the
    auto-replace band is not whether a fix exists but who decides: the swap is
    visible, so it is never pre-ticked."""
    from qc.fixer import is_fixable, tick_reason

    de = ciede2000(_hx("305E90"), NAVY)
    assert 5.0 < de <= 10.0
    prs = make_prs()
    _add_filled_shape(_blank_slide(prs), "305E90")
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = detect(ctx)
    assert len(recs) == 1
    rec = recs[0]
    assert rec.issue_type == "color_palette.off_palette_rgb"
    assert rec.severity == "warning"
    assert rec.confidence == "medium"
    assert rec.new_value is not None, "the nearest palette colour is the target"
    assert rec.old_value == "305E90"
    assert is_fixable(rec.to_dict())
    assert "visible change" in tick_reason(rec.to_dict())


def test_wild_color_is_error(make_prs, en_profile, tmp_path):
    prs = make_prs()
    _add_filled_shape(_blank_slide(prs), "FF00FF")
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = detect(ctx)
    assert len(recs) == 1
    assert recs[0].issue_type == "color_palette.off_palette_rgb"
    # Severity and confidence answer different questions here: the tool is
    # CERTAIN this colour is not in the palette (error) and only guessing that
    # the nearest one is what was meant (medium). Confidence is about the fix.
    assert recs[0].severity == "error"
    assert recs[0].confidence == "medium"
    assert recs[0].new_value is not None
    assert "no near palette match" in recs[0].message


def test_theme_accent1_fill_is_compliant(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1),
                                Inches(2), Inches(1))
    sh.fill.solid()
    sh.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert detect(ctx) == []


def test_removed_theme_slot_is_flagged(make_prs, en_profile, tmp_path):
    raw = copy.deepcopy(en_profile.raw)
    slots = raw["config"]["color_palette"]["theme_color_slots"]
    raw["config"]["color_palette"]["theme_color_slots"] = [
        s for s in slots if s != "accent1"]
    profile = Profile(raw)

    prs = make_prs()
    slide = _blank_slide(prs)
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1),
                                Inches(2), Inches(1))
    sh.fill.solid()
    sh.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    ctx = save_and_ctx(prs, tmp_path, profile)
    recs = detect(ctx)
    assert len(recs) == 1
    r = recs[0]
    assert r.issue_type == "color_palette.disallowed_theme_slot"
    assert r.severity == "warning"
    assert r.confidence == "high"
    assert r.old_value == "accent1"
    assert r.profile_rule_id == "color_palette.theme_color_slots"


def test_run_text_color_off_palette(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Quarterly summary"
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.color.rgb = RGBColor.from_string("FF00FF")
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = detect(ctx)
    assert len(recs) == 1
    r = recs[0]
    assert r.property == "rPr.solidFill"
    assert r.issue_type == "color_palette.off_palette_rgb"
    assert r.severity == "error"
    assert r.old_value == "FF00FF"
    assert r.arabic_flag is False


def test_arabic_shape_records_carry_arabic_flag(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    sh = _add_filled_shape(slide, "2A5A8A")
    sh.text_frame.text = "مرحبا"  # Arabic text
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = detect(ctx)
    assert len(recs) == 1
    assert recs[0].issue_type == "color_palette.off_palette_rgb"
    assert recs[0].arabic_flag is True
    assert recs[0].action == "flagged"


def test_clean_control_deck_zero_records(make_prs, en_profile, tmp_path):
    """No explicit solid fills off-palette, no explicit run colors, plus a
    gradient-free default-styled shape and a noFill textbox: zero records."""
    prs = make_prs()
    slide = _blank_slide(prs)
    _add_filled_shape(slide, "C00000", left_in=1)  # prezlab_red exact
    box = slide.shapes.add_textbox(Inches(4), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "On-brand text with inherited color"
    slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(1),
                           Inches(1), Inches(1))  # style-ref fill only
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert detect(ctx) == []
