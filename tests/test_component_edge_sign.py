"""A component LEFT of the thing it should line up with is still misaligned.

`synthesize` measured a signed offset and kept only positive ones. That is
right against the FRAME, where a component sticking out past the presentation
space is a margin breach owned by another rule. Against another COMPONENT it
threw away half the findings: the vision pass named the intended line
correctly, the geometry measured it correctly, and a negative number dropped
the record on the floor.

The case that found it: a two-column comparison whose left heading sat 9mm
left of the column body under it. The model answered "these two share a left
edge with the slide title" and the pass reported nothing (31/08/2026).
"""

import io

import pytest
from pptx import Presentation
from pptx.util import Emu, Pt

from qc.components import synthesize

IN = 914400
SLIDE_W, SLIDE_H = 12192000, 6858000


def _slide(boxes):
    """boxes: [(left_in, top_in, w_in, h_in, text)] -> (slide, [shape_id])."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(SLIDE_W), Emu(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ids = []
    for left, top, w, h, text in boxes:
        shape = slide.shapes.add_textbox(Emu(int(left * IN)), Emu(int(top * IN)),
                                         Emu(int(w * IN)), Emu(int(h * IN)))
        shape.text_frame.text = text
        shape.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
        ids.append(str(shape.shape_id))
    # reopen so the test reads geometry the way the pass does
    buf = io.BytesIO()
    prs.save(buf)
    return Presentation(buf).slides[0], ids


def _comparison_slide():
    """The real one: title and body on 0.9in, the heading out at 0.55in."""
    return _slide([
        (0.9, 0.5, 11.5, 0.8, "Comparing Methods"),
        (0.55, 1.6, 2.6, 0.4, "Method 1"),     # 0.35in left of its own column
        (0.9, 2.2, 5.0, 3.4, "1 2 3 4 5"),
    ])


def _layout(ids, anchor="Slide Title"):
    """The answer the vision model actually returned for that slide."""
    return {
        "components": [
            {"name": "Slide Title", "shape_ids": [ids[0]]},
            {"name": "Method 1 Heading", "shape_ids": [ids[1]]},
            {"name": "Method 1 Body", "shape_ids": [ids[2]]},
        ],
        "alignments": [{
            "axis": "left",
            "components": ["Method 1 Heading", "Method 1 Body"],
            "anchor": anchor,
            "rationale": "The column heading and body should align with the "
                         "start of the slide title.",
        }],
    }


def test_a_heading_left_of_its_column_is_reported():
    slide, ids = _comparison_slide()

    out = synthesize(slide, 0, _layout(ids), None, [], SLIDE_W, SLIDE_H)

    misaligned = [r for r in out
                  if r["issue_type"] == "margin_alignment.component_edge_misaligned"]
    assert len(misaligned) == 1, \
        "the heading sits 0.35in left of its anchor and was not reported"
    rec = misaligned[0]
    assert rec["shape_id"] == ids[1], "it named the wrong shape"
    # the offset is reported as a distance, never as a negative number
    assert "8.9mm" in rec["message"], rec["message"]
    assert "-" not in rec["message"].split("sits")[1][:12]
    # and it snaps to the anchor's edge, not to its own
    assert int(rec["new_value"]) == int(0.9 * IN)


def test_the_member_already_on_the_line_is_left_alone():
    """One record, for the one that is off. The body is already on 0.9in."""
    slide, ids = _comparison_slide()

    out = synthesize(slide, 0, _layout(ids), None, [], SLIDE_W, SLIDE_H)

    assert [r["shape_id"] for r in out
            if r["issue_type"].endswith("component_edge_misaligned")] == [ids[1]]


def test_a_heading_right_of_its_column_is_reported_too():
    """The direction was never the point; the distance is."""
    slide, ids = _slide([
        (0.9, 0.5, 11.5, 0.8, "Comparing Methods"),
        (1.25, 1.6, 2.6, 0.4, "Method 1"),     # 0.35in RIGHT this time
        (0.9, 2.2, 5.0, 3.4, "1 2 3 4 5"),
    ])

    out = synthesize(slide, 0, _layout(ids), None, [], SLIDE_W, SLIDE_H)

    misaligned = [r for r in out
                  if r["issue_type"].endswith("component_edge_misaligned")]
    assert len(misaligned) == 1
    assert "8.9mm" in misaligned[0]["message"]


def test_a_component_outboard_of_the_FRAME_is_still_left_to_the_margin_rule():
    """The one-sided test was right for the frame and stays. A shape sticking
    out past the presentation space is a margin breach, and reporting it here
    as well would give a designer two findings and two fixes for one defect."""
    slide, ids = _slide([
        (0.9, 0.5, 11.5, 0.8, "Comparing Methods"),
        (0.20, 1.6, 2.6, 0.4, "Method 1"),     # outside the frame's left edge
        (0.9, 2.2, 5.0, 3.4, "1 2 3 4 5"),
    ])
    frame = (int(0.9 * IN), int(0.4 * IN),
             SLIDE_W - int(0.9 * IN), SLIDE_H - int(0.4 * IN))

    out = synthesize(slide, 0, _layout(ids, anchor="frame"), frame, [],
                     SLIDE_W, SLIDE_H)

    assert not [r for r in out
                if r["issue_type"].endswith("component_edge_misaligned")], \
        "a breach of the frame belongs to the margin rule, not to this one"


def test_a_gap_too_big_to_be_drift_is_still_a_composition():
    """Past the window it is a layout decision, whichever side it falls."""
    slide, ids = _slide([
        (0.9, 0.5, 11.5, 0.8, "Comparing Methods"),
        (4.0, 1.6, 2.6, 0.4, "Method 1"),      # 3.1in out: nobody nudged that
        (0.9, 2.2, 5.0, 3.4, "1 2 3 4 5"),
    ])

    out = synthesize(slide, 0, _layout(ids), None, [], SLIDE_W, SLIDE_H)

    assert not [r for r in out
                if r["issue_type"].endswith("component_edge_misaligned")]


@pytest.mark.parametrize("offset_in", [0.02, 0.0])
def test_a_shape_on_the_line_is_never_reported(offset_in):
    """Inside the tolerance is on the line, on either side of it."""
    slide, ids = _slide([
        (0.9, 0.5, 11.5, 0.8, "Comparing Methods"),
        (0.9 - offset_in, 1.6, 2.6, 0.4, "Method 1"),
        (0.9, 2.2, 5.0, 3.4, "1 2 3 4 5"),
    ])

    out = synthesize(slide, 0, _layout(ids), None, [], SLIDE_W, SLIDE_H)

    assert not [r for r in out
                if r["issue_type"].endswith("component_edge_misaligned")]
