"""The component orchestration layer: the model decides what the things are and
which line they belong on, code measures and computes every target.

Two decisions in the geometry pipeline are not geometry questions, and this
layer exists because answering them with arithmetic is what produced the
alignment bugs:

WHAT IS ONE THING - guessed today from overlap and adjacency, which fails both
ways: a corner rule welded to a photo gets left behind, and a column of stacked
blocks decides it is carrying its own neighbours so every member moves twice.

WHICH ONE IS WRONG - taken today from the median, so on a slide where one
element sits on the master's stated line and three drifted off it, the majority
wins and the tool proposes pulling the correct one off the line to join them
(design lead, 24/08/2026).

The API is stubbed throughout. What is under test is the PRECISION GATE: that
code re-verifies everything the model says, drops what does not check out, and
never lets a model-supplied number reach the deck.
"""

import io

import pytest
from pptx import Presentation
from pptx.util import Emu

import qc.components as C
from qc.fixer import apply_fixes, is_fixable, needs_explicit_tick, tick_reason

IN = 914400
BLANK = 6
# a frame whose left edge is 0.5in and top edge 1.9in
SPACE = (int(0.5 * IN), int(1.9 * IN), int(12.83 * IN), int(6.8 * IN))


def _deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    return prs


def _box(slide, x, y, w, h):
    return slide.shapes.add_shape(1, Emu(int(x * IN)), Emu(int(y * IN)),
                                  Emu(int(w * IN)), Emu(int(h * IN)))


def _card(slide, x, y):
    """A card, its icon and its label: three shapes, one component."""
    card = _box(slide, x, y, 2.4, 1.6)
    icon = _box(slide, x + 0.15, y + 0.15, 0.4, 0.4)
    label = _box(slide, x + 0.15, y + 0.8, 2.1, 0.5)
    return card, icon, label


def _synth(prs, layout, space=SPACE):
    slide = prs.slides[0]
    return C.synthesize(slide, 0, layout, space, [],
                        prs.slide_width, prs.slide_height)


# ------------------------------------------------- what the model is asked


def test_the_inventory_carries_geometry_but_never_the_words(make_prs):
    """The image already shows the text; the question here is about boxes, and
    a slide's copy is the part a client would mind being sent."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    shape = _box(slide, 1, 2, 3, 1)
    shape.text_frame.text = "Confidential client wording"

    inv = C.inventory(slide, prs.slide_width, prs.slide_height)
    assert len(inv) == 1
    entry = inv[0]
    assert entry["id"] == str(shape.shape_id)
    assert entry["text"] is True, "whether it holds text is useful"
    assert "Confidential" not in repr(inv), "the words must not be sent"
    # normalized, so the numbers mean the same thing as the image
    assert 0 < entry["x"] < 1 and 0 < entry["y"] < 1


def test_the_frame_is_described_in_the_inventorys_own_units():
    """"frame" has to be a line the model can see, not a word it has to trust."""
    note = C._frame_note(SPACE, 12192000, 6858000)
    assert "left x=0.0375" in note and "top y=0.2533" in note

    assert "not an available anchor" in C._frame_note(None, 12192000, 6858000)


# --------------------------------------------------- the entity, moved as one


def test_a_component_is_measured_and_moved_as_one_thing(make_prs):
    """The card, its icon and its label share a bounding box and one delta, so
    the arrangement inside the component survives the move."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    left_card = _card(slide, 0.5, 2.0)          # on the frame's left edge
    low_card = _card(slide, 0.72, 4.0)          # 0.22in inboard of it

    layout = {
        "components": [
            {"name": "top card", "shape_ids": [str(s.shape_id) for s in left_card]},
            {"name": "lower card", "shape_ids": [str(s.shape_id) for s in low_card]},
        ],
        "alignments": [{"axis": "left", "anchor": "top card",
                        "components": ["top card", "lower card"],
                        "rationale": "one column, one left edge"}],
    }
    recs = _synth(prs, layout)
    assert len(recs) == 1, f"expected one decision, got {len(recs)}"
    rec = recs[0]
    assert rec["issue_type"] == "margin_alignment.component_edge_misaligned"
    assert int(rec["new_value"]) == int(0.5 * IN)
    members = set(rec["locator"].split(":", 2)[2].split(","))
    assert members == {str(s.shape_id) for s in low_card}, \
        "the whole component travels, and only it"

    buf = io.BytesIO()
    prs.save(buf)
    gaps = [(s.left - low_card[0].left, s.top - low_card[0].top)
            for s in low_card]
    out = apply_fixes(buf.getvalue(), recs, {rec["record_id"]})
    assert out.applied == 1, [o.reason for o in out.outcomes]
    after = {str(s.shape_id): s
             for s in Presentation(io.BytesIO(out.cleaned_bytes)).slides[0].shapes}
    moved = [after[str(s.shape_id)] for s in low_card]
    assert moved[0].left == int(0.5 * IN), "the component is not on the line"
    assert [(s.left - moved[0].left, s.top - moved[0].top) for s in moved] == gaps, \
        "the arrangement inside the component changed"


def test_nothing_in_a_component_moves_twice(make_prs):
    """The bug this layer exists to end: stacked members are each other's
    neighbours, so geometry-inferred carrying moved every one of them once for
    itself and again for the member beside it."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    anchor = _box(slide, 0.5, 2.0, 3.0, 0.6)
    stack = [_box(slide, 0.72, 3.0 + i * 0.7, 3.0, 0.6) for i in range(3)]

    layout = {
        "components": [
            {"name": "heading", "shape_ids": [str(anchor.shape_id)]},
            {"name": "body block", "shape_ids": [str(s.shape_id) for s in stack]},
        ],
        "alignments": [{"axis": "left", "anchor": "frame",
                        "components": ["heading", "body block"],
                        "rationale": "the column starts on the frame"}],
    }
    recs = _synth(prs, layout)
    buf = io.BytesIO()
    prs.save(buf)
    out = apply_fixes(buf.getvalue(), recs, {r["record_id"] for r in recs})
    after = Presentation(io.BytesIO(out.cleaned_bytes)).slides[0]
    lefts = sorted({s.left for s in after.shapes})
    assert lefts == [int(0.5 * IN)], \
        f"left edges {lefts}: something moved twice or was left behind"


# ------------------------------------------- the stated line beats the median


def test_the_anchor_is_the_reference_even_when_it_is_outvoted(make_prs):
    """One element on the frame, three drifted off it. The median says the
    three are right; the anchor says the frame is. Only the three move."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    on_line = _box(slide, 0.5, 2.0, 2.0, 0.6)
    drifted = [_box(slide, 0.72, 2.8 + i * 0.7, 2.0, 0.6) for i in range(3)]

    layout = {
        "components": [{"name": "first", "shape_ids": [str(on_line.shape_id)]}]
        + [{"name": f"block {i}", "shape_ids": [str(s.shape_id)]}
           for i, s in enumerate(drifted)],
        "alignments": [{"axis": "left", "anchor": "frame",
                        "components": ["first"] + [f"block {i}" for i in range(3)],
                        "rationale": "one column on the frame"}],
    }
    recs = _synth(prs, layout)
    moved = {r["shape_id"] for r in recs}
    assert moved == {str(s.shape_id) for s in drifted}, \
        "the element already on the line must not be asked to move"
    assert all(int(r["new_value"]) == int(0.5 * IN) for r in recs)


def test_a_component_used_as_the_anchor_is_never_asked_to_move(make_prs):
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    ref = _box(slide, 3.0, 2.0, 2.0, 0.6)
    other = _box(slide, 3.3, 3.0, 2.0, 0.6)
    layout = {
        "components": [{"name": "ref", "shape_ids": [str(ref.shape_id)]},
                       {"name": "other", "shape_ids": [str(other.shape_id)]}],
        "alignments": [{"axis": "left", "anchor": "ref",
                        "components": ["ref", "other"], "rationale": "x"}],
    }
    recs = _synth(prs, layout)
    assert {r["shape_id"] for r in recs} == {str(other.shape_id)}


# --------------------------------------------------------- the precision gate


@pytest.mark.parametrize("layout,why", [
    ({"components": [{"name": "a", "shape_ids": ["99999"]},
                     {"name": "b", "shape_ids": ["99998"]}],
      "alignments": [{"axis": "left", "anchor": "a", "components": ["a", "b"],
                      "rationale": "x"}]},
     "shape ids that are not on the slide"),
    ({"components": [], "alignments": [{"axis": "left", "anchor": "frame",
                                        "components": ["ghost"],
                                        "rationale": "x"}]},
     "components that were never defined"),
    ({"components": [{"name": "a", "shape_ids": []}], "alignments": []},
     "a component naming nothing"),
])
def test_what_does_not_check_out_is_dropped(layout, why, make_prs):
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    _box(slide, 1, 2, 2, 1)
    _box(slide, 4, 2, 2, 1)
    assert _synth(prs, layout) == [], f"{why} produced a record"


def test_an_anchor_naming_nothing_is_not_an_anchor(make_prs):
    """Including "frame" on a deck whose master states no frame."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    a, b = _box(slide, 0.5, 2, 2, 1), _box(slide, 0.9, 3, 2, 1)
    layout = {
        "components": [{"name": "a", "shape_ids": [str(a.shape_id)]},
                       {"name": "b", "shape_ids": [str(b.shape_id)]}],
        "alignments": [{"axis": "left", "anchor": "frame",
                        "components": ["a", "b"], "rationale": "x"}],
    }
    assert C.synthesize(prs.slides[0], 0, layout, None, [],
                        prs.slide_width, prs.slide_height) == []


def test_a_line_already_held_produces_nothing(make_prs):
    """Below the perceptual floor there is nothing to fix, whatever the model
    thought it saw."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    a = _box(slide, 0.5, 2, 2, 1)
    b = _box(slide, 0.5 + (C.TOL_EMU / 2) / IN, 3, 2, 1)
    layout = {
        "components": [{"name": "a", "shape_ids": [str(a.shape_id)]},
                       {"name": "b", "shape_ids": [str(b.shape_id)]}],
        "alignments": [{"axis": "left", "anchor": "a",
                        "components": ["a", "b"], "rationale": "x"}],
    }
    assert _synth(prs, layout) == []


def test_a_gap_too_big_to_move_is_reported_rather_than_dropped(make_prs):
    """Past the window the model's answer is SHOWN, not overruled.

    This used to assert the opposite - a 3in gap produced nothing, on the
    grounds that it must be a deliberate indent. That is a ruler overruling a
    judgment: the model looked at the slide and said these two were meant to
    share a left edge, and WINDOW_EMU is a threshold for deciding intent in the
    passes where nothing supplies it. Here something does. The further out a
    component sat, the more certain the silence, which is backwards.

    What the window still decides is whether it is a ONE-CLICK move. Past it
    the record carries no target and low confidence, so it is not tickable
    (qc.fixer.is_fixable) - the designer sees what the model saw and decides.
    """
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    a = _box(slide, 0.5, 2, 2, 1)
    b = _box(slide, 3.5, 3, 2, 1)          # 3in in
    layout = {
        "components": [{"name": "a", "shape_ids": [str(a.shape_id)]},
                       {"name": "b", "shape_ids": [str(b.shape_id)]}],
        "alignments": [{"axis": "left", "anchor": "a",
                        "components": ["a", "b"], "rationale": "x"}],
    }
    recs = _synth(prs, layout)
    assert len(recs) == 1, "the model's answer was thrown away by a ruler"
    rec = recs[0]
    assert rec["confidence"] == "low"
    assert rec["new_value"] is None, "no target: this is not a one-click move"
    assert "too far to move for you" in rec["message"]
    assert not is_fixable(rec)


def test_a_rotated_member_makes_its_component_unmeasurable(make_prs):
    """A rotated shape's stored box is not its rendered one. It still belongs
    to the component and still travels with it; it just cannot be the evidence
    for a move."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    a = _box(slide, 0.5, 2, 2, 1)
    b = _box(slide, 0.9, 3, 2, 1)
    b.rotation = 15
    layout = {
        "components": [{"name": "a", "shape_ids": [str(a.shape_id)]},
                       {"name": "b", "shape_ids": [str(b.shape_id)]}],
        "alignments": [{"axis": "left", "anchor": "a",
                        "components": ["a", "b"], "rationale": "x"}],
    }
    assert _synth(prs, layout) == []


def test_a_shape_sticking_out_past_the_frame_is_not_this_rules_business(
        make_prs):
    """That is a margin breach, and the safe-zone rule owns it. Reporting it
    here would put two findings and two different fixes on one shape."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    a = _box(slide, 0.5, 2, 2, 1)
    b = _box(slide, 0.2, 3, 2, 1)          # OUTSIDE the frame
    layout = {
        "components": [{"name": "a", "shape_ids": [str(a.shape_id)]},
                       {"name": "b", "shape_ids": [str(b.shape_id)]}],
        "alignments": [{"axis": "left", "anchor": "frame",
                        "components": ["a", "b"], "rationale": "x"}],
    }
    assert _synth(prs, layout) == []


# --------------------------------------------------------- who decides what


def test_a_model_judged_fix_is_offered_and_never_pre_ticked(make_prs):
    """The numbers are the tool's - code measured them. The claim that these
    shapes are one component and that this is the intended line came from
    the model, and a design judgment is the designer's to confirm."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    a = _box(slide, 0.5, 2, 2, 1)
    b = _box(slide, 0.72, 3, 2, 1)
    layout = {
        "components": [{"name": "a", "shape_ids": [str(a.shape_id)]},
                       {"name": "b", "shape_ids": [str(b.shape_id)]}],
        "alignments": [{"axis": "left", "anchor": "a",
                        "components": ["a", "b"], "rationale": "x"}],
    }
    rec = _synth(prs, layout)[0]
    assert is_fixable(rec), "a checked judgment should still be offerable"
    assert needs_explicit_tick(rec)
    assert "judgment is yours" in tick_reason(rec)


def test_the_model_never_supplies_a_coordinate():
    """The schema has no place to put one. Ask a model for EMU and you get
    plausible EMU; the only numbers in a record are the ones code measured."""
    fields = set()

    def walk(node):
        if isinstance(node, dict):
            for key, value in node.items():
                if key == "properties" and isinstance(value, dict):
                    fields.update(value)
                walk(value)
        elif isinstance(node, list):
            for item in node:
                walk(item)

    walk(C.LAYOUT_SCHEMA)
    # "group" is a yes or no about an entity the model already named, not a
    # measurement, so it belongs on this list rather than breaking the rule it
    # protects (design lead, 26/08/2026).
    assert fields == {"components", "alignments", "name", "shape_ids",
                      "group", "axis", "anchor", "rationale"}, fields


def test_the_axis_and_anchor_vocabularies_are_closed():
    axis = C.LAYOUT_SCHEMA["properties"]["alignments"]["items"]["properties"]["axis"]
    assert axis["enum"] == ["top", "left", "right"]
    for node in (C.LAYOUT_SCHEMA,
                 C.LAYOUT_SCHEMA["properties"]["components"]["items"],
                 C.LAYOUT_SCHEMA["properties"]["alignments"]["items"]):
        assert node["additionalProperties"] is False


# --------------------------------------------------------------- the run loop


def test_one_bad_slide_never_sinks_the_run(monkeypatch, make_prs):
    """An advisory layer over a pipeline that works without it."""
    prs = _deck()
    for _ in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
        _box(slide, 0.5, 2.0, 2.0, 0.6)
        _box(slide, 0.72, 3.0, 2.0, 0.6)
        _box(slide, 0.72, 3.8, 2.0, 0.6)
    buf = io.BytesIO()
    prs.save(buf)

    calls = []

    def _flaky(png, inv, frame_note):
        calls.append(1)
        if len(calls) == 2:
            raise RuntimeError("the API had a bad day")
        return {"components": [], "alignments": []}

    monkeypatch.setattr(C, "_ask_vision", _flaky)
    _recs, reviewed = C.run_components(
        buf.getvalue(), {i: b"png" for i in range(3)}, {"records": []}, SPACE)
    assert len(calls) == 3, "the run stopped at the failure"
    assert reviewed == 2, "a failed slide must not be counted as reviewed"


def test_a_slide_with_nothing_to_group_is_not_sent(monkeypatch, make_prs):
    """Two shapes is not a composition, and a vision call costs money."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    _box(slide, 1, 2, 2, 1)
    buf = io.BytesIO()
    prs.save(buf)

    monkeypatch.setattr(C, "_ask_vision",
                        lambda *a: pytest.fail("a 1-shape slide was sent"))
    _recs, reviewed = C.run_components(buf.getvalue(), {0: b"png"},
                                       {"records": []}, SPACE)
    assert reviewed == 0


def test_a_slide_with_no_render_is_skipped(monkeypatch, make_prs):
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    for i in range(3):
        _box(slide, 1 + i, 2, 0.8, 1)
    buf = io.BytesIO()
    prs.save(buf)
    monkeypatch.setattr(C, "_ask_vision",
                        lambda *a: pytest.fail("sent without a picture"))
    _recs, reviewed = C.run_components(buf.getvalue(), {}, {"records": []},
                                       SPACE)
    assert reviewed == 0


# ------------------------------------------------------------------ the route


def _client(monkeypatch):
    from fastapi.testclient import TestClient

    from qc import web

    monkeypatch.setattr(web, "AUTH_REQUIRED", False)
    web.app.state.auth_required = False
    return TestClient(web.app), web


def test_the_route_needs_a_key_and_says_so(monkeypatch, fixtures_dir):
    from tests.conftest import job_id_of

    client, web = _client(monkeypatch)
    # qc.llm, not qc.assist. The route used to gate on the assistant's own key
    # check while run_components went through qc.llm, so the button refused for
    # a key the pass does not use. Patching the seam the route ACTUALLY calls
    # is the test.
    import qc.llm as llm

    monkeypatch.setattr(llm, "api_configured", lambda: False)
    deck = (fixtures_dir / "mixed_layouts.pptx").read_bytes()
    r = client.post("/audit", files={"deck": ("d.pptx", deck, "app/x")},
                    data={"profile": "prezlab_en"})
    job = job_id_of(r)
    out = client.post(f"/components/{job}")
    assert out.status_code == 200
    assert "needs a model key" in out.text
    # and it names the variable this host would actually read
    assert "GEMINI_API_KEY" in out.text


def test_a_configured_host_is_not_asked_for_another_vendors_key(monkeypatch,
                                                                fixtures_dir):
    """The regression itself: a working Gemini key must open the button."""
    from tests.conftest import job_id_of

    client, web = _client(monkeypatch)
    import qc.llm as llm

    monkeypatch.setenv("GEMINI_API_KEY", "test-key")
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    assert llm.api_configured(), "the fixture must present a configured host"

    deck = (fixtures_dir / "mixed_layouts.pptx").read_bytes()
    r = client.post("/audit", files={"deck": ("d.pptx", deck, "app/x")},
                    data={"profile": "prezlab_en"})
    out = client.post(f"/components/{job_id_of(r)}")

    assert out.status_code == 200
    assert "ANTHROPIC_API_KEY" not in out.text
    assert "needs a model key" not in out.text


def test_the_report_offers_the_review_and_discloses_the_images(
        monkeypatch, fixtures_dir):
    from tests.conftest import job_id_of

    client, _web = _client(monkeypatch)
    deck = (fixtures_dir / "mixed_layouts.pptx").read_bytes()
    r = client.post("/audit", files={"deck": ("d.pptx", deck, "app/x")},
                    data={"profile": "prezlab_en"})
    report = client.get(f"/audit/{job_id_of(r)}").text
    assert "Component review" in report
    assert "/components/" in report
    assert "one thing" in report, "the two questions must be stated"
    assert "never supplies a coordinate" in report


# ------------------------------------------ the words the picture does not show
#
# Design lead, 26/08/2026: the model should be able to read white text. It
# cannot, and neither can anyone - white text on a white ground is not in the
# image. What it CAN be given is the fact, out of the file, beside the picture.


def _white_on_white():
    """A white panel with white text on it, and a legible line below. The first
    is really there and is invisible; the second is the control."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    panel = _box(slide, 1, 1, 5, 2)
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    hidden = slide.shapes.add_textbox(Emu(int(1.2 * IN)), Emu(int(1.2 * IN)),
                                      Emu(int(4 * IN)), Emu(int(0.5 * IN)))
    run = hidden.text_frame.paragraphs[0].add_run()
    run.text = "A heading nobody can see"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    shown = slide.shapes.add_textbox(Emu(IN), Emu(int(4 * IN)),
                                     Emu(int(4 * IN)), Emu(int(0.5 * IN)))
    run = shown.text_frame.paragraphs[0].add_run()
    run.text = "A line that reads perfectly well"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x11, 0x11, 0x11)

    buf = io.BytesIO()
    prs.save(buf)
    reopened = Presentation(io.BytesIO(buf.getvalue()))
    return reopened, str(hidden.shape_id), str(shown.shape_id)


def test_text_the_render_will_not_show_is_flagged_in_the_inventory():
    """A box of white-on-white text reads as EMPTY in the picture, so a
    component built from the picture alone leaves the label off the card it
    belongs to. The file knows better and says so."""
    prs, hidden_id, shown_id = _white_on_white()
    slide = prs.slides[0]
    inv = C.inventory(slide, prs.slide_width, prs.slide_height, prs)
    by_id = {entry["id"]: entry for entry in inv}

    assert by_id[hidden_id].get("invisible") is True
    assert by_id[hidden_id]["text"] is True, "it is still a text block"
    assert "invisible" not in by_id[shown_id], \
        "a legible line must not be flagged; the flag would then mean nothing"


def test_the_prompt_explains_what_the_flag_means():
    """A flag in the payload that nothing explains is noise, and noise in a
    prompt is worse than an omission: the model has to guess what it means."""
    prs, hidden_id, _shown = _white_on_white()
    inv = C.inventory(prs.slides[0], prs.slide_width, prs.slide_height, prs)
    note = C._hidden_note(inv)
    assert hidden_id in note
    assert "not in the picture" in note
    assert "not as empty boxes" in note

    assert C._hidden_note([{"id": "1"}]) == "", \
        "nothing hidden, nothing said"


def test_the_inventory_still_works_without_the_presentation():
    """The colour resolution needs the presentation. Without it the flag is
    ABSENT rather than wrong: a missing flag understates what is on the slide,
    and a wrong one sends the model looking for words that are plainly there."""
    prs, hidden_id, _shown = _white_on_white()
    inv = C.inventory(prs.slides[0], prs.slide_width, prs.slide_height)
    assert inv, "the inventory is still produced"
    assert all("invisible" not in entry for entry in inv)


def test_the_flag_carries_no_words_either():
    """Same rule as the rest of the inventory: the model is told a box holds
    text it cannot see, never what the text says."""
    prs, _hidden, _shown = _white_on_white()
    inv = C.inventory(prs.slides[0], prs.slide_width, prs.slide_height, prs)
    payload = repr(inv) + C._hidden_note(inv)
    assert "nobody can see" not in payload and "reads perfectly" not in payload


# ============================================ what should be ONE OBJECT
#
# Design lead, 26/08/2026: the model should also decide when to group elements
# and what to treat as an entity. It already answers the first half - naming the
# components is what this pass is for - and that answer was only ever used to
# move things TOGETHER. It was never written into the file, so the next person to
# open the deck drags the card and leaves the icon behind.
#
# The judgment is the model's and the write is code's, as everywhere here. What
# these tests hold is the verification in between, and that a group lands without
# moving anything.


def _card_deck(make_prs=None):
    """A card, its icon and its label: three shapes, one object."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    card, icon, label = _card(slide, 1, 1)
    label.text_frame.text = "Card label"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), [str(s.shape_id) for s in (card, icon, label)]


def _reopen(blob):
    return Presentation(io.BytesIO(blob)).slides[0]


def _records(slide, ids, group=True, name="Card"):
    layout = {"components": [{"name": name, "shape_ids": ids, "group": group}]}
    comps = C._component_boxes(layout["components"], C._boxes(slide))
    return C._grouping_records(slide, 0, layout, comps, set())


def test_a_component_the_model_calls_one_object_is_offered_as_a_group():
    blob, ids = _card_deck()
    recs = _records(_reopen(blob), ids)
    assert len(recs) == 1
    rec = recs[0]
    assert rec["issue_type"] == "margin_alignment.should_be_grouped"
    assert rec["new_value"] == ",".join(ids)
    assert "one object made of 3 shapes" in rec["message"]


def test_a_component_the_model_does_not_want_grouped_is_not_offered():
    """Two cards in a row are two objects that merely align. A group nobody
    wanted puts every element inside it one click further away."""
    blob, ids = _card_deck()
    assert _records(_reopen(blob), ids, group=False) == []


def test_a_single_shape_is_never_a_group():
    blob, ids = _card_deck()
    assert _records(_reopen(blob), ids[:1]) == []


def test_a_rotated_member_stops_the_group():
    """A group takes the bounding box of its members, and a rotated member's
    stored box is not its rendered one. Grouping around it would move what it
    contains."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    card, icon, label = _card(slide, 1, 1)
    label.text_frame.text = "Card label"
    icon.rotation = 15.0
    buf = io.BytesIO()
    prs.save(buf)
    ids = [str(s.shape_id) for s in (card, icon, label)]
    assert _records(_reopen(buf.getvalue()), ids) == []


def test_shapes_already_in_a_group_are_left_alone():
    """Regrouping a group is a different decision and not this one."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    group = slide.shapes.add_group_shape()
    a = group.shapes.add_shape(1, Emu(IN), Emu(IN), Emu(IN), Emu(IN))
    b = group.shapes.add_shape(1, Emu(2 * IN), Emu(IN), Emu(IN), Emu(IN))
    buf = io.BytesIO()
    prs.save(buf)
    ids = [str(a.shape_id), str(b.shape_id)]
    assert _records(_reopen(buf.getvalue()), ids) == []


def test_grouping_is_tickable_and_never_pre_ticked():
    """It changes the shape TREE, which is the most invasive thing this tool
    does to a slide."""
    blob, ids = _card_deck()
    rec = _records(_reopen(blob), ids)[0]
    assert is_fixable(rec)
    assert needs_explicit_tick(rec)
    why = tick_reason(rec)
    assert "one object" in why and "your approval" in why


def test_grouping_moves_nothing():
    """The whole test of a group: it lands and the slide looks identical. The
    child space is set equal to the parent space, so a child's offset means the
    same thing inside the group as it did outside it."""
    from qc.design import placed_shapes

    blob, ids = _card_deck()
    slide = _reopen(blob)
    before = {str(p.shape.shape_id): p.box for p in placed_shapes(slide)}
    rec = _records(slide, ids)[0]

    result = apply_fixes(blob, [rec], {rec["record_id"]})
    assert result.applied == 1, [(o.outcome, o.reason) for o in result.outcomes]

    after = _reopen(result.cleaned_bytes)
    assert len(after.shapes) == 1, "three shapes became one object"
    assert str(after.shapes[0].shape_type).startswith("GROUP")

    now = {str(p.shape.shape_id): p.box
           for p in placed_shapes(after) if p.grouped}
    assert len(now) == 3
    for shape_id, box in now.items():
        assert before[shape_id] == box, f"{shape_id} moved"


def test_grouping_keeps_the_drawing_order_of_its_topmost_member():
    """The group sits where its topmost part sat, so nothing changes what covers
    what."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    under = _box(slide, 0.5, 0.5, 6, 4)
    card, icon, label = _card(slide, 1, 1)
    label.text_frame.text = "Card label"
    over = _box(slide, 4, 3, 2, 2)
    buf = io.BytesIO()
    prs.save(buf)
    blob = buf.getvalue()

    ids = [str(s.shape_id) for s in (card, icon, label)]
    rec = _records(_reopen(blob), ids)[0]
    result = apply_fixes(blob, [rec], {rec["record_id"]})

    after = _reopen(result.cleaned_bytes)
    kinds = [str(s.shape_type).split(" ")[0] for s in after.shapes]
    assert kinds == ["AUTO_SHAPE", "GROUP", "AUTO_SHAPE"], kinds
    assert str(after.shapes[0].shape_id) == str(under.shape_id)
    assert str(after.shapes[2].shape_id) == str(over.shape_id)


def test_a_group_whose_shapes_have_gone_is_skipped_not_guessed():
    """The slide can move on between the review and the tick. Grouping what is
    left would be a guess at what the component was."""
    blob, ids = _card_deck()
    rec = _records(_reopen(blob), ids)[0]
    rec = dict(rec, new_value=",".join(ids + ["999999"]))
    result = apply_fixes(blob, [rec], {rec["record_id"]})
    assert result.applied == 0
    assert "no longer on the slide" in result.outcomes[0].reason
    assert len(_reopen(result.cleaned_bytes).shapes) == 3, "and nothing moved"


def test_the_schema_lets_the_model_say_group_and_nothing_more():
    """It answers a yes or no about an entity it already named. It does not get
    to say where the group goes or what it contains beyond the ids it was
    handed."""
    comp = C.LAYOUT_SCHEMA["properties"]["components"]["items"]
    assert set(comp["properties"]) == {"name", "shape_ids", "group"}
    assert comp["additionalProperties"] is False
    assert comp["properties"]["group"]["type"] == "boolean"


def test_the_call_budget_is_spent_on_calls_not_on_successes(monkeypatch):
    """MAX_SLIDES used to cap the slides that ANSWERED, not the calls made. That
    reads as generous and is the opposite: under an outage nothing answers, so
    the loop kept going and a 200-slide deck made 200 failing calls to review
    nothing (30/08/2026)."""
    import io

    from pptx import Presentation
    from pptx.util import Emu

    import qc.components as CP
    from qc.llm import LLMUnavailable

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    n_slides = CP.MAX_SLIDES + 12
    for _ in range(n_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for i in range(3):
            slide.shapes.add_shape(1, Emu(914400 * (i + 1)), Emu(914400),
                                   Emu(914400), Emu(914400))
    buf = io.BytesIO()
    prs.save(buf)

    calls = []

    def _always_down(**kwargs):
        calls.append(1)
        raise LLMUnavailable("the model could not be reached")

    monkeypatch.setattr(CP, "ask_json", _always_down)
    records, reviewed = CP.run_components(
        buf.getvalue(), {i: b"png" for i in range(n_slides)}, {"records": []})

    assert records == [] and reviewed == 0, (
        "nothing answered, so nothing is reviewed and nothing is clean")
    assert len(calls) <= CP.MAX_SLIDES, (
        f"an outage cost {len(calls)} calls on a {n_slides}-slide deck; the "
        f"budget is {CP.MAX_SLIDES}")
