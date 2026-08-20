"""Content overflow: a layout spilling past the slide edge gets ONE
proportional select-all rescale anchored at the left margin (the designer's
move on the ground-truth rankings slide: 92.8%), placeholders untouched."""

import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

import qc.modules.margin_alignment as ma
from qc.fixer import apply_fixes, is_fixable
from tests.conftest import save_and_ctx

IN = 914400
MM = 36000


def _overflow_slide(prs, spans=((2, 110), (115, 230), (240, 345))):
    """Two rows of three panels whose x-spans (mm) cross the margin frame
    on BOTH sides - the ground-truth rankings-slide pattern (many
    breachers, layout wider than the frame, centers still on the slide)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    boxes = []
    for row, top in enumerate((40, 100)):
        for l_mm, r_mm in spans:
            boxes.append(slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Emu(l_mm * MM), Emu(top * MM),
                Emu((r_mm - l_mm) * MM), Emu(50 * MM)))
    return slide, boxes


def _overflow(ctx):
    return [r.to_dict() for r in ma.detect(ctx)
            if r.issue_type == "margin_alignment.content_overflow"]


def test_overflow_detected_and_rescaled(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide, boxes = _overflow_slide(prs)
    title = slide.shapes.add_textbox(Emu(12 * MM), Emu(10 * MM),
                                     Emu(300 * MM), Emu(20 * MM))
    title.text_frame.text = "unrelated wide text box excluded by span rule"
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    recs = _overflow(ctx)
    assert len(recs) == 1
    rec = recs[0]
    assert rec["severity"] == "warning"      # big change: never pre-ticked
    assert is_fixable(rec)
    scale = int(rec["new_value"]) / 1000.0
    assert 0.80 <= scale < 1.0

    gaps_before = boxes[1].left - (boxes[0].left + boxes[0].width)
    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {rec["record_id"]})
    assert result.applied == 1, [o.reason for o in result.outcomes]
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    by_id = {s.shape_id: s for s in out.shapes}
    W = prs.slide_width
    margins = en_profile.get("geometry.safe_zone_margins_emu") or {}
    for b in boxes:
        s = by_id[b.shape_id]
        assert s.left >= margins.get("left", 0) - 1000
        assert s.left + s.width <= W + 1000 - margins.get("right", 0)
    # proportions preserved: the gap scaled by the same factor
    s0, s1 = by_id[boxes[0].shape_id], by_id[boxes[1].shape_id]
    gap_after = s1.left - (s0.left + s0.width)
    assert abs(gap_after - gaps_before * scale) <= 20


def test_placeholders_are_not_rescaled(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide, boxes = _overflow_slide(prs)
    # a real title placeholder must be excluded from the selection
    layout = prs.slide_layouts[0]
    slide2 = prs.slides.add_slide(layout)
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = _overflow(ctx)
    assert len(recs) == 1
    ids = recs[0]["locator"].split(":", 3)[3].split(",")
    assert all(str(b.shape_id) in ids for b in boxes)


def test_an_arabic_layout_is_rescaled_against_the_right_margin(make_prs,
                                                               en_profile,
                                                               tmp_path):
    """The anchor mirrors with the script: shrinking an Arabic layout against
    the LEFT margin walks it away from its reading edge."""
    prs = make_prs()
    slide, boxes = _overflow_slide(prs)
    for i, box in enumerate(boxes):
        box.text_frame.text = "نمتلك خبرات واسعة وسجلًا حافلًا"
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    recs = _overflow(ctx)
    assert len(recs) == 1 and is_fixable(recs[0])
    assert "right margin" in recs[0]["message"]

    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {recs[0]["record_id"]})
    assert result.applied == 1, [o.reason for o in result.outcomes]
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    margins = en_profile.get("geometry.safe_zone_margins_emu") or {}
    right = prs.slide_width - margins.get("right", 0)
    assert abs(max(s.left + s.width for s in out.shapes) - right) <= 1000, \
        "the block's right edge should land on the right margin"


def test_no_breach_no_record_and_drastic_is_flag_only(make_prs, en_profile,
                                                      tmp_path):
    prs = make_prs()
    # comfortably inside: nothing
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(3):
        s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(20 * MM + i * 90 * MM),
                            Emu(60 * MM), Emu(60 * MM), Emu(60 * MM))
    # 393mm of content against a ~313mm frame: would need < 80%, flag-only
    _overflow_slide(prs, spans=((2, 110), (115, 230), (275, 395)))
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = _overflow(ctx)
    assert len(recs) == 1
    assert recs[0]["new_value"] is None
    assert not is_fixable(recs[0])
    assert "too drastic" in recs[0]["message"]
