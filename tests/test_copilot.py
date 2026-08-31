"""Design copilot: the code-side precision gate over vision observations,
and the route wiring. The model API is never called in tests."""

import io

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Emu

import qc.copilot as copilot
from qc.web import app
from tests.conftest import job_id_of

client = TestClient(app)
IN = 914400


def _slide_with_row(gaps=(500000, 620000, 380000)):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w, h = 1000000, 500000
    boxes, x = [], IN
    for i, gap in enumerate((0,) + tuple(gaps)):
        x += gap + (w if i else 0)
        tb = slide.shapes.add_textbox(Emu(x), Emu(IN), Emu(w), Emu(h))
        tb.text_frame.text = f"card {i}"
        boxes.append(tb)
    return prs, slide, boxes


def test_inventory_shape():
    prs, slide, boxes = _slide_with_row()
    inv = copilot.inventory(slide, prs.slide_width, prs.slide_height)
    assert len(inv) == 4
    assert {"id", "kind", "x", "y", "w", "h", "text"} == set(inv[0])
    assert all(item["text"] for item in inv)


def test_synthesize_distribute_verified_and_emitted():
    prs, slide, boxes = _slide_with_row()
    obs = [{"action": "distribute_row",
            "shape_ids": [str(b.shape_id) for b in boxes],
            "rationale": "Cards in this row have uneven gaps."}]
    recs = copilot.synthesize(slide, 0, obs, existing=[])
    assert len(recs) == 1
    rec = recs[0]
    assert rec["issue_type"] == "margin_alignment.uneven_spacing"
    assert rec["locator"].startswith("dist-row:")
    assert rec["confidence"] == "medium"
    assert "Design copilot" in rec["message"]

    from qc.fixer import is_fixable

    assert is_fixable(rec)


def test_synthesize_drops_unverifiable_observations():
    prs, slide, boxes = _slide_with_row(gaps=(500000, 500000, 500000))
    even_ids = [str(b.shape_id) for b in boxes]
    obs = [
        # already even: nothing to distribute
        {"action": "distribute_row", "shape_ids": even_ids, "rationale": "x"},
        # invented shape ids: dropped
        {"action": "align_left", "shape_ids": ["991", "992", "993"],
         "rationale": "x"},
        # too few shapes: dropped
        {"action": "match_widths", "shape_ids": even_ids[:2], "rationale": "x"},
    ]
    assert copilot.synthesize(slide, 0, obs, existing=[]) == []


def test_synthesize_align_and_match_compute_targets():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    boxes = []
    for i, (left, w) in enumerate([(IN, 1000000), (IN, 1000000),
                                   (IN + 60000, 1300000)]):
        tb = slide.shapes.add_textbox(Emu(left), Emu(IN + i * 800000),
                                      Emu(w), Emu(500000))
        tb.text_frame.text = "x"
        boxes.append(tb)
    ids = [str(b.shape_id) for b in boxes]
    obs = [{"action": "align_left", "shape_ids": ids, "rationale": "r"},
           {"action": "match_widths", "shape_ids": ids, "rationale": "r"}]
    recs = copilot.synthesize(slide, 0, obs, existing=[])
    kinds = {r["issue_type"] for r in recs}
    assert kinds == {"margin_alignment.edge_misaligned",
                     "shape_size.size_mismatch"}
    edge = next(r for r in recs if "edge" in r["issue_type"])
    assert edge["shape_id"] == ids[2] and int(edge["new_value"]) == IN
    size = next(r for r in recs if "size" in r["issue_type"])
    assert size["shape_id"] == ids[2]
    assert size["new_value"] == "1000000x500000"  # width matched, height kept


def test_copilot_route_merges_records(fixtures_dir, monkeypatch):
    deck = (fixtures_dir / "mixed_layouts.pptx").read_bytes()
    r = client.post("/audit",
                    files={"deck": ("mixed_layouts.pptx", deck,
                                    "application/octet-stream")},
                    data={"profile": "prezlab_en"})
    job_id = job_id_of(r)
    before_total = client.get(f"/manifest/{job_id}").json()["summary"]["total"]

    import qc.web as web

    canned = [dict(record_id="cp1", job_id=None, slide_index=0, shape_id="99",
                   shape_path=None, module="margin_alignment",
                   issue_type="margin_alignment.edge_misaligned",
                   property="spPr.xfrm.off.x", old_value="1", new_value="2",
                   severity="warning", action="flagged", confidence="medium",
                   arabic_flag=False, profile_rule_id=None,
                   message="Design copilot: canned.", locator=None,
                   created_at="")]
    monkeypatch.setattr("qc.llm.api_configured", lambda: True)
    monkeypatch.setattr("qc.copilot.run_copilot",
                        lambda deck, thumbs, manifest: (canned, 3))
    monkeypatch.setattr(web, "_ensure_thumbs", lambda jid, job:
                        job.__setitem__("thumbs", {0: b"png"}))

    r = client.post(f"/copilot/{job_id}")
    assert r.status_code == 200
    assert "reviewed 3 slides and added 1 suggestion" in r.text
    after = client.get(f"/manifest/{job_id}").json()["summary"]["total"]
    assert after == before_total + 1


def test_copilot_route_without_key(fixtures_dir, monkeypatch):
    deck = (fixtures_dir / "mixed_layouts.pptx").read_bytes()
    r = client.post("/audit",
                    files={"deck": ("mixed_layouts.pptx", deck,
                                    "application/octet-stream")},
                    data={"profile": "prezlab_en"})
    job_id = job_id_of(r)
    monkeypatch.setattr("qc.llm.api_configured", lambda: False)
    r = client.post(f"/copilot/{job_id}")
    assert r.status_code == 200
    assert "needs an API key for the configured model provider" in r.text


def test_the_copilot_asks_through_the_one_door(monkeypatch):
    """It built its own model client until 30/08/2026, which made it the
    only judgment pass with no timeout, no retry and no truncation guard - and
    put a second model id in the config. What this pins is the seam: the pass
    states a system prompt, a schema and an image, and qc.llm decides who
    answers."""
    import qc.copilot as CP

    seen = {}

    def _fake(*, system, prompt, schema, images=None, max_tokens=8192):
        seen.update(system=system, prompt=prompt, schema=schema,
                    images=images)
        return {"observations": [{"action": "align_left",
                                  "shape_ids": ["1", "2", "3"],
                                  "rationale": "three cards, one edge"}]}

    monkeypatch.setattr(CP, "ask_json", _fake)
    out = CP._ask_vision(b"png", [{"id": "1"}])

    assert out == [{"action": "align_left", "shape_ids": ["1", "2", "3"],
                    "rationale": "three cards, one edge"}]
    assert seen["schema"] is CP.OBSERVATIONS_SCHEMA, "its own schema, unaltered"
    assert seen["images"] == [b"png"]
    assert not hasattr(CP, "anthropic"), "no SDK is imported in this module"


def test_an_unreachable_model_skips_the_slide_rather_than_sinking_the_run(
        monkeypatch):
    """qc.llm raises LLMUnavailable for "could not ask", which is never the same
    as "asked and found nothing clean". run_copilot has always skipped a bad
    slide; what changed is that a timeout and a 429 now ARRIVE as that, instead
    of hanging the worker or being counted as a clean slide."""
    import io

    from pptx import Presentation
    from pptx.util import Emu

    import qc.copilot as CP
    from qc.llm import LLMUnavailable

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(3):
        slide.shapes.add_shape(1, Emu(914400 * (i + 1)), Emu(914400),
                               Emu(914400), Emu(914400))
    buf = io.BytesIO()
    prs.save(buf)

    def _down(**kwargs):
        raise LLMUnavailable("the model could not be reached")

    monkeypatch.setattr(CP, "ask_json", _down)
    records, reviewed = CP.run_copilot(buf.getvalue(), {0: b"png"},
                                       {"records": []})

    assert records == [] and reviewed == 0, (
        "a slide that could not be asked about is skipped, never counted as "
        "reviewed and never recorded as clean")
