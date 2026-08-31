"""The copilot's precision gate must not throw away the defects a designer
sees from across the room.

Two real misses on a client deck (31/08/2026), both in the gate rather than in
the model:

    a label 0.35in left of the block it heads    dropped: a 0.15in ceiling
                                                 discarded anything further out
    the same, as a PAIR                          dropped: the gate required
                                                 three shapes per action

The model reported both. The verification step silently deleted them, so the
worse the defect, the more certain the silence.
"""


from pptx import Presentation
from pptx.util import Emu

import qc.copilot as copilot
from qc.fixer import is_fixable

IN = 914400


def _pair(label_left=0.55, block_left=0.9):
    """A column label over the block it heads, the label sitting left of it."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    label = slide.shapes.add_textbox(Emu(int(label_left * IN)), Emu(int(1.6 * IN)),
                                     Emu(int(2.6 * IN)), Emu(int(0.4 * IN)))
    label.text_frame.text = "Method 1"
    block = slide.shapes.add_textbox(Emu(int(block_left * IN)), Emu(int(2.2 * IN)),
                                     Emu(int(5.0 * IN)), Emu(int(3.4 * IN)))
    block.text_frame.text = "1 2 3 4 5"
    return slide, label, block


def _align(label, block):
    return [{"action": "align_left",
             "shape_ids": [str(label.shape_id), str(block.shape_id)],
             "rationale": "The column label should sit square with its body."}]


def test_a_pair_is_a_real_alignment_case():
    """The gate required three shapes per action, so the commonest real
    misalignment there is - a label over its block - could not be expressed."""
    slide, label, block = _pair()

    recs = copilot.synthesize(slide, 0, _align(label, block), existing=[])

    assert len(recs) == 1, "a two-shape alignment was dropped by the gate"
    assert recs[0]["issue_type"] == "margin_alignment.edge_misaligned"
    assert str(recs[0]["shape_id"]) == str(label.shape_id), \
        "the LABEL moves, not the block it heads"


def test_the_pair_snaps_to_the_larger_shape_not_to_whichever_sorts_second():
    """Two values have no majority. Taking a median of two would move the block
    sideways half the time, which is worse than doing nothing."""
    slide, label, block = _pair(label_left=0.55, block_left=0.9)

    recs = copilot.synthesize(slide, 0, _align(label, block), existing=[])

    assert int(recs[0]["new_value"]) == int(0.9 * IN)
    assert int(recs[0]["old_value"]) == int(0.55 * IN)


def test_the_same_pair_the_other_way_round_still_moves_the_label():
    """Order of ids, and which side the label strays to, must not decide it."""
    slide, label, block = _pair(label_left=1.4, block_left=0.9)
    obs = [{"action": "align_left",
            "shape_ids": [str(block.shape_id), str(label.shape_id)],
            "rationale": "x"}]

    recs = copilot.synthesize(slide, 0, obs, existing=[])

    assert len(recs) == 1
    assert str(recs[0]["shape_id"]) == str(label.shape_id)
    assert int(recs[0]["new_value"]) == int(0.9 * IN)


def test_a_third_of_an_inch_off_is_reported_and_fixable():
    """The old 0.15in ceiling dropped this. It is the exact defect the design
    lead reported, and it is well inside a move the tool can compute."""
    slide, label, block = _pair(label_left=0.55, block_left=0.9)   # 0.35in off

    rec = copilot.synthesize(slide, 0, _align(label, block), existing=[])[0]

    assert rec["confidence"] == "medium"
    assert is_fixable(rec), "it is a computed snap; it should be one press"


def test_a_gap_too_big_to_snap_is_reported_but_not_offered_as_a_move():
    """Never silently dropped, and never auto-moved either. A shape two inches
    out is more likely a mis-grouping than a designer's slip, so the designer
    decides."""
    slide, label, block = _pair(label_left=3.5, block_left=0.9)    # 2.6in off

    recs = copilot.synthesize(slide, 0, _align(label, block), existing=[])

    assert len(recs) == 1, "it was dropped instead of reported"
    rec = recs[0]
    assert rec["confidence"] == "low"
    assert rec["new_value"] is None
    assert not is_fixable(rec), "a move this big must not be one-click"
    assert "too far to move for you" in rec["message"]


def test_a_shape_already_on_the_line_says_nothing():
    slide, label, block = _pair(label_left=0.9, block_left=0.9)

    assert copilot.synthesize(slide, 0, _align(label, block), existing=[]) == []


def test_distribute_and_match_still_need_three():
    """Widening the alignment gate must not widen the others: there is no
    rhythm in two gaps and no sibling set in two shapes."""
    slide, label, block = _pair()
    ids = [str(label.shape_id), str(block.shape_id)]

    for action in ("distribute_row", "distribute_col",
                   "match_widths", "match_heights"):
        obs = [{"action": action, "shape_ids": ids, "rationale": "x"}]
        assert copilot.synthesize(slide, 0, obs, existing=[]) == [], action


def test_the_odd_one_in_a_row_of_five_is_still_found():
    """The circles case: four on a line, one dropped below it. With three or
    more the target is the median, so the four win and the one moves."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    circles = []
    for i in range(5):
        top = 1.75 + (0.12 if i == 2 else 0)
        circles.append(slide.shapes.add_textbox(
            Emu(int((0.9 + i * 2.3) * IN)), Emu(int(top * IN)),
            Emu(int(2.0 * IN)), Emu(int(2.0 * IN))))

    obs = [{"action": "align_top",
            "shape_ids": [str(c.shape_id) for c in circles],
            "rationale": "One goal sits below the rest of the row."}]
    recs = copilot.synthesize(slide, 0, obs, existing=[])

    assert len(recs) == 1, "only the stray should move"
    assert str(recs[0]["shape_id"]) == str(circles[2].shape_id)
    assert int(recs[0]["new_value"]) == int(1.75 * IN)
    assert is_fixable(recs[0])
