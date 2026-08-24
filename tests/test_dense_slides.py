"""Crowded slides: what the pairwise caps hide, and saying when they bite.

Overlap is O(n^2) so a limit has to exist, but the limit was set before either
pass was profiled and it was cutting the ANSWER rather than the cost. Measured
on a dense slide (24/08/2026): at 200 shapes a cap of 120 reported 54 overlaps
and a cap of 250 reported 204, so three quarters were being hidden to save 0.3s
of work - on exactly the slides where overlaps live (design lead, screenshots of
two dense Arabic slides, 24/08/2026).

The other half of the fix is that a cap which does bite now says so. A slide of
300 elements silently checked to 120 of them looks identical to a slide with
nothing wrong.
"""

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

import qc.design as D
import qc.modules.margin_alignment as ma
from qc.profile import Profile
from tests.conftest import save_and_ctx

IN = 914400
BLANK = 6


def _crowded(prs, n, overlap=True):
    """n text boxes laid out to overlap their neighbours, the shape of the
    dense slides in the screenshots."""
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    step = 0.45 if overlap else 0.75
    for i in range(n):
        col, row = i % 6, (i // 6) % 11
        box = slide.shapes.add_textbox(
            Emu(int((0.4 + col * 2.0) * IN)), Emu(int((0.4 + row * step) * IN)),
            Emu(int(1.9 * IN)), Emu(int(0.5 * IN)))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = f"label {i}"
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xEE, 0xF3, 0xFC)
    return slide


def _bytes(prs):
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_a_dense_slide_is_no_longer_checked_to_a_fifth_of_itself(make_prs):
    """The regression this exists for: the count used to stop growing with the
    slide, because the cap - not the slide - decided how many pairs were
    compared."""
    counts = {}
    for n in (120, 200):
        prs = make_prs()
        _crowded(prs, n)
        found = D.scan(_bytes(prs), {})
        counts[n] = sum(1 for f in found if f.kind == "overlap")
    assert counts[200] > counts[120] * 1.5, (
        f"200 shapes reported {counts[200]} overlaps against {counts[120]} for "
        f"120: the cap is still deciding the answer")


def test_the_cap_says_so_when_it_bites(make_prs, monkeypatch):
    """A silent cut reads as "nothing else here", and the slides that trip it
    are the crowded ones."""
    monkeypatch.setattr(D, "MAX_PAIRWISE", 20)
    prs = make_prs()
    _crowded(prs, 60)
    found = D.scan(_bytes(prs), {})
    capped = [f for f in found if "too crowded" in f.headline]
    assert len(capped) == 1, "the slide was truncated in silence"
    assert capped[0].severity == "info", "it is a disclosure, not a defect"
    assert capped[0].options == [], "there is nothing here to fix"
    assert "40" in capped[0].detail, "it must say how many were left out"


def test_a_slide_inside_the_cap_says_nothing(make_prs):
    prs = make_prs()
    _crowded(prs, 12)
    found = D.scan(_bytes(prs), {})
    assert not [f for f in found if "too crowded" in f.headline]


def test_the_text_bearing_shapes_are_the_ones_kept(make_prs, monkeypatch):
    """This check only ever reports a pair with text on at least one side, so a
    pair dropped from the no-text end could not have produced a finding. Taking
    the first N in z-order threw that away at random."""
    monkeypatch.setattr(D, "MAX_PAIRWISE", 8)
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    # 20 plain rectangles first, so z-order order would keep only those...
    for i in range(20):
        slide.shapes.add_shape(1, Emu(int((0.4 + (i % 5) * 2.0) * IN)),
                               Emu(int((0.4 + (i // 5) * 0.5) * IN)),
                               Emu(int(1.9 * IN)), Emu(int(0.5 * IN)))
    # ...and the overlapping text last. Stacked tightly enough that the WORDS
    # overlap, not just the boxes: the check measures text extents, so a step
    # bigger than a line height is two boxes touching and nothing colliding.
    for i in range(6):
        box = slide.shapes.add_textbox(Emu(int(1.0 * IN)),
                                       Emu(int((3.0 + i * 0.06) * IN)),
                                       Emu(int(4.0 * IN)), Emu(int(0.4 * IN)))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = f"stacked line {i}"
        run.font.size = Pt(18)

    found = D.scan(_bytes(prs), {})
    assert [f for f in found if f.kind == "overlap"], (
        "the text was dropped in favour of rectangles that can never collide")


def test_the_audits_own_overlap_cap_also_reports_itself(make_prs, en_profile,
                                                        tmp_path, monkeypatch):
    monkeypatch.setattr(ma, "MAX_PAIRWISE_SHAPES", 15)
    prs = make_prs()
    _crowded(prs, 40)
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = [r.to_dict() for r in ma.detect(ctx)
            if r.issue_type == "margin_alignment.overlap_check_capped"]
    assert len(recs) == 1
    assert recs[0]["severity"] == "info"
    assert recs[0]["shape_id"] == "-", "it is about the slide, not a shape"


def test_the_caps_agree_with_each_other():
    """Two passes looking at the same slide with different ideas of how much of
    it is worth reading would report different amounts of the same problem."""
    assert D.MAX_PAIRWISE == ma.MAX_PAIRWISE_SHAPES
