"""Deterministic extraction: the ground truth the judgment passes argue from.

What is protected here is not the JSON's shape but the four claims the pipeline
rests on, because each one is a way this layer can look right and be wrong:

    a colour is reported BOTH as written and as resolved, so "on palette via
    accent1" and "someone typed a near-navy" stay distinguishable;

    an INHERITED colour or font comes back with a value, not null - the
    python-pptx trap, and the reason this module borrows the audit's resolvers;

    geometry is in SLIDE coordinates, group children composed, so a card's icon
    is not reported at the wrong side of the slide;

    and the ids are the same currency every other pass validates against, so
    "does this id exist" is one rule rather than three.
"""

import io
import json

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from qc import extract
from qc.applymaster import plan_assignments
from qc.components import inventory
from qc.design import placed_shapes

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"


# ------------------------------------------------------------------ helpers


def _bytes(prs) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _deck(prs) -> dict:
    """Extract from BYTES, the way the web layer calls it: an upload is never
    written to disk, so the bytes path is the one that has to work."""
    return extract.extract_deck(_bytes(prs))


def _shapes(deck, slide_index: int = 0) -> dict:
    return {s["id"]: s for s in deck["slides"][slide_index]["shapes"]}


def _set_solid(shape, color_xml: str):
    """Replace a shape's fill with one written exactly this way. Planted as XML
    because the distinction under test - a schemeClr reference against a typed
    srgbClr - is a distinction python-pptx's fill API does not expose."""
    spPr = shape._element.spPr
    for tag in ("a:solidFill", "a:noFill", "a:blipFill", "a:gradFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    fill = etree.fromstring(
        f'<a:solidFill xmlns:a="{A}">{color_xml}</a:solidFill>')
    # a:solidFill goes after a:xfrm and the geometry, before a:ln.
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(fill)
    else:
        spPr.append(fill)


def _palette_slide(make_prs):
    """One rectangle on accent1, one on accent1 at 60% luminosity, and one on a
    hand-typed hex that is not in the theme. The three cases a palette review
    has to tell apart."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    boxes = [slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1 + i * 3),
                                    Inches(1), Inches(2), Inches(1))
             for i in range(3)]
    _set_solid(boxes[0], '<a:schemeClr val="accent1"/>')
    _set_solid(boxes[1], '<a:schemeClr val="accent1"><a:lumMod val="60000"/>'
                         '<a:lumOff val="40000"/></a:schemeClr>')
    _set_solid(boxes[2], '<a:srgbClr val="1F3864"/>')
    return prs, boxes


# ------------------------------------------------------ as written vs resolved


def test_a_theme_reference_and_a_typed_hex_are_not_the_same_fact(make_prs):
    """The distinction the palette review is built on. All three shapes are
    painted; only one of them is a colour somebody chose by hand."""
    prs, _boxes = _palette_slide(make_prs)
    fills = [s["fill"] for s in _deck(prs)["slides"][0]["shapes"]]

    themed = [f for f in fills if f["as_written"] == "theme_color"]
    typed = [f for f in fills if f["as_written"] == "explicit_rgb"]
    assert len(themed) == 2 and len(typed) == 1

    assert typed[0]["hex"] == "1F3864"
    assert typed[0]["written_hex"] == "1F3864"
    for f in themed:
        assert f["theme_slot"] == "accent1"
        assert f["hex"], "a theme reference still has to resolve to a value"


def test_a_transform_on_a_theme_colour_is_recorded_and_applied(make_prs):
    """accent1 at 60% luminosity is a deliberate lighter brand blue, not a
    second blue. The transform is reported so nothing downstream reads it as
    off-palette, and the resolved hex is the transformed one, not the slot's."""
    prs, _boxes = _palette_slide(make_prs)
    deck = _deck(prs)
    lightened = [s["fill"] for s in deck["slides"][0]["shapes"]
                 if s["fill"].get("transforms")]
    assert len(lightened) == 1
    fill = lightened[0]
    assert fill["transforms"] == {"lumMod": 0.6, "lumOff": 0.4}
    assert fill["theme_slot"] == "accent1"
    assert fill["hex"] != deck["theme"]["colors"]["accent1"], \
        "the transform has to reach the resolved value"


def test_an_inherited_text_colour_still_has_a_value(make_prs):
    """python-pptx answers None for every inherited colour and font, which is
    most of them. A reader that passes those on as null reports a deck with no
    typography in it, and every check downstream then has nothing to compare."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Nothing on this run states a font or a colour"

    assert run.font.color.type is None, "the run itself states no colour"
    assert run.font.name is None and run.font.size is None

    got = [s for s in _deck(prs)["slides"][0]["shapes"]
           if s.get("paragraphs")][0]["paragraphs"][0]["runs"][0]
    assert got["color"]["hex"], "the resolved colour has to be a value"
    assert got["color"]["as_written"] == "inherited"
    assert got["color"]["from"], "and it has to say which level supplied it"
    assert got["font_name"] and got["size_pt"]
    assert got["font_from"] and got["size_from"]


def test_a_run_that_states_its_own_colour_is_marked_as_written(make_prs):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Typed navy"
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    got = [s for s in _deck(prs)["slides"][0]["shapes"]
           if s.get("paragraphs")][0]["paragraphs"][0]["runs"][0]
    assert got["color"]["as_written"] == "explicit_rgb"
    assert got["color"]["hex"] == "1F3864"
    assert got["size_pt"] == 24.0


# ---------------------------------------------------------------- backgrounds


def test_the_background_says_which_level_stated_it(make_prs):
    """A hex typed onto one slide is a local override; the same hex on the
    master is the brand. The value alone cannot tell a designer which of the two
    they are looking at, so the level is reported beside it."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cSld = slide._element.find(qn("p:cSld"))
    bg = etree.Element(f"{{{P}}}bg")
    pr = etree.SubElement(bg, f"{{{P}}}bgPr")
    fill = etree.SubElement(pr, f"{{{A}}}solidFill")
    etree.SubElement(fill, f"{{{A}}}srgbClr").set("val", "0B1F2A")
    etree.SubElement(pr, f"{{{A}}}effectLst")
    cSld.insert(0, bg)   # p:bg is the first child of p:cSld

    bg_out = _deck(prs)["slides"][0]["background"]
    assert bg_out["stated_by"] == "slide"
    assert bg_out["as_written"] == "explicit_rgb"
    assert bg_out["hex"] == "0B1F2A"


def test_a_deck_that_states_no_background_falls_to_the_theme(make_prs):
    """No stated background is not "no background": PowerPoint paints the
    theme's bg1 slot, and that is the ground the text on the slide sits on."""
    prs = make_prs()
    prs.slides.add_slide(prs.slide_layouts[6])
    bg = _deck(prs)["slides"][0]["background"]
    assert bg["hex"], "the ground a slide is painted on is never unknown here"
    assert bg["stated_by"] in ("layout", "master", "theme")


# ------------------------------------------------------------------- geometry


def test_a_grouped_shape_is_reported_where_it_lands_on_the_slide(make_prs):
    """A shape inside a group carries an offset in the GROUP's child space. Hand
    that over as a slide position and a card's icon is reported at the wrong
    side of the slide, which is the one mistake a layout question cannot
    survive."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    child = group.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1),
                                   Inches(1), Inches(1))
    # python-pptx fits a new group's child space exactly to its children,
    # which makes the transform an identity and would prove nothing. Move
    # the child origin and halve its extent, which is what a group resized
    # on a slide actually looks like: it now offsets its children and scales
    # them 2x, so a child's stored offset is not where it lands.
    xfrm = group._element.find(qn("p:grpSpPr")).find(qn("a:xfrm"))
    ch_off, ch_ext = xfrm.find(qn("a:chOff")), xfrm.find(qn("a:chExt"))
    ch_off.set("x", "0")
    ch_off.set("y", "0")
    ch_ext.set("cx", str(int(ch_ext.get("cx")) // 2))
    ch_ext.set("cy", str(int(ch_ext.get("cy")) // 2))

    blob = _bytes(prs)
    reported = _shapes(extract.extract_deck(blob))[str(child.shape_id)]
    assert reported["in_group"] is True

    boxes = {str(p.shape.shape_id): p.box
             for p in placed_shapes(Presentation(io.BytesIO(blob)).slides[0])}
    expected = boxes[str(child.shape_id)]
    assert reported["position"]["left"]["emu"] == expected[0]
    assert reported["position"]["top"]["emu"] == expected[1]
    assert reported["position"]["left"]["emu"] != child.left, (
        "the child's stored offset is not its place on the slide")
    assert reported["position"]["width"]["emu"] != child.width, (
        "a group's scale reaches its children's size, not just their "
        "offset")


def test_the_group_itself_is_listed_as_well_as_its_children(make_prs):
    """A group is what a designer selects and drags, so it is a thing on the
    slide. Its children are things too, because the icon welded to a card is
    what a layout question is about."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    a = group.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1),
                               Inches(1), Inches(1))
    b = group.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), Inches(1),
                               Inches(1), Inches(1))

    got = _shapes(_deck(prs))
    assert got[str(group.shape_id)]["content"] == "group"
    assert got[str(group.shape_id)]["in_group"] is False
    assert got[str(a.shape_id)]["in_group"] and got[str(b.shape_id)]["in_group"]


# ------------------------------------------------------------------------ ids


def test_the_ids_are_the_currency_the_vision_passes_validate_against(make_prs):
    """Referential validation is one rule across every pass: a model may only
    name ids from the inventory it was handed. Two id schemes would mean two
    rules, and the second one would be the one nobody wrote."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(3):
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1 + i), Inches(1),
                               Inches(1), Inches(1))

    blob = _bytes(prs)
    inv_ids = {item["id"] for item in inventory(
        Presentation(io.BytesIO(blob)).slides[0],
        prs.slide_width, prs.slide_height)}
    top_level = {s["id"] for s in extract.extract_deck(blob)["slides"][0]["shapes"]
                 if not s["in_group"]}
    assert inv_ids and inv_ids <= top_level


# -------------------------------------------------------------------- layouts


def test_the_layout_inventory_feeds_the_planner_unchanged(fixtures_dir):
    """The master half of the pipeline input is the reader the format pass
    already plans from (qc.stylespec.extract_layouts). If this view needed
    massaging before it could be planned from, there would be two answers to
    "what layouts does this master have"."""
    out = extract.extract_master_layouts(str(fixtures_dir / "mixed_layouts.pptx"))
    assert out["layouts"], "a master with layouts has to report them"
    for entry in out["layouts"]:
        assert {"name", "type", "placeholders", "master_index"} <= set(entry)
        assert "background" in entry

    deck = Presentation(str(fixtures_dir / "clean.pptx"))
    plans = plan_assignments(deck, out["layouts"])
    assert len(plans) == len(list(deck.slides))
    assert all(p.target_layout for p in plans), \
        "every slide gets a target, or the planner could not read this view"


def test_the_layout_view_carries_no_embedded_assets(fixtures_dir):
    """This view exists to be reasoned over next to the rendered pictures. A
    base64 background image in the middle of it helps nobody and costs the
    context the slides need."""
    blob = json.dumps(extract.extract_master_layouts(
        str(fixtures_dir / "mixed_layouts.pptx")))
    assert "base64" not in blob and '"bytes"' not in blob


def test_the_pair_carries_both_halves(fixtures_dir):
    """Layout matching and the palette review are comparisons. A pass handed
    the two halves separately can be handed halves from different files."""
    pair = extract.extract_pair(str(fixtures_dir / "clean.pptx"),
                                str(fixtures_dir / "mixed_layouts.pptx"))
    assert pair["deck"]["slides"] and pair["master"]["layouts"]


# -------------------------------------------------------------------- palette


def test_the_palette_roll_up_counts_what_was_typed_by_hand(make_prs):
    """The designer's opening question. Two of these three shapes are on the
    theme; one is a hex somebody typed, and that is the one worth a
    conversation."""
    pal = extract.palette_inventory(_deck(_palette_slide(make_prs)[0]))
    assert pal["explicit_count"] == 1
    typed = [c for c in pal["colours"] if c["written"].get("explicit_rgb")]
    assert typed[0]["hex"] == "1F3864"
    assert "fill" in typed[0]["roles"]
    assert pal["theme_slots"]["accent1"], "the theme it is judged against"


def test_the_roll_up_separates_backgrounds_from_everything_else(make_prs):
    """"Which colours are explicit for backgrounds" is a different question from
    which colours appear at all, so the role a colour plays travels with it."""
    pal = extract.palette_inventory(_deck(_palette_slide(make_prs)[0]))
    roles = {r for c in pal["colours"] for r in c["roles"]}
    assert any(r.startswith("background") for r in roles)
    assert "fill" in roles


# ----------------------------------------------------------------- robustness


def test_a_chart_is_not_reported_as_a_text_box(make_prs):
    """A chart and a table are graphic frames that also answer to
    has_text_frame. Read the wrong property first and every chart in the deck
    is reported as text."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ["a", "b"]
    data.add_series("s", (1.0, 2.0))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1),
                           Inches(4), Inches(3), data)
    slide.shapes.add_table(2, 2, Inches(6), Inches(1), Inches(3), Inches(2))

    kinds = {s["content"] for s in _deck(prs)["slides"][0]["shapes"]}
    assert "chart" in kinds and "table" in kinds
    assert "text" not in kinds


def test_a_real_deck_extracts_without_losing_a_slide(fixtures_dir):
    """The corpus deck with the awkward content. Every slide has to appear, and
    nothing may come back as an exception swallowed into an empty list."""
    path = str(fixtures_dir / "heavy.pptx")
    deck = extract.extract_deck(path)
    assert len(deck["slides"]) == len(list(Presentation(path).slides))
    assert deck["schema_version"] == extract.SCHEMA_VERSION
    for slide in deck["slides"]:
        for shape in slide["shapes"]:
            assert "text_error" not in shape, shape.get("text_error")
            assert shape["fill"]["paint"] in ("solid", "opaque", "clear")
