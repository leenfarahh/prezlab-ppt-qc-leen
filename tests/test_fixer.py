"""Tests for qc.fixer.apply_fixes: the v1.5 deterministic fix tier.

Each fix test runs the real pipeline end to end: build a planted-violation
deck, audit it from disk via run_audit, select record ids, apply_fixes on the
original bytes, then verify both the mutated deck and a re-audit of the
cleaned bytes.
"""

import io
import json
from copy import deepcopy

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.util import Emu, Inches, Pt

from qc.engine import run_audit
from qc.fixer import apply_fixes
from qc.profile import PROFILES_DIR, Profile
from qc.records import make_record

FOOTER_TEXT = "Prezlab | Confidential"
DOM_W, DOM_H = 2000000, 1000000
DEV_W = 2100000  # 100000 EMU off, well beyond the 9525 tolerance
TOP = 1000000


def _save(prs, tmp_path, name):
    path = tmp_path / name
    prs.save(path)
    return path


def _apply(path, audit_result, record_ids):
    records = [r.to_dict() for r in audit_result.records]
    return apply_fixes(path.read_bytes(), records, record_ids)


def _reaudit(cleaned_bytes, tmp_path, profile, modules, name="cleaned.pptx"):
    path = tmp_path / name
    path.write_bytes(cleaned_bytes)
    return run_audit(path, profile, modules=modules)


def _add_rect(slide, left, width=DOM_W, height=DOM_H):
    return slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(left), Emu(TOP), Emu(width), Emu(height))


def _hf_profile(**template_overrides) -> Profile:
    data = json.loads(
        (PROFILES_DIR / "prezlab_en.json").read_text(encoding="utf-8"))
    data["config"]["header_footer"]["template"].update(template_overrides)
    return Profile(data)


def _add_slide_with_footer(prs, text):
    """python-pptx does not clone footer placeholders on add_slide, so copy
    the layout's footer element onto the slide before setting its text."""
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    for ph in layout.placeholders:
        if ph.placeholder_format.type == PP_PLACEHOLDER.FOOTER:
            slide.shapes._spTree.append(deepcopy(ph._element))
    footer = next(ph for ph in slide.placeholders
                  if ph.placeholder_format.type == PP_PLACEHOLDER.FOOTER)
    footer.text_frame.text = text
    # Caption-role latin family keeps the audit to the text mismatch only.
    footer.text_frame.paragraphs[0].runs[0].font.name = "Trebuchet MS"
    return footer


def _find_footer(prs, slide_index=0):
    return next(ph for ph in prs.slides[slide_index].placeholders
                if ph.placeholder_format.type == PP_PLACEHOLDER.FOOTER)


def test_font_family_fix_applied_and_finding_cleared(make_prs, en_profile, tmp_path):
    prs = make_prs()
    for text in ("First off-profile run", "Second off-profile run"):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        body = slide.placeholders[1]
        body.text_frame.text = text
        run = body.text_frame.paragraphs[0].runs[0]
        run.font.name = "Comic Sans MS"
        run.font.size = Pt(18)  # role-compliant size so only family is flagged
    path = _save(prs, tmp_path, "font.pptx")

    result = run_audit(path, en_profile, modules=["font"])
    fam = [r for r in result.records
           if r.issue_type == "font.family_out_of_set"]
    assert len(fam) == 2
    target = next(r for r in fam if r.slide_index == 0)
    other = next(r for r in fam if r.slide_index == 1)

    fix = _apply(path, result, {target.record_id})
    assert fix.applied == 1
    assert len(fix.outcomes) == 1
    assert fix.outcomes[0].record_id == target.record_id
    assert fix.outcomes[0].outcome == "changed"

    # Only the applied record flips to "changed"; the rest stay "flagged".
    by_id = {r["record_id"]: r for r in fix.records}
    assert by_id[target.record_id]["action"] == "changed"
    assert by_id[other.record_id]["action"] == "flagged"

    cleaned = Presentation(io.BytesIO(fix.cleaned_bytes))
    fixed_run = cleaned.slides[0].placeholders[1].text_frame.paragraphs[0].runs[0]
    assert fixed_run.font.name == "Trebuchet MS"

    re_result = _reaudit(fix.cleaned_bytes, tmp_path, en_profile, ["font"])
    remaining = [r for r in re_result.records
                 if r.issue_type == "font.family_out_of_set"]
    assert [r.slide_index for r in remaining] == [1]  # slide 0 finding is gone


def test_placeholder_geometry_fix_restores_inheritance(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    assert title.left is not None  # inherited geometry must resolve
    title.left = title.left + Inches(2)
    path = _save(prs, tmp_path, "geometry.pptx")

    result = run_audit(path, en_profile, modules=["master_slide"])
    geo = [r for r in result.records
           if r.issue_type == "master_slide.placeholder_geometry_off"]
    assert len(geo) == 1

    fix = _apply(path, result, {geo[0].record_id})
    assert fix.applied == 1
    assert fix.outcomes[0].outcome == "changed"

    re_result = _reaudit(fix.cleaned_bytes, tmp_path, en_profile, ["master_slide"])
    assert [r for r in re_result.records
            if r.issue_type == "master_slide.placeholder_geometry_off"] == []


def test_shape_size_fix_snaps_to_dominant(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(4):
        _add_rect(slide, left=500000 + i * 2200000)
    deviant = _add_rect(slide, left=500000, width=DEV_W)
    deviant_id = str(deviant.shape_id)
    path = _save(prs, tmp_path, "sizes.pptx")

    result = run_audit(path, en_profile, modules=["shape_size"])
    mismatch = [r for r in result.records
                if r.issue_type == "shape_size.size_mismatch"]
    assert len(mismatch) == 1
    assert mismatch[0].shape_id == deviant_id

    fix = _apply(path, result, {mismatch[0].record_id})
    assert fix.applied == 1
    assert fix.outcomes[0].outcome == "changed"

    cleaned = Presentation(io.BytesIO(fix.cleaned_bytes))
    fixed = next(s for s in cleaned.slides[0].shapes
                 if str(s.shape_id) == deviant_id)
    assert (fixed.width, fixed.height) == (DOM_W, DOM_H)

    re_result = _reaudit(fix.cleaned_bytes, tmp_path, en_profile, ["shape_size"])
    assert [r for r in re_result.records
            if r.issue_type == "shape_size.size_mismatch"] == []


def test_footer_text_fix_sets_profile_text(make_prs, tmp_path):
    profile = _hf_profile(footer_text=FOOTER_TEXT)
    prs = make_prs()
    _add_slide_with_footer(prs, "Wrong footer")
    path = _save(prs, tmp_path, "footer.pptx")

    result = run_audit(path, profile, modules=["header_footer"])
    mismatch = [r for r in result.records
                if r.issue_type == "header_footer.text_mismatch"]
    assert len(mismatch) == 1
    assert mismatch[0].new_value == FOOTER_TEXT

    fix = _apply(path, result, {mismatch[0].record_id})
    assert fix.applied == 1
    assert fix.outcomes[0].outcome == "changed"

    cleaned = Presentation(io.BytesIO(fix.cleaned_bytes))
    assert _find_footer(cleaned).text_frame.text == FOOTER_TEXT

    re_result = _reaudit(fix.cleaned_bytes, tmp_path, profile, ["header_footer"])
    assert [r for r in re_result.records
            if r.issue_type == "header_footer.text_mismatch"] == []


def test_arabic_font_substitution_applies_on_explicit_tick(make_prs,
                                                           tmp_path):
    """Since 12/08/2026 the Arabic font guard is tick-to-approve: an
    explicitly selected complex-script substitution applies (writing the
    cs typeface, never the text); text-editing fixes stay guarded (see
    tests/test_arabic_geometry.py)."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    body = slide.placeholders[1]
    body.text_frame.text = "عرض تقديمي"
    path = _save(prs, tmp_path, "arabic.pptx")
    original_bytes = path.read_bytes()

    rec = make_record(
        slide_index=0, shape_id=str(body.shape_id), module="font",
        issue_type="font.family_out_of_set", severity="error",
        action="flagged", confidence="deterministic", arabic_flag=True,
        property="rPr.cs.typeface", old_value="Arial", new_value="Dubai",
        locator="p0/r0", profile_rule_id="font.roles.body.complex_script",
        message="Arabic run, hand-built for the guard test.",
    ).to_dict()

    fix = apply_fixes(original_bytes, [rec], {rec["record_id"]})
    assert fix.applied == 1
    assert fix.outcomes[0].outcome == "changed"

    from pptx.oxml.ns import qn

    cleaned = Presentation(io.BytesIO(fix.cleaned_bytes))
    run = cleaned.slides[0].placeholders[1].text_frame.paragraphs[0].runs[0]
    cs = run._r.find(qn("a:rPr")).find(qn("a:cs"))
    assert cs.get("typeface") == "Dubai"
    assert run.text == "عرض تقديمي"  # the text itself is never touched


def test_unknown_record_id_is_skipped(make_prs, tmp_path):
    prs = make_prs()
    prs.slides.add_slide(prs.slide_layouts[6])
    path = _save(prs, tmp_path, "plain.pptx")

    fix = apply_fixes(path.read_bytes(), [], {"deadbeefdeadbeef"})
    assert fix.applied == 0
    assert len(fix.outcomes) == 1
    assert fix.outcomes[0].record_id == "deadbeefdeadbeef"
    assert fix.outcomes[0].outcome == "skipped"
    assert fix.outcomes[0].reason == "unknown record id"
    assert len(Presentation(io.BytesIO(fix.cleaned_bytes)).slides) == 1


def test_empty_selection_applies_nothing(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(4):
        _add_rect(slide, left=500000 + i * 2200000)
    deviant = _add_rect(slide, left=500000, width=DEV_W)
    deviant_id = str(deviant.shape_id)
    path = _save(prs, tmp_path, "noselect.pptx")

    result = run_audit(path, en_profile, modules=["shape_size"])
    assert result.records  # findings exist, but none are selected

    fix = _apply(path, result, set())
    assert fix.applied == 0
    assert fix.outcomes == []
    assert all(r["action"] == "flagged" for r in fix.records)

    cleaned = Presentation(io.BytesIO(fix.cleaned_bytes))
    untouched = next(s for s in cleaned.slides[0].shapes
                     if str(s.shape_id) == deviant_id)
    assert untouched.width == DEV_W


def _xfrm_invariant_violations(deck_bytes: bytes) -> list[str]:
    """Shapes whose xfrm has an off without an ext (or vice versa):
    PowerPoint renders such placeholders as degenerate slivers."""
    import zipfile

    from lxml import etree

    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    bad = []
    with zipfile.ZipFile(io.BytesIO(deck_bytes)) as z:
        for name in z.namelist():
            if not (name.startswith("ppt/slides/slide") and name.endswith(".xml")):
                continue
            root = etree.fromstring(z.read(name))
            for xfrm in root.iter(f"{A}xfrm"):
                off = xfrm.find(f"{A}off") is not None
                ext = xfrm.find(f"{A}ext") is not None
                if off != ext:
                    bad.append(f"{name}: off={off} ext={ext}")
    return bad


def test_geometry_snap_plus_edge_snap_leave_complete_xfrm(make_prs, en_profile, tmp_path):
    """Regression (real-deck corruption, 14/07/2026): geometry snap removes
    the xfrm, a later positional fix then wrote an off-without-ext transform
    that PowerPoint renders one letter per line. Whatever the record order,
    the cleaned deck must never contain a half transform."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    title = slide.shapes.title
    title.text = "A title moved off its layout position"
    title.left, title.top = Emu(3982563), Emu(384048)  # explicit override

    path = _save(prs, tmp_path, "conflict.pptx")
    records = [r.to_dict() for r in run_audit(
        path, "prezlab_en", modules=["master_slide"]).records]
    geo = [r for r in records
           if r["issue_type"] == "master_slide.placeholder_geometry_off"]
    assert geo, "planted geometry override must be detected"

    # a positional fix on the SAME shape, forced to sort after the snap
    edge = make_record(
        slide_index=0, shape_id=geo[0]["shape_id"], module="margin_alignment",
        issue_type="margin_alignment.edge_misaligned", severity="warning",
        confidence="medium", property="spPr.xfrm.off.x",
        old_value=3982563, new_value=450987,
        message="left edge off cluster median").to_dict()
    edge["record_id"] = "zzz-" + edge["record_id"]  # sorts after any uuid
    records.append(edge)

    fix = apply_fixes(path.read_bytes(), records,
                      {geo[0]["record_id"], edge["record_id"]})
    assert all(o.outcome == "changed" for o in fix.outcomes), \
        [(o.record_id, o.reason) for o in fix.outcomes]
    assert _xfrm_invariant_violations(fix.cleaned_bytes) == []

    # and the shape is where the edge fix asked, with a real extent
    cleaned = Presentation(io.BytesIO(fix.cleaned_bytes))
    t = cleaned.slides[0].shapes.title
    assert t.left == 450987
    assert t.width is not None and t.width > 0
