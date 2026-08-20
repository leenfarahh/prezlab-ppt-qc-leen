"""Tests for the font audit module (qc.modules.font)."""

import copy

from pptx.util import Pt

from qc.modules import font as font_module
from qc.profile import Profile
from spike.arabic import set_cs_typeface
from tests.conftest import save_and_ctx


def _add_body_slide(prs, text):
    """Slide from the Title and Content layout with `text` in the body
    placeholder; returns (slide, first run of the body)."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    body = slide.placeholders[1]
    body.text_frame.text = text
    return slide, body.text_frame.paragraphs[0].runs[0]


def _style(run, name="Trebuchet MS", size_pt=18, bold=None):
    run.font.name = name
    run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold


def test_wrong_latin_family_flagged_as_error(make_prs, en_profile, tmp_path):
    prs = make_prs()
    _, run = _add_body_slide(prs, "Quarterly results overview")
    _style(run, name="Comic Sans MS", size_pt=18)

    records = font_module.detect(save_and_ctx(prs, tmp_path, en_profile))
    fam = [r for r in records if r.issue_type == "font.family_out_of_set"]
    assert len(fam) == 1
    r = fam[0]
    assert r.severity == "error"
    assert r.action == "flagged"
    assert r.property == "rPr.latin.typeface"
    assert r.old_value == "Comic Sans MS"
    assert r.new_value == "Trebuchet MS"
    assert r.confidence == "deterministic"  # family came from run.rPr
    assert r.profile_rule_id == "font.roles.body.latin"
    assert r.arabic_flag is False
    assert r.slide_index == 0


def test_arabic_run_missing_cs_typeface(make_prs, en_profile, tmp_path):
    prs = make_prs()
    _, run = _add_body_slide(prs, "التقرير السنوي")
    _style(run, name="Trebuchet MS", size_pt=18)  # latin set, no a:cs

    records = font_module.detect(save_and_ctx(prs, tmp_path, en_profile))
    missing = [r for r in records if r.issue_type == "font.cs_typeface_missing"]
    assert len(missing) == 1
    r = missing[0]
    assert r.severity == "warning"
    assert r.arabic_flag is True
    assert r.action == "flagged"
    # since 12/08/2026 the record proposes the first allowed cs font and
    # the tick is the approval (never pre-selected)
    assert r.new_value is not None
    assert "ticking it is your approval" in r.message
    # The latin family on an Arabic run is never audited.
    assert not [x for x in records if x.issue_type == "font.family_out_of_set"]


def test_arabic_run_disallowed_cs_typeface(make_prs, en_profile, tmp_path):
    prs = make_prs()
    _, run = _add_body_slide(prs, "التقرير السنوي")
    _style(run, name="Trebuchet MS", size_pt=18)
    set_cs_typeface(run, "Arial")  # profile allows only Dubai

    records = font_module.detect(save_and_ctx(prs, tmp_path, en_profile))
    fam = [r for r in records if r.issue_type == "font.family_out_of_set"]
    assert len(fam) == 1
    r = fam[0]
    assert r.severity == "error"
    assert r.arabic_flag is True
    assert r.action == "flagged"
    assert r.property == "rPr.cs.typeface"
    assert r.old_value == "Arial"
    # since 12/08/2026: the first allowed cs font is proposed; the fix is
    # never pre-selected, ticking it is the designer's approval
    assert r.new_value == "Dubai"
    assert "ticking it is your approval" in r.message
    assert r.profile_rule_id == "font.roles.body.complex_script"


def test_size_off_role_beyond_tolerance(make_prs, en_profile, tmp_path):
    prs = make_prs()
    _, run = _add_body_slide(prs, "Body text at the wrong size")
    _style(run, name="Trebuchet MS", size_pt=30)  # target 18, tolerance 0.5

    records = font_module.detect(save_and_ctx(prs, tmp_path, en_profile))
    size = [r for r in records if r.issue_type == "font.size_off_role"]
    assert len(size) == 1
    r = size[0]
    assert r.severity == "warning"
    assert r.confidence == "high"  # placeholder shape, role is known
    assert r.property == "rPr.sz"
    assert float(r.old_value) == 30.0
    assert float(r.new_value) == 18.0
    assert r.profile_rule_id == "font.roles.body.size_pt"
    assert r.arabic_flag is False
    assert len(records) == 1  # family and weight are compliant


def test_mixed_bold_paragraph_one_record(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide, run1 = _add_body_slide(prs, "Bold lead ")
    _style(run1, size_pt=18, bold=True)
    para = slide.placeholders[1].text_frame.paragraphs[0]
    run2 = para.add_run()
    run2.text = "then regular tail"
    _style(run2, size_pt=18, bold=False)

    records = font_module.detect(save_and_ctx(prs, tmp_path, en_profile))
    mixed = [r for r in records if r.issue_type == "font.mixed_weight"]
    assert len(mixed) == 1
    r = mixed[0]
    assert r.severity == "info"
    assert r.confidence == "high"
    assert r.property == "rPr.b"
    assert r.arabic_flag is False
    assert len(records) == 1  # exactly one record for the whole paragraph


def test_theme_ref_disallowed_only_when_profile_says_so(make_prs, en_profile, tmp_path):
    prs = make_prs()
    _, run = _add_body_slide(prs, "Theme referenced text")
    run.font.name = "+mn-lt"  # explicit theme font reference on the run
    run.font.size = Pt(18)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    # Refs allowed (default profile): no theme_ref record, but the theme
    # resolved family (Calibri) is still checked against the allowed set.
    records = font_module.detect(ctx)
    assert not [r for r in records if r.issue_type == "font.theme_ref_disallowed"]
    fam = [r for r in records if r.issue_type == "font.family_out_of_set"]
    assert len(fam) == 1
    assert fam[0].old_value == "Calibri"

    # Same deck, refs disallowed: theme_ref_disallowed fires.
    data = copy.deepcopy(en_profile.raw)
    data["config"]["font"]["theme_font_refs_allowed"] = False
    ctx.profile = Profile(data)
    records = font_module.detect(ctx)
    theme = [r for r in records if r.issue_type == "font.theme_ref_disallowed"]
    assert len(theme) == 1
    assert theme[0].severity == "error"
    assert theme[0].profile_rule_id == "font.theme_font_refs_allowed"


def test_clean_control_deck_yields_zero_records(make_prs, en_profile, tmp_path):
    # Explicit run values matching prezlab_en; inherited template defaults
    # (Calibri 18/28/...) would otherwise be off-profile.
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text_frame.text = "Annual review"
    _style(title.text_frame.paragraphs[0].runs[0], name="Georgia", size_pt=40)
    body = slide.placeholders[1]
    body.text_frame.text = "Compliant body copy"
    _style(body.text_frame.paragraphs[0].runs[0], name="Trebuchet MS", size_pt=18)

    records = font_module.detect(save_and_ctx(prs, tmp_path, en_profile))
    assert records == []
