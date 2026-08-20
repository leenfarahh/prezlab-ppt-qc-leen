"""Tests for qc.modules.header_footer (planted-violation decks)."""

import json
from copy import deepcopy

from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Emu, Inches

from qc.modules import header_footer
from qc.profile import PROFILES_DIR, Profile
from tests.conftest import save_and_ctx

FOOTER_TEXT = "Prezlab | Confidential"


def hf_profile(**template_overrides) -> Profile:
    """prezlab_en with header_footer.template values switched on."""
    data = json.loads(
        (PROFILES_DIR / "prezlab_en.json").read_text(encoding="utf-8"))
    data["config"]["header_footer"]["template"].update(template_overrides)
    return Profile(data)


def add_slide_with_footer(prs, text, latin="Trebuchet MS"):
    """Slide from the 'Title and Content' layout plus a real footer ph copied
    from the layout (python-pptx does not clone footer phs on add_slide)."""
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    for ph in layout.placeholders:
        if ph.placeholder_format.type == PP_PLACEHOLDER.FOOTER:
            slide.shapes._spTree.append(deepcopy(ph._element))
    footer = next(ph for ph in slide.placeholders
                  if ph.placeholder_format.type == PP_PLACEHOLDER.FOOTER)
    footer.text_frame.text = text
    if latin is not None:
        footer.text_frame.paragraphs[0].runs[0].font.name = latin
    return slide, footer


def test_missing_footer_placeholder(make_prs, tmp_path):
    prs = make_prs()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank layout, no footer ph
    ctx = save_and_ctx(prs, tmp_path, hf_profile(footer_text=FOOTER_TEXT))
    records = header_footer.detect(ctx)
    assert len(records) == 1
    rec = records[0]
    assert rec.issue_type == "header_footer.missing"
    assert rec.severity == "error"
    assert rec.confidence == "deterministic"
    assert rec.property == "ph.footer"
    assert rec.slide_index == 0
    assert rec.action == "flagged"
    assert rec.profile_rule_id == "header_footer.template.footer_text"


def test_footer_text_mismatch(make_prs, tmp_path):
    prs = make_prs()
    add_slide_with_footer(prs, "Wrong footer")
    ctx = save_and_ctx(prs, tmp_path, hf_profile(footer_text=FOOTER_TEXT))
    records = header_footer.detect(ctx)
    assert len(records) == 1
    rec = records[0]
    assert rec.issue_type == "header_footer.text_mismatch"
    assert rec.severity == "error"
    assert rec.confidence == "deterministic"
    assert rec.old_value == "Wrong footer"
    assert rec.new_value == FOOTER_TEXT
    assert rec.action == "flagged"


def test_clean_deck_zero_records(make_prs, tmp_path):
    prs = make_prs()
    add_slide_with_footer(prs, FOOTER_TEXT)  # right text, caption latin font
    ctx = save_and_ctx(prs, tmp_path, hf_profile(footer_text=FOOTER_TEXT))
    assert header_footer.detect(ctx) == []


def test_fake_footer_textbox_is_medium_confidence_missing(make_prs, tmp_path):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Bottom 20% of a 7.5in slide starts at 6.0in; place the box below that.
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0),
                                   Inches(5.0), Inches(0.35))
    box.text_frame.text = FOOTER_TEXT
    ctx = save_and_ctx(prs, tmp_path, hf_profile(footer_text=FOOTER_TEXT))
    records = header_footer.detect(ctx)
    assert len(records) == 1
    rec = records[0]
    assert rec.issue_type == "header_footer.missing"
    assert rec.severity == "error"
    assert rec.confidence == "medium"
    assert rec.property == "ph.footer"
    assert rec.shape_id == str(box.shape_id)
    assert "text box at the footer position" in rec.message


def test_slide_number_required_and_absent(make_prs, tmp_path):
    prs = make_prs()
    prs.slides.add_slide(prs.slide_layouts[6])
    ctx = save_and_ctx(prs, tmp_path, hf_profile(slide_number=True))
    records = header_footer.detect(ctx)
    assert len(records) == 1
    rec = records[0]
    assert rec.issue_type == "header_footer.missing"
    assert rec.severity == "error"
    assert rec.property == "ph.sldNum"
    assert rec.profile_rule_id == "header_footer.template.slide_number"


def test_position_mismatch(make_prs, tmp_path):
    prs = make_prs()
    _, footer = add_slide_with_footer(prs, FOOTER_TEXT)
    footer.left = Emu(914400)
    footer.top = Emu(6400800)
    profile = hf_profile(
        footer_text=FOOTER_TEXT,
        position_emu={"left": 914400 + 50000, "top": 6400800})
    ctx = save_and_ctx(prs, tmp_path, profile)
    records = header_footer.detect(ctx)
    assert len(records) == 1
    rec = records[0]
    assert rec.issue_type == "header_footer.position_mismatch"
    assert rec.severity == "warning"
    assert rec.confidence == "high"
    assert rec.arabic_flag is False
    assert "Arabic" not in rec.message


def test_arabic_footer_position_mismatch_flagged(make_prs, tmp_path):
    arabic_footer = "عرض تقديمي"
    prs = make_prs()
    _, footer = add_slide_with_footer(prs, arabic_footer, latin=None)
    footer.left = Emu(914400)
    footer.top = Emu(6400800)
    profile = hf_profile(
        footer_text=arabic_footer,  # expected text matches, only position off
        position_emu={"left": 914400 + 50000, "top": 6400800})
    ctx = save_and_ctx(prs, tmp_path, profile)
    records = header_footer.detect(ctx)
    assert len(records) == 1
    rec = records[0]
    assert rec.issue_type == "header_footer.position_mismatch"
    assert rec.arabic_flag is True
    assert "Arabic content, manual review" in rec.message
    # Arabic runs must not be latin-family checked: no font_mismatch emitted.
    assert all(r.issue_type != "header_footer.font_mismatch" for r in records)


def test_footer_font_mismatch(make_prs, tmp_path):
    prs = make_prs()
    add_slide_with_footer(prs, FOOTER_TEXT, latin="Comic Sans MS")
    ctx = save_and_ctx(prs, tmp_path, hf_profile(footer_text=FOOTER_TEXT))
    records = header_footer.detect(ctx)
    assert len(records) == 1
    rec = records[0]
    assert rec.issue_type == "header_footer.font_mismatch"
    assert rec.severity == "warning"
    assert rec.confidence == "medium"
    assert rec.old_value == "Comic Sans MS"
    assert rec.new_value == "Trebuchet MS"
    assert rec.profile_rule_id == "font.roles.caption.latin"


def test_profile_without_template_not_enforced(make_prs, tmp_path, en_profile):
    prs = make_prs()
    prs.slides.add_slide(prs.slide_layouts[6])  # would be missing if enforced
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert header_footer.detect(ctx) == []
