"""A heading past the margin: flagged, never acted on; the BOX is measured.

Two rules from the design lead (19/08/2026):

1. A title or standfirst that runs past its margin gets no action anywhere -
   not in the formatting pass, not in the audit, not in the fix engine. It is
   flagged so the designer can ask the client which they want, because whether
   a heading may break the frame is house style, not a defect.
2. What a margin is measured against is the shape's BOX, on all four sides.
   Text insets and glyph extents are rendering, not geometry the deck states.
"""

import io

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

import qc.modules.margin_alignment as ma
from qc.fixer import FIXABLE_ISSUES, is_fixable
from qc.migrate import migrate_deck
from tests.conftest import save_and_ctx
from tests.test_migrate import _plant_guides

HEADING = "margin_alignment.heading_past_margin"
SAFE_ZONE = "margin_alignment.outside_safe_zone"
OVERFLOW = "margin_alignment.content_overflow"

IN = 914400
BLANK_LAYOUT = 6
TITLE_LAYOUT = 0

# prezlab_en, against a 12192000 x 6858000 canvas.
M_LEFT = M_RIGHT = 457200
M_TOP = 274638
M_BOTTOM = 365125
SLIDE_W, SLIDE_H = 12192000, 6858000


def _of_type(ctx, issue_type):
    return [r for r in ma.detect(ctx) if r.issue_type == issue_type]


def _title_slide(prs, *, left, top, width, height, text="The Heading"):
    """A slide whose TITLE placeholder is placed exactly where the test wants
    it. Placeholder geometry is overridden rather than inherited: the stock
    template centres its title mid-canvas, which breaches nothing."""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
    title = None
    for ph in slide.placeholders:
        ph.text_frame.clear()
        if str(ph.placeholder_format.type).startswith(("TITLE", "CENTER_TITLE")):
            ph.left, ph.top, ph.width, ph.height = (
                Emu(left), Emu(top), Emu(width), Emu(height))
            ph.text_frame.text = text
            title = ph
    return slide, title


# ------------------------------------------------------- the audit's answer


@pytest.mark.parametrize("side,geometry", [
    ("left", dict(left=200000, top=1000000, width=6000000, height=800000)),
    ("top", dict(left=1000000, top=100000, width=6000000, height=800000)),
    ("right", dict(left=6000000, top=1000000, width=6000000, height=800000)),
    ("bottom", dict(left=1000000, top=6200000, width=6000000, height=500000)),
])
def test_every_side_flags_the_heading_and_fixes_nothing(
        side, geometry, make_prs, en_profile, tmp_path):
    """All four margins, measured the same way and answered the same way."""
    prs = make_prs()
    _slide, title = _title_slide(prs, **geometry)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = _of_type(ctx, HEADING)
    assert len(records) == 1, f"{side} margin"
    rec = records[0]
    assert rec.shape_id == str(title.shape_id)
    assert side in rec.message
    assert rec.action == "flagged"
    assert rec.new_value is None, "a heading gets no computed target"
    assert "ask them" in rec.message
    assert not is_fixable(rec.to_dict())
    # and the ordinary breach record does not ALSO fire: one shape, one answer
    assert not [r for r in _of_type(ctx, SAFE_ZONE)
                if r.shape_id == str(title.shape_id)]


def test_the_heading_issue_type_is_not_in_the_fixable_set():
    """The guarantee behind 'take no action', stated where it cannot drift."""
    assert HEADING not in FIXABLE_ISSUES


def test_a_heading_inside_its_margins_never_flags_however_long_the_text(
        make_prs, en_profile, tmp_path):
    """The BOX is what is measured. A line long enough to render past the frame
    says nothing about where the deck placed the box, and estimating glyph
    widths would swap a fact for a guess about the reader's fonts."""
    prs = make_prs()
    _title_slide(
        prs, left=M_LEFT, top=M_TOP, width=6000000, height=500000,
        text="A second LLM integration will summarize each team member's day, "
             "every morning, well before review, at a length no box this size "
             "could ever hold on one line")
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    assert _of_type(ctx, HEADING) == []
    assert _of_type(ctx, SAFE_ZONE) == []


def test_an_ordinary_text_box_still_gets_the_plain_breach_record(
        make_prs, en_profile, tmp_path):
    """Only headings are excused. Body content past the safe zone is the
    existing finding, unchanged."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    box = slide.shapes.add_textbox(Emu(0), Emu(2000000), Emu(3000000),
                                   Emu(600000))
    box.text_frame.text = "body copy that runs past the left margin"
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    assert _of_type(ctx, HEADING) == []
    breaches = _of_type(ctx, SAFE_ZONE)
    assert [r.shape_id for r in breaches] == [str(box.shape_id)]


# --------------------------------------- and no cross-slide pin either way


PINNED = "margin_alignment.recurring_off_position"


def _repeating_title_deck(prs, *, stray_left):
    """A title bar in the same spot on four slides, with the last one moved.
    That is what the cross-slide anchor check is for - and the one fix that
    could still reposition a heading behind a designer's tick."""
    strays = []
    for i in range(4):
        slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
        for ph in slide.placeholders:
            ph.text_frame.clear()
            if str(ph.placeholder_format.type).startswith(
                    ("TITLE", "CENTER_TITLE")):
                left = stray_left if i == 3 else 1000000
                ph.left, ph.top = Emu(left), Emu(1000000)
                ph.width, ph.height = Emu(6000000), Emu(800000)
                ph.text_frame.text = "The deck's running title"
                if i == 3:
                    strays.append(ph)
    return strays[0]


def test_a_heading_past_the_margin_is_never_pinned_back_into_line(
        make_prs, en_profile, tmp_path):
    """The pin is a real fix, pre-selectable and confident. A heading already
    outside the frame must reach the report instead of the fix engine."""
    prs = make_prs()
    stray = _repeating_title_deck(prs, stray_left=100000)  # past the left
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    assert [r.shape_id for r in _of_type(ctx, HEADING)] == [str(stray.shape_id)]
    assert _of_type(ctx, PINNED) == []


def test_a_heading_inside_the_margins_is_still_pinned_back_into_line(
        make_prs, en_profile, tmp_path):
    """Only the margin case is excused. A title that simply wandered off the
    line it holds on every other slide is still a defect worth fixing."""
    prs = make_prs()
    stray = _repeating_title_deck(prs, stray_left=1900000)  # inside, but off
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    assert _of_type(ctx, HEADING) == []
    assert [r.shape_id for r in _of_type(ctx, PINNED)] == [str(stray.shape_id)]


# ------------------------------------- headings drive no rescale either way


def _text(slide, *, left, top, width, height, text, pt):
    shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width),
                                     Emu(height))
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(pt)
    return shape


def _export_slide(prs, *, second_row):
    """An export-tool slide: no placeholders anywhere, a heading set large at
    the top, and panels of body copy. Breaching panels are the evidence the
    LAYOUT is too wide; the heading must never be counted among them."""
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    head = _text(slide, left=100000, top=200000, width=4000000, height=500000,
                 text="The Heading", pt=40)
    panels = []
    rows = [2000000] + ([4000000] if second_row else [])
    for top in rows:
        for left, width in ((100000, 4000000), (4300000, 4000000),
                            (8000000, 4200000)):
            panels.append(_text(slide, left=left, top=top, width=width,
                                height=1500000, text="Body copy", pt=12))
    if not second_row:  # one more left-breaching panel, so 3 panels breach
        panels.append(_text(slide, left=100000, top=4000000, width=4000000,
                            height=1500000, text="Body copy", pt=12))
    return slide, head, panels


def test_a_heading_is_never_evidence_that_the_layout_is_too_wide(
        make_prs, en_profile, tmp_path):
    """Three panels cross the frame; the heading would be a fourth and tip the
    slide into a full-deck rescale. It must not count."""
    prs = make_prs()
    _slide, head, _panels = _export_slide(prs, second_row=False)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    assert _of_type(ctx, OVERFLOW) == []
    # and the heading itself is recognised and reported, on an export deck
    # that carries no placeholder to say so
    flagged = _of_type(ctx, HEADING)
    assert [r.shape_id for r in flagged] == [str(head.shape_id)]


def test_a_heading_is_never_part_of_the_rescale_selection(
        make_prs, en_profile, tmp_path):
    """With enough panels breaching on their own the rescale is offered, and
    the heading is still left out of the shapes it would move."""
    prs = make_prs()
    _slide, head, panels = _export_slide(prs, second_row=True)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = _of_type(ctx, OVERFLOW)
    assert len(records) == 1
    ids = records[0].locator.split(":", 3)[3].split(",")
    assert str(head.shape_id) not in ids
    assert all(str(p.shape_id) in ids for p in panels)


def test_a_lone_text_box_is_not_promoted_to_a_heading(make_prs, en_profile,
                                                      tmp_path):
    """Inference needs hierarchy to read. One box on a slide is not a title
    just because nothing competes with it, and calling it one would quietly
    excuse it from every check a heading is excused from."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    box = slide.shapes.add_textbox(Emu(0), Emu(500000), Emu(3000000),
                                   Emu(600000))
    box.text_frame.text = "one line, high on the slide"
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    assert _of_type(ctx, HEADING) == []
    assert [r.shape_id for r in _of_type(ctx, SAFE_ZONE)] == [str(box.shape_id)]


# --------------------------------------------- the formatting pass's answer


def _guided_heading_deck(*, title_box, guides=(0.60, 12.73)):
    """A deck whose master DRAWS its margins, and whose title placeholder is
    placed where the test wants relative to them."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(SLIDE_W), Emu(SLIDE_H)
    _plant_guides(prs.slide_masters[0], vertical_in=guides,
                  horizontal_in=(0.42, 7.05))
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        if kind.startswith(("TITLE", "CENTER_TITLE")):
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in title_box)
            ph.text_frame.text = "Next: Personalized Daily Digests"
        elif kind.startswith("SUBTITLE"):
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(0.60 * IN)), Emu(int(1.40 * IN)),
                Emu(int(11.0 * IN)), Emu(int(0.35 * IN)))
            ph.text_frame.text = "A second LLM integration."
    for i in range(3):
        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(int((0.60 + i * 4.0) * IN)),
            Emu(int(2.4 * IN)), Emu(int(3.6 * IN)), Emu(int(1.9 * IN)))
        card.text_frame.text = f"Card {i + 1}"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _title_box_of(deck_bytes):
    slide = Presentation(io.BytesIO(deck_bytes)).slides[0]
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False) and str(
                shape.placeholder_format.type).startswith(
                    ("TITLE", "CENTER_TITLE")):
            return (shape.left, shape.top, shape.width, shape.height)
    return None


def test_the_formatting_pass_reports_a_wide_heading_and_moves_nothing():
    source = _guided_heading_deck(title_box=(0.30, 0.42, 12.60, 0.92))
    before = _title_box_of(source)

    out, changes = migrate_deck(source)

    notes = [c for c in changes if c.action == "heading past the margin"]
    assert len(notes) == 1
    assert "left, right" in notes[0].detail
    assert "ask them" in notes[0].detail
    assert notes[0].severity == "alert"
    assert _title_box_of(out) == before, "the heading must not have moved"


def test_a_heading_inside_the_drawn_frame_is_not_reported():
    source = _guided_heading_deck(title_box=(0.60, 0.45, 12.13, 0.92))
    _out, changes = migrate_deck(source)
    assert [c for c in changes if c.action == "heading past the margin"] == []


def test_a_stated_frame_seats_the_block_and_reports_the_bottom_overflow():
    """All four sides are read, and they divide the way a translate forces them
    to: the stated TOP binds, the bottom reports.

    This test used to assert the opposite - that a stated bottom margin holds
    the block back from its top line. That rule made alignment conditional on
    how tall a slide's content happened to be, so a deck came back with some
    slides on the frame and some not, which is the complaint it caused (design
    lead, 20-21/08/2026, twice: first for the header band, then for the
    presentation space). A 0.10in overflow that is reported beats a top edge
    that wanders."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(SLIDE_W), Emu(SLIDE_H)
    _plant_guides(prs.slide_masters[0], vertical_in=(0.60, 12.73),
                  horizontal_in=(0.42, 6.00))  # bottom margin at 6.00in
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        if kind.startswith(("TITLE", "CENTER_TITLE")):
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(0.60 * IN)), Emu(int(0.42 * IN)),
                Emu(int(12.13 * IN)), Emu(int(0.92 * IN)))
            ph.text_frame.text = "A heading"
        elif kind.startswith("SUBTITLE"):
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(0.60 * IN)), Emu(int(1.40 * IN)),
                Emu(int(12.13 * IN)), Emu(int(0.35 * IN)))
            ph.text_frame.text = "A standfirst"
    # A block just above the body margin, tall enough that the 0.10in it would
    # be pushed down would carry its foot past the bottom guide. Clear of the
    # header placeholders on purpose, so the collision pass plays no part and
    # the bottom margin is the only thing that can hold it.
    top, height = int(1.80 * IN), int(4.20 * IN)
    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(0.60 * IN)),
                                   Emu(top), Emu(int(6.0 * IN)), Emu(height))
    block.text_frame.text = "BLOCK"
    buf = io.BytesIO()
    prs.save(buf)

    out, changes = migrate_deck(buf.getvalue())
    moved = next(s for s in Presentation(io.BytesIO(out)).slides[0].shapes
                 if s.has_text_frame and s.text_frame.text == "BLOCK")
    assert moved.top == int(1.90 * IN), "seated on the line the master states"
    assert moved.top + moved.height > int(6.00 * IN), "so it overflows"
    spill = [c for c in changes if c.action == "content does not fit"]
    assert spill, "and the overflow past the bottom guide has to be reported"
    assert "past the bottom margin" in spill[0].detail
