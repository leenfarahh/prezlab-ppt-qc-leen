"""The pause between an upload and a rewritten file.

Applying a master was one press that guessed at the slides it could not place,
and the only way to see what it had guessed was to open the result - by which
point the guess was in the file. The step in between is where a designer says
which layout each of those slides belongs on.

What the route has to get right:

  - the first press REBUILDS NOTHING, so a wrong deck or a wrong profile costs a
    moment rather than a rebuild;
  - the picks posted are the layouts applied, or the page is decoration;
  - a host with no renderer still gets a usable page, in words;
  - and an expired plan says so rather than 500ing, because the layout page is
    the one page in the tool a designer sits on for a while.

apply_master is stubbed throughout: everything asserted here has to hold on a
machine with no PowerPoint.
"""

import io

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Emu, Inches

from qc import web
from qc.applymaster import ApplyResult


# ------------------------------------------------------------------ fixtures


def _bytes(prs) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _master() -> bytes:
    return _bytes(Presentation())


def _deck(slides=2) -> bytes:
    """A deck of blank slides carrying two columns of content. Blank matches the
    master's own 'Blank' layout by name, so nothing is uncertain until a test
    makes it so."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    for _ in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for left in (0.8, 7.0):
            box = slide.shapes.add_textbox(Inches(left), Inches(2.0),
                                           Inches(5.0), Inches(3.0))
            box.text_frame.text = "Point"
    return _bytes(prs)


@pytest.fixture()
def client(monkeypatch):
    monkeypatch.setattr(web, "AUTH_REQUIRED", False)
    monkeypatch.setattr(web.app.state, "auth_required", False)
    web._plans.clear()
    return TestClient(web.app)


@pytest.fixture()
def no_powerpoint(monkeypatch):
    seen = {}

    def _fake(deck_bytes, master_bytes, plans):
        seen["plans"] = plans
        out = Presentation(io.BytesIO(deck_bytes))
        buf = io.BytesIO()
        out.save(buf)
        return ApplyResult(deck=buf.getvalue(), plans=plans, masters=1)

    monkeypatch.setattr("qc.applymaster.apply_master", _fake)
    monkeypatch.setattr("qc.engine.run_audit", lambda *a, **k: type(
        "R", (), {"to_manifest": lambda self: {
            "slides": 2, "summary": {"total": 0, "by_severity": {},
                                     "arabic_flagged": 0},
            "records": [], "profile_id": "p", "profile_version": 1}})())
    return seen


@pytest.fixture()
def no_renderer(monkeypatch):
    def _boom(*a, **k):
        raise RuntimeError("PowerPoint is not on this machine")

    monkeypatch.setattr("qc.render.export_decks_png", _boom)


@pytest.fixture()
def profile_with_master(monkeypatch):
    """A profile carrying a real master, saved into the test's own store."""
    from qc.profile import PROFILES_DIR
    from qc.templates import save_master
    import json

    pid = "layout_step_client"
    (PROFILES_DIR / f"{pid}.json").write_text(json.dumps({
        "id": pid, "name": "Layout Step Client", "version": 1,
        "config": {}}), encoding="utf-8")
    save_master(pid, _master())
    return pid


def _uncertain(monkeypatch):
    """Force every slide to be a fallback, so the page has something to ask."""
    import qc.applymaster as AM

    real = AM.plan_assignments

    def _all_fallback(deck_prs, layouts):
        plans = real(deck_prs, layouts)
        for plan in plans:
            plan.match_rule = "fallback"
            plan.note = "no layout named 'Blank'"
        return plans

    monkeypatch.setattr(AM, "plan_assignments", _all_fallback)


def _start(client, pid, slides=2):
    return client.post("/prep", data={"profile": pid},
                       files={"deck": ("rough.pptx", _deck(slides), "app/x")})


# --------------------------------------------------------- the first press


def test_the_first_press_rebuilds_nothing(client, profile_with_master,
                                          no_renderer, monkeypatch):
    """The whole point of the pause. A wrong deck or a wrong profile has to
    cost a moment, not a rebuild."""
    def _never(*a, **k):
        raise AssertionError("the master was applied before it was approved")

    monkeypatch.setattr("qc.applymaster.apply_master", _never)
    _uncertain(monkeypatch)

    page = _start(client, profile_with_master)
    assert page.status_code == 200
    assert "Step 2 of 3" in page.text
    assert "slides to place" in page.text
    assert len(web._plans) == 1


def test_the_page_offers_only_the_masters_own_layouts(client,
                                                      profile_with_master,
                                                      no_renderer, monkeypatch):
    _uncertain(monkeypatch)
    page = _start(client, profile_with_master, slides=1).text

    assert 'value="Two Content"' in page
    assert 'name="pick_0"' in page
    # and the answer that is always available
    assert 'value="__leave__"' in page


def test_a_host_with_no_renderer_still_gets_a_usable_page(client,
                                                          profile_with_master,
                                                          no_renderer,
                                                          monkeypatch):
    """The choice is readable from the files alone - what the slide holds and
    what each layout offers are both arithmetic - so a missing renderer costs
    the pictures and nothing else."""
    _uncertain(monkeypatch)
    page = _start(client, profile_with_master, slides=1).text

    assert "could not be rendered" in page
    assert "content box" in page, "what each layout offers is still stated"
    assert 'name="pick_0"' in page, "and the choice is still makeable"


def test_a_matched_deck_with_no_renderer_says_so_rather_than_404ing(
        client, profile_with_master, no_renderer):
    """The page hands out an image URL per slide and renders them on demand, so
    the ONE render a fully-matched deck attempts - the layout catalogue - is
    also the only evidence about whether this host can render at all. Swallow
    its failure and every card points at a 404."""
    page = _start(client, profile_with_master).text

    assert "could not be rendered" in page
    assert "/plan-img/" not in page, "image URLs on a host that cannot render"
    assert 'name="pick_0"' in page, "the choice is still makeable"


def test_every_slide_matching_still_lists_every_slide(
        client, profile_with_master, no_renderer):
    """Certainty is not a reason to rebuild without being asked, and it is not a
    reason to hide the deck either.

    This used to assert a page with no slides on it at all: "Every slide matched
    a layout in this master... so there is nothing to choose", and a press. The
    press was right and the empty page was not - a designer approving a rebuild
    of their client's deck is approving what happens to each slide, and they
    could not see any of them (design lead, 02/09/2026)."""
    page = _start(client, profile_with_master)
    flat = " ".join(page.text.split())

    assert page.status_code == 200
    assert "Every slide matched" in flat
    assert "Apply master to all 2 slides" in flat
    # both slides are on the page, both changeable
    assert 'name="pick_0"' in page.text and 'name="pick_1"' in page.text
    assert flat.count("Matched") >= 2
    assert "not listed because there is nothing to decide" not in flat


def test_a_deck_whose_choices_cannot_be_built_still_gets_a_page_and_a_press(
        client, profile_with_master, no_renderer, monkeypatch):
    """The degraded path, which is now the ONLY thing that reaches the
    no-choices page.

    It used to be the "every slide matched" page as well, and that conflation
    is what let a working deck arrive at a page with no slides on it. Now a
    matched deck gets the ordinary page and only a plan whose per-slide review
    could not be computed lands here - qc.prep.plan swallows that rather than
    raising, because a deck that rebuilds on the targets plan_assignments picked
    is still a rebuilt deck."""
    monkeypatch.setattr("qc.layoutpick.choices",
                        lambda *a, **k: (_ for _ in ()).throw(RuntimeError("x")))

    page = _start(client, profile_with_master)

    assert page.status_code == 200
    assert "could not be built for this deck" in page.text
    assert "Apply master to 2 slides" in " ".join(page.text.split())


def test_showing_every_slide_does_not_render_every_slide(
        client, profile_with_master, monkeypatch):
    """The first press has to stay cheap. Every slide is on the page, and the
    matched cards are folded with lazy images, so nothing is rendered for them
    until somebody opens one - otherwise a 200-slide deck would go through
    PowerPoint before the page appeared."""
    calls = []

    def _count(decks, indices, *a, **k):
        calls.append(sorted(indices))
        return {f"s:{i}": b"png" for i in indices}

    monkeypatch.setattr("qc.render.export_decks_png", _count)
    monkeypatch.setattr("qc.render.layout_catalogue",
                        lambda *a, **k: (b"", [], None))

    page = _start(client, profile_with_master)
    plan_id = next(reversed(web._plans))

    assert web._plans[plan_id].get("shots", {}) == {},         "a matched deck went through PowerPoint before anyone asked"
    slide_calls = [c for c in calls if c and max(c) < 2]
    assert slide_calls == [], f"slides rendered up front: {slide_calls}"
    # and the cards still point at a picture, so the browser can fetch one
    assert f"/plan-img/{plan_id}/slide-0.png" in page.text
    assert 'loading="lazy"' in page.text


def test_opening_a_folded_card_renders_its_window(
        client, profile_with_master, monkeypatch):
    """And when somebody does ask, the picture arrives - a window at a time,
    because PowerPoint charges per call rather than per slide."""
    monkeypatch.setattr("qc.render.export_decks_png",
                        lambda decks, indices, *a, **k:
                        {f"s:{i}": b"png" for i in indices})
    monkeypatch.setattr("qc.render.layout_catalogue",
                        lambda *a, **k: (b"", [], None))
    _start(client, profile_with_master)
    plan_id = next(reversed(web._plans))

    got = client.get(f"/plan-img/{plan_id}/slide-1.png")

    assert got.status_code == 200
    assert got.content == b"png"
    # the window came with it, so the neighbour is free
    assert 0 in web._plans[plan_id]["shots"]


def test_a_matched_slide_says_where_it_is_going_and_can_be_moved(
        client, profile_with_master, no_renderer):
    """The two things the old page could not do: show the layout each slide is
    headed for, and let a designer disagree with one."""
    page = _start(client, profile_with_master).text

    assert "Going onto" in page
    assert "Where it is now" in page, "the pre-selected option says what it is"
    assert 'value="Blank" checked' in page,         "a matched slide pre-selects where it already is"


def test_scrolling_past_a_matched_slide_decides_nothing(
        client, profile_with_master, no_powerpoint, no_renderer):
    """Every slide posts a radio now, and a matched card's radio is pre-set to
    where the slide already is. Passing those through to apply_picks would
    stamp "chosen by the designer" onto a whole deck nobody opened, and the
    coverage report is built on that distinction."""
    _start(client, profile_with_master)
    plan_id = next(reversed(web._plans))
    plan = web._plans[plan_id]["plan"]
    posted = {f"pick_{c.slide_index}": c.current for c in plan.choices}

    client.post(f"/prep/{plan_id}/layouts", data=posted)

    rules = {p.match_rule for p in no_powerpoint["plans"]}
    assert "chosen" not in rules, "a slide nobody looked at was recorded as decided"
    assert all(not p.review for p in no_powerpoint["plans"])


def test_moving_a_matched_slide_is_recorded_as_the_designers(
        client, profile_with_master, no_powerpoint, no_renderer):
    """The capability that did not exist before: overruling a match the file
    made. It is a real decision and it is recorded as one."""
    _start(client, profile_with_master)
    plan_id = next(reversed(web._plans))
    plan = web._plans[plan_id]["plan"]
    first = plan.choices[0]
    assert first.settled and first.current != "Two Content"

    client.post(f"/prep/{plan_id}/layouts",
                data={f"pick_{first.slide_index}": "Two Content"})

    moved = no_powerpoint["plans"][first.slide_index]
    assert moved.target_layout == "Two Content"
    assert moved.match_rule == "chosen"
    assert moved.review == "chosen by the designer"


# -------------------------------------------------------- the second press


def test_the_picks_posted_are_the_layouts_applied(client, profile_with_master,
                                                  no_powerpoint, no_renderer,
                                                  monkeypatch):
    _uncertain(monkeypatch)
    _start(client, profile_with_master)
    plan_id = next(reversed(web._plans))

    done = client.post(f"/prep/{plan_id}/layouts",
                       data={"pick_0": "Two Content",
                             "pick_1": "Section Header"})
    assert done.status_code == 200

    applied = {p.slide_index: p.target_layout for p in no_powerpoint["plans"]}
    assert applied == {0: "Two Content", 1: "Section Header"}


def test_the_select_beats_the_radio(client, profile_with_master,
                                    no_powerpoint, no_renderer, monkeypatch):
    """The radio carries a default nobody had to touch; the select is only
    reachable by opening it and choosing."""
    _uncertain(monkeypatch)
    _start(client, profile_with_master, slides=1)
    plan_id = next(reversed(web._plans))

    client.post(f"/prep/{plan_id}/layouts",
                data={"pick_0": "Two Content", "other_0": "Picture with Caption"})
    assert no_powerpoint["plans"][0].target_layout == "Picture with Caption"


def test_leaving_a_slide_rebuilds_it_and_reports_the_gap(client,
                                                         profile_with_master,
                                                         no_powerpoint,
                                                         no_renderer,
                                                         monkeypatch):
    """"None of these fit" is a real answer and the page has to carry it all the
    way through: the slide still rebuilds on whatever the file gave it, and the
    refusal is recorded so the coverage report can count it as a gap in the
    master rather than a slide nobody opened."""
    from qc.layoutpick import LEAVE

    _uncertain(monkeypatch)
    _start(client, profile_with_master, slides=1)
    plan_id = next(reversed(web._plans))

    done = client.post(f"/prep/{plan_id}/layouts", data={"pick_0": LEAVE})
    plan = no_powerpoint["plans"][0]
    assert plan.match_rule == "fallback"
    assert plan.review == "no fit"
    assert "no layout in this master that fits" in done.text


def test_when_nothing_fits_the_default_is_to_say_so(client,
                                                    profile_with_master,
                                                    no_renderer, monkeypatch):
    """A master with no layout that can hold the slide should not have one
    pre-selected anyway. The gap is the answer, and pre-picking a near-miss
    hides it behind a choice nobody made."""
    import qc.layoutpick as LP

    _uncertain(monkeypatch)
    # Every candidate reports that it does not fit.
    real_rank = LP.rank

    def _never_fits(*a, **kw):
        cands, sig = real_rank(*a, **kw)
        for c in cands:
            c.fits = False
        return cands, sig

    monkeypatch.setattr(LP, "rank", _never_fits)

    page = _start(client, profile_with_master, slides=1).text
    leave = page.split(f'value="{LP.LEAVE}"')[1][:40]
    assert "checked" in leave, "the refusal is not the default when it should be"


def test_pressing_apply_with_no_picks_uses_the_suggestions(
        client, profile_with_master, no_powerpoint, no_renderer, monkeypatch):
    _uncertain(monkeypatch)
    _start(client, profile_with_master, slides=1)
    plan_id = next(reversed(web._plans))

    done = client.post(f"/prep/{plan_id}/layouts", data={})
    assert done.status_code == 200
    assert no_powerpoint["plans"][0].target_layout


def test_the_run_says_what_was_decided(client, profile_with_master,
                                       no_powerpoint, no_renderer, monkeypatch):
    """The layout step's own record, on the page that comes out of it. A
    decision nobody can see afterwards is a decision nobody will trust."""
    _uncertain(monkeypatch)
    _start(client, profile_with_master, slides=1)
    plan_id = next(reversed(web._plans))

    done = client.post(f"/prep/{plan_id}/layouts",
                       data={"pick_0": "Section Header"})
    assert "needed a layout decision" in done.text


def test_the_plan_is_dropped_once_it_has_been_spent(client,
                                                    profile_with_master,
                                                    no_powerpoint, no_renderer,
                                                    monkeypatch):
    """Its bytes live on the job afterwards. Keeping both is a second copy of a
    client deck in memory for no reason."""
    _uncertain(monkeypatch)
    _start(client, profile_with_master, slides=1)
    plan_id = next(reversed(web._plans))

    client.post(f"/prep/{plan_id}/layouts", data={})
    assert plan_id not in web._plans


# ------------------------------------------------------------ what goes wrong


def test_an_expired_plan_says_start_again(client):
    page = client.get("/prep/deadbeef/layouts")
    assert page.status_code == 404
    assert "expired" in page.text

    posted = client.post("/prep/deadbeef/layouts", data={})
    assert posted.status_code == 404
    assert "expired" in posted.text


def test_an_image_for_a_plan_that_is_gone_is_a_404_not_a_crash(client):
    assert client.get("/plan-img/deadbeef/slide-0.png").status_code == 404


def test_a_junk_image_key_is_refused(client, profile_with_master, no_renderer,
                                     monkeypatch):
    _uncertain(monkeypatch)
    _start(client, profile_with_master, slides=1)
    plan_id = next(reversed(web._plans))

    for key in ("slide-nope", "layout-99", "whatever"):
        assert client.get(f"/plan-img/{plan_id}/{key}.png").status_code == 404


def test_the_layout_page_can_be_reloaded(client, profile_with_master,
                                         no_renderer, monkeypatch):
    """A designer sits on this page. It has to survive a refresh, and it must
    not re-render the slides to do it."""
    _uncertain(monkeypatch)
    _start(client, profile_with_master, slides=1)
    plan_id = next(reversed(web._plans))

    again = client.get(f"/prep/{plan_id}/layouts")
    assert again.status_code == 200
    assert "Step 2 of 3" in again.text
