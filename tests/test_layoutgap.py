"""The deck-level layout question: what does this deck need that the master
does not have.

qc.layoutpick answers it per slide, and eleven fallback rows are eleven
questions. This pass turns them into the answer a designer acts on - "six of
these want the same two-column layout, and this master has none" - so what the
tests protect is that the count is trustworthy:

  - slides asking for the same thing are ONE gap, not six;
  - a matched slide is never a gap, whatever it looks like;
  - a group counts as one thing to place, because that is what it is;
  - "nobody looked at these" and "the review found no home for these" are
    counted apart, since only the second is evidence about the master;
  - and the closest layout named is one the master really has.
"""

import io

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from qc import web
from qc.applymaster import SlidePlan, plan_assignments
from qc.layoutgap import Coverage, describe, headline, report, signature
from qc.stylespec import dominant_master, extract_layouts


# ------------------------------------------------------------------ helpers


def _prs():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    return prs


def _bytes(prs) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _reopen(prs):
    return Presentation(io.BytesIO(_bytes(prs)))


def _title(slide, text="A heading"):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11),
                                  Inches(0.9))
    box.text_frame.text = text
    return box


def _block(slide, left, top=2.0, width=5.0, height=3.0, text="Point one"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                   Inches(height))
    box.text_frame.text = text
    return box


def _two_column_deck(slides=4):
    prs = _prs()
    for n in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _title(slide, f"Comparison {n + 1}")
        _block(slide, 0.8)
        _block(slide, 7.0)
    return prs


def _lean_master():
    """A master offering only a section header and a title-only layout: nothing
    named the way a deck's layouts are named, and nothing with two columns."""
    master = Presentation()
    return [entry for entry in extract_layouts(dominant_master(master),
                                               embed_assets=False)
            if entry["type"] in ("titleOnly", "secHead")]


def _coverage(prs, layouts=None, review=None):
    layouts = layouts if layouts is not None else _lean_master()
    deck = _reopen(prs)
    plans = plan_assignments(deck, layouts)
    for plan in plans:
        if review and plan.match_rule == "fallback":
            plan.review = review
    return report(deck, layouts, plans), plans


# ---------------------------------------------------------------- signatures


def test_the_same_request_stated_four_times_is_one_gap():
    """The whole point. Four two-column slides are one missing layout, and a
    designer who reads that adds one layout instead of opening four slides."""
    cov, _plans = _coverage(_two_column_deck(4))
    assert cov.unplaced == 4
    assert len(cov.gaps) == 1
    gap = cov.gaps[0]
    assert gap.places == 4 and gap.slides == [0, 1, 2, 3]
    assert "2 columns" in gap.label


def test_a_stack_is_one_column_and_a_row_is_several():
    """Adjacency down the page is not a column count. A stack of four blocks is
    one column of content; four across the slide are four."""
    stacked = _prs()
    slide = stacked.slides.add_slide(stacked.slide_layouts[6])
    _title(slide)
    for i in range(4):
        _block(slide, 1.0, top=1.6 + i * 1.3, width=5.0, height=1.1)
    reopened = _reopen(stacked)
    sig = signature(reopened.slides[0], reopened.slide_width,
                    reopened.slide_height)
    assert sig["columns"] == 1 and sig["blocks"] == 4

    across = _prs()
    slide = across.slides.add_slide(across.slide_layouts[6])
    _title(slide)
    for i in range(4):
        _block(slide, 0.5 + i * 3.2, top=2.0, width=2.8, height=3.0)
    reopened = _reopen(across)
    sig = signature(reopened.slides[0], reopened.slide_width,
                    reopened.slide_height)
    assert sig["columns"] == 4 and sig["blocks"] == 4


def test_a_group_is_one_thing_to_place():
    """A card with its icon and its label is a single block the rebuild has to
    find a home for. Counting its members separately turns a two-card slide
    into a request for six boxes."""
    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide)
    for left in (1.0, 7.0):
        group = slide.shapes.add_group_shape()
        group.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left),
                               Inches(2.0), Inches(4.5), Inches(3.0))
        group.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.3), Inches(2.3),
                               Inches(0.8), Inches(0.8))
        label = group.shapes.add_textbox(Inches(left + 0.3), Inches(3.5),
                                         Inches(3.8), Inches(0.8))
        label.text_frame.text = "Card"

    reopened = _reopen(prs)
    sig = signature(reopened.slides[0], reopened.slide_width,
                    reopened.slide_height)
    assert sig["blocks"] == 2, "two cards, not six shapes"
    assert sig["columns"] == 2
    assert sig["groups"] == 2


def test_a_full_bleed_panel_is_a_ground_not_a_block():
    """A backdrop covering the slide is what the content sits ON. Counted as
    content, every slide that has one looks like it needs an extra box."""
    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                           Inches(13.333), Inches(7.5))
    _block(slide, 2.0, top=3.0, width=9.0, height=1.5, text="One statement")

    reopened = _reopen(prs)
    sig = signature(reopened.slides[0], reopened.slide_width,
                    reopened.slide_height)
    assert sig["blocks"] == 1


def test_the_label_reads_as_a_sentence():
    assert describe({"title": True, "blocks": 2, "columns": 2, "text": 2,
                     "images": 0, "charts": 0, "tables": 0, "groups": 0}) == \
        "a title over 2 columns of 2 text blocks"
    assert describe({"title": False, "blocks": 1, "columns": 1, "text": 0,
                     "images": 1, "charts": 0, "tables": 0, "groups": 0}) == \
        "no title, 1 picture"


# ------------------------------------------------------------------ coverage


def test_a_matched_slide_is_never_a_gap():
    """Matching by name is a fact the file states. A deck rebuilt on the master
    it came from has nothing missing, whatever its slides look like."""
    prs = _two_column_deck(3)
    layouts = extract_layouts(dominant_master(Presentation()),
                              embed_assets=False)
    cov, _plans = _coverage(prs, layouts)
    assert cov.by_rule.get("name") == 3
    assert cov.unplaced == 0 and cov.gaps == []
    assert "Every one of the 3 slides" in headline(cov)


def test_the_report_names_the_closest_layout_the_master_really_has():
    """"The master has nothing for this" is only actionable next to what it does
    have, and the name has to be one a designer can go and open."""
    layouts = _lean_master()
    cov, _plans = _coverage(_two_column_deck(2), layouts)
    names = {entry["name"] for entry in layouts}
    gap = cov.gaps[0]
    assert gap.closest in names
    assert "wants 2 blocks in 2 columns" in gap.closest_note


def test_layouts_nothing_lands_on_are_listed():
    """A layout the deck never uses is worth seeing: when it is plainly what one
    of the unplaced slides was asking for, the names did not line up."""
    cov, _plans = _coverage(_two_column_deck(2))
    assert cov.used_layouts, "the fallbacks all landed somewhere"
    assert cov.unused_layouts
    assert set(cov.unused_layouts).isdisjoint(cov.used_layouts)


def test_nobody_looking_is_not_the_same_as_no_layout_fitting():
    """Only one of the two is evidence about the master, and a report that
    conflates them overstates its own case."""
    cold, _plans = _coverage(_two_column_deck(2))
    assert cold.review_ran is False
    assert cold.not_reviewed == 2
    assert cold.gaps[0].refused == 0

    looked, _plans = _coverage(_two_column_deck(2), review="no fit")
    assert looked.review_ran is True
    assert looked.not_reviewed == 0
    assert looked.gaps[0].refused == 2
    assert looked.gaps[0].reviewed == 2

    invented, _plans = _coverage(_two_column_deck(2), review="unusable answer")
    assert invented.gaps[0].refused == 0, (
        "an invented layout name says nothing about what the master offers")
    assert invented.gaps[0].reviewed == 2


def test_the_headline_leads_with_the_biggest_gap():
    cov, _plans = _coverage(_two_column_deck(5))
    line = headline(cov)
    assert "5 slides of 5" in line and "want the same thing" in line


def test_an_empty_deck_says_so_rather_than_dividing_by_zero():
    assert "no slides" in headline(Coverage(slides=0))


def test_a_slide_the_reader_cannot_read_still_counts():
    """Totals that do not add up to the deck are worse than a coarse label. A
    plan pointing past the end of the deck is dropped, and one whose slide
    cannot be read goes in its own bucket."""
    layouts = _lean_master()
    deck = _reopen(_two_column_deck(1))
    plans = plan_assignments(deck, layouts)
    plans.append(SlidePlan(slide_index=99, source_layout="Ghost",
                           source_type=None, target_layout=layouts[0]["name"],
                           match_rule="fallback"))
    cov = report(deck, layouts, plans)
    assert cov.slides == 2, "the plan count is the deck's own"
    assert sum(g.places for g in cov.gaps) == 1, "the ghost slide is not invented"


# ---------------------------------------------------------------- the page


def _client(monkeypatch):
    monkeypatch.setattr(web, "AUTH_REQUIRED", False)
    web.app.state.auth_required = False
    return TestClient(web.app)


def _profile_with_master(client, name="Coverage Client"):
    from qc.store import add_user

    add_user("Lead", "lead")
    client.post("/whoami", json={"name": "Lead"})
    client.post("/master", files={"master": ("brand.pptx",
                                             _bytes(Presentation()), "app/x")})
    # The LAST spec, and popped once used. web._specs is module state shared by
    # every test in the run: reading the first entry picks up whatever an
    # earlier module left there, and leaving this one behind breaks the next
    # test that reads the first entry expecting its own.
    spec_id = next(reversed(web._specs))
    from urllib.parse import parse_qs, urlparse

    location = client.post(f"/spec/{spec_id}/profile", data={"name": name},
                           follow_redirects=False).headers["location"]
    web._specs.pop(spec_id, None)
    # Saving lands back on Prepare a deck with the new profile picked, rather
    # than in a profile editor: /prep?saved=<pid>.
    return parse_qs(urlparse(location).query)["saved"][0]


def _recheck(client, pid, deck_bytes, master=None,
             master_name="brand.pptx", check_id="covjob"):
    """Coverage, through the one route that still reports it.

    There is no pre-flight page any more: coverage is reported on the Prepare a
    deck result as a matter of course, and the only way to ask for it again is
    the "check the revised master" form on that page. So these drive that form,
    with the deck seeded exactly as a prepare run seeds it.
    """
    from qc.templates import load_master

    web._check_jobs[check_id] = {"deck": deck_bytes, "filename": "rough.pptx",
                                 "profile": pid}
    return client.post(
        f"/check/{check_id}/again",
        files={"master": (master_name,
                          load_master(pid) if master is None else master,
                          "app/x")})


def test_there_is_no_pre_flight_page(monkeypatch):
    """The coverage report has one door: the prepared deck it belongs to."""
    client = _client(monkeypatch)
    assert client.get("/check").status_code == 404
    assert client.post("/check", data={"profile": "prezlab_en"},
                       files={"deck": ("d.pptx", b"x", "app/x")}
                       ).status_code == 404


def test_a_non_pptx_master_is_refused(monkeypatch):
    from qc.profile import PROFILES_DIR

    client = _client(monkeypatch)
    pid = _profile_with_master(client, "Junk Master Client")
    try:
        r = _recheck(client, pid, _bytes(_two_column_deck(1)),
                     master=b"x", master_name="notes.txt")
        assert r.status_code == 400
        assert "must be a .pptx" in r.text
    finally:
        (PROFILES_DIR / f"{pid}.json").unlink(missing_ok=True)
        web._check_jobs.clear()


def test_a_profile_with_no_master_has_nothing_to_check_against():
    """A pure read of the helper: no route reaches this branch now that the
    re-check always carries the file, and the message still has to be right."""
    _blob, _name, error = web._read_master(None, "prezlab_en")
    assert "carries no master file" in error

    _blob, _name, error = web._read_master(None, "")
    assert "upload the master" in error


def test_the_check_reports_coverage_without_touching_powerpoint(monkeypatch):
    """The reason this page exists. Applying a master is a rewrite that needs
    desktop PowerPoint and about a second a slide; the question of whether the
    master can build the deck at all is answerable by reading both files."""
    from qc.profile import PROFILES_DIR

    client = _client(monkeypatch)
    pid = _profile_with_master(client)
    try:
        r = _recheck(client, pid, _bytes(_two_column_deck(3)))
        assert r.status_code == 200
        assert "What this master can build" in r.text
        assert "Nothing has been changed" in r.text
        # The whole answer is read from the two files, with no model and no
        # renderer, which is why the page can be produced at all on a host with
        # neither. There is nothing weaker about it to disclose any more.
        assert "could not be looked at" not in r.text
    finally:
        (PROFILES_DIR / f"{pid}.json").unlink(missing_ok=True)


# ------------------------------------- the page when there IS something missing
#
# The tests above check a deck the master can build. This one forces a gap by
# renaming the deck's own layout, which is what an export tool or another
# agency's template does: the structure is fine and the NAME matches nothing, so
# the slides fall back.


def _deck_on_an_unknown_layout(slides=3) -> bytes:
    from pptx.oxml.ns import qn

    prs = _two_column_deck(slides)
    for slide in prs.slides:
        layout = slide.slide_layout._element
        cSld = layout.find(qn("p:cSld"))
        cSld.set("name", "Agency Two Column")
        # The archetype goes too. Matching runs by name and THEN by archetype,
        # and a layout that still says type="blank" matches the master's blank
        # one - which is exactly the right answer and no gap at all.
        if layout.get("type"):
            del layout.attrib["type"]
    return _bytes(prs)


def test_a_deck_the_master_cannot_build_gets_the_gap_and_the_reason(monkeypatch):
    """Coverage says what is missing. Proposing the layout to add needs a model,
    and when there is none the page says so rather than leaving a silence where
    the proposal would be."""
    from qc.profile import PROFILES_DIR

    client = _client(monkeypatch)
    pid = _profile_with_master(client, "Gap Client")
    try:
        r = _recheck(client, pid, _deck_on_an_unknown_layout())
        assert r.status_code == 200
        assert "no layout in this master" in r.text
        # no key on a test host, so the proposal cannot run and says why
        assert "Suggested layouts" in r.text
        assert "needs a model" in r.text or "No model key" in r.text
    finally:
        (PROFILES_DIR / f"{pid}.json").unlink(missing_ok=True)


def test_the_suggestion_window_appears_when_a_model_answers(monkeypatch):
    """The whole point of the window. Stubbed at qc.layoutsuggest's own seam, so
    what is under test is the route wiring rather than the model."""
    import qc.layoutsuggest as LS
    from qc.profile import PROFILES_DIR

    monkeypatch.setattr(web, "AI_ENABLED", True, raising=False)
    monkeypatch.setattr("qc.llm.api_configured", lambda: True)
    monkeypatch.setattr(LS, "ask_json", lambda **k: {
        "name": "Agency comparison", "archetype": "twoObj", "columns": 2,
        "why": "For slides that set two propositions against each other.",
        "boxes": [{"kind": "title", "column": 0, "label": "Heading"},
                  {"kind": "body", "column": 1, "label": "Left"},
                  {"kind": "body", "column": 2, "label": "Right"}]})
    # The review would render slides; this test is about the proposal, so the
    # renderer is switched off and the slides go unlooked-at.
    monkeypatch.setattr("qc.render.RENDERER", "none")

    client = _client(monkeypatch)
    pid = _profile_with_master(client, "Gap Client 2")
    try:
        r = _recheck(client, pid, _deck_on_an_unknown_layout())
        assert r.status_code == 200
        assert "What to add to the master" in r.text
        assert "Agency comparison" in r.text and "twoObj" in r.text
        assert "<svg" in r.text, "a wireframe, not a list of placeholder types"
        assert "does not edit a client's master" in r.text
    finally:
        (PROFILES_DIR / f"{pid}.json").unlink(missing_ok=True)


# ==================================== placed, and the wrong shape for it
#
# Matching runs by layout NAME first and a name match was never questioned: the
# designer called both layouts the same thing and meant them to correspond. That
# is a claim about INTENT, and it says nothing about whether the content fits the
# boxes. A two-column slide on a one-body layout is reported as matched and loses
# a column when PowerPoint remaps it (design lead, 26/08/2026).


def _named_deck(layout_index=1, columns=2, slides=3):
    """A deck built on one of the stock layouts, with its placeholders stripped
    and free shapes put in their place. The layout NAME still matches the
    master's, so matching settles it and the content decides nothing."""
    prs = _prs()
    layout = prs.slide_layouts[layout_index]
    for n in range(slides):
        slide = prs.slides.add_slide(layout)
        for ph in list(slide.placeholders):
            ph._element.getparent().remove(ph._element)
        _title(slide, f"Heading {n + 1}")
        for col in range(columns):
            _block(slide, 0.8 + col * 6.2, top=2.0, width=5.0, height=3.0)
    return prs


def _cov_of(prs, layouts=None):
    layouts = (layouts if layouts is not None
               else extract_layouts(dominant_master(Presentation()),
                                    embed_assets=False))
    deck = _reopen(prs)
    plans = plan_assignments(deck, layouts)
    return report(deck, layouts, plans), plans


def test_a_name_match_whose_content_does_not_fit_is_reported():
    """"Title and Content" offers one body box in one column. Three two-column
    slides land on it by name, and nothing before this said so."""
    cov, _plans = _cov_of(_named_deck())
    assert cov.unplaced == 0, "they were all placed"
    assert len(cov.misfits) == 3
    m = cov.misfits[0]
    assert m.layout == "Title and Content" and m.rule == "name"
    assert "2 columns" in m.label
    assert "1 column" in m.reason and "2 columns" in m.reason
    assert "1 content box" in m.offers


def test_the_headline_says_so_even_though_nothing_is_unplaced():
    """A deck can be fully placed and still wrong, and a headline that only
    counts unplaced slides would call this one clean."""
    cov, _plans = _cov_of(_named_deck())
    line = headline(cov)
    assert "Every slide has a layout" in line
    assert "does not fit" in line and "3 slides" in line


def test_a_layout_that_offers_no_boxes_is_never_a_misfit():
    """Every export-tool deck is a hundred slides of free shapes on Blank, and
    Blank makes no promises. A check that fires on every deck gets switched
    off."""
    cov, _plans = _cov_of(_named_deck(layout_index=6))   # Blank
    assert cov.misfits == []


def test_a_slide_whose_content_fits_is_not_a_misfit():
    """One column of content on a one-body layout is exactly right."""
    cov, _plans = _cov_of(_named_deck(columns=1))
    assert cov.misfits == []


def test_a_review_that_confirms_the_layout_clears_the_misfit():
    """The signature comparison reads the shapes, and the shapes are not the
    whole story - which is the entire reason for looking. A model that looked and
    said it fits overrules the comparison that raised it."""
    prs = _named_deck()
    layouts = extract_layouts(dominant_master(Presentation()),
                              embed_assets=False)
    deck = _reopen(prs)
    plans = plan_assignments(deck, layouts)
    for plan in plans:
        plan.review = "confirmed"
    assert report(deck, layouts, plans).misfits == []


def test_misfits_cluster_like_gaps_so_one_proposal_serves_them():
    """Six slides on a one-column layout that need two are one missing layout
    stated six times, exactly as six unplaced slides are."""
    cov, _plans = _cov_of(_named_deck(slides=4))
    assert len(cov.misfit_clusters) == 1
    gap = cov.misfit_clusters[0]
    assert gap.places == 4 and gap.slides == [0, 1, 2, 3]
    # `closest` is the layout they are ON, which is the useful comparison here
    assert gap.closest == "Title and Content"
    assert "1 content box" in gap.closest_note


def test_the_proposal_pass_answers_a_misfit_cluster(monkeypatch):
    import qc.layoutsuggest as LS

    monkeypatch.setattr(LS, "ask_json", lambda **k: {
        "name": "Two-column comparison", "archetype": "twoObj", "columns": 2,
        "why": "For slides that set two things against each other.",
        "boxes": [{"kind": "title", "column": 0},
                  {"kind": "body", "column": 1},
                  {"kind": "body", "column": 2}]})
    prs = _named_deck()
    cov, _plans = _cov_of(prs)
    assert not cov.gaps and cov.misfit_clusters

    out, asked, _unreach = LS.suggest(cov, _reopen(prs),
                            extract_layouts(dominant_master(Presentation()),
                                            embed_assets=False))
    assert asked == 1 and len(out) == 1
    assert out[0].name == "Two-column comparison" and out[0].places == 3


# ------------------------------------------------------- the page, and the loop


def test_the_page_lists_the_misfits_and_asks_for_the_revised_master(monkeypatch):
    from qc.profile import PROFILES_DIR

    client = _client(monkeypatch)
    pid = _profile_with_master(client, "Misfit Client")
    try:
        r = _recheck(client, pid, _bytes(_named_deck()))
        assert r.status_code == 200
        assert "Slides on a layout that does not fit them" in r.text
        assert "Title and Content" in r.text
        # and the loop: build it, come back with the file
        assert "Check the revised master" in r.text
        assert "/again" in r.text
    finally:
        (PROFILES_DIR / f"{pid}.json").unlink(missing_ok=True)


def test_the_recheck_reads_the_new_master_against_the_same_deck(monkeypatch):
    """The deck is held for it. A designer revising a layout three times should
    upload three masters, not three masters and three decks."""
    from qc.profile import PROFILES_DIR

    client = _client(monkeypatch)
    pid = _profile_with_master(client, "Misfit Client 2")
    try:
        # The "revised" master here is the stock template, which has a
        # two-content layout the deck's slides could use. The deck is the one
        # the prepare run already held under this id, so only the master is
        # uploaded.
        r = _recheck(client, pid, _bytes(_named_deck()),
                     master=_bytes(Presentation()), master_name="brand v2.pptx")
        assert r.status_code == 200
        assert "brand v2.pptx" in r.text, "it says which file it checked"
        assert "What this master can build" in r.text
    finally:
        (PROFILES_DIR / f"{pid}.json").unlink(missing_ok=True)
        web._check_jobs.clear()


def test_a_recheck_for_a_forgotten_deck_says_so(monkeypatch):
    client = _client(monkeypatch)
    r = client.post("/check/deadbeef/again",
                    files={"master": ("m.pptx", _bytes(Presentation()), "app/x")})
    assert r.status_code == 404
    assert "no longer held in memory" in r.text


def test_the_uploaded_master_wins_over_the_profile_stored_one(monkeypatch):
    """The re-check is ABOUT the file the designer just changed, so the upload
    beats the copy the profile carries. Checking the stored one would answer a
    question nobody asked."""
    from qc.profile import PROFILES_DIR

    client = _client(monkeypatch)
    pid = _profile_with_master(client, "Upload Wins Client")
    try:
        r = _recheck(client, pid, _bytes(_named_deck()),
                     master=_bytes(Presentation()), master_name="brand v2.pptx")
        assert r.status_code == 200
        assert "brand v2.pptx" in r.text
    finally:
        (PROFILES_DIR / f"{pid}.json").unlink(missing_ok=True)
        web._check_jobs.clear()


def test_the_check_route_says_what_is_still_undecided(monkeypatch):
    """The re-check names the slides that still need a layout picked.

    This used to render those slides and ask a vision model to place them, and
    the pass was DARK: `export_decks_png` was called without being imported into
    scope, every call raised NameError, and a bare `except Exception` degraded
    it to "the slides could not be looked at". The page said a true-sounding
    sentence and the feature did nothing (30/08/2026).

    The model came out on 31/08/2026 and the guess with it - a re-check answers
    "did the layout I just built close the gap?" with arithmetic. What survives
    from that bug is the shape of the test: drive the route, and assert the
    page says something SPECIFIC rather than a sentence that would read the
    same if the code behind it had never run.
    """
    from qc.profile import PROFILES_DIR

    client = _client(monkeypatch)

    # Force every slide to be a fallback, so there is definitely something
    # undecided. Whether THIS deck happens to match THIS master by name is not
    # what is under test.
    import qc.applymaster as AM

    _plan_for_real = AM.plan_assignments

    def _all_fallback(deck_prs, layouts):
        plans = _plan_for_real(deck_prs, layouts)
        for p in plans:
            p.match_rule = "fallback"
        return plans

    monkeypatch.setattr(AM, "plan_assignments", _all_fallback)

    pid = _profile_with_master(client)
    try:
        r = _recheck(client, pid, _bytes(_two_column_deck(3)))
        assert r.status_code == 200
        assert "still need a layout decision" in r.text, (
            "the re-check must count what is undecided, not go quiet about it")
        assert "NameError" not in r.text
    finally:
        (PROFILES_DIR / f"{pid}.json").unlink(missing_ok=True)
