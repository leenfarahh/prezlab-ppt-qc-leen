"""Layout matching for the master application: the one step of the format pass
that knows it does not know.

qc.applymaster places a slide by NAME, then by ARCHETYPE, then falls back. The
first two are facts read off the file. The fallback is a guess, and the slide is
rebuilt onto whatever came first in a preference list while PowerPoint moves its
content into boxes nobody chose - which is why the plan already labels it "content
may be orphaned, check this slide". A designer with both files open answers that
in two seconds, because they can SEE the structure.

What these tests protect is that looking can only ever IMPROVE a plan:

  - a name or archetype match is never sent, never revisited, never downgraded;
  - a layout name the master does not have is discarded, not applied;
  - anything unanswerable stays the fallback it was, still labelled as one;
  - an uncertain answer is marked uncertain rather than presented as a match.

The model is stubbed throughout (qc.llm.ask_json is the seam).
"""

import io

import pytest
from pptx import Presentation
from pptx.util import Emu

import qc.layoutmatch as LM
from qc.applymaster import SlidePlan

IN = 914400


def _plan(idx, rule, target="Content", source="Our approach", stype=None):
    return SlidePlan(slide_index=idx, source_layout=source, source_type=stype,
                     target_layout=target, match_rule=rule,
                     note="content may be orphaned" if rule == "fallback" else "")


def _deck(slides=2):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    for _ in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Emu(IN), Emu(IN), Emu(3 * IN), Emu(IN))
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture()
def stub(monkeypatch):
    """A master whose layouts are named, and a recording stub for the call."""
    monkeypatch.setattr(LM, "_layout_sheet",
                        lambda master: ([b"png1", b"png2", b"png3"],
                                        ["Title Slide", "Two Content",
                                         "Section Header"]))
    asked = []

    def _answer(reply):
        def _ask(**kwargs):
            asked.append(kwargs)
            return reply
        monkeypatch.setattr(LM, "ask_json", _ask)
        return asked

    return _answer


# ------------------------------------------------------- only fallbacks go


def test_a_matched_slide_is_never_sent(stub):
    """A name match is a fact the file states. Asking anyway spends a vision
    call to be told what we already knew."""
    asked = stub({"layout": "Two Content", "confident": True,
                  "rationale": "two columns"})
    plans = [_plan(0, "name"), _plan(1, "archetype")]
    out, reviewed = LM.review_fallbacks(_deck(), b"master", plans,
                                        {0: b"a", 1: b"b"})
    assert reviewed == 0 and asked == []
    assert [p.match_rule for p in out] == ["name", "archetype"]


def test_nothing_happens_when_there_are_no_fallbacks(stub):
    stub({"layout": "Two Content", "confident": True, "rationale": "x"})
    plans = [_plan(0, "name")]
    out, reviewed = LM.review_fallbacks(_deck(1), b"master", plans, {0: b"a"})
    assert reviewed == 0 and out[0].target_layout == "Content"


# ------------------------------------------------------------ placing one


def test_a_fallback_is_placed_and_labelled_as_reviewed(stub):
    asked = stub({"layout": "Two Content", "confident": True,
                  "rationale": "a two column comparison"})
    plans = [_plan(0, "name"), _plan(1, "fallback")]
    out, reviewed = LM.review_fallbacks(_deck(), b"master", plans,
                                        {0: b"a", 1: b"b"})
    assert reviewed == 1
    assert len(asked) == 1, "only the fallback was sent"
    assert out[0].match_rule == "name", "the matched slide is untouched"
    placed = out[1]
    assert placed.target_layout == "Two Content"
    assert placed.match_rule == "reviewed"
    assert "two column comparison" in placed.note


def test_an_unsure_answer_says_so_rather_than_passing_as_a_match(stub):
    """A false on `confident` is not a failure - it tells the designer to look,
    which beats a confident guess that quietly orphans a column of text."""
    stub({"layout": "Two Content", "confident": False,
          "rationale": "either of two would serve"})
    out, _reviewed = LM.review_fallbacks(_deck(1), b"master",
                                         [_plan(0, "fallback")], {0: b"a"})
    assert out[0].match_rule == "reviewed (uncertain)"
    assert "NOT confident" in out[0].note
    assert out[0].target_layout == "Two Content"


def test_the_name_is_matched_case_and_space_insensitively(stub):
    stub({"layout": "  two content  ", "confident": True, "rationale": "x"})
    out, _r = LM.review_fallbacks(_deck(1), b"master", [_plan(0, "fallback")],
                                  {0: b"a"})
    assert out[0].target_layout == "Two Content", \
        "the master's own spelling is what gets applied"


# --------------------------------------------------------- the closed set


def test_a_layout_the_master_does_not_have_is_discarded(stub):
    """The closed set is exactly what catches an invented name. A hallucinated
    layout must not reach PowerPoint."""
    stub({"layout": "Comparison With Images", "confident": True,
          "rationale": "it looks like a comparison"})
    out, reviewed = LM.review_fallbacks(_deck(1), b"master",
                                        [_plan(0, "fallback")], {0: b"a"})
    assert reviewed == 1, "it was asked and answered"
    assert out[0].match_rule == "fallback", "but the answer was not usable"
    assert out[0].target_layout == "Content"


def test_none_of_them_fit_is_a_real_answer(stub):
    stub({"layout": "", "confident": False, "rationale": "no counterpart"})
    out, _r = LM.review_fallbacks(_deck(1), b"master", [_plan(0, "fallback")],
                                  {0: b"a"})
    assert out[0].match_rule == "fallback"


def test_the_schema_offers_no_place_to_put_geometry():
    """The model picks a name. It never sees or sets a coordinate; assigning
    the layout still runs PowerPoint's own placeholder matching."""
    assert set(LM.MATCH_SCHEMA["properties"]) == {"layout", "confident",
                                                  "rationale"}
    assert LM.MATCH_SCHEMA["additionalProperties"] is False


# ------------------------------------------------------------- failure paths


def test_a_failed_call_leaves_the_fallback_alone(monkeypatch):
    monkeypatch.setattr(LM, "_layout_sheet",
                        lambda m: ([b"png"], ["Two Content"]))

    def _boom(**kwargs):
        raise RuntimeError("the model was unreachable")

    monkeypatch.setattr(LM, "ask_json", _boom)
    out, reviewed = LM.review_fallbacks(_deck(1), b"master",
                                        [_plan(0, "fallback")], {0: b"a"})
    assert reviewed == 0, "a failure is not a review"
    assert out[0].match_rule == "fallback" and out[0].target_layout == "Content"


def test_no_renderable_layouts_means_no_review(monkeypatch):
    """Without pictures of the master's layouts there is nothing to choose
    from, and guessing from names alone is what the archetype rule already
    does better."""
    monkeypatch.setattr(LM, "_layout_sheet", lambda m: (None, []))
    monkeypatch.setattr(LM, "ask_json",
                        lambda **k: pytest.fail("asked with no layouts"))
    out, reviewed = LM.review_fallbacks(_deck(1), b"master",
                                        [_plan(0, "fallback")], {0: b"a"})
    assert reviewed == 0 and out[0].match_rule == "fallback"


def test_a_slide_with_no_render_is_skipped(stub):
    asked = stub({"layout": "Two Content", "confident": True, "rationale": "x"})
    out, reviewed = LM.review_fallbacks(_deck(1), b"master",
                                        [_plan(0, "fallback")], {})
    assert reviewed == 0 and asked == []
    assert out[0].match_rule == "fallback"


def test_the_slides_structure_is_handed_over_with_the_picture(stub):
    """"Two columns" is a claim the numbers can support, so they go too."""
    asked = stub({"layout": "Two Content", "confident": True, "rationale": "x"})
    LM.review_fallbacks(_deck(1), b"master", [_plan(0, "fallback")], {0: b"a"})
    prompt = asked[0]["prompt"]
    assert "text_blocks" in prompt
    assert "Our approach" in prompt, "the source layout's name is context"
    assert "Two Content" in prompt, "the catalogue has to be in the prompt"
    # the slide first, then the layout sheet, in the order the prompt describes
    assert asked[0]["images"][0] == b"a"
    assert asked[0]["images"][1:] == [b"png1", b"png2", b"png3"]
