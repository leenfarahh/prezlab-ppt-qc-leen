"""Choosing a layout without asking a model.

This replaced a vision call, so what the tests protect is the property that
motivated the swap: the answer is a function of two files and nothing else. Same
deck, same master, same answer, on a machine with no network and no PowerPoint.

Four claims:

  - EVERY slide is put in front of a designer and only the uncertain ones are
    counted as questions, because a page that asks forty questions to surface
    four gets pressed through unread and a page that hides thirty six of them
    cannot be checked against the deck (02/09/2026);
  - a layout that FITS outranks a layout that is merely close, whatever the
    arithmetic says, since "this can hold your content" is a different class of
    answer from "this is nearly the right shape";
  - a pick is applied verbatim or dropped, never interpreted; and
  - nothing here reaches a model, a renderer or PowerPoint.
"""

import io

import pytest
from pptx import Presentation
from pptx.util import Inches

from qc.applymaster import SlidePlan, plan_assignments
from qc.layoutpick import (LEAVE, apply_picks, choices, note, rank,
                          undecided)
from qc.stylespec import dominant_master, extract_layouts


# ------------------------------------------------------------------ fixtures


def _prs():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    return prs


def _reopen(prs):
    buf = io.BytesIO()
    prs.save(buf)
    return Presentation(io.BytesIO(buf.getvalue()))


def _title(slide, text="A heading"):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11),
                                   Inches(0.9))
    box.text_frame.text = text


def _block(slide, left, top=2.0, width=5.0, height=3.0):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                   Inches(height))
    box.text_frame.text = "Point one"


def _two_column_deck(slides=3):
    prs = _prs()
    for n in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])   # Blank
        _title(slide, f"Comparison {n + 1}")
        _block(slide, 0.8)
        _block(slide, 7.0)
    return prs


def _layouts():
    return extract_layouts(dominant_master(Presentation()), embed_assets=False)


def _plans(prs, layouts, rule=None):
    deck = _reopen(prs)
    plans = plan_assignments(deck, layouts)
    if rule:
        for plan in plans:
            plan.match_rule = rule
    return deck, plans


# ------------------------------------------------- what gets asked about


def test_every_slide_is_offered_and_only_some_are_questions():
    """The deck is on the page; the questions are flagged.

    This used to assert the opposite - a slide that matched by name onto a
    layout its content fits produced NO Choice at all, and the page replaced it
    with a line saying how many were "not listed because there is nothing to
    decide about them". A designer approving a rebuild wants to see the deck
    (design lead, 02/09/2026).

    The reasoning that motivated the omission is intact and now lives in
    `settled`: nobody reads forty rows to find four, so the four are what the
    count and the styling lead on. What changed is that the other thirty six
    exist, can be checked, and can be moved."""
    layouts = _layouts()
    deck, plans = _plans(_two_column_deck(3), layouts)
    for plan in plans:
        assert plan.match_rule == "name", "the fixture must match by name"

    offered = choices(deck, layouts, plans)

    assert [c.slide_index for c in offered] == [0, 1, 2]
    assert all(c.settled for c in offered)
    assert undecided(offered) == 0, "a matched slide is not a question"


def test_a_settled_slide_is_pre_selected_where_it_already_is():
    """Not where the ranking would put it. The file's name match is a designer's
    stated intent and the score is arithmetic; a page that silently moves
    matched slides on load is a page whose defaults cannot be trusted."""
    layouts = _layouts()
    deck, plans = _plans(_two_column_deck(1), layouts)

    choice = choices(deck, layouts, plans)[0]

    assert choice.settled
    assert choice.suggested == choice.current
    assert choice.candidates[0].name == choice.current,         "where the slide is going has to be visible, not buried in the select"


def test_the_layout_a_settled_slide_is_on_is_always_one_of_its_options():
    """The shortlist is by score and a name match is not, so the layout a slide
    is already going onto can rank seventh. Pre-selecting something the card
    does not show would leave a designer looking at five wrong options and no
    indication of where the slide actually goes."""
    layouts = _layouts()
    deck, plans = _plans(_two_column_deck(2), layouts)

    for choice in choices(deck, layouts, plans):
        assert choice.current in {c.name for c in choice.candidates}
        assert len(choice.candidates) <= 5, "the shortlist is still a shortlist"


def test_a_slide_nothing_matched_is_offered_with_a_reason():
    layouts = _layouts()
    deck, plans = _plans(_two_column_deck(2), layouts, rule="fallback")

    offered = choices(deck, layouts, plans)
    assert [c.slide_index for c in offered] == [0, 1]
    for choice in offered:
        assert choice.rule == "fallback"
        assert choice.reason
        assert "2 columns" in choice.wants
        assert choice.candidates, "a choice with no options is not a choice"
        assert choice.suggested in {c.name for c in choice.candidates}


def test_every_option_is_a_layout_the_master_really_has():
    """The closed set is the safety property. A name that is not in the master
    cannot be applied, so the picker must never be able to produce one."""
    layouts = _layouts()
    real = {l["name"] for l in layouts}
    deck, plans = _plans(_two_column_deck(1), layouts, rule="fallback")

    for choice in choices(deck, layouts, plans):
        for candidate in choice.candidates:
            assert candidate.name in real


# ------------------------------------------------------------- the ranking


def test_a_layout_that_fits_outranks_one_that_is_merely_close():
    """Sorted on fit first and score second. A layout that can hold the content
    is a different kind of answer from one that is nearly the right shape, and
    burying it under a closer miss is how a designer picks the wrong one."""
    layouts = _layouts()
    deck = _reopen(_two_column_deck(1))
    ranked, _sig = rank(deck.slides[0], layouts, deck.slide_width,
                        deck.slide_height)

    fitting = [c.fits for c in ranked]
    assert fitting == sorted(fitting, reverse=True), (
        "a non-fitting layout sorted above a fitting one")


def test_the_ranking_is_stable_between_runs():
    """The property the vision pass could not offer. Two runs over the same two
    files must produce the same order, or a designer who reloads gets a
    different suggestion and cannot tell which one to trust."""
    layouts = _layouts()
    deck = _reopen(_two_column_deck(1))
    first, _ = rank(deck.slides[0], layouts, deck.slide_width,
                    deck.slide_height)
    second, _ = rank(deck.slides[0], layouts, deck.slide_width,
                     deck.slide_height)
    assert [c.name for c in first] == [c.name for c in second]


def test_the_source_archetype_earns_a_discount():
    """The OOXML type token is the file's own statement of what a layout is for,
    so a slide that came off a 'twoObj' layout should prefer the master's."""
    layouts = _layouts()
    deck = _reopen(_two_column_deck(1))
    plain, _ = rank(deck.slides[0], layouts, deck.slide_width,
                    deck.slide_height)
    biased, _ = rank(deck.slides[0], layouts, deck.slide_width,
                     deck.slide_height, source_type="twoObj")

    target = next(l["name"] for l in layouts if l["type"] == "twoObj")
    assert {c.name: c.score for c in biased}[target] \
        < {c.name: c.score for c in plain}[target]


def test_a_shared_word_in_the_layout_names_breaks_a_tie():
    """Weaker evidence than the archetype and still evidence: a deck's
    'Comparison' slide belongs on a master's 'Comparison' layout."""
    layouts = _layouts()
    deck = _reopen(_two_column_deck(1))
    plain, _ = rank(deck.slides[0], layouts, deck.slide_width,
                    deck.slide_height)
    biased, _ = rank(deck.slides[0], layouts, deck.slide_width,
                     deck.slide_height, source_name="Comparison grid")

    assert {c.name: c.score for c in biased}["Comparison"] \
        < {c.name: c.score for c in plain}["Comparison"]


def test_a_word_every_master_uses_is_not_evidence():
    """"Title" and "Slide" appear in half the layout names of any master.
    Matching on them would hand every layout the same discount, which is the
    same as handing none of them one, but noisier."""
    layouts = _layouts()
    deck = _reopen(_two_column_deck(1))
    plain, _ = rank(deck.slides[0], layouts, deck.slide_width,
                    deck.slide_height)
    biased, _ = rank(deck.slides[0], layouts, deck.slide_width,
                     deck.slide_height, source_name="Title Slide")

    assert {c.name: c.score for c in biased} == {c.name: c.score for c in plain}


# ---------------------------------------------------------------- the picks


def test_a_pick_is_written_onto_the_plan_and_recorded_as_a_decision():
    layouts = _layouts()
    _deck, plans = _plans(_two_column_deck(2), layouts, rule="fallback")
    wanted = layouts[-1]["name"]

    moved = apply_picks(plans, {0: wanted}, layouts)
    assert moved == 1
    assert plans[0].target_layout == wanted
    assert plans[0].match_rule == "chosen"
    assert plans[0].review == "chosen by the designer"
    # Untouched slides keep what the file gave them, and keep saying so.
    assert plans[1].match_rule == "fallback"


def test_confirming_the_suggestion_is_recorded_even_though_nothing_moved():
    """"The designer looked and agreed" and "nobody looked" are different facts
    about a slide, and the coverage report reads match_rule to tell them
    apart."""
    layouts = _layouts()
    _deck, plans = _plans(_two_column_deck(1), layouts, rule="fallback")
    already = plans[0].target_layout

    moved = apply_picks(plans, {0: already}, layouts)
    assert moved == 0
    assert plans[0].match_rule == "chosen"


def test_leaving_a_slide_is_recorded_as_a_refusal():
    """The most informative answer on the page, and the one easiest to throw
    away. A designer who reads the master against this slide and says none of
    these fit has established that the master is missing a layout; a slide
    nobody opened has established nothing, and the coverage report counts them
    apart (qc.layoutgap.Gap.refused)."""
    layouts = _layouts()
    _deck, plans = _plans(_two_column_deck(2), layouts, rule="fallback")
    before = plans[0].target_layout

    assert apply_picks(plans, {0: LEAVE}, layouts) == 0
    assert plans[0].review == "no fit"
    # It still rebuilds, on exactly what the file gave it.
    assert plans[0].target_layout == before
    assert plans[0].match_rule == "fallback"
    # and the slide nobody answered about is still unanswered
    assert plans[1].review == ""


def test_an_invented_layout_name_is_dropped_not_applied():
    """The only way to send one is to edit the form. A slide that keeps its
    computed target still rebuilds; a slide pointed at a layout that does not
    exist does not."""
    layouts = _layouts()
    _deck, plans = _plans(_two_column_deck(1), layouts, rule="fallback")
    before = plans[0].target_layout

    assert apply_picks(plans, {0: "A Layout Nobody Built"}, layouts) == 0
    assert plans[0].target_layout == before
    assert plans[0].match_rule == "fallback", (
        "a rejected pick must not be recorded as a decision")


def test_a_pick_for_a_slide_that_is_not_in_the_deck_is_ignored():
    layouts = _layouts()
    _deck, plans = _plans(_two_column_deck(1), layouts, rule="fallback")
    assert apply_picks(plans, {97: layouts[0]["name"]}, layouts) == 0


def test_the_name_matching_on_a_pick_is_forgiving_about_case_and_spacing():
    """The value posts back from a select the tool rendered, so it should match
    exactly - but a master with trailing whitespace in a layout name is common
    enough that failing closed on it would drop a legitimate pick."""
    layouts = _layouts()
    _deck, plans = _plans(_two_column_deck(1), layouts, rule="fallback")
    wanted = layouts[-1]["name"]

    apply_picks(plans, {0: f"  {wanted.upper()}  "}, layouts)
    assert plans[0].target_layout == wanted, "the stored spelling wins"


# ------------------------------------------------------------------ the note


def test_the_note_counts_what_happened_rather_than_asserting_it_went_well():
    assert "nothing to choose" in note(0, 0, 0)

    line = note(5, 5, 2)
    assert "5 slides needed a layout decision" in line
    assert "2 moved to a layout you picked" in line
    assert "3 kept the suggestion" in line

    left = note(4, 1, 1)
    assert "3 were left on the fallback" in left

    # A refusal is said as a gap in the master, not as an unanswered slide.
    refused = note(4, 1, 1, refused=3)
    assert "3 had no layout in this master that fits" in refused
    assert "left on the fallback" not in refused


def test_moving_a_slide_the_file_matched_is_counted_apart():
    """A designer overruling a match nobody flagged is not one of the questions
    the page asked, and it only became possible when every slide went on the
    page (02/09/2026). Counted separately so a run where the tool was never in
    doubt does not read as a run with open questions."""
    line = note(4, 3, 2, 0, overridden=1)

    assert "4 slides needed a layout decision" in line
    assert "2 moved to a layout you picked" in line
    assert "1 of the moves was a slide the file had already matched" in line
    # 4 questions = 1 moved + 1 kept + 2 nobody answered. The override is one of
    # the 2 moves and is NOT one of the 4.
    assert "1 kept the suggestion" in line
    assert "2 were left on the fallback" in line


def test_an_override_on_a_deck_with_no_questions_still_says_something():
    """Every slide matched and the designer moved one anyway. "There was
    nothing to choose" would be a run note that contradicts the file."""
    line = note(0, 1, 1, 0, overridden=1)

    assert "nothing to choose" not in line
    assert "1 moved to a layout you picked" in line
    assert "already matched" in line


# --------------------------------------------------------- what it never does


def test_nothing_here_asks_a_model_or_renders_anything(monkeypatch):
    """The property the whole swap was for. Applying a master is arithmetic
    over two files; a rebuild that depends on a network call produces a
    different deck on a bad afternoon."""
    def _never(*a, **k):
        raise AssertionError("the layout picker reached outside the process")

    monkeypatch.setattr("qc.llm.ask_json", _never)
    monkeypatch.setattr("qc.render.export_decks_png", _never)

    layouts = _layouts()
    deck, plans = _plans(_two_column_deck(3), layouts, rule="fallback")
    offered = choices(deck, layouts, plans)
    apply_picks(plans, {c.slide_index: c.suggested for c in offered}, layouts)
    assert offered


def test_an_unreadable_slide_costs_its_own_choice_and_nothing_else():
    """One bad slide must not blank the page. It keeps whatever the planner
    gave it, which is the outcome that existed before this step."""
    layouts = _layouts()
    deck, plans = _plans(_two_column_deck(2), layouts, rule="fallback")
    plans.append(SlidePlan(slide_index=99, source_layout="Ghost",
                           source_type=None, target_layout=None,
                           match_rule="fallback"))

    offered = choices(deck, layouts, plans)
    assert [c.slide_index for c in offered] == [0, 1]


@pytest.mark.parametrize("picks", [None, {}])
def test_no_picks_at_all_leaves_every_plan_exactly_as_it_was(picks):
    """Pressing Apply without touching a radio is legitimate and means "use the
    suggestions". It must not be able to corrupt a plan on the way through."""
    layouts = _layouts()
    _deck, plans = _plans(_two_column_deck(2), layouts, rule="fallback")
    before = [(p.target_layout, p.match_rule) for p in plans]

    assert apply_picks(plans, picks, layouts) == 0
    assert [(p.target_layout, p.match_rule) for p in plans] == before
