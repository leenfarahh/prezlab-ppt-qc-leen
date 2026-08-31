"""What the master is missing, drawn so a designer can build it.

qc.layoutgap says eleven slides had nowhere to go and clusters them. This is the
other half: the layout to ADD. What the tests hold is the division of labour that
makes a proposal usable rather than plausible.

  - the model names the STRUCTURE and never a coordinate; the boxes are placed on
    the master's own stated frame by code, so a wireframe lands on the client's
    margins instead of on invented numbers;
  - a proposal that does not answer the group it was asked about is DISCARDED,
    because a one-column layout offered for a two-column gap reads as an answer
    and is not one;
  - an archetype token or a box kind outside the closed sets is dropped, since
    the token is what the format pass matches on;
  - a name the master already has is reported as a COLLISION rather than offered
    as new, because that is a different finding: the layout exists and the slides
    did not reach it.

And nothing is built. A client's master is not a file this tool edits.
"""

import io

import pytest
from pptx import Presentation
from pptx.util import Emu, Inches

import qc.layoutsuggest as LS
from qc.applymaster import plan_assignments
from qc.layoutgap import report
from qc.stylespec import dominant_master, extract_layouts

IN = 914400


# ------------------------------------------------------------------ fixtures


def _two_column_deck(slides=4) -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    for n in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        head = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11),
                                        Inches(0.9))
        head.text_frame.text = f"Comparison {n + 1}"
        for col in (0.8, 7.0):
            box = slide.shapes.add_textbox(Inches(col), Inches(2), Inches(5),
                                           Inches(3))
            box.text_frame.text = "Point one"
    buf = io.BytesIO()
    prs.save(buf)
    return Presentation(io.BytesIO(buf.getvalue()))


def _lean_master() -> list:
    master = Presentation()
    return [l for l in extract_layouts(dominant_master(master),
                                       embed_assets=False)
            if l["type"] in ("titleOnly", "secHead")]


def _coverage(prs=None, layouts=None):
    prs = prs or _two_column_deck()
    layouts = layouts if layouts is not None else _lean_master()
    return report(prs, layouts, plan_assignments(prs, layouts)), prs, layouts


_GOOD = {
    "name": "Two-column comparison",
    "archetype": "twoObj",
    "columns": 2,
    "why": "For slides that set two propositions against each other.",
    "boxes": [
        {"kind": "title", "column": 0, "label": "Heading"},
        {"kind": "body", "column": 1, "label": "Left column"},
        {"kind": "body", "column": 2, "label": "Right column"},
    ],
}


@pytest.fixture()
def stub(monkeypatch):
    seen = []

    def _answer(reply):
        def _ask(**kwargs):
            seen.append(kwargs)
            return reply

        monkeypatch.setattr(LS, "ask_json", _ask)
        return seen

    return _answer


# ------------------------------------------------------------ the proposal


def test_a_gap_gets_one_proposal_that_answers_it(stub):
    seen = stub(_GOOD)
    cov, prs, layouts = _coverage()
    out, asked, _unreach = LS.suggest(cov, prs, layouts)

    assert asked == len(cov.gaps) and len(out) == 1
    s = out[0]
    assert s.name == "Two-column comparison" and s.archetype == "twoObj"
    assert s.places == 4 and s.serves == [0, 1, 2, 3]
    assert s.gap_label == cov.gaps[0].label
    # the group it was asked about travels in the prompt, so the answer is about
    # this deck rather than about two-column layouts in general
    assert "columns" in seen[0]["prompt"] and "2 columns" in seen[0]["prompt"]


def test_a_matched_deck_is_never_asked(stub):
    """A slide the master already places is not a question."""
    seen = stub(_GOOD)
    prs = _two_column_deck(2)
    layouts = extract_layouts(dominant_master(Presentation()),
                              embed_assets=False)
    cov, _prs, _l = _coverage(prs, layouts)
    assert not cov.gaps
    out, asked, _unreach = LS.suggest(cov, prs, layouts)
    assert (out, asked, seen) == ([], 0, [])


# ------------------------------------------------- code owns the geometry


def test_the_boxes_land_on_the_masters_own_frame(stub):
    """The model never returns a coordinate. A proposal placed on invented
    numbers would show a designer a layout at the wrong margins, which is worse
    than showing them a list of placeholder types."""
    stub(_GOOD)
    cov, prs, layouts = _coverage()
    frame = LS.frame_of(prs, None, {"box_emu": (IN, IN, 12 * IN, 6 * IN)})
    out, _asked, _unreach = LS.suggest(cov, prs, layouts,
                             {"box_emu": (IN, IN, 12 * IN, 6 * IN)})
    s = out[0]

    left, top, right, bottom = frame
    for box in s.boxes:
        bl, bt, br, bb = box["box"]
        assert left <= bl < br <= right, box
        assert top <= bt < bb <= bottom, box

    title = next(b for b in s.boxes if b["kind"] == "title")
    assert title["box"][0] == left and title["box"][2] == right, \
        "a title spans the frame"
    columns = sorted(b["box"][0] for b in s.boxes if b["column"])
    assert columns[0] < columns[1], "and the columns sit side by side"


def test_a_master_with_no_stated_frame_still_gets_a_wireframe(stub):
    """Not a guess at the brand's margins - somewhere honest to draw."""
    stub(_GOOD)
    cov, prs, layouts = _coverage()
    out, _asked, _unreach = LS.suggest(cov, prs, layouts, None)
    assert out and all(b["box"][2] <= int(prs.slide_width) for b in out[0].boxes)


def test_the_wireframe_is_drawn_at_the_slides_own_shape(stub):
    stub(_GOOD)
    cov, prs, layouts = _coverage()
    out, _asked, _unreach = LS.suggest(cov, prs, layouts)
    svg = LS.wireframe(out[0], prs)
    assert svg.startswith("<svg viewBox=\"0 0 12192000 6858000\"") or \
        "viewBox=\"0 0 " in svg
    assert svg.count("<rect") == 3 and "Left column" in svg
    assert "<script" not in svg


# --------------------------------------------------------------- validation


def test_a_proposal_that_does_not_answer_the_gap_is_discarded(stub):
    """One column offered for a two-column gap reads as an answer and is not
    one."""
    stub({**_GOOD, "columns": 1,
          "boxes": [{"kind": "title", "column": 0},
                    {"kind": "body", "column": 1}]})
    cov, prs, layouts = _coverage()
    out, asked, _unreach = LS.suggest(cov, prs, layouts)
    assert asked == 1 and out == []


def test_a_proposal_with_no_title_for_a_titled_gap_is_discarded(stub):
    stub({**_GOOD,
          "boxes": [{"kind": "body", "column": 1},
                    {"kind": "body", "column": 2}]})
    cov, prs, layouts = _coverage()
    out, _asked, _unreach = LS.suggest(cov, prs, layouts)
    assert out == []


def test_an_invented_archetype_is_refused(stub):
    """The token is what the format pass matches on, so an invented one produces
    a layout that never gets used."""
    stub({**_GOOD, "archetype": "twoColumnComparison"})
    cov, prs, layouts = _coverage()
    out, _asked, _unreach = LS.suggest(cov, prs, layouts)
    assert out == []


def test_an_invented_box_kind_is_dropped(stub):
    """PowerPoint has a menu of placeholder types and this is not on it. The
    box goes rather than the proposal, when what is left still serves."""
    stub({**_GOOD, "boxes": _GOOD["boxes"] + [{"kind": "sidebar", "column": 2}]})
    cov, prs, layouts = _coverage()
    out, _asked, _unreach = LS.suggest(cov, prs, layouts)
    assert out and all(b["kind"] in LS.KINDS for b in out[0].boxes)


def test_a_name_the_master_already_has_is_a_collision_not_a_proposal(stub):
    """A different finding: the layout exists and the slides did not reach it,
    which is a naming problem. A designer told to add a layout the master
    already has stops trusting the page."""
    stub({**_GOOD, "name": "Title Only"})
    cov, prs, layouts = _coverage()
    out, _asked, _unreach = LS.suggest(cov, prs, layouts)
    assert out and out[0].collides_with == "Title Only"


def test_an_unanswerable_gap_simply_gets_no_proposal(monkeypatch):
    """The gap still stands in the report, which is the state before this pass
    ran. A failure here must not cost the coverage."""
    def _boom(**kwargs):
        raise RuntimeError("unreachable")

    monkeypatch.setattr(LS, "ask_json", _boom)
    cov, prs, layouts = _coverage()
    out, asked, _unreach = LS.suggest(cov, prs, layouts)
    assert out == [] and asked == 1
    assert cov.gaps, "the diagnosis is untouched"


def test_the_schema_offers_no_place_to_put_geometry():
    props = set(LS.SPEC_SCHEMA["properties"])
    assert props == {"name", "archetype", "columns", "boxes", "why"}
    box = LS.SPEC_SCHEMA["properties"]["boxes"]["items"]
    assert set(box["properties"]) == {"kind", "column", "label"}
    assert box["additionalProperties"] is False
    assert LS.SPEC_SCHEMA["additionalProperties"] is False


def test_nothing_here_writes_to_the_master():
    """A client's master is not a file this tool edits, and a layout carries type
    styles, guides and brand furniture that are a designer's to add."""
    import inspect

    source = inspect.getsource(LS)
    for forbidden in (".save(", "add_slide", "SaveAs", "write_bytes"):
        assert forbidden not in source, forbidden


# --------------------------------------------------------------- the window


def test_the_window_shows_the_proposal_beside_its_evidence(stub):
    from qc.ui_check import render_coverage

    stub(_GOOD)
    cov, prs, layouts = _coverage()
    out, _asked, _unreach = LS.suggest(cov, prs, layouts)
    html = render_coverage(cov, standalone=True, suggestions=out,
                           pictures={0: LS.wireframe(out[0], prs)})

    assert "What to add to the master" in html
    assert "Two-column comparison" in html and "twoObj" in html
    assert "<svg" in html and "Left column" in html
    assert "does not edit a client's master" in html, (
        "the page has to say it proposes rather than builds")
    # the gap it answers is still on the page: the proposal is not a replacement
    # for the evidence
    assert cov.gaps[0].label in html
    assert html.index(cov.gaps[0].label) < html.index("What to add to the master")


def test_the_window_says_why_it_is_empty_rather_than_vanishing(stub):
    from qc.ui_check import render_coverage

    cov, _prs, _layouts = _coverage()
    html = render_coverage(cov, standalone=True, suggestions=[],
                           suggest_note="Proposing a layout needs a model, and "
                                        "none is configured here.")
    assert "Suggested layouts" in html and "needs a model" in html


# ------------------------------------ an outage is not a rejected proposal
#
# Both used to produce the same sentence: "a proposal that did not answer the
# group it was asked about is discarded rather than shown." That asserts the
# model ANSWERED. Under a 429 it never answered, and telling a designer their
# proposal was rejected on quality when the truth is an exhausted quota sends
# them to look at the gaps for a fault that is not there (30/08/2026).


def test_an_unreachable_model_is_reported_as_one(monkeypatch):
    from qc.llm import LLMUnavailable

    def _down(**kwargs):
        raise LLMUnavailable("the model's rate limit was reached")

    monkeypatch.setattr(LS, "ask_json", _down)
    cov, prs, layouts = _coverage()
    out, _asked, unreachable = LS.suggest(cov, prs, layouts)

    assert out == []
    assert unreachable, (
        "the caller cannot tell an outage from a rejection without this, and "
        "the two need opposite sentences")
    assert "rate limit" in unreachable


def test_asking_stops_once_the_provider_is_refusing(monkeypatch):
    """Five more calls into a refusing provider buy nothing and cost quota."""
    calls = []

    def _down(**kwargs):
        calls.append(1)
        raise RuntimeError("429")

    monkeypatch.setattr(LS, "ask_json", _down)
    cov, prs, layouts = _coverage()
    LS.suggest(cov, prs, layouts)
    assert len(calls) == 1, f"kept asking after a refusal ({len(calls)} calls)"


def test_a_rejected_proposal_is_still_reported_as_a_rejection(stub):
    """The other half: when the model DID answer and the answer does not serve
    the gap, that is a real result and must keep saying so."""
    stub({"name": "One big box", "archetype": "obj", "columns": 1,
          "boxes": [{"kind": "body", "column": 1}], "why": "one column"})
    cov, prs, layouts = _coverage()
    out, asked, unreachable = LS.suggest(cov, prs, layouts)

    assert out == [], "a one-column layout does not serve a two-column gap"
    assert asked, "it was asked"
    assert not unreachable, (
        "the model answered; this is a discarded proposal, not an outage")
