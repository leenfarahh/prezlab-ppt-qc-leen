"""Tests for qc.modules.margin_alignment (audit-only, top-level shapes)."""

from tests.conftest import save_and_ctx

from qc.modules.margin_alignment import detect

# Profile geometry (prezlab_en): margins L/R 457200, T 274638, B 365125;
# edge and spacing tolerance 9525 EMU.

BLANK_LAYOUT = 6
IN = 914400  # one inch in EMU


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])


def _box(slide, left, top, width=1000000, height=500000, text=""):
    tb = slide.shapes.add_textbox(left, top, width, height)
    if text:
        tb.text_frame.text = text
    return tb


def test_outside_safe_zone_left_breach(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    # safe zones govern text content, so the breaching shape carries text
    bad = _box(slide, left=0, top=IN, text="breaching text")
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = detect(ctx)
    assert len(records) == 1
    rec = records[0]
    assert rec.issue_type == "margin_alignment.outside_safe_zone"
    assert rec.severity == "warning"
    assert rec.confidence == "deterministic"
    assert rec.action == "flagged"
    assert rec.slide_index == 0
    assert rec.shape_id == str(bad.shape_id)
    assert rec.property == "spPr.xfrm.off"
    assert "left" in rec.message
    assert rec.arabic_flag is False
    assert rec.profile_rule_id == "geometry.safe_zone_margins_emu"


def test_rotated_breaching_shape_is_skipped(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    bad = _box(slide, left=0, top=IN, text="rotated text")
    bad.rotation = 45
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = detect(ctx)
    assert len(records) == 1
    rec = records[0]
    assert rec.issue_type == "margin_alignment.outside_safe_zone"
    assert rec.action == "skipped"
    assert rec.confidence == "low"
    assert rec.severity == "warning"
    assert "rotated shape, stored bounding box unreliable" in rec.message


def test_edge_misaligned_flags_only_the_deviant(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    # Vertical stack: two at X, one 20000 EMU off. Within the 3x tolerance
    # clustering window (28575) but past the 9525 tolerance from the median.
    _box(slide, left=IN, top=IN)
    _box(slide, left=IN, top=2 * IN)
    deviant = _box(slide, left=IN + 20000, top=3 * IN)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = detect(ctx)
    assert len(records) == 1
    rec = records[0]
    assert rec.issue_type == "margin_alignment.edge_misaligned"
    assert rec.severity == "warning"
    assert rec.confidence == "medium"
    assert rec.action == "flagged"
    assert rec.shape_id == str(deviant.shape_id)
    assert rec.profile_rule_id == "geometry.alignment.edge_tolerance_emu"


def test_uneven_spacing_names_the_odd_gap(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    w = 1000000
    # Row of three, gaps 500000 then 560000 (odd by 60000 > tolerance 9525).
    _box(slide, left=IN, top=IN, width=w)
    _box(slide, left=IN + w + 500000, top=IN, width=w)
    third = _box(slide, left=IN + 2 * w + 500000 + 560000, top=IN, width=w)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = detect(ctx)
    assert len(records) == 1
    rec = records[0]
    # same-size lines are claimed by the stronger cluster_rhythm check
    assert rec.issue_type == "margin_alignment.cluster_rhythm"
    assert rec.severity == "warning"
    assert rec.confidence == "medium"
    assert rec.action == "flagged"
    assert rec.shape_id == str(third.shape_id)
    assert "560000" in rec.message
    assert rec.profile_rule_id == "geometry.alignment.spacing_tolerance_emu"


def test_arabic_shape_record_is_flagged_for_manual_review(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    _box(slide, left=0, top=IN, text="مرحبا")  # Arabic
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = detect(ctx)
    assert len(records) == 1
    rec = records[0]
    assert rec.issue_type == "margin_alignment.outside_safe_zone"
    assert rec.arabic_flag is True
    assert "Arabic content, manual review" in rec.message


def test_clean_control_deck_yields_zero_records(make_prs, en_profile, tmp_path):
    prs = make_prs()
    # Slide 1: perfectly aligned vertical stack inside the safe zone.
    s1 = _blank_slide(prs)
    _box(s1, left=IN, top=IN)
    _box(s1, left=IN, top=2 * IN)
    _box(s1, left=IN, top=3 * IN)
    # Slide 2: row of three with identical gaps.
    s2 = _blank_slide(prs)
    w = 1000000
    _box(s2, left=IN, top=IN, width=w)
    _box(s2, left=IN + w + 500000, top=IN, width=w)
    _box(s2, left=IN + 2 * (w + 500000), top=IN, width=w)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    assert detect(ctx) == []


def test_full_bleed_background_exempt(make_prs, en_profile, tmp_path):
    """Real-deck tuning: a background covering most of the slide crosses
    every margin by design and must not flag."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    from qc.modules import margin_alignment as mod
    from tests.conftest import save_and_ctx

    prs = make_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0),
                       prs.slide_width, prs.slide_height)  # full-bleed bg
    s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(2_000_000),
                       prs.slide_width, Emu(300_000))  # full-width band
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert [r for r in mod.detect(ctx)
            if r.issue_type == "margin_alignment.outside_safe_zone"] == []


def test_small_breaching_shape_still_flagged(make_prs, en_profile, tmp_path):
    """Safe zones govern text: a small text shape at the edge flags; the
    same shape without text is decorative and exempt (ground-truth rule)."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu, Inches

    from qc.modules import margin_alignment as mod
    from tests.conftest import save_and_ctx

    prs = make_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Inches(2),
                             Inches(2), Inches(1))  # small shape at left edge
    shp.text_frame.text = "edge text"
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    hits = [r for r in mod.detect(ctx)
            if r.issue_type == "margin_alignment.outside_safe_zone"]
    assert len(hits) == 1 and "left" in hits[0].message


def test_textless_graphic_exempt_from_safe_zone(make_prs, en_profile, tmp_path):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu, Inches

    from qc.modules import margin_alignment as mod
    from tests.conftest import save_and_ctx

    prs = make_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Inches(2),
                       Inches(2), Inches(1))  # decorative, no text
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert [r for r in mod.detect(ctx)
            if r.issue_type == "margin_alignment.outside_safe_zone"] == []


# ---- real-deck feedback round (consulting deck, 14/07/2026) ----------------


def test_top_edge_misaligned_in_a_row(make_prs, en_profile, tmp_path):
    """Cards in a row whose tops drift (the screenshots' slide 12 case)."""
    prs = make_prs()
    slide = _blank_slide(prs)
    _box(slide, left=IN, top=IN)
    _box(slide, left=3 * IN, top=IN)
    deviant = _box(slide, left=5 * IN, top=IN + 20000)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = detect(ctx)
    assert len(records) == 1
    rec = records[0]
    assert rec.issue_type == "margin_alignment.edge_misaligned"
    assert rec.property == "spPr.xfrm.off.y"
    assert rec.shape_id == str(deviant.shape_id)
    assert "top edge" in rec.message


def test_visible_misalignment_beyond_old_window_is_caught(make_prs, en_profile, tmp_path):
    """A 1mm drift (36000 EMU) was invisible to the old 3x-tolerance window;
    the intent window must catch it. With only 3 shapes the evidence is
    thin, so it stays a warning."""
    prs = make_prs()
    slide = _blank_slide(prs)
    _box(slide, left=IN, top=IN)
    _box(slide, left=IN, top=2 * IN)
    deviant = _box(slide, left=IN + 36000, top=3 * IN)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = [r for r in detect(ctx)
               if r.issue_type == "margin_alignment.edge_misaligned"]
    assert len(records) == 1
    assert records[0].shape_id == str(deviant.shape_id)
    assert records[0].severity == "warning"


def test_obvious_misalignment_escalates_to_error(make_prs, en_profile, tmp_path):
    """Evidence-based severity: one shape visibly (>=1mm) off a tight
    line-up of several others is confidently wrong."""
    prs = make_prs()
    slide = _blank_slide(prs)
    for i in range(5):
        _box(slide, left=IN, top=IN + i * 700000)
    deviant = _box(slide, left=IN + 50000, top=IN + 5 * 700000)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = [r for r in detect(ctx)
               if r.issue_type == "margin_alignment.edge_misaligned"]
    assert len(records) == 1
    rec = records[0]
    assert rec.shape_id == str(deviant.shape_id)
    assert rec.severity == "error"
    assert rec.confidence == "high"


def test_group_interior_misalignment_detected(make_prs, en_profile, tmp_path):
    """Consulting decks group everything; alignment must see inside."""
    prs = make_prs()
    slide = _blank_slide(prs)
    group = slide.shapes.add_group_shape()
    tb1 = group.shapes.add_textbox(IN, IN, 1000000, 500000)
    tb2 = group.shapes.add_textbox(IN, 2 * IN, 1000000, 500000)
    deviant = group.shapes.add_textbox(IN + 20000, 3 * IN, 1000000, 500000)
    for tb in (tb1, tb2, deviant):
        tb.text_frame.text = "card"
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = [r for r in detect(ctx)
               if r.issue_type == "margin_alignment.edge_misaligned"]
    assert len(records) == 1
    assert records[0].shape_id == str(deviant.shape_id)
    assert "inside group" in records[0].message


def test_squeezed_text_box_flagged(make_prs, en_profile, tmp_path):
    """A text box crushed into a strip (the vertical-letters break)."""
    prs = make_prs()
    slide = _blank_slide(prs)
    bad = _box(slide, left=IN, top=IN, width=300000, height=2000000,
               text="Our real estate practice")
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = detect(ctx)
    assert len(records) == 1
    rec = records[0]
    assert rec.issue_type == "margin_alignment.squeezed_text"
    assert rec.shape_id == str(bad.shape_id)
    assert rec.severity == "error"  # a letter-per-line strip is broken, period
    assert rec.new_value is None  # flag-only: widening is a design decision


def test_overlapping_text_boxes_flagged(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    _box(slide, left=IN, top=IN, width=2000000, height=1000000, text="under")
    over = _box(slide, left=IN + 200000, top=IN + 100000,
                width=2000000, height=1000000, text="over")
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = [r for r in detect(ctx)
               if r.issue_type == "margin_alignment.text_overlap"]
    assert len(records) == 1
    assert records[0].shape_id == str(over.shape_id)
    assert records[0].severity == "error"  # survivors of the exemptions
    assert records[0].new_value is None


def test_multiple_odd_gaps_offer_distribute_fix(make_prs, en_profile, tmp_path):
    """Several odd gaps: offered as an even distribution (tickable, never
    pre-selected), applied with first and last shapes anchored."""
    import io

    from pptx import Presentation

    from qc.engine import run_audit
    from qc.fixer import apply_fixes, is_fixable

    prs = make_prs()
    slide = _blank_slide(prs)
    w = 1000000
    x = IN
    boxes = []
    for gap in (0, 500000, 560000, 620000):  # two gaps off the median
        x += gap + (w if gap else 0)
        boxes.append(_box(slide, left=x, top=IN, width=w))
    path = tmp_path / "ragged.pptx"
    prs.save(path)

    result = run_audit(path, "prezlab_en", modules=["margin_alignment"])
    records = [r.to_dict() for r in result.records]
    spacing = [r for r in records
               if r["issue_type"] == "margin_alignment.uneven_spacing"]
    assert len(spacing) == 1
    rec = spacing[0]
    assert rec["locator"].startswith("dist-row:")
    assert is_fixable(rec)
    assert "distributes them evenly" in rec["message"]

    fix = apply_fixes(path.read_bytes(), records, {rec["record_id"]})
    assert [o.outcome for o in fix.outcomes] == ["changed"], \
        [(o.record_id, o.reason) for o in fix.outcomes]
    cleaned = Presentation(io.BytesIO(fix.cleaned_bytes))
    by_id = {str(sh.shape_id): sh for sh in cleaned.slides[0].shapes}
    lefts = sorted(by_id[str(b.shape_id)].left for b in boxes)
    assert lefts[0] == boxes[0].left            # first anchored
    assert lefts[-1] == boxes[-1].left          # last anchored
    gaps = [lefts[i + 1] - (lefts[i] + w) for i in range(3)]
    assert max(gaps) - min(gaps) <= 2           # equal within rounding


def test_top_snap_fix_applies(make_prs, en_profile, tmp_path):
    import io

    from pptx import Presentation

    from qc.engine import run_audit
    from qc.fixer import apply_fixes

    prs = make_prs()
    slide = _blank_slide(prs)
    _box(slide, left=IN, top=IN)
    _box(slide, left=3 * IN, top=IN)
    deviant = _box(slide, left=5 * IN, top=IN + 20000)
    path = tmp_path / "row.pptx"
    prs.save(path)

    result = run_audit(path, "prezlab_en", modules=["margin_alignment"])
    records = [r.to_dict() for r in result.records]
    ids = {r["record_id"] for r in records
           if r["issue_type"] == "margin_alignment.edge_misaligned"}
    assert ids
    fix = apply_fixes(path.read_bytes(), records, ids)
    assert [o.outcome for o in fix.outcomes] == ["changed"]
    cleaned = Presentation(io.BytesIO(fix.cleaned_bytes))
    by_id = {str(sh.shape_id): sh for sh in cleaned.slides[0].shapes}
    assert by_id[str(deviant.shape_id)].top == IN  # snapped to the median


def test_vertical_spacing_fix_shifts_column_tail(make_prs, en_profile, tmp_path):
    import io

    from pptx import Presentation

    from qc.engine import run_audit
    from qc.fixer import apply_fixes

    prs = make_prs()
    slide = _blank_slide(prs)
    h = 500000
    _box(slide, left=IN, top=IN, height=h)
    _box(slide, left=IN, top=IN + h + 500000, height=h)
    third = _box(slide, left=IN, top=IN + 2 * h + 500000 + 560000, height=h)
    path = tmp_path / "stack.pptx"
    prs.save(path)

    result = run_audit(path, "prezlab_en", modules=["margin_alignment"])
    records = [r.to_dict() for r in result.records]
    spacing = [r for r in records
               if r["issue_type"] == "margin_alignment.uneven_spacing"]
    assert len(spacing) == 1
    assert spacing[0]["locator"].startswith("col:")
    assert "vertical" in spacing[0]["message"]

    fix = apply_fixes(path.read_bytes(), records,
                      {spacing[0]["record_id"]})
    assert [o.outcome for o in fix.outcomes] == ["changed"]
    cleaned = Presentation(io.BytesIO(fix.cleaned_bytes))
    by_id = {str(sh.shape_id): sh for sh in cleaned.slides[0].shapes}
    # odd gap 560000 closed to the median 500000: tail moved up 60000
    assert by_id[str(third.shape_id)].top == IN + 2 * h + 2 * 500000


def test_panel_snap_carries_its_contents(make_prs, en_profile, tmp_path):
    """Moving a container must move what's visually inside it: a panel
    snapped into line takes its labels/photos along (real-deck finding:
    a photo-grid panel aligned to its neighbors, orphaning the photos)."""
    import io

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    from qc.engine import run_audit
    from qc.fixer import apply_fixes

    prs = make_prs()
    slide = _blank_slide(prs)
    pw, ph = 2000000, 2500000
    # four tightly aligned panels + one visibly off (error escalation)
    for i in range(4):
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(IN),
                               Emu(IN + i * (ph + 200000)), Emu(pw), Emu(ph))
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(IN + 60000),
        Emu(IN + 4 * (ph + 200000)), Emu(pw), Emu(ph))
    # contents inside the deviant panel
    inner1 = _box(slide, left=panel.left + 200000, top=panel.top + 200000,
                  width=600000, height=300000, text="label")
    inner2 = _box(slide, left=panel.left + 200000, top=panel.top + 700000,
                  width=600000, height=300000, text="photo")
    path = tmp_path / "panel.pptx"
    prs.save(path)

    result = run_audit(path, "prezlab_en", modules=["margin_alignment"])
    records = [r.to_dict() for r in result.records]
    edge = [r for r in records
            if r["issue_type"] == "margin_alignment.edge_misaligned"
            and r["shape_id"] == str(panel.shape_id)]
    assert len(edge) == 1 and edge[0]["severity"] == "error"

    fix = apply_fixes(path.read_bytes(), records, {edge[0]["record_id"]})
    assert [o.outcome for o in fix.outcomes] == ["changed"], \
        [(o.record_id, o.reason) for o in fix.outcomes]

    cleaned = Presentation(io.BytesIO(fix.cleaned_bytes))
    by_id = {str(sh.shape_id): sh for sh in cleaned.slides[0].shapes}
    moved_panel = by_id[str(panel.shape_id)]
    assert moved_panel.left == IN  # snapped to the cluster median
    delta = IN - (IN + 60000)
    # contents travelled with the panel, composition intact
    assert by_id[str(inner1.shape_id)].left == inner1.left + delta
    assert by_id[str(inner2.shape_id)].left == inner2.left + delta


def test_contained_icons_never_join_global_clusters(make_prs, en_profile, tmp_path):
    """Slide-12 regression: icons composed identically with their cards
    must not be clustered against the card tabs above them and snapped out
    of their composition."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    prs = make_prs()
    slide = _blank_slide(prs)
    # four aligned tabs; the last carries an icon slightly lower than the
    # tab tops (sitting inside its tab, as designed)
    for i in range(4):
        tab = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(IN + i * 2000000), Emu(IN),
            Emu(1200000), Emu(400000))
    icon = _box(slide, left=IN + 3 * 2000000 + 400000, top=IN + 80000,
                width=300000, height=250000)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = [r for r in detect(ctx)
               if r.issue_type == "margin_alignment.edge_misaligned"]
    assert records == []  # the icon belongs to its tab, not the tab cluster


def test_relative_misalignment_recomposed_with_own_card(make_prs, en_profile, tmp_path):
    """An icon sitting differently INSIDE its card than its peers sit in
    theirs is flagged, and the snap target is its own card's inset."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    prs = make_prs()
    slide = _blank_slide(prs)
    cards, icons = [], []
    for i in range(4):
        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(IN + i * 2200000), Emu(IN),
            Emu(1800000), Emu(1400000))
        cards.append(card)
        inset = 300000 if i == 3 else 200000  # the last icon is off-inset
        icons.append(_box(slide, left=card.left + inset,
                          top=card.top + 200000,
                          width=400000, height=300000))
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = [r for r in detect(ctx)
               if r.issue_type == "margin_alignment.edge_misaligned"]
    assert len(records) == 1
    rec = records[0]
    assert rec.shape_id == str(icons[3].shape_id)
    assert rec.severity == "error"  # 4-cluster, visibly off the shared inset
    assert int(rec.new_value) == cards[3].left + 200000  # its OWN card's inset
    assert "its own" in rec.message


def test_mosaic_fragments_exempt_from_clusters(make_prs, en_profile, tmp_path):
    """>=5 touching shapes compose ONE visual object (a wheel, a podium);
    its fragments must not be clustered against outside edges."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    prs = make_prs()
    slide = _blank_slide(prs)
    # a stepped 'podium' of six abutting blocks whose lefts stagger within
    # the intent window of three aligned free boxes
    for i in range(6):
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Emu(IN + i * 20000), Emu(IN + i * 300000),
                               Emu(400000), Emu(300000))
    _box(slide, left=5 * IN, top=IN)
    _box(slide, left=5 * IN, top=2 * IN)
    _box(slide, left=5 * IN, top=3 * IN)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = [r for r in detect(ctx)
               if r.issue_type == "margin_alignment.edge_misaligned"]
    assert records == []  # neither podium steps nor cross-contamination


def test_wall_spacing_is_advisory_only(make_prs, en_profile, tmp_path):
    """A photo/logo wall row (>=6 same-size items) with one odd gap gets an
    honest flag, never a computed nudge."""
    from qc.fixer import is_fixable

    prs = make_prs()
    slide = _blank_slide(prs)
    w, h = 700000, 700000
    x = IN
    gaps = [0, 300000, 300000, 300000, 300000, 360000, 300000]
    for i, gap in enumerate(gaps):
        x += gap + (w if i else 0)
        _box(slide, left=x, top=IN, width=w, height=h)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    records = [r for r in detect(ctx)
               if r.issue_type == "margin_alignment.uneven_spacing"]
    assert len(records) == 1  # exactly one, no double emit
    rec = records[0]
    assert rec.new_value is None
    assert "adjust by eye" in rec.message
    assert not is_fixable(rec.to_dict())


def test_spacing_fix_reverts_if_it_would_create_a_collision(make_prs, en_profile, tmp_path):
    """A fix must never make the slide worse: closing an odd gap that walks
    the shape into a bystander is refused and reverted (real-deck finding:
    a spacing snap pushed a logo into its neighbor)."""
    import io

    from pptx import Presentation

    from qc.engine import run_audit
    from qc.fixer import apply_fixes

    prs = make_prs()
    slide = _blank_slide(prs)
    w, h = 1000000, 500000
    _box(slide, left=IN, top=IN, width=w, height=h)
    _box(slide, left=IN + w + 500000, top=IN, width=w, height=h)
    # odd gap of 800000 before the third box (median 500000)
    third = _box(slide, left=IN + 2 * w + 500000 + 800000, top=IN,
                 width=w, height=h)
    # bystander sitting exactly where the third box would land, offset
    # vertically so it joins neither the row nor any alignment cluster
    _box(slide, left=third.left - 400000, top=IN + 300000,
         width=300000, height=h)

    path = tmp_path / "collision.pptx"
    prs.save(path)
    result = run_audit(path, "prezlab_en", modules=["margin_alignment"])
    records = [r.to_dict() for r in result.records]
    spacing = [r for r in records
               if r["issue_type"] in ("margin_alignment.uneven_spacing",
                                      "margin_alignment.cluster_rhythm")
               and r["new_value"] is not None]
    assert len(spacing) == 1

    fix = apply_fixes(path.read_bytes(), records, {spacing[0]["record_id"]})
    assert fix.outcomes[0].outcome == "skipped"
    assert "into shape" in fix.outcomes[0].reason

    cleaned = Presentation(io.BytesIO(fix.cleaned_bytes))
    by_id = {str(sh.shape_id): sh for sh in cleaned.slides[0].shapes}
    assert by_id[str(third.shape_id)].left == third.left  # untouched


# --------------------------------------- the text inside an aligned box
#
# The one alignment defect the box-edge rules cannot see: boxes that share a
# top edge to the EMU, holding text that sits at different heights inside them
# because their vertical anchors disagree. The anchor is stated in the file, so
# this stays a fact about the deck rather than a rendering estimate.


def _anchored(slide, left, top, anchor, text="Label text",
              width=1000000, height=500000, autofit=False):
    """A fixed-height text box, which is the only kind whose anchor draws
    anything. python-pptx (like PowerPoint's own Insert > Text Box) writes
    a:spAutoFit, so the box hugs its copy and sits wherever the anchor puts
    nothing; `autofit=False` is the designer's "do not autofit"."""
    from pptx.oxml.ns import qn

    tb = _box(slide, left, top, width=width, height=height, text=text)
    bodyPr = tb.text_frame._txBody.find(qn("a:bodyPr"))
    if not autofit:
        for el in bodyPr.findall(qn("a:spAutoFit")):
            bodyPr.remove(el)
    if anchor:
        bodyPr.set("anchor", anchor)
    return tb


def _anchor_records(records):
    return [r for r in records
            if r.issue_type == "margin_alignment.text_anchor_mismatch"]


def test_one_box_in_an_aligned_row_anchors_its_text_differently(
        make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    _anchored(slide, IN, IN, "t")
    _anchored(slide, 3 * IN, IN, "t")
    odd = _anchored(slide, 5 * IN, IN, "b")
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    found = _anchor_records(detect(ctx))
    assert len(found) == 1
    rec = found[0]
    assert rec.shape_id == str(odd.shape_id)
    assert rec.property == "bodyPr.anchor"
    assert rec.old_value == "b"
    assert rec.new_value == "t"
    assert rec.action == "flagged"
    assert rec.severity == "warning"
    assert "bottom" in rec.message and "top" in rec.message


def test_a_row_that_agrees_says_nothing(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    for i in range(3):
        _anchored(slide, IN + i * 2 * IN, IN, "ctr")
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    assert _anchor_records(detect(ctx)) == []


def test_an_unstated_anchor_is_the_ooxml_default_not_a_deviation(
        make_prs, en_profile, tmp_path):
    """A box that states nothing is anchored top, which is what the two that
    say so are: three agreeing boxes, no finding."""
    prs = make_prs()
    slide = _blank_slide(prs)
    _anchored(slide, IN, IN, "t")
    _anchored(slide, 3 * IN, IN, "t")
    _anchored(slide, 5 * IN, IN, None)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    assert _anchor_records(detect(ctx)) == []


def test_boxes_of_different_heights_are_not_one_row_of_labels(
        make_prs, en_profile, tmp_path):
    """Different-size boxes starting on the same line are assorted content,
    and their anchors are not comparable."""
    prs = make_prs()
    slide = _blank_slide(prs)
    _anchored(slide, IN, IN, "t")
    _anchored(slide, 3 * IN, IN, "t")
    _anchored(slide, 5 * IN, IN, "b", height=2000000)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    assert _anchor_records(detect(ctx)) == []


def test_a_row_with_no_majority_is_left_alone(make_prs, en_profile, tmp_path):
    """Three boxes, three anchors: there is no line to have deviated from."""
    prs = make_prs()
    slide = _blank_slide(prs)
    _anchored(slide, IN, IN, "t")
    _anchored(slide, 3 * IN, IN, "ctr")
    _anchored(slide, 5 * IN, IN, "b")
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    assert _anchor_records(detect(ctx)) == []


def test_a_box_that_hugs_its_text_is_excluded(make_prs, en_profile, tmp_path):
    """With a:spAutoFit the box height tracks the copy, so the anchor draws
    nothing and every such row would report a difference nobody can see."""
    prs = make_prs()
    slide = _blank_slide(prs)
    _anchored(slide, IN, IN, "t")
    _anchored(slide, 3 * IN, IN, "t")
    _anchored(slide, 5 * IN, IN, "b", autofit=True)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    assert _anchor_records(detect(ctx)) == []


def test_the_anchor_mismatch_is_never_offered_as_a_fix(make_prs, en_profile,
                                                       tmp_path):
    """A deliberately bottom-anchored caption is a real design, so the tool
    states the difference and the designer settles it."""
    from qc.fixer import FIXABLE_ISSUES

    assert "margin_alignment.text_anchor_mismatch" not in FIXABLE_ISSUES


def test_arabic_text_anchors_are_compared_the_same_way(make_prs, en_profile,
                                                       tmp_path):
    """The anchor is vertical, so it is direction-neutral: an Arabic row is
    judged by the same rule and the record is marked for review."""
    prs = make_prs()
    slide = _blank_slide(prs)
    for i, anchor in enumerate(("t", "t", "b")):
        _anchored(slide, IN + i * 2 * IN, IN, anchor, text="نص عربي هنا")
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    found = _anchor_records(detect(ctx))
    assert len(found) == 1
    assert found[0].arabic_flag is True
