"""Tests for qc.modules.master_slide (layout outliers, placeholder geometry)."""

from pptx.util import Inches

from qc.modules.master_slide import detect
from tests.conftest import save_and_ctx

GEO = "master_slide.placeholder_geometry_off"
OUTLIER = "master_slide.layout_outlier"


def _mixed_layout_deck(make_prs):
    """7 slides on layout A, 3 on layout B. Both layouts get slides with
    freshly cloned placeholders (no explicit xfrm), so only layout checks
    can fire."""
    prs = make_prs()
    layout_a = prs.slide_layouts[1]  # Title and Content
    layout_b = prs.slide_layouts[6]  # Blank
    for _ in range(7):
        prs.slides.add_slide(layout_a)
    for _ in range(3):
        prs.slides.add_slide(layout_b)
    return prs, layout_a, layout_b


def test_moved_title_placeholder_is_geometry_off_deterministic(
        make_prs, en_profile, tmp_path):
    prs = make_prs()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    prs.slides.add_slide(layout)  # untouched control slide

    title = slide.shapes.title
    assert title.left is not None  # inherited geometry must resolve
    title.left = title.left + Inches(2)
    moved_id = str(title.shape_id)

    ctx = save_and_ctx(prs, tmp_path, en_profile)
    records = detect(ctx)

    assert len(records) == 1
    rec = records[0]
    assert rec.issue_type == GEO
    assert rec.severity == "warning"
    assert rec.confidence == "deterministic"
    assert rec.action == "flagged"
    assert rec.slide_index == 0
    assert rec.shape_id == moved_id
    assert rec.property == "spPr.xfrm"
    assert rec.arabic_flag is False
    assert rec.old_value is not None and rec.new_value is not None


def test_census_flags_minority_layout_when_allowlist_empty(
        make_prs, en_profile, tmp_path):
    prs, layout_a, layout_b = _mixed_layout_deck(make_prs)
    assert en_profile.get("master_slide.layout_allowlist") == []

    ctx = save_and_ctx(prs, tmp_path, en_profile)
    records = detect(ctx)

    assert all(r.issue_type == OUTLIER for r in records)
    assert len(records) == 3
    assert sorted(r.slide_index for r in records) == [7, 8, 9]
    for rec in records:
        assert rec.severity == "warning"
        assert rec.confidence == "medium"
        assert rec.action == "flagged"
        assert rec.shape_id == "-"
        assert rec.property == "slideLayout"
        assert rec.old_value == layout_b.name
        assert rec.new_value == layout_a.name
        assert "inferred" in rec.message


def test_allowlist_flags_outliers_with_high_confidence(
        make_prs, en_profile, tmp_path):
    prs, layout_a, layout_b = _mixed_layout_deck(make_prs)
    en_profile.config["master_slide"]["layout_allowlist"] = [layout_a.name]

    ctx = save_and_ctx(prs, tmp_path, en_profile)
    records = detect(ctx)

    assert all(r.issue_type == OUTLIER for r in records)
    assert len(records) == 3
    assert sorted(r.slide_index for r in records) == [7, 8, 9]
    for rec in records:
        assert rec.severity == "warning"
        assert rec.confidence == "high"
        assert rec.shape_id == "-"
        assert rec.old_value == layout_b.name
        assert rec.profile_rule_id == "master_slide.layout_allowlist"


def test_clean_deck_yields_zero_records(make_prs, en_profile, tmp_path):
    prs = make_prs()
    layout = prs.slide_layouts[1]
    for _ in range(4):
        prs.slides.add_slide(layout)

    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert detect(ctx) == []


def test_arabic_placeholder_geometry_off_carries_guard(
        make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text_frame.text = "عرض تقديمي"
    assert title.left is not None
    title.left = title.left + Inches(2)

    ctx = save_and_ctx(prs, tmp_path, en_profile)
    records = detect(ctx)

    geo = [r for r in records if r.issue_type == GEO]
    assert len(geo) == 1
    rec = geo[0]
    assert rec.arabic_flag is True
    # geometry-only fix: since 12/08/2026 the Arabic note is transparency,
    # not a manual-review demand (the inherit snap never opens the text)
    assert "text untouched" in rec.message
    assert rec.action == "flagged"  # audit-only, never mutated
    assert rec.severity == "warning"
    assert rec.confidence == "deterministic"
