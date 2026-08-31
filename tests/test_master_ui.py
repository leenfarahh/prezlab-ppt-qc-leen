"""Step 1 in the web tier: submit a master, review the spec, save it as a
profile the audit engine can read.

All of it happens on Prepare a deck. There is no master page and no profile
editor: the master is dropped on that page, the spec renders under the form
that read it, and the profile is saved (or an existing one re-pointed at the
new file) without leaving."""

import io
import json

from fastapi.testclient import TestClient
from pptx import Presentation

from qc import web
from qc.profile import Profile
from qc.stylespec import extract_style_spec, spec_to_profile


def _master_bytes(**kw) -> bytes:
    """A master-only file: the Stage 1 submission shape, no slides."""
    prs = Presentation(**kw)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _client(tmp_path, monkeypatch):
    monkeypatch.setattr(web, "AUTH_REQUIRED", False)
    web.app.state.auth_required = False
    return TestClient(web.app)


def _as_lead(client):
    from qc.store import add_user

    add_user("Lead Person", "lead")
    client.post("/whoami", json={"name": "Lead Person"})


# ------------------------------------------------------------------- routes


def _saved_pid(response) -> str:
    """The profile id out of the redirect a save lands on."""
    from urllib.parse import parse_qs, urlparse

    parsed = urlparse(response.headers["location"])
    assert parsed.path == "/prep", parsed.path
    return parse_qs(parsed.query)["saved"][0]


def _read_master(client, blob, filename="brand.pptx"):
    """Step 1: read a master, and hand back the spec id it was held under."""
    r = client.post("/master", files={"master": (filename, blob, "app/x")})
    # The NEWEST spec. web._specs is module state shared by every test in the
    # session, so the first entry is whatever some earlier test read.
    return r, next(reversed(web._specs))


def test_the_master_step_is_on_the_prepare_page(tmp_path, monkeypatch):
    """There is no /master page to render. Reading a master is step 1 of the
    one page, and a GET must not resurrect a second door onto it."""
    client = _client(tmp_path, monkeypatch)
    assert client.get("/master").status_code == 405

    r = client.get("/prep")
    assert r.status_code == 200
    assert "Save a master as a profile" in r.text
    assert 'action="/master"' in r.text


def test_reading_a_master_shows_the_spec(tmp_path, monkeypatch):
    client = _client(tmp_path, monkeypatch)
    r = client.post("/master", files={"master": ("brand.pptx", _master_bytes(),
                                                 "application/vnd.ms-powerpoint")})
    assert r.status_code == 200
    assert "brand.pptx" in r.text
    # ...on the page that asked for it, with step 2 still below it, rather than
    # on a page of its own.
    assert "Apply it to a messy deck" in r.text
    assert 'action="/prep"' in r.text
    # The theme is shown by role, with real swatches.
    assert "accent1" in r.text
    # No slides is the normal case and the page says so rather than warning.
    assert "carries no slides" in r.text
    # Layout archetypes are surfaced, since they drive Stage 2 matching.
    assert "secHead" in r.text


def test_non_pptx_is_refused(tmp_path, monkeypatch):
    client = _client(tmp_path, monkeypatch)
    r = client.post("/master", files={"master": ("notes.txt", b"nope", "text/plain")})
    assert r.status_code == 400
    assert "Only .pptx" in r.text


def test_unreadable_file_fails_with_a_message_not_a_traceback(tmp_path, monkeypatch):
    client = _client(tmp_path, monkeypatch)
    r = client.post("/master", files={"master": ("broken.pptx", b"PK\x03\x04junk",
                                                 "application/vnd.ms-powerpoint")})
    assert r.status_code == 422
    assert "Could not read that master" in r.text


def test_spec_json_downloads_and_round_trips(tmp_path, monkeypatch):
    client = _client(tmp_path, monkeypatch)
    client.post("/master", files={"master": ("brand.pptx", _master_bytes(), "app/x")})
    # The NEWEST spec. web._specs is module state shared by every test in the
    # session, so the first entry is whatever some earlier test read.
    spec_id = next(reversed(web._specs))

    r = client.get(f"/spec/{spec_id}.json")
    assert r.status_code == 200
    assert "brand-stylespec.json" in r.headers["content-disposition"]
    spec = r.json()
    assert spec["spec_version"] == 1
    assert spec["theme"]["colors"]["accent1"]


def test_unknown_spec_id_is_a_clean_404(tmp_path, monkeypatch):
    r = _client(tmp_path, monkeypatch).get("/spec/deadbeef.json")
    assert r.status_code == 404


# ------------------------------------------------------------ saving a spec


def test_saving_a_profile_needs_a_lead(tmp_path, monkeypatch):
    client = _client(tmp_path, monkeypatch)
    client.post("/master", files={"master": ("brand.pptx", _master_bytes(), "app/x")})
    # The NEWEST spec. web._specs is module state shared by every test in the
    # session, so the first entry is whatever some earlier test read.
    spec_id = next(reversed(web._specs))

    r = client.post(f"/spec/{spec_id}/profile", data={"name": "Client X"},
                    follow_redirects=False)
    assert r.status_code == 403
    assert "lead or admin" in r.text


def test_lead_saves_a_spec_as_a_usable_profile(tmp_path, monkeypatch):
    from qc.profile import PROFILES_DIR

    client = _client(tmp_path, monkeypatch)
    _as_lead(client)
    client.post("/master", files={"master": ("brand.pptx", _master_bytes(), "app/x")})
    # The NEWEST spec. web._specs is module state shared by every test in the
    # session, so the first entry is whatever some earlier test read.
    spec_id = next(reversed(web._specs))

    r = client.post(f"/spec/{spec_id}/profile", data={"name": "Client Zed master"},
                    follow_redirects=False)
    assert r.status_code == 303
    pid = _saved_pid(r)
    written = PROFILES_DIR / f"{pid}.json"
    try:
        assert written.exists()
        # It has to load through the ordinary profile path, or the audit
        # engine cannot use it.
        profile = Profile.load(pid)
        assert profile.get("font.roles.title.latin")
        assert profile.get("color_palette.named_colors")
        assert profile.get("theme.colors.accent1")
    finally:
        written.unlink(missing_ok=True)


# --------------------------------------------- replacing a stored master
#
# A profile's master is stored once, when the profile is created, and the format
# step hands PowerPoint that stored copy. Everything the designer adds to their
# master afterwards - a presentation-space rectangle, a moved guide, a new
# layout - reached no deck at all before this existed (design lead, 21/08/2026:
# the presentation-space box was missing from every formatted deck).


def _master_with_space(left_in=1.2, top_in=0.4, w_in=10.5, h_in=6.4) -> bytes:
    from pptx.oxml.shapes.autoshape import CT_Shape

    IN = 914400
    prs = Presentation()
    prs.slide_width, prs.slide_height = 12192000, 6858000
    master = prs.slide_masters[0]
    sp = CT_Shape.new_autoshape_sp(950, "Presentation space", "rect",
                                   int(left_in * IN), int(top_in * IN),
                                   int(w_in * IN), int(h_in * IN))
    master.shapes._spTree.insert_element_before(sp, "p:extLst")
    shape = next(s for s in master.shapes if s.name == "Presentation space")
    shape.fill.background()
    shape.line.fill.background()
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _saved_profile(client) -> str:
    _, spec_id = _read_master(client, _master_bytes())
    r = client.post(f"/spec/{spec_id}/profile", data={"name": "Frame Client"},
                    follow_redirects=False)
    return _saved_pid(r)


def _replace_master(client, pid, blob, filename="brand-v2.pptx"):
    """Point an existing profile at a revised master, the way the page does:
    read it in step 1, then pick that profile as the target."""
    read, spec_id = _read_master(client, blob, filename)
    if read.status_code != 200:
        return read
    return client.post(f"/spec/{spec_id}/profile", data={"target": pid},
                       follow_redirects=False)


def test_replacing_the_master_stores_the_new_file_and_re_reads_its_frame(
        tmp_path, monkeypatch):
    from qc.profile import PROFILES_DIR
    from qc.templates import load_master

    client = _client(tmp_path, monkeypatch)
    _as_lead(client)
    pid = _saved_profile(client)
    try:
        before = Profile.load(pid).get("geometry.safe_zone_margins_emu")

        new_master = _master_with_space()
        r = _replace_master(client, pid, new_master)
        assert r.status_code == 303
        assert _saved_pid(r) == pid, "it must be the SAME profile, not a copy"

        landed = client.get(r.headers["location"])
        assert "Master replaced" in landed.text

        # the FILE is what gets applied, so that is what must have changed
        assert load_master(pid) == new_master
        after = Profile.load(pid).get("geometry.safe_zone_margins_emu")
        assert after != before
        assert after["left"] == int(1.2 * 914400)
        assert Profile.load(pid).get("style_spec_source.grid_source") == \
            "presentation_space"
    finally:
        (PROFILES_DIR / f"{pid}.json").unlink(missing_ok=True)


def test_replacing_the_master_keeps_hand_edited_rules(tmp_path, monkeypatch):
    """Fonts, palette and tolerances are decisions about the client, not
    readings of the file. Reverting them to a fresh projection would punish
    every edit ever made in the editor."""
    from qc.profile import PROFILES_DIR

    client = _client(tmp_path, monkeypatch)
    _as_lead(client)
    pid = _saved_profile(client)
    path = PROFILES_DIR / f"{pid}.json"
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
        data["config"]["geometry"]["alignment"]["edge_tolerance_emu"] = 28575
        data["config"]["font"]["roles"]["title"]["latin"] = ["Prezlab Display"]
        path.write_text(json.dumps(data, indent=2), encoding="utf-8")

        _replace_master(client, pid, _master_with_space())
        after = Profile.load(pid)
        assert after.get("geometry.alignment.edge_tolerance_emu") == 28575
        assert after.get("font.roles.title.latin") == ["Prezlab Display"]
        assert after.version > 1, "the version must move: the rules changed"
    finally:
        path.unlink(missing_ok=True)


def test_replacing_the_master_needs_a_lead(tmp_path, monkeypatch):
    from qc.profile import PROFILES_DIR
    from qc.templates import load_master

    client = _client(tmp_path, monkeypatch)
    _as_lead(client)
    pid = _saved_profile(client)
    path = PROFILES_DIR / f"{pid}.json"
    try:
        stored = load_master(pid)
        anonymous = _client(tmp_path, monkeypatch)   # its own cookie jar
        r = _replace_master(anonymous, pid, _master_with_space())
        assert r.status_code == 403
        assert "lead or admin" in r.text
        assert load_master(pid) == stored, "the stored file must be untouched"
    finally:
        path.unlink(missing_ok=True)


def test_a_broken_master_upload_is_refused_without_losing_the_old_one(
        tmp_path, monkeypatch):
    from qc.profile import PROFILES_DIR
    from qc.templates import load_master

    client = _client(tmp_path, monkeypatch)
    _as_lead(client)
    pid = _saved_profile(client)
    path = PROFILES_DIR / f"{pid}.json"
    try:
        stored = load_master(pid)
        # The junk never gets as far as the profile: step 1 refuses to read it,
        # so there is no spec to point anything at.
        r = _replace_master(client, pid, b"not a pptx", "junk.pptx")
        assert r.status_code == 422
        assert "Could not read that master" in r.text
        assert load_master(pid) == stored
    finally:
        path.unlink(missing_ok=True)


def test_expired_spec_says_so_instead_of_500(tmp_path, monkeypatch):
    client = _client(tmp_path, monkeypatch)
    _as_lead(client)
    r = client.post("/spec/gone/profile", data={"name": "X"},
                    follow_redirects=False)
    assert r.status_code == 404
    assert "expired" in r.text


# --------------------------------------------------------- the spec -> profile


def test_profile_takes_major_font_for_headings_and_minor_for_body():
    spec = extract_style_spec(Presentation(), source="m.pptx")
    theme = spec["theme"]["fonts"]
    cfg = spec_to_profile(spec, "p", "P")["config"]

    assert theme["major"]["latin"] in cfg["font"]["roles"]["title"]["latin"]
    assert theme["minor"]["latin"] in cfg["font"]["roles"]["body"]["latin"]


def test_palette_entries_carry_their_theme_slot():
    """The whole point of sourcing the palette from the theme: a later theme
    edit has to carry the palette with it."""
    spec = extract_style_spec(Presentation(), source="m.pptx")
    named = spec_to_profile(spec, "p", "P")["config"]["color_palette"]["named_colors"]

    assert named
    assert all(c["theme_ref"] == c["name"] for c in named)
    assert not any(c["name"] in ("hlink", "folHlink") for c in named)


def test_grid_columns_enable_the_profile_grid():
    spec = extract_style_spec(Presentation(), source="m.pptx")
    spec["grid"] = {"guides": {"vertical_emu": [], "horizontal_emu": []},
                    "margins_emu": {"left": 1, "right": 2, "top": 3, "bottom": 4},
                    "columns": 6, "gutter_emu": 304800, "source": "guides"}
    geo = spec_to_profile(spec, "p", "P")["config"]["geometry"]

    assert geo["safe_zone_margins_emu"]["left"] == 1
    assert geo["grid"] == {"columns": 6, "gutter_emu": 304800, "enabled": True}


def test_profile_falls_back_to_canvas_ratios_without_a_grid():
    spec = extract_style_spec(Presentation(), source="m.pptx")
    spec["grid"] = {"guides": {"vertical_emu": [], "horizontal_emu": []},
                    "margins_emu": None, "columns": None,
                    "gutter_emu": None, "source": None}
    geo = spec_to_profile(spec, "p", "P")["config"]["geometry"]

    assert geo["safe_zone_margins_emu"]["left"] > 0
    assert geo["grid"]["enabled"] is False


def test_profile_is_json_serialisable():
    spec = extract_style_spec(Presentation(), source="m.pptx")
    profile = spec_to_profile(spec, "p", "P")
    assert json.loads(json.dumps(profile)) == profile


# ------------------------------- writing the presentation space into a master
#
# Without a stated frame, the content area is inferred from where the master's
# own placeholders happen to sit, and every deck formatted against it is seated
# on that inference (qc.stylespec.infer_grid). This turns the inference into a
# statement, and it is the one edit this tool will make to a master - an
# invisible rectangle carrying a decision the designer has already made, handed
# back as a COPY to check.

IN = 914400


def _stamped(box=None):
    from qc.pspace import stamp_master
    from qc.stylespec import dominant_master, infer_grid

    data = _master_bytes()
    prs = Presentation(io.BytesIO(data))
    if box is None:
        m = infer_grid(prs, dominant_master(prs))["margins_emu"]
        box = (m["left"], m["top"],
               prs.slide_width - m["right"], prs.slide_height - m["bottom"])
    return data, stamp_master(data, box)


def test_stamping_turns_an_inferred_frame_into_a_stated_one():
    from qc.stylespec import dominant_master, infer_grid, read_presentation_space

    data, (out, note) = _stamped()
    before = infer_grid(Presentation(io.BytesIO(data)),
                        dominant_master(Presentation(io.BytesIO(data))))
    assert before["source"] != "presentation_space", "the fixture starts inferred"

    prs = Presentation(io.BytesIO(out))
    master = dominant_master(prs)
    assert infer_grid(prs, master)["source"] == "presentation_space", (
        "the whole point: the frame is now read, not inferred")
    space = read_presentation_space(prs, master)
    assert space and space["marker"] == "alt_text", (
        "written the way ToolsToo writes it, so the designer's own add-in "
        "finds it too")
    assert not space.get("prints"), (
        "a marker that prints appears on every slide of every delivered deck")
    assert "ToolsToo_PS" in note


def test_the_original_master_is_not_touched():
    """A client's master is not a file this tool edits. What comes back is a
    copy; the bytes that went in are unchanged."""
    data, (out, _note) = _stamped()
    assert out != data, "a copy was produced"
    from qc.stylespec import dominant_master, read_presentation_space

    prs = Presentation(io.BytesIO(data))
    assert read_presentation_space(prs, dominant_master(prs)) is None, (
        "the submitted bytes still state no presentation space")


def test_a_master_that_already_states_one_is_left_alone():
    _data, (out, _note) = _stamped()
    from qc.pspace import stamp_master

    again, note = stamp_master(out, (IN, IN, 2 * IN, 2 * IN))
    assert again == out, "nothing was added and nothing was rewritten"
    assert "already states a presentation space" in note


def test_a_frame_off_the_canvas_is_refused():
    """It would be read back AS the frame and seat every deck against it."""
    import pytest

    from qc.pspace import stamp_master

    data = _master_bytes()
    prs = Presentation(io.BytesIO(data))
    for box, why in (((0, 0, prs.slide_width + IN, prs.slide_height), "too wide"),
                     ((-IN, 0, IN, IN), "starts off the left"),
                     ((2 * IN, 0, IN, IN), "right edge before its left")):
        with pytest.raises(ValueError):
            stamp_master(data, box)


def test_the_route_hands_back_a_deck_and_says_what_it_did(tmp_path, monkeypatch):
    client = _client(tmp_path, monkeypatch)
    r = client.post("/master", files={"master": ("brand.pptx", _master_bytes(),
                                                 "application/octet-stream")})
    assert r.status_code == 200
    spec_id = next(reversed(web._specs))
    try:
        r = client.post(f"/spec/{spec_id}/pspace",
                        data={"left": "0.5", "top": "0.4",
                              "right": "9.5", "bottom": "7.0"})
        assert r.status_code == 200
        assert r.content[:2] == b"PK", "a real .pptx came back"
        assert "presentation-space.pptx" in r.headers["content-disposition"]
        assert "ToolsToo_PS" in r.headers["X-QC-Note"]

        from qc.stylespec import dominant_master, read_presentation_space

        prs = Presentation(io.BytesIO(r.content))
        space = read_presentation_space(prs, dominant_master(prs))
        assert space, "and the frame it wrote is readable"
        assert round(space["box_emu"][0] / IN, 2) == 0.5
    finally:
        web._specs.pop(spec_id, None)


def test_the_route_refuses_junk_and_an_expired_master(tmp_path, monkeypatch):
    client = _client(tmp_path, monkeypatch)
    r = client.post("/master", files={"master": ("brand.pptx", _master_bytes(),
                                                 "application/octet-stream")})
    spec_id = next(reversed(web._specs))
    try:
        bad = client.post(f"/spec/{spec_id}/pspace",
                          data={"left": "wide", "top": "0.4",
                                "right": "9.5", "bottom": "7.0"})
        assert bad.status_code == 400 and "inches" in bad.json()["error"]

        off = client.post(f"/spec/{spec_id}/pspace",
                          data={"left": "0.5", "top": "0.4",
                                "right": "99", "bottom": "7.0"})
        assert off.status_code == 400 and "outside the slide" in off.json()["error"]
    finally:
        web._specs.pop(spec_id, None)

    gone = client.post("/spec/deadbeef/pspace",
                       data={"left": "0.5", "top": "0.4",
                             "right": "9.5", "bottom": "7.0"})
    assert gone.status_code == 410


def test_the_offer_appears_only_where_it_helps(tmp_path, monkeypatch):
    """On a master with no stated frame it is offered, pre-filled with the
    inferred numbers. On one that already states a frame there is nothing to
    add and the form would be a dead button."""
    client = _client(tmp_path, monkeypatch)
    r = client.post("/master", files={"master": ("brand.pptx", _master_bytes(),
                                                 "application/octet-stream")})
    spec_id = next(reversed(web._specs))
    web._specs.pop(spec_id, None)
    assert "Write these margins into the master" in r.text
    assert "/pspace" in r.text

    _data, (out, _note) = _stamped()
    r2 = client.post("/master", files={"master": ("brand.pptx", out,
                                                  "application/octet-stream")})
    spec_id2 = next(reversed(web._specs))
    web._specs.pop(spec_id2, None)
    assert "Write these margins into the master" not in r2.text
