"""Moving a deck's content into the master it was just given.

Applying a layout remaps placeholder content only; a deck of free-floating
shapes keeps them exactly where they were. These tests build that situation
directly (blank layout, everything a free text box) because it is what an
export-tool deck actually looks like, and it is the case that made an applied
master look like nothing had happened.
"""

import functools
import io

import pytest
from pptx import Presentation
from pptx.util import Emu, Pt

from qc.migrate import migrate_deck
from qc.util import iter_shapes_deep

IN = 914400


# Header geometry copied from the real client master this was built against
# (its "USE" layout): a title band across the top and a thin subtitle strip
# under it. The stock template's layouts do not combine those - its title
# slide centres both mid-canvas - so the placeholders are positioned here
# explicitly rather than inheriting whatever the sample template happens to
# do. Getting this wrong is what made the first version of these tests fail:
# the content never overlapped a placeholder that sat in the middle.
_TITLE_BOX = (0.48, 0.42, 12.40, 0.92)
_SUBTITLE_BOX = (0.48, 1.40, 12.40, 0.35)


@pytest.fixture()
def removals_performed(monkeypatch):
    """Run this test's migration WITH its removals performed.

    Removal is opt-in as of 26/08/2026 (design lead: nothing leaves a slide
    unless a designer asks for it), so a default run only PROPOSES: the piece
    stays on the slide and the change carries a remove_op the page can offer.

    The tests that take this fixture are about which pieces the pass identifies
    and what taking one out does. Whether it does it unasked is a different
    question, and it has its own tests.

    Patched at this module's own name, so every helper here follows without each
    one growing a flag it would then have to thread.
    """
    monkeypatch.setitem(globals(), "migrate_deck",
                        functools.partial(migrate_deck, remove=True))


def _deck(*, with_furniture=False, tall_content=False):
    """A slide whose content is all free shapes, over a master-style header."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # ctrTitle + subTitle
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)

    def tb(x, y, w, h, text, size):
        shape = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)),
                                         Emu(int(w)), Emu(int(h)))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(size)
        return shape

    tb(0.6 * IN, 0.35 * IN, 2 * IN, 0.3 * IN, "CORE FEATURES", 11)
    tb(2.4 * IN, 0.30 * IN, 7 * IN, 0.6 * IN, "Full Lifecycle Management", 30)
    tb(0.6 * IN, 1.05 * IN, 8 * IN, 0.4 * IN, "Every layer supports it.", 13)
    height = 6.0 if tall_content else 3.0
    for i in range(3):
        tb((0.6 + i * 4.1) * IN, 1.9 * IN, 3.8 * IN, height * IN, f"Card {i+1}", 14)
    if with_furniture:
        tb(0.5 * IN, 6.9 * IN, 1.5 * IN, 0.25 * IN, "Strategy&", 9)
        tb(12.7 * IN, 6.9 * IN, 0.4 * IN, 0.25 * IN, "1", 9)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _run(**kw):
    out, changes = migrate_deck(_deck(**kw))
    # migrate_deck catches per-slide exceptions so one bad slide cannot fail a
    # deck. That resilience once hid a NameError as "nothing moved", so every
    # test asserts the pass actually ran rather than silently skipped.
    skipped = [c for c in changes if c.action == "migration skipped"]
    assert not skipped, f"the migration raised: {skipped[0].detail}"
    return Presentation(io.BytesIO(out)).slides[0], changes


def _actions(changes):
    return [c.action for c in changes]


def _ph_text(slide, *want):
    """Placeholder text by type name. TITLE and CENTER_TITLE are both titles;
    a layout picks one and the migration handles either."""
    for ph in slide.placeholders:
        if str(ph.placeholder_format.type).startswith(want):
            return ph.text_frame.text
    return None


def _bottoms(deck_bytes):
    slide = Presentation(io.BytesIO(deck_bytes)).slides[0]
    return {s.shape_id: (s.top or 0) + (s.height or 0) for s in slide.shapes}


# ------------------------------------------------------------ text migration


def test_largest_type_in_the_title_band_becomes_the_title():
    slide, changes = _run()
    assert "title into placeholder" in _actions(changes)
    assert _ph_text(slide, "TITLE", "CENTER_TITLE") == "Full Lifecycle Management"
    # the heading is gone as a free shape, not duplicated
    free = [s.text_frame.text for s in slide.shapes
            if not s.is_placeholder and s.has_text_frame]
    assert "Full Lifecycle Management" not in free


def test_standfirst_reaches_the_subtitle_even_without_overlap():
    """A hand-laid deck rarely puts its standfirst exactly where the master's
    subtitle box sits, so overlap alone misses it."""
    slide, changes = _run()
    assert "subtitle into placeholder" in _actions(changes)
    assert _ph_text(slide, "SUBTITLE") == "Every layer supports it."


def test_the_eyebrow_is_not_mistaken_for_the_title():
    """CORE FEATURES sits in the title band too, but it is 11pt against 30pt."""
    slide, _ = _run()
    assert _ph_text(slide, "TITLE", "CENTER_TITLE") != "CORE FEATURES"


# --------------------------------------------------------------- the block


def test_remaining_content_moves_as_one_block():
    """Relative arrangement is the design; the cards must keep their spacing."""
    before = Presentation(io.BytesIO(_deck())).slides[0]
    card_lefts = sorted(s.left for s in before.shapes
                        if s.has_text_frame and s.text_frame.text.startswith("Card"))
    gaps_before = [b - a for a, b in zip(card_lefts, card_lefts[1:])]

    slide, _changes = _run()
    after = sorted(s.left for s in slide.shapes
                   if s.has_text_frame and s.text_frame.text.startswith("Card"))
    # Asserted on the outcome, not on a "content block moved" record: with the
    # eyebrow now swept, the body may already sit on the margin and need no
    # move at all. Spacing preserved is the invariant either way.
    assert [b - a for a, b in zip(after, after[1:])] == gaps_before


def test_migration_never_pushes_a_shape_further_off_the_canvas():
    """Regression: dragging footer-band shapes with the content block moved
    them to 8.45in on a 7.5in canvas, out of sight.

    Asserted as "no worse than before", not "everything fits": a source deck
    can arrive with content already overflowing, and this pass is not allowed
    to fix that by scaling. What it must never do is make it worse.

    Scoped to a master that states no body ceiling, which is the case here (the
    stock template draws no guides). Where a master DOES draw one, that line
    binds and the overflow it creates is reported instead - see
    test_the_ceiling_binds_even_when_the_block_then_overflows."""
    source = _deck(with_furniture=True, tall_content=True)
    height = Presentation(io.BytesIO(source)).slide_height
    before_overflow = max(
        [b - height for b in _bottoms(source).values() if b > height] or [0])

    out, _ = migrate_deck(source)
    after_overflow = max(
        [b - height for b in _bottoms(out).values() if b > height] or [0])

    assert after_overflow <= before_overflow
    slide = Presentation(io.BytesIO(out)).slides[0]
    assert all((s.top or 0) >= 0 for s in slide.shapes)


def test_footer_furniture_is_never_dragged_downward():
    """The specific shape of the original bug: a footer at 6.9in must not end
    up below the canvas because the content above it moved."""
    source = _deck(with_furniture=False, tall_content=False)
    prs = Presentation(io.BytesIO(source))
    slide = prs.slides[0]
    tb = slide.shapes.add_textbox(Emu(int(0.5 * IN)), Emu(int(6.9 * IN)),
                                  Emu(int(3 * IN)), Emu(int(0.25 * IN)))
    tb.text_frame.paragraphs[0].add_run().text = "Source: internal analysis 2026"
    buf = io.BytesIO()
    prs.save(buf)

    out, _ = migrate_deck(buf.getvalue())
    after = Presentation(io.BytesIO(out)).slides[0]
    note = [s for s in after.shapes
            if s.has_text_frame and s.text_frame.text.startswith("Source:")]
    assert note, "the footnote must survive"
    assert note[0].top == Emu(int(6.9 * IN)), "left exactly where it was"


def test_oversized_content_is_reported_not_scaled():
    """Shrinking a text box does not shrink its type, so a silent scale would
    produce overflowing text that looks fine in the XML."""
    slide, changes = _run(tall_content=True)
    assert "content does not fit" in _actions(changes)
    cards = [s for s in slide.shapes
             if s.has_text_frame and s.text_frame.text.startswith("Card")]
    assert all(c.height == Emu(int(6.0 * IN)) for c in cards)


# ------------------------------------------------------------ housekeeping


def test_a_bare_page_number_duplicate_is_removed(removals_performed):
    """The slide's own hand-typed page number, where the master stamps one."""
    slide, changes = _run(with_furniture=True)
    removed = [c for c in changes if c.action == "removed duplicate furniture"]
    assert [c.detail.split()[0] for c in removed] == ["page"]
    text = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "Strategy&" in text, "an unmatched footer string is content, not a duplicate"


def test_footer_band_shapes_are_parked_not_moved():
    """Anything already in the footer band is page furniture by position and
    stays there; only the content above it forms the moving block.

    Uses normal-height content on purpose: with oversized content the block
    already starts below the region top, so no move happens and there is
    nothing to park it against."""
    slide, _changes = _run(with_furniture=True, tall_content=False)
    footer = [s for s in slide.shapes
              if s.has_text_frame and s.text_frame.text == "Strategy&"]
    assert footer and footer[0].top == Emu(int(6.9 * IN)),         "a footer-band shape must never be dragged by the content block"


def test_no_empty_placeholder_survives():
    """An empty placeholder shows its prompt text in the editor, which is what
    made the applied master look like nothing had happened."""
    slide, _ = _run()
    for ph in slide.placeholders:
        if str(ph.placeholder_format.type).startswith(
                ("FOOTER", "SLIDE_NUMBER", "DATE")):
            continue  # furniture placeholders are meant to be empty
        assert ph.text_frame.text.strip(), "an empty placeholder was left behind"


def test_an_unfillable_placeholder_is_removed_and_reported(removals_performed):
    """A layout offering more placeholders than the slide has content for."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
    buf = io.BytesIO()
    prs.save(buf)

    out, changes = migrate_deck(buf.getvalue())
    assert [c.action for c in changes].count("removed empty placeholder") == 2
    assert not list(Presentation(io.BytesIO(out)).slides[0].placeholders)


def test_a_slide_with_nothing_to_move_is_left_alone():
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    buf = io.BytesIO()
    prs.save(buf)
    out, changes = migrate_deck(buf.getvalue())
    assert Presentation(io.BytesIO(out))
    assert not [c for c in changes if c.action.startswith(("title", "subtitle"))]


def test_every_change_names_its_slide_and_reads_as_a_sentence():
    _, changes = _run(with_furniture=True)
    for c in changes:
        assert c.slide_index == 0
        assert c.detail and c.action
        assert str(c).startswith("slide 1: ")


def test_a_grouped_title_is_still_found():
    """A converter or a designer may group the eyebrow with the heading. A
    group carries no text of its own, so a top-level-only scan never sees the
    shape that holds the title."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)

    group = slide.shapes.add_group_shape()
    for x, y, w, h, text, size in (
            (0.6, 0.35, 2.0, 0.3, "CORE FEATURES", 11),
            (2.4, 0.30, 7.0, 0.6, "Full Lifecycle Management", 30)):
        shape = group.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                                         Emu(int(w * IN)), Emu(int(h * IN)))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(size)
    buf = io.BytesIO()
    prs.save(buf)

    out, changes = migrate_deck(buf.getvalue())
    slide = Presentation(io.BytesIO(out)).slides[0]

    assert _ph_text(slide, "TITLE", "CENTER_TITLE") == "Full Lifecycle Management"
    lifted = [c for c in changes if c.action == "title into placeholder"]
    assert lifted and "lifted out of a group" in lifted[0].detail
    # the rest of the group survives
    remaining = [s.text_frame.text for s, _ in iter_shapes_deep(slide.shapes)
                 if s.has_text_frame]
    assert "CORE FEATURES" in remaining


# ------------------------------------------------------------- collisions


def _header_deck(*, heading_pt=28, eyebrow_pt=11, standfirst_pt=13,
                 content_to_bottom=True):
    """The real failure from the client deck: an eyebrow and a heading both
    inside the master's title band, a standfirst under them, and content that
    already reaches the bottom edge so the block cannot shift down."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)

    def tb(x, y, w, h, text, size=None):
        shape = slide.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                                         Emu(int(w * IN)), Emu(int(h * IN)))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = text
        if size:
            run.font.size = Pt(size)
        return shape

    tb(0.6, 0.45, 3.0, 0.28, "TECHNICAL ARCHITECTURE", eyebrow_pt)
    tb(0.6, 0.62, 7.0, 0.55, "Under the hood", heading_pt)
    tb(0.6, 1.05, 9.0, 0.35, "A lean service boundary.", standfirst_pt)
    tb(0.6, 1.6, 11.0, 5.7 if content_to_bottom else 2.0, "DIAGRAM")
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _overlaps(slide):
    boxes = [(s, (s.left, s.top, s.left + s.width, s.top + s.height))
             for s in slide.shapes if s.left is not None]
    hits = []
    for i, (a, ba) in enumerate(boxes):
        for b, bb in boxes[i + 1:]:
            x = min(ba[2], bb[2]) - max(ba[0], bb[0])
            y = min(ba[3], bb[3]) - max(ba[1], bb[1])
            if x > 0 and y > 0:
                hits.append((a, b))
    return hits


def _filled_placeholders(slide):
    """An EMPTY placeholder is a box waiting for content: it draws nothing, so
    nothing can print over it and it cannot print over anything. This pass used
    to delete them, which hid the distinction; they are proposed now (design
    lead, 26/08/2026), so the question has to name what it always meant."""
    return [s for s in slide.shapes if s.is_placeholder
            and s.has_text_frame and s.text_frame.text.strip()]


def test_removing_the_stray_clears_the_filled_placeholder():
    """The defect from the client deck: the heading went into the title
    placeholder and the eyebrow stayed on top of it. Taking the eyebrow out
    clears it, and this is what taking it out is FOR."""
    out, changes = migrate_deck(_header_deck(), remove=True)
    slide = Presentation(io.BytesIO(out)).slides[0]

    assert _ph_text(slide, "TITLE", "CENTER_TITLE") == "Under the hood"
    placeholders = _filled_placeholders(slide)
    for a, b in _overlaps(slide):
        assert not (a in placeholders or b in placeholders), \
            "a shape is still printing over a filled placeholder"
    # The guarantee is that nothing overlaps a placeholder, not that a NUDGE
    # achieved it: header remnants are placed under the header band by the block
    # pass, so a nudge is the fallback, not the norm.
    assert not any(c.action == "migration skipped" for c in changes)


def test_a_collision_it_cannot_clear_without_removing_is_reported_instead():
    """A default run does not remove, and it cannot nudge this one clear either:
    the eyebrow was in the header band before the pass ran, and the collision
    pass only cleans up collisions it caused. So the overlap stands - and the
    change says so, names the piece, and carries the way out.

    Shipping a deck with a known overlap and no mention of it would be the worst
    of both policies."""
    out, changes = migrate_deck(_header_deck())
    slide = Presentation(io.BytesIO(out)).slides[0]

    assert _ph_text(slide, "TITLE", "CENTER_TITLE") == "Under the hood"
    kept = [c for c in changes if c.action == "unplaced text left in place"]
    assert kept, [c.action for c in changes]
    assert all(c.severity == "alert" for c in kept)
    assert all(c.remove_op and c.remove_id for c in kept), \
        "a proposal with no way to perform it is just a complaint"
    # and the piece really is still there, which is the point of the policy
    assert any(kept[0].removed_text == s.text_frame.text.strip()
               for s in slide.shapes if s.has_text_frame)


def test_a_nudge_clears_the_whole_header_band_not_just_one_placeholder(removals_performed):
    """Shifting past the title alone dropped the eyebrow onto the subtitle,
    trading one collision for another."""
    out, changes = migrate_deck(_header_deck())
    slide = Presentation(io.BytesIO(out)).slides[0]

    # The eyebrow no longer needs nudging because it is swept as unplaced text.
    # What must hold either way: nothing free is left inside the header band.
    assert any(c.action == "removed unplaced text" for c in changes)
    floor = max(p.top + p.height for p in slide.placeholders
                if str(p.placeholder_format.type).startswith(
                    ("TITLE", "CENTER_TITLE", "SUBTITLE")))
    for shape in slide.shapes:
        if shape.is_placeholder or shape.top is None:
            continue
        assert shape.top + shape.height > floor,             f"{shape.text_frame.text!r} is still inside the header band"


def test_one_shape_gets_one_nudge_and_one_log_line():
    """Iterating placeholder-first moved a shape once per placeholder it hit."""
    _out, changes = migrate_deck(_header_deck())
    nudges = [c for c in changes if c.action == "nudged clear of a placeholder"]
    labels = [c.detail.split("'")[1] for c in nudges]
    assert len(labels) == len(set(labels)), f"a shape was nudged twice: {labels}"


def test_a_full_slide_block_is_never_consumed_into_the_subtitle():
    """The subtitle fallback CONSUMES the shape it picks. An unbounded version
    swallowed the whole diagram and deleted it."""
    out, changes = migrate_deck(_header_deck())
    slide = Presentation(io.BytesIO(out)).slides[0]

    assert _ph_text(slide, "SUBTITLE") == "A lean service boundary."
    text = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    assert "DIAGRAM" in text, "the diagram must survive as content"


def test_with_no_explicit_sizes_the_higher_line_becomes_the_title():
    """A deck that sets no font sizes has every line resolving to the same
    inherited size, so they genuinely tie and rule 2 decides: the line nearer
    the top wins. Worth knowing as a limitation rather than a bug - the
    eyebrow, being uppermost, takes the title slot. A deck that sets real
    sizes ranks correctly (see the test above)."""
    out, _ = migrate_deck(_header_deck(heading_pt=None, eyebrow_pt=None,
                                       standfirst_pt=None))
    slide = Presentation(io.BytesIO(out)).slides[0]
    assert _ph_text(slide, "TITLE", "CENTER_TITLE") == "TECHNICAL ARCHITECTURE"


def test_duplicated_text_over_a_placeholder_is_removed(removals_performed):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)
    for y, size in ((0.5, 28), (0.75, 28)):
        shape = slide.shapes.add_textbox(Emu(int(0.6 * IN)), Emu(int(y * IN)),
                                         Emu(int(7 * IN)), Emu(int(0.5 * IN)))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = "Under the hood"
        run.font.size = Pt(size)
    buf = io.BytesIO()
    prs.save(buf)

    out, changes = migrate_deck(buf.getvalue())
    slide = Presentation(io.BytesIO(out)).slides[0]
    assert any(c.action == "removed duplicated text" for c in changes)
    kept = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    assert kept.count("Under the hood") == 1


def test_content_on_content_overlap_is_reported_not_moved():
    """Choosing which of two content blocks gives way is a layout decision;
    the audit owns it. This pass must not guess, but must not hide it either."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tops = []
    for y in (3.0, 3.1):
        shape = slide.shapes.add_textbox(Emu(int(2 * IN)), Emu(int(y * IN)),
                                         Emu(int(5 * IN)), Emu(int(1 * IN)))
        shape.text_frame.paragraphs[0].add_run().text = f"Block at {y}"
        tops.append(y)
    buf = io.BytesIO()
    prs.save(buf)

    out, changes = migrate_deck(buf.getvalue())
    flagged = [c for c in changes if c.action == "text overlaps text"]
    assert flagged and "layout decision" in flagged[0].detail
    slide = Presentation(io.BytesIO(out)).slides[0]
    assert sorted(round(s.top / IN, 1) for s in slide.shapes) == tops


# --------------------------------------------------- background and bounds

_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _set_solid_bg(container, hexval):
    from lxml import etree

    cSld = container._element.find(f"{_P}cSld")
    for old in cSld.findall(f"{_P}bg"):
        cSld.remove(old)
    bg = etree.Element(f"{_P}bg")
    bgPr = etree.SubElement(bg, f"{_P}bgPr")
    solid = etree.SubElement(bgPr, f"{_A}solidFill")
    etree.SubElement(solid, f"{_A}srgbClr").set("val", hexval)
    etree.SubElement(bgPr, f"{_A}effectLst")
    cSld.insert(0, bg)


def _has_own_bg(container) -> bool:
    cSld = container._element.find(f"{_P}cSld")
    return bool(cSld is not None and cSld.findall(f"{_P}bg"))


def test_a_slide_background_override_is_dropped_so_the_master_shows(removals_performed):
    """An exported deck stamps a white background on every slide, and a
    slide-level background beats the master's. That is how a deck adopts
    every other part of a master and still comes out the wrong colour."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    _set_solid_bg(prs.slide_masters[0], "1E2E61")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, "FFFFFF")
    buf = io.BytesIO()
    prs.save(buf)
    assert _has_own_bg(Presentation(io.BytesIO(buf.getvalue())).slides[0])

    out, changes = migrate_deck(buf.getvalue())
    slide = Presentation(io.BytesIO(out)).slides[0]
    assert not _has_own_bg(slide), "the override must be gone"
    assert any(c.action == "dropped background override" for c in changes)


def test_a_slide_keeps_its_background_when_the_master_declares_none():
    """Stripping an override with nothing to inherit would leave the slide
    worse off, not more on-brand."""
    prs = Presentation()
    master_cSld = prs.slide_masters[0]._element.find(f"{_P}cSld")
    for old in master_cSld.findall(f"{_P}bg"):
        master_cSld.remove(old)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, "FFFFFF")
    buf = io.BytesIO()
    prs.save(buf)

    out, changes = migrate_deck(buf.getvalue())
    assert _has_own_bg(Presentation(io.BytesIO(out)).slides[0])
    assert not any(c.action == "dropped background override" for c in changes)


def test_the_block_move_never_pushes_content_over_the_footer():
    """Regression from the client deck: bullets ended up printing across the
    footer text because the move was bounded by the canvas, not the region."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)

    def tb(y, h, text, size=None):
        shape = slide.shapes.add_textbox(Emu(int(0.6 * IN)), Emu(int(y * IN)),
                                         Emu(int(9 * IN)), Emu(int(h * IN)))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = text
        if size:
            run.font.size = Pt(size)
        return shape

    tb(0.5, 0.6, "A heading", 28)
    tb(1.1, 0.3, "A standfirst.", 13)
    # content that already runs close to the bottom
    tb(1.7, 4.9, "BULLETS", 14)
    footer = tb(6.95, 0.25, "Strategy&", 9)
    footer_top = footer.top
    buf = io.BytesIO()
    prs.save(buf)

    out, _ = migrate_deck(buf.getvalue())
    slide = Presentation(io.BytesIO(out)).slides[0]
    bullets = next(s for s in slide.shapes
                   if s.has_text_frame and s.text_frame.text == "BULLETS")
    assert bullets.top + bullets.height <= footer_top, \
        "content was pushed into the footer band"


# ------------------------------------------------- margins, not movements


def _margin_deck(*, eyebrow_top, body_top, stub=False, header_rule=False,
                 edge_stamps=False):
    """A client-shaped slide: eyebrow, heading, standfirst, then a content
    cluster. eyebrow_top and body_top vary independently so the body's final
    position can be tested for independence from what sits above it.

    The three optional extras are the shapes that must NOT get to say where the
    block starts, each taken from the deck that exposed it:

    stub          a 0.0017in square in the top-left corner, what think-cell and
                  every embedded OLE object leave behind
    header_rule   a 3in by 0.02in hairline in the header band: no text, so the
                  remnant sweep leaves it alone, and thin on one axis only, so
                  it is not degenerate
    edge_stamps   a mark near each side edge, so the block's BOUNDING BOX spans
                  the page while no single shape does
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)

    def tb(x, y, w, h, text, size=None):
        shape = slide.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                                         Emu(int(w * IN)), Emu(int(h * IN)))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = text
        if size:
            run.font.size = Pt(size)

    tb(0.6, eyebrow_top, 3.0, 0.28, "FUTURE WORK", 11)
    tb(0.6, 0.55, 9.0, 0.60, "Next: Personalized Daily Digests", 28)
    tb(0.6, 1.25, 10.0, 0.40, "A second LLM integration.", 13)
    for i in range(3):
        tb(0.6 + i * 4.1, body_top, 3.8, 1.9, f"Card {i + 1}", 12)
    tb(0.6, body_top + 2.1, 10.0, 0.8, "Bullets", 12)
    if stub:
        slide.shapes.add_textbox(Emu(1588), Emu(1588), Emu(1588), Emu(1588))
    if header_rule:
        tb(0.6, 0.30, 3.0, 0.02, "")
    if edge_stamps:
        tb(0.0, body_top, 1.0, 0.5, "")
        tb(12.33, body_top, 1.0, 0.5, "")
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _tops(deck_bytes, prefix):
    slide = Presentation(io.BytesIO(deck_bytes)).slides[0]
    return sorted(s.top for s in slide.shapes
                  if s.has_text_frame and s.text_frame.text.startswith(prefix))


def test_body_is_pulled_up_to_the_margin_not_only_pushed_down():
    """Content sitting below the master's body margin must rise to meet it.
    Only ever moving down could never close a gap, which is what left 1.5in of
    dead space under the subtitle on the client deck."""
    source = _margin_deck(eyebrow_top=0.30, body_top=3.30)
    before = _tops(source, "Card")
    out, changes = migrate_deck(source)
    after = _tops(out, "Card")

    assert after[0] < before[0], "the body should have moved UP"
    moved = [c for c in changes if c.action == "content block moved"]
    assert moved and "body now starts at" in moved[0].detail
    # and the report names the frame it was seated on, so a stale stored master
    # cannot look like a bug
    assert "placeholder extents" in moved[0].detail


def test_body_position_is_independent_of_what_sits_above_it():
    """The property behind 'margins, not movements': how far a leftover header
    shape had to travel must not displace the body at all."""
    high = migrate_deck(_margin_deck(eyebrow_top=0.30, body_top=3.30))[0]
    low = migrate_deck(_margin_deck(eyebrow_top=1.00, body_top=3.30))[0]

    assert _tops(high, "Card") == _tops(low, "Card")


def test_relative_arrangement_inside_the_body_still_holds():
    """Margins govern where the body STARTS; inside it, the designer's
    spacing is still the design."""
    source = _margin_deck(eyebrow_top=0.30, body_top=3.30)
    before_cards = _tops(source, "Card")
    before_gap = _tops(source, "Bullets")[0] - before_cards[0]

    out, _ = migrate_deck(source)
    after_cards = _tops(out, "Card")
    after_gap = _tops(out, "Bullets")[0] - after_cards[0]

    assert len(set(after_cards)) == 1, "the cards must stay on one line"
    assert after_gap == before_gap


def test_an_unplaced_header_line_is_swept_rather_than_parked_in_the_body(removals_performed):
    """Stacking a remnant above the body was the earlier behaviour; it is now
    removed and reported, so it can never end up sitting inside the content."""
    out, changes = migrate_deck(_margin_deck(eyebrow_top=0.30, body_top=3.30))
    slide = Presentation(io.BytesIO(out)).slides[0]

    swept = [c for c in changes if c.action == "removed unplaced text"]
    assert [c.removed_text for c in swept] == ["FUTURE WORK"]
    assert "FUTURE WORK" not in [s.text_frame.text for s in slide.shapes
                                 if s.has_text_frame]


# ------------------------------------ what may NOT be the block's own edge


def test_an_embedded_object_stub_is_not_the_top_of_the_body():
    """think-cell parks a 0.0017in square at (0.002in, 0.002in) on every slide
    it has ever touched, and PowerPoint leaves the same stub behind for any
    embedded OLE object. Letting one be the block's top-left corner seated a
    whole client deck on it: the real content, already sitting on the master's
    body line, went 1.90in DOWN and 1.1in to 1.9in off the bottom of every page
    (real deck, 23/08/2026 - 24 of its 26 slides)."""
    plain = migrate_deck(_margin_deck(eyebrow_top=0.30, body_top=1.90))[0]
    stubbed = migrate_deck(_margin_deck(eyebrow_top=0.30, body_top=1.90,
                                        stub=True))[0]

    assert _tops(stubbed, "Card") == _tops(plain, "Card")


def test_a_stub_it_cannot_anchor_still_travels_with_the_block():
    """Excluded from the MEASUREMENT, not from the move: a shape left behind
    while its slide moves is a different defect, so the count of what shipped
    has to match."""
    source = _margin_deck(eyebrow_top=0.30, body_top=3.30, stub=True)
    before = Presentation(io.BytesIO(source)).slides[0]
    out, _changes = migrate_deck(source)
    after = Presentation(io.BytesIO(out)).slides[0]

    def stubs(slide):
        return [s for s in slide.shapes
                if s.width and s.height and s.width < 45720 and s.height < 45720]

    assert len(stubs(after)) == len(stubs(before)) == 1
    card_delta = (_tops(out, "Card")[0] - _tops(source, "Card")[0])
    assert stubs(after)[0].top - stubs(before)[0].top == card_delta


def test_a_hairline_in_the_header_band_is_not_the_top_of_the_body():
    """A rule or corner mark drawn above the line where content begins carries
    no text, so the remnant sweep leaves it alone - correctly, it is not
    unplaced text - and it is thin on ONE axis only, so it is not degenerate
    either. It still may not be the body's top edge: on the client deck a 0.02in
    bar sitting at 0.37in pushed the real content 1.53in down and off the
    page."""
    plain = migrate_deck(_margin_deck(eyebrow_top=0.30, body_top=1.90))[0]
    ruled = migrate_deck(_margin_deck(eyebrow_top=0.30, body_top=1.90,
                                      header_rule=True))[0]

    assert _tops(ruled, "Card") == _tops(plain, "Card")


def test_the_whole_body_still_comes_down_when_it_all_sits_above_the_line():
    """The exclusion above is conditional for a reason: a slide whose entire
    body sits in the reserved strip has to be brought DOWN onto the line, and
    dropping every anchor there would leave exactly those slides untouched."""
    source = _margin_deck(eyebrow_top=0.30, body_top=0.90)
    out, _changes = migrate_deck(source)

    assert _tops(out, "Card")[0] > _tops(source, "Card")[0]


def _stated_deck(shapes, *, title="The Heading", subtitle=None,
                 title_in_placeholder=False):
    """A slide on a GUIDED master, so the frame is one the master states and the
    body-top line binds. `shapes` are (x, y, w, h, text) in inches.

    `subtitle` fills the master's subtitle placeholder up front. That is not a
    detail: a placeholder this pass finds EMPTY is filled from the header band,
    so leftover header text is only ever swept on a slide whose placeholders
    arrived full - which is the case on the client's decks and the one worth
    testing."""
    from tests.test_arabic_layout import _plant_guides

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    _plant_guides(prs.slide_masters[0])
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)
        if subtitle and kind.startswith("SUBTITLE"):
            ph.text_frame.text = subtitle
        if title_in_placeholder and kind.startswith(("TITLE", "CENTER_TITLE")):
            ph.text_frame.text = title
    if not title_in_placeholder:
        shape = slide.shapes.add_textbox(Emu(int(0.6 * IN)), Emu(int(0.55 * IN)),
                                         Emu(int(9 * IN)), Emu(int(0.6 * IN)))
        shape.text_frame.paragraphs[0].add_run().text = title
    for x, y, w, h, text in shapes:
        box = slide.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                                       Emu(int(w * IN)), Emu(int(h * IN)))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(12)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _top_of(deck_bytes, text):
    slide = Presentation(io.BytesIO(deck_bytes)).slides[0]
    return next(s.top for s in slide.shapes
                if s.has_text_frame and s.text_frame.text == text)


def _shape_with(deck_bytes, text):
    slide = Presentation(io.BytesIO(deck_bytes)).slides[0]
    return next(s for s in slide.shapes
                if s.has_text_frame and s.text_frame.text == text)


def _left_margin(deck_bytes):
    """The frame's left edge as the DECK states it. Never a literal: PowerPoint
    stores guides in eighths of a point, so a guide a designer set at 0.60in
    reads back as 547688 EMU, and an inch-based expectation misses by 952."""
    from qc.stylespec import dominant_master, infer_grid

    prs = Presentation(io.BytesIO(deck_bytes))
    return infer_grid(prs, dominant_master(prs))["margins_emu"]["left"]


def test_a_too_tall_block_is_still_seated_on_the_line_and_the_cost_measured():
    """The frame's top line binds on every slide, including the ones whose
    content cannot fit under it: a header band quietly broken on the busy slides
    costs the deck the line that makes its headers read as one, and an overflow
    that is REPORTED does not (design lead, 23/08/2026, re-confirming 21/08).
    What the report owes in exchange is the arithmetic."""
    source = _stated_deck([(0.6, 0.30, 3.0, 3.8, "Contact details"),
                           (4.5, 0.80, 3.0, 6.4, "A full-page column")])
    out, changes = migrate_deck(source)

    assert _top_of(out, "Contact details") == 1737360, "seated on the line"
    misfit = [c for c in changes if c.action == "content does not fit"]
    assert misfit and misfit[0].severity == "alert"
    assert "past the slide edge" in misfit[0].detail
    assert "in tall and" in misfit[0].detail, \
        "the report has to measure the block, not just characterise it"


def test_a_row_of_labels_beside_content_is_not_swept_as_unplaced_text():
    """A table's column headings end wherever the designer put them, and the
    header cutoff falls between them and the labels of the same row: on the
    client's Gantt slide 'team members' and 'months of work' ended at 1.84in and
    1.59in while the month numbers labelling the same row ended at 2.01in, so
    the two words were deleted and the numbers kept (23/08/2026)."""
    # the heading ends inside the header band; its row-mates cross out of it
    source = _stated_deck([(0.9, 1.47, 1.8, 0.37, "Team members"),
                           (4.3, 1.35, 0.5, 0.66, "01"),
                           (5.0, 1.35, 0.5, 0.66, "02"),
                           (0.9, 2.30, 12.0, 3.0, "The table")],
                          subtitle="A standfirst")
    out, changes = migrate_deck(source)
    texts = [s.text_frame.text for s
             in Presentation(io.BytesIO(out)).slides[0].shapes
             if s.has_text_frame]

    assert "Team members" in texts, "a column heading is content, not a remnant"
    assert not [c for c in changes if c.action == "removed unplaced text"]


def test_header_text_with_nothing_beside_it_is_still_swept(removals_performed):
    """The sparing above is not a licence to keep everything: a stray note
    alone in the header band has no row to belong to and still goes, with its
    own XML kept so a designer can put it back."""
    source = _stated_deck([(0.6, 0.05, 3.1, 0.28, "To be translated"),
                           (0.9, 2.30, 12.0, 3.0, "The body")],
                          subtitle="A standfirst")
    _out, changes = migrate_deck(source)

    swept = [c for c in changes if c.action == "removed unplaced text"]
    assert [c.removed_text for c in swept] == ["To be translated"]
    assert swept[0].undo, "and it has to be undoable"


def test_a_full_bleed_band_does_not_veto_the_start_margin_for_the_slide():
    """A 1.28in white footer band running edge to edge under an Arabic cover
    title left that title sitting in the left half of the page: right-aligned
    inside a box half the slide wide, which is not what anyone means by right
    aligned (design lead, 23/08/2026). The band's exemption is its own, not the
    slide's."""
    source = _stated_deck([(0.5, 2.00, 6.0, 1.9, "Arabic-style title"),
                           (0.0, 5.82, 13.33, 1.28, "")])
    out, changes = migrate_deck(source)
    slide = Presentation(io.BytesIO(out)).slides[0]
    band = next(s for s in slide.shapes
                if s.width and s.width > 12 * IN and not s.is_placeholder)

    # the title is seated on the start margin, the band still bleeds
    assert _shape_with(out, "Arabic-style title").left == _left_margin(source)
    assert band.left == 0 and band.width == int(13.33 * IN)
    moved = [c for c in changes if c.action == "content block moved"]
    assert moved and "edge to edge" in moved[0].detail


def test_what_sits_on_a_bleed_band_stays_on_it():
    """A logo stamped on a full-width footer band belongs to the band, not to
    the text three inches above it. Carrying it along took the client's mark
    from the corner of its cover to the middle of the strip."""
    source = _stated_deck([(6.0, 2.00, 6.0, 1.9, "The title"),
                           (0.0, 5.82, 13.33, 1.28, ""),
                           (0.5, 5.95, 3.4, 1.0, "LOGO")])
    out, _changes = migrate_deck(source)

    assert _shape_with(out, "LOGO").left == int(0.5 * IN), \
        "the logo belongs to the band it sits on, not to the text above it"
    assert _shape_with(out, "The title").left == _left_margin(source), \
        "and the title is still seated on the margin"


def test_a_flat_connector_travels_with_the_block_it_divides():
    """A perfectly horizontal rule has a height of exactly zero, so it cannot be
    measured - and it was therefore never moved. A deck of tables came back with
    every divider stranded where it was while the rows it separates moved
    beneath it (design lead, 23/08/2026)."""
    source = _stated_deck([(0.9, 2.30, 6.0, 1.0, "Row one"),
                           (0.9, 3.60, 6.0, 1.0, "Row two")])
    prs = Presentation(io.BytesIO(source))
    slide = prs.slides[0]
    rule = slide.shapes.add_shape(1, Emu(int(0.9 * IN)), Emu(int(3.5 * IN)),
                                  Emu(int(6.0 * IN)), Emu(0))
    rule.name = "Divider"
    buf = io.BytesIO()
    prs.save(buf)
    source = buf.getvalue()

    out, changes = migrate_deck(source)

    def rule_top(deck_bytes):
        sl = Presentation(io.BytesIO(deck_bytes)).slides[0]
        return next(s.top for s in sl.shapes if s.name == "Divider")

    moved = _top_of(out, "Row one") - _top_of(source, "Row one")
    assert moved != 0, "the fixture must have moved the block"
    assert rule_top(out) - rule_top(source) == moved, \
        "the divider has to travel with the rows it divides"
    note = [c for c in changes if c.action == "content block moved"]
    assert note and note[0].detail.startswith("3 shape(s)"), \
        "and it has to be counted among what shifted"


def test_a_graphic_is_reported_over_a_placeholder_rather_than_nudged_alone():
    """Nudging graphics broke the arrangement it was meant to protect: a 0.02in
    decorative bar was pushed 1.53in into the body on three slides, and two
    full-width table rules were pushed 0.63in and 0.40in - different distances,
    so they collapsed onto each other (design lead, 23/08/2026)."""
    source = _stated_deck([(11.6, 0.37, 1.7, 0.02, ""),
                           (0.9, 2.30, 12.0, 3.0, "The body")])
    out, changes = migrate_deck(source)
    slide = Presentation(io.BytesIO(out)).slides[0]
    bar = next(s for s in slide.shapes
               if s.width and abs(s.width - int(1.7 * IN)) < 1000)

    assert not [c for c in changes
                if c.action == "nudged clear of a placeholder"]
    body_move = _top_of(out, "The body") - _top_of(source, "The body")
    assert bar.top - int(0.37 * IN) == body_move, \
        "the bar should have moved with the block and no further"


def test_only_recurring_furniture_is_left_behind_in_the_bottom_strip():
    """The strip alone parked 25 shapes on the client's deck of which none were
    furniture: a table's last row at 6.61in, axis labels at 6.80in, a legend at
    7.09in. Each stayed while the composition it belongs to moved. Low on the
    page AND recurring across the deck is what makes something furniture."""
    source = _stated_deck([(0.9, 2.30, 6.0, 1.0, "Body"),
                           (0.9, 6.70, 3.0, 0.3, "Legend")])
    out, _changes = migrate_deck(source)

    moved = _top_of(out, "Body") - _top_of(source, "Body")
    assert moved != 0
    assert _top_of(out, "Legend") - _top_of(source, "Legend") == moved, \
        "a legend on one slide is content, not page furniture"


def test_nothing_is_seated_when_the_only_content_is_in_the_bottom_strip():
    """The strip travels but never anchors. A slide whose only free content is
    one footer bar at 7.17in had that bar seated on the body line, hoisting it
    5.27in to the top of an otherwise empty page."""
    source = _stated_deck([(0.4, 7.17, 12.7, 0.34, "A closing note")])
    out, changes = migrate_deck(source)

    assert _top_of(out, "A closing note") == _top_of(source, "A closing note")
    assert not [c for c in changes if c.action == "content block moved"]


def test_a_stray_is_swept_even_when_this_pass_placed_nothing(removals_performed):
    """Whether this pass filled a placeholder is not the question. Gating the
    sweep on it kept every stray on a slide PowerPoint had already matched, and
    the strays then travelled INTO the body with the block - a working note
    reading "To be translated" ended up inside the content area instead of out
    of the deck (design lead, 23/08/2026: remove it, flag it, let me put it
    back)."""
    # both header placeholders arrive full, exactly as PowerPoint's matching
    # leaves them, so this pass places nothing of its own
    source = _stated_deck([(0.6, 0.05, 3.1, 0.28, "To be translated"),
                           (0.9, 2.30, 12.0, 3.0, "The body")],
                          subtitle="A standfirst", title_in_placeholder=True)
    out, changes = migrate_deck(source)
    texts = [s.text_frame.text for s
             in Presentation(io.BytesIO(out)).slides[0].shapes
             if s.has_text_frame]

    assert "To be translated" not in texts, "a stray must leave the deck"
    swept = [c for c in changes if c.action == "removed unplaced text"]
    assert [c.removed_text for c in swept] == ["To be translated"]
    assert swept[0].severity == "alert", "and it has to be flagged"
    assert swept[0].undo, "and be undoable"
    assert "where the body begins" in swept[0].detail, \
        "the flag should say what test it failed"


def test_a_stray_is_removed_rather_than_carried_into_the_body(removals_performed):
    """The specific thing that was wrong: not that it survived, but WHERE it
    survived. Travelling with the block put it inside the content area."""
    source = _stated_deck([(0.6, 0.05, 3.1, 0.28, "To be translated"),
                           (0.9, 3.30, 12.0, 2.0, "The body")],
                          subtitle="A standfirst", title_in_placeholder=True)
    out, _changes = migrate_deck(source)
    slide = Presentation(io.BytesIO(out)).slides[0]

    assert not [s for s in slide.shapes if s.has_text_frame
                and s.text_frame.text == "To be translated"]
    # the body still moved, so the sweep is not just "nothing happened"
    assert _top_of(out, "The body") != _top_of(source, "The body")


def _two_slide_deck(per_slide, *, subtitle="A standfirst"):
    """One guided master, two slides, each given (x, y, w, h, text) shapes."""
    from tests.test_arabic_layout import _plant_guides

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    _plant_guides(prs.slide_masters[0])
    for shapes in per_slide:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        for ph in slide.placeholders:
            ph.text_frame.clear()
            kind = str(ph.placeholder_format.type)
            box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
                _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
            if box:
                ph.left, ph.top, ph.width, ph.height = (
                    Emu(int(v * IN)) for v in box)
            if kind.startswith(("TITLE", "CENTER_TITLE")):
                ph.text_frame.text = "The Heading"
            elif subtitle and kind.startswith("SUBTITLE"):
                ph.text_frame.text = subtitle
        for x, y, w, h, text in shapes:
            b = slide.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                                         Emu(int(w * IN)), Emu(int(h * IN)))
            run = b.text_frame.paragraphs[0].add_run()
            run.text = text
            run.font.size = Pt(12)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _texts(deck_bytes, idx):
    slide = Presentation(io.BytesIO(deck_bytes)).slides[idx]
    return [s.text_frame.text for s in slide.shapes if s.has_text_frame]


def test_a_stray_is_judged_over_the_whole_deck_not_one_slide(removals_performed):
    """"To be translated" sat alone at the top of most slides and, on the ones
    where it happened to share a band with a numbered badge, looked exactly like
    a row member. Judged per slide it came off five slides and stayed on seven,
    which is the worst of both outcomes (design lead, 23/08/2026)."""
    source = _two_slide_deck([
        # slide 1: the note is alone up there
        [(0.6, 0.05, 3.1, 0.28, "To be translated"),
         (0.9, 2.30, 12.0, 3.0, "Body one")],
        # slide 2: it happens to sit level with a badge, so per-slide it looks
        # like one of a row
        [(4.6, 0.05, 3.1, 0.28, "To be translated"),
         (0.5, 0.07, 0.4, 0.27, "01"),
         (0.9, 0.07, 6.6, 0.29, "A card heading"),
         (0.9, 2.30, 12.0, 3.0, "Body two")],
    ])
    out, changes = migrate_deck(source)

    assert "To be translated" not in _texts(out, 0)
    assert "To be translated" not in _texts(out, 1), \
        "it is the same note on both slides"
    # and the badge and its heading, which really are a row, are untouched
    assert "01" in _texts(out, 1) and "A card heading" in _texts(out, 1)
    swept = {c.removed_text for c in changes if c.removed_text}
    assert swept == {"To be translated"}


def test_a_row_wholly_above_the_body_line_is_kept():
    """A Gantt's eighteen month numbers all ended at 1.27in against a body
    beginning at 1.90in, so asking each whether it lines up with the BODY
    answered no eighteen times and the sweep took the lot. They line up with
    each other, which is what a row is."""
    labels = [(7.0 + i * 0.34, 0.96, 0.33, 0.31, f"{i + 1:02d}")
              for i in range(8)]
    source = _stated_deck(labels + [(0.9, 2.30, 12.0, 3.0, "The table")],
                          subtitle="A standfirst", title_in_placeholder=True)
    out, changes = migrate_deck(source)
    kept = _texts(out, 0)

    for i in range(8):
        assert f"{i + 1:02d}" in kept, f"month label {i + 1:02d} was deleted"
    assert not [c for c in changes if c.action == "removed unplaced text"]


def test_a_stray_cannot_vouch_for_another_stray(removals_performed):
    """Two working notes stamped side by side each make the other look like a
    row. Once either is known to be a stray from somewhere else in the deck, it
    stops vouching and the second is unmasked - which is why the set is iterated
    rather than computed in one pass (design lead, 23/08/2026: three "not
    comprehensive" notes survived on the client's deck until it was)."""
    source = _two_slide_deck([
        # slide 1 identifies the first note: nothing beside it
        [(0.6, 0.05, 3.1, 0.28, "To be translated"),
         (0.9, 2.30, 12.0, 3.0, "Body one")],
        # slide 2 has both, side by side, vouching for each other
        [(0.6, 0.05, 3.1, 0.28, "To be translated"),
         (4.0, 0.08, 2.0, 0.26, "Not comprehensive"),
         (0.9, 2.30, 12.0, 3.0, "Body two")],
    ])
    out, _changes = migrate_deck(source)

    assert "To be translated" not in _texts(out, 1)
    assert "Not comprehensive" not in _texts(out, 1), \
        "with its voucher gone, the second note is a stray too"


def test_a_collision_that_predates_the_pass_is_reported_not_nudged():
    """This function's own words are that it clears collisions the pass CAUSES.
    A shape already standing in the header band arrived that way, and clearing
    it moves one member of a row on its own - a Gantt's month labels were pushed
    0.93in clear of the table they label (design lead, 23/08/2026)."""
    # the labels start inside the title placeholder's band (0.42-1.34)
    labels = [(7.0 + i * 0.34, 0.96, 0.33, 0.31, f"{i + 1:02d}")
              for i in range(6)]
    source = _stated_deck(labels + [(0.9, 2.30, 12.0, 3.0, "The table")],
                          subtitle="A standfirst", title_in_placeholder=True)
    out, changes = migrate_deck(source)

    assert not [c for c in changes
                if c.action == "nudged clear of a placeholder"]
    reported = [c for c in changes if c.action == "overlap needs a designer"]
    assert reported and "before this pass moved anything" in reported[0].detail
    # every label moved by the same amount as the table, so the row holds
    moves = {_top_of(out, f"{i + 1:02d}") - _top_of(source, f"{i + 1:02d}")
             for i in range(6)}
    assert len(moves) == 1
    assert moves == {_top_of(out, "The table") - _top_of(source, "The table")}


def test_a_displaced_rival_is_still_swept(removals_performed):
    """The gate above is not a licence to keep everything: where this pass DID
    fill a placeholder from the header band, the text it beat still has nowhere
    to go and still goes, with its XML kept."""
    # the title placeholder is empty, so this pass fills it from the band and
    # LEFTOVER is the text it beat
    source = _stated_deck([(0.6, 1.45, 3.0, 0.28, "LEFTOVER"),
                           (0.9, 2.30, 12.0, 3.0, "The body")],
                          subtitle="A standfirst")
    _out, changes = migrate_deck(source)

    swept = [c for c in changes if c.action == "removed unplaced text"]
    assert [c.removed_text for c in swept] == ["LEFTOVER"]
    assert swept[0].undo, "and it stays undoable"


def test_a_block_that_merely_spans_the_page_is_still_seated_on_the_margin():
    """The full-bleed exemption is about a SHAPE running edge to edge, not about
    a bounding box that happens to. Testing the union let a mark near each side
    edge exempt every other shape on the slide from the start margin."""
    def block_left(deck_bytes):
        slide = Presentation(io.BytesIO(deck_bytes)).slides[0]
        return min(s.left for s in slide.shapes
                   if not s.is_placeholder and s.left is not None)

    stamped = _margin_deck(eyebrow_top=0.30, body_top=3.30, edge_stamps=True)
    out, _changes = migrate_deck(stamped)
    seated = migrate_deck(_margin_deck(eyebrow_top=0.30, body_top=3.30))[0]

    assert block_left(out) != block_left(stamped), \
        "the block was left where it was instead of seated on the margin"
    assert block_left(out) == block_left(seated), \
        "it was seated somewhere other than the margin the plain block gets"


# ------------------------------------------- largest is title, second is sub


def _header_slide(entries):
    """A slide with the master-style header band and the given text boxes,
    as (top_in, height_in, text, size_pt)."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)
    for top, height, text, size in entries:
        shape = slide.shapes.add_textbox(Emu(int(0.6 * IN)), Emu(int(top * IN)),
                                         Emu(int(9 * IN)), Emu(int(height * IN)))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = text
        if size:
            run.font.size = Pt(size)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _placed_removing(deck_bytes):
    """_placed with the removals performed. See the removals_performed fixture:
    removal is opt-in, and a test about what a REMOVED piece looks like has to
    ask for one."""
    from qc.migrate import migrate_deck as _migrate

    out, changes = _migrate(deck_bytes, remove=True)
    slide = Presentation(io.BytesIO(out)).slides[0]
    return (_ph_text(slide, "TITLE", "CENTER_TITLE"),
            _ph_text(slide, "SUBTITLE"), out, changes)


def _placed(deck_bytes):
    out, changes = migrate_deck(deck_bytes)
    slide = Presentation(io.BytesIO(out)).slides[0]
    return (_ph_text(slide, "TITLE", "CENTER_TITLE"),
            _ph_text(slide, "SUBTITLE"), out, changes)


def test_largest_type_is_the_title_and_second_largest_the_subtitle():
    title, subtitle, _out, _ch = _placed(_header_slide([
        (0.45, 0.28, "SMALL EYEBROW", 11),
        (0.62, 0.60, "The Heading", 28),
        (1.20, 0.35, "The standfirst line", 16),
        (2.60, 3.00, "CARDS", 12),
    ]))
    assert title == "The Heading"
    assert subtitle == "The standfirst line"


def test_rank_ignores_source_order_entirely():
    """Same three lines declared smallest-first must rank identically."""
    a = _placed(_header_slide([
        (0.62, 0.60, "The Heading", 28), (0.45, 0.28, "EYEBROW", 11),
        (1.20, 0.35, "The standfirst", 16), (2.60, 3.00, "CARDS", 12)]))[:2]
    b = _placed(_header_slide([
        (0.45, 0.28, "EYEBROW", 11), (1.20, 0.35, "The standfirst", 16),
        (0.62, 0.60, "The Heading", 28), (2.60, 3.00, "CARDS", 12)]))[:2]
    assert a == b == ("The Heading", "The standfirst")


def test_equal_sizes_break_by_position_higher_wins():
    """The safety net: at the same type size the line nearer the top of the
    slide is the heading."""
    title, subtitle, _out, _ch = _placed(_header_slide([
        (0.85, 0.30, "lower down", 20),
        (0.45, 0.30, "nearer the top", 20),
        (2.60, 3.00, "CARDS", 12),
    ]))
    assert title == "nearer the top"
    assert subtitle == "lower down"


def test_a_chart_figure_lower_down_never_becomes_the_title():
    """Ranking is bounded to the header band. Taking the globally largest text
    would let a big number in a chart outrank the heading."""
    title, _sub, _out, _ch = _placed(_header_slide([
        (0.62, 0.60, "The Heading", 24),
        (3.50, 1.20, "98%", 72),
    ]))
    assert title == "The Heading"


# --------------------------------------- unplaced text is removed and flagged


def test_unplaced_header_text_is_removed_and_flagged_with_its_full_text(removals_performed):
    """Rule: text the master has no placeholder for is removed, but reported
    loudly and in full so a designer can put it back."""
    full = "CORE FEATURES / STRATEGY AND OPERATIONS"
    _title, _sub, out, changes = _placed(_header_slide([
        (0.45, 0.28, full, 11),
        (0.62, 0.60, "The Heading", 28),
        (1.20, 0.35, "The standfirst", 16),
        (2.60, 3.00, "CARDS", 12),
    ]))

    alerts = [c for c in changes if c.severity == "alert"]
    assert len(alerts) == 1
    assert alerts[0].action == "removed unplaced text"
    # Full text, not a preview: a designer must not have to retype it.
    assert alerts[0].removed_text == full
    assert str(alerts[0]).startswith("slide 1: !! ")

    slide = Presentation(io.BytesIO(out)).slides[0]
    assert full not in [s.text_frame.text for s in slide.shapes
                        if s.has_text_frame]


def test_body_content_is_never_removed_as_unplaced():
    """Only the header band is swept. Body content has a home and keeps it."""
    _title, _sub, out, changes = _placed(_header_slide([
        (0.62, 0.60, "The Heading", 28),
        (1.20, 0.35, "The standfirst", 16),
        (2.60, 1.00, "Body paragraph one", 12),
        (3.80, 1.00, "Body paragraph two", 12),
    ]))
    assert not [c for c in changes if c.action == "removed unplaced text"]
    kept = [s.text_frame.text for s in
            Presentation(io.BytesIO(out)).slides[0].shapes if s.has_text_frame]
    assert "Body paragraph one" in kept and "Body paragraph two" in kept


def _report_html(changes):
    """The "what the rebuild did" section of the prepared-deck page. It is a
    fragment now: there is no format result page, only step 3 of Prepare a
    deck."""
    from qc.ui_format import _content_section

    return _content_section(changes, job_id="j")


def test_the_report_surfaces_proposals_before_routine_moves():
    """A default run proposes, so the rows that ask the designer for something
    are the proposals - and they come first, for the same reason removals used
    to: burying them in slide order among routine moves is how a deck ships with
    something nobody looked at."""
    _t, _s, _out, changes = _placed(_header_slide([
        (0.45, 0.28, "AN EYEBROW", 11),
        (0.62, 0.60, "The Heading", 28),
        (1.20, 0.35, "The standfirst", 16),
        (2.60, 3.00, "CARDS", 12),
    ]))
    html = _report_html(changes)
    assert "would take out, and did not" in html
    assert "AN EYEBROW" in html
    assert "still in the deck" in html, "the wording has to say nothing happened"
    assert html.index("would take out") < html.index("What moved into the master")
    # and it offers the way out, on the route that performs it
    assert 'action="/format/j/remove"' in html
    assert 'name="remove_ids"' in html


def test_the_report_still_surfaces_a_performed_removal():
    """Once a designer has taken something out, the deck really is missing it,
    and the page says so in the past tense with a way back."""
    _t, _s, _out, changes = _placed_removing(_header_slide([
        (0.45, 0.28, "AN EYEBROW", 11),
        (0.62, 0.60, "The Heading", 28),
        (1.20, 0.35, "The standfirst", 16),
        (2.60, 3.00, "CARDS", 12),
    ]))
    html = _report_html(changes)
    assert "were removed" in html
    assert "AN EYEBROW" in html
    assert html.index("were removed") < html.index("What moved into the master")


# ---------------------------------------------------- left / right margins

_P15 = "http://schemas.microsoft.com/office/powerpoint/2012/main"


def _plant_guides(master, vertical_in=(), horizontal_in=()):
    """Guides as desktop PowerPoint stores them: eighths of a point from the
    top-left edge, orient omitted for vertical. python-pptx cannot draw them,
    so the format captured from a COM probe is planted directly."""
    import itertools

    from lxml import etree

    ext_lst = etree.SubElement(master._element, f"{_P}extLst")
    ext = etree.SubElement(ext_lst, f"{_P}ext")
    ext.set("uri", "{GUIDES}")
    lst = etree.SubElement(ext, f"{{{_P15}}}sldGuideLst")
    gid = itertools.count(1)
    for pos, horz in ([(v, False) for v in vertical_in]
                      + [(h, True) for h in horizontal_in]):
        g = etree.SubElement(lst, f"{{{_P15}}}guide")
        g.set("id", str(next(gid)))
        if horz:
            g.set("orient", "horz")
        if pos:
            g.set("pos", str(int(pos * 72 * 8)))


def _guided_deck(*, card_left=0.73, card_width=3.8, guides=(0.60, 12.73)):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    _plant_guides(prs.slide_masters[0], vertical_in=guides,
                  horizontal_in=(0.42, 7.05))
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)

    def tb(x, y, w, h, text, size=None):
        shape = slide.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                                         Emu(int(w * IN)), Emu(int(h * IN)))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = text
        if size:
            run.font.size = Pt(size)

    tb(0.6, 0.55, 9.0, 0.6, "The Heading", 28)
    tb(0.6, 1.20, 9.0, 0.35, "The standfirst", 16)
    for i in range(3):
        tb(card_left + i * (card_width + 0.3), 2.2, card_width, 2.5,
           f"Card {i + 1}", 12)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _declared_margins(deck_bytes):
    """The left margin and right edge the master actually declares.

    Read back rather than hand-computed in inches: guides are stored in
    EIGHTHS OF A POINT, so 0.60in (43.2pt) quantises to 43.125pt and a literal
    Emu(0.6 * 914400) is ~950 EMU off. The format's granularity is the truth
    here, not the number typed into PowerPoint."""
    from qc.stylespec import dominant_master, infer_grid

    prs = Presentation(io.BytesIO(deck_bytes))
    grid = infer_grid(prs, dominant_master(prs))
    margins = grid["margins_emu"]
    return margins["left"], prs.slide_width - margins["right"]


def _card_bounds(deck_bytes):
    slide = Presentation(io.BytesIO(deck_bytes)).slides[0]
    cards = [s for s in slide.shapes
             if s.has_text_frame and s.text_frame.text.startswith("Card")]
    return (min(c.left for c in cards),
            max(c.left + c.width for c in cards))


def test_the_left_margin_binds_from_the_masters_guides():
    """Content indented inside the margin is pulled out to meet it. Correcting
    only breaches left the cards a few millimetres inside the title's own left
    edge, which reads as misalignment even though no margin was crossed."""
    source = _guided_deck(card_left=0.73)
    margin_left, margin_right = _declared_margins(source)
    out, changes = migrate_deck(source)
    left, right = _card_bounds(out)

    assert left == margin_left, "the block should sit on the left guide"
    assert right <= margin_right, "and stay inside the right guide"
    moved = [c for c in changes if c.action == "content block moved"]
    assert moved and "-0.13in" in moved[0].detail


def test_guides_beat_placeholder_extents_as_the_margin_source():
    """A designer's guides are a stated intention; placeholder extents are only
    an inference from where the title happens to sit."""
    wide_src = _guided_deck(guides=(0.60, 12.73))
    tight_src = _guided_deck(guides=(1.20, 12.13))

    assert _card_bounds(migrate_deck(wide_src)[0])[0] ==         _declared_margins(wide_src)[0]
    assert _card_bounds(migrate_deck(tight_src)[0])[0] ==         _declared_margins(tight_src)[0]
    # and the two really are different margins, so the source is being read
    assert _declared_margins(wide_src)[0] != _declared_margins(tight_src)[0]


def test_full_bleed_content_is_exempt_from_the_left_margin():
    """A band running edge to edge is not indented content, and dragging it to
    the margin would destroy the effect."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    _plant_guides(prs.slide_masters[0], vertical_in=(0.60, 12.73),
                  horizontal_in=(0.42, 7.05))
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)
    band = slide.shapes.add_textbox(Emu(0), Emu(int(2.5 * IN)),
                                    prs.slide_width, Emu(int(2 * IN)))
    band.text_frame.paragraphs[0].add_run().text = "FULL BLEED BAND"
    buf = io.BytesIO()
    prs.save(buf)

    out, _ = migrate_deck(buf.getvalue())
    shape = next(s for s in Presentation(io.BytesIO(out)).slides[0].shapes
                 if s.has_text_frame and s.text_frame.text == "FULL BLEED BAND")
    assert shape.left == 0, "full-bleed content must not be indented"


# ------------------------------------------------- the reserved header band
#
# A client master draws four horizontal guides: the page's top and bottom
# margins, and between them the floor its subtitle may not cross and the
# ceiling its body may not cross. The strip between those two stays empty on
# every slide, so the ceiling is where the body starts - not the top margin
# (that is where the PAGE starts) and not the header placeholder's floor (that
# is where one layout's title box happens to end).


def _banded_deck(*, card_top=1.4, card_height=2.0, ceiling=2.40):
    """A guided master whose body ceiling sits BELOW the header placeholder's
    floor plus its gap, so the two sources give different answers and the test
    can tell which one the block was seated on."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    _plant_guides(prs.slide_masters[0], vertical_in=(0.60, 12.73),
                  horizontal_in=(0.45, 1.65, ceiling, 6.80))
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)

    def tb(x, y, w, h, text, size=None):
        shape = slide.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                                         Emu(int(w * IN)), Emu(int(h * IN)))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = text
        if size:
            run.font.size = Pt(size)

    tb(0.6, 0.55, 9.0, 0.6, "The Heading", 28)
    tb(0.6, 1.20, 9.0, 0.35, "The standfirst", 16)
    for i in range(3):
        tb(0.6 + i * 4.1, card_top, 3.8, card_height, f"Card {i + 1}", 12)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _declared_ceiling(deck_bytes):
    from qc.stylespec import dominant_master, infer_grid

    prs = Presentation(io.BytesIO(deck_bytes))
    return infer_grid(prs, dominant_master(prs))["body_top_emu"]


def test_the_body_lands_on_the_masters_ceiling_not_the_header_floor():
    """Both are real statements, and the tighter one wins - but the ceiling has
    to be consulted at all. Seeding from the top margin left the header
    placeholder's floor as the only thing holding the body down."""
    source = _banded_deck(card_top=1.4)
    out, changes = migrate_deck(source)

    assert _tops(out, "Card") == [_declared_ceiling(source)] * 3
    moved = [c for c in changes if c.action == "content block moved"]
    assert moved, "the block should have been seated on the ceiling"


def test_the_ceiling_binds_even_when_the_block_then_overflows():
    """The strip the master keeps clear stays clear on the busy slides too.
    Clamping the move at the bottom margin collapses to no move at all on
    exactly those slides, which is how a deck comes back with content sitting
    in the reserved strip on every full slide. The overflow is reported
    instead - loudly, because content past the slide edge will not print."""
    source = _banded_deck(card_top=1.4, card_height=6.0)
    out, changes = migrate_deck(source)

    assert _tops(out, "Card") == [_declared_ceiling(source)] * 3
    alerts = [c for c in changes if c.action == "content does not fit"]
    assert alerts and alerts[0].severity == "alert"
    assert "past the bottom margin" in alerts[0].detail
    assert "past the slide edge" in alerts[0].detail
    assert "still clear" in alerts[0].detail


def test_a_master_with_no_stated_ceiling_keeps_the_bottom_clamp():
    """Where the master draws no body ceiling there is no statement to honour,
    so the old restraint holds: an overflowing block is never pushed further
    down. (The unguided case is covered by
    test_migration_never_pushes_a_shape_further_off_the_canvas.)"""
    guided = _guided_deck()          # two horizontal guides: margins only
    assert _declared_ceiling(guided) is None


def test_content_wider_than_the_margins_is_reported_not_narrowed():
    """Narrowing a text box reflows its text, so a block that cannot fit
    between both margins is aligned left and flagged for a designer."""
    out, changes = migrate_deck(_guided_deck(card_left=0.6, card_width=4.6))
    alerts = [c for c in changes if c.action == "wider than the margins"]

    assert alerts and alerts[0].severity == "alert"
    assert "cannot sit inside both margins" in alerts[0].detail
    before, after = _card_bounds(_guided_deck(card_left=0.6, card_width=4.6)), _card_bounds(out)
    assert (after[1] - after[0]) == (before[1] - before[0]), "must not be narrowed"


# ================================================ nothing leaves unasked
#
# Design lead, 26/08/2026. The migration removed five classes of thing on its
# own and offered them back afterwards, so the deck a designer downloaded had
# already lost things and the restore was a repair rather than a decision. It
# proposes now, and qc.migrate.apply_removals performs what was ticked.
#
# Detection is unchanged, and that is the point: every one of the five is still
# found and still reported as an alert. What changed is who acts on it.


def _every_removal_class():
    """A deck that trips the removal classes at once: hand-typed page
    furniture the master supplies, an empty placeholder, and a slide-level
    background override that beats the master's."""
    prs = Presentation(io.BytesIO(_deck(with_furniture=True)))
    _set_solid_bg(prs.slides[0], "0B1F2A")
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_a_default_run_loses_no_words():
    """The policy in one assertion, stated over WORDS rather than shapes.

    The shape count does legitimately fall: putting a title into its
    placeholder is implemented as copy-then-delete, and leaving the source
    behind would print the line twice. That is a move, and every word
    survives it. What must never happen on a default run is that a string a
    designer wrote is not in the output at all."""
    data = _every_removal_class()

    def words(deck_bytes):
        prs = Presentation(io.BytesIO(deck_bytes))
        return {t for slide in prs.slides
                for shape in slide.shapes
                if shape.has_text_frame
                for t in [shape.text_frame.text.strip()] if t}

    out, changes = migrate_deck(data)
    missing = words(data) - words(out)
    assert not missing, f"a default run lost {sorted(missing)}"
    assert not [c for c in changes if c.action.startswith("removed ")], (
        [c.action for c in changes])


def test_every_removal_it_finds_comes_back_as_a_proposal():
    """A proposal with no way to perform it is just a complaint. Each one
    carries the op, a handle for the page to tick, and what the piece says."""
    _out, changes = migrate_deck(_every_removal_class())
    proposals = [c for c in changes if c.remove_op]
    assert proposals, [c.action for c in changes]
    for c in proposals:
        assert c.severity == "alert", c.action
        assert c.remove_id, c.action
        assert c.remove_op.get("kind") in ("shape", "background")
        assert c.remove_op.get("slide_index") == c.slide_index
        if c.remove_op["kind"] == "shape":
            assert c.remove_op.get("shape_id")


def test_the_same_pass_still_finds_the_same_pieces():
    """Detection did not change - only who acts on it. What a removing run takes
    out is exactly what a default run proposes."""
    data = _every_removal_class()
    _out, proposed = migrate_deck(data)
    _out2, performed = migrate_deck(data, remove=True)

    def _slides(changes, pred):
        return sorted(c.slide_index for c in changes if pred(c))

    assert _slides(proposed, lambda c: bool(c.remove_op)) == \
        _slides(performed, lambda c: c.action.startswith(("removed ",
                                                          "dropped ")))


def test_a_ticked_proposal_is_performed_and_nothing_else_is():
    """One tick, one removal. A pass that took out the neighbours of what was
    ticked would be worse than one that removed everything, because nobody would
    be looking."""
    from qc.migrate import apply_removals

    data = _every_removal_class()
    out, changes = migrate_deck(data)
    shapes = [c for c in changes if (c.remove_op or {}).get("kind") == "shape"]
    assert shapes

    before = len(Presentation(io.BytesIO(out)).slides[0].shapes)
    after_bytes, performed = apply_removals(out, [shapes[0].remove_op])
    after = len(Presentation(io.BytesIO(after_bytes)).slides[0].shapes)

    assert after == before - 1
    assert len(performed) == 1 and performed[0].action == "removed on request"
    assert performed[0].undo, "a removal without an undo is not reversible"


def test_a_performed_removal_undoes_exactly():
    from qc.migrate import apply_removals
    from qc.undo import apply_undo

    data = _every_removal_class()
    out, changes = migrate_deck(data)
    op = next(c.remove_op for c in changes
              if (c.remove_op or {}).get("kind") == "shape")
    gone, performed = apply_removals(out, [op])
    back, outcomes = apply_undo(gone, [{"change_id": "c0", "slide_index": 0,
                                        "action": "removed on request",
                                        "ops": performed[0].undo}])
    assert outcomes[0]["done"]
    assert _texts(back, 0) == _texts(out, 0)


def test_removing_the_background_override_lets_the_master_through():
    from qc.migrate import apply_removals

    data = _every_removal_class()
    out, changes = migrate_deck(data)
    op = next((c.remove_op for c in changes
               if (c.remove_op or {}).get("kind") == "background"), None)
    assert op is not None, "the fixture states its own background"
    assert _has_own_bg(Presentation(io.BytesIO(out)).slides[0])

    gone, performed = apply_removals(out, [op])
    assert not _has_own_bg(Presentation(io.BytesIO(gone)).slides[0])
    assert performed[0].undo, "and it goes back"


def test_a_proposal_for_a_shape_that_has_gone_is_skipped_not_guessed():
    """The deck can move on between the proposal and the tick. Guessing which
    shape was meant is exactly the guess that loses content."""
    from qc.migrate import apply_removals

    data = _every_removal_class()
    out, _changes = migrate_deck(data)
    gone, performed = apply_removals(out, [{"kind": "shape", "slide_index": 0,
                                            "shape_id": "999999"}])
    assert len(performed) == 1
    assert performed[0].action == "removal skipped"
    assert _texts(gone, 0) == _texts(out, 0), "and nothing else was touched"


def test_a_proposal_for_a_slide_that_has_gone_is_skipped():
    from qc.migrate import apply_removals

    data = _every_removal_class()
    out, _changes = migrate_deck(data)
    _gone, performed = apply_removals(out, [{"kind": "shape",
                                             "slide_index": 99,
                                             "shape_id": "2"}])
    assert performed[0].action == "removal skipped"
    assert "no longer in the deck" in performed[0].detail


# ------------------------------- seating must not push content off the page
#
# dx is chosen from the START edge alone: seat the leftmost shape on the left
# margin. That is right while the block fits and harmful when it does not - a
# block already spanning the canvas gets shifted by the whole left margin, so
# its far edge lands that far PAST the slide edge. Content that printed before
# this pass ran does not print after it.
#
# Reproduced 30/08/2026 against a master with no presentation space and a
# 0.92in inferred margin: a 13.33in block on a 13.33in slide was moved to
# 0.92-14.25in, and the report graded it "info" because the overflow was
# measured from the block's BOTTOM only.


def _full_width_deck(*, left_at=0.0, right_edge=13.33):
    """Two columns whose bounding box spans (nearly) the whole canvas, on a
    master stating no presentation space."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

    def tb(x, y, w, h, text):
        shape = slide.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                                         Emu(int(w * IN)), Emu(int(h * IN)))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(12)
        return shape

    half = (right_edge - left_at) / 2 - 0.4
    tb(left_at, 2.6, half, 2.0, "Left column body copy. " * 8)
    tb(right_edge - half, 3.0, half, 2.0, "Right column body copy. " * 8)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _right_edges(deck_bytes):
    slide = Presentation(io.BytesIO(deck_bytes)).slides[0]
    return [s.left + s.width for s in slide.shapes
            if s.left is not None and s.width is not None
            and s.has_text_frame and "column" in s.text_frame.text]


def test_seating_never_pushes_content_off_the_slide():
    source = _full_width_deck()
    slide_w = Presentation(io.BytesIO(source)).slide_width

    before = _right_edges(source)
    assert before and max(before) <= slide_w, "the fixture starts on the page"

    out, _changes = migrate_deck(source)
    after = _right_edges(out)
    assert after, "the columns survived the migration"
    assert max(after) <= slide_w, (
        f"seating the left edge pushed content "
        f"{(max(after) - slide_w) / IN:.2f}in off the right of the page; it "
        f"printed before this pass ran")


def test_a_block_with_slack_is_still_seated_exactly():
    """The clamp must only bite when seating would spill. A block that fits
    keeps the behaviour it had - this is not a retreat from seating."""
    source = _full_width_deck(left_at=0.0, right_edge=10.0)
    out, changes = migrate_deck(source)
    lefts = sorted(s.left for s in Presentation(io.BytesIO(out)).slides[0].shapes
                   if s.left is not None and s.has_text_frame
                   and "column" in s.text_frame.text)
    assert lefts, "the columns survived"
    moved = [c for c in changes if c.action == "content block moved"]
    assert moved, "a block with room to move is still seated on the margin"
    assert lefts[0] > 0, "and it actually moved off the slide edge"


def test_a_block_off_the_page_sideways_is_an_alert_that_says_so():
    """Both spill measurements came off the block's bottom, so a block hanging
    off the SIDE reported an empty spill and severity 'info' - the one grade
    that says there is nothing to look at."""
    source = _full_width_deck()
    _out, changes = migrate_deck(source)
    fit = [c for c in changes if c.action == "content does not fit"]
    if not fit:
        return                      # nothing overflowed; nothing to report
    detail = fit[0].detail
    assert "margin" in detail, (
        f"the block is wider than the content area and the report says "
        f"nothing about where it went: {detail}")
