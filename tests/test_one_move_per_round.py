"""One move per shape per apply round, and one judgment per collection.

Real-deck finding (12/08/2026, RTL strategy map): two row-spacing fixes moved a
chip row and its goals-box row by DIFFERENT horizontal deltas, tearing the
vertical pairs apart - and the box also rode the chip's fix as a satellite
before being moved again by its own.

That is now prevented twice over. A chip stacked on its box is a SATELLITE of
it, so the pair is judged once, as one collection, and the fix moves both
(design lead, 20/08/2026). And where two fixes genuinely do share a shape, the
first claims its movers and the rest wait for the next re-audit round.
"""

import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

import qc.modules.margin_alignment as ma
from qc.fixer import apply_fixes, is_fixable

from tests.conftest import save_and_ctx

MM = 36000

_SPACING = ("margin_alignment.uneven_spacing", "margin_alignment.cluster_rhythm")


def test_a_collection_is_judged_once_and_moves_as_one(make_prs, en_profile,
                                                      tmp_path):
    """Three chip+box column pairs with one odd gap between columns.

    The chips are satellites of their boxes, so the row is judged ONCE rather
    than as two rows that could be corrected by different deltas, and every
    chip must still be centered on its box afterwards."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chips, boxes = [], []
    lefts = (30, 100, 185)  # 15mm odd gap between columns 2 and 3
    for left in lefts:
        chip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left * MM),
                                      Emu(40 * MM), Emu(40 * MM),
                                      Emu(10 * MM))
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left * MM),
                                     Emu(53 * MM), Emu(40 * MM),
                                     Emu(30 * MM))
        chips.append(chip)
        boxes.append(box)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    recs = [r.to_dict() for r in ma.detect(ctx)
            if r.issue_type in _SPACING and is_fixable(r.to_dict())]
    assert len(recs) == 1, "the pair row is one collection, so one finding"
    assert recs[0]["shape_id"] in {str(b.shape_id) for b in boxes}, \
        "judged on the anchor, not on the satellite riding it"

    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {r["record_id"] for r in recs})
    assert result.applied == 1, [o.reason for o in result.outcomes]

    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    by_id = {s.shape_id: s for s in out.shapes}
    moved = 0
    for chip, box in zip(chips, boxes):
        c, b = by_id[chip.shape_id], by_id[box.shape_id]
        assert abs((c.left + c.width // 2) - (b.left + b.width // 2)) <= 20, (
            "chip drifted off its box: pair torn apart")
        moved += 1 if b.left != box.left else 0
    assert moved, "the odd gap should have been closed"


def test_two_fixes_sharing_a_shape_apply_once_per_round(make_prs, en_profile,
                                                        tmp_path):
    """One shape can be the tail of a row AND of a column, and the two fixes
    carry different deltas. The first claims its movers; the second is skipped
    with a reason and comes back on the next round."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def rect(left, top):
        return slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left * MM),
                                      Emu(top * MM), Emu(40 * MM),
                                      Emu(25 * MM))

    corner = None
    for left, top in ((30, 100), (100, 100), (185, 100)):
        corner = rect(left, top)          # row: gaps 30 / 45mm
    for top in (20, 55):
        rect(185, top)                    # column through the same shape
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    recs = [r.to_dict() for r in ma.detect(ctx)
            if r.issue_type in _SPACING and is_fixable(r.to_dict())]
    shared = [r for r in recs if r["shape_id"] == str(corner.shape_id)]
    assert len(shared) >= 2, "the corner shape ends both a row and a column"

    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {r["record_id"] for r in recs})
    skipped = [o for o in result.outcomes if o.outcome == "skipped"]
    assert result.applied == 1
    assert any("already applied this round" in o.reason for o in skipped)
