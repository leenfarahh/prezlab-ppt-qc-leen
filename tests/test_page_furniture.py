"""Page furniture: fake slide numbers and off-canvas footer text are judged
against the deck's OWN layout/master (no profile needed) and fixed the way
the ground-truth designer did: replace the fake with a real inheriting
placeholder; align the off-canvas footer's bottom to the layout baseline."""

import copy
import io

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Emu

import qc.modules.header_footer as hf
from qc.fixer import apply_fixes, is_fixable
from tests.conftest import save_and_ctx

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _furniture(ctx, kind):
    return [r.to_dict() for r in hf.detect(ctx)
            if r.issue_type == f"header_footer.{kind}"]


def _sldnum_phs(slide):
    return [p for p in slide.placeholders
            if p.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER]


def _textbox(slide, text, left, top, w=500000, h=140000):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tb.text_frame.text = text
    return tb


def test_fake_number_detected_and_replaced(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _slide(prs)
    H = prs.slide_height
    fake = _textbox(slide, "1", 11000000, int(H * 0.92))
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    recs = _furniture(ctx, "fake_slide_number")
    assert len(recs) == 1
    rec = recs[0]
    assert rec["severity"] == "error"
    assert rec["shape_id"] == str(fake.shape_id)
    assert is_fixable(rec)

    deck = ctx.deck_path.read_bytes()
    result = apply_fixes(deck, recs, {rec["record_id"]})
    assert result.applied == 1

    cleaned = Presentation(io.BytesIO(result.cleaned_bytes))
    out = cleaned.slides[0]
    assert not any(str(s.shape_id) == rec["shape_id"] for s in out.shapes)
    phs = _sldnum_phs(out)
    assert len(phs) == 1
    layout_ph = next(p for p in out.slide_layout.placeholders
                     if p.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER)
    # xfrm stripped: geometry inherits live from the layout
    assert phs[0].placeholder_format.idx == layout_ph.placeholder_format.idx
    assert (phs[0].left, phs[0].top) == (layout_ph.left, layout_ph.top)
    assert 'type="slidenum"' in phs[0]._element.xml

    # converged: a re-audit finds no furniture issues
    ctx2 = save_and_ctx(cleaned, tmp_path, en_profile, name="cleaned.pptx")
    assert _furniture(ctx2, "fake_slide_number") == []


def test_fake_number_guards(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _slide(prs)
    H = prs.slide_height
    _textbox(slide, "01", 11000000, 100000, h=200000)     # top strip, padded
    _textbox(slide, "1", 5000000, int(H * 0.5))           # mid-slide: content
    _textbox(slide, "7", 11000000, int(H * 0.92))         # wrong number
    _textbox(slide, "1", 11000000, int(H * 0.92),         # too tall for furniture
             h=600000)
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = _furniture(ctx, "fake_slide_number")
    assert len(recs) == 1
    assert recs[0]["old_value"] == "01"


def test_slidenum_field_outside_placeholder_is_fake(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _slide(prs)
    H = prs.slide_height
    tb = _textbox(slide, "", 11000000, int(H * 0.92))
    p = tb.text_frame.paragraphs[0]._p
    fld = etree.SubElement(p, f"{{{A}}}fld")
    fld.set("id", "{11111111-2222-3333-4444-555555555555}")
    fld.set("type", "slidenum")
    t = etree.SubElement(fld, f"{{{A}}}t")
    t.text = "9"
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = _furniture(ctx, "fake_slide_number")
    assert len(recs) == 1 and recs[0]["severity"] == "error"


def test_fake_number_without_layout_placeholder_flag_only(make_prs, en_profile,
                                                          tmp_path):
    prs = make_prs()
    slide = _slide(prs)
    H = prs.slide_height
    for source in (prs.slide_layouts[6], prs.slide_masters[0]):
        for ph in list(source.placeholders):
            if ph.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                ph._element.getparent().remove(ph._element)
    _textbox(slide, "1", 11000000, int(H * 0.92))
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = _furniture(ctx, "fake_slide_number")
    assert len(recs) == 1
    assert recs[0]["severity"] == "warning"
    assert recs[0]["new_value"] is None
    assert not is_fixable(recs[0])


def test_two_fakes_yield_one_placeholder(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _slide(prs)
    H = prs.slide_height
    _textbox(slide, "1", 11000000, int(H * 0.92))
    _textbox(slide, "1", 200000, 100000, h=200000)
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = _furniture(ctx, "fake_slide_number")
    assert len(recs) == 2
    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {r["record_id"] for r in recs})
    assert result.applied == 2
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    assert len(_sldnum_phs(out)) == 1
    texts = [s.text_frame.text for s in out.shapes
             if getattr(s, "has_text_frame", False) and not s.is_placeholder]
    assert "1" not in texts


def test_off_canvas_footer_moved_to_baseline(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _slide(prs)
    H = prs.slide_height
    src = _textbox(slide, "Source: analysis", 2500000, H - 1000, h=300000)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    recs = _furniture(ctx, "footer_off_canvas")
    assert len(recs) == 1
    rec = recs[0]
    assert rec["severity"] == "error" and is_fixable(rec)
    ftr = next(p for p in prs.slide_layouts[6].placeholders
               if p.placeholder_format.type == PP_PLACEHOLDER.FOOTER)
    target = ftr.top + ftr.height - 300000
    assert int(rec["new_value"]) == target

    result = apply_fixes(ctx.deck_path.read_bytes(), recs, {rec["record_id"]})
    assert result.applied == 1
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    moved = next(s for s in out.shapes if str(s.shape_id) == rec["shape_id"])
    assert moved.top == target
    assert moved.top + moved.height <= H


def test_on_canvas_footer_not_flagged(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _slide(prs)
    H = prs.slide_height
    _textbox(slide, "Source: analysis", 2500000, int(H * 0.9), h=300000)
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert _furniture(ctx, "footer_off_canvas") == []
