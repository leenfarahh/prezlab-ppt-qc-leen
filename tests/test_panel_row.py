"""Panel-row alignment: sibling container panels whose row is broken are
moved as blocks (panel + contents), the way designers actually work.
Ground truth 20/07/2026: a 69-shape Consultants panel sat 6.2mm off its
row; the per-shape intent window never saw it."""

import io

from pptx import Presentation
from pptx.util import Emu

import qc.modules.margin_alignment as ma
from qc.fixer import apply_fixes, is_fixable
from tests.conftest import save_and_ctx

IN = 914400


def _panel(slide, left, top, w=3 * IN, h=2 * IN, label="x"):
    """A rectangle panel containing one small text box (so it registers as
    a container)."""
    from pptx.enum.shapes import MSO_SHAPE

    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top),
                                  Emu(w), Emu(h))
    inner = slide.shapes.add_textbox(Emu(left + 200000), Emu(top + 150000),
                                     Emu(800000), Emu(300000))
    inner.text_frame.text = label
    return rect, inner


def _rows(ctx, issue="panel_row_misaligned"):
    return [r.to_dict() for r in ma.detect(ctx)
            if r.issue_type == f"margin_alignment.{issue}"]


def test_broken_panel_row_detected_and_block_moved(make_prs, en_profile,
                                                   tmp_path):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    T = IN
    _panel(slide, IN, T)
    _panel(slide, int(4.5 * IN), T)
    bad, inner = _panel(slide, IN * 8, T + 228600)  # 6.35mm off the row
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    recs = _rows(ctx)
    assert len(recs) == 1
    rec = recs[0]
    assert rec["shape_id"] == str(bad.shape_id)
    assert rec["severity"] == "error"
    assert int(rec["new_value"]) == T
    assert is_fixable(rec)

    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {rec["record_id"]})
    assert result.applied == 1
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    by_id = {s.shape_id: s for s in out.shapes}
    assert by_id[bad.shape_id].top == T
    # the contents traveled with the panel: composition preserved
    assert by_id[inner.shape_id].top == T + 150000


def test_panel_row_suppresses_pershape_duplicate(make_prs, en_profile,
                                                 tmp_path):
    """A small drift sits inside BOTH the panel window and the per-shape
    intent window; the shape must get exactly one record (the panel one)."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    T = IN
    _panel(slide, IN, T)
    _panel(slide, int(4.5 * IN), T)
    bad, _ = _panel(slide, IN * 8, T + 100000)  # 2.8mm: both windows see it
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    all_recs = [r.to_dict() for r in ma.detect(ctx)]
    mine = [r for r in all_recs
            if r["shape_id"] == str(bad.shape_id)
            and r["property"] == "spPr.xfrm.off.y"]
    assert len(mine) == 1
    assert mine[0]["issue_type"] == "margin_alignment.panel_row_misaligned"


def test_panel_row_guards(make_prs, en_profile, tmp_path):
    prs = make_prs()

    # two panels only: no row evidence
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    _panel(s1, IN, IN)
    _panel(s1, IN * 8, IN + 228600)

    # deviation beyond the panel window: assumed intentional
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    _panel(s2, IN, IN)
    _panel(s2, int(4.5 * IN), IN)
    _panel(s2, IN * 8, IN + 700000)

    # empty rectangles contain nothing: not panels
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    from pptx.enum.shapes import MSO_SHAPE
    for i, top in enumerate((IN, IN, IN + 228600)):
        s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(IN + i * 4 * IN),
                            Emu(top), Emu(3 * IN), Emu(2 * IN))

    # different height classes: not one row
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    _panel(s4, IN, IN, h=2 * IN)
    _panel(s4, int(4.5 * IN), IN, h=2 * IN)
    _panel(s4, IN * 8, IN + 228600, h=4 * IN)

    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert _rows(ctx) == []


def test_bottom_line_alignment_when_bottoms_agree(make_prs, en_profile,
                                                  tmp_path):
    """Different heights but a shared bottom line: the row aligns on
    bottoms, and the target preserves the deviant's height."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bottom = 4 * IN
    _panel(slide, IN, bottom - 2 * IN, h=2 * IN)
    _panel(slide, int(4.5 * IN), bottom - int(2.2 * IN), h=int(2.2 * IN))
    bad, _ = _panel(slide, IN * 8, bottom - 2 * IN + 228600, h=2 * IN)
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = _rows(ctx)
    assert len(recs) == 1
    assert recs[0]["shape_id"] == str(bad.shape_id)
    assert int(recs[0]["new_value"]) == bottom - 2 * IN


def test_inset_fix_then_panel_move_converge(make_prs, en_profile, tmp_path):
    """The slide-5 lesson as a regression test: band 1 sits off the band
    row AND its icon has the wrong inset. Applying both fixes must land the
    bands on one row and the icons on one row. This only works because
    contained-relative fixes apply before panel moves."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    from pptx.enum.shapes import MSO_SHAPE

    T = 2 * IN
    inset = 52636
    band_w, band_h = int(2.5 * IN), 506557
    icons = []
    for i, band_top in enumerate((T - 40986, T, T)):
        left = IN + i * 3 * IN
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left),
                                      Emu(band_top), Emu(band_w), Emu(band_h))
        band.text_frame.text = f"band {i}"
        icon = slide.shapes.add_textbox(Emu(left + 90000),
                                        Emu(T + inset),  # one ABSOLUTE row
                                        Emu(384619), Emu(401284))
        icon.text_frame.text = "i"
        icons.append(icon)
    # icons share an absolute row, so icon 0's INSET is off by 40986;
    # band 0's TOP is off the band row by the same amount.
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = [r.to_dict() for r in ma.detect(ctx)]
    panel = [r for r in recs
             if r["issue_type"] == "margin_alignment.panel_row_misaligned"]
    inset_recs = [r for r in recs
                  if r["issue_type"] == "margin_alignment.edge_misaligned"
                  and r["property"] == "spPr.xfrm.off.y"
                  and r["shape_id"] == str(icons[0].shape_id)]
    assert len(panel) == 1 and int(panel[0]["new_value"]) == T
    assert len(inset_recs) == 1

    both = panel + inset_recs
    result = apply_fixes(ctx.deck_path.read_bytes(), both,
                         {r["record_id"] for r in both})
    assert result.applied == 2
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    by_id = {s.shape_id: s for s in out.shapes}
    band_tops = {by_id[r].top for r in
                 [s.shape_id for s in out.shapes
                  if getattr(s, "has_text_frame", False)
                  and s.text_frame.text.startswith("band")]}
    icon_tops = {by_id[i.shape_id].top for i in icons}
    assert band_tops == {T}, band_tops
    assert icon_tops == {T + inset}, icon_tops
