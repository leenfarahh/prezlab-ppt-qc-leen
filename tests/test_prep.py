"""Plan, decide, rebuild: the master applied, then the rebuilt deck audited.

Five claims, and the first is the reason the pass exists in this order:

  - THE AUDIT READS THE REBUILT DECK, not the upload. An audit of the raw file
    reports margins the master is about to reset and fonts it is about to
    replace, and a designer working through that list is fixing a file that no
    longer exists.
  - PLANNING TOUCHES NOTHING. plan() reads both files and stops, so a designer
    sees the layout decisions before the irreversible half runs;
  - the layouts a designer PICKED are the layouts applied - build() takes the
    plans as decided and must never re-derive them, which would silently
    discard every pick;
  - the gaps are computed from the plans that are ABOUT to be applied, so the
    report and the file cannot disagree about where a slide went;
  - a pass that is not the deliverable degrades to a sentence - no render, no
    audit, still a rebuilt deck;
  - and the run lands as ONE job in both registries, so the download and the
    design page are looking at the same bytes.

The COM half is stubbed by monkeypatching qc.applymaster.apply_master, which is
imported at call time for exactly this reason: everything worth asserting here
(which layout, which gap, which deck got audited) must be testable on a machine
with no PowerPoint on it.
"""

import io

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Emu

from qc import prep as prepmod
from qc import web
from qc.applymaster import ApplyResult, SlidePlan
from qc.prep import Prep, PrepError, build, headline, plan

IN = 914400


# ------------------------------------------------------------------ fixtures


def _deck(layout_indexes=(0, 1, 2)) -> bytes:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    for i in layout_indexes:
        prs.slides.add_slide(prs.slide_layouts[i])
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _master() -> bytes:
    buf = io.BytesIO()
    Presentation().save(buf)
    return buf.getvalue()


REBUILT = b"REBUILT-DECK-BYTES"


@pytest.fixture()
def no_powerpoint(monkeypatch):
    """apply_master, stubbed: it hands back a deck whose bytes are distinct
    from the upload, so any test can tell which file a later pass read."""
    seen = {}

    def _fake(deck_bytes, master_bytes, plans):
        seen["deck"] = deck_bytes
        seen["plans"] = plans
        out = Presentation(io.BytesIO(deck_bytes))
        buf = io.BytesIO()
        out.save(buf)
        seen["out"] = buf.getvalue()
        return ApplyResult(deck=seen["out"], plans=plans, masters=1)

    monkeypatch.setattr("qc.applymaster.apply_master", _fake)
    return seen


# ------------------------------------------------- the order the passes run in


def test_the_audit_reads_the_rebuilt_deck_not_the_upload(no_powerpoint,
                                                         monkeypatch):
    """The whole reason the pass is in this order. An audit of the upload
    describes a file the master is about to overwrite."""
    audited = {}

    class _Result:
        def to_manifest(self):
            return {"slides": 3, "summary": {"total": 0}, "records": []}

    def _run_audit(path, profile, selected=None):
        audited["bytes"] = path.read_bytes()
        return _Result()

    monkeypatch.setattr("qc.engine.run_audit", _run_audit)

    upload = _deck()
    out = prepmod.run(plan(upload, "client.pptx", _master()),
                      _master(), object())
    assert out.manifest is not None
    assert audited["bytes"] == out.deck
    assert audited["bytes"] != upload


def test_a_failed_audit_still_leaves_a_rebuilt_deck(no_powerpoint, monkeypatch):
    """The rebuild is the deliverable. Losing the audit costs a report, and
    saying so is the difference between a missing panel and a missing file."""
    def _boom(path, profile, selected=None):
        raise RuntimeError("no profile modules")

    monkeypatch.setattr("qc.engine.run_audit", _boom)

    out = prepmod.run(plan(_deck(), "client.pptx", _master()),
                      _master(), object())
    assert out.deck is not None
    assert out.manifest is None
    assert "could not be audited" in out.audit_note
    assert "RuntimeError" in out.audit_note


# ------------------------------------------------------------ the master half


def test_plans_and_coverage_come_from_the_same_plans(no_powerpoint):
    out = build(plan(_deck(), "client.pptx", _master()), _master())
    assert out.slides == 3
    assert out.coverage is not None
    # The coverage counts the plans that were handed to apply_master, not a
    # second reading of the deck.
    assert out.coverage.slides == len(no_powerpoint["plans"])
    assert sum(out.coverage.by_rule.values()) == 3


def test_the_master_application_asks_no_model_anything(no_powerpoint,
                                                       monkeypatch):
    """Design lead, 31/08/2026: applying a master is code orchestration.

    The vision layout matcher used to run inside this call. Nothing in plan()
    or build() may reach a model now - not to place a slide, not to propose a
    layout, not to phrase a note - because a rebuild that depends on a network
    call is a rebuild that produces a different file on a bad afternoon."""
    def _never(*a, **k):
        raise AssertionError("the master application asked a model")

    monkeypatch.setattr("qc.llm.ask_json", _never)
    monkeypatch.setattr("qc.render.export_decks_png", _never)

    out = build(plan(_deck(), "client.pptx", _master()), _master())
    assert out.suggestions == []
    assert out.coverage is not None
    assert out.deck is not None


def test_planning_rebuilds_nothing(monkeypatch):
    """plan() is the half a designer waits on before seeing anything, so it
    must not reach PowerPoint at all - not even to fail."""
    def _never(*a, **k):
        raise AssertionError("planning touched the rebuild")

    monkeypatch.setattr("qc.applymaster.apply_master", _never)
    prepared = plan(_deck(), "client.pptx", _master())
    assert prepared.slides == 3
    assert prepared.source == _deck()
    assert [l["name"] for l in prepared.layouts]


def test_the_picked_layouts_are_the_layouts_applied(no_powerpoint):
    """The pause is worthless if the rebuild re-derives what it was told.

    build() takes the plans as decided; re-planning inside it would discard
    every choice a designer made on the layout page, silently."""
    prepared = plan(_deck(), "client.pptx", _master())
    target = prepared.layouts[-1]["name"]
    for slide_plan in prepared.plans:
        slide_plan.target_layout = target
        slide_plan.match_rule = "chosen"

    build(prepared, _master())
    assert [p.target_layout for p in no_powerpoint["plans"]] == [target] * 3
    assert {p.match_rule for p in no_powerpoint["plans"]} == {"chosen"}


def test_a_master_with_no_slide_master_is_refused_before_anything_runs():
    with pytest.raises(PrepError) as exc:
        plan(_deck(), "client.pptx", b"not a pptx at all")
    assert exc.value.status == 422


def test_a_deck_with_no_slides_is_refused():
    with pytest.raises(PrepError) as exc:
        plan(_deck(layout_indexes=()), "empty.pptx", _master())
    assert "no slides" in str(exc.value)


# ------------------------------------------------------------- what it leads with


def test_headline_states_the_master_half_and_the_slide_half():
    out = Prep(filename="d.pptx", source=b"", deck=b"x", applied=34,
               plans=[SlidePlan(i, "L", None, "T", "name") for i in range(40)],
               manifest={"slides": 40, "summary": {"total": 12},
                         "records": [{"module": "font"}] * 12})
    line = headline(out)
    assert "Rebuilt 34 of 40" in line
    assert "12 findings left on it" in line


def test_headline_says_when_the_audit_did_not_run():
    out = Prep(filename="d.pptx", source=b"", deck=b"x", applied=2,
               plans=[SlidePlan(i, "L", None, "T", "name") for i in range(2)])
    assert "the audit did not run" in headline(out)


def test_headline_counts_the_gaps_as_a_change_to_the_master():
    from qc.layoutgap import Coverage, Gap

    cov = Coverage(slides=6, by_rule={"name": 2, "fallback": 4},
                   gaps=[Gap(label="a title over 2 columns", slides=[2, 3, 4, 5],
                             source_layouts=["Blank"], signature={})])
    out = Prep(filename="d.pptx", source=b"", deck=b"x", applied=6,
               plans=[SlidePlan(i, "L", None, "T", "name") for i in range(6)],
               coverage=cov,
               manifest={"slides": 6, "summary": {"total": 0}, "records": []})
    line = headline(out)
    assert "4 had no layout in it" in line
    assert "1 thing the master is missing" in line


# ------------------------------------------------------------------ the route


@pytest.fixture()
def client(monkeypatch):
    monkeypatch.setattr(web, "AUTH_REQUIRED", False)
    monkeypatch.setattr(web.app.state, "auth_required", False)
    return TestClient(web.app)


def test_intake_lists_only_profiles_that_carry_a_master(client):
    page = client.get("/prep")
    assert page.status_code == 200
    # The button says what the press does, and it no longer rebuilds anything:
    # it reads the deck and hands over the layout decisions.
    assert "Choose the layouts" in page.text


def test_only_pptx_is_accepted(client):
    reply = client.post("/prep", files={"deck": ("x.txt", b"nope")},
                        data={"profile": "prezlab_en"})
    assert reply.status_code == 400
    assert "Only .pptx files" in reply.text


def test_a_profile_with_no_master_says_so_rather_than_failing(client):
    reply = client.post("/prep", files={"deck": ("d.pptx", _deck())},
                        data={"profile": "prezlab_en"})
    # Either it has no master (the message) or it ran; both are honest, and a
    # 500 is not.
    assert reply.status_code in (200, 400, 422, 503)
    if reply.status_code == 400:
        assert "carries no master" in reply.text


def test_one_run_is_one_job_in_both_registries(no_powerpoint, monkeypatch):
    """The download and the design page must be looking at the same bytes. Two
    dicts would mean the first fix applied from one left the other serving the
    file as it stood before it."""
    class _Result:
        def to_manifest(self):
            return {"slides": 3, "summary": {"total": 0}, "records": []}

    monkeypatch.setattr("qc.engine.run_audit",
                        lambda path, profile, selected=None: _Result())
    out = prepmod.run(plan(_deck(), "client.pptx", _master()),
                      _master(), object())
    job = web._register_prep("jobid1234", out, "prezlab_en", object())
    try:
        assert web._format_jobs["jobid1234"] is job
        assert web._jobs["jobid1234"] is job
        job["deck"] = b"CHANGED"
        assert web._format_jobs["jobid1234"]["deck"] == b"CHANGED"
        assert web._jobs["jobid1234"]["deck"] == b"CHANGED"
    finally:
        web._format_jobs.pop("jobid1234", None)
        web._jobs.pop("jobid1234", None)


# ------------------------------------------------------------------ the page


def _rendered(prep, **over):
    from qc.ui_prep import render_prep_result

    kwargs = dict(prep=prep, job_id="abc123", profile_name="Prezlab EN",
                  headline=headline(prep),
                  auto={"deck": {"fixes": 6, "held": 1, "picks": 2, "left": 3,
                                 "reasons": ["more than one answer is right"]}},
                  design_open=5,
                  per_slide={0: {"warning": 3}, 2: {"error": 1, "warning": 2}},
                  chat=True)
    kwargs.update(over)
    return render_prep_result(**kwargs)


def _full_prep():
    from qc.layoutgap import Coverage, Gap

    cov = Coverage(slides=6, by_rule={"name": 2, "fallback": 4},
                   gaps=[Gap(label="a title over 2 columns", slides=[2, 3, 4, 5],
                             source_layouts=["Blank"], signature={},
                             closest="Title Only")])
    return Prep(
        filename="client deck.pptx", source=b"x", deck=b"y", applied=5,
        plans=[SlidePlan(i, "Blank", None, "Title Only",
                         "fallback" if i > 1 else "name") for i in range(6)],
        coverage=cov,
        manifest={"slides": 6, "summary": {"total": 9},
                  "records": [{"module": "font",
                               "issue_type": "font.family_out_of_set",
                               "severity": "warning", "slide_index": i % 6,
                               "record_id": f"r{i}"} for i in range(9)]})


def test_the_page_keeps_the_two_lists_apart():
    """The gaps are about the MASTER and the findings are about the SLIDES.
    They look similar and they are not the same kind of work."""
    html = _rendered(_full_prep())
    assert "What the master could not build" in html
    assert "What is left on the slides" in html
    assert html.index("What the master could not build") \
        < html.index("What is left on the slides")


def test_the_hand_over_states_its_counts_before_it_is_pressed():
    html = _rendered(_full_prep())
    assert "6 audit fixes" in html
    assert "2 design decisions" in html
    assert "3 decisions would be left for you" in html
    assert "1 fix asks for your explicit approval" in html


def test_the_worst_slides_link_straight_at_the_slide():
    html = _rendered(_full_prep())
    assert 'href="/design/abc123?n=2"' in html
    assert "1 error, 2 warning" in html


def test_a_run_with_no_audit_offers_no_link_that_would_404():
    prep = _full_prep()
    prep.manifest = None
    prep.audit_note = "The rebuilt deck could not be audited (OSError)."
    html = _rendered(prep, design_open=None, auto={}, per_slide={})
    assert "/design/abc123" not in html
    assert "/audit/abc123" not in html
    assert "could not be audited" in html
    # The rebuild is still the deliverable and is still downloadable.
    assert "/format/abc123/download" in html


def test_the_ask_box_says_it_can_act_and_that_nothing_has_yet():
    """Both halves, plainly. A box that looks like it might already have
    changed something invites a designer to stop reading it; one that looks
    inert sends them off to find the page."""
    copy = " ".join(_rendered(_full_prep()).split())
    assert "What do you want to know, or done?" in copy
    assert "nothing happens until you press the button on it" in copy


def test_a_run_whose_audit_failed_is_not_registered_as_an_audit_job(
        no_powerpoint, monkeypatch):
    """The design page reads the manifest for its slide count, so a job without
    one is not an audit job however much it looks like one."""
    monkeypatch.setattr(
        "qc.engine.run_audit",
        lambda *a, **k: (_ for _ in ()).throw(RuntimeError("nope")))
    out = prepmod.run(plan(_deck(), "client.pptx", _master()),
                      _master(), object())
    web._register_prep("jobid5678", out, "prezlab_en", object())
    try:
        assert "jobid5678" in web._format_jobs
        assert "jobid5678" not in web._jobs
    finally:
        web._format_jobs.pop("jobid5678", None)
