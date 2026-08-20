"""Writing the presentation space into a formatted deck.

Reading the marker is tested in test_stylespec.py. This is the other half: the
OUTPUT deck has to carry it on every slide master, because a marker the designer
drew on one layout serves that layout only - and a slide that could not be
rebuilt keeps the deck's ORIGINAL design alive, which never had one at all.
"""

import io
import itertools

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from qc.pspace import PS_ALT, ensure_presentation_space, frame_in
from qc.stylespec import dominant_master, infer_grid, read_presentation_space

IN = 914400
_SHAPE_IDS = itertools.count(1400)


def _bytes(prs) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _draw_marker(container, left_in, top_in, w_in, h_in, alt=PS_ALT,
                 name="Rectangle 2"):
    """The rectangle the add-in leaves behind: PowerPoint's own shape name, the
    marker in the alt text. python-pptx cannot add shapes to a master or a
    layout, so the element goes in directly."""
    from pptx.oxml.shapes.autoshape import CT_Shape

    sp = CT_Shape.new_autoshape_sp(next(_SHAPE_IDS), name, "rect",
                                   int(left_in * IN), int(top_in * IN),
                                   int(w_in * IN), int(h_in * IN))
    container.shapes._spTree.insert_element_before(sp, "p:extLst")
    shape = next(s for s in container.shapes if s._element is sp)
    if alt:
        sp.find(qn("p:nvSpPr")).find(qn("p:cNvPr")).set("descr", alt)
    shape.fill.background()
    shape.line.fill.background()
    return shape


def _deck_with_marker_on_a_layout():
    """The shape of the real client master: one layout carries the frame, the
    master itself does not, and slides sit on several layouts."""
    prs = Presentation()
    _draw_marker(prs.slide_masters[0].slide_layouts[1], 0.5, 2.0, 9.0, 4.0)
    for i in (0, 1, 2):
        prs.slides.add_slide(prs.slide_layouts[i])
    return prs


# --------------------------------------------------------------- reading it


def test_the_frame_a_deck_states_is_found_on_a_layout():
    box, where = frame_in(_bytes(_deck_with_marker_on_a_layout()))
    assert box == [int(0.5 * IN), 2 * IN, int(9.5 * IN), 6 * IN]
    assert "layout" in where


def test_a_deck_with_no_marker_states_no_frame():
    assert frame_in(_bytes(Presentation())) == (None, None)


# --------------------------------------------------------------- writing it


def test_the_marker_lands_on_the_master_so_every_layout_inherits_it():
    out, notes = ensure_presentation_space(_bytes(_deck_with_marker_on_a_layout()))

    prs = Presentation(io.BytesIO(out))
    master = dominant_master(prs)
    space = read_presentation_space(prs, master)
    assert space["source"] == "master", "the master itself now states the frame"
    assert space["marker"] == "alt_text"
    assert space["alt_text"] == PS_ALT
    assert space["box_emu"] == [int(0.5 * IN), 2 * IN, int(9.5 * IN), 6 * IN]
    assert any("added to 1 slide master" in n for n in notes)


def test_the_marker_written_is_invisible():
    """A marker that prints appears on every slide of the delivered deck."""
    out, _notes = ensure_presentation_space(_bytes(_deck_with_marker_on_a_layout()))

    prs = Presentation(io.BytesIO(out))
    space = read_presentation_space(prs, dominant_master(prs))
    assert space["prints"] is False


def test_a_master_that_already_states_the_frame_is_left_alone():
    prs = Presentation()
    _draw_marker(prs.slide_masters[0], 0.5, 2.0, 9.0, 4.0)
    raw = _bytes(prs)

    out, notes = ensure_presentation_space(raw)
    assert out == raw, "nothing to add, so nothing was rewritten"
    assert any("already carries" in n for n in notes)


def test_a_deck_with_nothing_to_read_and_no_fallback_is_untouched():
    raw = _bytes(Presentation())
    out, notes = ensure_presentation_space(raw)
    assert out == raw
    assert any("states no presentation space" in n for n in notes)


def test_the_master_file_is_the_fallback_when_the_output_states_nothing():
    raw = _bytes(Presentation())
    box = [int(0.5 * IN), 2 * IN, int(9.5 * IN), 6 * IN]

    out, notes = ensure_presentation_space(raw, fallback_box=box,
                                           fallback_size=(10 * IN, int(7.5 * IN)))
    prs = Presentation(io.BytesIO(out))
    assert read_presentation_space(prs, dominant_master(prs))["box_emu"] == box
    assert not any("scaled" in n for n in notes)


def test_a_fallback_from_another_slide_size_is_scaled_and_said_so():
    """PowerPoint resizes a loaded design to the deck's slide size. A frame
    taken from the master file at a different size would otherwise be planted
    at 16:9 numbers on a 4:3 deck and called the frame."""
    raw = _bytes(Presentation())            # 10 x 7.5in
    box = [int(0.5 * IN), 2 * IN, int(12.5 * IN), 6 * IN]

    out, notes = ensure_presentation_space(
        raw, fallback_box=box, fallback_size=(int(13.333 * IN), int(7.5 * IN)))
    prs = Presentation(io.BytesIO(out))
    got = read_presentation_space(prs, dominant_master(prs))["box_emu"]
    assert got[2] < 10 * IN, "the frame has to fit the deck it is written into"
    assert any("scaled" in n for n in notes)


def test_the_written_marker_reads_back_as_a_body_frame():
    """End to end: what is written is what the next pass reads, including the
    page-margin-versus-body-ceiling distinction the frame's top depends on."""
    out, _notes = ensure_presentation_space(_bytes(_deck_with_marker_on_a_layout()))

    prs = Presentation(io.BytesIO(out))
    grid = infer_grid(prs, dominant_master(prs))
    assert grid["source"] == "presentation_space"
    assert grid["space_states"] == "body", "it starts below the title box"
    assert grid["body_top_emu"] == 2 * IN
    assert grid["margins_emu"]["top"] < 2 * IN, "the page top is not the body top"


# ------------------------------------------------------------------- Arabic


def test_the_frame_is_written_the_same_way_for_an_arabic_deck():
    """The marker is a rectangle, so it is direction-neutral by construction -
    and that is worth pinning: an RTL deck seats its body on the RIGHT edge of
    this same frame, so a frame that came out mirrored would move every Arabic
    slide."""
    prs = Presentation()
    _draw_marker(prs.slide_masters[0].slide_layouts[1], 0.5, 2.0, 9.0, 4.0)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1))
    box.text_frame.text = "مرحبا بكم في العرض"
    for para in box.text_frame.paragraphs:
        para._p.get_or_add_pPr().set("rtl", "1")

    out, _notes = ensure_presentation_space(_bytes(prs))
    prs2 = Presentation(io.BytesIO(out))
    space = read_presentation_space(prs2, dominant_master(prs2))
    assert space["box_emu"] == [int(0.5 * IN), 2 * IN, int(9.5 * IN), 6 * IN]
    # The Arabic run survived the rewrite intact.
    text = "".join(sh.text_frame.text for sh in prs2.slides[0].shapes
                   if sh.has_text_frame)
    assert "مرحبا بكم في العرض" in text
