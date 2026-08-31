"""The review page for a format run, and taking one change back.

Two things are being protected here and they are different. The PAGE must work
on any machine: rendering needs desktop PowerPoint or LibreOffice, neither is
guaranteed, and a designer who cannot see the pictures must still be able to
read the change list and press Undo. The UNDO must be exact: it replays the
state qc.migrate stored, so a shape comes back at the coordinates it held, not
near them.
"""

import io
from pathlib import Path

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Emu, Pt

from qc import web
from qc.migrate import ContentChange, migrate_deck
from qc.render import layout_catalogue
from qc.undo import apply_undo, expand

IN = 914400
_TITLE_BOX = (0.48, 0.42, 12.40, 0.92)
_SUBTITLE_BOX = (0.48, 1.40, 12.40, 0.35)


def _bytes(prs) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _client(monkeypatch):
    monkeypatch.setattr(web, "AUTH_REQUIRED", False)
    web.app.state.auth_required = False
    return TestClient(web.app)


def _deck() -> bytes:
    """A slide shaped like the decks this pass exists for: a header band the
    master owns, an eyebrow with no slot to go to, and a content block sitting
    below where the master says the body begins."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        ph.text_frame.clear()
        kind = str(ph.placeholder_format.type)
        box = _TITLE_BOX if kind.startswith(("TITLE", "CENTER_TITLE")) else (
            _SUBTITLE_BOX if kind.startswith("SUBTITLE") else None)
        if box:
            ph.left, ph.top, ph.width, ph.height = (
                Emu(int(v * IN)) for v in box)

    def tb(x, y, w, h, text, size):
        shape = slide.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                                         Emu(int(w * IN)), Emu(int(h * IN)))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(size)

    tb(0.6, 0.30, 3.0, 0.28, "FUTURE WORK", 11)
    tb(0.6, 0.55, 9.0, 0.60, "Next: Personalized Daily Digests", 28)
    tb(0.6, 1.25, 10.0, 0.40, "A second LLM integration.", 13)
    for i in range(3):
        tb(0.6 + i * 4.1, 3.30, 3.8, 1.9, f"Card {i + 1}", 12)
    return _bytes(prs)


def _offsets(deck_bytes):
    """{text: (left, top)} for every shape carrying text, placeholders too."""
    slide = Presentation(io.BytesIO(deck_bytes)).slides[0]
    out = {}
    for s in slide.shapes:
        if s.has_text_frame and s.text_frame.text.strip():
            out[s.text_frame.text.strip()] = (s.left, s.top)
    return out


def _migrated_with_removals(source):
    """The migration with its removals performed.

    Removal is opt-in (design lead, 26/08/2026: nothing leaves a slide unless a
    designer asks), so a default run only proposes. The tests below are about
    undoing a removal, which needs one to exist.
    """
    return migrate_deck(source, remove=True)


def _job(monkeypatch, job_id="reviewjob", *,
         remove: bool = False):
    """A finished prepare job over _deck(), with the source kept as the page
    keeps it.

    It carries a Prep because every format job is a prepared deck: the remove
    and restore routes answer on the prepared-deck page, which is drawn from
    the run rather than from the job's loose keys."""
    from qc.prep import Prep

    source = _deck()
    deck, changes = migrate_deck(source, remove=remove)
    prep = Prep(filename="d.pptx", source=source, deck=deck, applied=1,
                changes=changes)
    web._format_jobs[job_id] = {
        "deck": deck, "source": source, "filename": "d.pptx",
        "profile": "prezlab_en", "plans": [], "errors": {}, "applied": 1,
        "changes": changes, "restored": [], "undone": [], "undo_notes": {},
        "removed": [], "prep": prep, "manifest": None,
    }
    return job_id, source, changes


# ------------------------------------------------------------- the undo itself


def test_undoing_the_block_move_puts_every_shape_back_exactly():
    """The whole point of storing state rather than a delta: the shapes land on
    the coordinates they held, to the EMU, not near them."""
    source = _deck()
    deck, changes = migrate_deck(source)
    move = next(c for c in changes if c.action == "content block moved")
    assert move.undo, "the move must carry the coordinates it moved from"

    before = _offsets(source)
    after = _offsets(deck)
    assert after["Card 1"] != before["Card 1"], "the fixture must have moved"

    undone, outcomes = apply_undo(deck, [{"change_id": move.change_id,
                                          "slide_index": move.slide_index,
                                          "ops": move.undo}])
    assert outcomes[0]["done"]
    back = _offsets(undone)
    for card in ("Card 1", "Card 2", "Card 3"):
        assert back[card] == before[card]


def test_undoing_a_removal_brings_the_text_back():
    source = _deck()
    deck, changes = _migrated_with_removals(source)
    swept = next(c for c in changes if c.action == "removed unplaced text")
    assert "FUTURE WORK" not in _offsets(deck)

    undone, outcomes = apply_undo(deck, [{"change_id": swept.change_id,
                                          "slide_index": swept.slide_index,
                                          "ops": swept.undo}])
    assert outcomes[0]["done"]
    assert "FUTURE WORK" in _offsets(undone)
    assert _offsets(undone)["FUTURE WORK"] == _offsets(source)["FUTURE WORK"]


def test_undoing_a_placeholder_fill_empties_it_and_returns_the_source():
    """Both halves or neither: undoing the fill without returning the shape
    loses the wording, and returning the shape without emptying the
    placeholder prints it twice."""
    source = _deck()
    deck, changes = migrate_deck(source)
    fill = next(c for c in changes if c.action == "title into placeholder")
    heading = "Next: Personalized Daily Digests"

    slide = Presentation(io.BytesIO(deck)).slides[0]
    assert any(s.is_placeholder and s.text_frame.text.strip() == heading
               for s in slide.shapes)

    undone, outcomes = apply_undo(deck, [{"change_id": fill.change_id,
                                          "slide_index": fill.slide_index,
                                          "ops": fill.undo}])
    assert outcomes[0]["done"]
    slide = Presentation(io.BytesIO(undone)).slides[0]
    holders = [s for s in slide.shapes
               if s.has_text_frame and s.text_frame.text.strip() == heading]
    assert len(holders) == 1, "the heading must exist exactly once"
    assert not holders[0].is_placeholder, "it belongs back in its own box"


def test_undoing_a_removal_also_takes_back_the_move_it_enabled():
    """The pass swept a shape and THEN seated the block, so the block's position
    is a consequence of the sweep. Putting the shape back on its own returned it
    to its original coordinates while everything else stayed where the sweep let
    it go, and the returned text landed on top of the content instead of beside
    it (design lead, 23/08/2026, "the kept text should also go back to its place
    instead of overlapping them")."""
    source = _deck()
    deck, changes = _migrated_with_removals(source)
    swept = next(c for c in changes if c.action == "removed unplaced text")
    move = next(c for c in changes if c.action == "content block moved")

    pulled = expand(changes, [swept.change_id])
    assert move.change_id in [c.change_id for c in pulled], \
        "the move has to come back with the removal that preceded it"

    items = [{"change_id": c.change_id, "slide_index": c.slide_index,
              "ops": c.undo} for c in pulled]
    undone, _outcomes = apply_undo(deck, items)

    before, after = _offsets(source), _offsets(undone)
    for text in ("FUTURE WORK", "Card 1", "Card 2", "Card 3"):
        assert after[text] == before[text], f"{text} is not back where it was"


def test_undoing_a_later_change_leaves_the_earlier_ones_alone():
    """Peeling back from the end, not resetting the slide: undoing the LAST
    change must not undo the sweep that ran before it."""
    source = _deck()
    _deck_out, changes = migrate_deck(source)
    move = next(c for c in changes if c.action == "content block moved")

    pulled = [c.action for c in expand(changes, [move.change_id])]
    assert "removed unplaced text" not in pulled
    assert pulled == ["content block moved"]


def test_a_change_on_another_slide_is_never_dragged_in():
    """Slides are independent: what happened on one says nothing about the
    other, and taking a whole deck back is what the original upload is for."""
    changes = [
        ContentChange(0, "removed unplaced text", "x", change_id="c0",
                      undo=[{"op": "insert", "xml": "<x/>"}]),
        ContentChange(1, "content block moved", "y", change_id="c1",
                      undo=[{"op": "offset", "shape_id": "1",
                             "left": 0, "top": 0}]),
        ContentChange(0, "content block moved", "z", change_id="c2",
                      undo=[{"op": "offset", "shape_id": "2",
                             "left": 0, "top": 0}]),
    ]
    assert [c.change_id for c in expand(changes, ["c0"])] == ["c0", "c2"]


def test_a_report_only_change_is_not_dragged_in_and_does_not_block():
    """A row with nothing to undo is skipped rather than counted: it would
    otherwise show up as a change that came back without changing anything."""
    changes = [
        ContentChange(0, "removed unplaced text", "x", change_id="c0",
                      undo=[{"op": "insert", "xml": "<x/>"}]),
        ContentChange(0, "content does not fit", "y", change_id="c1"),
        ContentChange(0, "content block moved", "z", change_id="c2",
                      undo=[{"op": "offset", "shape_id": "2",
                             "left": 0, "top": 0}]),
    ]
    assert [c.change_id for c in expand(changes, ["c0"])] == ["c0", "c2"]


def test_overlapping_changes_are_undone_last_first():
    """The block move and the collision nudge both touch the same shape, the
    nudge running on the position the move left it in. Replaying them front to
    back puts the shape back and then straight forward again, so the undo
    silently does nothing (real deck, 23/08/2026: 5 of 721 shapes)."""
    shape_id = "77"
    ops_move = [{"op": "offset", "shape_id": shape_id, "left": 0, "top": 0}]
    ops_nudge = [{"op": "offset", "shape_id": shape_id,
                  "left": 0, "top": 2 * IN}]

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(0), Emu(4 * IN), Emu(IN), Emu(IN))
    box._element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr"
    ).find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr"
    ).set("id", shape_id)

    out, outcomes = apply_undo(_bytes(prs), [
        {"change_id": "move", "slide_index": 0, "ops": ops_move},
        {"change_id": "nudge", "slide_index": 0, "ops": ops_nudge},
    ])
    assert [o["change_id"] for o in outcomes] == ["move", "nudge"], \
        "outcomes must come back in the caller's order"
    landed = Presentation(io.BytesIO(out)).slides[0].shapes[0].top
    assert landed == 0, "the earliest change must be the last one replayed"


def test_an_undo_that_finds_nothing_says_so_rather_than_claiming_success():
    deck = _deck()
    _out, outcomes = apply_undo(deck, [{"change_id": "c0", "slide_index": 0,
                                        "ops": [{"op": "offset",
                                                 "shape_id": "99999",
                                                 "left": 0, "top": 0}]}])
    assert outcomes[0]["done"] is False
    assert "no longer" in outcomes[0]["detail"]


def test_an_undo_for_a_slide_the_deck_no_longer_has_is_reported():
    _out, outcomes = apply_undo(_deck(), [{"change_id": "c0",
                                           "slide_index": 40, "ops": []}])
    assert outcomes[0]["done"] is False


def test_the_deck_still_opens_after_an_undo():
    """Every operation writes into a shape tree whose child order the schema
    fixes; getting it wrong is a repair prompt, not an exception."""
    source = _deck()
    deck, changes = migrate_deck(source)
    items = [{"change_id": c.change_id, "slide_index": c.slide_index,
              "ops": c.undo} for c in changes if c.undo]
    out, _outcomes = apply_undo(deck, items)
    slide = Presentation(io.BytesIO(out)).slides[0]
    ids = [s.shape_id for s in slide.shapes]
    assert len(ids) == len(set(ids)), "duplicate shape ids read as a damaged file"


# -------------------------------------------------------------- the layouts


def test_the_layout_catalogue_is_one_blank_slide_per_layout():
    """How a layout gets photographed at all: PowerPoint exports slides, so
    each layout needs an empty slide of its own to be seen on."""
    prs = Presentation()
    catalogue, entries, skipped = layout_catalogue(_bytes(prs))
    assert skipped == 0
    built = Presentation(io.BytesIO(catalogue))
    assert len(built.slides) == len(entries) == len(prs.slide_layouts)
    assert [s.slide_layout.name for s in built.slides] == \
        [e["layout"] for e in entries]
    # and none of the original deck's own slides are left to confuse the view
    with_slides = Presentation()
    with_slides.slides.add_slide(with_slides.slide_layouts[0])
    catalogue, entries, _ = layout_catalogue(_bytes(with_slides))
    assert len(Presentation(io.BytesIO(catalogue)).slides) == len(entries)


# ------------------------------------------------------------------ the page


# --------------------------------------------- surviving a dead renderer


def test_a_render_failure_does_not_lose_the_layout_list(monkeypatch):
    """Which layouts a deck arrived with and which it has now is read out of the
    file by python-pptx and needs no PowerPoint. Letting an export failure take
    the whole answer down made the page say "No layouts to show" about a master
    carrying twelve of them (design lead, 23/08/2026)."""
    from qc import render

    def boom(*_a, **_kw):
        raise RuntimeError("Presentations.Open : Failed.")

    monkeypatch.setattr(render, "export_decks_png", boom)
    out = render.layout_previews(_deck(), _deck())

    assert out["before"] and out["after"], "the entries must survive"
    assert out["images"] == {}
    assert "Failed" in out["error"]


def test_a_render_failure_does_not_lose_the_slide_list(monkeypatch):
    from qc import render

    monkeypatch.setattr(render, "export_decks_png",
                        lambda *a, **k: (_ for _ in ()).throw(RuntimeError("no")))
    out = render.slide_previews(_deck(), _deck(), [0])

    assert out["indices"] == [0]
    assert out["images"] == {}
    assert out["error"]


def test_every_picture_the_page_asks_for_is_actually_served(monkeypatch):
    """The deck view caches per page ("review_deck_0"), and a route that looked
    for the bucket by its literal name found nothing: the page rendered eighteen
    pictures and served none of them."""
    import re

    from qc import render

    fake = b"\x89PNG\r\n\x1a\nfake"
    monkeypatch.setattr(render, "export_decks_png",
                        lambda decks, idx, width=None: {
                            f"{name}:{i}": fake for name in decks for i in idx})
    client = _client(monkeypatch)
    job_id, _source, _changes = _job(monkeypatch)
    try:
        for view in ("master", "deck"):
            page = client.get(f"/format/{job_id}/review?view={view}")
            asked = re.findall(rf"/review-img/{job_id}/([A-Za-z0-9_]+)\.png",
                               page.text)
            assert asked, f"the {view} view should ask for pictures"
            for key in asked:
                got = client.get(f"/review-img/{job_id}/{key}.png")
                assert got.status_code == 200, f"{view}: {key} was not served"
                assert got.content == fake
    finally:
        web._format_jobs.pop(job_id, None)


def test_the_master_view_lists_layouts_with_no_renderer(monkeypatch):
    monkeypatch.setattr("qc.render.RENDERER", "none")
    client = _client(monkeypatch)
    job_id, _source, _changes = _job(monkeypatch)
    try:
        r = client.get(f"/format/{job_id}/review?view=master")
        assert "No layouts to show" not in r.text
        assert "Title Slide" in r.text, "the layout names come from the file"
        assert "read from the deck itself" in r.text
    finally:
        web._format_jobs.pop(job_id, None)


# ------------------------------------------- not leaking a PowerPoint


def test_a_wedged_instance_is_terminated_rather_than_left_behind(monkeypatch):
    """Quit cannot close an instance sitting on a repair prompt or halfway
    through a failed Open. The call returns, the exception is swallowed, and a
    windowless POWERPNT.EXE survives; every later attempt on the host then fails
    against it and leaks another (three in one afternoon, 23/08/2026)."""
    from qc import unify

    killed = []
    monkeypatch.setattr(unify, "automation_pids", lambda: {1, 2, 99})
    monkeypatch.setattr(unify, "_terminate", lambda pid, ms: killed.append(pid))
    monkeypatch.setattr(unify.os, "name", "nt")

    class Wedged:
        def Quit(self):
            raise OSError("wedged")

    unify.force_quit(Wedged(), started={1, 2})
    assert killed == [99], "only the instance that appeared since is ours to end"


def test_a_powerpoint_the_designer_opened_is_never_touched(monkeypatch):
    """The snapshot is what makes the kill safe. Anything already running when
    the run started is somebody else's, and it may hold unsaved work."""
    from qc import unify

    killed = []
    monkeypatch.setattr(unify, "automation_pids", lambda: {7})
    monkeypatch.setattr(unify, "_terminate", lambda pid, ms: killed.append(pid))
    monkeypatch.setattr(unify.os, "name", "nt")

    unify.force_quit(type("A", (), {"Quit": lambda self: None})(), started={7})
    assert killed == []


def test_force_quit_without_a_snapshot_only_asks_nicely(monkeypatch):
    """No snapshot means no way to tell our instance from anyone else's, so it
    must not guess."""
    from qc import unify

    killed = []
    monkeypatch.setattr(unify, "automation_pids", lambda: {1, 2, 3})
    monkeypatch.setattr(unify, "_terminate", lambda pid, ms: killed.append(pid))
    quit_called = []
    unify.force_quit(
        type("A", (), {"Quit": lambda self: quit_called.append(True)})())

    assert quit_called == [True]
    assert killed == []


def test_a_leftover_from_an_earlier_run_is_cleared_before_the_next_one(
        monkeypatch):
    """The hole the teardown snapshot cannot cover.

    An instance that was alive BEFORE a run is in `started`, so it can never be
    in the difference force_quit kills - and DispatchEx does not start a second
    PowerPoint, it attaches to whatever /AUTOMATION instance is already there
    (two DispatchEx calls, one pid, measured 24/08/2026). So a leftover is not
    something the next run works around: it inherits it, wedges, and leaves the
    leftover behind for the run after that. A real host spent a day answering
    "PowerPoint would not start" to one instance orphaned at 10:05.
    """
    from qc import unify

    alive = {45572}
    listings = []
    monkeypatch.setattr(unify.os, "name", "nt")
    monkeypatch.setattr(unify, "automation_pids",
                        lambda: listings.append(1) or set(alive))
    monkeypatch.setattr(unify, "_terminate",
                        lambda pid, ms: alive.discard(pid))

    # It hands back what SURVIVED, which is the snapshot force_quit wants: an
    # empty one here, so the instance this run goes on to create IS in that
    # difference and cannot become the next leftover.
    assert unify.sweep_automation() == set()
    assert alive == set(), "the leftover survived the sweep"
    assert len(listings) == 2, "a sweep that killed something has to re-check"


def test_the_sweep_reports_what_would_not_die(monkeypatch):
    """A leftover that refuses to close has to come back as still alive: the run
    attaches to it either way, and the teardown must not then mistake it for the
    instance it created and try to kill it twice."""
    from qc import unify

    monkeypatch.setattr(unify.os, "name", "nt")
    monkeypatch.setattr(unify, "automation_pids", lambda: {7})
    monkeypatch.setattr(unify, "_terminate", lambda pid, ms: None)  # immortal

    assert unify.sweep_automation() == {7}


def test_the_sweep_is_a_no_op_off_windows(monkeypatch):
    from qc import unify

    monkeypatch.setattr(unify.os, "name", "posix")
    monkeypatch.setattr(unify, "automation_pids",
                        lambda: (_ for _ in ()).throw(AssertionError("asked")))
    assert unify.sweep_automation() == set()


def test_the_server_execution_advice_names_the_leftover_first():
    """The message said "an update or repair prompt waiting to be answered" and
    nothing else, and sent a designer looking for a prompt that was not there
    while a leftover instance was the actual cause (24/08/2026)."""
    from qc.unify import com_failure_advice

    advice = com_failure_advice(Exception(-2146959355, "Server execution failed",
                                          None, None))
    assert "POWERPNT.EXE" in advice and "no window" in advice
    assert "repair" in advice, "the other cause must still be named"
    assert "Server execution failed" in advice, "the raw error is the bug report"


def test_the_review_page_works_without_a_renderer(monkeypatch):
    """Rendering needs PowerPoint or LibreOffice on the host. The change list
    and the Undo buttons must not."""
    monkeypatch.setattr("qc.render.RENDERER", "none")
    client = _client(monkeypatch)
    job_id, _source, _changes = _job(monkeypatch)
    try:
        r = client.get(f"/format/{job_id}/review?view=deck")
        assert r.status_code == 200
        assert "pictures are missing" in r.text
        assert "content block moved" in r.text
        assert 'name="change_ids"' in r.text
    finally:
        web._format_jobs.pop(job_id, None)


def test_both_views_are_reachable_and_named(monkeypatch):
    monkeypatch.setattr("qc.render.RENDERER", "none")
    client = _client(monkeypatch)
    job_id, _source, _changes = _job(monkeypatch)
    try:
        master = client.get(f"/format/{job_id}/review?view=master")
        deck = client.get(f"/format/{job_id}/review?view=deck")
        assert "Layouts the deck arrived with" in master.text
        assert "Layouts the deck has now" in master.text
        assert "Slide 1" in deck.text
    finally:
        web._format_jobs.pop(job_id, None)


def test_a_report_only_row_offers_no_undo_and_says_why(monkeypatch):
    monkeypatch.setattr("qc.render.RENDERER", "none")
    client = _client(monkeypatch)
    job_id = "reportonly"
    web._format_jobs[job_id] = {
        "deck": _deck(), "source": _deck(), "filename": "d.pptx",
        "profile": "prezlab_en", "plans": [], "errors": {}, "applied": 1,
        "changes": [ContentChange(0, "heading past the margin", "runs wide",
                                  severity="alert", change_id="c0")],
        "restored": [], "undone": [], "undo_notes": {},
    }
    try:
        r = client.get(f"/format/{job_id}/review?view=deck")
        assert "whether a heading may break the margin" in r.text
        assert 'value="c0"' not in r.text
    finally:
        web._format_jobs.pop(job_id, None)


def test_undo_from_the_page_changes_the_downloaded_deck(monkeypatch):
    monkeypatch.setattr("qc.render.RENDERER", "none")
    client = _client(monkeypatch)
    job_id, source, changes = _job(monkeypatch)
    move = next(c for c in changes if c.action == "content block moved")
    try:
        r = client.post(f"/format/{job_id}/undo",
                        data={"change_ids": [move.change_id]})
        assert r.status_code == 200
        assert "undone" in r.text
        out = client.get(f"/format/{job_id}/download").content
        assert _offsets(out)["Card 1"] == _offsets(source)["Card 1"]
    finally:
        web._format_jobs.pop(job_id, None)


def test_resubmitting_an_undo_does_not_apply_it_twice(monkeypatch):
    """A refresh or a back button resends the POST. Replaying an insert would
    put a second copy of the returned shape on the slide."""
    monkeypatch.setattr("qc.render.RENDERER", "none")
    client = _client(monkeypatch)
    job_id, _source, changes = _job(monkeypatch, remove=True)
    swept = next(c for c in changes if c.action == "removed unplaced text")
    try:
        client.post(f"/format/{job_id}/undo",
                    data={"change_ids": [swept.change_id]})
        client.post(f"/format/{job_id}/undo",
                    data={"change_ids": [swept.change_id]})
        slide = Presentation(io.BytesIO(
            client.get(f"/format/{job_id}/download").content)).slides[0]
        copies = [s for s in slide.shapes if s.has_text_frame
                  and s.text_frame.text.strip() == "FUTURE WORK"]
        assert len(copies) == 1
    finally:
        web._format_jobs.pop(job_id, None)


def test_undoing_nothing_leaves_the_deck_alone(monkeypatch):
    monkeypatch.setattr("qc.render.RENDERER", "none")
    client = _client(monkeypatch)
    job_id, _source, _changes = _job(monkeypatch)
    before = web._format_jobs[job_id]["deck"]
    try:
        r = client.post(f"/format/{job_id}/undo", data={})
        assert r.status_code == 200
        assert web._format_jobs[job_id]["deck"] == before
    finally:
        web._format_jobs.pop(job_id, None)


def test_the_deck_view_pages_rather_than_stopping(monkeypatch):
    """A review that ends at slide 20 of a 26-slide deck and says nothing reads
    as six slides the tool declined to show (design lead, 23/08/2026). The deck
    always has all of them; only the review is paginated."""
    monkeypatch.setattr("qc.render.RENDERER", "none")
    client = _client(monkeypatch)
    job_id = "pagedjob"
    changes = [ContentChange(i, "content block moved", f"slide {i+1}",
                             change_id=f"c{i}",
                             undo=[{"op": "offset", "shape_id": "1",
                                    "left": 0, "top": 0}])
               for i in range(26)]
    web._format_jobs[job_id] = {
        "deck": _deck(), "source": _deck(), "filename": "d.pptx",
        "profile": "prezlab_en", "errors": {}, "applied": 26,
        "plans": [type("P", (), {"slide_index": i, "source_layout": "a",
                                 "source_type": None, "target_layout": "b",
                                 "match_rule": "name", "note": ""})()
                  for i in range(26)],
        "changes": changes, "restored": [], "undone": [], "undo_notes": {},
    }
    try:
        first = client.get(f"/format/{job_id}/review?view=deck")
        assert "Slides 1&ndash;20 of 26 with changes" in first.text
        assert "page 1 of 2" in first.text
        assert "Next" in first.text
        assert "The deck has all 26 slides" in first.text

        second = client.get(f"/format/{job_id}/review?view=deck&page=1")
        assert "Slides 21&ndash;26 of 26" in second.text
        assert "Slide 26" in second.text, "the last slide must be reachable"
    finally:
        web._format_jobs.pop(job_id, None)


def test_review_and_undo_on_an_unknown_job_are_clean_404s(monkeypatch):
    client = _client(monkeypatch)
    assert client.get("/format/deadbeef/review").status_code == 404
    assert client.post("/format/deadbeef/undo", data={}).status_code == 404
    assert client.get("/review-img/deadbeef/layout_before_0.png").status_code == 404


def test_the_result_page_points_at_the_review(monkeypatch):
    """The prepared deck's page is the result page now, and before/after is one
    of the four things it hands off to."""
    from qc.prep import Prep
    from qc.ui_prep import render_prep_result

    prep = Prep(filename="d.pptx", source=b"", deck=b"", applied=1)
    html = render_prep_result(prep=prep, job_id="j1", profile_name="P",
                              headline="Rebuilt 1 of 1 slides on the master.")
    assert "/format/j1/review" in html
    assert "/format/j1/download" in html


def test_a_clean_host_costs_one_process_listing_not_two(monkeypatch):
    """Each listing is a WMI or PowerShell round trip that reads every
    process's command line: 300-500ms. Four of them per render was 2s of a
    7.7s wait spent asking Windows the same question (measured 24/08/2026)."""
    from qc import unify

    calls, killed = [], []
    monkeypatch.setattr(unify.os, "name", "nt")
    monkeypatch.setattr(unify, "automation_pids",
                        lambda: calls.append(1) or set())
    monkeypatch.setattr(unify, "_terminate",
                        lambda pid, ms: killed.append(pid))

    assert unify.sweep_automation() == set()
    assert killed == [], "nothing was running, so nothing was there to end"
    assert len(calls) == 1, f"asked {len(calls)} times with nothing to sweep"


def test_the_process_listing_prefers_wmi_but_survives_without_it(monkeypatch):
    """PowerShell costs ~510ms a call because starting PowerShell does;
    in-process WMI answers the same question in ~330ms. Either way the answer
    has to be the same, and a host where WMI is unreachable still needs one."""
    from qc import unify

    monkeypatch.setattr(unify.os, "name", "nt")
    monkeypatch.setattr(unify, "_automation_pids_wmi", lambda: {11, 22})
    assert unify.automation_pids() == {11, 22}

    monkeypatch.setattr(unify, "_automation_pids_wmi", lambda: None)
    fake = type("R", (), {"stdout": "33\n44\n"})()
    monkeypatch.setattr(unify.subprocess, "run", lambda *a, **k: fake)
    assert unify.automation_pids() == {33, 44}, "the fallback did not run"


# ------------------------------------------- the button that performs a removal
#
# Removal is opt-in as of 26/08/2026, so the deck a designer downloads still has
# everything the pass found. This is the other half of that: the tick, and what
# happens when the same POST arrives twice.


def _texts_of(deck_bytes, idx=0):
    slide = Presentation(io.BytesIO(deck_bytes)).slides[idx]
    return [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]


def test_ticking_a_proposal_takes_that_piece_out(monkeypatch):
    monkeypatch.setattr("qc.render.RENDERER", "none")
    client = _client(monkeypatch)
    job_id, _source, changes = _job(monkeypatch, "removejob")
    proposal = next(c for c in changes
                    if (c.remove_op or {}).get("kind") == "shape")
    piece = proposal.removed_text

    assert piece in _texts_of(web._format_jobs[job_id]["deck"]), \
        "the pass left it in place, which is the policy"

    r = client.post(f"/format/{job_id}/remove",
                    data={"remove_ids": proposal.remove_id})
    assert r.status_code == 200
    assert piece not in _texts_of(web._format_jobs[job_id]["deck"])
    # and it is recorded as a change, with an undo, like everything else
    performed = [c for c in web._format_jobs[job_id]["changes"]
                 if c.action == "removed on request"]
    assert len(performed) == 1 and performed[0].undo


def test_an_unticked_proposal_is_left_alone(monkeypatch):
    """One tick, one removal. A pass that took out the neighbours of what was
    ticked would be worse than one that removed everything, because nobody would
    be looking for it."""
    monkeypatch.setattr("qc.render.RENDERER", "none")
    client = _client(monkeypatch)
    job_id, _source, changes = _job(monkeypatch, "removejob2")
    proposals = [c for c in changes
                 if (c.remove_op or {}).get("kind") == "shape"]
    ticked = proposals[0]
    others = [c.removed_text for c in proposals[1:] if c.removed_text]

    r = client.post(f"/format/{job_id}/remove",
                    data={"remove_ids": ticked.remove_id})
    assert r.status_code == 200
    assert ticked.removed_text not in _texts_of(web._format_jobs[job_id]["deck"])
    left = _texts_of(web._format_jobs[job_id]["deck"])
    for text in others:
        assert text in left, f"{text!r} was not ticked and came out anyway"


def test_resubmitting_the_removal_does_not_report_a_failure(monkeypatch):
    """A refresh or a back button resends the POST. The second one must not
    complain about a piece that is correctly gone."""
    monkeypatch.setattr("qc.render.RENDERER", "none")
    client = _client(monkeypatch)
    job_id, _source, changes = _job(monkeypatch, "removejob3")
    proposal = next(c for c in changes
                    if (c.remove_op or {}).get("kind") == "shape")

    client.post(f"/format/{job_id}/remove", data={"remove_ids": proposal.remove_id})
    before = list(web._format_jobs[job_id]["changes"])
    r = client.post(f"/format/{job_id}/remove",
                    data={"remove_ids": proposal.remove_id})
    assert r.status_code == 200
    assert web._format_jobs[job_id]["changes"] == before, \
        "the second submission changed something"
    assert "removal skipped" not in r.text


def test_removing_from_an_unknown_job_is_a_clean_404(monkeypatch):
    client = _client(monkeypatch)
    r = client.post("/format/deadbeef/remove", data={"remove_ids": "r0"})
    assert r.status_code == 404


# ---------------------------------------------- the LibreOffice page split
#
# The cloud renderer converts a deck to PDF once and then rasterises pages. It
# used to launch one pdftoppm PER PAGE, each of which re-parses the whole PDF -
# and _ensure_thumbs asks for every slide of the deck, so a 200-slide deck was
# 200 process launches over the same document (30/08/2026).


def test_consecutive_pages_are_rendered_in_one_pass():
    from qc.render import _contiguous

    assert _contiguous([0, 1, 2, 3]) == [(0, 3)], "one call, not four"
    assert _contiguous(list(range(200))) == [(0, 199)]


def test_gaps_split_into_runs_and_nothing_is_lost():
    from qc.render import _contiguous

    assert _contiguous([0, 1, 5, 6, 7, 9]) == [(0, 1), (5, 7), (9, 9)]
    # Unsorted and duplicated input is what callers actually pass (a window of
    # slides plus the one being viewed).
    assert _contiguous([7, 5, 6, 5]) == [(5, 7)]
    assert _contiguous([]) == []


def test_every_wanted_page_comes_back_keyed_by_its_slide_index(tmp_path,
                                                               monkeypatch):
    """The multi-page branch matches files by sorted order, because pdftoppm
    picks its own suffix width. What must hold is that slide 5 is keyed 5."""
    import qc.render as R

    def _fake_run(args, **kwargs):
        first = int(args[args.index("-f") + 1])
        last = int(args[args.index("-l") + 1])
        stem = Path(args[-1])
        if "-singlefile" in args:
            stem.with_suffix(".png").write_bytes(b"page%d" % first)
        else:
            for n in range(first, last + 1):
                (stem.parent / f"{stem.name}-{n:03d}.png").write_bytes(
                    b"page%d" % n)
        return None

    monkeypatch.setattr(R.subprocess, "run", _fake_run)
    out = R._pages_to_png("pdftoppm", tmp_path / "d.pdf", "deck",
                          [3, 4, 5, 9], tmp_path, 1360)

    assert sorted(out) == ["deck:3", "deck:4", "deck:5", "deck:9"]
    assert out["deck:3"] == b"page4"     # 1-based page 4 is slide index 3
    assert out["deck:9"] == b"page10"


def test_a_shape_hanging_off_the_slide_is_highlighted_where_it_shows():
    """Origin and size were clamped independently - max(0, left/w) alongside
    min(1, width/w) - which is correct only while the shape is on the canvas. A
    shape hanging off the left edge had its origin pulled to 0 and kept its full
    width, so the highlight covered the wrong part of the slide: on exactly the
    shapes an audit flags for sitting outside the frame (30/08/2026)."""
    from qc.render import _fraction_box

    W = H = 1000

    # Half off the left edge: the visible half starts at 0 and is half as wide.
    box = _fraction_box(-200, 100, 400, 200, W, H)
    assert box["x"] == 0.0
    assert box["w"] == pytest.approx(0.2), (
        "the clamped origin kept the full 0.4 width and shifted the rectangle "
        "over the wrong shapes")

    # Fully on the slide: unchanged.
    assert _fraction_box(100, 200, 300, 400, W, H) == pytest.approx(
        {"x": 0.1, "y": 0.2, "w": 0.3, "h": 0.4})

    # Running off the right edge: clipped at the edge, never past 1.0.
    box = _fraction_box(800, 0, 400, 100, W, H)
    assert box["x"] == pytest.approx(0.8)
    assert box["x"] + box["w"] <= 1.0

    # Entirely off the canvas: no rectangle rather than a wrong one.
    assert _fraction_box(-500, 0, 200, 100, W, H)["w"] == 0.0
