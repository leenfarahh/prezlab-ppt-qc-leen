"""Three fixes from the 12/08/2026 RTL-deck review batch: deepened text
overlaps are collisions (slide 16's merged columns), cross-slide anchors
are pinned instead of locally nudged (titles aligned then misaligned),
and table cells get font auditing (28 tables escaped it)."""

import io

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

import qc.modules.font as font_mod
import qc.modules.margin_alignment as ma
from qc.fixer import apply_fixes, is_fixable
from qc.profile import Profile
from tests.conftest import save_and_ctx

MM = 36000

PROFILE = Profile({"id": "t", "config": {"font": {"roles": {"body": {
    "latin": ["Arial"], "complex_script": ["Sakkal Majalla"],
    "allowed_weights": ["regular", "bold"]}}}}})


def test_deepened_text_overlap_is_a_collision(make_prs, en_profile,
                                              tmp_path):
    """Two wide text boxes already overlap as rectangles; a fix pushing one
    substantially deeper into the other must be refused even though the
    overlap is not new."""
    from qc.fixer import _collision_created, _slide_rects

    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = slide.shapes.add_textbox(Emu(20 * MM), Emu(50 * MM),
                                 Emu(80 * MM), Emu(40 * MM))
    a.text_frame.text = "column one body text"
    b = slide.shapes.add_textbox(Emu(85 * MM), Emu(50 * MM),
                                 Emu(80 * MM), Emu(40 * MM))
    b.text_frame.text = "column two body text"
    # they already overlap by 15mm (RTL text boxes routinely do), so the
    # old new-overlap-only rule would wave the shove through
    before = _slide_rects(slide)
    b.left = Emu(50 * MM)  # a fix shoves column two 35mm deeper
    reason = _collision_created(slide, {str(b.shape_id)}, before)
    assert reason is not None and "deeper into" in reason


def _deck_with_anchor(make_prs, off_mm=0):
    """Four slides sharing a title bar at one spot; the last slide's copy
    sits off by off_mm."""
    prs = make_prs()
    bars = []
    for i in range(4):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = 40 * MM + (off_mm * MM if i == 3 else 0)
        bar = slide.shapes.add_textbox(Emu(left), Emu(12 * MM),
                                       Emu(250 * MM), Emu(14 * MM))
        bar.text_frame.text = f"section title {i}"
        bar.name = "TitleBar"
        bars.append(bar)
    return prs, bars


def test_recurring_anchor_straggler_is_pinned(make_prs, en_profile,
                                              tmp_path):
    prs, bars = _deck_with_anchor(make_prs, off_mm=12)
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = [r.to_dict() for r in ma.detect(ctx)
            if r.issue_type == "margin_alignment.recurring_off_position"]
    assert len(recs) == 1
    rec = recs[0]
    assert rec["slide_index"] == 3
    assert is_fixable(rec)

    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {rec["record_id"]})
    assert result.applied == 1
    out = Presentation(io.BytesIO(result.cleaned_bytes))
    fixed = next(s for s in out.slides[3].shapes if s.name == "TitleBar")
    assert fixed.left == 40 * MM


def test_anchors_never_join_local_clusters(make_prs, en_profile, tmp_path):
    """A pinned anchor must not be nudged toward slide-local clusters: no
    local positional record may target it."""
    prs, bars = _deck_with_anchor(make_prs, off_mm=0)
    slide = prs.slides[0]
    # a local column of boxes 2mm off the title bar's left edge would,
    # without the anchor exclusion, pull the title toward its cluster
    for i in range(3):
        box = slide.shapes.add_textbox(Emu(42 * MM), Emu((40 + i * 30) * MM),
                                       Emu(60 * MM), Emu(20 * MM))
        box.text_frame.text = f"box {i}"
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    bar_ids = {str(b.shape_id) for b in bars}
    local = [r for r in ma.detect(ctx)
             if r.shape_id in bar_ids
             and r.issue_type != "margin_alignment.recurring_off_position"]
    assert local == []


def test_table_cells_get_font_findings_and_fixes(make_prs, tmp_path):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    frame = slide.shapes.add_table(2, 2, Emu(30 * MM), Emu(40 * MM),
                                   Emu(150 * MM), Emu(40 * MM))
    table = frame.table
    ar = table.cell(0, 0).text_frame.paragraphs[0]
    run_ar = ar.add_run()
    run_ar.text = "دراسة الجدوى"
    run_ar.font.name = "Calibri"  # explicit latin, no cs: forgotten Arabic
    lat = table.cell(1, 1).text_frame.paragraphs[0]
    run_lat = lat.add_run()
    run_lat.text = "KPI 2031"
    run_lat.font.name = "Roboto"

    ctx = save_and_ctx(prs, tmp_path, PROFILE)
    recs = [r.to_dict() for r in font_mod.detect(ctx)
            if str(r.shape_id) == str(frame.shape_id)]
    kinds = {r["issue_type"]: r for r in recs}
    assert "font.cs_typeface_missing" in kinds
    assert "font.family_out_of_set" in kinds
    cs_rec = kinds["font.cs_typeface_missing"]
    lat_rec = kinds["font.family_out_of_set"]
    assert cs_rec["locator"] == "t0,0/p0/r0"
    assert cs_rec["new_value"] == "Sakkal Majalla"
    assert lat_rec["locator"] == "t1,1/p0/r0"
    assert lat_rec["new_value"] == "Arial"
    assert is_fixable(cs_rec) and is_fixable(lat_rec)

    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {r["record_id"] for r in recs})
    assert result.applied == 2
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    otable = next(s for s in out.shapes
                  if getattr(s, "has_table", False)).table
    o_ar = otable.cell(0, 0).text_frame.paragraphs[0].runs[0]
    cs = o_ar._r.find(qn("a:rPr")).find(qn("a:cs"))
    assert cs.get("typeface") == "Sakkal Majalla"
    assert o_ar.text == "دراسة الجدوى"
    o_lat = otable.cell(1, 1).text_frame.paragraphs[0].runs[0]
    assert o_lat.font.name == "Arial"
