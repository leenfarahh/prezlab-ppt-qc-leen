"""The before/after review shows the whole deck, not only what changed.

It used to build its slide list from the applied records, so a slide nothing
was fixed on was simply absent. A designer counting down the review found
1, 2, 3, 5 and had no way to tell "slide 4 was fine" from "slide 4 was never
looked at" - and those are opposite facts about a deck about to go to a client
(design lead, 31/08/2026).
"""

import io

from pptx import Presentation
from pptx.util import Emu, Pt

from qc.render import build_diff
from qc.ui import render_diff


def _deck(n=4) -> bytes:
    prs = Presentation()
    for i in range(n):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Emu(914400), Emu(914400),
                                       Emu(3657600), Emu(457200))
        box.text_frame.text = f"Slide {i + 1}"
        box.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _fixed_on(*slide_indexes) -> list[dict]:
    return [{"slide_index": i, "issue_type": "font.family_out_of_set",
             "record_id": f"r{i}", "shape_id": "0"} for i in slide_indexes]


def _no_rendering(monkeypatch):
    """The renderer is PowerPoint; this is about which slides are ASKED for."""
    asked = {}

    def _export(decks, indices, **kw):
        asked["indices"] = list(indices)
        return {f"{name}:{i}": b"png" for name in decks for i in indices}

    monkeypatch.setattr("qc.render.export_decks_png", _export)
    monkeypatch.setattr("qc.render.shape_rects", lambda blob, wanted: {})
    return asked


def test_a_slide_with_no_fixes_is_still_in_the_review(monkeypatch):
    asked = _no_rendering(monkeypatch)
    deck = _deck(4)

    diff = build_diff(deck, deck, _fixed_on(0, 1, 2))

    assert [sl["index"] for sl in diff["slides"]] == [0, 1, 2, 3], \
        "slide 4 was fixed on nothing and dropped out of the review"
    assert asked["indices"] == [0, 1, 2, 3], "it was never even rendered"
    assert diff["slides"][3]["changes"] == 0
    assert diff["slides"][0]["changes"] == 1


def test_the_untouched_slide_says_so_rather_than_looking_examined(monkeypatch):
    _no_rendering(monkeypatch)
    deck = _deck(4)

    html = render_diff("d.pptx", "j1", build_diff(deck, deck, _fixed_on(0)))

    assert "Slide 4" in html
    assert "no changes" in html
    assert "Nothing was applied to this slide" in html
    # and it is not dressed up as a change: no outline, and the count is honest
    assert "All 4 slides, 1 changed" in html


def test_an_untouched_slide_carries_no_highlight_rectangles(monkeypatch):
    """Outlining a shape on a slide nothing happened to would be a lie about
    what the fix pass did."""
    asked = _no_rendering(monkeypatch)
    deck = _deck(3)

    diff = build_diff(deck, deck, _fixed_on(1))

    untouched = [sl for sl in diff["slides"] if sl["changes"] == 0]
    assert len(untouched) == 2
    assert all(not sl["before_rects"] and not sl["after_rects"]
               for sl in untouched)
    assert asked["indices"] == [0, 1, 2]


def test_a_deck_that_cannot_be_reopened_still_reviews_what_was_fixed(
        monkeypatch):
    """The slide count comes from the deck. If that read fails, the review
    falls back to the records rather than collapsing to nothing."""
    _no_rendering(monkeypatch)

    diff = build_diff(b"not a pptx", b"not a pptx", _fixed_on(0, 2))

    assert [sl["index"] for sl in diff["slides"]] == [0, 1, 2]
    assert [sl["changes"] for sl in diff["slides"]] == [1, 0, 1]


def test_nothing_applied_at_all_is_still_nothing_to_review(monkeypatch):
    _no_rendering(monkeypatch)
    assert build_diff(b"junk", b"junk", []) == {"slides": [], "images": {}}
