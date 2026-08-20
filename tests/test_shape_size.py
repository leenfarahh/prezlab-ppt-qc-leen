"""Tests for the shape_size detection module (planted-violation decks)."""

import copy

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

from qc.modules.shape_size import detect
from qc.profile import Profile
from tests.conftest import save_and_ctx

DOM_W, DOM_H = 2000000, 1000000
DEV_W = 2100000  # 100000 EMU off, well beyond the 9525 tolerance
TOP = 1000000


def _add_rect(slide, left, width=DOM_W, height=DOM_H, kind=MSO_SHAPE.RECTANGLE):
    return slide.shapes.add_shape(kind, Emu(left), Emu(TOP), Emu(width), Emu(height))


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_size_mismatch_flags_only_deviant(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    for i in range(4):
        _add_rect(slide, left=500000 + i * 2200000)
    deviant = _add_rect(slide, left=500000, width=DEV_W)
    deviant_id = str(deviant.shape_id)

    ctx = save_and_ctx(prs, tmp_path, en_profile)
    records = detect(ctx)

    assert len(records) == 1
    rec = records[0]
    assert rec.issue_type == "shape_size.size_mismatch"
    assert rec.severity == "warning"
    assert rec.confidence == "high"
    assert rec.action == "flagged"
    assert rec.shape_id == deviant_id
    assert rec.old_value == f"{DEV_W}x{DOM_H}"
    assert rec.new_value == f"{DOM_W}x{DOM_H}"
    assert str(DOM_W) in rec.message and "rect" in rec.message
    assert rec.arabic_flag is False


def test_cohort_below_min_size_not_flagged(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    _add_rect(slide, left=500000, width=DOM_W)
    _add_rect(slide, left=3000000, width=DEV_W)  # differs, but cohort of 2 < 3

    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert detect(ctx) == []


def test_different_presets_never_cross_matched(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    for i in range(3):
        _add_rect(slide, left=500000 + i * 2200000)
    # An oval at a very different size must not be judged against the
    # rectangle cohort (exact prst match only in v1).
    _add_rect(slide, left=500000, width=4000000, height=3000000,
              kind=MSO_SHAPE.OVAL)

    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert detect(ctx) == []


def test_off_grid_fires_only_when_grid_enabled(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    margin_left = en_profile.get("geometry.safe_zone_margins_emu.left")
    margin_right = en_profile.get("geometry.safe_zone_margins_emu.right")
    columns = en_profile.get("geometry.grid.columns")
    col_width = (prs.slide_width - margin_left - margin_right) / columns

    on_grid = _add_rect(slide, left=margin_left + int(round(2 * col_width)))
    off_grid = _add_rect(slide, left=margin_left + int(round(2.5 * col_width)))
    off_grid_id = str(off_grid.shape_id)

    # Default profile has the grid disabled: no off_grid records at all.
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert [r for r in detect(ctx) if r.issue_type == "shape_size.off_grid"] == []

    grid_data = copy.deepcopy(en_profile.raw)
    grid_data["config"]["geometry"]["grid"]["enabled"] = True
    grid_profile = Profile(grid_data)

    ctx = save_and_ctx(prs, tmp_path, grid_profile, name="grid.pptx")
    off = [r for r in detect(ctx) if r.issue_type == "shape_size.off_grid"]
    assert len(off) == 1
    rec = off[0]
    assert rec.shape_id == off_grid_id
    assert rec.shape_id != str(on_grid.shape_id)
    assert rec.severity == "warning"
    assert rec.confidence == "low"
    assert rec.action == "flagged"
    assert rec.profile_rule_id == "geometry.grid"


def test_arabic_shape_carries_arabic_flag(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = _blank_slide(prs)
    for i in range(4):
        _add_rect(slide, left=500000 + i * 2200000)
    deviant = _add_rect(slide, left=500000, width=DEV_W)
    deviant.text_frame.text = "مرحبا"  # Arabic text

    ctx = save_and_ctx(prs, tmp_path, en_profile)
    records = detect(ctx)
    assert len(records) == 1
    assert records[0].issue_type == "shape_size.size_mismatch"
    assert records[0].arabic_flag is True
    assert records[0].action == "flagged"


def test_clean_uniform_deck_yields_zero(make_prs, en_profile, tmp_path):
    prs = make_prs()
    for _ in range(2):
        slide = _blank_slide(prs)
        for i in range(3):
            _add_rect(slide, left=500000 + i * 2200000)
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert detect(ctx) == []


def test_intentional_size_variety_not_flagged(make_prs, en_profile, tmp_path):
    """Real-deck tuning: same preset at clearly different sizes is design
    intent, not drift. Only near-duplicates are compared."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu, Inches

    from qc.modules import shape_size as mod
    from tests.conftest import save_and_ctx

    prs = make_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    for i, w in enumerate((1_000_000, 2_000_000, 4_000_000)):
        s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1 + i * 3), Inches(1),
                           Emu(w), Emu(1_000_000))
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert mod.detect(ctx) == []


def test_near_duplicates_still_flagged(make_prs, en_profile, tmp_path):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu, Inches

    from qc.modules import shape_size as mod
    from tests.conftest import save_and_ctx

    prs = make_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    for i, w in enumerate((2_000_000, 2_000_000, 2_000_000, 2_080_000)):
        s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1 + i * 2.2), Inches(1),
                           Emu(w), Emu(1_000_000))
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    records = mod.detect(ctx)
    assert len(records) == 1
    assert records[0].old_value == "2080000x1000000"
