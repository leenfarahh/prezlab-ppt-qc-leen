"""Stage 1: reading a submitted master's design surface into a Style Spec.

The guide fixtures inject the exact XML desktop PowerPoint writes, captured
from a COM-authored probe deck on 18/08/2026: p15:sldGuideLst in the
slideMaster's extLst, positions in eighths of a point from the top-left slide
edge, missing pos meaning 0, orient="horz" meaning horizontal and a missing
orient meaning vertical. python-pptx cannot draw guides, so the format is
planted rather than produced.
"""

import io
import itertools
import struct
import zlib

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml.shapes.picture import CT_Picture
from pptx.util import Emu

from qc.stylespec import (GUIDE_UNIT_EMU, P15, extract_brand, extract_layouts,
                          extract_style_spec, extract_theme, find_brand_marks,
                          infer_grid, read_guides)

IN = 914400
HALF = IN // 2
PT8 = GUIDE_UNIT_EMU  # one guide unit (an eighth of a point) in EMU

_SHAPE_IDS = itertools.count(700)


# ------------------------------------------------------------------ helpers


def _png(rgb=(0, 0, 0)) -> io.BytesIO:
    def chunk(tag, data):
        return (struct.pack(">I", len(data)) + tag + data
                + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF))

    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    raw = bytes([0]) + bytes(rgb)
    return io.BytesIO(b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr)
                      + chunk(b"IDAT", zlib.compress(raw)) + chunk(b"IEND", b""))


def _stamp(container, png, left=IN, top=IN, w=HALF, h=HALF, name="brandmark"):
    _part, rId = container.part.get_or_add_image_part(png)
    pic = CT_Picture.new_pic(next(_SHAPE_IDS), name, name, rId, left, top, w, h)
    container._element.spTree.append(pic)
    return pic


def _plant_guides(master, vertical_pt=(), horizontal_pt=()):
    """Plant guides the way PowerPoint stores them: eighths of a point, and
    orient omitted for vertical guides."""
    ext_lst = etree.SubElement(
        master._element, "{http://schemas.openxmlformats.org/presentationml/2006/main}extLst")
    ext = etree.SubElement(
        ext_lst, "{http://schemas.openxmlformats.org/presentationml/2006/main}ext")
    ext.set("uri", "{GUIDE-TEST}")
    lst = etree.SubElement(ext, f"{{{P15}}}sldGuideLst")
    gid = itertools.count(1)
    for pos in vertical_pt:
        g = etree.SubElement(lst, f"{{{P15}}}guide")
        g.set("id", str(next(gid)))
        if pos:  # PowerPoint omits pos entirely for a guide at 0
            g.set("pos", str(int(pos * 8)))
    for pos in horizontal_pt:
        g = etree.SubElement(lst, f"{{{P15}}}guide")
        g.set("id", str(next(gid)))
        g.set("orient", "horz")
        if pos:
            g.set("pos", str(int(pos * 8)))
    return lst


def _master_only():
    """A presentation with a master and zero slides: Stage 1's normal input."""
    prs = Presentation()
    assert len(prs.slides) == 0
    return prs


# ------------------------------------------------------------------- guides


def test_guides_convert_from_eighths_of_a_point_to_emu():
    prs = _master_only()
    master = prs.slide_masters[0]
    _plant_guides(master, vertical_pt=(0, 480), horizontal_pt=(100,))

    guides = read_guides(master)
    assert guides["vertical_emu"] == [0, int(round(480 * 8 * PT8))]
    assert guides["horizontal_emu"] == [int(round(100 * 8 * PT8))]
    # 480pt is 12700*480 EMU; the conversion must survive the round trip.
    assert guides["vertical_emu"][1] == 480 * 12700


def test_missing_orient_means_vertical():
    prs = _master_only()
    _plant_guides(prs.slide_masters[0], vertical_pt=(10, 20, 30))
    guides = read_guides(prs.slide_masters[0])
    assert len(guides["vertical_emu"]) == 3
    assert guides["horizontal_emu"] == []


def test_grid_margins_come_from_the_outermost_guides():
    prs = _master_only()
    sw, sh = prs.slide_width, prs.slide_height
    _plant_guides(prs.slide_masters[0], vertical_pt=(36, 684), horizontal_pt=(27, 513))

    grid = infer_grid(prs, prs.slide_masters[0])
    assert grid["source"] == "guides"
    assert grid["margins_emu"]["left"] == 36 * 12700
    assert grid["margins_emu"]["right"] == sw - 684 * 12700
    assert grid["margins_emu"]["top"] == 27 * 12700
    assert grid["margins_emu"]["bottom"] == sh - 513 * 12700


def test_even_column_grid_is_inferred():
    prs = _master_only()
    # Six 100pt columns separated by 20pt gutters, starting at 36pt.
    edges, x = [], 36
    for _ in range(6):
        edges += [x, x + 100]
        x += 120
    _plant_guides(prs.slide_masters[0], vertical_pt=edges, horizontal_pt=(27, 513))

    grid = infer_grid(prs, prs.slide_masters[0])
    assert grid["columns"] == 6
    assert grid["gutter_emu"] == 20 * 12700


def test_uneven_guides_get_no_column_guess():
    """A wrong column count is worse for Stage 2 than an absent one."""
    prs = _master_only()
    _plant_guides(prs.slide_masters[0], vertical_pt=(36, 120, 200, 400, 410, 600),
                  horizontal_pt=(27, 513))
    grid = infer_grid(prs, prs.slide_masters[0])
    assert grid["source"] == "guides"
    assert grid["columns"] is None
    assert grid["gutter_emu"] is None


# ------------------------------------------------------- the header band
#
# The client masters draw four horizontal guides: the page's top and bottom
# margins, and between them the floor the subtitle may not cross and the ceiling
# the body may not cross. The strip between those two stays empty on every
# slide, which is what makes the headers read as one line down a deck.


def test_the_two_interior_guides_are_the_subtitle_floor_and_body_ceiling():
    prs = _master_only()
    _plant_guides(prs.slide_masters[0], vertical_pt=(36, 684),
                  horizontal_pt=(27, 119, 137, 513))

    grid = infer_grid(prs, prs.slide_masters[0])
    assert grid["subtitle_floor_emu"] == 119 * 12700
    assert grid["body_top_emu"] == 137 * 12700
    # and the outermost guides still answer for the page margins
    assert grid["margins_emu"]["top"] == 27 * 12700


def test_a_centre_guide_is_not_a_body_ceiling():
    """Masters carry a centre guide on each axis as a placement aid. Read as a
    ceiling it would put the start of the body half way down the slide."""
    prs = _master_only()
    centre_pt = prs.slide_height / 12700 / 2
    _plant_guides(prs.slide_masters[0], vertical_pt=(36, 684),
                  horizontal_pt=(27, 119, 137, centre_pt, 513))

    grid = infer_grid(prs, prs.slide_masters[0])
    assert grid["subtitle_floor_emu"] == 119 * 12700
    assert grid["body_top_emu"] == 137 * 12700


def test_a_single_interior_guide_states_a_ceiling_but_no_strip():
    prs = _master_only()
    _plant_guides(prs.slide_masters[0], vertical_pt=(36, 684),
                  horizontal_pt=(27, 137, 513))

    grid = infer_grid(prs, prs.slide_masters[0])
    assert grid["subtitle_floor_emu"] is None
    assert grid["body_top_emu"] == 137 * 12700


def test_three_interior_guides_name_no_band_at_all():
    """Guessing which pair a designer meant would seat content on a line
    nobody drew."""
    prs = _master_only()
    _plant_guides(prs.slide_masters[0], vertical_pt=(36, 684),
                  horizontal_pt=(27, 90, 119, 137, 513))

    grid = infer_grid(prs, prs.slide_masters[0])
    assert grid["subtitle_floor_emu"] is None
    assert grid["body_top_emu"] is None


def test_the_band_reaches_the_audit_profile():
    prs = _master_only()
    _plant_guides(prs.slide_masters[0], vertical_pt=(36, 684),
                  horizontal_pt=(27, 119, 137, 513))
    spec = extract_style_spec(prs, embed_assets=False)

    from qc.stylespec import spec_to_profile

    geo = spec_to_profile(spec, "p", "P")["config"]["geometry"]
    assert geo["body_band_emu"] == {"subtitle_floor": 119 * 12700,
                                    "body_top": 137 * 12700}


def test_a_master_with_no_band_carries_none_rather_than_a_default():
    """There is no sensible default for where a body begins: a guessed ceiling
    is a line the client never drew, and every slide would be measured
    against it."""
    prs = _master_only()
    _plant_guides(prs.slide_masters[0], vertical_pt=(36, 684),
                  horizontal_pt=(27, 513))
    spec = extract_style_spec(prs, embed_assets=False)

    from qc.stylespec import spec_to_profile

    assert spec_to_profile(spec, "p", "P")["config"]["geometry"][
        "body_band_emu"] is None


# -------------------------------------------------- declared title sizes
#
# "Some titles have bigger font sizes than others, why is that?" Nothing
# resizes them: a slide's title inherits whatever its LAYOUT declares, and the
# client master declares 24, 25, 28 and 32pt across its twelve layouts. The read
# has to show that, because it is fixable once in the master rather than per
# deck (design lead, 20/08/2026).


def _declare_title_size(layout, sz_pt=None, autofit=None):
    from pptx.oxml.ns import qn

    ph = next(p for p in layout.placeholders
              if str(p.placeholder_format.type).startswith(
                  ("TITLE", "CENTER_TITLE")))
    txBody = ph._element.find(qn("p:txBody"))
    if sz_pt:
        lst = txBody.makeelement(qn("a:lstStyle"), {})
        lvl = etree.SubElement(lst, qn("a:lvl1pPr"))
        etree.SubElement(lvl, qn("a:defRPr")).set("sz", str(int(sz_pt * 100)))
        txBody.insert(1, lst)
    if autofit:
        bodyPr = txBody.find(qn("a:bodyPr"))
        etree.SubElement(bodyPr, qn(f"a:{autofit}"))
    return ph


def test_a_layout_reports_the_title_size_it_declares():
    """Read from the layout's own lstStyle. The stock template already declares
    sizes on some of its layouts, so this asserts on the two it changes."""
    prs = _master_only()
    layouts = prs.slide_masters[0].slide_layouts
    _declare_title_size(layouts[0], sz_pt=28)
    _declare_title_size(layouts[1], sz_pt=32, autofit="normAutofit")

    spec = extract_style_spec(prs, embed_assets=False)

    def title(lay):
        return next(p for p in lay["placeholders"]
                    if p["type"] in ("title", "ctrTitle"))

    assert title(spec["layouts"][0])["size_pt"] == 28
    assert title(spec["layouts"][1])["size_pt"] == 32
    assert title(spec["layouts"][1])["autofit"] == "normAutofit"


def _layouts(*specs):
    """Layout dicts as the spec carries them: (name, size_pt, autofit)."""
    return [{"name": name,
             "placeholders": [{"type": "title", "size_pt": size,
                               "autofit": autofit}]}
            for name, size, autofit in specs]


def test_the_read_names_the_layouts_that_disagree_about_the_title():
    from qc.ui_master import _title_size_note

    note = _title_size_note(_layouts(("Section", 28, None),
                                     ("Content", 32, "normAutofit"),
                                     ("Quote", None, None)))
    assert "2 different title sizes" in note
    assert "28pt" in note and "32pt" in note
    assert "Section" in note and "Content" in note
    assert "shrink text on overflow" in note
    assert "Quote" not in note, "a layout that inherits declares nothing"


def test_no_note_when_every_layout_agrees():
    from qc.ui_master import _title_size_note

    assert _title_size_note(_layouts(("A", 28, None), ("B", 28, "noAutofit"),
                                     ("C", None, None))) == ""


# ------------------------------------------------- the presentation space
#
# A rectangle the designer draws on the master and names. It outranks the guides
# because it is not an interpretation: a master can carry an outer page margin,
# a column grid, a header band and a bleed line, and choosing among them is the
# guess this ends (design lead, 21/08/2026: "some cases have multiple margins,
# so presentation space is safer").


def _draw_space(container, left_in, top_in, w_in, h_in,
                name="Presentation space", fill=False, alt=None):
    """Plant the rectangle a designer would draw. python-pptx cannot add shapes
    to a master or a layout, so the element goes in directly - the same approach
    the guide fixtures take.

    `alt` is the shape's alt text (cNvPr/@descr), which is where the ToolsToo
    add-in stamps its marker: the shape name stays whatever PowerPoint made it.
    """
    from pptx.oxml.shapes.autoshape import CT_Shape

    sp = CT_Shape.new_autoshape_sp(next(_SHAPE_IDS), name, "rect",
                                   int(left_in * IN), int(top_in * IN),
                                   int(w_in * IN), int(h_in * IN))
    container.shapes._spTree.insert_element_before(sp, "p:extLst")
    shape = next(s for s in container.shapes if s.name == name)
    if alt is not None:
        shape._element.find(qn("p:nvSpPr")).find(qn("p:cNvPr")).set("descr", alt)
    if not fill:
        shape.fill.background()
        shape.line.fill.background()
    return shape


def test_the_presentation_space_states_the_frame():
    prs = _master_only()
    sw, sh = prs.slide_width, prs.slide_height
    _draw_space(prs.slide_masters[0], 0.5, 1.0, 9.0, 5.5)

    grid = infer_grid(prs, prs.slide_masters[0])
    assert grid["source"] == "presentation_space"
    m = grid["margins_emu"]
    assert m["left"] == int(0.5 * IN)
    assert m["top"] == int(1.0 * IN)
    assert m["right"] == sw - int(9.5 * IN)
    assert m["bottom"] == sh - int(6.5 * IN)


def test_the_presentation_space_beats_the_guides():
    """Both are stated, but only one of them needs interpreting."""
    prs = _master_only()
    _plant_guides(prs.slide_masters[0], vertical_pt=(36, 684),
                  horizontal_pt=(27, 513))
    _draw_space(prs.slide_masters[0], 1.0, 1.0, 8.0, 5.0)

    grid = infer_grid(prs, prs.slide_masters[0])
    assert grid["source"] == "presentation_space"
    assert grid["margins_emu"]["left"] == int(1.0 * IN)
    # the guides are still read, and still answer for the header band
    assert grid["guides"]["vertical_emu"]


def test_the_header_band_is_still_read_from_the_guides():
    """A single rectangle cannot state a pair of lines inside itself."""
    prs = _master_only()
    _plant_guides(prs.slide_masters[0], vertical_pt=(36, 684),
                  horizontal_pt=(27, 119, 137, 513))
    _draw_space(prs.slide_masters[0], 0.5, 0.4, 9.0, 6.0)

    grid = infer_grid(prs, prs.slide_masters[0])
    assert grid["source"] == "presentation_space"
    assert grid["subtitle_floor_emu"] == 119 * 12700
    assert grid["body_top_emu"] == 137 * 12700


def test_a_marker_that_would_print_is_reported():
    prs = _master_only()
    _draw_space(prs.slide_masters[0], 0.5, 1.0, 9.0, 5.5, fill=True)

    space = infer_grid(prs, prs.slide_masters[0])["presentation_space"]
    assert space["prints"] is True


def test_a_marker_on_a_layout_is_read_and_its_home_reported():
    """A designer trying this out will put it wherever seems natural."""
    prs = _master_only()
    layout = prs.slide_masters[0].slide_layouts[0]
    _draw_space(layout, 0.75, 1.0, 8.5, 5.0)

    grid = infer_grid(prs, prs.slide_masters[0])
    assert grid["source"] == "presentation_space"
    assert grid["margins_emu"]["left"] == int(0.75 * IN)
    assert layout.name in grid["presentation_space"]["source"]


def test_the_name_match_is_forgiving_but_not_loose():
    prs = _master_only()
    master = prs.slide_masters[0]
    _draw_space(master, 0.5, 1.0, 9.0, 5.5, name="PRESENTATION-SPACE")
    assert infer_grid(prs, master)["source"] == "presentation_space"

    prs2 = _master_only()
    _draw_space(prs2.slide_masters[0], 0.5, 1.0, 9.0, 5.5,
                name="Presentation notes")
    assert infer_grid(prs2, prs2.slide_masters[0])["source"] != \
        "presentation_space"


def test_a_marker_declared_by_alt_text_is_read():
    """What the add-in actually writes. ToolsToo stamps the alt text and leaves
    the shape name as PowerPoint made it, so a name-only read missed the only
    marker on the client's own master."""
    prs = _master_only()
    master = prs.slide_masters[0]
    _draw_space(master, 0.5, 1.0, 9.0, 5.5, name="Rectangle 2",
                alt="ToolsToo_PS")

    grid = infer_grid(prs, master)
    assert grid["source"] == "presentation_space"
    assert grid["margins_emu"]["left"] == HALF
    space = grid["presentation_space"]
    assert space["marker"] == "alt_text"
    assert space["alt_text"] == "ToolsToo_PS"
    assert space["name"] == "Rectangle 2"


def test_the_alt_text_marker_outranks_a_typed_name():
    """A name is typed and survives copying, so it can be a leftover; the alt
    text is on the rectangle the add-in is currently driving."""
    prs = _master_only()
    master = prs.slide_masters[0]
    _draw_space(master, 2.0, 1.0, 6.0, 5.0, name="Presentation space")
    _draw_space(master, 0.5, 1.0, 9.0, 5.5, name="Rectangle 9",
                alt="ToolsToo_PS")

    grid = infer_grid(prs, master)
    assert grid["presentation_space"]["marker"] == "alt_text"
    assert grid["margins_emu"]["left"] == HALF, "the add-in's rectangle binds"
    rivals = grid["presentation_space"]["rivals"]
    assert [r["name"] for r in rivals] == ["Presentation space"]


def test_two_markers_that_agree_are_not_reported_as_a_disagreement():
    prs = _master_only()
    master = prs.slide_masters[0]
    _draw_space(master, 0.5, 1.0, 9.0, 5.5, name="Presentation space")
    _draw_space(master, 0.5, 1.0, 9.0, 5.5, name="Rectangle 9",
                alt="ToolsToo_PS")

    assert infer_grid(prs, master)["presentation_space"]["rivals"] == []


def test_markers_on_two_layouts_that_disagree_are_both_named():
    """Which rectangle won has to be answerable: the frame governs every slide
    of every deck on the profile, and layout order is nobody's decision."""
    prs = _master_only()
    master = prs.slide_masters[0]
    first, second = master.slide_layouts[0], master.slide_layouts[1]
    _draw_space(first, 0.5, 1.0, 9.0, 5.5, name="Rectangle 2",
                alt="ToolsToo_PS")
    _draw_space(second, 1.5, 1.0, 7.0, 5.0, name="Rectangle 3",
                alt="ToolsToo_PS")

    space = infer_grid(prs, master)["presentation_space"]
    assert space["box_emu"][0] == HALF, "the first layout's marker states it"
    assert [r["name"] for r in space["rivals"]] == ["Rectangle 3"]
    assert second.name in space["rivals"][0]["source"]


# ------------------------------------------- what the rectangle is a frame AROUND
#
# Designers draw it both ways. Around the whole page, its top is the page's top
# margin. Around the BODY - which is how the client master draws it, landing on
# the master's own body placeholder to the EMU - its top is where CONTENT
# begins, and the header sits above it by design. Read as a page margin, that
# top becomes the deck's safe-zone ceiling and every title on every slide is
# reported as breaking a margin nobody drew there.


def test_a_frame_around_the_body_does_not_become_the_page_margin():
    prs = _master_only()
    master = prs.slide_masters[0]          # its title box ends at 1.55in
    _plant_guides(master, vertical_pt=(36, 684), horizontal_pt=(36, 504))
    _draw_space(master, 0.5, 2.0, 9.0, 4.0)

    grid = infer_grid(prs, master)
    assert grid["source"] == "presentation_space"
    assert grid["space_states"] == "body"
    assert grid["margins_emu"]["top"] == HALF, "the page top comes from the guides"
    assert grid["body_top_emu"] == 2 * IN, "the rectangle states where body starts"
    # The sides and the floor are still the rectangle's: those it does state.
    assert grid["margins_emu"]["left"] == HALF
    assert grid["margins_emu"]["bottom"] == int(1.5 * IN)


def test_a_frame_around_the_whole_page_keeps_its_top_margin():
    prs = _master_only()
    master = prs.slide_masters[0]
    _plant_guides(master, vertical_pt=(36, 684), horizontal_pt=(36, 504))
    _draw_space(master, 0.5, 1.0, 9.0, 5.5)   # above the title's floor

    grid = infer_grid(prs, master)
    assert grid["space_states"] == "page"
    assert grid["margins_emu"]["top"] == IN
    assert grid["body_top_emu"] is None, "no band guides, so nothing states one"


def test_a_body_frame_leaves_the_guides_band_alone():
    """The client master states the same line twice - the rectangle's top and
    the second interior guide are both 1.90in - and the guides also carry the
    subtitle floor above it, which one rectangle cannot state."""
    prs = _master_only()
    master = prs.slide_masters[0]
    _plant_guides(master, vertical_pt=(36, 684),
                  horizontal_pt=(36, 119, 137, 504))
    _draw_space(master, 0.5, 137 * 12700 / IN, 9.0, 4.0)

    grid = infer_grid(prs, master)
    assert grid["space_states"] == "body"
    assert grid["subtitle_floor_emu"] == 119 * 12700
    assert grid["body_top_emu"] == 137 * 12700
    assert grid["margins_emu"]["top"] == HALF


def test_a_body_frame_falls_back_to_placeholders_for_the_page_top():
    """No guides to read the page's own top margin from, so it comes from the
    next-best statement rather than from the body frame."""
    prs = _master_only()
    master = prs.slide_masters[0]
    _draw_space(master, 0.5, 2.0, 9.0, 4.0)

    title_top = min(ph.top for ph in master.placeholders
                    if ph.top is not None and ph.top < 2 * IN)

    grid = infer_grid(prs, master)
    assert grid["space_states"] == "body"
    assert grid["margins_emu"]["top"] == title_top, "the title box's own top"
    assert grid["body_top_emu"] == 2 * IN


def test_a_rectangle_off_the_canvas_is_reported_and_not_used():
    """A frame has to be ON the page. One hanging off the canvas - left over
    from a resized master, or drawn at another slide size - would hand every
    downstream pass a negative margin."""
    prs = _master_only()
    _plant_guides(prs.slide_masters[0], vertical_pt=(36, 684),
                  horizontal_pt=(27, 513))
    _draw_space(prs.slide_masters[0], 0.5, 1.0, 14.0, 5.5)   # wider than 10in

    grid = infer_grid(prs, prs.slide_masters[0])
    assert grid["source"] == "guides", "the bad rectangle must not be used"
    assert grid["presentation_space"]["problem"]
    assert all(v >= 0 for v in grid["margins_emu"].values())


def test_grid_falls_back_to_master_content_placeholders():
    prs = _master_only()
    grid = infer_grid(prs, prs.slide_masters[0])
    assert grid["source"] == "placeholders"
    assert grid["guides"] == {"vertical_emu": [], "horizontal_emu": []}
    # Footer and slide number sit low in the margin band; if they had counted
    # as content the bottom margin would collapse toward zero.
    assert grid["margins_emu"]["bottom"] > 0


# ---------------------------------------------------------------- layouts


def test_layout_archetype_and_placeholder_geometry_source():
    """The archetype token is Stage 2's template-bank matching key, and
    explicit-vs-inherited says which layouts follow a master edit."""
    prs = _master_only()
    layouts = {lay["name"]: lay for lay in extract_layouts(prs.slide_masters[0])}

    assert layouts["Title Slide"]["type"] == "title"
    assert layouts["Title and Content"]["type"] == "obj"
    assert layouts["Section Header"]["type"] == "secHead"
    assert layouts["Two Content"]["type"] == "twoObj"

    title_ph = [p for p in layouts["Title Slide"]["placeholders"]
                if p["type"] == "ctrTitle"]
    assert title_ph and title_ph[0]["geometry_source"] == "explicit"

    inherited = [p for p in layouts["Title and Content"]["placeholders"]
                 if p["type"] == "title"]
    assert inherited and inherited[0]["geometry_source"] == "inherited"
    # Inherited geometry still resolves to real numbers, via the master.
    assert inherited[0]["position_emu"]["width"] > 0


def test_body_placeholders_default_to_the_body_token():
    """OOXML omits p:ph/@type for body placeholders; absent must not read as
    'no type', or Stage 2 loses every content placeholder."""
    prs = _master_only()
    two = [lay for lay in extract_layouts(prs.slide_masters[0])
           if lay["name"] == "Two Content"][0]
    bodies = [p for p in two["placeholders"] if p["type"] == "body"]
    assert len(bodies) == 2


# ------------------------------------------------------- master furniture


def test_slide_number_field_is_not_reported_as_text():
    """The sldNum placeholder holds an a:fld that python-pptx renders as
    private-use marker characters; that string is never content."""
    prs = _master_only()
    spec = extract_style_spec(prs)
    sldnum = spec["master"]["slide_number"]
    assert sldnum["present"]
    assert sldnum["field"] == "slidenum"
    assert sldnum["text"] is None


def test_master_background_resolves_a_theme_reference():
    prs = _master_only()
    bg = extract_style_spec(prs)["master"]["background"]
    assert bg["kind"] == "theme_ref"
    assert bg["hex"] == "FFFFFF"


# ------------------------------------------------------------------- brand


def test_master_submission_ignores_slide_stamps():
    """Stage 1 reads the design surface only. A picture repeated on slides is
    content in this context, however often it recurs."""
    prs = Presentation()
    for _ in range(5):
        prs.slides.add_slide(prs.slide_layouts[6]).shapes.add_picture(
            _png((3, 3, 3)), Emu(IN), Emu(IN), Emu(HALF), Emu(HALF))

    assert find_brand_marks(prs, include_slides=False) == []
    assert extract_brand(prs, include_slides=False)["logo"] is None
    # The deck-survey path still sees it; the two questions differ.
    assert find_brand_marks(prs, include_slides=True)[0]["scope"] == "slides"


def test_master_logo_is_found_on_the_design_surface():
    prs = _master_only()
    _stamp(prs.slide_masters[0], _png((11, 22, 33)), left=6 * IN, top=IN // 4)
    spec = extract_style_spec(prs)
    assert spec["brand"]["logo"]["scope"] == "master"
    assert spec["brand"]["logo"]["position_emu"]["left"] == 6 * IN


# -------------------------------------------------------------------- spec


def test_spec_is_complete_for_a_file_with_no_slides():
    """The whole premise of Stage 1: a designer submits a master, not a deck."""
    prs = _master_only()
    spec = extract_style_spec(prs, source="master.pptx")

    assert spec["spec_version"] == 1
    assert spec["meta"]["slide_count"] == 0
    assert spec["meta"]["source_file"] == "master.pptx"
    assert spec["meta"]["slide_size_emu"]["width"] == prs.slide_width
    assert spec["theme"]["colors"]["accent1"]
    assert spec["theme"]["fonts"]["major"]["latin"]
    assert len(spec["layouts"]) == len(list(prs.slide_masters[0].slide_layouts))
    assert spec["master"]["placeholders"]
    assert spec["grid"]["source"] == "placeholders"


def test_spec_is_json_serialisable():
    """It is the artifact every later stage consumes, so it has to round-trip
    through JSON with no custom encoder."""
    import json

    prs = _master_only()
    _stamp(prs.slide_masters[0], _png((5, 5, 5)))
    _plant_guides(prs.slide_masters[0], vertical_pt=(36, 684), horizontal_pt=(27, 513))
    spec = extract_style_spec(prs, source="m.pptx")

    round_tripped = json.loads(json.dumps(spec))
    assert round_tripped == spec


def test_theme_still_reads_from_a_master_with_no_slides():
    theme = extract_theme(_master_only())
    assert {"dk1", "lt1", "accent1", "accent6"} <= set(theme["colors"])
    assert theme["color_map"]["bg1"] == "lt1"
