"""When the model and a measured rule say the same thing, the model's record is
the one that survives.

"The LLM output is not being taken into account; it should weigh more than the
code reasoning" (design lead, 02/09/2026). Nothing was ignoring the model. Four
separate places QUIETLY PREFERRED THE RULER, and each one looked locally
reasonable:

1. qc.copilot.synthesize deduplicated its new records against EVERY record on
   the manifest, so a measured edge finding that arrived first took the claim
   and the model's answer about that edge was never emitted.

2. qc.components.synthesize did the same, one issue type along.

3. qc.web._surviving_vision dropped any suggestion the re-audit had
   independently found - explicitly, with a comment saying so.

4. qc.components.synthesize threw away any component sitting more than 0.5in
   off the line the model named, on the grounds that a gap that big must be a
   deliberate composition. That is a threshold for INFERRING intent, used in a
   pass where the model had already stated it.

The net effect: exactly one card per shape per move, which is right, and it was
always the one without the reason on it. These tests fix which half wins.

The model is stubbed throughout. What is being tested is what happens to its
answer after it arrives.
"""

import io

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Emu

import qc.components as components
import qc.copilot as copilot
import qc.llm as llm
import qc.web as web
from qc.records import claim_keys, make_record, vision_wins
from qc.web import app
from tests.conftest import job_id_of

client = TestClient(app)
IN = 914400
SLIDE_W, SLIDE_H = 12192000, 6858000


def _rec(source, issue, shape_id="7", prop="spPr.xfrm.off.y", slide=0,
         message="x"):
    return make_record(
        slide_index=slide, shape_id=shape_id, module="margin_alignment",
        issue_type=issue, property=prop, old_value=100, new_value=200,
        confidence="medium", message=message, source=source).to_dict()


# ------------------------------------------------------- the rule itself


def test_two_names_for_one_question_are_one_claim():
    """The measured module and the component review have different names for
    "this shape is off the line it should sit on". A key that read the issue
    type would never see them collide."""
    measured = _rec("measured", "margin_alignment.edge_misaligned")
    seen = _rec("vision", "margin_alignment.component_edge_misaligned")
    assert claim_keys(measured) == claim_keys(seen)


def test_the_axis_still_separates_two_different_moves():
    """A shape off its top edge and off its left edge are two claims, not one:
    fixing either leaves the other true."""
    up = _rec("vision", "margin_alignment.edge_misaligned",
              prop="spPr.xfrm.off.y")
    across = _rec("vision", "margin_alignment.edge_misaligned",
                  prop="spPr.xfrm.off.x")
    assert not (claim_keys(up) & claim_keys(across))


def test_a_reposition_on_both_axes_collides_with_either_one():
    """spPr.xfrm.off names no axis because it moves the shape outright. It has
    to collide with a record that moves the same shape on either axis alone, or
    a designer gets one card that moves it and another that moves it back."""
    both = _rec("measured", "margin_alignment.recurring_off_position",
                prop="spPr.xfrm.off")
    one = _rec("vision", "margin_alignment.component_edge_misaligned",
               prop="spPr.xfrm.off.y")
    assert claim_keys(both) & claim_keys(one)


def test_a_different_question_about_the_same_shape_is_a_different_claim():
    """THE NARROWNESS IS THE SAFETY. A footer sitting off the page and a card
    3mm off its line both write spPr.xfrm.off.y on the same shape, and they are
    not the same fact. Keyed on the property alone, an alignment suggestion
    would silently delete a footer that runs off the slide."""
    footer = make_record(slide_index=0, shape_id="7", module="header_footer",
                         issue_type="header_footer.footer_off_canvas",
                         property="spPr.xfrm.off.y", old_value=1, new_value=2,
                         message="x").to_dict()
    seen = _rec("vision", "margin_alignment.component_edge_misaligned",
                prop="spPr.xfrm.off.y")

    assert not (claim_keys(footer) & claim_keys(seen))
    assert len(vision_wins([footer, seen])) == 2


def test_a_finding_with_no_family_claims_nothing():
    """A font record cannot be superseded and cannot supersede: there is no
    question a vision pass and the font module both answer."""
    font = make_record(slide_index=0, shape_id="7", module="font",
                       issue_type="font.family_out_of_set",
                       property="rPr.latin", old_value="Arial",
                       new_value="Georgia", message="x").to_dict()
    assert claim_keys(font) == set()


def test_the_measured_twin_is_the_one_that_goes():
    measured = _rec("measured", "margin_alignment.edge_misaligned",
                    message="Top edge 3mm off the median of its row.")
    seen = _rec("vision", "margin_alignment.component_edge_misaligned",
                message="The three cards were meant to share a top edge.")

    kept = vision_wins([measured, seen])

    assert len(kept) == 1
    assert kept[0]["source"] == "vision", "the ruler won"
    assert "meant to share" in kept[0]["message"]


def test_a_measured_finding_the_model_said_nothing_about_is_untouched():
    """The eviction is narrow on purpose. Most of an audit is fonts, colours
    and margins the model was never asked about, and losing those would be a
    much worse bug than the one being fixed."""
    font = make_record(slide_index=0, shape_id="7", module="font",
                       issue_type="font.family_out_of_set",
                       property="rPr.latin", old_value="Arial",
                       new_value="Georgia", message="x").to_dict()
    other_shape = _rec("measured", "margin_alignment.edge_misaligned",
                       shape_id="9")
    seen = _rec("vision", "margin_alignment.component_edge_misaligned")

    kept = vision_wins([font, other_shape, seen])

    assert len(kept) == 3
    assert {r["record_id"] for r in kept} == {font["record_id"],
                                              other_shape["record_id"],
                                              seen["record_id"]}


def test_with_no_vision_records_nothing_moves():
    records = [_rec("measured", "margin_alignment.edge_misaligned"),
               _rec("measured", "margin_alignment.edge_misaligned",
                    shape_id="9")]
    assert vision_wins(records) == records


# -------------------------------------------- the passes stop silencing it


def _pair_slide():
    """Two headers meant to share a top, one 0.3in high. The measured module
    would flag the same edge on the same shape."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(SLIDE_W), Emu(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = slide.shapes.add_textbox(Emu(int(0.9 * IN)), Emu(int(1.54 * IN)),
                                    Emu(int(5.0 * IN)), Emu(int(0.9 * IN)))
    left.text_frame.text = "Method 1"
    right = slide.shapes.add_textbox(Emu(int(6.75 * IN)), Emu(int(1.84 * IN)),
                                     Emu(int(5.1 * IN)), Emu(int(0.9 * IN)))
    right.text_frame.text = "Method 2"
    buf = io.BytesIO()
    prs.save(buf)
    return Presentation(io.BytesIO(buf.getvalue())).slides[0]


def test_a_measured_record_no_longer_silences_the_copilot():
    """The measured module had already flagged this exact edge. The model's
    answer about it still has to be emitted - it is the one that will win."""
    slide = _pair_slide()
    left = list(slide.shapes)[0]
    obs = [{"action": "align_top",
            "shape_ids": [str(s.shape_id) for s in slide.shapes],
            "rationale": "The two column headers should share a top edge."}]

    alone = copilot.synthesize(slide, 0, obs, existing=[])
    assert len(alone) == 1, "the pass found nothing to compare against"

    measured = _rec("measured", "margin_alignment.edge_misaligned",
                    shape_id=str(left.shape_id), prop="spPr.xfrm.off.y")
    with_rival = copilot.synthesize(slide, 0, obs, existing=[measured])

    assert len(with_rival) == 1, \
        "a measured record took the claim and the model's answer vanished"
    assert with_rival[0]["source"] == "vision"


def test_another_vision_record_does_still_silence_it():
    """The de-duplication is not being switched off, only re-pointed. Two
    vision passes noticing one edge must still produce one card."""
    slide = _pair_slide()
    obs = [{"action": "align_top",
            "shape_ids": [str(s.shape_id) for s in slide.shapes],
            "rationale": "The two column headers should share a top edge."}]
    first = copilot.synthesize(slide, 0, obs, existing=[])
    again = copilot.synthesize(slide, 0, obs, existing=first)
    assert again == []


def test_a_measured_record_no_longer_silences_the_component_review():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(SLIDE_W), Emu(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = slide.shapes.add_textbox(Emu(int(0.9 * IN)), Emu(2 * IN),
                                 Emu(2 * IN), Emu(IN))
    a.text_frame.text = "a"
    b = slide.shapes.add_textbox(Emu(int(1.2 * IN)), Emu(3 * IN),
                                 Emu(2 * IN), Emu(IN))
    b.text_frame.text = "b"
    buf = io.BytesIO()
    prs.save(buf)
    slide = Presentation(io.BytesIO(buf.getvalue())).slides[0]
    ids = [str(s.shape_id) for s in slide.shapes]
    layout = {
        "components": [{"name": "a", "shape_ids": [ids[0]]},
                       {"name": "b", "shape_ids": [ids[1]]}],
        "alignments": [{"axis": "left", "anchor": "a",
                        "components": ["a", "b"],
                        "rationale": "The two blocks share a left edge."}],
    }

    alone = components.synthesize(slide, 0, layout, None, [], SLIDE_W, SLIDE_H)
    assert len(alone) == 1

    measured = _rec("measured", "margin_alignment.edge_misaligned",
                    shape_id=ids[1], prop="spPr.xfrm.off.x")
    with_rival = components.synthesize(slide, 0, layout, None, [measured],
                                       SLIDE_W, SLIDE_H)
    assert len(with_rival) == 1 and with_rival[0]["source"] == "vision"


# ------------------------------------------------------ end to end, on a job


def _two_column_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(SLIDE_W), Emu(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Emu(int(0.9 * IN)), Emu(int(0.4 * IN)),
                                     Emu(int(11.5 * IN)), Emu(int(1.45 * IN)))
    title.text_frame.text = "Two ways to do it"
    left = slide.shapes.add_textbox(Emu(int(0.9 * IN)), Emu(int(1.54 * IN)),
                                    Emu(int(5.0 * IN)), Emu(int(0.9 * IN)))
    left.text_frame.text = "Method 1"
    right = slide.shapes.add_textbox(Emu(int(6.75 * IN)), Emu(int(1.84 * IN)),
                                     Emu(int(5.1 * IN)), Emu(int(0.9 * IN)))
    right.text_frame.text = "Method 2"
    buf = io.BytesIO()
    prs.save(buf)
    return left, right, buf.getvalue()


def test_the_model_wins_on_the_manifest_and_the_page_says_so(monkeypatch):
    """The whole path: a measured record about one edge, the model's answer
    about the same edge, and one card left at the end - the model's."""
    left, right, data = _two_column_deck()
    r = client.post("/audit", files={"deck": ("d.pptx", data,
                                              "application/octet-stream")},
                    data={"profile": "prezlab_en"})
    job_id = job_id_of(r)
    job = web._job(job_id)

    # A measured record naming exactly the edge the model is about to name.
    planted = _rec("measured", "margin_alignment.edge_misaligned",
                   shape_id=str(left.shape_id), prop="spPr.xfrm.off.y",
                   message="Top edge off the median of its row.")
    job["manifest"]["records"].append(planted)

    obs = [{"action": "align_top",
            "shape_ids": [str(left.shape_id), str(right.shape_id)],
            "rationale": "The two column headers should share a top edge."}]
    monkeypatch.setattr(copilot, "ask_json", lambda **kw: {"observations": obs})
    monkeypatch.setattr(llm, "api_configured", lambda: True)
    monkeypatch.setattr(web, "_ensure_thumbs", lambda jid, j: j.__setitem__(
        "thumbs", {i: b"png" for i in range(j["manifest"]["slides"])}))

    page = client.post(f"/copilot/{job_id}").text
    records = web._job(job_id)["manifest"]["records"]

    mine = [r for r in records
            if str(r["shape_id"]) == str(left.shape_id)
            and r.get("property") == "spPr.xfrm.off.y"]
    assert len(mine) == 1, f"two cards for one move: {[r['source'] for r in mine]}"
    assert mine[0]["source"] == "vision"
    assert planted["record_id"] not in {r["record_id"] for r in records}
    assert "replaced by the model" in page, \
        "a record disappeared from the count with nothing said about it"


def test_the_summary_counts_what_the_manifest_holds(monkeypatch):
    """The eviction happens inside _merge_records so the summary is recomputed
    after it. A total that still counts an evicted record is how the numbers on
    the report start disagreeing with the rows under them."""
    left, right, data = _two_column_deck()
    r = client.post("/audit", files={"deck": ("d.pptx", data,
                                              "application/octet-stream")},
                    data={"profile": "prezlab_en"})
    job_id = job_id_of(r)
    job = web._job(job_id)
    job["manifest"]["records"].append(
        _rec("measured", "margin_alignment.edge_misaligned",
             shape_id=str(left.shape_id), prop="spPr.xfrm.off.y"))

    obs = [{"action": "align_top",
            "shape_ids": [str(left.shape_id), str(right.shape_id)],
            "rationale": "Shared top edge."}]
    monkeypatch.setattr(copilot, "ask_json", lambda **kw: {"observations": obs})
    monkeypatch.setattr(llm, "api_configured", lambda: True)
    monkeypatch.setattr(web, "_ensure_thumbs", lambda jid, j: j.__setitem__(
        "thumbs", {i: b"png" for i in range(j["manifest"]["slides"])}))
    client.post(f"/copilot/{job_id}")

    manifest = web._job(job_id)["manifest"]
    assert manifest["summary"]["total"] == len(manifest["records"])


def test_a_suggestion_the_reaudit_rediscovers_is_no_longer_the_one_dropped(
        monkeypatch):
    """_surviving_vision used to drop exactly this record. It is the one to
    keep: the re-audit can always re-derive the measured half, and nothing can
    re-derive the model's without paying for another call."""
    left, right, data = _two_column_deck()
    r = client.post("/audit", files={"deck": ("d.pptx", data,
                                              "application/octet-stream")},
                    data={"profile": "prezlab_en"})
    job_id = job_id_of(r)
    job = web._job(job_id)

    obs = [{"action": "align_top",
            "shape_ids": [str(left.shape_id), str(right.shape_id)],
            "rationale": "The two column headers should share a top edge."}]
    monkeypatch.setattr(copilot, "ask_json", lambda **kw: {"observations": obs})
    monkeypatch.setattr(llm, "api_configured", lambda: True)
    monkeypatch.setattr(web, "_ensure_thumbs", lambda jid, j: j.__setitem__(
        "thumbs", {i: b"png" for i in range(j["manifest"]["slides"])}))
    client.post(f"/copilot/{job_id}")

    seen = [r for r in job["manifest"]["records"] if r.get("source") == "vision"]
    assert len(seen) == 1

    survivors = web._surviving_vision(job)
    assert [s["record_id"] for s in survivors] == [seen[0]["record_id"]]

    # and the re-audit's own twin is what goes when the two land together
    rival = _rec("measured", "margin_alignment.edge_misaligned",
                 shape_id=str(seen[0]["shape_id"]), prop=seen[0]["property"])
    job["manifest"] = {"records": [rival], "slides": 1, "deck": "d.pptx",
                       "summary": {}}
    web._merge_records(job, survivors)
    kept = job["manifest"]["records"]
    assert [r["source"] for r in kept] == ["vision"]
