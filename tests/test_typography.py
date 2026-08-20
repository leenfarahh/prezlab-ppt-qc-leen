"""Typography: label case conventions, redundant size overrides, sibling
size agreement. Ground truth 20/07/2026: the designer's text changes were
rule-shaped (40 sibling-set uppercases by retype, 11 no-op override
removals, sibling rescales)."""

import io

from pptx import Presentation
from pptx.util import Emu, Pt

import qc.modules.typography as ty
from qc.fixer import apply_fixes, is_fixable
from qc.profile import Profile
from tests.conftest import save_and_ctx

IN = 914400


def _label(slide, text, left, top, w=1500000, h=400000, pt=12):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tb.text_frame.text = text
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(pt)
    return tb


def _records(ctx, kind):
    return [r.to_dict() for r in ty.detect(ctx)
            if r.issue_type == f"typography.{kind}"]


def test_sibling_caps_convention_flags_the_stray(make_prs, en_profile,
                                                 tmp_path):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _label(slide, "ALPHA SECTION", IN, IN)
    _label(slide, "BETA SECTION", 3 * IN, IN)
    stray = _label(slide, "Gamma section", 5 * IN, IN)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    recs = _records(ctx, "case_inconsistent")
    assert len(recs) == 1
    rec = recs[0]
    assert rec["shape_id"] == str(stray.shape_id)
    assert rec["old_value"] == "Gamma section"
    assert rec["new_value"] == "GAMMA SECTION"
    assert is_fixable(rec)

    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {rec["record_id"]})
    assert result.applied == 1
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    fixed = next(s for s in out.shapes
                 if str(s.shape_id) == rec["shape_id"])
    assert fixed.text_frame.text == "GAMMA SECTION"


def test_no_caps_majority_no_flag(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _label(slide, "ALPHA SECTION", IN, IN)
    _label(slide, "Beta section", 3 * IN, IN)
    _label(slide, "Gamma section", 5 * IN, IN)
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert _records(ctx, "case_inconsistent") == []


def test_profile_convention_flags_all_lowercase_labels(make_prs, tmp_path):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _label(slide, "Beta section", 3 * IN, IN)
    profile = Profile({"id": "t", "config": {
        "typography": {"label_case": "upper"}}})
    ctx = save_and_ctx(prs, tmp_path, profile)
    recs = _records(ctx, "case_inconsistent")
    assert len(recs) == 1
    assert "profile label case" in recs[0]["message"]


def test_long_text_and_sentences_are_not_labels(make_prs, en_profile,
                                                tmp_path):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _label(slide, "THIS HEADER", IN, IN)
    _label(slide, "OTHER HEADER", 3 * IN, IN)
    _label(slide, "A sentence explaining the point.", 5 * IN, IN)
    _label(slide, "far too many words to read as a label here", 7 * IN, IN)
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert _records(ctx, "case_inconsistent") == []


def test_redundant_size_override_detected_and_stripped(make_prs, en_profile,
                                                       tmp_path):
    from pptx.oxml.ns import qn
    from spike.resolver import resolve_run

    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text_frame.text = "Deck title"
    run = title.text_frame.paragraphs[0].runs[0]
    para = title.text_frame.paragraphs[0]
    inherited = resolve_run(run, para, title, slide, prs).size_pt
    assert inherited.source != "hard-default"
    run.font.size = Pt(inherited.value)  # restates what it inherits

    ctx = save_and_ctx(prs, tmp_path, en_profile)
    recs = _records(ctx, "redundant_size_override")
    mine = [r for r in recs if r["locator"] == "p0/r0"]
    assert len(mine) == 1
    rec = mine[0]
    assert rec["confidence"] == "deterministic"
    assert is_fixable(rec)

    result = apply_fixes(ctx.deck_path.read_bytes(), [rec],
                         {rec["record_id"]})
    assert result.applied == 1
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    rpr = out.shapes.title.text_frame.paragraphs[0].runs[0]._r.find(qn("a:rPr"))
    assert rpr is None or rpr.get("sz") is None
    # and the effective size is unchanged: the removal was a no-op
    out_run = out.shapes.title.text_frame.paragraphs[0].runs[0]
    out_para = out.shapes.title.text_frame.paragraphs[0]
    eff = resolve_run(out_run, out_para, out.shapes.title, out,
                      Presentation(io.BytesIO(result.cleaned_bytes)))
    assert abs(eff.size_pt.value - inherited.value) < 0.01


def test_meaningful_override_not_flagged(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text_frame.text = "Deck title"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(13.5)  # deliberate
    ctx = save_and_ctx(prs, tmp_path, en_profile)
    assert [r for r in _records(ctx, "redundant_size_override")
            if r["shape_id"] == str(title.shape_id)] == []


def test_sibling_size_stray_unified(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _label(slide, "ALPHA", IN, IN, pt=12)
    _label(slide, "BETA", 3 * IN, IN, pt=12)
    stray = _label(slide, "GAMMA", 5 * IN, IN, pt=10)
    ctx = save_and_ctx(prs, tmp_path, en_profile)

    recs = _records(ctx, "size_inconsistent")
    assert len(recs) == 1
    rec = recs[0]
    assert rec["shape_id"] == str(stray.shape_id)
    assert int(rec["old_value"]) == 1000 and int(rec["new_value"]) == 1200

    result = apply_fixes(ctx.deck_path.read_bytes(), recs,
                         {rec["record_id"]})
    assert result.applied == 1
    out = Presentation(io.BytesIO(result.cleaned_bytes)).slides[0]
    fixed = next(s for s in out.shapes if str(s.shape_id) == rec["shape_id"])
    assert fixed.text_frame.paragraphs[0].runs[0].font.size.pt == 12


def test_learn_label_case(make_prs):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(8):
        _label(slide, f"HEADER {'ABCDEFGH'[i]}", IN + i * 200000,
               IN + i * 500000)
    for i in range(4):
        _label(slide, f"Chip {'abcd'[i]}", 8 * IN, IN + i * 500000)
    assert ty.learn_label_case(prs) == "upper"   # 8/12 = 67%

    prs2 = make_prs()
    s2 = prs2.slides.add_slide(prs2.slide_layouts[6])
    for i in range(6):
        _label(s2, f"Chip {'abcdef'[i]}", IN, IN + i * 500000)
    for i in range(6):
        _label(s2, f"TAG {'ABCDEF'[i]}", 4 * IN, IN + i * 500000)
    assert ty.learn_label_case(prs2) is None     # 50%: no convention
