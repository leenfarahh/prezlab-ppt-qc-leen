"""Master unification: clone detection, package-level dedup, COM re-apply.

The two-master decks are crafted at package level (no PowerPoint needed), so
the whole clone path is exercised on any host; only the COM re-apply test is
gated on a desktop PowerPoint install.
"""

import io
import zipfile

import pytest
from lxml import etree
from pptx import Presentation

from qc.engine import run_audit
from qc.fixer import apply_fixes, is_fixable
from qc.unify import analyze, com_available, dedup

P = "http://schemas.openxmlformats.org/presentationml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CT_MASTER = ("application/vnd.openxmlformats-officedocument."
             "presentationml.slideMaster+xml")
CT_LAYOUT = ("application/vnd.openxmlformats-officedocument."
             "presentationml.slideLayout+xml")
CT_THEME = "application/vnd.openxmlformats-officedocument.theme+xml"


def _clone_master_deck(src_bytes: bytes, foreign_slides: list[int]) -> bytes:
    """Duplicate master1 (+layouts +theme) as a second master and repoint the
    given zero-based slides to the clone's layouts: the copy-paste pollution
    case, built deterministically."""
    zin = zipfile.ZipFile(io.BytesIO(src_bytes))
    parts = {n: zin.read(n) for n in zin.namelist()}
    zin.close()

    m_rels = etree.fromstring(parts["ppt/slideMasters/_rels/slideMaster1.xml.rels"])
    layout_map = {}  # original layout part -> clone layout part
    next_layout = 1 + max(
        int(n.rsplit("slideLayout", 1)[1].split(".")[0])
        for n in parts if n.startswith("ppt/slideLayouts/slideLayout")
        and n.endswith(".xml") and "_rels" not in n)
    theme_clone = "ppt/theme/theme900.xml"

    for rel in m_rels:
        tgt = rel.get("Target")
        if rel.get("Type").endswith("/slideLayout"):
            orig = "ppt/slideLayouts/" + tgt.rsplit("/", 1)[1]
            clone = f"ppt/slideLayouts/slideLayout{next_layout}.xml"
            next_layout += 1
            layout_map[orig] = clone
            rel.set("Target", "../slideLayouts/" + clone.rsplit("/", 1)[1])
        elif rel.get("Type").endswith("/theme"):
            rel.set("Target", "../theme/" + theme_clone.rsplit("/", 1)[1])

    # clone parts: master xml near-verbatim (its rIds still resolve via the
    # cloned rels, but sldLayoutId ids must be unique document-wide),
    # layouts verbatim with master rel retargeted, theme verbatim
    m2 = etree.fromstring(parts["ppt/slideMasters/slideMaster1.xml"])
    all_ids = [int(e.get("id"))
               for e in m2.iter(f"{{{P}}}sldLayoutId") if e.get("id")]
    pres_tmp = etree.fromstring(parts["ppt/presentation.xml"])
    all_ids += [int(e.get("id"))
                for e in pres_tmp.iter(f"{{{P}}}sldMasterId") if e.get("id")]
    next_id = max(all_ids) + 2  # +1 is reserved for the new sldMasterId below
    for e in m2.iter(f"{{{P}}}sldLayoutId"):
        e.set("id", str(next_id))
        next_id += 1
    parts["ppt/slideMasters/slideMaster2.xml"] = etree.tostring(
        m2, xml_declaration=True, encoding="UTF-8", standalone=True)
    parts["ppt/slideMasters/_rels/slideMaster2.xml.rels"] = etree.tostring(m_rels)
    theme1 = next(n for n in parts if n.startswith("ppt/theme/theme")
                  and "_rels" not in n)
    parts[theme_clone] = parts[theme1]
    for orig, clone in layout_map.items():
        parts[clone] = parts[orig]
        lrels = etree.fromstring(parts[f"ppt/slideLayouts/_rels/{orig.rsplit('/', 1)[1]}.rels"])
        for rel in lrels:
            if rel.get("Type").endswith("/slideMaster"):
                rel.set("Target", "../slideMasters/slideMaster2.xml")
        parts[f"ppt/slideLayouts/_rels/{clone.rsplit('/', 1)[1]}.rels"] = etree.tostring(lrels)

    # register: content types, presentation rels, sldMasterIdLst
    ctypes = etree.fromstring(parts["[Content_Types].xml"])
    for part, ct in ([("/" + "ppt/slideMasters/slideMaster2.xml", CT_MASTER),
                      ("/" + theme_clone, CT_THEME)] +
                     [("/" + c, CT_LAYOUT) for c in layout_map.values()]):
        o = etree.SubElement(ctypes, f"{{{CT_NS}}}Override")
        o.set("PartName", part)
        o.set("ContentType", ct)
    parts["[Content_Types].xml"] = etree.tostring(ctypes, xml_declaration=True,
                                                  encoding="UTF-8", standalone=True)

    pres_rels = etree.fromstring(parts["ppt/_rels/presentation.xml.rels"])
    new_rid = "rId900"
    rel = etree.SubElement(pres_rels, f"{{{REL_NS}}}Relationship")
    rel.set("Id", new_rid)
    rel.set("Type", f"{R}/slideMaster")
    rel.set("Target", "slideMasters/slideMaster2.xml")
    parts["ppt/_rels/presentation.xml.rels"] = etree.tostring(
        pres_rels, xml_declaration=True, encoding="UTF-8", standalone=True)

    pres = etree.fromstring(parts["ppt/presentation.xml"])
    lst = pres.find(f"{{{P}}}sldMasterIdLst")
    ids = [int(e.get("id")) for e in lst] + [
        int(e.get("id"))
        for n in parts if n.startswith("ppt/slideMasters/slideMaster")
        and "_rels" not in n
        for e in etree.fromstring(parts[n]).iter(f"{{{P}}}sldLayoutId")
        if e.get("id")]
    entry = etree.SubElement(lst, f"{{{P}}}sldMasterId")
    entry.set("id", str(max(ids) + 1))
    entry.set(f"{{{R}}}id", new_rid)
    parts["ppt/presentation.xml"] = etree.tostring(
        pres, xml_declaration=True, encoding="UTF-8", standalone=True)

    # repoint the chosen slides to the clone's layouts
    for s_idx in foreign_slides:
        rels_name = f"ppt/slides/_rels/slide{s_idx + 1}.xml.rels"
        srels = etree.fromstring(parts[rels_name])
        for rel in srels:
            if rel.get("Type").endswith("/slideLayout"):
                orig = "ppt/slideLayouts/" + rel.get("Target").rsplit("/", 1)[1]
                rel.set("Target",
                        "../slideLayouts/" + layout_map[orig].rsplit("/", 1)[1])
        parts[rels_name] = etree.tostring(srels, xml_declaration=True,
                                          encoding="UTF-8", standalone=True)

    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        for name, data in parts.items():
            z.writestr(name, data)
    return out.getvalue()


CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def _mutate_clone_theme(deck_bytes: bytes) -> bytes:
    """Turn the clone into a genuinely different master by editing its theme
    (different major font), so structural identity no longer holds."""
    zin = zipfile.ZipFile(io.BytesIO(deck_bytes))
    parts = {n: zin.read(n) for n in zin.namelist()}
    zin.close()
    theme = etree.fromstring(parts["ppt/theme/theme900.xml"])
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    latin = theme.find(f".//{{{A}}}majorFont/{{{A}}}latin")
    latin.set("typeface", "Comic Sans MS")
    parts["ppt/theme/theme900.xml"] = etree.tostring(theme)
    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        for name, data in parts.items():
            z.writestr(name, data)
    return out.getvalue()


@pytest.fixture(scope="module")
def clone_deck(fixtures_dir):
    return _clone_master_deck(
        (fixtures_dir / "mixed_layouts.pptx").read_bytes(), [1, 3])


@pytest.fixture(scope="module")
def nonclone_deck(clone_deck):
    return _mutate_clone_theme(clone_deck)


def test_single_master_deck_has_no_foreign_slides(fixtures_dir):
    a = analyze((fixtures_dir / "mixed_layouts.pptx").read_bytes())
    assert not a.multiple_masters
    assert a.foreign == []


def test_analyze_detects_clone_and_twins(clone_deck):
    a = analyze(clone_deck)
    assert a.multiple_masters
    assert a.dominant == "ppt/slideMasters/slideMaster1.xml"
    assert a.clone_masters == {"ppt/slideMasters/slideMaster2.xml"}
    assert sorted(f.slide_index for f in a.foreign) == [1, 3]
    assert all(f.twin_layout_part for f in a.foreign)
    # the crafted deck really has two masters when opened normally
    assert len(Presentation(io.BytesIO(clone_deck)).slide_masters) == 2


def test_nonclone_master_offers_name_match_not_twin(nonclone_deck):
    a = analyze(nonclone_deck)
    assert a.clone_masters == set()
    assert all(f.twin_layout_part is None for f in a.foreign)
    assert all(f.name_match_layout for f in a.foreign)  # same layout names


def test_audit_emits_foreign_master_records(clone_deck, tmp_path):
    deck = tmp_path / "two_master.pptx"
    deck.write_bytes(clone_deck)
    result = run_audit(deck, "prezlab_en", modules=["master_slide"])
    recs = [r for r in result.records
            if r.issue_type == "master_slide.foreign_master"]
    assert sorted(r.slide_index for r in recs) == [1, 3]
    for r in recs:
        assert r.confidence == "deterministic"
        assert r.locator.startswith("dedup:ppt/slideLayouts/")
        assert is_fixable(r.to_dict())


def test_dedup_fix_end_to_end(clone_deck, tmp_path):
    deck = tmp_path / "two_master.pptx"
    deck.write_bytes(clone_deck)
    result = run_audit(deck, "prezlab_en", modules=["master_slide"])
    records = [r.to_dict() for r in result.records]
    ids = {r["record_id"] for r in records
           if r["issue_type"] == "master_slide.foreign_master"}
    assert ids

    before = Presentation(io.BytesIO(clone_deck))
    titles_before = [s.shapes.title.text if s.shapes.title else None
                     for s in before.slides]

    fix = apply_fixes(clone_deck, records, ids)
    assert all(o.outcome == "changed" for o in fix.outcomes), \
        [(o.record_id, o.reason) for o in fix.outcomes]

    after = Presentation(io.BytesIO(fix.cleaned_bytes))
    assert len(after.slide_masters) == 1
    titles_after = [s.shapes.title.text if s.shapes.title else None
                    for s in after.slides]
    assert titles_after == titles_before  # slide content untouched

    # no orphaned parts left in the package
    with zipfile.ZipFile(io.BytesIO(fix.cleaned_bytes)) as z:
        masters = [n for n in z.namelist()
                   if n.startswith("ppt/slideMasters/") and n.endswith(".xml")
                   and "_rels" not in n]
    assert masters == ["ppt/slideMasters/slideMaster1.xml"]

    # verify-after-write: re-audit is clean of the issue
    cleaned = tmp_path / "cleaned.pptx"
    cleaned.write_bytes(fix.cleaned_bytes)
    again = run_audit(cleaned, "prezlab_en", modules=["master_slide"])
    assert "master_slide.foreign_master" not in again.summary["by_issue_type"]


def test_nonclone_records_are_medium_confidence_com_route(nonclone_deck, tmp_path):
    deck = tmp_path / "nonclone.pptx"
    deck.write_bytes(nonclone_deck)
    result = run_audit(deck, "prezlab_en", modules=["master_slide"])
    recs = [r.to_dict() for r in result.records
            if r.issue_type == "master_slide.foreign_master"]
    assert sorted(r["slide_index"] for r in recs) == [1, 3]
    for r in recs:
        assert r["confidence"] == "medium"
        assert r["locator"].startswith("com:")
        assert is_fixable(r)  # tickable, but never pre-selected until promoted


def _two_slide_arabic_deck() -> bytes:
    """Slide 0 English, slide 1 Arabic: the guard fixture for unification."""
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    s1.shapes.title.text = "English title"
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "عنوان عربي"
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def test_arabic_slides_keep_guard_on_com_route(tmp_path):
    src = _two_slide_arabic_deck()
    # non-clone master under the Arabic slide -> COM route, guard must hold
    deck_bytes = _mutate_clone_theme(_clone_master_deck(src, [1]))
    deck = tmp_path / "ar_nonclone.pptx"
    deck.write_bytes(deck_bytes)
    result = run_audit(deck, "prezlab_bilingual", modules=["master_slide"])
    recs = [r.to_dict() for r in result.records
            if r.issue_type == "master_slide.foreign_master"]
    assert len(recs) == 1 and recs[0]["slide_index"] == 1
    assert recs[0]["arabic_flag"] and recs[0]["locator"].startswith("com:")
    assert not is_fixable(recs[0])

    # but the clone/dedup route stays fixable even on Arabic slides: the
    # slide's own XML is untouched, so there is nothing RTL to break
    deck2 = tmp_path / "ar_clone.pptx"
    deck2.write_bytes(_clone_master_deck(src, [1]))
    result2 = run_audit(deck2, "prezlab_bilingual", modules=["master_slide"])
    recs2 = [r.to_dict() for r in result2.records
             if r.issue_type == "master_slide.foreign_master"]
    assert len(recs2) == 1
    assert not recs2[0]["arabic_flag"] and is_fixable(recs2[0])


def test_dedup_only_removes_emptied_masters(clone_deck):
    a = analyze(clone_deck)
    one = next(f for f in a.foreign if f.slide_index == 1)
    # repoint only ONE of the two foreign slides: the clone master keeps a
    # slide, so it must survive
    out = dedup(clone_deck, {1: one.twin_layout_part})
    assert len(Presentation(io.BytesIO(out)).slide_masters) == 2


@pytest.mark.skipif(not com_available(), reason="desktop PowerPoint required")
def test_com_unify_end_to_end(nonclone_deck, tmp_path):
    deck = tmp_path / "nonclone.pptx"
    deck.write_bytes(nonclone_deck)
    result = run_audit(deck, "prezlab_en", modules=["master_slide"])
    records = [r.to_dict() for r in result.records]
    ids = {r["record_id"] for r in records
           if r["issue_type"] == "master_slide.foreign_master"}

    before = Presentation(io.BytesIO(nonclone_deck))
    titles_before = [s.shapes.title.text if s.shapes.title else None
                     for s in before.slides]

    fix = apply_fixes(nonclone_deck, records, ids)
    changed = [o for o in fix.outcomes if o.outcome == "changed"]
    assert len(changed) == len(ids), \
        [(o.record_id, o.reason) for o in fix.outcomes]

    after = Presentation(io.BytesIO(fix.cleaned_bytes))
    assert len(after.slide_masters) == 1
    titles_after = [s.shapes.title.text if s.shapes.title else None
                    for s in after.slides]
    assert titles_after == titles_before  # PowerPoint preserved the content
