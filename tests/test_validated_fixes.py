"""Validation plumbing: geometry fixes end-to-end, promotion gating, the
local store, identity/comment endpoints, and report pre-check behavior."""

import io

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches

import qc.promotion as promotion
import qc.store
import qc.triage as triage
from qc.engine import run_audit
from qc.fixer import apply_fixes, is_fixable
from qc.records import make_record
from qc.util import iter_shapes_deep
from qc.web import app

BLANK_LAYOUT = 6
IN = 914400  # one inch in EMU
EDGE_TYPE = "margin_alignment.edge_misaligned"
SPACING_TYPE = "margin_alignment.uneven_spacing"


def _box(slide, left, top, width=1000000, height=500000):
    return slide.shapes.add_textbox(left, top, width, height)


def _shape_by_id(prs_bytes: bytes, slide_index: int, shape_id: str):
    prs = Presentation(io.BytesIO(prs_bytes))
    for shape, _path in iter_shapes_deep(prs.slides[slide_index].shapes):
        if str(shape.shape_id) == shape_id:
            return shape
    raise AssertionError(f"shape {shape_id} not found on slide {slide_index}")


def _audit_records(path, profile):
    return run_audit(path, profile, ["margin_alignment"]).to_manifest()["records"]


@pytest.fixture()
def store(monkeypatch, tmp_path):
    """qc.store isolated to tmp_path; never touches the real data/ dir."""
    monkeypatch.setattr(qc.store, "DATA_DIR", tmp_path)
    monkeypatch.setattr(qc.store, "DB_PATH", tmp_path / "qc.db")
    return qc.store


@pytest.fixture()
def triage_tmp(monkeypatch, tmp_path):
    """qc.triage isolated to tmp_path; never touches the real data/ dir."""
    monkeypatch.setattr(triage, "DATA_DIR", tmp_path)
    monkeypatch.setattr(triage, "TRIAGE_LOG", tmp_path / "triage-log.jsonl")
    return triage


# --- geometry fixes end-to-end ------------------------------------------


def test_edge_misaligned_fix_end_to_end(make_prs, en_profile, tmp_path):
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    x = Inches(2)
    # Deviation must stay inside the 3x-tolerance clustering window (28575
    # EMU) to be seen as near-aligned intent, yet exceed the 9525 tolerance.
    _box(slide, left=x, top=IN)
    _box(slide, left=x, top=2 * IN)
    deviant = _box(slide, left=x + 20000, top=3 * IN)
    path = tmp_path / "edge.pptx"
    prs.save(path)

    records = _audit_records(path, en_profile)
    edge = [r for r in records if r["issue_type"] == EDGE_TYPE]
    assert len(edge) == 1
    rec = edge[0]
    assert rec["shape_id"] == str(deviant.shape_id)

    fx = apply_fixes(path.read_bytes(), records, {rec["record_id"]})
    assert fx.applied == 1
    assert [o.outcome for o in fx.outcomes] == ["changed"]

    fixed = _shape_by_id(fx.cleaned_bytes, 0, rec["shape_id"])
    assert fixed.left == int(rec["new_value"]) == x

    clean_path = tmp_path / "edge.cleaned.pptx"
    clean_path.write_bytes(fx.cleaned_bytes)
    after = _audit_records(clean_path, en_profile)
    assert [r for r in after if r["issue_type"] == EDGE_TYPE] == []

    updated = {r["record_id"]: r for r in fx.records}
    assert updated[rec["record_id"]]["action"] == "changed"


def test_uneven_spacing_fix_translates_row_tail(make_prs, en_profile, tmp_path):
    """Same-size rows are claimed by the stronger cluster_rhythm check
    (lift-the-tail semantics, satellites riding); the translate-the-tail
    behavior itself is unchanged."""
    prs = make_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    w, g = 1000000, 500000
    x0 = IN
    x1 = x0 + w + g
    x2 = x1 + w + g + 200000  # the odd gap sits before this shape
    x3 = x2 + w + g
    boxes = [_box(slide, left=x, top=IN, width=w) for x in (x0, x1, x2, x3)]
    path = tmp_path / "row.pptx"
    prs.save(path)

    records = _audit_records(path, en_profile)
    spacing = [r for r in records
               if r["issue_type"] == "margin_alignment.cluster_rhythm"]
    assert len(spacing) == 1
    rec = spacing[0]
    assert rec["shape_id"] == str(boxes[2].shape_id)
    assert rec["locator"] == "lift-row:" + ",".join(
        str(b.shape_id) for b in boxes[2:])
    assert int(rec["old_value"]) == g + 200000
    assert int(rec["new_value"]) == g

    fx = apply_fixes(path.read_bytes(), records, {rec["record_id"]})
    assert fx.applied == 1

    lefts = {str(b.shape_id): _shape_by_id(fx.cleaned_bytes, 0, str(b.shape_id)).left
             for b in boxes}
    ids = [str(b.shape_id) for b in boxes]
    # Head of the row untouched.
    assert lefts[ids[0]] == x0
    assert lefts[ids[1]] == x1
    # The odd gap now equals the median gap...
    assert lefts[ids[2]] - (lefts[ids[1]] + w) == g
    # ...and the FOLLOWING gap is unchanged: the tail moved as one unit.
    assert lefts[ids[3]] - (lefts[ids[2]] + w) == g
    assert lefts[ids[2]] == x2 - 200000
    assert lefts[ids[3]] == x3 - 200000

    clean_path = tmp_path / "row.cleaned.pptx"
    clean_path.write_bytes(fx.cleaned_bytes)
    after = _audit_records(clean_path, en_profile)
    assert [r for r in after if r["issue_type"] == SPACING_TYPE] == []


# --- is_fixable confidence and Arabic gates ------------------------------


def _rec(issue_type, confidence, arabic=False, action="flagged", new_value="1"):
    return {"issue_type": issue_type, "confidence": confidence,
            "arabic_flag": arabic, "action": action, "new_value": new_value}


def test_uneven_spacing_fixable_at_low_confidence():
    assert is_fixable(_rec(SPACING_TYPE, "low")) is True


def test_edge_misaligned_fixable_at_medium_but_not_low():
    assert is_fixable(_rec(EDGE_TYPE, "medium")) is True
    assert is_fixable(_rec(EDGE_TYPE, "low")) is False


def test_arabic_guard_scoped_to_text_touching_fixes():
    """Since 12/08/2026 the Arabic guard is scoped: pure-geometry fixes
    (like an edge snap) are script-neutral and stay fixable; anything that
    edits runs or text keeps the guard (see tests/test_arabic_geometry.py)."""
    assert is_fixable(_rec(EDGE_TYPE, "medium", arabic=True)) is True
    # font substitution is tick-to-approve (fixable, never pre-selected)
    assert is_fixable(_rec("font.family_out_of_set", "deterministic",
                           arabic=True)) is True
    # text-editing fixes stay guarded outright
    assert is_fixable(_rec("header_footer.text_mismatch", "deterministic",
                           arabic=True)) is False


# --- promotion from designer triage --------------------------------------


def _synthetic_record(i: int, issue_type: str = EDGE_TYPE) -> dict:
    return {"record_id": f"synthetic-{i}", "issue_type": issue_type,
            "module": "margin_alignment", "severity": "warning",
            "confidence": "medium", "arabic_flag": False,
            "message": "synthetic triage entry"}


def test_promotion_after_enough_confirmed_reviews(monkeypatch, triage_tmp):
    monkeypatch.setattr(promotion, "MIN_REVIEWED", 3)
    for i in range(3):
        triage_tmp.log_triage(_synthetic_record(i), "confirmed",
                              deck="d.pptx", profile_id="prezlab_en")
    assert EDGE_TYPE in promotion.promoted_issue_types()
    status = promotion.promotion_status(EDGE_TYPE)
    assert status["promoted"] is True
    assert status["reviewed"] == 3
    assert status["needed"] == 0


def test_promotion_blocked_by_false_positive_rate(monkeypatch, triage_tmp):
    monkeypatch.setattr(promotion, "MIN_REVIEWED", 3)
    for i in range(3):
        triage_tmp.log_triage(_synthetic_record(i), "confirmed",
                              deck="d.pptx", profile_id="prezlab_en")
    # One false alarm on a fourth record: fp_rate 0.25 > MAX_FP_RATE 0.05.
    triage_tmp.log_triage(_synthetic_record(99), "false_positive",
                          deck="d.pptx", profile_id="prezlab_en")
    assert EDGE_TYPE not in promotion.promoted_issue_types()
    status = promotion.promotion_status(EDGE_TYPE)
    assert status["promoted"] is False
    assert status["fp_rate"] > promotion.MAX_FP_RATE


def test_promotion_status_unseen_type_needs_full_quota(triage_tmp):
    status = promotion.promotion_status("nonexistent.check")
    assert status == {"reviewed": 0, "fp_rate": 0.0, "promoted": False,
                      "needed": promotion.MIN_REVIEWED}
    assert status["needed"] > 0


# --- store: users and comments --------------------------------------------


def test_store_users_roundtrip_and_invalid_role(store):
    user = store.add_user("Rana", "designer")
    assert user["name"] == "Rana" and user["role"] == "designer"
    assert [u["name"] for u in store.list_users()] == ["Rana"]
    assert store.get_user("Rana")["role"] == "designer"
    assert store.get_user("Nobody") is None
    with pytest.raises(ValueError):
        store.add_user("Sami", "intern")


def test_store_comments_roundtrip_and_counts(store):
    store.add_user("Rana", "designer")
    c1 = store.add_comment("deck.pptx", 0, "Rana", "check the header")
    store.add_comment("deck.pptx", 0, "Rana", "and the footer",
                      record_id="rec-1")
    store.add_comment("deck.pptx", 2, "Rana", "slide three note")
    store.add_comment("other.pptx", 0, "Rana", "different deck")

    slide0 = store.comments_for("deck.pptx", 0)
    assert [c["text"] for c in slide0] == ["check the header", "and the footer"]
    assert slide0[0]["id"] == c1["id"]
    assert slide0[1]["record_id"] == "rec-1"
    assert len(store.comments_for("deck.pptx")) == 3
    assert store.comment_counts("deck.pptx") == {0: 2, 2: 1}
    with pytest.raises(ValueError):
        store.add_comment("deck.pptx", 0, "Rana", "   ")


def test_store_delete_comment_permissions(store):
    store.add_user("Rana", "designer")
    store.add_user("Sami", "designer")
    store.add_user("Lubna", "lead")
    comment = store.add_comment("deck.pptx", 0, "Rana", "mine")

    # Another designer cannot delete someone else's comment.
    assert store.delete_comment(comment["id"], "Sami") is False
    assert len(store.comments_for("deck.pptx", 0)) == 1
    # The author can.
    assert store.delete_comment(comment["id"], "Rana") is True
    assert store.comments_for("deck.pptx", 0) == []
    # A lead can delete anyone's.
    other = store.add_comment("deck.pptx", 0, "Sami", "theirs")
    assert store.delete_comment(other["id"], "Lubna") is True
    assert store.comments_for("deck.pptx", 0) == []


# --- identity and comment endpoints ---------------------------------------


def test_identity_endpoints_cookie_flow(store):
    client = TestClient(app)
    body = client.get("/me").json()
    assert body["user"] is None
    assert body["users"] == []

    assert client.post("/whoami", json={"name": "Ghost"}).status_code == 404

    store.add_user("Rana", "designer")
    resp = client.post("/whoami", json={"name": "Rana"})
    assert resp.status_code == 200
    assert resp.cookies.get("qc_user") == "Rana"
    # The client carries the cookie forward; /me now resolves the user.
    body = client.get("/me").json()
    assert body["user"]["name"] == "Rana"
    assert body["user"]["role"] == "designer"


def test_comments_endpoints_require_identity(store):
    client = TestClient(app)
    payload = {"deck": "deck.pptx", "slide_index": 0, "text": "tighten this"}
    assert client.post("/comments", json=payload).status_code == 401

    store.add_user("Rana", "designer")
    client.post("/whoami", json={"name": "Rana"})
    resp = client.post("/comments", json=payload)
    assert resp.status_code == 200
    assert resp.json()["comment"]["author"] == "Rana"

    got = client.get("/comments", params={"deck": "deck.pptx", "slide": 0}).json()
    assert [c["text"] for c in got["comments"]] == ["tighten this"]
    other = client.get("/comments", params={"deck": "x.pptx", "slide": 0}).json()
    assert other["comments"] == []


# --- report pre-check honors promotion -------------------------------------


def test_report_precheck_honors_promotion():
    from qc.ui import render_report

    rec = make_record(
        slide_index=0, shape_id="7", module="margin_alignment",
        issue_type=EDGE_TYPE, severity="warning", action="flagged",
        confidence="medium", property="spPr.xfrm.off",
        old_value=1848800, new_value=1828800,
        profile_rule_id="geometry.alignment.edge_tolerance_emu",
        arabic_flag=False,
        message="left edge off cluster median by 20000 EMU",
    ).to_dict()
    manifest = {
        "deck": "demo.pptx", "profile_id": "prezlab_en", "profile_version": 1,
        "slides": 1,
        "summary": {"by_severity": {"warning": 1},
                    "by_issue_type": {EDGE_TYPE: 1},
                    "by_module": {"margin_alignment": 1},
                    "arabic_flagged": 0, "total": 1},
        "records": [rec],
    }
    needle = f'value="{rec["record_id"]}" form="applyform"'

    html = render_report(manifest, "job1", can_fix=True, promoted=set())
    assert needle in html  # the checkbox is offered as a suggestion
    assert needle + " checked" not in html

    html = render_report(manifest, "job1", can_fix=True, promoted={EDGE_TYPE})
    assert needle + " checked" in html


def test_report_prechecks_confident_errors_without_promotion():
    """Evidence-escalated errors (confidently wrong + safe fix) come
    pre-ticked; the tick is reserved for ambiguous warnings."""
    from qc.ui import render_report

    rec = make_record(
        slide_index=0, shape_id="7", module="margin_alignment",
        issue_type=EDGE_TYPE, severity="error", action="flagged",
        confidence="high", property="spPr.xfrm.off",
        old_value=1888800, new_value=1828800,
        profile_rule_id="geometry.alignment.edge_tolerance_emu",
        arabic_flag=False,
        message="left edge off cluster median by 60000 EMU",
    ).to_dict()
    manifest = {
        "deck": "demo.pptx", "profile_id": "prezlab_en", "profile_version": 1,
        "slides": 1,
        "summary": {"by_severity": {"error": 1},
                    "by_issue_type": {EDGE_TYPE: 1},
                    "by_module": {"margin_alignment": 1},
                    "arabic_flagged": 0, "total": 1},
        "records": [rec],
    }
    needle = f'value="{rec["record_id"]}" form="applyform"'
    html = render_report(manifest, "job1", can_fix=True, promoted=set())
    assert needle + " checked" in html
