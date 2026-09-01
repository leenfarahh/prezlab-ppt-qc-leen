"""The whole vision path, end to end: model -> record -> page -> deck.

"The model detects it but the app neither shows it nor fixes it" (design lead,
01/09/2026). Four separate faults on one path, none of them in the model and
none of them in the prompt. Each is protected here.

1. A GRAZE WAS TREATED AS A WELD. Any overlap at all made two shapes one
   object, and text boxes overlap constantly without their text touching: a
   full-width title's box lapped 0.3in over the header below it. Asked to nudge
   that header down, the fixer carried the title too, the title landed on the
   next column's header, the collision guard correctly refused - and the page
   said "Applied 0 fixes" (qc.util.rides_with).

2. THE REASON WAS THROWN AWAY. "Applied 0 fixes. Skipped 1." The sentence the
   fixer had already written - which shape it would have collided with - was
   counted and discarded (qc.web._skip_note).

3. THE FIRST FIX DESTROYED EVERY SUGGESTION. Applying one unrelated font fix
   re-audits the deck and replaces the manifest, and the re-audit only knows
   the measured modules, so every copilot and component record on the page
   vanished (qc.web._surviving_vision).

4. A DAILY QUOTA WAS REPORTED AS A RATE LIMIT. "Run the pass again shortly"
   cannot work when the day's twenty requests are spent, and it sent people
   looking at their network (qc.llm._quota_message).

The model is stubbed throughout. What is being tested is everything that
happens to its answer AFTER it arrives.
"""

import io

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Emu

import qc.copilot as copilot
import qc.llm as llm
import qc.web as web
from qc.web import app
from tests.conftest import job_id_of

client = TestClient(app)
IN = 914400


def _two_column_slide():
    """A title whose BOX laps over the left column's header, which is what an
    ordinary text box looks like: the box is taller than the words in it."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Emu(int(0.9 * IN)), Emu(int(0.4 * IN)),
                                     Emu(int(11.5 * IN)), Emu(int(1.45 * IN)))
    title.text_frame.text = "Two ways to do it"
    # left header: 0.3in higher than the right one, and its box starts inside
    # the title's box
    left = slide.shapes.add_textbox(Emu(int(0.9 * IN)), Emu(int(1.54 * IN)),
                                    Emu(int(5.0 * IN)), Emu(int(0.9 * IN)))
    left.text_frame.text = "Method 1"
    # Wider than the left one on purpose: _shared_edge makes the LARGER shape
    # of a pair the spine, so the line is the right header's top and the shape
    # that moves is the left one - the header the title's box grazes, which is
    # the case being protected.
    right = slide.shapes.add_textbox(Emu(int(6.75 * IN)), Emu(int(1.84 * IN)),
                                     Emu(int(5.1 * IN)), Emu(int(0.9 * IN)))
    right.text_frame.text = "Method 2"
    body = slide.shapes.add_textbox(Emu(int(0.9 * IN)), Emu(int(2.74 * IN)),
                                    Emu(int(5.0 * IN)), Emu(int(3.0 * IN)))
    body.text_frame.text = "1 2 3"
    buf = io.BytesIO()
    prs.save(buf)
    return slide, title, left, right, buf.getvalue()


def _observation(a, b):
    return [{"action": "align_top",
             "shape_ids": [str(a.shape_id), str(b.shape_id)],
             "rationale": "The two column headers should share a top edge."}]


def test_a_title_whose_box_grazes_a_header_is_not_carried_by_it():
    """The specific fault: the fix moved the header AND the title, the title
    collided, and the whole thing was rolled back."""
    from qc.fixer import _carried_contents, apply_fixes

    slide, title, left, right, data = _two_column_slide()
    recs = copilot.synthesize(slide, 0, _observation(left, right), existing=[])
    assert len(recs) == 1 and str(recs[0]["shape_id"]) == str(left.shape_id)

    carried = [str(c.shape_id) for c in _carried_contents(slide, left, "y")]
    assert str(title.shape_id) not in carried, \
        "the title is a neighbour of the header, not its passenger"

    result = apply_fixes(data, recs, {recs[0]["record_id"]})
    assert result.applied == 1, [o.reason for o in result.outcomes]

    after = {str(s.shape_id): s.top
             for s in Presentation(io.BytesIO(result.cleaned_bytes)).slides[0].shapes}
    assert after[str(left.shape_id)] == right.top, "the header did not land"
    assert after[str(right.shape_id)] == right.top, "the spine moved"
    assert after[str(title.shape_id)] == title.top, "the title was dragged"


def test_a_label_lying_on_its_card_is_still_welded_to_it():
    """The weld rule still has to do its job: a shape sitting ON another is one
    object with it whatever axis the move is on."""
    from qc.fixer import _carried_contents

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    card = slide.shapes.add_textbox(Emu(IN), Emu(IN), Emu(3 * IN), Emu(2 * IN))
    card.text_frame.text = "card"
    label = slide.shapes.add_textbox(Emu(int(1.2 * IN)), Emu(int(1.2 * IN)),
                                     Emu(int(1.0 * IN)), Emu(int(0.4 * IN)))
    label.text_frame.text = "label"

    carried = [str(c.shape_id) for c in _carried_contents(slide, card, "y")]
    assert str(label.shape_id) in carried


def _job_with_a_suggestion(monkeypatch, deck_bytes, obs):
    """Audit a deck, then run the copilot over it with the model stubbed.

    Through monkeypatch, never by assignment: these are module globals shared
    by the whole session, and a hand-assigned _ensure_thumbs that hands out
    b"png" outlives this file and breaks the first later test that renders a
    real image (which is exactly what it did - test_web's PDF route died in
    PIL, three files away).
    """
    r = client.post("/audit", files={"deck": ("d.pptx", deck_bytes,
                                             "application/octet-stream")},
                    data={"profile": "prezlab_en"})
    job_id = job_id_of(r)
    monkeypatch.setattr(copilot, "ask_json", lambda **kw: {"observations": obs})
    monkeypatch.setattr(llm, "api_configured", lambda: True)
    monkeypatch.setattr(web, "_ensure_thumbs", lambda jid, job: job.__setitem__(
        "thumbs", {i: b"png" for i in range(job["manifest"]["slides"])}))
    client.post(f"/copilot/{job_id}")
    return job_id


def _vision(job_id):
    return [r for r in web._job(job_id)["manifest"]["records"]
            if r.get("source") == "vision"]


def test_a_suggestion_survives_a_fix_applied_to_something_else(monkeypatch):
    _slide, _t, left, right, data = _two_column_slide()
    job_id = _job_with_a_suggestion(monkeypatch, data, _observation(left, right))
    assert len(_vision(job_id)) == 1, "the copilot's record never reached the page"

    other = next((r for r in web._job(job_id)["manifest"]["records"]
                  if r.get("source") != "vision"
                  and r["issue_type"] == "font.family_out_of_set"), None)
    if other is None:                     # nothing unrelated to apply
        return
    client.post("/apply", data={"job_id": job_id,
                                "record_ids": [other["record_id"]]})

    assert len(_vision(job_id)) == 1, \
        "an unrelated fix wiped the model's suggestion off the page"


def test_the_page_says_why_a_fix_was_skipped():
    """A skip with a reason nobody prints is the tool doing nothing and
    declining to say what would have to change."""
    from qc.fixer import FixOutcome

    note = web._skip_note([
        FixOutcome("a", "skipped", "snapping would push shape 4 into shape 7"),
        FixOutcome("b", "skipped", "snapping would push shape 4 into shape 7"),
        FixOutcome("c", "skipped", "shape not found")])
    assert "Skipped 3" in note
    assert "shape 4 into shape 7" in note
    assert "shape not found" in note
    assert note.count("shape 4 into shape 7") == 1, "reasons are grouped"
    assert web._skip_note([]) == ""


def test_a_daily_quota_is_not_reported_as_a_rate_limit():
    """The free tier allows twenty requests A DAY for a flash model. Told to
    "run the pass again shortly", a designer retries all afternoon against a
    quota that only resets at midnight."""
    from qc.llm import _quota_message

    body = ("429 RESOURCE_EXHAUSTED. {'error': {'code': 429, 'message': "
            "'You exceeded your current quota. * Quota exceeded for metric: "
            "generativelanguage.googleapis.com/generate_content_free_tier_"
            "requests, limit: 20, model: gemini-3.5-flash. Please retry in "
            "25.1s.', 'status': 'RESOURCE_EXHAUSTED', 'details': [{'@type': "
            "'type.googleapis.com/google.rpc.QuotaFailure', 'violations': "
            "[{'quotaId': 'GenerateRequestsPerDayPerProjectPerModel-FreeTier',"
            " 'quotaValue': '20'}]}]}}")
    said = _quota_message(RuntimeError(body), 429)

    assert "DAILY" in said and "20 requests" in said
    assert "again shortly" not in said, "that advice cannot work on a day quota"
    assert "midnight" in said

    minute = ("429 RESOURCE_EXHAUSTED. {'error': {'message': 'Quota exceeded "
              "for metric: generate_content_free_tier_requests, limit: 10, "
              "model: gemini-3.5-flash. Please retry in 4s.', 'details': "
              "[{'quotaId': 'GenerateRequestsPerMinutePerProjectPerModel'}]}}")
    said = _quota_message(RuntimeError(minute), 429)
    assert "again shortly" in said and "DAILY" not in said


def test_a_stale_suggestion_is_dropped_rather_than_carried():
    """A vision record that was measured against a shape which has since moved
    names a target read off geometry that no longer exists."""
    from qc.fixer import still_open

    slide, _title, left, right, _data = _two_column_slide()
    rec = copilot.synthesize(slide, 0, _observation(left, right),
                             existing=[])[0]

    assert still_open(rec, slide) is True
    left.top = left.top + int(0.5 * IN)          # somebody else moved it
    assert still_open(rec, slide) is False

    left.top = int(rec["new_value"])             # already where it belongs
    assert still_open(rec, slide) is False
